#!/usr/bin/env python3
"""images.py - harvest, compress and base64-embed property imagery.

Hero images come from the brochure pages (embedded rasters). Each image is
resized and JPEG-compressed to a per-image byte budget so the assembled
single-file dashboard stays a sane size (the reference was 11 MB; we target
~80-120 KB/image). Returns data: URIs ready for the canonical 'photo'/'plan'
fields. EMF/WMF vectors that Pillow cannot read are skipped (caller falls back
to a page raster or the placeholder).
"""
from __future__ import annotations

import base64
import io
import json
import math
import re
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
try:
    import fitz  # PyMuPDF
except Exception:  # sandbox without PyMuPDF: pypdfium2/pdfplumber shim (same 9-call surface)
    import fitz_shim as fitz
try:
    fitz.TOOLS.mupdf_display_errors(False)  # silence MuPDF C-level warnings (broker quiet mode)
except Exception:
    pass
try:
    from PIL import Image
    _HAS_PIL = True
except ImportError:  # Cowork sandbox without Pillow and no pip: degrade, never crash at import.
    # Pillow has no pure-python shim - but the skill MAY ship a bundled Pillow wheel; try
    # it (a NO-OP unless it matches this interpreter). If still absent, the whole image
    # ladder no-ops to the pre-baked placeholder asset. `from __future__ import annotations`
    # (top of file) keeps every `Image.Image` type hint a lazy string, so rebinding is safe.
    try:
        import _vendor_wheels as _vw
        _vw.ensure("PIL", "pillow")
        from PIL import Image
        _HAS_PIL = True
    except Exception:
        Image = None  # type: ignore[assignment]
        _HAS_PIL = False

if _HAS_PIL:
    try:
        import pillow_heif  # noqa: registers HEIC opener
        pillow_heif.register_heif_opener()
    except Exception:
        pass

HERO_MAX_EDGE = 1280
PLAN_MAX_EDGE = 1100
DEFAULT_BUDGET_KB = 110
GALLERY_MAX = 6  # max photos attached per property for the carousel (hero + up to 5 more),
#                  filled BEST-FIRST by photographic_score; the rest are noted in the Gaps Report.
#                  The single self-contained HTML embeds every image as base64, so this cap keeps
#                  the file portable (shareable/emailable) even on an image-heavy deck.
MIN_HERO_W, MIN_HERO_H = 320, 200  # reject logos/icons (PHOTO hero floor)
# plans get their OWN, lower floor: a usable site plan is often small, and the
# photo floor silently discarded it before any scoring happened (a real
# placeholder-shipped-despite-usable-plan failure). Below even this = icon.
MIN_PLAN_W, MIN_PLAN_H = 220, 160
# score at/above which an embedded image (or a cropped image REGION) is a real,
# if unspectacular, photo and ALWAYS beats the whole-page raster. With the
# flatness multiplier (see photographic_score), real photos measure 9-27 and
# plans/maps/logos 0-13 on a real Spanish deck (TEDi calibration); the multiplier
# (not this floor) is what reranks photo-over-plan when both are present.
# Without this floor, a busy page could out-score its own photo and ship a
# cluttered full-page tile as the hero.
MODEST_PHOTO = 6.0


def to_data_uri(jpeg_bytes: bytes, mime: str = "image/jpeg") -> str:
    return f"data:{mime};base64," + base64.b64encode(jpeg_bytes).decode("ascii")


def compress(img: Image.Image, max_edge: int = HERO_MAX_EDGE,
             budget_kb: int = DEFAULT_BUDGET_KB) -> bytes:
    """Resize to max_edge and step quality down until under budget."""
    img = img.convert("RGB")
    w, h = img.size
    scale = min(1.0, max_edge / max(w, h))
    if scale < 1.0:
        img = img.resize((max(1, int(w * scale)), max(1, int(h * scale))), Image.LANCZOS)
    for q in (78, 70, 62, 54, 46, 38):
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=q, optimize=True)
        data = buf.getvalue()
        if len(data) <= budget_kb * 1024 or q == 38:
            return data
    return data  # type: ignore[return-value]


def _open(image_bytes: bytes) -> Image.Image | None:
    if Image is None:
        return None  # no Pillow: no decoder; callers treat None as 'skip this candidate'
    try:
        return Image.open(io.BytesIO(image_bytes))
    except Exception:
        return None


def page_embedded_images(doc: "fitz.Document", page_index: int) -> list[dict]:
    """All decodable raster images on a page, largest first."""
    out = []
    page = doc[page_index]
    for info in page.get_images(full=True):
        xref = info[0]
        try:
            ext = doc.extract_image(xref)
            img = _open(ext["image"])
            if img is None:
                continue
            w, h = img.size
            out.append({"img": img, "w": w, "h": h, "area": w * h, "xref": xref})
        except Exception:
            continue
    out.sort(key=lambda d: -d["area"])
    return out


def candidates_for_page(path: Path, page_index: int) -> list[dict]:
    """The page's HERO-SIZE embedded image candidates, in a STABLE 0-based order -
    the ONE list both sides of the LLM-hero contract share. interpret_prep writes a
    thumbnail per entry (the sub-agent LOOKS at them and references the `index`); merge
    re-derives the SAME list (embedded_by_index) to bind the chosen `index`. Both call
    THIS function on the same filtered order, so the index can never disagree.

    The order is page_embedded_images() (largest-area first) FILTERED to the hero size
    floor (>= MIN_HERO_W x MIN_HERO_H, which screens out logos/icons), and the 0-based
    POSITION in that filtered list IS the candidate's `index`. A PDF page reads its
    embedded rasters; a .pptx slide reads its slide pictures; any other kind (or an open
    failure) yields [] so the caller falls back to the deterministic ladder gracefully.

    Returns [{index, img, w, h}] - `img` is a live PIL image (callers that only need the
    metadata, e.g. a manifest, ignore it). Never raises."""
    try:
        path = Path(path)
        if path.suffix.lower() == ".pptx":
            raw = slide_pictures(path, page_index)
        else:
            raw = page_embedded_images(_get_doc(path), page_index)
    except Exception:
        return []
    out: list[dict] = []
    for im in raw:
        if im.get("w", 0) >= MIN_HERO_W and im.get("h", 0) >= MIN_HERO_H:
            out.append({"index": len(out), "img": im["img"], "w": im["w"], "h": im["h"]})
    return out


def embedded_by_index(path: Path, page_index: int, index: int,
                      budget_kb: int = DEFAULT_BUDGET_KB) -> str | None:
    """The candidate at `index` of candidates_for_page(path, page_index), compressed to a
    hero data URI - the binding half of the LLM-hero contract (merge passes the sub-agent's
    chosen __meta.heroRef here). None when the index is out of range, Pillow is absent, or
    anything fails (the caller then falls back to the deterministic ladder). Never raises."""
    if Image is None:
        return None  # no Pillow: no decoder; merge falls back to the deterministic ladder
    try:
        cands = candidates_for_page(path, page_index)
        if not (isinstance(index, int) and 0 <= index < len(cands)):
            return None
        return to_data_uri(compress(cands[index]["img"], HERO_MAX_EDGE, budget_kb))
    except Exception:
        return None


def page_raster(doc: "fitz.Document", page_index: int, dpi: int = 150) -> Image.Image:
    """Render a whole page to a raster (EMF/vector fallback)."""
    if Image is None:
        return None  # no Pillow: no raster; callers None-check or wrap in try/except
    pix = doc[page_index].get_pixmap(dpi=dpi)
    return Image.open(io.BytesIO(pix.tobytes("png")))


def photographic_score(img) -> float:
    """A ~0..60 'is this a real photograph?' score from PURE PIXEL STATISTICS, used
    to pick a property hero over a floor plan, logo, branded divider or map. Nothing
    about any client, brand, colour, language or region is encoded - so it works the
    same on a Spanish, German, Polish or UK deck. Combines:
      * colourfulness (Hasler-Susstrunk) - photos are colourful; line art/logos/plans are not,
      * non-white fraction - floor plans and paper docs are mostly white,
      * an aspect-ratio sanity penalty - logos/banners are extreme, photos are not,
      * a FLATNESS multiplier - the share of pixels in the 5 dominant quantised
        colours. A COLOURFUL site plan (blue floor fill + green landscaping) beat a
        real aerial photo on colourfulness alone in a real run; but plans/maps are
        dominated by flat fills (flat5 0.55-0.95 measured) while photos spread
        across shades (0.19-0.37), so (1 - flat5) reranks photo-over-plan without
        rejecting 3D renders (~0.7) when they are the page's only marketing image,
      * a SINGLE-DOMINANT-COLOUR penalty - the share of the ONE most common colour: a
        logo, a solid 'photo pending' holding card, a road-map screenshot and a
        text-on-white page are each dominated by one colour, while a photo, a render and
        a satellite/aerial are not, so it demotes the non-property imagery the harvester
        used to mistake for a hero."""
    try:
        im = img.convert("RGB")
    except Exception:
        return 0.0
    w0, h0 = im.size
    if w0 < 2 or h0 < 2:
        return 0.0
    im.thumbnail((96, 96))  # small sample - the stats are scale-stable and this keeps it fast
    px = list(im.getdata())
    n = len(px) or 1
    s_rg = s_yb = s_rg2 = s_yb2 = 0.0
    white = 0
    from collections import Counter
    bins: Counter = Counter()
    for r, g, b in px:
        a = r - g
        c = 0.5 * (r + g) - b
        s_rg += a; s_yb += c; s_rg2 += a * a; s_yb2 += c * c
        if r >= 235 and g >= 235 and b >= 235:
            white += 1
        bins[(r >> 4, g >> 4, b >> 4)] += 1
    m_rg, m_yb = s_rg / n, s_yb / n
    std = math.sqrt(max(0.0, s_rg2 / n - m_rg * m_rg) + max(0.0, s_yb2 / n - m_yb * m_yb))
    mean = math.sqrt(m_rg * m_rg + m_yb * m_yb)
    colourfulness = std + 0.3 * mean
    white_frac = white / n
    flat5 = sum(c for _, c in bins.most_common(5)) / n
    flat1 = bins.most_common(1)[0][1] / n if bins else 0.0
    # SINGLE-DOMINANT-COLOUR penalty (added for the live-run feedback): a logo, a solid
    # 'photo pending' holding card, a ROAD-map screenshot (dominant map-paper colour) and a
    # text page on white all have ONE quantised colour over a large share of the image; a
    # real photo, a 3D render and a SATELLITE/aerial never do. So this demotes exactly the
    # non-property imagery Tier A used to pick, WITHOUT touching a genuine hero. Calibrated
    # conservatively - no penalty below 0.45 (a photo with a big sky or wall is safe), full
    # demotion by 0.80 (logos/solids) - and it is pure pixel statistics, so language/client agnostic.
    single_pen = 1.0 if flat1 <= 0.45 else max(0.0, 1.0 - (flat1 - 0.45) / 0.35)
    ratio = w0 / max(1, h0)
    aspect_pen = 1.0 if 0.45 <= ratio <= 2.4 else (0.6 if 0.3 <= ratio <= 3.5 else 0.3)
    return colourfulness * (1.0 - white_frac) * aspect_pen * (1.0 - flat5) * single_pen


# --- hero-kind classifier (the hero ladder + G-images gate) ------------------- #
# A card's hero must be the REAL photographic content on the page - a warehouse/site
# PHOTO, an AERIAL, or a 3-D RENDER - never a road MAP, a flat PLAN diagram, or a
# whole-slide screenshot/text raster. photographic_score (colourfulness) alone mis-ranks
# these on real CBRE decks: a Google-Maps screenshot scored 18.85 and a real warehouse
# photo 3.66. The discriminator is CONTINUOUS TONE, not colour: a photo/aerial/render
# spreads luminance smoothly (high luminance entropy, few flat-fill blocks); a map, plan,
# icon or text raster is piecewise-flat (low entropy and/or large uniform blocks) or a
# single dominant colour. Pure pixel statistics - client/brand/language-agnostic.
# Thresholds calibrated on the real TEDi ES decks (photos/aerials/renders lumEnt >= ~6.6,
# site plan ~6.06, road map ~4.3, icons <= ~4).
PHOTO_LUM_ENT = 6.3     # luminance entropy (bits) at/above which an image is continuous-tone
PHOTO_FLAT_MAX = 0.18   # max share of near-uniform 8x8 luminance blocks for a photo
MAP_PALE_MIN = 0.40     # share of pale low-saturation non-white pixels = road-map palette
LOGO_FLAT1_MAX = 0.55   # single dominant quantised-colour share = logo / solid / holding card
TEXT_WHITE_MIN = 0.55   # mostly-white + low tonal variety = a text/doc raster


def _hero_signals(img) -> dict:
    """Cheap continuous-tone statistics on a 96px thumbnail for classify_image."""
    im = img.convert("RGB")
    im.thumbnail((96, 96))
    w, h = im.size
    n = w * h or 1
    from collections import Counter
    bins: Counter = Counter()
    white = pale = 0
    for r, g, b in im.getdata():
        bins[(r >> 4, g >> 4, b >> 4)] += 1
        if r >= 235 and g >= 235 and b >= 235:
            white += 1
        else:
            mx, mn = max(r, g, b), min(r, g, b)
            if mx >= 150 and (mx - mn) <= 40:  # bright + desaturated = map land/road paper
                pale += 1
    L = im.convert("L")
    hist = L.histogram()
    lum_ent = -sum((c / n) * math.log2(c / n) for c in hist if c)
    # near-uniform 8x8 luminance blocks on a 48x48 grid (flat fills of plans/maps/icons)
    Ls = list(L.resize((48, 48)).getdata())
    flat = tot = 0
    for by in range(0, 48, 8):
        for bx in range(0, 48, 8):
            blk = [Ls[(by + dy) * 48 + (bx + dx)] for dy in range(8) for dx in range(8)]
            tot += 1
            if max(blk) - min(blk) <= 12:
                flat += 1
    return {"white": white / n, "pale": pale / n, "lum_ent": lum_ent,
            "flat": flat / (tot or 1),
            "flat1": (bins.most_common(1)[0][1] / n if bins else 0.0)}


def classify_image(img, sig: dict | None = None) -> str:
    """Coarse hero-relevant kind from pure pixel statistics:
      'photo' - a real photograph, aerial or 3-D render (continuous tone): the ONLY kind
                that should lead a card. A near-monochrome warehouse interior, a satellite
                aerial and a marketing render all qualify; colour is NOT required.
      'plan'  - a flat-fill site/floor diagram (line art, saturated fills).
      'map'   - a road-map screenshot (pale desaturated land/road palette).
      'text'  - a text/doc raster (mostly white, little tonal variety) - incl. a slide screenshot.
      'logo'  - a solid / holding card / single-colour mark.
    Never raises (a stats failure returns 'photo' so the harvest is never blocked)."""
    if Image is None:
        return "photo"
    try:
        s = sig or _hero_signals(img)
    except Exception:
        return "photo"
    if s["flat1"] >= LOGO_FLAT1_MAX:
        return "logo"
    if s["lum_ent"] >= PHOTO_LUM_ENT and s["flat"] <= PHOTO_FLAT_MAX:
        return "photo"
    if s["white"] >= TEXT_WHITE_MIN and s["lum_ent"] < 5.0:
        return "text"
    if s["pale"] >= MAP_PALE_MIN:
        return "map"
    return "plan"


# hero ladder: a photo/aerial/render leads; a plan beats a map beats text/logo; only drop
# a tier when nothing higher exists on the property's pages (plans + maps still go in the
# gallery + the Site Plan toggle, they just stop being the first impression).
HERO_TIER = {"photo": 0, "plan": 2, "map": 3, "text": 4, "logo": 5}


def is_photo_kind(img) -> bool:
    """True when an image is a real photo/aerial/render (the only valid silent hero)."""
    return classify_image(img) == "photo"


def classify_data_uri(uri: str) -> str:
    """classify_image for a 'data:image/...;base64,...' hero URI (used by the G-images
    gate to BLOCK a map/plan/screenshot hero). Returns 'photo' on any decode/stats
    failure - a gate must never crash or block on an unreadable URI."""
    if Image is None or not (isinstance(uri, str) and "base64," in uri):
        return "photo"
    try:
        import base64
        raw = base64.b64decode(uri.split("base64,", 1)[1])
        return classify_image(Image.open(io.BytesIO(raw)))
    except Exception:
        return "photo"


_DOC_CACHE: dict[str, "fitz.Document"] = {}


def _get_doc(pdf_path: Path) -> "fitz.Document":
    """Open a brochure PDF once and reuse the handle. Image harvesting hits the
    same PDF once per property, so caching avoids reopening/reparsing it N times
    in the merge loop. Call close_doc_cache() when a merge run is done."""
    key = str(Path(pdf_path).resolve())
    doc = _DOC_CACHE.get(key)
    if doc is None:
        doc = fitz.open(pdf_path)
        _DOC_CACHE[key] = doc
    return doc


def close_doc_cache() -> None:
    """Close and forget every cached PDF handle (call at the end of merge), and
    drop the per-document geometry/crop memos that go with them."""
    for doc in _DOC_CACHE.values():
        try:
            doc.close()
        except Exception:
            pass
    _DOC_CACHE.clear()
    try:
        _PLACED_CACHE.clear()
        _CROPS_CACHE.clear()
        _PPTX_CACHE.clear()
    except Exception:
        pass


def _cache_file(pdf_path, page_index, budget_kb, kind, cache_dir, ext=".uri"):
    """Cache path for a (deck, page, budget, kind) unit. ext is `.uri` for the visual
    hero/plan/gallery data URIs and `.json` for intermediate per-page geometry/photo
    caches (so the two never collide in a `*.uri` count and stay self-describing)."""
    if not cache_dir:
        return None
    try:
        import hashlib
        st = Path(pdf_path).stat()
        key = hashlib.sha1(f"v2|{Path(pdf_path).name}|{st.st_size}|{st.st_mtime_ns}|"
                           f"{page_index}|{budget_kb}|{kind}".encode()).hexdigest()
        cdir = Path(cache_dir)
        cdir.mkdir(parents=True, exist_ok=True)
        return cdir / f"{key}{ext}"
    except Exception:
        return None  # cache trouble must never break the harvest


def _cache_read(cf):
    """data URI, '' for a cached negative, or None when there is no cache entry.
    A TRUNCATED entry (a kill mid-write before the atomic rename below existed) must
    not be served: base64 from compress() is always %4==0, so a payload failing that
    is incomplete -> treat as a miss and recompute."""
    if cf is None or not cf.exists():
        return None
    try:
        v = cf.read_text(encoding="ascii")
        if v == "NONE":
            return ""
        if v.startswith("data:image/") and "base64," in v:
            b64 = v.split("base64,", 1)[1]
            if len(b64) > 32 and len(b64) % 4 == 0:
                return v
        return None
    except Exception:
        return None


def _cache_write(cf, val):
    if cf is not None:
        try:  # atomic: write a temp then rename, so a kill mid-write can never leave
            # a truncated .uri that the prefix check would happily serve as a hero
            import os
            tmp = cf.with_suffix(cf.suffix + ".tmp")
            tmp.write_text(val if val else "NONE", encoding="ascii")
            os.replace(tmp, cf)
        except Exception:
            pass


def _cache_read_json(cf):
    """A cached JSON value (list/dict), or None when absent / truncated (a kill mid-write
    before the atomic rename) - a parse failure is a miss, so the unit is recomputed."""
    if cf is None or not cf.exists():
        return None
    try:
        return json.loads(cf.read_text(encoding="utf-8"))
    except Exception:
        return None


def _cache_write_json(cf, obj):
    if cf is not None:
        try:  # atomic temp+rename, exactly like _cache_write - a shell-cap kill mid-write
            import os
            tmp = cf.with_suffix(cf.suffix + ".tmp")
            tmp.write_text(json.dumps(obj, ensure_ascii=False), encoding="utf-8")
            os.replace(tmp, cf)
        except Exception:
            pass


def page_hero_and_plan(pdf_path: Path, page_index: int,
                       budget_kb: int = DEFAULT_BUDGET_KB,
                       cache_dir: Path | str | None = None) -> tuple[str | None, str | None]:
    """(hero_uri, plan_uri) for one brochure page - the broker's combination rules:
      * a photo exists           -> hero = the photo; plan slot = the site plan (or None)
      * no photo, a plan exists  -> hero = the PLAN, and the plan slot carries it too
      * neither                  -> hero = legacy page-render comparison (or None); plan None

    Hero ladder (engine-agnostic): A) best decodable EMBEDDED image - a real photo
    (>= MODEST_PHOTO with the flatness multiplier) wins outright, the page raster
    can never beat it; B) the most photographic image REGION cropped from the page
    render via pdfplumber geometry (location-map and boilerplate boxes excluded) -
    covers backends that cannot decode the image streams; C) page render vs
    sub-modest embedded; D) None (placeholder upstream). The site plan comes from
    page_plan(). Both results are memoised on disk per (source, page, budget) -
    re-rastering + the JPEG ladder were the dominant re-run wall-clock cost."""
    if Image is None:
        return (None, None)  # no Pillow: no hero/plan; merge fills the placeholder
    cf_h = _cache_file(pdf_path, page_index, budget_kb, "hero", cache_dir)
    cf_p = _cache_file(pdf_path, page_index, budget_kb, "plan", cache_dir)
    ch, cp = _cache_read(cf_h), _cache_read(cf_p)
    if ch is not None and cp is not None:
        return (ch or None), (cp or None)

    hero = plan = None
    # Tier A: embedded images ranked by KIND TIER first (a real photo/aerial/render leads;
    # a plan/map/text never beats a photo even when it scores higher on colourfulness),
    # then by photographic_score within a tier. A PHOTO-kind embedded image is the hero
    # OUTRIGHT - even a near-monochrome warehouse interior the colourfulness score
    # under-rates. A non-photo best (a slide's map/plan) is NOT made the hero here: it
    # falls through to the crop/plan/render tiers so the page's real photo (or the plan
    # slot) wins, and the G-images gate BLOCKS a non-photo hero that still survives.
    doc = _get_doc(pdf_path)
    cands = []  # (tier, -score, img)
    for im in page_embedded_images(doc, page_index):
        if im["w"] >= MIN_HERO_W and im["h"] >= MIN_HERO_H:
            cands.append((HERO_TIER.get(classify_image(im["img"]), 9),
                          -photographic_score(im["img"]), im["img"]))
    cands.sort(key=lambda c: (c[0], c[1]))
    best = cands[0][2] if cands else None
    best_is_photo = bool(cands) and cands[0][0] == HERO_TIER["photo"]
    if best is not None and best_is_photo:
        hero = to_data_uri(compress(best, HERO_MAX_EDGE, budget_kb))
    # Tier B: photographic region cropped from the page render (stream-decode-proof, and
    # it rescues a photo baked into a flattened slide alongside text/plan/map)
    if hero is None:
        crop = bbox_crop_hero(pdf_path, page_index, cache_dir=cache_dir)
        if crop is not None:
            hero = to_data_uri(compress(crop, HERO_MAX_EDGE, budget_kb))
    # the site plan, independent of the hero
    plan_img = page_plan(pdf_path, page_index, cache_dir=cache_dir)
    if plan_img is not None:
        plan = to_data_uri(compress(plan_img, PLAN_MAX_EDGE, budget_kb))
    # plan-only page: the plan IS the hero (above a map / slide screenshot / placeholder)
    # and stays in the plan slot too
    if hero is None and plan is not None:
        hero = plan
    # Tier C: a TEXT-BEARING page is a SLIDE - its whole-page render is a "screenshot of
    # the slide", NEVER a hero. Only an IMAGE-ONLY page (no text layer) may use its render,
    # and only when the render itself is a real photo (a full-bleed photo page). Otherwise
    # fall to the best embedded as a LAST RESORT (a plan/map - the gate then BLOCKS the
    # non-photo hero for sign-off), else leave hero None -> the honest placeholder.
    if hero is None:
        try:
            page_text = (doc[page_index].get_text() or "").strip()
        except Exception:
            page_text = ""
        if len(page_text) < 200:  # image-only page: its render may BE the property image
            try:
                raster = page_raster(doc, page_index)
            except Exception:  # no renderer in this sandbox tier
                raster = None
            if (raster is not None and raster.width >= MIN_HERO_W
                    and classify_image(raster) == "photo"):
                hero = to_data_uri(compress(raster, HERO_MAX_EDGE, budget_kb))
        if hero is None and best is not None:
            hero = to_data_uri(compress(best, HERO_MAX_EDGE, budget_kb))

    _cache_write(cf_h, hero)
    _cache_write(cf_p, plan)
    return hero, plan


def hero_for_pdf_page(pdf_path: Path, page_index: int,
                      budget_kb: int = DEFAULT_BUDGET_KB,
                      cache_dir: Path | str | None = None) -> str | None:
    """Back-compat wrapper: the hero half of page_hero_and_plan()."""
    return page_hero_and_plan(pdf_path, page_index, budget_kb, cache_dir)[0]


def best_hero_in_deck(path: Path, budget_kb: int = DEFAULT_BUDGET_KB,
                      cache_dir: Path | str | None = None, max_pages: int = 80) -> str | None:
    """Best photographic HERO across a WHOLE deck (PDF pages OR PPTX slides). Used when
    a 0-record brochure has been MATCHED to a known property (from a tracker, email or
    other deck) but carries no page_no, so we scan the deck for its best photo. CHEAP by
    design - the embedded-image / slide-picture tier ONLY (no pdfplumber geometry, no
    page render) - so it never hits the _placed_layout speed cliff. Returns a hero data
    URI or None (no usable photo -> the property keeps its honest placeholder).

    DELEGATES to _deck_photo_index so the result is EXACTLY gallery_for_deck()[0]. The
    photo-match path in merge sets photo=this AND gallery=gallery_for_deck(), and the
    images gate requires gallery[0]==photo: deriving the hero from a SEPARATE ranking
    (the earlier code compared UNROUNDED scores while the index rounds to 3dp + tie-breaks
    by page/sig) could diverge on a near-tie and HARD-BLOCK the gate. Shares the index's
    on-disk cache, so the photo-match hero + gallery are one computation."""
    if Image is None:
        return None  # no Pillow: no hero; the property keeps its honest placeholder
    idx = _deck_photo_index(Path(path), budget_kb, cache_dir, max_pages)
    return idx[0]["uri"] if idx else None


def _photo_sig(img) -> str:
    """Cheap content signature (8x8 greyscale) to dedup the SAME photo repeated across
    pages/slides (a reused hero, a logo) so the gallery never shows a near-duplicate.
    Deterministic; '' on failure (treated as unique)."""
    import hashlib
    try:
        return hashlib.sha1(img.convert("L").resize((8, 8)).tobytes()).hexdigest()
    except Exception:
        return ""


def _deck_photo_index(path: Path, budget_kb: int, cache_dir, max_pages: int = 80) -> list[dict]:
    """Every REAL photo in a deck (>= hero size AND >= MODEST_PHOTO score), one entry per
    distinct image (sig-deduped, highest score kept, lowest page kept), compressed to a
    data URI and tagged with its 0-BASED page/slide, ranked best-first. Cached on disk as
    JSON per source so a re-run/resume is free. The page tag is what lets a MULTI-PROPERTY
    deck contribute only a given property's photos to its gallery, never a neighbour's.
    The MODEST_PHOTO + size floors keep the index to genuine photos (no logos/icons/maps),
    so even a 400-image vector page yields a short list."""
    if Image is None:
        return []
    path = Path(path)
    cf = _cache_file(path, "deck", budget_kb, "galleryidx", cache_dir)
    if cf is not None and cf.exists():
        try:
            return json.loads(cf.read_text(encoding="utf-8"))
        except Exception:
            pass
    # RESUMABLE: scan page-by-page through a per-page cache (each page's photos are
    # compressed + cached on their own), so a shell-cap kill mid-deck loses at most the
    # page in flight - the next run continues instead of re-scanning the whole deck (the
    # all-or-nothing whole-deck scan was an infinite re-run trap on a deck too big for one
    # ~40s window). Assembly below reproduces the old dedup/ranking exactly.
    try:
        if path.suffix.lower() == ".pptx":
            n_pages = min(len(list(_get_pptx(path).slides)), max_pages)
        else:
            n_pages = min(_get_doc(path).page_count, max_pages)
    except Exception:
        n_pages = 0
    by_sig: dict[str, dict] = {}
    for pno in range(n_pages):
        for e in _deck_page_photos(path, pno, budget_kb, cache_dir):
            prev = by_sig.get(e["sig"])
            if prev is None:
                by_sig[e["sig"]] = dict(e)
            else:  # same image across pages: EARLIEST page, BEST score (its uri)
                prev["page"] = min(prev["page"], e["page"])
                if e["score"] > prev["score"]:
                    prev["score"], prev["uri"] = e["score"], e["uri"]
    # rank by KIND TIER first (a real photo/aerial/render leads best_hero_in_deck and the
    # gallery; a plan/map never leads) then by score - so the photo-match hero is a photo,
    # not a map/plan, while plans/maps still appear later in the gallery. (kind defaults to
    # 'photo' for an older cache written before the classifier, preserving its score order.)
    ranked = sorted(by_sig.values(),
                    key=lambda d: (HERO_TIER.get(d.get("kind", "photo"), 9),
                                   -d["score"], d["page"], d["sig"]))
    index = [{"page": d["page"], "score": d["score"], "sig": d["sig"], "uri": d["uri"],
              "kind": d.get("kind", "photo")} for d in ranked]
    _cache_write_json(cf, index)  # whole-deck index (cheap once the per-page caches exist)
    return index


def _deck_page_photos(path: Path, page_index: int, budget_kb: int, cache_dir) -> list[dict]:
    """One page/slide's qualifying photos as [{page, score, sig, uri}] - the highest-
    scoring instance per distinct image (sig), >= hero size AND >= MODEST_PHOTO, compressed
    to a data URI. Cached per (deck, page, budget) so _deck_photo_index resumes mid-deck."""
    cf = _cache_file(path, page_index, budget_kb, "gidxpage", cache_dir, ext=".json")
    cached = _cache_read_json(cf)
    if cached is not None:
        return cached
    best: dict[str, dict] = {}  # sig -> {"_sc": float, "_img": Image}
    try:
        if Path(path).suffix.lower() == ".pptx":
            items = slide_pictures(path, page_index)
        else:
            items = page_embedded_images(_get_doc(path), page_index)
        for im in items:
            if im["w"] >= MIN_HERO_W and im["h"] >= MIN_HERO_H:
                kind = classify_image(im["img"])
                sc = photographic_score(im["img"])
                # keep a real PHOTO even when the colourfulness score under-rates it
                # (grey/industrial), and keep a plan/map that clears the score (the gallery
                # + Site Plan toggle keep them - the tier rank just never lets them LEAD)
                if kind != "photo" and sc < MODEST_PHOTO:
                    continue
                sig = _photo_sig(im["img"])
                prev = best.get(sig)
                if prev is None or sc > prev["_sc"]:
                    best[sig] = {"_sc": sc, "_img": im["img"], "_kind": kind}
    except Exception:
        pass
    out: list[dict] = []
    for sig, d in best.items():
        try:
            uri = to_data_uri(compress(d["_img"], HERO_MAX_EDGE, budget_kb))
        except Exception:
            continue
        out.append({"page": page_index, "score": round(d["_sc"], 3), "sig": sig,
                    "uri": uri, "kind": d["_kind"]})
    _cache_write_json(cf, out)
    return out


def gallery_for_pages(path: Path, page_nos, budget_kb: int = DEFAULT_BUDGET_KB,
                      cache_dir: Path | str | None = None, max_n: int = GALLERY_MAX) -> list[str]:
    """Up to max_n photo data URIs from the given 0-based pages of a deck, best-first.
    PAGE-SCOPED so a multi-property deck contributes only THIS property's photos. Returns
    a 2-tuple (uris, total_available) so the caller can note in the Gaps Report when more
    photos existed than the cap allowed."""
    pages = set(page_nos or [])
    idx = _deck_photo_index(Path(path), budget_kb, cache_dir)
    items = [e for e in idx if e["page"] in pages] if pages else idx
    return [e["uri"] for e in items[:max_n]], len(items)


def gallery_for_deck(path: Path, budget_kb: int = DEFAULT_BUDGET_KB,
                     cache_dir: Path | str | None = None, max_n: int = GALLERY_MAX) -> list[str]:
    """Up to max_n best photo data URIs across a WHOLE deck (the photo-match case: a single
    brochure matched to a tracker property IS that property). (uris, total_available)."""
    idx = _deck_photo_index(Path(path), budget_kb, cache_dir)
    return [e["uri"] for e in idx[:max_n]], len(idx)


def page_image_audit(pdf_path: Path, page_index: int, out_dir: Path, tag: str) -> list[str]:
    """The PLACEHOLDER AUDIT: dump EVERY image candidate on the page - all embedded
    images regardless of size, plus the geometry crops - as labelled thumbnails.
    A placeholder is never a silent default: when the pickers found nothing, a
    human/reviewer must be able to SEE the discard pile and sign off that nothing
    in it was a usable photo or plan (the failure this audits was a real site plan
    filtered out twice - by the size floor, then by the photo scorer - with nobody
    ever shown what was discarded). Returns the written file paths."""
    if Image is None:
        return []  # no Pillow: no audit montage (placeholders are surfaced honestly elsewhere)
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    files: list[str] = []

    def _save(img, kind: str, idx: int):
        try:
            im = img.convert("RGB")
            w, h = im.size
            im.thumbnail((480, 480))
            name = f"{tag}_p{page_index + 1}_{kind}{idx}_{w}x{h}.png"
            im.save(out_dir / name, "PNG")
            files.append(str((out_dir / name).resolve()))
        except Exception:
            pass

    try:
        doc = _get_doc(pdf_path)
        for i, im in enumerate(page_embedded_images(doc, page_index), start=1):
            _save(im["img"], "embedded", i)
    except Exception:
        pass
    try:
        for i, c in enumerate(_page_crops(pdf_path, page_index), start=1):
            _save(c["crop"], "region", i)
    except Exception:
        pass
    return files


_MAP_URI = re.compile(r"maps\.google|google\.[a-z.]{2,8}/maps|goo\.gl/maps|openstreetmap\.org|bing\.com/maps", re.I)
_PLACED_CACHE: dict[str, dict] = {}
# a page with more placed images than this is vector-art / a tiled background, not a
# normal brochure page (a real page has a handful of photos/plans/logos). pdfplumber's
# page.images costs 20-35s to build on such a page - OVER the ~45s shell cap. We detect
# it CHEAPLY via fitz get_images (~0.01s/page) and SKIP that page's geometry: Tier A
# (embedded photo) + Tier C (page render) in page_hero_and_plan still yield a hero; only
# the site-plan slot for that one page is sacrificed. Verified on a real deck: a 2,652-
# image page (page.images = 23.5s) where exactly 1 image is >=1.2% of the page.
_PATHOLOGICAL_IMAGES = 400


def _link_near_box(hl: dict, x0, top, x1, bot) -> bool:
    """A link overlapping the box, or sitting as its caption just below/above
    (brochures put the 'click for location' link under the map image)."""
    lx0, lt = float(hl["x0"]), float(hl["top"])
    lx1, lb = float(hl["x1"]), float(hl["bottom"])
    if min(x1, lx1) - max(x0, lx0) <= 0:
        return False  # no horizontal overlap
    if min(bot, lb) - max(top, lt) > 0:
        return True  # overlaps the image itself
    return min(abs(lt - bot), abs(top - lb)) <= 60  # caption proximity


def _placed_cache_file(pdf_path, cache_dir):
    """Disk path for the per-document geometry cache, keyed on bytes+mtime (P0-5)."""
    if not cache_dir:
        return None
    try:
        import hashlib
        st = Path(pdf_path).stat()
        h = hashlib.sha1(f"placed-v1|{Path(pdf_path).name}|{st.st_size}|{st.st_mtime_ns}"
                         .encode()).hexdigest()
        cdir = Path(cache_dir)
        cdir.mkdir(parents=True, exist_ok=True)
        return cdir / f"{h}.placed.json"
    except Exception:
        return None


def _extract_geom_for_page(page) -> list[dict]:
    """RAW image-geometry entries for ONE pdfplumber page: [{bbox, key, frac, aspect, map}].
    `key` is a STABLE per-image-object identity STRING (same PDF image object reused across
    pages -> same key) for cross-page boilerplate detection - never the placed box, because
    templated brochures put photo/plan/map at identical positions on every page. The boiler
    bool is resolved later (it needs the whole deck). Pure per-page, so it parallelises."""
    import hashlib
    entries: list[dict] = []
    map_links = [hl for hl in page.hyperlinks if _MAP_URI.search(str(hl.get("uri", "")))]
    pw, ph = float(page.width), float(page.height)
    for im in page.images:
        x0, top = float(im["x0"]), float(im["top"])
        x1, bot = float(im["x1"]), float(im["bottom"])
        w, h = x1 - x0, bot - top
        if w <= 4 or h <= 4:
            continue
        stream = im.get("stream")
        oid = getattr(stream, "objid", None)
        if oid is not None:
            ks = f"o:{oid}"
        else:
            try:
                raw = stream.rawdata or b""
                ks = f"h:{len(raw)}:{hashlib.sha1(raw[:64]).hexdigest()[:16]}"
            except Exception:
                ks = f"b:{round(x0)},{round(top)},{round(x1)},{round(bot)}"
        entries.append({"bbox": [x0, top, x1, bot], "key": ks,
                        "frac": (w * h) / (pw * ph), "aspect": w / h,
                        "map": any(_link_near_box(hl, x0, top, x1, bot) for hl in map_links)})
    return entries


def _page_pathological(pdf_path: Path, page_index: int) -> bool:
    """True for a vector-art/tiled page (thousands of placed images) - its 20-35s
    pdfplumber page.images access is SKIPPED (Tier A/C still give a hero). O(1) fitz check."""
    try:
        doc = _get_doc(pdf_path)
        return (page_index < doc.page_count
                and len(doc[page_index].get_images(full=True)) > _PATHOLOGICAL_IMAGES)
    except Exception:
        return False


def _placed_page(pdf_path: Path, page_index: int, cache_dir: Path | str | None = None) -> list[dict]:
    """One page's RAW geometry (boiler not yet resolved), cached per (deck, page) so the
    per-document layout RESUMES mid-deck after a shell-cap kill and PARALLELISES across
    pages in the pre-warm. A pathological page caches []."""
    cf = _cache_file(pdf_path, page_index, 0, "placedpage", cache_dir, ext=".json")
    cached = _cache_read_json(cf)
    if cached is not None:
        return cached
    entries: list[dict] = []
    try:
        if not _page_pathological(pdf_path, page_index):
            import pdfplumber
            with pdfplumber.open(str(pdf_path)) as pl:
                if page_index < len(pl.pages):
                    entries = _extract_geom_for_page(pl.pages[page_index])
    except Exception:
        entries = []
    _cache_write_json(cf, entries)
    return entries


def _placed_layout(pdf_path: Path, cache_dir: Path | str | None = None) -> dict:
    """Per-document image GEOMETRY via pdfplumber (engine-agnostic): for every page the
    placed image boxes, each flagged when a maps-service hyperlink sits on/under it
    (= the location map) or when the same image object repeats on >=3 pages (= boilerplate).

    RESUMABLE: each page's raw geometry is cached individually (_placed_page), then the deck
    is assembled by resolving the boiler flag across the per-page caches. The whole-deck
    layout is checkpointed ONLY when every page is present, so a shell-cap kill mid-deck
    loses at most the page in flight and the next run continues (the old all-or-nothing
    whole-deck parse was an infinite re-run trap on a deck too big for one ~40s window).
    The geometry is intermediate bbox data, so determinism/chrome are unaffected."""
    key = str(Path(pdf_path).resolve())
    if key in _PLACED_CACHE:
        return _PLACED_CACHE[key]
    cf = _placed_cache_file(pdf_path, cache_dir)
    if cf is not None and cf.exists():
        try:
            disk = json.loads(cf.read_text(encoding="utf-8"))
            if isinstance(disk, dict) and "pages" in disk:
                _PLACED_CACHE[key] = disk
                return disk
        except Exception:
            pass
    # page count (cheap, fitz) + per-page pathological set
    n = 0
    skip_pages: set = set()
    try:
        doc = _get_doc(pdf_path)
        n = doc.page_count
        for p in range(n):
            if len(doc[p].get_images(full=True)) > _PATHOLOGICAL_IMAGES:
                skip_pages.add(p)
    except Exception:
        pass
    # read whatever per-page caches already exist; compute the rest with ONE pdfplumber open
    raw: list = [_cache_read_json(_cache_file(pdf_path, p, 0, "placedpage", cache_dir, ext=".json"))
                 for p in range(n)]
    missing = [p for p in range(n) if raw[p] is None]
    if missing:
        try:
            import pdfplumber
            with pdfplumber.open(str(pdf_path)) as pl:
                for p in missing:
                    entries = ([] if (p in skip_pages or p >= len(pl.pages))
                               else _extract_geom_for_page(pl.pages[p]))
                    raw[p] = entries
                    _cache_write_json(_cache_file(pdf_path, p, 0, "placedpage", cache_dir, ext=".json"), entries)
        except Exception:
            pass
    # assemble: boiler = an image object placed on >=3 pages (deck >=4 pages), from the
    # pages we have - identical to the old single-scan result once every page is present
    content_pages: dict = {}
    for p in range(n):
        for e in (raw[p] or []):
            content_pages.setdefault(e["key"], set()).add(p)
    boiler = ({k for k, ps in content_pages.items() if len(ps) >= 3} if n >= 4 else set())
    out = {"pages": []}
    for p in range(n):
        out["pages"].append([{"bbox": e["bbox"], "frac": e["frac"], "aspect": e["aspect"],
                              "map": e["map"], "boiler": e["key"] in boiler}
                             for e in (raw[p] or [])])
    _PLACED_CACHE[key] = out
    complete = n > 0 and all(raw[p] is not None for p in range(n))
    if cf is not None and complete and out["pages"]:
        _cache_write_json(cf, out)  # whole-deck checkpoint ONLY when every page is in
    return out


def _unit_cached(spec) -> bool:
    """True when a pre-warm work unit's atomic cache already exists (so it is skipped)."""
    kind, path_str, page, budget, cache_str = spec
    if kind == "placedpage":
        cf = _cache_file(path_str, page, 0, "placedpage", cache_str, ext=".json")
    elif kind == "gidxpage":
        cf = _cache_file(path_str, page, budget, "gidxpage", cache_str, ext=".json")
    else:  # 'hero' (PDF) / 'slidehero' (PPTX) -> the hero .uri is the primary artefact
        cf = _cache_file(path_str, page, budget,
                         "slide_hero" if kind == "slidehero" else "hero", cache_str)
    return cf is not None and cf.exists()


def _prewarm_unit(spec):
    """ProcessPool worker: compute + CACHE one image unit; the on-disk atomic cache IS the
    result (the return is just ok/err). Top-level + stdlib-importable so it pickles to a
    child on both fork (Linux/Cowork) and spawn (Windows). Never raises."""
    try:
        kind, path_str, page, budget, cache_str = spec
        path = Path(path_str)
        if kind == "placedpage":
            _placed_page(path, page, cache_str)
        elif kind == "gidxpage":
            _deck_page_photos(path, page, budget, cache_str)
        elif kind == "slidehero":
            slide_hero_and_plan(path, page, budget, cache_dir=cache_str)
        else:  # 'hero'
            page_hero_and_plan(path, page, budget, cache_dir=cache_str)
        return True
    except Exception:
        return False


def _crop_stats(crop) -> tuple[float, float]:
    """(white_frac, balance) of a crop - the plan signature is a BALANCED mix of
    white paper and drawn ink (balance peaks at white_frac 0.5); photos and map
    tiles sit at the extremes."""
    im = crop.convert("RGB")
    im.thumbnail((96, 96))
    px = list(im.getdata())
    n = len(px) or 1
    white = sum(1 for r, g, b in px if r >= 235 and g >= 235 and b >= 235) / n
    return white, 4.0 * white * (1.0 - white)


_CROPS_CACHE: dict[tuple, list] = {}


def _page_crops(pdf_path: Path, page_index: int, dpi: int = 150,
                cache_dir: Path | str | None = None) -> list[dict]:
    """The page's content-image regions cropped from the render: each candidate
    placed box (>=4%% of the page, sane aspect, not boilerplate) with its crop,
    photographic score and plan stats, plus the map flag. [] when no renderer or
    no geometry - callers fall through their ladders. Memoised per (path, page)
    since the hero tier-B and the plan picker both consume it."""
    memo_key = (str(Path(pdf_path).resolve()), page_index, dpi)
    if memo_key in _CROPS_CACHE:
        return _CROPS_CACHE[memo_key]
    _CROPS_CACHE[memo_key] = []  # set early so a failure is not recomputed per caller
    layout = _placed_layout(pdf_path, cache_dir)
    if page_index >= len(layout["pages"]):
        return []
    # fraction floor only screens out icon-scale boxes; the ABSOLUTE pixel floors
    # (MIN_PLAN / MIN_HERO) do the real gating downstream. 0.04 was high enough
    # to discard genuinely small site plans before they were ever scored.
    boxes = [b for b in layout["pages"][page_index]
             if b["frac"] >= 0.012 and 0.45 <= b["aspect"] <= 3.0
             and not b.get("boiler")]
    if not boxes:
        return []
    try:
        raster = page_raster(_get_doc(pdf_path), page_index, dpi=dpi)
    except Exception:
        return []  # renderer-less sandbox tier
    # raster width / page width in points -> px per point
    try:
        import pdfplumber
        with pdfplumber.open(str(pdf_path)) as pl:
            pw = float(pl.pages[page_index].width)
    except Exception:
        return []
    scale = raster.width / pw
    out = []
    for b in boxes:
        x0, top, x1, bot = b["bbox"]
        crop = raster.crop((max(0, int(x0 * scale)), max(0, int(top * scale)),
                            min(raster.width, int(x1 * scale)),
                            min(raster.height, int(bot * scale))))
        # the PLAN floor, not the photo floor: small site plans must reach the
        # scorers (photo candidacy re-applies MIN_HERO downstream)
        if crop.width < MIN_PLAN_W or crop.height < MIN_PLAN_H:
            continue
        white, balance = _crop_stats(crop)
        out.append({"crop": crop, "map": b["map"], "score": photographic_score(crop),
                    "white": white, "balance": balance,
                    "rank": balance * math.sqrt(crop.width * crop.height)})
    _CROPS_CACHE[memo_key] = out
    return out


def bbox_crop_hero(pdf_path: Path, page_index: int, dpi: int = 150,
                   cache_dir: Path | str | None = None):
    """Tier-B hero: the most photographic embedded-image REGION cropped out of
    the page raster (geometry via pdfplumber - works even when the backend cannot
    DECODE the image streams, the failure that shipped whole cluttered pages as
    heroes). Location-map and boilerplate boxes are excluded. Returns a PIL image
    scoring >= MODEST_PHOTO, or None."""
    if Image is None:
        return None  # no Pillow: no geometry crops
    cands = [c for c in _page_crops(pdf_path, page_index, dpi, cache_dir)
             if not c["map"] and c["crop"].width >= MIN_HERO_W
             and c["crop"].height >= MIN_HERO_H]  # photo candidacy keeps the photo floor
    best = max(cands, key=lambda c: c["score"], default=None)
    return best["crop"] if best and best["score"] >= MODEST_PHOTO else None


def page_plan(pdf_path: Path, page_index: int, dpi: int = 150,
              cache_dir: Path | str | None = None):
    """The page's SITE PLAN, or None. A plan is a content image region that is
    not the page's PHOTO (only the box that actually wins the hero is excluded -
    a COLOURFUL site plan scores photo-ish ~13 and must stay a plan candidate),
    not the hyperlinked location map, not boilerplate, with the plan signature: a
    balanced mix of white paper and drawn ink (white_frac 0.15-0.90; balance x
    area ranks the best). Calibrated on a real Spanish deck where the colourful
    site plan, the grey line plan, the Google location map and the photos all had
    to separate correctly."""
    if Image is None:
        return None  # no Pillow: no geometry crops
    cands = [c for c in _page_crops(pdf_path, page_index, dpi, cache_dir) if not c["map"]]
    if not cands:
        return None
    hero = max(cands, key=lambda c: c["score"])
    if hero["score"] < MODEST_PHOTO:
        hero = None  # no photo on this page - every region stays a plan candidate
    plans = [c for c in cands
             if c is not hero and 0.15 <= c["white"] <= 0.90]
    best = max(plans, key=lambda c: c["rank"], default=None)
    return best["crop"] if best else None


# --- WHOLE-PAGE RENDERED SITE PLAN (vector line-art a placed-image crop cannot reach) -- #
# A site plan is often VECTOR graphics drawn straight into the page, not a placed raster.
# Pulled as an embedded image it rasterises to solid black, so page_plan()/planRef (which
# crop placed-image boxes only) find nothing. page_raster() renders the vector content
# correctly, so we RENDER the page and crop to its ink bbox. The LLM names the plan page
# (__meta.plan_page) and a deterministic render+classify detector is the universal
# fallback + verifier. The render binds the PLAN SLOT ONLY - a vector plan is NEVER made
# the card hero (it would trip the G-images gate and change the card look; intentional).
PLAN_RENDER_DPI = 150


def _rendered_plan_crop(path: Path, page_index: int, dpi: int = PLAN_RENDER_DPI,
                        cache: Path | str | None = None) -> tuple:
    """Render the whole page and crop to its INK bounding box (Pillow getbbox on a
    thresholded non-white copy, with a small margin), then classify the crop. Returns
    (crop, signals, kind) - (None, {}, None) when Pillow / the renderer is unavailable or
    the page is effectively blank. Pure helper for page_render_plan + best_plan_page_render;
    cache is accepted for signature parity (the rendered URI is cached by the callers)."""
    if Image is None:
        return (None, {}, None)
    try:
        path = Path(path)
        if path.suffix.lower() == ".pptx":
            pdf = soffice_pdf(path, cache)
            if pdf is None:
                return (None, {}, None)
            doc = _get_doc(pdf)
        else:
            doc = _get_doc(path)
        if not (0 <= page_index < doc.page_count):
            return (None, {}, None)
        raster = page_raster(doc, page_index, dpi=dpi)
    except Exception:
        return (None, {}, None)  # renderer-less sandbox tier / open failure -> honest None
    if raster is None:
        return (None, {}, None)
    try:
        rgb = raster.convert("RGB")
        # ink mask: anything that is NOT near-white. point() on the greyscale gives a
        # crisp foreground; getbbox() returns the tight box of the drawn content so the
        # paper margins (which read as 'text'/blank) do not dominate the classification.
        grey = rgb.convert("L")
        ink = grey.point(lambda p: 0 if p >= 238 else 255)
        bbox = ink.getbbox()
        if bbox is None:
            return (None, {}, None)  # all white = a blank page, never a plan
        x0, y0, x1, y1 = bbox
        mw = max(8, int((x1 - x0) * 0.03))
        mh = max(8, int((y1 - y0) * 0.03))
        x0 = max(0, x0 - mw); y0 = max(0, y0 - mh)
        x1 = min(rgb.width, x1 + mw); y1 = min(rgb.height, y1 + mh)
        crop = rgb.crop((x0, y0, x1, y1))
        if crop.width < 2 or crop.height < 2:
            return (None, {}, None)
        sig = _hero_signals(crop)
        kind = classify_image(crop, sig)
        return (crop, sig, kind)
    except Exception:
        return (None, {}, None)


def page_render_plan(path: Path, page_index: int, budget_kb: int = DEFAULT_BUDGET_KB,
                     cache_dir: Path | str | None = None) -> str | None:
    """The LLM-HINTED plan page (__meta.plan_page): render+ink-crop the page and bind it
    to the plan slot, LENIENTLY verified - bind UNLESS it is an obvious photo (classify
    'photo') or near-blank/all-white. The agent picked this page by LOOKING at its render,
    so the verify only screens out a clearly-wrong pick, never demands the plan signature.
    Returns a compressed plan data URI or None. Cached per (source, page, budget) under
    kind='planpage' so a resume is byte-deterministic. Degrades to None without Pillow /
    a renderer (honest null)."""
    if Image is None:
        return None
    if not (isinstance(page_index, int) and not isinstance(page_index, bool) and page_index >= 0):
        return None
    cf = _cache_file(path, page_index, budget_kb, "planpage", cache_dir)
    cached = _cache_read(cf)
    if cached is not None:
        return cached or None
    crop, sig, kind = _rendered_plan_crop(path, page_index, cache=cache_dir)
    uri = None
    if crop is not None:
        # LENIENT: reject only an obvious photo or a near-blank page; bind anything else
        # (the LLM looked at the render and chose THIS page as the plan).
        white = sig.get("white", 0.0)
        if kind != "photo" and white <= 0.985:
            uri = to_data_uri(compress(crop, PLAN_MAX_EDGE, budget_kb))
    _cache_write(cf, uri)
    return uri


def best_plan_page_render(path: Path, page_nos, budget_kb: int = DEFAULT_BUDGET_KB,
                          cache_dir: Path | str | None = None) -> tuple:
    """DETERMINISTIC fallback (no LLM hint): over the given (per-property) pages, render+
    ink-crop+classify and pick the most plan-like page. CONSERVATIVE so it never fabricates
    a plan on a real deck - a page qualifies ONLY when classify_image == 'plan' AND its
    white fraction is balanced (0.15-0.90, the page_plan signature - paper + drawn ink),
    explicitly NOT a photo and NOT a map. Returns (uri, page_no) for the best (highest
    balance) qualifying page, else (None, None). The page set is SORTED so the result is a
    pure function of (source, pages). Each candidate page's URI is cached per
    (source, page, budget) under kind='planpage'. Degrades to (None, None) without Pillow /
    a renderer."""
    if Image is None:
        return (None, None)
    best = None  # (balance, page_no, uri)
    for pno in sorted({p for p in (page_nos or [])
                       if isinstance(p, int) and not isinstance(p, bool) and p >= 0}):
        # CONSERVATIVE GATE 1: only a page with NO hero-size EMBEDDED image is eligible.
        # A genuine vector site-plan page is image-light (the plan is drawn vector content,
        # not a placed raster); a page that already carries a hero-size embedded image is a
        # property/photo page whose plan, if any, comes via the embedded planRef / page_plan
        # tier - never this whole-page render fallback. This keeps the detector off ordinary
        # photo pages (so an existing offline fixture never spuriously gains a plan).
        try:
            if candidates_for_page(path, pno):
                continue
        except Exception:
            pass
        crop, sig, kind = _rendered_plan_crop(path, pno, cache=cache_dir)
        if crop is None:
            continue
        white = sig.get("white", 0.0)
        if kind != "plan" or not (0.15 <= white <= 0.90):
            continue  # not the plan signature (a photo/map/text/blank page never binds)
        balance = 4.0 * white * (1.0 - white)
        cf = _cache_file(path, pno, budget_kb, "planpage", cache_dir)
        uri = _cache_read(cf)
        if uri is None:
            uri = to_data_uri(compress(crop, PLAN_MAX_EDGE, budget_kb))
            _cache_write(cf, uri)
        # a cached "" is a prior negative for this (source,page,budget) - leave it; never
        # recompute/overwrite (that would defeat resume + the shared planpage-cache contract).
        # A "" is falsy, so the `if uri` below skips it - the page simply does not bind.
        if uri and (best is None or balance > best[0]):
            best = (balance, pno, uri)
    if best is None:
        return (None, None)
    return (best[2], best[1])


_PLACEHOLDER: str | None = None


def placeholder() -> str:
    """A neutral CBRE-green 16:9 placeholder data URI (cached)."""
    global _PLACEHOLDER
    if _PLACEHOLDER is None:
        if Image is None:  # no Pillow: serve the pre-baked placeholder asset (a valid
            # data:image/jpeg URI the chrome + the images gate accept). Integrity-tracked,
            # so a truncated copy is caught by preflight, not shipped.
            _PLACEHOLDER = (Path(__file__).resolve().parent.parent
                            / "assets" / "placeholder.uri").read_text(encoding="utf-8").strip()
        else:
            img = Image.new("RGB", (1280, 720), (0, 63, 45))  # CBRE dark green
            _PLACEHOLDER = to_data_uri(compress(img, 1280, 40))
    return _PLACEHOLDER


# --------------------------------------------------------------------------- #
# PPTX slides - the PDF page ladder's twin. A vision transcription of a slide
# deck carries source_type "pptx" (per the vision contract), but the harvest
# above is PDF-only, so every such record silently degraded to the placeholder
# AND bypassed the placeholder audit (tried_pages was PDF-fed only). These give
# slides the same hero/plan/audit treatment.
# --------------------------------------------------------------------------- #

_PPTX_CACHE: dict[str, object] = {}


def _get_pptx(path: Path):
    key = str(Path(path).resolve())
    prs = _PPTX_CACHE.get(key)
    if prs is None:
        from pptx import Presentation
        prs = Presentation(str(path))
        _PPTX_CACHE[key] = prs
    return prs


def slide_pictures(pptx_path: Path, slide_index: int) -> list[dict]:
    """All decodable raster pictures on one slide, largest first (undecodable
    WMF/EMF vectors are skipped, never abort the harvest)."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = _get_pptx(pptx_path)
        slides = list(prs.slides)
        if not (0 <= slide_index < len(slides)):
            return []
        out = []
        for shape in slides[slide_index].shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = None
                try:
                    img = _open(shape.image.blob)
                except Exception:
                    img = None
                if img is not None:
                    w, h = img.size
                    out.append({"img": img, "w": w, "h": h, "area": w * h})
        out.sort(key=lambda d: -d["area"])
        return out
    except Exception:
        return []


_SOFFICE: object = False  # False = not probed yet; then str | None


def _find_soffice() -> str | None:
    """Locate a headless LibreOffice (the only reliable slide renderer -
    python-pptx cannot rasterise). Memoised; None when absent."""
    global _SOFFICE
    if _SOFFICE is not False:
        return _SOFFICE  # type: ignore[return-value]
    import shutil
    cand = shutil.which("soffice") or shutil.which("libreoffice")
    if not cand and sys.platform == "win32":
        for p in (Path(r"C:\Program Files\LibreOffice\program\soffice.exe"),
                  Path(r"C:\Program Files (x86)\LibreOffice\program\soffice.exe")):
            if p.exists():
                cand = str(p)
                break
    _SOFFICE = cand
    return cand


def soffice_pdf(src: Path, cache_dir: Path | str | None) -> Path | None:
    """Convert a slide deck to PDF via headless LibreOffice (slides map 1:1 to
    pages), memoised on disk per (name, size, mtime) so one deck converts once
    per run history. None when LibreOffice is absent or conversion fails -
    callers fall back to embedded slide pictures."""
    exe = _find_soffice()
    if exe is None:
        return None
    src = Path(src)
    try:
        st = src.stat()
    except OSError:
        return None
    import hashlib
    import subprocess
    import tempfile
    cdir = Path(cache_dir) if cache_dir else Path(tempfile.gettempdir()) / "cbre_longlist_soffice"
    try:
        cdir.mkdir(parents=True, exist_ok=True)
    except Exception:
        return None
    key = hashlib.sha1(f"soffice|{src.name}|{st.st_size}|{st.st_mtime_ns}".encode()).hexdigest()[:16]
    target = cdir / f"{src.stem}.{key}.pdf"
    if target.exists() and target.stat().st_size > 0:
        return target
    try:
        with tempfile.TemporaryDirectory(prefix="soffice_") as td:
            # a private user profile so a running desktop LibreOffice (or a
            # parallel conversion) can never lock the conversion out
            profile = Path(td) / "profile"
            profile.mkdir()
            subprocess.run(
                [exe, "--headless", "--norestore",
                 f"-env:UserInstallation={profile.as_uri()}",
                 "--convert-to", "pdf", "--outdir", td, str(src)],
                check=True, capture_output=True, timeout=180)
            produced = Path(td) / (src.stem + ".pdf")
            if not produced.exists() or produced.stat().st_size == 0:
                return None
            import shutil
            shutil.move(str(produced), str(target))
        return target
    except Exception:
        return None


def slide_hero_and_plan(pptx_path: Path, slide_index: int,
                        budget_kb: int = DEFAULT_BUDGET_KB,
                        cache_dir: Path | str | None = None) -> tuple[str | None, str | None]:
    """(hero_uri, plan_uri) for one slide, same combination rules as the PDF
    page ladder. PREFERRED: LibreOffice renders the deck to PDF (slides map 1:1
    to pages) and the FULL page ladder runs on the converted page - embedded
    tiers, geometry crops, the plan picker, everything. Fallback (no
    LibreOffice): the slide's decodable embedded pictures - the best
    photo-scoring picture is the hero, a plan-signature picture fills the plan
    slot, and a plan-only slide promotes the plan to hero. Memoised on disk per
    (source, slide, budget) like the PDF path."""
    if Image is None:
        return (None, None)  # no Pillow: no slide hero/plan; merge fills the placeholder
    cf_h = _cache_file(pptx_path, slide_index, budget_kb, "slide_hero", cache_dir)
    cf_p = _cache_file(pptx_path, slide_index, budget_kb, "slide_plan", cache_dir)
    ch, cp = _cache_read(cf_h), _cache_read(cf_p)
    if ch is not None and cp is not None:
        return (ch or None), (cp or None)

    hero = plan = None
    pdf = soffice_pdf(pptx_path, cache_dir)
    if pdf is not None:
        try:
            hero, plan = page_hero_and_plan(pdf, slide_index, budget_kb, cache_dir=cache_dir)
        except Exception:
            hero = plan = None
    if hero is None and plan is None:
        pics = slide_pictures(pptx_path, slide_index)
        best, best_score = None, -1.0
        for im in pics:
            if im["w"] >= MIN_HERO_W and im["h"] >= MIN_HERO_H:
                sc = photographic_score(im["img"])
                if sc > best_score:
                    best, best_score = im["img"], sc
        if best is not None and best_score >= MODEST_PHOTO:
            hero = to_data_uri(compress(best, HERO_MAX_EDGE, budget_kb))
        plan_cands = []
        for im in pics:
            if im["img"] is best and hero is not None:
                continue  # the hero photo is never also the plan
            if im["w"] < MIN_PLAN_W or im["h"] < MIN_PLAN_H:
                continue
            try:
                white, balance = _crop_stats(im["img"].convert("RGB"))
            except Exception:
                continue
            if 0.15 <= white <= 0.90:  # the plan signature (page_plan's bounds)
                plan_cands.append((balance * math.sqrt(im["w"] * im["h"]), im["img"]))
        if plan_cands:
            plan_cands.sort(key=lambda t: -t[0])
            plan = to_data_uri(compress(plan_cands[0][1], PLAN_MAX_EDGE, budget_kb))
        if hero is None and plan is not None:
            hero = plan  # plan-only slide: the plan IS the hero (broker rule)
        elif hero is None and best is not None:
            # sub-modest but real picture still beats the placeholder (the PDF
            # ladder's renderer-less tier C does the same)
            hero = to_data_uri(compress(best, HERO_MAX_EDGE, budget_kb))

    _cache_write(cf_h, hero)
    _cache_write(cf_p, plan)
    return hero, plan


def slide_image_audit(pptx_path: Path, slide_index: int, out_dir: Path, tag: str,
                      cache_dir: Path | str | None = None) -> list[str]:
    """page_image_audit's PPTX twin: dump every embedded slide picture (plus the
    LibreOffice-rendered slide, when a renderer is available) as labelled
    thumbnails, so a slide-sourced placeholder is a reviewed conclusion too."""
    if Image is None:
        return []  # no Pillow: no audit montage
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    files: list[str] = []

    def _save(img, kind: str, idx: int):
        try:
            im = img.convert("RGB")
            w, h = im.size
            im.thumbnail((480, 480))
            name = f"{tag}_s{slide_index + 1}_{kind}{idx}_{w}x{h}.png"
            im.save(out_dir / name, "PNG")
            files.append(str((out_dir / name).resolve()))
        except Exception:
            pass

    for i, im in enumerate(slide_pictures(pptx_path, slide_index), start=1):
        _save(im["img"], "picture", i)
    try:
        pdf = soffice_pdf(pptx_path, cache_dir)
        if pdf is not None:
            _save(page_raster(_get_doc(pdf), slide_index), "sliderender", 1)
    except Exception:
        pass
    return files
