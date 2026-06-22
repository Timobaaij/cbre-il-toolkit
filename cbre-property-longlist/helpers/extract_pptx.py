#!/usr/bin/env python3
"""extract_pptx.py - candidate records + hero imagery from a brochure PPTX.

Reuses the PDF label machinery (one property per slide). PPTX is the preferred
IMAGE source (slide pictures are higher resolution than the PDF's embedded
rasters); PDF stays the preferred FIELD source. Emits the same record shape so
merge.py can dedupe a slide against its matching PDF page by city+developer+park.

CLI:
  python extract_pptx.py <file.pptx> --region Budapest --country HU [--out records.json]
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.path.insert(0, str(Path(__file__).resolve().parent))
import extract_pdf as P
import images as IMG


def slide_text(slide) -> str:
    chunks = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    chunks.append(line)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    t = cell.text.strip()
                    if t:
                        chunks.append(t)
    return "\n".join(chunks)


def slide_hero(slide, budget_kb: int) -> str | None:
    # the most PHOTOGRAPHIC picture on the slide (not merely the largest), so a logo
    # or branded element does not win by area - same scorer as the PDF path
    best, best_score = None, -1.0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img = IMG._open(shape.image.blob)
            except Exception:
                continue
            if img and img.width >= IMG.MIN_HERO_W and img.height >= IMG.MIN_HERO_H:
                sc = IMG.photographic_score(img)
                if sc > best_score:
                    best, best_score = img, sc
    return IMG.to_data_uri(IMG.compress(best, IMG.HERO_MAX_EDGE, budget_kb)) if best else None


def extract(path: Path, region: str, country: str, with_images: bool = True,
            budget_kb: int = IMG.DEFAULT_BUDGET_KB) -> list[dict]:
    prs = Presentation(str(path))
    slides = [(i, slide, slide_text(slide)) for i, slide in enumerate(prs.slides)]

    # OWN-LINE path (>=2 own-line labels on a slide, any language) - unchanged. If
    # NO slide has own-line labels (inline "Label value" deck), fall back to
    # same-line parsing (>=2 inline labels). Mirrors extract_pdf's document gate.
    own = [(i, s, t) for i, s, t in slides if len(P._find_labels(t)) >= 2]
    inline_mode = not own
    use = own if own else [(i, s, t) for i, s, t in slides if len(P._find_inline_labels(t)) >= 2]

    records = []
    for i, slide, text in use:
        rec = P.parse_property_page(text, None, region, country, Path(path).name, i)
        if inline_mode and not P._has_core(rec):  # coreless spec-only slide -> vision fallback
            continue
        rec["__meta"]["source_type"] = "pptx"
        rec["__meta"]["locator_base"] = f"slide {i + 1}"
        if with_images:
            hero = slide_hero(slide, budget_kb)
            if hero:
                rec["photo"] = hero
                rec["__meta"]["prov"]["photo"] = f"slide {i + 1}"
        records.append(rec)
    return records


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--region", required=True)
    ap.add_argument("--country", required=True)
    ap.add_argument("--no-images", action="store_true")
    ap.add_argument("--out")
    args = ap.parse_args()
    recs = extract(Path(args.pptx), args.region, args.country, with_images=not args.no_images)
    sys.stdout.reconfigure(encoding="utf-8")
    if args.out:
        Path(args.out).write_text(json.dumps(recs, ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"OK {len(recs)} records -> {args.out}")
    else:
        print(json.dumps([{k: v for k, v in r.items() if k != 'photo'} for r in recs],
                         ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
