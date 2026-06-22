"""Story-led SCENE COMPOSER for CBRE decks - the default way to build.

THE STORY MAKES THE LAYOUT. You do not pick a recipe and pour content into it.
You declare what each slide must SAY (a deck of slides; each slide a `scene` of
ordered rows; each row split into cells; each cell one styled primitive) and this
composer lays the scene on the safe CBRE grid and sizes every cell's text UP to
fill its space. Slide count and shape flex with the story; two decks should look
different; two slides should rarely look the same.

It is the same engine the cbre-il-account-briefing skill uses, generalised for any
deck. It sits on top of `build.py` (the CBRE visual system): it draws the CBRE
chrome (eyebrow, fit-to-text serif headline, lead line, footer, wordmark) and then
composes the scene from the cells, reusing the `build` primitives and helpers as
the palette each cell draws from.

Slide kinds (the `kind` field):
  cover    -> the polished build.cover (giant serif title + optional themes strip)
  section  -> build.section_divider (giant numeral + title + optional lead/items)
  scene    -> a freely composed slide: chrome + an ordered list of rows/cells
  closing  -> build.thank_you (optional contact cards)

Cell kinds (inside a scene row's `cells`):
  prose  - a full paragraph (text; optional bold lead `label`). The explainer.
  stat   - a hero value + caption label. Several stat cells in a row = a KPI strip.
  list   - items [{title,text}], numbered or bulleted.
  table  - headers + rows; the last column is the analytical read.
  panel  - a dark accent side box; title + items [{label,value}] (or text).
  quote  - a pull quote + attribution.
  heading- a small section label. rule - a thin divider line.
  callout- the CBRE expert-note box (title + body, optional tag).
  chips  - a row of rounded pills (tags/countries/status).
  card   - one roman/decimal numbered card (a row of card cells = a card grid).
  image  - a picture from a path, fit within the cell.

THE RULES (same spirit as the il-account-briefing skill):
  1. Story-led, not recipe-led. Compose each slide from what it must say.
  2. Density from substance, never from tricks. Fill a slide with MORE real,
     relevant content if it makes sense, never with spacing or ballooned fonts.
     Leftover space is a signal to write more, not to stretch what is there.
  3. Explain, do not tabulate. Prose carries the narrative; tables are evidence.
  4. No lazy repetition. A scene layout should appear at most twice, ideally once.
  5. Readable, not a billboard. Text sizes UP to a readable cap, never tiny.
  6. No em/en dashes (swept before save). Box grows to fit text; font never shrinks.

Usage (Python):
  import compose
  compose.render(plan, "Deck.pptx")        # plan = a dict (below) or a .json path

Usage (CLI):
  python compose.py plan.json Deck.pptx [--no-resolve] [--no-label-bake]

Plan shape:
  {
    "deck_meta": {"eyebrow": "CBRE | ADVISORY"},   # optional defaults
    "slides": [
      {"kind": "cover", "title": "...", "subtitle": "...", "eyebrow": "...",
       "date": "...", "themes": ["...", "..."]},
      {"kind": "scene", "tone": "dark", "eyebrow": "01 | CONTEXT",
       "headline": "...", "lead": "...", "footer": "...",
       "scene": [
         {"weight": 1.3, "cells": [{"kind": "prose", "label": "THE SHIFT", "text": "..."}]},
         {"weight": 0.8, "cells": [{"kind": "stat", "value": "46%", "label": "..."},
                                   {"kind": "stat", "value": "17", "label": "..."}]}
       ]},
      {"kind": "section", "number": 1, "title": "...", "lead": "...", "items": ["..."]},
      {"kind": "closing", "title": "Thank you.", "contacts": [{"name": "...", "email": "..."}]}
    ]
  }
"""
from __future__ import annotations

import argparse
import json
import os
import shutil
import sys
import tempfile
from pathlib import Path

# Locate the cbre-corporate-pptx build library (this file lives next to it).
_HERE = Path(__file__).resolve().parent
for _p in (_HERE, Path.home() / ".claude/skills/cbre-corporate-pptx/scripts"):
    if (_p / "build.py").exists():
        sys.path.insert(0, str(_p.resolve()))
        break
import build  # noqa: E402
from build import COLORS as C, FONTS as Fz, ED_X, ED_W, ED_SAFE_BOT, SLIDE_W  # noqa: E402
from pptx.util import Inches  # noqa: E402

# Cell kinds this composer can draw (kept in sync with the CELL dict below).
CELL_KINDS = {"prose", "stat", "list", "table", "panel", "quote", "heading",
              "rule", "callout", "chips", "card", "image"}

# ---------------------------------------------------------------------------
# Text hygiene + size estimators
# ---------------------------------------------------------------------------

def _clean(s):
    if not isinstance(s, str):
        return s
    return s.replace("—", "-").replace("–", "-")

def _ink(tone):
    return C["white"] if tone == "dark" else C["ink"]

def _muted(tone):
    return C["mint"] if tone == "dark" else C["ink_2"]

def _accent(tone):
    return C["mint"] if tone == "dark" else C["mint_dark"]

def _para_h(text, w, size, ls=1.34):
    """Conservative (slightly tall) wrapped-paragraph height in inches."""
    cpl = max(16, int(w / (0.0080 * size)))
    n = len(str(text))
    return max(1, (n + cpl - 1) // cpl) * (size / 72.0) * ls

def _fit_fill(text, w, avail_h, lo=12.0, hi=20.0, ls=1.34):
    """Largest body size in [lo, hi] whose wrapped text fits avail_h. Picking the
    largest size that fits is what fills the cell and enlarges the text (density)."""
    size = hi
    while size > lo and _para_h(text, w, size, ls) > avail_h:
        size -= 0.5
    return max(size, lo)

# ---------------------------------------------------------------------------
# Slide chrome (the consistent CBRE frame; the diversity lives in the scene)
# ---------------------------------------------------------------------------

def _chrome(deck, plan, slide):
    """Background, footer/wordmark, and the eyebrow + serif headline + lead, drawn by
    build.editorial_header so the TITLE FITS ITS TEXT and returns the exact y where
    the body begins (no reserved-whitespace gap under the title).
    Return (slide, body_top, body_bot, tone)."""
    tone = slide.get("tone", "light")
    s = build.blank(deck, tone=tone)
    default_eyebrow = (plan.get("deck_meta", {}) or {}).get("eyebrow") or "CBRE"
    title = _clean(slide.get("headline") or "")
    lead = _clean(slide.get("lead")) if slide.get("lead") else None
    y = build.editorial_header(s, eyebrow_text=_clean(slide.get("eyebrow") or default_eyebrow),
                               title=title, tone=tone, intro=lead,
                               title_size=slide.get("headline_size", 30))
    foot = slide.get("footer")
    if foot:
        build._text(s, _clean(foot), x=ED_X, y=ED_SAFE_BOT + 0.02, w=SLIDE_W - ED_X - 0.55, h=0.26,
                    font=Fz["sans_l"], size=8.5,
                    color=C["ink_2"] if tone == "light" else C["mint"], anchor="top")
    body_top = y + 0.18
    body_bot = ED_SAFE_BOT - (0.34 if foot else 0.0)
    return s, body_top, body_bot, tone

# ---------------------------------------------------------------------------
# Cell primitives (each draws one styled element to fill its rect)
# ---------------------------------------------------------------------------

def c_prose(s, cell, x, y, w, h, tone):
    yy = y
    label = _clean(cell.get("label") or "")
    if label:
        build._text(s, label, x=x, y=yy, w=w, h=0.28, font=Fz["sans_sb"], size=12,
                    color=_accent(tone), bold=True, uppercase=True, letter_spacing=1.5, anchor="top")
        yy += 0.40
    text = _clean(cell.get("text") or "")
    avail = (y + h) - yy
    size = _fit_fill(text, w, avail, lo=12.5, hi=cell.get("max_size", 16.0))
    build.body(s, [text], x=x, y=yy, w=w, h=avail, size=size, color=_ink(tone),
               tone=tone, line_spacing=1.36)

def c_stat(s, cell, x, y, w, h, tone):
    value = _clean(cell.get("value") or "")
    label = _clean(cell.get("label") or "")
    # Scale the hero number to the cell, but cap by WIDTH so a long value never
    # wraps to a stray second line.
    wcap = w * 130.0 / max(4, len(value))
    vsize = max(24.0, min(60.0, h * 30.0, wcap))
    build.serif_title(s, value, x=x, y=y, w=w, h=vsize / 72 * 1.25, size=vsize,
                      tone=tone, line_spacing=1.0)
    ly = y + vsize / 72.0 * 1.18
    if label:
        build._text(s, label, x=x, y=ly, w=w, h=max(0.3, (y + h) - ly),
                    font=Fz["sans_l"], size=12, color=_muted(tone), anchor="top", line_spacing=1.22)

def c_list(s, cell, x, y, w, h, tone):
    items = [it for it in cell.get("items", []) if it]
    if not items:
        return
    numbered = cell.get("numbered", True)
    gap = 0.16
    ih = (h - gap * (len(items) - 1)) / len(items)
    for i, it in enumerate(items):
        iy = y + i * (ih + gap)
        title = _clean(it.get("title") or "")
        text = _clean(it.get("text") or "")
        if numbered:
            build._text(s, f"{i + 1:02d}", x=x, y=iy, w=0.85, h=0.5, font=Fz["serif"],
                        size=23, color=_accent(tone), anchor="top")
            tx, tw = x + 1.0, w - 1.0
        else:
            build._rect(s, x, iy + 0.06, 0.16, 0.16, fill=_accent(tone))
            tx, tw = x + 0.34, w - 0.34
        ty = iy
        if title:
            build._text(s, title, x=tx, y=ty, w=tw, h=0.30, font=Fz["sans_sb"], size=13.5,
                        color=_ink(tone), bold=True, anchor="top")
            ty += 0.36
        if text:
            avail = (iy + ih) - ty
            tsize = _fit_fill(text, tw, avail, lo=11.0, hi=14.5)
            build.body(s, [text], x=tx, y=ty, w=tw, h=avail, size=tsize,
                       color=_ink(tone), tone=tone, line_spacing=1.30)

def c_table(s, cell, x, y, w, h, tone):
    headers = [_clean(z) for z in cell.get("headers", [])]
    rows = [[_clean(z) for z in r] for r in cell.get("rows", [])]
    if not headers and not rows:
        return
    build.table(s, headers, rows, x=x, y=y, w=w, h=max(1.0, h), tone=tone,
                font_size=max(10.5, cell.get("font_size", 12)), col_aligns=cell.get("aligns"))

def c_panel(s, cell, x, y, w, h, tone):
    """A filled accent panel (dark side box). Items stack by MEASURED height and the
    value font shrinks to fit, so a wrapped value never collides with the next label."""
    build._rect(s, x, y, w, h, fill=C["green_3"])
    build._rect(s, x, y, w, 0.055, fill=C["gold"])
    pad = 0.28
    iw = w - 2 * pad
    yy = y + pad + 0.05
    title = _clean(cell.get("title") or "")
    if title:
        build._text(s, title, x=x + pad, y=yy, w=iw, h=0.34, font=Fz["sans_sb"], size=13,
                    color=C["gold"], bold=True, uppercase=True, letter_spacing=1.6, anchor="top")
        yy += 0.40 + (0.20 if len(title) > 28 else 0.0)
    inner_bot = y + h - pad
    items = [it for it in cell.get("items", []) if it]
    if items:
        gap, lab_h = 0.20, 0.30
        avail = inner_bot - yy
        def content_h(vs):
            return sum(lab_h + _para_h(_clean(it.get("value") or ""), iw, vs, 1.22) + gap for it in items)
        vsize = 17.5
        while vsize > 10.5 and content_h(vsize) > avail:
            vsize -= 0.5
        for it in items:
            build._text(s, _clean(it.get("label") or ""), x=x + pad, y=yy, w=iw, h=0.26,
                        font=Fz["sans_sb"], size=11, color=C["mint"], bold=True,
                        uppercase=True, letter_spacing=1.0, anchor="top")
            vy = yy + lab_h
            vh = _para_h(_clean(it.get("value") or ""), iw, vsize, 1.22)
            build._text(s, _clean(it.get("value") or ""), x=x + pad, y=vy, w=iw, h=vh + 0.05,
                        font=Fz["sans_l"], size=vsize, color=C["off_white"], anchor="top", line_spacing=1.25)
            yy += lab_h + vh + gap
    elif cell.get("text"):
        txt = _clean(cell["text"])
        avail = inner_bot - yy
        size = _fit_fill(txt, iw, avail, lo=11.5, hi=15.0)
        build.body(s, [txt], x=x + pad, y=yy, w=iw, h=avail, size=size,
                   color=C["off_white"], tone="dark", line_spacing=1.34)

def c_quote(s, cell, x, y, w, h, tone):
    text = _clean(cell.get("text") or "")
    attrib = _clean(cell.get("attrib") or "")
    qh = h - (0.34 if attrib else 0.0)
    size = _fit_fill(text, w, qh, lo=15.0, hi=30.0, ls=1.14)
    build.serif_title(s, "“" + text + "”", x=x, y=y, w=w, h=qh, size=size,
                      tone=tone, line_spacing=1.12)
    if attrib:
        build._text(s, attrib, x=x, y=y + qh, w=w, h=0.3, font=Fz["sans_sb"], size=11,
                    color=_accent(tone), uppercase=True, letter_spacing=1.5, anchor="top")

def c_heading(s, cell, x, y, w, h, tone):
    build._text(s, _clean(cell.get("text") or ""), x=x, y=y, w=w, h=h, font=Fz["sans_sb"],
                size=cell.get("size", 13), color=_accent(tone), bold=True,
                uppercase=cell.get("uppercase", True), letter_spacing=1.6, anchor="middle")

def c_rule(s, cell, x, y, w, h, tone):
    build._rect(s, x, y + h / 2, w, 0.014, fill=_accent(tone))

def c_callout(s, cell, x, y, w, h, tone):
    """The CBRE expert-note box. Sized to the cell; the body autofits and the
    callout primitive keeps its bg rect in sync via the resolve pass."""
    title = _clean(cell.get("title") or "CBRE VIEW")
    text = _clean(cell.get("text") or "")
    tag = _clean(cell.get("tag")) if cell.get("tag") else None
    build.callout(s, title=title, body_text=text, x=x, y=y, w=w, h=max(1.0, h),
                  tone=tone, tag=tag)

def c_chips(s, cell, x, y, w, h, tone):
    """A row of rounded pills that wrap within the cell width."""
    items = [it for it in cell.get("items", []) if it]
    if not items:
        return
    chip_h, gap = 0.40, 0.16
    cx, cy = x, y
    line = (C["rule_light"] if tone == "light" else None)
    for it in items:
        txt = _clean(it if isinstance(it, str) else (it.get("text") or ""))
        cw = min(w, 0.34 + len(txt) * 0.095)
        if cx + cw > x + w + 0.01 and cx > x:
            cx = x
            cy += chip_h + gap
        if cy + chip_h > y + h + 0.05:
            break
        build.chip(s, txt, x=cx, y=cy, w=cw, h=chip_h, tone=tone, line=line)
        cx += cw + gap

def c_card(s, cell, x, y, w, h, tone):
    """One numbered card filling the cell. style='roman' -> roman_card (bullet body);
    style='decimal' -> decimal_card (paragraph body). A row of card cells = a grid."""
    style = cell.get("style", "decimal")
    n = cell.get("n", 1)
    title = _clean(cell.get("title") or "")
    accent = cell.get("accent", "mint")
    if style == "roman":
        items = [_clean(it if isinstance(it, str) else (it.get("text") or ""))
                 for it in (cell.get("items") or []) if it]
        if not items and cell.get("text"):
            items = [_clean(cell["text"])]
        build.roman_card(s, int(n) if str(n).isdigit() else 1, x, y, w, max(0.8, h),
                         title=title, body_lines=items or [""], tone=tone, accent=accent,
                         subtitle=_clean(cell.get("subtitle")) if cell.get("subtitle") else None)
    else:
        build.decimal_card(s, n, x, y, w, max(0.8, h), title=title,
                           body_text=_clean(cell.get("text") or ""), tone=tone)

def c_image(s, cell, x, y, w, h, tone):
    """A picture fit within the cell rect, preserving aspect ratio, centred."""
    path = cell.get("path")
    if not (path and Path(path).exists()):
        build._rect(s, x, y, w, h, fill=C["green_2"] if tone == "dark" else C["off_white"])
        build._text(s, _clean(cell.get("alt") or "image"), x=x, y=y + h / 2 - 0.2, w=w, h=0.4,
                    font=Fz["sans_l"], size=11, color=_muted(tone), align="center", anchor="middle")
        return
    pic = s.shapes.add_picture(str(path), Inches(x), Inches(y))
    cell_w, cell_h = Inches(w), Inches(h)
    scale = min(cell_w / pic.width, cell_h / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int(Inches(x) + (cell_w - pic.width) / 2)
    pic.top = int(Inches(y) + (cell_h - pic.height) / 2)

CELL = {
    "prose": c_prose, "stat": c_stat, "list": c_list, "table": c_table,
    "panel": c_panel, "quote": c_quote, "heading": c_heading, "rule": c_rule,
    "callout": c_callout, "chips": c_chips, "card": c_card, "image": c_image,
}

# ---------------------------------------------------------------------------
# Scene layout (rows fill the body by weight; cells split the row by span)
# ---------------------------------------------------------------------------

def _render_scene(s, scene, x, y, w, h, tone):
    rows = [r for r in scene if r.get("cells")]
    if not rows:
        return
    gap_v = 0.22
    total_wt = sum(float(r.get("weight", 1.0)) for r in rows) or 1.0
    avail_h = h - gap_v * (len(rows) - 1)
    cy = y
    for r in rows:
        rh = avail_h * (float(r.get("weight", 1.0)) / total_wt)
        cells = r["cells"]
        gap_h = 0.40
        total_span = sum(float(c.get("span", 1.0)) for c in cells) or 1.0
        avail_w = w - gap_h * (len(cells) - 1)
        cx = x
        for c in cells:
            cw = avail_w * (float(c.get("span", 1.0)) / total_span)
            CELL.get(c.get("kind", "prose"), c_prose)(s, c, cx, cy, cw, rh, tone)
            cx += cw + gap_h
        cy += rh + gap_v

def r_scene(deck, plan, slide):
    scene = slide.get("scene", [])
    if not any(r.get("cells") for r in scene):
        # Empty scene degrades to a clean placeholder callout, never a bare header.
        s = build.blank(deck, tone=slide.get("tone", "light"))
        if slide.get("eyebrow"):
            build.eyebrow(s, _clean(slide["eyebrow"]), tone=slide.get("tone", "light"))
        note = _clean(slide.get("placeholder") or "Content to be added.")
        build.callout(s, title="PLACEHOLDER", body_text=note, x=ED_X, y=1.2, w=ED_W,
                      h=max(1.05, build.predict_callout_h(note, w=ED_W)),
                      tone=slide.get("tone", "light"))
        return s
    s, top, bot, tone = _chrome(deck, plan, slide)
    _render_scene(s, scene, ED_X, top, ED_W, max(1.0, bot - top), tone)
    return s

# ---------------------------------------------------------------------------
# Dedicated chrome slides (delegate to the polished build recipes)
# ---------------------------------------------------------------------------

def r_cover(deck, plan, slide):
    build.cover(deck, title=_clean(slide.get("title", "")),
                subtitle=_clean(slide.get("subtitle")) if slide.get("subtitle") else None,
                presenter=slide.get("presenter"), org=slide.get("org"), date=slide.get("date"),
                tone=slide.get("tone", "dark"),
                eyebrow_text=_clean(slide.get("eyebrow")) if slide.get("eyebrow") else None,
                themes=[_clean(t) for t in slide.get("themes", [])] or None)

def r_section(deck, plan, slide):
    build.section_divider(deck, number=slide.get("number", 1), title=_clean(slide.get("title", "")),
                          eyebrow_text=_clean(slide.get("eyebrow")) if slide.get("eyebrow") else None,
                          tone=slide.get("tone", "dark"),
                          lead=_clean(slide.get("lead")) if slide.get("lead") else None,
                          items=[_clean(i) for i in slide.get("items", [])] or None)

def r_closing(deck, plan, slide):
    build.thank_you(deck, title=_clean(slide.get("title", "Thank you.")),
                    subtitle=_clean(slide.get("subtitle")) if slide.get("subtitle") else None,
                    contacts=slide.get("contacts"), tone=slide.get("tone", "dark"))

KIND = {"cover": r_cover, "section": r_section, "divider": r_section,
        "closing": r_closing, "thank_you": r_closing, "scene": r_scene}

# ---------------------------------------------------------------------------
# Dash sweep + render
# ---------------------------------------------------------------------------

def _sweep_dashes(deck):
    n = 0
    def walk(shapes):
        nonlocal n
        for sh in shapes:
            if sh.shape_type == 6:
                walk(sh.shapes)
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if "—" in r.text or "–" in r.text:
                            r.text = _clean(r.text); n += 1
            if sh.has_table:
                for row in sh.table.rows:
                    for cl in row.cells:
                        for p in cl.text_frame.paragraphs:
                            for r in p.runs:
                                if "—" in r.text or "–" in r.text:
                                    r.text = _clean(r.text); n += 1
    for slide in deck.slides:
        walk(slide.shapes)
    return n

def render(plan, out_path, *, resolve=None, label_and_bake=True, audit=True):
    """Compose and render a story-led scene deck. `plan` is a dict (see module
    docstring) or a path to a .json file. Saves via build.save, so on Windows it
    runs the resolve pass, inherits the org sensitivity label and bakes
    fit-to-text (disable with label_and_bake=False)."""
    if isinstance(plan, (str, Path)):
        plan = json.loads(Path(plan).read_text(encoding="utf-8"))
    deck = build.new_deck()
    for slide in plan.get("slides", []):
        KIND.get(slide.get("kind", "scene"), r_scene)(deck, plan, slide)
    _sweep_dashes(deck)
    tmp = Path(tempfile.gettempdir()) / Path(out_path).name
    # build.save() inherits the org sensitivity label and bakes fit-to-text by
    # default (when the resolve pass runs); label_and_bake=False disables both.
    if label_and_bake:
        build.save(deck, str(tmp), resolve=resolve, audit=audit)
    else:
        build.save(deck, str(tmp), resolve=resolve, label_from=False, bake=False, audit=audit)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(tmp, out_path)
    try:
        os.remove(tmp)
    except OSError:
        pass
    return out_path

def main():
    p = argparse.ArgumentParser(description="Render a story-led scene plan into a CBRE deck.")
    p.add_argument("plan", help="path to the plan .json")
    p.add_argument("out", help="output .pptx path")
    p.add_argument("--no-resolve", action="store_true", help="skip the render-and-measure pass")
    p.add_argument("--no-label-bake", action="store_true",
                   help="skip the sensitivity-label inherit + autofit bake (no PowerPoint)")
    a = p.parse_args()
    out = render(a.plan, a.out, resolve=(False if a.no_resolve else None),
                 label_and_bake=not a.no_label_bake)
    print(f"deck -> {out}")

if __name__ == "__main__":
    main()
