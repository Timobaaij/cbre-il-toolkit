"""CBRE Corporate Deck Builder — v2.

Builds high-density, visually engaging CBRE-branded decks inspired by the
"CBRE - Slides I like" reference set. The previous version of this skill cloned
slide layouts from the official CBRE Full Layout Library; that approach
produced safe but sparse decks. This rewrite draws every shape from scratch,
so the output matches the editorial density and rhythm of the reference deck:

    - eyebrow tag + serif headline + intro paragraph
    - mint cards, gold-accented dark cards
    - Roman-numeral and decimal-numbered frameworks
    - mint-headed comparison tables with colored summary rows
    - inline KPI strips, callout boxes with dotted-mint underlines
    - white "CBRE" wordmark bottom-right, copyright footer bottom-left
    - aggressive dark/light alternation across the deck

The deck is 16:9 (13.33 x 7.50"). Every slide pattern is a single function call.

Quick start
-----------

    import sys
    from pathlib import Path
    # Portable: works in both Claude Code (Windows) and Claude.ai (Linux sandbox),
    # and for any user (derives the install path from the home directory).
    for _p in (
        Path("scripts"),
        Path.home() / ".claude/skills/cbre-corporate-pptx/scripts",
    ):
        if _p.exists():
            sys.path.insert(0, str(_p.resolve()))
            break
    import build

    deck = build.new_deck()

    build.cover(deck,
        title="Project Horizon",
        subtitle="A 5-year European industrial site selection strategy",
        presenter="T. Baaij | Senior Consultant",
        org="CBRE Supply Chain Advisory",
        date="MAY 2026")

    build.contents(deck, items=[
        "Our Understanding & Value Proposition",
        "CBRE Industrial & Logistics in Europe",
        "Location Selection",
        "Funding & Financing",
        "Project Management",
        "Project Team & Coordination",
    ])

    build.case_study(deck,
        eyebrow_text="CASE STUDY",
        title="CATL in Europe: The Cost of Getting Location Wrong",
        intro="CATL, the world's largest EV battery maker (~37% global share), "
              "selected Arnstadt, Germany for its first overseas gigafactory in "
              "2018, investing over $2 billion. The original vision: 100 GWh of "
              "European capacity from Germany. By 2024 the plant was capped at "
              "just 14 GWh, and expansion was cancelled.",
        framework_title="What went wrong in Germany",
        framework=[
            ("01", "Energy Costs", "German industrial energy among the highest in the EU. At plant level, margins eroded to point where scaling was unviable."),
            ("02", "Labour & Culture", "Cultural hurdles, labour disputes, strikes. Local workforce training slower and costlier than planned."),
            ("03", "Permitting", "Multi-year regulatory delays. Lawsuits, environmental reviews, and animal protection laws blocked progress."),
            ("04", "Demand Shift", "VW cut EV production nearby. The anchor customer thesis for Germany collapsed, killing the expansion case."),
        ],
        table_headers=["", "GERMANY", "HUNGARY"],
        table_rows=[
            ["Capacity",     "14 GWh (capped)",      "100 GWh"],
            ["Investment",   "~$2 billion",          "EUR 7.34 billion"],
            ["Energy",       "Among highest in EU",  "~40% lower than Germany"],
            ["Corporate Tax","~30% effective",       "9% (EU lowest)"],
            ["Labour",       "Strikes, friction",    "Favourable, flexible"],
            ["Permitting",   "Slow, litigious",      "Streamlined, govt. backed"],
            ["Outcome",      "Expansion cancelled",  "Primary European hub"],
        ],
        stat_strip_title="The cost of misalignment",
        stats=[
            ("86%",   "Capacity shortfall vs Germany plan"),
            ("EUR 7.3B", "EU investment lost"),
            ("4 yrs", "Behind schedule vs course-correction"),
        ],
        callout_title="LESSON FOR NIO",
        callout_body="Every location variable that forced CATL's costly pivot is "
                     "identifiable before ground was broken. Energy, permitting, "
                     "labour, incentives: none were surprises, they were blind "
                     "spots. For NIO's European manufacturing decision, "
                     "independent location advisory eliminates those blind spots "
                     "before the billion-euro lessons.",
        callout_tag="BG / RO / PL / HU SHORTLIST")

    build.thank_you(deck)
    build.save(deck, "Project-Horizon.pptx")


Pattern catalogue
-----------------

NOTE: every recipe below is keyword-only after `deck` (note the `*`), and the
eyebrow argument is named `eyebrow_text` (not `eyebrow`). Calling positionally
or with `eyebrow=` raises TypeError. The signatures here mirror the real `def`s
in this file — keep them in sync if you change a signature.

Cover & navigation
    cover(deck, *, title, subtitle=None, presenter=None, org=None,
          date=None, eyebrow_text=None, themes=None, tone="dark")
    contents(deck, *, items, eyebrow_text=None, title="Contents", tone="dark")
    section_divider(deck, *, number, title, eyebrow_text=None,
                    lead=None, items=None, tone="dark")
    thank_you(deck, *, title="Thank you.", subtitle=None, contacts=None,
              tone="dark")

Story & narrative
    case_study(deck, *, eyebrow_text, title, intro, framework_title,
               framework, table_headers, table_rows, stat_strip_title, stats,
               callout_title=None, callout_body=None, callout_tag=None,
               tone="dark")
    statement(deck, *, text, attribution=None, eyebrow_text=None,
              support=None, support_label=None, pillars=None, tone="dark")
    stat_hero(deck, *, eyebrow_text, title, stat, label, footnote=None,
              tone="dark")
    stat_strip(deck, *, eyebrow_text, title, subtitle=None, stats,
               body_text=None, coverage=None, tone="dark")

Framework grids (the "I, II, III..." cards we like)
    framework_roman(deck, *, eyebrow_text, title, items, intro=None,
                    side_callout=None, accent="mint", columns=3, tone="dark")

Tables & matrices
    worksheet_table(deck, *, eyebrow_text, title, intro, assumptions,
                    table_headers, table_rows, kpi_strip_items,
                    footnote=None, tone="light")
    comparison_table(deck, *, eyebrow_text, title, subtitle=None, columns,
                     sections, footer_label=None, footer_values=None,
                     tone="light")
    decision_matrix(deck, *, eyebrow_text, title, left_label, gate_label,
                    right_label, rows, callout_title=None, callout_body=None,
                    tone="dark")
    why_columns(deck, *, eyebrow_text, title, columns, intro=None,
                takeaway=None, takeaway_label=None, numbered=False,
                tone="light", accent="bright_green")
    why_two_col(deck, *, eyebrow_text, title, intro, drivers_label, drivers,
                right_eyebrow, cards, tone="dark")

Every pattern accepts `tone="dark"` or `tone="light"`. Where the references
have an obvious natural tone (case studies dark, worksheets light), the
default matches. Always mix tones across the deck (see SKILL.md rhythm rule).

Custom slides
-------------

For one-off slides, use `blank(deck, tone="dark")` which paints the
background, footer, and wordmark, then returns the slide for free drawing.
The low-level helpers (eyebrow, serif_title, body, roman_card, decimal_card,
kpi_block, kpi_strip, table, callout, etc.) are all exposed.
"""

from __future__ import annotations

from pathlib import Path
from dataclasses import dataclass, field
from typing import Iterable, Optional, Sequence, Tuple, Union, List, Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from lxml import etree


# ---------------------------------------------------------------------------
# Canvas
# ---------------------------------------------------------------------------

SLIDE_W = 13.333  # inches
SLIDE_H = 7.5
SAFE_L = 0.55
SAFE_R = SLIDE_W - 0.55
SAFE_T = 0.45
SAFE_B = SLIDE_H - 0.32

# ---------------------------------------------------------------------------
# Palette (extracted from the reference deck)
# ---------------------------------------------------------------------------

COLORS = {
    # CBRE core — palette sampled directly from the "CBRE - Slides I like" PDF.
    # The dark bg is a teal-shifted near-black, not the lifted #003F2D from the
    # corporate template. The "gold" accent is actually a pale wheat/cream, not
    # the warm orange-gold from the brand guide.
    "green":         RGBColor(0x01, 0x2A, 0x2C),   # primary dark teal-green bg
    "green_2":       RGBColor(0x10, 0x38, 0x38),   # lifted card surface
    "green_3":       RGBColor(0x18, 0x44, 0x40),   # callout / row-alt surface
    "green_4":       RGBColor(0x00, 0x38, 0x28),   # darker variant for callouts
    "bright_green":  RGBColor(0x17, 0xE8, 0x8F),   # lime accent — used sparingly
    "mint":          RGBColor(0x80, 0xB8, 0xA8),   # cooler mint/seafoam
    "mint_dark":     RGBColor(0x53, 0x8F, 0x86),   # deeper mint for body on light
    "mint_pale":     RGBColor(0xC0, 0xD0, 0xC8),   # pale mint
    "gold":          RGBColor(0xD8, 0xD8, 0x98),   # pale cream/wheat — the real "gold" of the reference deck
    "gold_warm":     RGBColor(0xCB, 0xA2, 0x58),   # legacy warm gold — kept for compatibility
    "blue":          RGBColor(0x38, 0x78, 0xA0),   # steel-blue accent for some card stripes

    # Neutral
    "white":         RGBColor(0xFF, 0xFF, 0xFF),
    "off_white":     RGBColor(0xF7, 0xF7, 0xF5),
    "page_light":    RGBColor(0xFF, 0xFF, 0xFF),
    "page_dark":     RGBColor(0x01, 0x2A, 0x2C),

    # Text on light
    "ink":           RGBColor(0x0C, 0x1C, 0x1E),   # near-black for body on white
    "ink_2":         RGBColor(0x43, 0x52, 0x54),   # dark teal body
    "charcoal":      RGBColor(0x7F, 0x84, 0x81),   # captions, footers
    "rule_light":    RGBColor(0xD7, 0xDB, 0xDA),   # row separators on light tables
    "rule_dark":     RGBColor(0x18, 0x48, 0x48),   # row separators on dark tables

    # ----- Official CBRE 2026 brand palette (v17) — ADDITIVE -----------------
    # The keys above are the editorial "Slides I like" set and remain the
    # DEFAULTS, so the deck's look is unchanged. The entries below add the
    # official brand colours and brand-named ALIASES so brand names resolve in
    # code. Skill-name -> official-brand-name mapping (canonical doc:
    # references/brand-guidelines.md):
    #   gold ≈ Wheat (#DBD99A)      mint ≈ Celadon (#80BBAD)
    #   mint_dark ≈ Sage (#538184)  ink_2 = Dark Grey (#435254, exact)
    #   charcoal ≈ Cement (#7F8480) bright_green = Accent Green (#17E88F, exact)
    #   green ≈ Dark Green (#012A2D; skill uses #012A2C — 1 digit, identical)
    # Missing official primaries (now AVAILABLE, but not defaults):
    "cbre_green":    RGBColor(0x00, 0x3F, 0x2D),   # #003F2D main brand / logo colour
    "midnight":      RGBColor(0x03, 0x28, 0x42),   # #032842 primary
    "light_grey":    RGBColor(0xCA, 0xD1, 0xD3),   # #CAD1D3 primary
    # Brand-named aliases (resolve to the existing editorial values):
    "wheat":         RGBColor(0xD8, 0xD8, 0x98),   # = gold
    "celadon":       RGBColor(0x80, 0xB8, 0xA8),   # = mint
    "sage":          RGBColor(0x53, 0x8F, 0x86),   # = mint_dark
    "dark_grey":     RGBColor(0x43, 0x52, 0x54),   # = ink_2  (official #435254)
    "cement":        RGBColor(0x7F, 0x84, 0x81),   # = charcoal (official #7F8480)
    "accent_green":  RGBColor(0x17, 0xE8, 0x8F),   # = bright_green (official #17E88F)
    "dark_green":    RGBColor(0x01, 0x2A, 0x2C),   # = green (official Dark Green #012A2D)
    # Approved tints (other tints must meet WCAG contrast):
    "midnight_tint": RGBColor(0x77, 0x8F, 0x9C),   # #778F9C
    "sage_tint":     RGBColor(0x96, 0xB3, 0xB6),   # #96B3B6
    "celadon_tint":  RGBColor(0xC0, 0xD4, 0xCB),   # #C0D4CB
    "wheat_tint":    RGBColor(0xEF, 0xEC, 0xD2),   # #EFECD2
    "cement_tint":   RGBColor(0xCB, 0xCD, 0xCB),   # #CBCDCB
    # Data-visualisation only — CHARTS & GRAPHS ONLY (see CHART_COLORS):
    "negative_red":  RGBColor(0xAD, 0x2A, 0x2A),   # #AD2A2A — negative values only
    "data_orange":   RGBColor(0xD2, 0x78, 0x5A),   # #D2785A
    "data_purple":   RGBColor(0x88, 0x50, 0x73),   # #885073
    "data_lpurple":  RGBColor(0xA3, 0x88, 0xBF),   # #A388BF
    "data_blue":     RGBColor(0x1F, 0x37, 0x65),   # #1F3765
    "data_lblue":    RGBColor(0x3E, 0x7C, 0xA6),   # #3E7CA6
    # NOTE: "blue" (#3878A0) above is OFF-BRAND (not a 2026 brand colour). It is
    # retained for back-compat; prefer "midnight" for card-stripe cycling.
}

# Hex strings without "#" prefix for python-pptx XML usage
def _hex(name: str) -> str:
    c = COLORS[name]
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


# Data-visualisation palette — CHARTS & GRAPHS ONLY (2026 guidelines).
# Ordered for categorical series. `negative_red` is deliberately NOT in this
# list: use it only to mark a negative value. Do not use these colours outside
# charts/graphs. There is no chart engine in this skill yet — these are
# constants for when one is added or a chart is hand-built.
CHART_COLORS = [
    COLORS["celadon"], COLORS["dark_grey"], COLORS["accent_green"],
    COLORS["wheat"], COLORS["data_orange"], COLORS["data_purple"],
    COLORS["data_lpurple"], COLORS["data_blue"], COLORS["data_lblue"],
    COLORS["light_grey"],
]


# ---------------------------------------------------------------------------
# Fonts (CBRE corporate; viewer falls back if not installed)
# ---------------------------------------------------------------------------

FONTS = {
    # NOTE: only weights that exist on a standard CBRE-installed Windows are
    # referenced here. "Financier Display Light" / "Calibre Bold" are NOT
    # commonly installed — use "Financier Display" + "Calibre Semibold" instead.
    "serif":    "Financier Display",     # editorial headlines (regular weight)
    "serif_m":  "Financier Display Medium",
    "serif_sb": "Financier Display Semibold",
    "serif_l":  "Financier Display",     # alias — no Light variant; renders as regular
    "sans":     "Calibre",                # body & UI
    "sans_l":   "Calibre Light",
    "sans_m":   "Calibre Medium",
    "sans_sb":  "Calibre Semibold",
    "sans_b":   "Calibre Semibold",       # alias — no Bold variant on CBRE Win
    "mono":     "Space Mono",             # dates, eyebrow tags, small UI labels
}

# Brand typography rules (2026 v17 guidelines) — see references/brand-guidelines.md.
#   Financier Display : headlines ONLY, MINIMUM 20 pt, TITLE CASE, never all caps.
#                       If serif < 20 pt is wanted, the brand says use Calibre.
#   Calibre           : sentence case (no forced title case); ALL CAPS only for
#                       short hero moments / eyebrows. Body leading 120%
#                       (skill default is 1.30 — kept to preserve the look).
#   Space Mono        : detail only — ALL CAPS, <=5 words, do NOT alter spacing.
#   Fallbacks (when brand fonts absent): Times -> Financier, Tahoma -> Calibre
#   (PowerPoint substitutes automatically; do not hard-swap in code).
FINANCIER_MIN_PT = 20


# ---------------------------------------------------------------------------
# Deck construction
# ---------------------------------------------------------------------------

def new_deck() -> Presentation:
    """Open a fresh blank 16:9 deck (13.33 x 7.50"). No master chrome — every
    slide paints its own background, footer, and CBRE wordmark."""
    deck = Presentation()
    deck.slide_width = Inches(SLIDE_W)
    deck.slide_height = Inches(SLIDE_H)
    # Tone log — each layout function appends its tone string here so
    # audit_tones() can verify the dark/light balance at the end.
    deck._cbre_tone_log = []  # type: ignore[attr-defined]
    return deck


def save(deck: Presentation, path: Union[str, Path],
         *, audit: bool = True,
         resolve: Optional[bool] = None,
         resolve_verbose: bool = True,
         bake: Optional[bool] = None,
         label_from: Optional[Union[str, Path]] = None) -> Path:
    """Save the deck. By default, also runs audit_tones() and prints the
    dark/light balance report — pass audit=False to suppress.

    Resolve pass
    ------------

    `resolve` controls the render-and-measure correction pass:

      resolve=None     — auto-detect. Runs the pass on Windows when
                         PowerPoint COM is available; skips on Linux /
                         when COM is unreachable (LibreOffice substitutes
                         fonts, so measurements would be misleading).
      resolve=True     — always run. Raises RuntimeError if COM is
                         unavailable.
      resolve=False    — skip. Identical to the pre-resolve behaviour.

    When resolve runs:
      1. The deck is saved to a draft path next to `path`.
      2. Each slide is rendered to PNG via `to_png.ps1 -SlideIndex N`.
      3. For each text element registered via Flow / CardFlow, the
         actual rendered height is pixel-walked and dependents are
         shifted to use the real height.
      4. The deck is re-saved (in-memory shape positions reflect the
         shifts) to the final `path`.
      5. Draft + PNGs are cleaned up.

    See `resolve_slide()` for the per-slide choreography.

    Sensitivity label + autofit bake (default ON when resolve runs)
    ---------------------------------------------------------------
    After the final save, the deck inherits the org sensitivity label and its
    text frames are baked to fit-to-text, so it opens labelled (no labelling
    prompt) and with every box already fitting its text. python-pptx writes
    spAutoFit (SHAPE_TO_FIT_TEXT) but PowerPoint only applies it when the text is
    edited, so without the bake boxes do not fit their text on first open.

      label_from=None (default) — inherit the bundled org label
                           (`assets/sensitivity_label.xml`) if present. This is a
                           file-level write (no PowerPoint), so it is cheap and
                           runs on every save, including resolve=False and Linux.
      label_from=<path>  — inherit the label from that LabelInfo.xml, or any file
                           already carrying an unencrypted MIP sensitivity label.
      label_from=False   — skip labelling.

      bake=None (default) — run the PowerPoint-COM autofit bake when the resolve
                           pass ran (Windows + COM) AND the file is labelled
                           (under a mandatory-labelling policy an unlabelled file
                           cannot be edited: E_ACCESSDENIED). So a plain
                           `save(deck, path)` on Windows labels + bakes;
                           `save(resolve=False)` stays fast (label only, no COM);
                           Linux labels only.
      bake=True / False  — force / skip the bake.

    The bake is anchor-aware: top-anchored frames are baked to fit-to-text;
    middle/bottom slots (table cells, centred headings) are pinned fixed so their
    alignment is preserved. Neither step ever raises: if PowerPoint or the label
    asset is unavailable the deck is left exactly as python-pptx saved it. See
    `apply_sensitivity_label()`, `_bake_autofit_com()` and `bake_autofit.ps1`.
    """
    out = Path(path).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)

    # Decide whether to resolve.
    if resolve is None:
        resolve = _resolve_available()
    elif resolve is True and not _resolve_available():
        raise RuntimeError(
            "save(resolve=True) was requested but PowerPoint COM is "
            "not available on this platform. Either run on Windows "
            "with PowerPoint installed, or pass resolve=False to skip "
            "the render-and-measure pass."
        )

    def _label_and_bake(target: Path):
        # Inherit the org sensitivity label by default: a file-level write (no
        # PowerPoint), policy-immune, so it runs on every save (including
        # resolve=False and Linux). label_from=None -> the bundled org asset if
        # present; label_from=False -> skip; a path -> use it.
        labelled = False
        src = label_from
        if src is None:
            _asset = Path(__file__).resolve().parent / "assets" / "sensitivity_label.xml"
            src = str(_asset) if _asset.exists() else None
        elif src is False:
            src = None
        if src:
            labelled = apply_sensitivity_label(target, src, verbose=resolve_verbose)
        # Bake fit-to-text. Default: run when the resolve pass ran (Windows + COM)
        # AND the file is labelled (an unlabelled file is uneditable under a
        # mandatory-labelling policy). bake=True forces; bake=False skips.
        do_bake = bake if bake is not None else (bool(resolve) and labelled)
        if do_bake:
            _bake_autofit_com(target, verbose=resolve_verbose)

    if not resolve:
        deck.save(str(out))
        _label_and_bake(out)
        if audit:
            audit_tones(deck, verbose=True)
            audit_line_of_sight(deck, verbose=True)
        return out

    # Resolve pass: save draft, render+measure each slide, re-save final.
    draft = out.with_suffix(".draft.pptx")
    png_dir = out.parent / ".cbre_resolve"
    png_dir.mkdir(exist_ok=True)
    resolve_succeeded = False
    try:
        deck.save(str(draft))
        any_shifts = False
        for idx, slide in enumerate(deck.slides, start=1):
            if not getattr(slide, "_cbre_resolve_plan", None):
                continue
            report = resolve_slide(slide, draft, idx, png_dir=png_dir,
                                    verbose=resolve_verbose)
            if any(abs(d) >= 0.008 for d in report.values()):
                any_shifts = True
        # Re-save with the corrected positions.
        deck.save(str(out))
        resolve_succeeded = True
        if resolve_verbose:
            if any_shifts:
                print(f"[resolve] applied shifts; final deck saved to {out}")
            else:
                print(f"[resolve] no significant shifts; deck saved to {out}")
    finally:
        # Clean up draft + PNGs only when resolve succeeded. If resolve
        # raised (e.g. CardOverflowError), preserve the draft .pptx and
        # the PNGs so the user can inspect the failing slide visually.
        if resolve_succeeded:
            try:
                if draft.exists():
                    draft.unlink()
                for png in png_dir.glob("*.png"):
                    png.unlink()
                if png_dir.exists() and not any(png_dir.iterdir()):
                    png_dir.rmdir()
            except OSError:
                pass
        else:
            # Print where the diagnostics live so the user can find them
            if resolve_verbose:
                print(f"[resolve] FAILED — draft preserved at {draft}",
                      file=__import__('sys').stderr)
                print(f"[resolve] FAILED — slide renders preserved in "
                      f"{png_dir}", file=__import__('sys').stderr)

    _label_and_bake(out)
    if audit:
        audit_tones(deck, verbose=True)
        audit_line_of_sight(deck, verbose=True)
    return out


def _resolve_available() -> bool:
    """True if we can run a render-and-measure resolve pass — i.e. we're
    on Windows with PowerShell + PowerPoint COM accessible."""
    import platform
    if platform.system() != "Windows":
        return False
    # We don't probe COM here — it's expensive and the to_png.ps1 invocation
    # in resolve_slide() will surface a clean error if PowerPoint isn't
    # installed. Treat Windows as "probably available".
    return True


def _bake_autofit_com(out_path: Union[str, Path], verbose: bool = True) -> bool:
    """Bake text-frame autofit into the saved deck via PowerPoint COM.

    python-pptx writes the autofit element (spAutoFit for SHAPE_TO_FIT_TEXT) but
    PowerPoint only recomputes the box geometry when a frame is edited, not on
    file open — so a freshly built deck shows boxes that do not fit their text
    until clicked. This opens the deck in PowerPoint, forces every text frame to
    apply its autofit (anchor-aware: top-anchored -> fit-to-text; middle/bottom
    slots pinned fixed), and re-saves, so the file is correct on first open.

    Windows + PowerPoint COM only. Returns True if the bake ran; on any other
    platform, or if PowerPoint is unavailable or the file is locked, it logs and
    returns False, leaving the python-pptx-saved deck intact. Never raises.
    """
    import platform
    import subprocess
    if platform.system() != "Windows":
        if verbose:
            print("[bake] skipped: autofit bake needs Windows + PowerPoint COM")
        return False
    script = Path(__file__).parent / "bake_autofit.ps1"
    if not script.exists():
        if verbose:
            print(f"[bake] skipped: {script.name} not found next to build.py")
        return False
    cmd = ["powershell", "-ExecutionPolicy", "Bypass", "-File", str(script),
           "-In", str(Path(out_path).resolve())]
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
    except Exception as e:  # noqa: BLE001
        if verbose:
            print(f"[bake] skipped (PowerPoint unavailable): {e}")
        return False
    if result.returncode != 0:
        if verbose:
            print(f"[bake] FAILED (exit {result.returncode}); deck left un-baked "
                  f"(open it once and save to apply autofit).\n"
                  f"  {result.stderr.strip()}")
        return False
    if verbose:
        print(f"[bake] {result.stdout.strip()} -> {out_path}")
    return True


def _load_label_xml(label_from: Union[str, Path]) -> Optional[bytes]:
    """Return the bytes of a docMetadata/LabelInfo.xml sensitivity-label part.

    `label_from` may be the LabelInfo.xml itself, or any .pptx/.docx/.xlsx that
    already carries an (unencrypted) MIP sensitivity label, from which the label
    part is extracted. Returns None if no label part is found."""
    import zipfile
    p = Path(label_from)
    if not p.exists():
        return None
    if p.suffix.lower() == ".xml":
        return p.read_bytes()
    try:
        with zipfile.ZipFile(str(p)) as z:
            if "docMetadata/LabelInfo.xml" in z.namelist():
                return z.read("docMetadata/LabelInfo.xml")
    except Exception:  # noqa: BLE001
        return None
    return None


def apply_sensitivity_label(pptx_path: Union[str, Path],
                            label_from: Union[str, Path],
                            verbose: bool = True) -> bool:
    """Inject an (unencrypted) MIP sensitivity label into a saved .pptx at the
    file level, so the deck opens as labelled (no labelling prompt) AND becomes
    editable under a mandatory-labelling policy, which is what lets the autofit
    bake (and any later automation) modify it.

    Adds the `docMetadata/LabelInfo.xml` part and the package-level
    `.../classificationlabels` relationship, matching how Office stores it. No
    encryption is involved, so this is a pure file write (no PowerPoint). Returns
    True if the label was added, False if already present or unavailable. Never
    raises: a failure leaves the deck unlabelled rather than losing it."""
    import re as _re
    import shutil as _shutil
    import zipfile
    PART = "docMetadata/LabelInfo.xml"
    REL_TYPE = ("http://schemas.microsoft.com/office/2020/02/relationships/"
                "classificationlabels")
    try:
        label_xml = _load_label_xml(label_from)
        if not label_xml:
            if verbose:
                print(f"[label] skipped: no label part in {label_from}")
            return False
        pptx_path = Path(pptx_path)
        with zipfile.ZipFile(str(pptx_path), "r") as zin:
            items = {n: zin.read(n) for n in zin.namelist()}
        if PART in items:
            return True  # already labelled (file is labelled afterward)
        items[PART] = label_xml
        rels = items.get("_rels/.rels", b"").decode("utf-8")
        existing = set(_re.findall(r'Id="(rId\d+)"', rels))
        n = 1
        while f"rId{n}" in existing:
            n += 1
        new_rel = f'<Relationship Id="rId{n}" Type="{REL_TYPE}" Target="{PART}"/>'
        rels = rels.replace("</Relationships>", new_rel + "</Relationships>")
        items["_rels/.rels"] = rels.encode("utf-8")
        tmp = pptx_path.with_suffix(".labeling.tmp")
        # [Content_Types].xml must lead the package; LabelInfo.xml rides the
        # default xml content type, matching how Office writes it.
        order = ["[Content_Types].xml"] + [k for k in items if k != "[Content_Types].xml"]
        with zipfile.ZipFile(str(tmp), "w", zipfile.ZIP_DEFLATED) as zout:
            for name in order:
                zout.writestr(name, items[name])
        _shutil.move(str(tmp), str(pptx_path))
        if verbose:
            print(f"[label] applied sensitivity label -> {pptx_path}")
        return True
    except Exception as e:  # noqa: BLE001
        if verbose:
            print(f"[label] skipped (could not apply): {e}")
        return False


def resolve_slide(slide, pptx_path: Union[str, Path], slide_index: int,
                  *, png_dir: Optional[Path] = None,
                  verbose: bool = True) -> dict:
    """Run the render-and-measure correction pass on a single slide.

    Workflow
    --------

    1. Render the slide to PNG via `scripts/to_png.ps1 -SlideIndex N`.
    2. Pixel-walk each ResolveElement in `slide._cbre_resolve_plan`,
       comparing actual rendered text height against the predictor's
       estimate (`h_predicted`).
    3. For each element where `|actual − predicted| ≥ 0.01"`, shift every
       dependent shape's top by `delta` (cumulative — later elements
       in the same Flow already received the prior deltas).
    4. For CardFlow members with a `card_box` constraint, after shift,
       verify the element still fits within `card_y + card_h −
       bottom_pad`; raise `CardOverflowError` if not.

    Parameters
    ----------
    slide : python-pptx slide
    pptx_path : path to the draft .pptx (must already be on disk)
    slide_index : 1-based slide number
    png_dir : where to write the rendered PNG. Defaults to a sibling
              `.cbre_resolve` folder of `pptx_path`.
    verbose : print per-element deltas when True.

    Returns
    -------
    dict[element_name → delta_inches]
    """
    pptx_path = Path(pptx_path)
    if png_dir is None:
        png_dir = pptx_path.parent / ".cbre_resolve"
    png_dir = Path(png_dir)
    png_dir.mkdir(exist_ok=True)

    plan = getattr(slide, "_cbre_resolve_plan", None) or []
    if not plan:
        return {}

    # 1. Render this slide to PNG.
    _render_single_slide(pptx_path, slide_index, png_dir)
    png_path = png_dir / f"{slide_index:02d}.png"
    if not png_path.exists():
        raise RuntimeError(
            f"resolve_slide: expected PNG at {png_path} after render but "
            f"none was written. Check that scripts/to_png.ps1 -SlideIndex "
            f"{slide_index} succeeds on this host."
        )

    # 2. Load and pixel-walk.
    from render_measure import RenderedSlide   # local import: optional dep
    rs = RenderedSlide(png_path, tolerance=20)

    # Slide-tone bg fallback for elements that didn't pass an explicit
    # bg_hint. Lets corner sampling skip the noisy local-sample step when
    # the element sits directly on the slide bg (most Flow.title / .body
    # calls). _paint_bg_native sets `slide._cbre_tone` to the tone name.
    slide_tone = getattr(slide, "_cbre_tone", None)
    slide_bg_rgb: Optional[Tuple[int, int, int]] = None
    if slide_tone == "dark":
        c = COLORS["page_dark"]
        slide_bg_rgb = (c[0], c[1], c[2])
    elif slide_tone == "light":
        c = COLORS["page_light"]
        slide_bg_rgb = (c[0], c[1], c[2])

    # Single threshold for both shape-resize and dependent-shift, so we
    # never enter a band where the box shrinks but downstream doesn't
    # move (which would silently open a sliver gap below the resized
    # shape). 0.008" = ~1 px at 120 DPI — below visual perception.
    RESOLVE_THRESHOLD = 0.008

    # Linked-height growth clamp: a linked bg rect can't grow past the
    # slide's safe-bottom band (where the wordmark + footer live). This
    # prevents a callout whose body actually under-predicted from
    # ballooning into the wordmark when its linked rect resizes downward.
    SAFE_BOT_LIMIT = SLIDE_H - 0.65

    # Probe height: default is h_predicted (the declared box) — safe but
    # misses cases where the predictor UNDER-estimates and text overflows
    # below the declared box. Callers that know their element is inside a
    # fixed-height container (e.g. callout body inside a callout) can set
    # `max_probe_h` to extend the measurement region without risk of
    # picking up neighbouring shapes' text. Without max_probe_h, a too-
    # generous probe drags in the next shape's text and the measurement
    # returns a spuriously large height that then shifts dependents the
    # wrong direction.
    report: dict = {}
    card_shifts: dict = {}     # cumulative shift per card_id
    for elem in plan:
        probe_h = elem.h_predicted
        if elem.max_probe_h is not None and elem.max_probe_h > probe_h:
            probe_h = elem.max_probe_h
        bg = elem.bg_hint if elem.bg_hint is not None else slide_bg_rgb
        bounds = rs.measure_text_bounds(elem.x, elem.y, elem.w, probe_h,
                                        bg=bg)
        if bounds is None:
            actual_h = elem.h_predicted    # no text detected — assume predictor
        else:
            actual_h = bounds.height
        delta = actual_h - elem.h_predicted
        report[elem.name] = delta

        # CardFlow constraint check (post-shift bottom must fit in card).
        # Keyed by `card_id` (a stable per-CardFlow integer) rather than
        # the card_box tuple — two cards with identical (x, y, w, h, pad)
        # geometry would collide on a tuple key.
        if elem.card_box is not None:
            cx, cy, cw, ch, bot_pad = elem.card_box
            key = elem.card_id if elem.card_id is not None else elem.card_box
            shift_in_card = card_shifts.get(key, 0.0)
            new_top = elem.y + shift_in_card
            new_bottom = new_top + actual_h
            limit = cy + ch - bot_pad
            if new_bottom > limit + 0.005:
                overshoot = new_bottom - limit
                raise CardOverflowError(
                    f"resolve_slide: '{elem.name}' renders at "
                    f"{actual_h:.3f}\" (predictor said "
                    f"{elem.h_predicted:.3f}\"); cumulative card content "
                    f"now reaches y={new_bottom:.3f}\", past the card's "
                    f"inner limit y={limit:.3f}\" by "
                    f"{overshoot:.3f}\". The card is too short for the "
                    f"actual rendered content. Reduce items in the grid "
                    f"(raises card_h), trim this card's copy, or split "
                    f"onto two slides."
                )
            card_shifts[key] = shift_in_card + delta

        # Resize the measured shape itself to match the actual rendered
        # text height. SHAPE_TO_FIT_TEXT (spAutoFit) is set on every text
        # shape, but python-pptx writes the originally-declared height to
        # the XML — PowerPoint's autofit only kicks in on click/edit, not
        # on file open. So a shape whose predictor under-estimated (e.g.
        # callout body declared h=0.27" but rendering 3 lines / 0.57")
        # would show the text overflowing the visible box rectangle when
        # the user opens the .pptx. Setting the height explicitly here
        # makes the saved file's box geometry honest, matching what the
        # autofit would do after a click.
        if abs(delta) >= RESOLVE_THRESHOLD and bounds is not None:
            elem.shape.height = Inches(actual_h)
            # Linked-height shapes (parent container rects, accent stripes,
            # etc.) absorb the same delta — so the parent rect grows when
            # the body grows, keeping the visible container consistent
            # with its rendered content. Clamp growth so a linked shape
            # can't cross the slide's safe-bottom band (wordmark + footer).
            if elem.linked_height_shapes:
                shift_emu = Inches(delta)
                for linked in elem.linked_height_shapes:
                    new_h_emu = linked.height + shift_emu
                    if delta > 0:
                        # Growing — clamp at safe-bottom
                        linked_top_in = linked.top / 914400
                        max_h_in = max(0.05, SAFE_BOT_LIMIT - linked_top_in)
                        max_h_emu = Inches(max_h_in)
                        if new_h_emu > max_h_emu:
                            new_h_emu = max_h_emu
                    linked.height = new_h_emu

            # Apply shift to every dependent shape (same threshold band so
            # the resized shape's bottom and downstream content stay in
            # lockstep — no sliver-drift gap).
            if elem.dependents:
                shift_emu = Inches(delta)
                for dep_shape in elem.dependents:
                    dep_shape.top = dep_shape.top + shift_emu

        if verbose and abs(delta) >= RESOLVE_THRESHOLD:
            print(f"[resolve] slide {slide_index}: {elem.name:50s} "
                  f"predicted {elem.h_predicted:.3f}\"  actual "
                  f"{actual_h:.3f}\"  delta {delta:+.3f}\" "
                  f"({len(elem.dependents)} dep)")

    return report


def _render_single_slide(pptx_path: Path, slide_index: int,
                         out_dir: Path) -> None:
    """Invoke scripts/to_png.ps1 to render a single slide to PNG.

    Raises RuntimeError if PowerShell / PowerPoint isn't available or the
    export fails. Blocks until the PNG is written.
    """
    import subprocess
    script = Path(__file__).parent / "to_png.ps1"
    if not script.exists():
        raise RuntimeError(f"resolve: to_png.ps1 not found at {script}")
    cmd = [
        "powershell", "-ExecutionPolicy", "Bypass",
        "-File", str(script),
        "-In", str(pptx_path),
        "-OutDir", str(out_dir),
        "-SlideIndex", str(slide_index),
    ]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
    if result.returncode != 0:
        raise RuntimeError(
            f"resolve: to_png.ps1 failed (exit {result.returncode}):\n"
            f"  STDOUT: {result.stdout.strip()}\n"
            f"  STDERR: {result.stderr.strip()}"
        )


def _log_tone(deck: Presentation, tone: str) -> None:
    """Record the tone of the slide that was just added. Safe to call even if
    the deck was not created through new_deck() (no-op in that case)."""
    log = getattr(deck, "_cbre_tone_log", None)
    if log is not None:
        log.append(tone)


def _blank_slide(deck: Presentation):
    """Add a slide using the blank layout (index 6 in the default python-pptx
    deck)."""
    layout = deck.slide_layouts[6]  # Blank
    slide = deck.slides.add_slide(layout)
    # Back-reference so _paint_bg can log the tone to the deck's tone log
    slide._cbre_deck = deck  # type: ignore[attr-defined]
    # Resolve plan: list of ResolveElement appended by Flow / CardFlow /
    # registered primitives. Used by resolve_slide() to drive the
    # render-and-measure correction pass. Pre-initialise so callers can
    # check `if slide._cbre_resolve_plan` without lazy-init shenanigans.
    slide._cbre_resolve_plan = []          # type: ignore[attr-defined]
    slide._cbre_tone = None                # type: ignore[attr-defined]
    return slide


def audit_tones(deck: Presentation, *, verbose: bool = True) -> dict:
    """Scan the tone log accumulated during deck construction and warn if the
    dark/light mix has drifted outside the recommended band.

    Recommended mix (for any deck of >= 6 slides):
      - 50%–70% dark slides
      - 30%–50% light slides (counting split-tone as 'dark' since the bulk
        of the visual area is dark)

    Returns a dict like:
      {
        "total": 15,
        "dark": 9,
        "light": 6,
        "dark_pct": 0.60,
        "light_pct": 0.40,
        "warnings": [...],   # human-readable strings
        "ok": True/False,
      }

    If verbose=True (default), warnings are printed to stdout. Either way they
    are returned in the dict so the caller can act on them.
    """
    log = list(getattr(deck, "_cbre_tone_log", []) or [])
    n = len(log)
    # Normalise: "split" counts as dark (top half is the visual anchor)
    dark = sum(1 for t in log if t in ("dark", "split"))
    light = sum(1 for t in log if t == "light")
    warnings: List[str] = []

    if n == 0:
        warnings.append("No slides logged. Did you build with build.new_deck()?")
    elif n < 3:
        warnings.append(f"Only {n} slide(s) — tone balance not meaningful yet.")
    else:
        dark_pct = dark / n
        light_pct = light / n
        if dark_pct > 0.80:
            warnings.append(
                f"Too dark: {dark}/{n} slides ({dark_pct:.0%}) are dark. "
                f"Recommended 50-70%. Consider switching a worksheet, "
                f"comparison_table, or stat_strip to tone='light'."
            )
        elif dark_pct < 0.50:
            warnings.append(
                f"Too light: only {dark}/{n} slides ({dark_pct:.0%}) are dark. "
                f"Recommended 50-70%. The CBRE deck should feel anchored in "
                f"dark teal-green - consider keeping cover, contents, "
                f"section_divider, case_study, statement, and thank_you on dark."
            )

    result = {
        "total": n,
        "dark": dark,
        "light": light,
        "dark_pct": (dark / n) if n else 0.0,
        "light_pct": (light / n) if n else 0.0,
        "log": log,
        "warnings": warnings,
        "ok": not warnings,
    }
    if verbose:
        # ASCII-only output — Windows consoles default to cp1252 which can't
        # encode unicode glyphs without explicit reconfiguration.
        print(f"[tone audit] {n} slides - dark={dark} ({result['dark_pct']:.0%}), "
              f"light={light} ({result['light_pct']:.0%})")
        # If any slide is split-tone, break it out so the user can see it
        split_n = sum(1 for t in log if t == "split")
        if split_n:
            print(f"             (of which {split_n} split-tone, counted as dark)")
        for w in warnings:
            print(f"  [warn] {w}")
        if not warnings and n >= 3:
            print("  [ok] Tone mix looks balanced.")
    return result


# ---------------------------------------------------------------------------
# Low-level drawing primitives
# ---------------------------------------------------------------------------

def _assert_pos_dims(w, h, where: str) -> None:
    """Guard against non-positive shape dimensions.

    A shape (textbox / autoshape / rect) with w<=0 or h<=0 is accepted by
    python-pptx but makes PowerPoint reject the WHOLE file as "corrupted and
    unreadable" (HRESULT 0x80070570) on open — the resolve/render pass then
    fails with no pointer to the culprit. The usual cause is a back-solved
    width like `w = panel_w - bar_w - gap` going negative when a sibling is at
    its max. Fail loudly here instead. (Zero-height *connectors* from `_line`
    are fine — those go through add_connector, not this path.)
    """
    if w <= 0 or h <= 0:
        raise ValueError(
            f"{where}: width and height must be positive (got w={w:.3f}, "
            f"h={h:.3f}). A non-positive dimension produces XML that PowerPoint "
            f"rejects as a corrupt file. Check any back-solved w/h that can go "
            f"negative — put variable-width siblings in a fixed column rather "
            f"than the leftover space.")


def _rect(slide, x, y, w, h, *, fill: Optional[RGBColor] = None,
          line: Optional[RGBColor] = None, line_w: float = 0.75):
    """Draw a filled rectangle. x/y/w/h in inches."""
    _assert_pos_dims(w, h, "_rect")
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


_SHRINK_BANNED_MSG = (
    "shrink=True / mode='shrink' / TEXT_TO_FIT_SHAPE / normAutoFit are "
    "BANNED in this skill. The library always uses SHAPE_TO_FIT_TEXT "
    "(spAutoFit) — the box grows to fit the text, the font never shrinks. "
    "Minimum font size is 9 pt. If your content overflows a container, "
    "the fix is to grow the container or trim the copy (drop a row, split "
    "into two slides, choose a denser composition). Other Python "
    "presentation libraries expose this mode — this one deliberately does "
    "not. See SKILL.md 'STOP — read this before writing any slide code'."
)


def _reject_banned_kwargs(kwargs: dict, fn_name: str) -> None:
    """Raise a loud, skill-specific TypeError if a caller passes any banned
    argument (shrink, mode='shrink', etc.). Generic Python TypeErrors don't
    teach the model why the argument is wrong — this one does."""
    if "shrink" in kwargs:
        raise TypeError(
            f"{fn_name}() does not accept `shrink`. {_SHRINK_BANNED_MSG}"
        )
    if kwargs.get("mode") == "shrink":
        raise TypeError(
            f"{fn_name}(mode='shrink') is not supported. {_SHRINK_BANNED_MSG}"
        )
    leftover = [k for k in kwargs if k != "mode"]
    if leftover:
        raise TypeError(
            f"{fn_name}() got unexpected keyword argument(s): {leftover}"
        )


def _text(slide, text: Union[str, List[str]], x, y, w, h, *,
          font: str = FONTS["sans"], size: float = 10,
          color: RGBColor = None, bold: bool = False, italic: bool = False,
          align: str = "left", anchor: str = "top",
          line_spacing: float = 1.15, letter_spacing: Optional[float] = None,
          uppercase: bool = False, fit: bool = False, autofit: bool = True,
          **_banned_kwargs):
    """Add a text box with the given properties. `text` may be a string or
    list of strings (one paragraph each).

    align: "left" | "center" | "right"
    anchor: "top" | "middle" | "bottom"
    letter_spacing: in 1/100 of a point if you want tracking

    Auto-size modes (mutually exclusive):
      fit=True     -> MSO_AUTO_SIZE.NONE  (box and font fixed — use only for
                      short, visually validated text like hero numerals).
      autofit=True -> MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT  (DEFAULT)
                      Box grows to fit the text. Font never shrinks.
                      Minimum font size is 9 pt.

    TEXT_TO_FIT_SHAPE ("shrink text on overflow") is not available — it
    produces unreadably small text and still causes visual overflow.
    If content overflows a container, increase the container or restructure
    the content (drop a row, split slides, choose a denser layout).
    """
    _reject_banned_kwargs(_banned_kwargs, "_text")
    _assert_pos_dims(w, h, "_text")
    color = color or COLORS["ink_2"]
    size = max(size, 9)  # Hard floor — never render below 9 pt.
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                 "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
    anchor_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                  "bottom": MSO_ANCHOR.BOTTOM}

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor_map[anchor]
    if fit:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    elif autofit:
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    paras = text if isinstance(text, list) else [text]
    for i, t in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_map[align]
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = (t.upper() if uppercase else t)
        f = run.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
        if letter_spacing is not None:
            # Apply character spacing via raw XML (python-pptx doesn't expose it)
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(int(letter_spacing * 100)))
    return box


def container_text(slide, text: Union[str, List[str]], x, y, w, h, *,
                   font: str = FONTS["sans"], size: float = 10,
                   color: RGBColor = None, bold: bool = False,
                   italic: bool = False, align: str = "left",
                   anchor: str = "top", line_spacing: float = 1.20,
                   letter_spacing: Optional[float] = None,
                   uppercase: bool = False, mode: str = "autofit",
                   **_banned_kwargs):
    """Text inside a container (card body, table cell, callout, KPI strip
    caption, status pill, panel body).

    Modes:
      "autofit" (default) -> MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT.
                             The box grows downward to fit the text.
                             If it overflows the container, increase the
                             container height or reduce copy length — never
                             shrink the font.
      "fit"               -> MSO_AUTO_SIZE.NONE.
                             Box and font both fixed. Use only for short,
                             known-fitting content (single words, hero stat
                             numerals, status pills) where size is already
                             visually validated.

    Font size is enforced at a minimum of 9 pt regardless of mode.
    PowerPoint's TEXT_TO_FIT_SHAPE ("shrink text on overflow") is not
    available — if a container overflows, resize the container or trim
    copy.

    All other args are identical to `_text`.
    """
    # Reject shrink/banned kwargs with a skill-specific error BEFORE the
    # mode check, so mode='shrink' produces the teaching message rather
    # than a generic ValueError.
    _banned_with_mode = dict(_banned_kwargs)
    if mode == "shrink":
        _banned_with_mode["mode"] = "shrink"
    _reject_banned_kwargs(_banned_with_mode, "container_text")
    if mode not in ("autofit", "fit"):
        raise ValueError(
            f"container_text mode must be 'autofit' or 'fit', got {mode!r}."
        )
    return _text(slide, text, x=x, y=y, w=w, h=h,
                 font=font, size=size, color=color, bold=bold, italic=italic,
                 align=align, anchor=anchor,
                 line_spacing=line_spacing, letter_spacing=letter_spacing,
                 uppercase=uppercase,
                 autofit=(mode == "autofit"),
                 fit=(mode == "fit"))


def assert_within(child_box, parent_x, parent_y, parent_w, parent_h,
                  *, name: str = "text"):
    """Sanity check: assert that a child shape sits fully inside a parent
    rectangle. Raises AssertionError with a clear message if not.

    Useful as a guard right after drawing text into a card body or table
    cell — catches the silent-overflow bug at build time:

        card_rect = _rect(s, cx, cy, cw, ch, fill=...)
        cap = container_text(s, caption, x=cx+0.20, y=cy+0.65, w=cw-0.4,
                             h=0.35, size=9)
        assert_within(cap, cx, cy, cw, ch, name="card caption")

    Compares declared box bounds, not rendered text bounds. Because all
    text helpers now autofit (SHAPE_TO_FIT_TEXT), the rendered box grows
    beyond the declared `h` when text wraps — this check catches the
    *declared* footprint, not the post-render footprint. It's still useful
    as a build-time guard against passing nonsensical coordinates, but it
    will NOT catch silent vertical growth. For overflow on real renders,
    use the visual QA pass (PNG/PDF export) called out in SKILL.md.
    """
    cx = child_box.left / 914400.0   # EMU -> inches
    cy = child_box.top / 914400.0
    cw = child_box.width / 914400.0
    ch = child_box.height / 914400.0
    overflows = []
    if cx < parent_x - 0.001:
        overflows.append(f"left {cx:.3f} < parent {parent_x:.3f}")
    if cy < parent_y - 0.001:
        overflows.append(f"top {cy:.3f} < parent {parent_y:.3f}")
    if cx + cw > parent_x + parent_w + 0.001:
        overflows.append(
            f"right {cx + cw:.3f} > parent {parent_x + parent_w:.3f}"
        )
    if cy + ch > parent_y + parent_h + 0.001:
        overflows.append(
            f"bottom {cy + ch:.3f} > parent {parent_y + parent_h:.3f}"
        )
    if overflows:
        raise AssertionError(
            f"{name} overflows parent: " + "; ".join(overflows)
        )


def grid_card_geometry(*, available_h: float, n_items: int,
                       row_gap: float, inner_content_h: float,
                       top_pad: float = 0.10, bottom_pad: float = 0.10,
                       name: str = "grid") -> float:
    """Compute the per-card height for an N-item single-column grid and
    raise AssertionError if the math would force the cards to cram.

    A "crammed" card is one where the floating card height (computed by
    dividing the available area among N items) leaves less than the
    requested `bottom_pad` of clearance below the inner content. This is
    the silent failure mode that produced the over-tight Approach: text on
    Stellantis Part 2 slide 7 — the cards "fit" mathematically but read
    as squashed because no minimum bottom margin was enforced.

    Args:
        available_h: total vertical space the grid occupies (e.g. SAFE_BOT
            minus the post-Flow content_top).
        n_items: number of cards in the column.
        row_gap: inter-card gap (inches).
        inner_content_h: height (in inches) that the card's inner content
            consumes, measured from the card's interior top — i.e. the
            distance from the card's top edge to the bottom of the last
            text element, INCLUDING the desired `top_pad`.
        top_pad: desired inner top padding (advisory only, baked into the
            inner_content_h figure the caller provides).
        bottom_pad: required clearance below the last inner element. The
            assertion fails if the computed card_h leaves less than this.
        name: label included in the AssertionError message.

    Returns:
        card_h in inches.

    Example:
        # Six ranked cards, each holding title (0.08+0.24) + body
        # (gap 0.02 + 0.20) + approach (gap 0.04 + 0.16) = 0.74" of inner
        # content. Demand 0.12" bottom clearance.
        card_h = grid_card_geometry(
            available_h=4.82, n_items=6, row_gap=0.06,
            inner_content_h=0.74, bottom_pad=0.12,
            name="opportunity grid",
        )
        # AssertionError: opportunity grid: 6 cards in 4.82" with 0.06" gaps
        # → card_h=0.753" leaves only 0.013" below 0.74" of content (need 0.12").
    """
    card_h = (available_h - (n_items - 1) * row_gap) / n_items
    headroom = card_h - inner_content_h
    if headroom < bottom_pad - 0.005:
        raise AssertionError(
            f"{name}: {n_items} cards in {available_h:.2f}\" with "
            f"{row_gap:.2f}\" gaps -> card_h={card_h:.3f}\" leaves only "
            f"{headroom:.3f}\" below {inner_content_h:.2f}\" of content "
            f"(need {bottom_pad:.2f}\"). Fix by: (a) reducing n_items, "
            f"(b) freeing space above the grid (shorter title / drop "
            f"intro body), (c) restructuring to a 2-column grid, or "
            f"(d) splitting onto two slides."
        )
    return card_h


def _line(slide, x1, y1, x2, y2, *, color: RGBColor, width_pt: float = 1.5,
          dash: bool = False):
    """Draw a straight line between (x1,y1) and (x2,y2) in inches."""
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1),
                                    Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width_pt)
    if dash:
        # round dot dash via XML
        lnEl = ln.line._get_or_add_ln()
        prstDash = lnEl.find(qn("a:prstDash"))
        if prstDash is None:
            prstDash = etree.SubElement(lnEl, qn("a:prstDash"))
        prstDash.set("val", "sysDot")
    return ln


def _vbar(slide, x, y, h, *, color: RGBColor, width_in: float = 0.04):
    """Thin vertical accent bar."""
    return _rect(slide, x, y, width_in, h, fill=color)


def _hbar(slide, x, y, w, *, color: RGBColor, height_in: float = 0.04):
    return _rect(slide, x, y, w, height_in, fill=color)


def _dotted_line(slide, x, y, w, *, color: RGBColor = None):
    """Mint dotted underline (signature CBRE callout element)."""
    color = color or COLORS["mint"]
    n_dots = int(w / 0.085)
    for i in range(n_dots):
        cx = x + i * 0.085
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(y),
                                   Inches(0.035), Inches(0.035))
        c.fill.solid()
        c.fill.fore_color.rgb = color
        c.line.fill.background()


# ---------------------------------------------------------------------------
# Line of Sight — the signature CBRE graphic device (2026 v17 guidelines)
# ---------------------------------------------------------------------------
# A single rule that reinforces DEPTH (vertical) or BREADTH (horizontal).
# Brand rules encoded here:
#   * weight from the preset set, in px (== pt: "pixel sizing is synonymous
#     with point sizing"); on a slide use 1-20.
#   * length >= the brand minimum for that weight;
#   * AT MOST ONE per layout (slide) — overuse is explicitly off-brand.
# Forms (see references/brand-guidelines.md): standalone, type-connector,
# border. Portal (picture-in-picture) is documented but not built — needs imagery.
LINE_OF_SIGHT_WEIGHTS_PX = (1, 2, 5, 10, 20, 50, 100, 200)
_LOS_MIN_LEN_PX = {1: 100, 2: 100, 5: 100, 10: 200, 20: 200,
                   50: 300, 100: 400, 200: 600}


def line_of_sight(slide, *, orientation: str = "horizontal",
                  x: float, y: float, length: float,
                  weight_px: int = 5, color: Optional[RGBColor] = None,
                  tone: str = "dark", form: str = "standalone",
                  strict: bool = False):
    """Draw a CBRE 'Line of Sight' rule and return the shape.

    Args:
        orientation: 'horizontal' (breadth) or 'vertical' (depth).
        x, y: top-left of the line, in inches.
        length: line length (the long dimension), in inches.
        weight_px: thickness from {1,2,5,10,20,50,100,200} (px == pt). On a
            slide, use 1-20; thicker weights are for large-format layouts.
        color: override; defaults to a high-contrast accent per tone
            (Accent Green on dark — used sparingly; CBRE Green on light).
        tone: 'dark' | 'light' (selects the default colour).
        form: advisory metadata — 'standalone' | 'type-connector' | 'border'.
        strict: if True, a 2nd line on the same slide RAISES; else warns.

    Brand guards: weight in the preset set; length >= the per-weight minimum;
    at most ONE Line of Sight per slide. See references/brand-guidelines.md.
    """
    if orientation not in ("horizontal", "vertical"):
        raise ValueError("orientation must be 'horizontal' or 'vertical'")
    if weight_px not in _LOS_MIN_LEN_PX:
        raise ValueError(
            f"weight_px={weight_px} is not a CBRE Line-of-Sight weight. Use "
            f"one of {LINE_OF_SIGHT_WEIGHTS_PX} (px == pt).")
    min_len_in = _LOS_MIN_LEN_PX[weight_px] / 72.0
    if length < min_len_in - 1e-6:
        raise ValueError(
            f"Line of Sight at weight {weight_px}px must be >= "
            f"{_LOS_MIN_LEN_PX[weight_px]}px ({min_len_in:.2f}\") long; got "
            f"{length:.2f}\". A stub reads as a dash/button — lengthen it or "
            f"drop to a thinner weight.")
    # Max one per layout.
    count = getattr(slide, "_cbre_los_count", 0)
    if count >= 1:
        msg = ("Second Line of Sight on one slide: the CBRE rule is AT MOST "
               "ONE per layout (more reads as frenetic / off-brand). Keep one "
               "or split across slides.")
        if strict:
            raise ValueError(msg)
        import sys
        print(f"[brand] {msg}", file=sys.stderr)
    try:
        slide._cbre_los_count = count + 1
    except Exception:
        pass

    if color is None:
        color = COLORS["accent_green"] if tone == "dark" else COLORS["cbre_green"]
    weight_in = weight_px / 72.0
    if orientation == "horizontal":
        return _rect(slide, x, y, length, weight_in, fill=color)
    return _rect(slide, x, y, weight_in, length, fill=color)


def audit_line_of_sight(deck: Presentation, *, verbose: bool = True) -> dict:
    """Summarise Line-of-Sight usage per slide at save time. The draw-time
    guard already holds it to <=1 per slide unless strict was overridden."""
    counts = [getattr(s, "_cbre_los_count", 0) for s in deck.slides]
    over = sum(1 for c in counts if c > 1)
    total = sum(counts)
    if verbose and total:
        msg = (f"[line-of-sight] {total} across "
               f"{sum(1 for c in counts if c)} slide(s)")
        if over:
            msg += f"; WARNING: {over} slide(s) exceed the 1-per-layout rule"
        print(msg)
    return {"per_slide": counts, "total": total, "slides_over": over}


# ---------------------------------------------------------------------------
# Text measurement and flow stacking
# ---------------------------------------------------------------------------
#
# Why this exists: PowerPoint's SHAPE_TO_FIT_TEXT resizes the SHAPE based on
# text content, but python-pptx has no way to read the post-render height.
# So if a layout hardcodes the y of the NEXT element as "title.y + title.h",
# and the title autofits to less than its declared h, a dead-space gap opens
# below the title.
#
# `measure_text` estimates the rendered height of a string within a fixed-
# width box using a calibrated per-font character-width model. It is not
# pixel-perfect (it slightly over-estimates for safety so flow text never
# overflows downstream elements) but it is close enough that vertical
# rhythm composed via the `Flow` class lands without visible dead space.

# Average character width factor per font family, expressed as a fraction
# of the font size in points. Calibrated empirically against rendered output
# from PowerPoint on Windows with CBRE's licensed fonts installed.
#
# Recalibrated 2026-05-19 against Stellantis I&L Brief renders: previous
# serif factor (0.50) over-predicted Financier Display width by ~16%, causing
# 1-line titles to be reserved as 2-line boxes and creating ~0.35" of dead
# space below short titles in Flow layouts. The current values are tuned
# so a 95-char title at 22pt and a 190-char body at 11pt both predict the
# same 1-line / 2-line outcome PowerPoint actually renders.
#
# Combined with Flow.title()/body() rendering with autofit=True, this gives
# visual safety against the residual ~1 char over-prediction at the wrap
# boundary: the *cursor* may advance a touch too far, but the *shape* itself
# stays tight to the rendered text — no opaque dead band inside the block.
_CHAR_W_FACTORS = {
    "serif":    0.40,   # Financier Display — editorial serif (was 0.50)
    "serif_m":  0.40,
    "serif_sb": 0.42,
    "serif_l":  0.40,
    "sans":     0.44,   # Calibre — moderately compressed humanist sans
    "sans_l":   0.41,   # Calibre Light — narrower
    "sans_m":   0.44,
    "sans_sb":  0.46,   # Calibre Semibold — slightly wider
    "sans_b":   0.46,
    "mono":     0.60,   # Space Mono — wide monospace
}


def measure_text(text: Union[str, List[str]], *, size: float, w: float,
                 font: str = "sans", line_spacing: float = 1.20,
                 uppercase: bool = False,
                 letter_spacing: Optional[float] = None) -> float:
    """Estimate the rendered height (in inches) of `text` inside a box of
    width `w` inches at `size` pt with the given font and line spacing.

    `font` may be either a key into FONTS (e.g. "sans_l") or the resolved
    font name (e.g. "Calibre Light"); both work.

    The estimate is calibrated to be slightly conservative — for typical
    deck typography it lands within ±0.04" of the actual rendered height,
    with a bias toward over-estimation so flow layouts never overlap.

    Returns 0.0 for empty input.
    """
    if not text:
        return 0.0
    # Resolve font key -> family
    family = font.lower()
    factor = None
    for key, val in _CHAR_W_FACTORS.items():
        if key in family or FONTS[key].lower() in family:
            factor = val
            break
    if factor is None:
        factor = 0.46  # default to sans
    if uppercase:
        factor *= 1.05   # caps are slightly wider on average
    # Letter spacing (1/100 pt) widens each char
    extra_w_pt = (letter_spacing or 0.0)
    char_w_in = (size * factor + extra_w_pt) / 72.0
    line_h_in = size * line_spacing / 72.0

    paras = text if isinstance(text, list) else text.split("\n")
    total_lines = 0
    for para in paras:
        if not para:
            total_lines += 1
            continue
        # Conservative: assume worst-case word-wrap (chars per line ≈ box-width)
        chars_per_line = max(1, int(w / char_w_in))
        n_chars = len(para)
        para_lines = max(1, -(-n_chars // chars_per_line))
        total_lines += para_lines

    # Total = lines * line height + small descender padding
    # 0.06 (was 0.10) — recalibrated alongside _CHAR_W_FACTORS to reduce
    # cumulative dead space across stacked Flow elements.
    descender_pad = size / 72.0 * 0.06
    return total_lines * line_h_in + descender_pad


# ---------------------------------------------------------------------------
# Resolve plan — render-and-measure layout corrections
# ---------------------------------------------------------------------------
#
# Every text shape that participates in a Flow / CardFlow / registered direct
# primitive call gets a ResolveElement appended to `slide._cbre_resolve_plan`.
# After the deck is built, `build.save(resolve=True)` (or `resolve_slide()`
# directly) renders each slide to PNG, pixel-walks the actual rendered text
# height for each registered element, and shifts every dependent shape by
# `actual − predicted`. The predictor (`measure_text`) becomes a first-draft
# seed; PowerPoint's render is the ground truth.

@dataclass
class ResolveElement:
    """One text shape participating in the resolve pass.

    shape         — the python-pptx shape (the textbox) to measure
    x, y, w       — declared box position/width in inches
    h_predicted   — predictor's height estimate; what dependents were positioned against
    dependents    — list of python-pptx shapes that move when h_actual ≠ h_predicted
    name          — diagnostic label (printed in the resolve report)
    bg_hint       — optional explicit (R, G, B) bg colour for pixel-walk; if None,
                    render_measure samples the local bg from box corners
    card_box      — for CardFlow members: (cx, cy, cw, ch, bottom_pad). After
                    shift, resolve enforces `top + actual_h ≤ cy + ch − bottom_pad`
                    and raises CardOverflowError if not.
    max_probe_h   — optional explicit ceiling for the resolve probe height.
                    Use when h_predicted is known to UNDER-estimate (e.g. a
                    callout body inside a fixed-height container). Default
                    None ⇒ resolve probes only up to h_predicted, which is
                    safe but misses under-prediction. Set this to the
                    container's inner height to catch text that overflows
                    the declared box without dragging neighbouring shapes
                    into the measurement.
    linked_height_shapes — list of shapes whose .height adjusts by the same
                    delta as this element's actual_h − h_predicted. Use this
                    to keep a parent container (callout bg rect, card surface
                    rect) in sync with its child text shape — so when resolve
                    resizes the body to its real height, the parent rect
                    grows or shrinks to match. Without this, the child
                    textbox resizes but the parent rect keeps its declared
                    height, leaving the child either overflowing below the
                    parent or floating with empty parent space below.
    card_id        — stable integer id for the card this element belongs to.
                    Resolve keys `card_shifts` by this id (not the card_box
                    tuple) so two cards on the same slide with identical
                    geometry don't collide in the cumulative-shift dict.
                    Set automatically by CardFlow via id(self).
    """
    shape: Any
    x: float
    y: float
    w: float
    h_predicted: float
    dependents: List[Any] = field(default_factory=list)
    name: str = ""
    bg_hint: Optional[Tuple[int, int, int]] = None
    card_box: Optional[Tuple[float, float, float, float, float]] = None
    max_probe_h: Optional[float] = None
    linked_height_shapes: List[Any] = field(default_factory=list)
    card_id: Optional[int] = None


def _resolve_plan(slide) -> List[ResolveElement]:
    """Return the slide's resolve plan list, creating it on first access."""
    plan = getattr(slide, "_cbre_resolve_plan", None)
    if plan is None:
        plan = []
        slide._cbre_resolve_plan = plan
    return plan


def _register_resolve_element(slide, shape, *, x: float, y: float, w: float,
                              h_predicted: float, name: str = "",
                              prior_elements: Optional[List[ResolveElement]] = None,
                              bg_hint: Optional[Tuple[int, int, int]] = None,
                              card_box: Optional[Tuple[float, float, float, float, float]] = None,
                              max_probe_h: Optional[float] = None,
                              linked_height_shapes: Optional[List[Any]] = None,
                              card_id: Optional[int] = None,
                              ) -> ResolveElement:
    """Append a new ResolveElement to the slide's plan, and add its shape as
    a dependent of every prior element in `prior_elements` (so when an earlier
    element grows by delta, this one shifts down by delta).

    Returns the new element so the caller (Flow / CardFlow) can append it to
    its own local element list.
    """
    elem = ResolveElement(shape=shape, x=x, y=y, w=w, h_predicted=h_predicted,
                          name=name, bg_hint=bg_hint, card_box=card_box,
                          max_probe_h=max_probe_h,
                          linked_height_shapes=list(linked_height_shapes or []),
                          card_id=card_id)
    if prior_elements:
        for prior in prior_elements:
            prior.dependents.append(shape)
    _resolve_plan(slide).append(elem)
    return elem


class Flow:
    """Vertical stacking helper. Tracks a y-cursor that advances by the
    MEASURED content height after each element — eliminating the dead-
    space problem where a fixed-height title box (autofit-shrunk to its
    actual content) leaves a gap below it because the next element's y
    was hardcoded.

    Usage:

        f = build.Flow(s, x=0.55, y=0.55, w=12.15, tone="light")
        f.eyebrow("ACQUISITION FRAMEWORK")
        f.gap(0.25)
        f.title("Six routes assessed. Three progressed.", size=32)
        f.gap(0.20)
        f.body("Six structures sit on the menu for a 15,000 sqm...")
        f.gap(0.35)
        # f.y is now exactly below the body text + 0.35" gap
        # ... build the table at y=f.y ...

    The cursor `f.y` is always pointing at the next available y after
    the last element + gap. To anchor an element at an absolute y
    (e.g. a bottom callout), use `f.gap_to(6.02)` to jump the cursor.
    """

    def __init__(self, slide, *, x: float, y: float, w: float,
                 tone: str = "dark"):
        self.slide = slide
        self.x = x
        self.y = y
        self.w = w
        self.tone = tone
        # Resolve registration: every text element added by this Flow gets
        # tracked here. When a NEW element is added, it becomes a dependent
        # of every prior one — so if an earlier element's actual height
        # differs from predicted, this one shifts to compensate.
        self._elements: List[ResolveElement] = []

    # ---- cursor controls ----

    def gap(self, h: float) -> "Flow":
        """Advance cursor by h inches (explicit blank space)."""
        self.y += h
        return self

    def gap_to(self, y: float) -> "Flow":
        """Jump cursor to absolute y in inches."""
        self.y = y
        return self

    def reserve_bottom(self, h: float) -> float:
        """Return the y at which content of height `h` would bottom out at
        the safe-bottom edge. Useful for anchoring a callout above the
        wordmark band:
            callout_y = f.reserve_bottom(0.78)  # callout 0.78" tall
        """
        return SLIDE_H - 0.65 - h   # 0.65 reserves footer + wordmark

    # ---- content elements ----

    def eyebrow(self, text: str, *, accent: str = "gold",
                underline_w: float = 1.65) -> float:
        """Eyebrow text + thin underline. Returns height consumed.

        Eyebrow is NOT registered for resolve: it's a single line of small
        text whose predictor accuracy is fine, and the trailing rule sits at
        a fixed offset that's stable across renders.
        """
        # Visual block: 10pt text (~0.16") + 0.36 offset + 0.018 rule = ~0.40"
        eyebrow(self.slide, text, tone=self.tone, x=self.x, y=self.y,
                accent=accent, underline_w=underline_w)
        h = 0.40
        self.y += h
        return h

    def title(self, text: str, *, size: float = 36,
              color: Optional[RGBColor] = None,
              line_spacing: float = 1.05,
              font: str = "serif",
              register: bool = True) -> float:
        """Serif headline. Auto-measured height — no dead space.

        Renders with autofit=True (SHAPE_TO_FIT_TEXT) so the shape itself
        shrinks tight to the actual rendered text. If `register=True`
        (default), the shape is added to the slide's resolve plan so the
        actual rendered height drives downstream element positions during
        `build.save(resolve=True)`.
        """
        h = measure_text(text, size=size, w=self.w, font=font,
                         line_spacing=line_spacing)
        shape = serif_title(self.slide, text, x=self.x, y=self.y,
                            w=self.w, h=h, size=size, tone=self.tone,
                            color=color, line_spacing=line_spacing,
                            autofit=True)
        if register and shape is not None:
            elem = _register_resolve_element(
                self.slide, shape, x=self.x, y=self.y, w=self.w,
                h_predicted=h, name=f"flow.title[{text[:40]}]",
                prior_elements=self._elements)
            self._elements.append(elem)
        self.y += h
        return h

    def body(self, text: Union[str, List[str]], *, size: float = 11.5,
             color: Optional[RGBColor] = None,
             line_spacing: float = 1.30,
             font: str = "sans_l",
             register: bool = True) -> float:
        """Body paragraph or bullet list. Auto-measured height.

        If `register=True` (default), the shape is added to the slide's
        resolve plan so the actual rendered height drives downstream
        element positions during `build.save(resolve=True)`.
        """
        h = measure_text(text, size=size, w=self.w, font=font,
                         line_spacing=line_spacing)
        shape = body(self.slide, text, x=self.x, y=self.y, w=self.w, h=h,
                     size=size, color=color, tone=self.tone,
                     line_spacing=line_spacing, autofit=True)
        if register and shape is not None:
            label = text if isinstance(text, str) else " / ".join(text)
            elem = _register_resolve_element(
                self.slide, shape, x=self.x, y=self.y, w=self.w,
                h_predicted=h, name=f"flow.body[{label[:40]}]",
                prior_elements=self._elements)
            self._elements.append(elem)
        self.y += h
        return h

    def subhead(self, text: str, *, size: float = 12,
                color: Optional[RGBColor] = None,
                line_spacing: float = 1.20,
                register: bool = True) -> float:
        """Sub-section label. Auto-measured height."""
        h = measure_text(text, size=size, w=self.w, font="sans_sb",
                         line_spacing=line_spacing)
        shape = subhead(self.slide, text, x=self.x, y=self.y, w=self.w, h=h,
                        size=size, color=color, tone=self.tone, autofit=True)
        if register and shape is not None:
            elem = _register_resolve_element(
                self.slide, shape, x=self.x, y=self.y, w=self.w,
                h_predicted=h, name=f"flow.subhead[{text[:40]}]",
                prior_elements=self._elements)
            self._elements.append(elem)
        self.y += h
        return h

    def rule(self, *, color: Optional[RGBColor] = None,
             width_pt: float = 0.5, width_in: Optional[float] = None) -> float:
        """Horizontal divider rule. Returns 0 — does not advance cursor on
        its own; pair with .gap() if you want space around it."""
        c = color or (COLORS["rule_dark"] if self.tone == "dark"
                      else COLORS["rule_light"])
        w = width_in if width_in is not None else self.w
        _line(self.slide, self.x, self.y, self.x + w, self.y,
              color=c, width_pt=width_pt)
        return 0.0

    def append(self, draw_fn, *, height: float) -> float:
        """Custom element. `draw_fn` receives (slide, x, y, w) and is
        expected to consume `height` inches of vertical space."""
        draw_fn(self.slide, self.x, self.y, self.w)
        self.y += height
        return height

    def attach_dependent(self, shape) -> "Flow":
        """Register `shape` as a dependent of every prior Flow element on
        this slide. When any of those elements resolves to a different
        height than predicted, `shape` shifts by the cumulative delta.

        Use this when you draw a non-Flow shape after Flow content — e.g.
        a table or KPI strip at `f.y` after Flow.title + Flow.body — and
        you want it to track Flow's measurement corrections:

            f.title("…")
            f.body("…")
            f.gap(0.32)
            table_shape = build.table(s, …, x=f.x, y=f.y, w=f.w, h=…)
            f.attach_dependent(table_shape)

        Multiple shapes can be attached individually for a multi-shape
        layout (e.g. a table with several adjacent helpers).
        """
        if shape is None:
            return self
        for prior in self._elements:
            prior.dependents.append(shape)
        return self

    def attach_dependents(self, shapes: Iterable) -> "Flow":
        """Convenience: attach multiple dependent shapes at once."""
        for shp in shapes:
            self.attach_dependent(shp)
        return self

    def attach_all_below(self, slide=None) -> "Flow":
        """One-line retrofit: attach every shape drawn on the slide AFTER
        the last Flow element as a dependent of every prior Flow element.

        Call this at the end of a slide-building function. It scans the
        slide's current shape list, finds the position of the last shape
        Flow registered, and attaches every shape after that position as
        a dependent. Net effect: any non-Flow shape drawn at `f.y` (or
        below) — table, KPI strip, callout, sub-Flow content, hand-laid
        column text — moves with Flow's resolve corrections.

        Equivalent to wrapping the entire post-Flow content in
        `with f.absorb_below():` but without re-indentation. Use this
        when retrofitting existing scripts; use `absorb_below` for new
        code where re-indenting is cheap.

        No-op if Flow has no registered elements yet.
        """
        target = slide if slide is not None else self.slide
        if not self._elements:
            return self
        # python-pptx creates a NEW Shape wrapper each time slide.shapes is
        # iterated, so `shp is last_shape` returns False even when both
        # point at the same underlying XML element. Compare via the
        # underlying `_element` (the <p:sp> / <p:graphicFrame>) instead.
        last_elem = self._elements[-1].shape._element
        shape_list = list(target.shapes)
        idx = None
        for i, shp in enumerate(shape_list):
            if shp._element is last_elem:
                idx = i
                break
        if idx is None:
            return self
        new_shapes = shape_list[idx + 1:]
        if new_shapes:
            self.attach_dependents(new_shapes)
        return self

    def absorb_below(self, slide=None):
        """Context manager that snapshots the slide's shape count on enter
        and attaches every shape added during the `with` block as a
        dependent of every prior Flow element.

        This is the ergonomic version of the manual snapshot pattern:

            with f.absorb_below():
                build.table(s, ..., x=f.x, y=f.y, w=f.w, h=...)
                build.callout(s, ..., x=f.x, y=..., ...)

        Without this, every downstream shape stays at its predictor-driven
        y while Flow.title / Flow.body tighten upward, leaving the visible
        dead space between Flow header and downstream content that the
        resolve report flags as "0 dep" entries.

        `slide` defaults to `self.slide` — pass an explicit slide only if
        you're drawing onto a different slide than the Flow's own.
        """
        target = slide if slide is not None else self.slide

        class _AbsorbBelowCtx:
            def __init__(_self, flow, slide):
                _self.flow = flow
                _self.slide = slide
                _self.n_before = 0

            def __enter__(_self):
                _self.n_before = len(_self.slide.shapes)
                return _self

            def __exit__(_self, exc_type, exc_val, tb):
                if exc_type is None:
                    new_shapes = list(_self.slide.shapes)[_self.n_before:]
                    _self.flow.attach_dependents(new_shapes)
                return False

        return _AbsorbBelowCtx(self, target)


# ---------------------------------------------------------------------------
# CBRE chrome (footer, wordmark) — drawn on every slide
# ---------------------------------------------------------------------------

def _paint_bg_native(slide, tone: str, *, log: bool = True):
    """Set the slide's native background colour. tone='dark' or 'light'.

    Replaces the previous overlay-rectangle approach (`_paint_bg`) — the
    background is now a true slide property, not an extra rectangle in
    the shape tree. Editors can't accidentally select / move / delete a
    "background" shape because there isn't one; clicking empty area
    selects the slide itself.

    For split-tone slides (e.g. `value_prop_intro` with a dark top and
    light bottom), call this with the *top* tone for the slide bg, then
    overlay a single `_rect(...)` for the bottom band — far cleaner than
    painting two overlapping full-bleed rects.

    Also records the tone on `slide._cbre_tone` so `resolve_slide()`
    knows which bg colour to use during pixel-walk corner sampling.

    If log=True (default), the tone is appended to the deck's tone log
    so audit_tones() can verify the dark/light balance later.
    """
    rgb = COLORS["page_dark"] if tone == "dark" else COLORS["page_light"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb
    try:
        slide._cbre_tone = tone   # type: ignore[attr-defined]
    except AttributeError:
        pass
    if log:
        deck = getattr(slide, "_cbre_deck", None)
        if deck is not None:
            _log_tone(deck, tone)


def _paint_bg(slide, tone: str, *, log: bool = True):
    """DEPRECATED — kept as a compatibility shim that forwards to
    `_paint_bg_native`. The old overlay-rectangle approach polluted the
    shape tree with a full-bleed rect that editors could accidentally
    select. See spacing-and-rules.md §5.

    Existing call sites have been migrated. This shim exists for any
    user-side scripts that import `_paint_bg` directly. New code should
    call `_paint_bg_native` (or simply use `blank()` / a recipe which
    will pick it up automatically).
    """
    import warnings
    warnings.warn(
        "_paint_bg is deprecated; use _paint_bg_native (native slide "
        "background, no overlay rectangle). See spacing-and-rules.md §5.",
        DeprecationWarning, stacklevel=2,
    )
    _paint_bg_native(slide, tone, log=log)
    return None


# --- Official logo artwork --------------------------------------------------
# The CBRE wordmark must be the supplied logo ARTWORK, never typed (2026 v17
# rule: "do not type the logo yourself"). Drop the official files in
# scripts/assets/ (see assets/README.md):
#   cbre-logo-white.(emf|png)  -> DARK backgrounds
#   cbre-logo-green.(emf|png)  -> LIGHT backgrounds
# Until present, _paint_footer falls back to a typed wordmark and warns once.
LOGO_DIR = Path(__file__).resolve().parent / "assets"
_LOGO_WARNED = False


def _logo_path(tone: str) -> Optional[Path]:
    """Return the artwork path for this tone (white on dark, colour on light),
    preferring vector EMF over PNG, or None if no artwork is installed."""
    variant = "white" if tone == "dark" else "green"
    for ext in (".emf", ".png"):
        p = LOGO_DIR / f"cbre-logo-{variant}{ext}"
        if p.exists():
            return p
    return None


def _logo(slide, tone: str, *, w: float = 1.05):
    """Place the official logo bottom-right at width `w` (height follows the
    artwork's native aspect). Returns the picture, or None if no artwork is
    installed. Clear space (>= logo height) is held by the empty footer band."""
    p = _logo_path(tone)
    if p is None:
        return None
    right = SLIDE_W - 0.45
    pic = slide.shapes.add_picture(str(p), Inches(right - w), Inches(0.0),
                                   width=Inches(w))
    h_in = pic.height / 914400.0          # EMU -> inches
    pic.top = Inches(SLIDE_H - 0.30 - h_in)
    return pic


def _paint_footer(slide, tone: str, *, page_no: Optional[int] = None):
    """Standard CBRE footer: copyright bottom-left, logo artwork bottom-right."""
    fg = COLORS["charcoal"] if tone == "light" else COLORS["mint_dark"]
    _text(slide, "Confidential & Proprietary | (c) 2026 CBRE, Inc.",
          x=0.55, y=7.15, w=6.0, h=0.22,
          font=FONTS["sans"], size=9, color=fg, anchor="top")

    # CBRE wordmark — official artwork bottom-right (brand: never type the logo).
    if _logo(slide, tone) is None:
        global _LOGO_WARNED
        if not _LOGO_WARNED:
            import sys
            print(
                f"[brand] No CBRE logo artwork in {LOGO_DIR} - falling back to a "
                "TYPED wordmark, which the 2026 brand guidelines forbid. Add "
                "cbre-logo-white.(emf|png) + cbre-logo-green.(emf|png) to make "
                "decks brand-compliant.", file=sys.stderr)
            _LOGO_WARNED = True
        wm_color = COLORS["white"] if tone == "dark" else COLORS["cbre_green"]
        _text(slide, "CBRE", x=12.10, y=6.92, w=1.10, h=0.45,
              font=FONTS["sans_b"], size=22, color=wm_color,
              bold=True, align="right", letter_spacing=-0.05, anchor="middle")


def blank(deck: Presentation, *, tone: str = "dark"):
    """Add an empty branded slide and return it (for custom layouts).
    Paints background, footer, and wordmark — you draw the rest."""
    s = _blank_slide(deck)
    _paint_bg_native(s, tone)
    _paint_footer(s, tone)
    return s


# ---------------------------------------------------------------------------
# Reusable visual atoms (the "look")
# ---------------------------------------------------------------------------

def eyebrow(slide, text: str, *, tone: str = "dark", x: float = 0.55,
            y: float = 0.55, color: Optional[RGBColor] = None,
            accent: str = "gold", underline_w: float = 1.65):
    """Uppercase eyebrow tag with a thin underline rule below it (matches the
    reference deck — see 'CBRE VALUE PROPOSITION', 'DEBT & STRUCTURED FINANCE',
    'REAL ESTATE CAPITAL STRATEGY' eyebrows).

    Accent defaults to cream/wheat ('gold' in the reference deck), which is the
    primary eyebrow colour. Use accent='mint' for secondary eyebrows on the
    same slide (e.g. the right-side 'MOST COMMON FUNDING STRATEGIES AT A GLANCE'
    pair on the why-occupiers slide).
    """
    if color is not None:
        c = color
    elif tone == "dark":
        c = COLORS["gold"] if accent == "gold" else COLORS["mint"]
    else:
        c = COLORS["ink"] if accent == "gold" else COLORS["mint_dark"]
    _text(slide, text, x=x, y=y, w=8.0, h=0.28,
          font=FONTS["sans_sb"], size=10, color=c,
          bold=True, uppercase=True, letter_spacing=1.5, anchor="top")
    # Underline rule directly below the text
    _rect(slide, x, y + 0.36, underline_w, 0.018, fill=c)


def _enforce_financier_headline(text: str, size: float) -> None:
    """Brand rule (2026 v17): Financier Display headlines must be >= 20 pt and
    title case, never all caps. Size is a hard error (there is always a clear
    alternative — Calibre); all-caps is a warning, since acronym-only strings
    (e.g. "CBRE", "ESG") are legitimate."""
    if size < FINANCIER_MIN_PT:
        raise ValueError(
            f"serif_title size={size}pt violates the CBRE brand rule that "
            f"Financier Display must be >= {FINANCIER_MIN_PT}pt. Headlines are "
            f"serif; anything smaller should be Calibre — use subhead()/body() "
            f"for small text, or kpi_block()/kpi_strip() for numerals."
        )
    letters = [ch for ch in text if ch.isalpha()]
    if len(letters) > 4 and text == text.upper():
        import sys
        print(
            f"[brand] serif_title received ALL-CAPS text ('{text[:40]}'): the "
            f"CBRE rule is TITLE CASE for Financier Display, never all caps. "
            f"Use title case (Calibre may be all-caps for hero moments).",
            file=sys.stderr,
        )


def serif_title(slide, text: str, *, x: float, y: float, w: float, h: float,
                size: float = 36, color: Optional[RGBColor] = None,
                tone: str = "dark", line_spacing: float = 1.05,
                autofit: bool = True, **_banned_kwargs):
    """Big serif headline. Default colors honor the tone.

    The title box always autofits its rendered text (box grows to fit).
    Pass a generous `w` and let autofit handle `h`. The `h` argument is a
    hint, not a constraint: a short title produces a short box, a long title
    produces a tall one.

    Brand rule (2026 v17): Financier Display is >= 20 pt and title case (never
    all caps). `size < 20` raises; all-caps warns. See brand-guidelines.md.

    Does NOT accept `shrink`. See SKILL.md.
    """
    _reject_banned_kwargs(_banned_kwargs, "serif_title")
    _enforce_financier_headline(text, size)
    c = color or (COLORS["white"] if tone == "dark" else COLORS["green"])
    return _text(slide, text, x=x, y=y, w=w, h=h,
                 font=FONTS["serif"], size=size, color=c,
                 line_spacing=line_spacing, anchor="top",
                 autofit=autofit)


def apply_real_bullets(shape, *, char="•", bullet_font="Arial",
                       marL=171450, space_before_pt=4):
    """Turn each non-empty paragraph of a text shape into a REAL PowerPoint
    bullet (an ``a:buChar`` with a hanging indent and breathing space).

    ALWAYS use this, or ``body(..., bullets=True)``, for bullet lists. NEVER
    type a bullet glyph (the dot, ``-``, ``*``) into the run text: that is a
    FAKE bullet with no hanging indent, so wrapped lines collapse back under the
    glyph instead of aligning under the text. ``marL`` 171450 EMU = 0.1875in;
    ``indent`` is the equal negative hang. Returns the shape."""
    from pptx.oxml.ns import qn  # python-pptx is always present here
    if shape is None or not getattr(shape, "has_text_frame", False):
        return shape
    for p in shape.text_frame.paragraphs:
        if not "".join(r.text for r in p.runs).strip():
            continue
        try:
            p.space_before = Pt(space_before_pt)
        except Exception:  # noqa
            pass
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(int(marL)))
        pPr.set("indent", str(-int(marL)))
        for tag in ("a:buNone", "a:buAutoNum", "a:buChar", "a:buFont"):
            el = pPr.find(qn(tag))
            if el is not None:
                pPr.remove(el)
        pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": bullet_font}))
        pPr.append(pPr.makeelement(qn("a:buChar"), {"char": char}))
    return shape


# Leading glyphs that mean someone is FAKING a bullet by typing it into the text.
# (The hyphen/asterisk are excluded so legitimate content like "-1%" is allowed.)
_FAKE_BULLET_GLYPHS = ("•", "‣", "▪", "◦", "·", "⁃", "∙")


def body(slide, text: Union[str, List[str]], *, x: float, y: float,
         w: float, h: float, size: float = 10.5,
         color: Optional[RGBColor] = None, tone: str = "dark",
         line_spacing: float = 1.30, autofit: bool = True,
         bullets: bool = False, **_banned_kwargs):
    """Paragraph or bullet body text. Default autofit=True — the box always
    grows to fit the text. If body text overflows a card, increase the card
    height or reduce copy length rather than shrinking the font.

    bullets=True -> render as a REAL PowerPoint bullet list (``a:buChar`` with a
    hanging indent). ALWAYS use this for lists. NEVER type a bullet glyph into
    the text to fake a list; ``body()`` raises ``ValueError`` if it sees one.
    See SKILL.md.

    Returns the python-pptx text shape so callers (Flow, CardFlow,
    resolve_slide) can register it for measurement-driven layout correction.

    Does NOT accept `shrink`. See SKILL.md.
    """
    _reject_banned_kwargs(_banned_kwargs, "body")
    for _t in (text if isinstance(text, list) else [text]):
        if isinstance(_t, str) and _t.lstrip()[:1] in _FAKE_BULLET_GLYPHS:
            raise ValueError(
                "Fake bullet detected: a bullet glyph was typed into body text "
                f"({_t.strip()[:40]!r}). Use body(..., bullets=True) for real "
                "PowerPoint bullets; never type the glyph into the text. See SKILL.md.")
    c = color or (COLORS["white"] if tone == "dark" else COLORS["ink_2"])
    box = _text(slide, text, x=x, y=y, w=w, h=h,
                font=FONTS["sans_l"], size=size, color=c,
                line_spacing=line_spacing, anchor="top",
                autofit=autofit)
    if bullets:
        apply_real_bullets(box)
    return box


def subhead(slide, text: str, *, x: float, y: float, w: float, h: float,
            size: float = 12, color: Optional[RGBColor] = None,
            tone: str = "dark", autofit: bool = True,
            **_banned_kwargs):
    _reject_banned_kwargs(_banned_kwargs, "subhead")
    c = color or (COLORS["mint"] if tone == "dark" else COLORS["green"])
    return _text(slide, text, x=x, y=y, w=w, h=h,
                 font=FONTS["sans_sb"], size=size, color=c, bold=True,
                 anchor="top", line_spacing=1.2, autofit=autofit)


# ---------------------------------------------------------------------------
# Card sizing & bounded layout — defensive guards
# ---------------------------------------------------------------------------
#
# All four guards below address the same failure mode: a card with fixed
# inner offsets and a floating outer height where the offsets nearly equal
# the height, leaving sub-0.05" bottom clearance. The numbers fit, the
# bounds-check passes, but the slide reads as crammed. Guards introduced
# 2026-05-19 after the Stellantis Part 2 slide 7 incident.

class CardOverflowError(AssertionError):
    """Raised by CardFlow / card primitives when content would exceed the
    card's bounded inner area. Distinct subclass so callers can catch
    cramming separately from other build assertions."""


def min_h_for_roman_card(n_bullets: int, *, subtitle: bool = False,
                         body_size: float = 9,
                         body_line_spacing: float = 1.32,
                         bottom_pad: float = 0.18) -> float:
    """Minimum card_h needed for a `roman_card` with the given content.

    Use this to pre-validate card_h before laying out a grid, or compare
    against a back-solved card height. The roman_card primitive itself
    raises CardOverflowError if you build at a smaller h — but pre-checking
    lets you reduce item count *before* building the slide rather than
    catching the assertion after the fact."""
    body_y_offset = 1.10 if subtitle else 0.92
    line_h_in = body_size * body_line_spacing / 72.0
    descender = body_size / 72.0 * 0.06
    body_h_needed = n_bullets * line_h_in + descender
    return body_y_offset + body_h_needed + bottom_pad


def min_h_for_decimal_card(body_text: Union[str, List[str]], *,
                           w: float, body_size: float = 9,
                           body_line_spacing: float = 1.32,
                           bottom_pad: float = 0.18) -> float:
    """Minimum card_h needed for a `decimal_card` rendering the given body
    text at the given inner width.

    decimal_card starts the body at y + 1.18 from the card top. The card
    therefore needs at least 1.18 + measured body height + bottom_pad to
    avoid the rendered text crashing into the card's bottom edge."""
    body_h_needed = measure_text(body_text, size=body_size, w=w - 0.10,
                                 font="sans_l",
                                 line_spacing=body_line_spacing)
    return 1.18 + body_h_needed + bottom_pad


def _assert_card_room(h: float, required_h: float, *,
                      kind: str, n: Optional[int] = None) -> None:
    """Internal: raise CardOverflowError if the caller asked for a card
    height that can't accommodate the requested content with healthy
    bottom padding. Message names the four legitimate fixes."""
    if h + 0.01 < required_h:
        label = f"{kind} card" + (f" #{n}" if n is not None else "")
        raise CardOverflowError(
            f"{label}: h={h:.2f}\" is below the minimum {required_h:.2f}\" "
            f"needed for the content + 0.18\" bottom padding. Cards built "
            f"at this h will visually cram against the bottom edge. Fix "
            f"structurally by: (a) reducing items so each card gets more "
            f"height, (b) trimming the copy, (c) restructuring to a "
            f"multi-column grid, or (d) splitting onto two slides. Do "
            f"NOT silence this by passing a smaller bottom_pad — that's "
            f"the original cramming bug wearing a number."
        )


class CardFlow:
    """Bounded vertical stacking helper for the *interior* of a card.

    Like `Flow`, but the cursor is clamped to the area `[top_pad,
    card_h - bottom_pad]` inside a card rectangle. Every text addition
    raises CardOverflowError if it would push the cursor past
    `card_h - bottom_pad` — eliminating the silent cramming pattern
    where fixed inner y-offsets compress against a thin floating card.

    The card rectangle itself is NOT drawn by CardFlow — draw the rect
    first (so you can pick the fill colour), then construct CardFlow
    against its geometry, then add inner content. Use `with_left_inset`
    to shift the inner column right (e.g. past a numeral block) without
    losing the bottom-pad guarantee.

    Example — the slide 7 ranked-card pattern, but cramming-proof:

        _rect(s, cx, cy, cw, ch, fill=COLORS["off_white"])
        _vbar(s, cx, cy, ch, color=accent, width_in=0.06)
        # Numeral and tag pill drawn at fixed positions...
        c = CardFlow(s, x=cx, y=cy, w=cw, h=ch,
                     left_inset=0.78, right_inset=tag_w + 0.30,
                     top_pad=0.16, bottom_pad=0.16, tone="light")
        c.text(title, size=10.5, font="sans_sb",
               color=COLORS["green"], bold=True)
        c.gap(0.06)
        c.text(body_text, size=9, font="sans_l", color=COLORS["ink_2"])
        c.gap(0.06)
        c.text("Approach: " + approach, size=9, font="sans_sb",
               color=COLORS["mint_dark"], italic=True)
        # If any text() pushes past cy + ch - 0.16, CardOverflowError fires
        # with a structural-fix message.
    """

    def __init__(self, slide, *, x: float, y: float, w: float, h: float,
                 top_pad: float = 0.12, bottom_pad: float = 0.14,
                 left_inset: float = 0.20, right_inset: float = 0.20,
                 tone: str = "dark",
                 bg_hint: Optional[Tuple[int, int, int]] = None):
        self.slide = slide
        self.card_x = x
        self.card_y = y
        self.card_w = w
        self.card_h = h
        self.tone = tone
        self.top_pad = top_pad
        self.bottom_pad = bottom_pad
        self.inner_x = x + left_inset
        self.inner_w = w - left_inset - right_inset
        self.y = y + top_pad
        self.bottom_limit = y + h - bottom_pad
        # Resolve registration. bg_hint helps render_measure: cards sit on
        # the slide bg but have their own fill; if the caller knows the
        # card fill (e.g. COLORS["off_white"] or COLORS["green_2"]),
        # passing it eliminates corner-sampling drift from card stripes.
        self.bg_hint = bg_hint
        self._elements: List[ResolveElement] = []

    def _check(self, advance: float, *, what: str) -> None:
        if self.y + advance > self.bottom_limit + 0.005:
            overshoot = (self.y + advance) - self.bottom_limit
            raise CardOverflowError(
                f"CardFlow.{what}: would push cursor to "
                f"{self.y + advance:.3f}\", past inner limit "
                f"{self.bottom_limit:.3f}\" (overshoot "
                f"{overshoot:.3f}\"). The card is too short for the "
                f"requested content + {self.bottom_pad:.2f}\" bottom "
                f"pad. Reduce items in the grid (raises card_h), trim "
                f"this card's copy, or split onto two slides."
            )

    def gap(self, h: float) -> "CardFlow":
        self._check(h, what="gap")
        self.y += h
        return self

    def text(self, text, *, size: float, font: str = "sans_l",
             color: Optional[RGBColor] = None, bold: bool = False,
             italic: bool = False, align: str = "left",
             line_spacing: float = 1.30,
             uppercase: bool = False,
             extra_right_inset: float = 0.0,
             register: bool = True) -> float:
        """Add a text block at the current cursor and advance by its
        measured height. Raises CardOverflowError on overflow.

        `extra_right_inset` shrinks this single element's width on the
        right edge (e.g. to clear a tag pill in the top-right corner
        without affecting subsequent elements' widths). The cursor still
        advances by the wrapped height computed at the narrower width.

        If `register=True` (default), the shape is added to the slide's
        resolve plan with the card's bounding box recorded. After resolve,
        if the actual rendered height of any inner element would push the
        cumulative cursor past `card_h − bottom_pad`, a CardOverflowError
        is raised — closing the predictor-vs-render gap at primitive
        level.
        """
        elem_w = self.inner_w - extra_right_inset
        h = measure_text(text, size=size, w=elem_w, font=font,
                         line_spacing=line_spacing, uppercase=uppercase)
        self._check(h, what="text")
        font_name = FONTS.get(font, font)
        shape = _text(self.slide, text, x=self.inner_x, y=self.y,
                      w=elem_w, h=h,
                      font=font_name, size=size, color=color, bold=bold,
                      italic=italic, align=align, line_spacing=line_spacing,
                      uppercase=uppercase, autofit=True)
        if register and shape is not None:
            label = text if isinstance(text, str) else " / ".join(text)
            card_box = (self.card_x, self.card_y, self.card_w, self.card_h,
                        self.bottom_pad)
            elem = _register_resolve_element(
                self.slide, shape, x=self.inner_x, y=self.y, w=elem_w,
                h_predicted=h, name=f"cardflow.text[{label[:30]}]",
                prior_elements=self._elements,
                bg_hint=self.bg_hint, card_box=card_box,
                card_id=id(self))
            self._elements.append(elem)
        self.y += h
        return h

    def headroom(self) -> float:
        """Inches of vertical space remaining below the cursor."""
        return self.bottom_limit - self.y


# ---------------------------------------------------------------------------
# Numbered framework atoms
# ---------------------------------------------------------------------------

_ROMAN = ["I", "II", "III", "IV", "V", "VI", "VII", "VIII", "IX", "X",
          "XI", "XII"]


def roman_card(slide, n: int, x: float, y: float, w: float, h: float, *,
               title: str, body_lines: List[str],
               accent: Union[str, RGBColor] = "mint", tone: str = "dark",
               subtitle: Optional[str] = None,
               accent_fill: Optional[RGBColor] = None):
    """Card with Roman numeral, title, optional subtitle, and bullet body.
    Used in 'framework_roman' patterns. Matches the reference deck — coloured
    top stripe, large serif numeral in the stripe colour, white bold title,
    mint dash bullets.

    accent_fill (optional): replaces the default card fill with a custom
        RGBColor. Use this to highlight one card in a grid by swapping its
        surface colour rather than stacking extra accents (a top stripe + a
        left bar + a tag pill, etc.). Example:

            roman_card(s, 3, ..., accent_fill=build.COLORS["green_3"])

        When set, the body / title text colours stay tone-correct; consider
        passing an explicit `accent` (the stripe and numeral colour) that
        contrasts well with the new fill.

    Raises CardOverflowError if `h` is too small for the bullet count plus
    healthy bottom padding. To pre-size a grid, call
    `min_h_for_roman_card(len(body_lines), subtitle=bool(subtitle))`.
    """
    # Cramming guard — fires before any shape is drawn.
    required_h = min_h_for_roman_card(
        len(body_lines), subtitle=bool(subtitle))
    _assert_card_room(h, required_h, kind="roman", n=n)

    # Resolve accent
    if isinstance(accent, RGBColor):
        stripe_color = accent
        numeral_color = accent
    else:
        stripe_color = {
            "gold": COLORS["gold"],
            "mint": COLORS["mint"],
            "blue": COLORS["blue"],
        }.get(accent, COLORS["mint"])
        numeral_color = stripe_color
    default_fill = COLORS["green_2"] if tone == "dark" else COLORS["off_white"]
    fill = accent_fill if accent_fill is not None else default_fill
    body_color = COLORS["white"] if tone == "dark" else COLORS["ink_2"]
    title_color = COLORS["white"] if tone == "dark" else COLORS["green"]
    bullet_color = COLORS["mint"] if tone == "dark" else COLORS["mint_dark"]

    _rect(slide, x, y, w, h, fill=fill)
    # Top accent stripe — matches reference deck card framework
    _rect(slide, x, y, w, 0.045, fill=stripe_color)
    # Large serif Roman numeral (kept compact)
    _text(slide, _ROMAN[n - 1], x=x + 0.26, y=y + 0.16, w=1.0, h=0.62,
          font=FONTS["serif_l"], size=32, color=numeral_color,
          line_spacing=1.0, anchor="top")
    # Title — sits to the right of the numeral. Pulled left from x+1.05
    # to x+0.80 to close the dead gap flagged in the design audit
    # between Roman numerals (~0.4" wide) and the following title.
    title_x = x + 0.80
    title_w = w - 0.80 - 0.22
    _text(slide, title, x=title_x, y=y + 0.26, w=title_w, h=0.55,
          font=FONTS["sans_sb"], size=12, color=title_color,
          bold=True, line_spacing=1.16, anchor="top")
    if subtitle:
        _text(slide, subtitle, x=title_x, y=y + 0.78, w=title_w, h=0.24,
              font=FONTS["sans"], size=9, color=numeral_color,
              italic=True, anchor="top")
    # Body bullets — dash markers in mint, body text white/ink
    body_y = y + (1.10 if subtitle else 0.92)
    body_lines_fmt = ["–   " + ln for ln in body_lines]
    body(slide, body_lines_fmt, x=x + 0.26, y=body_y,
         w=w - 0.45, h=h - (body_y - y) - 0.18,
         size=9, color=body_color, tone=tone, line_spacing=1.32)


def decimal_card(slide, n: int, x: float, y: float, w: float, h: float, *,
                 title: str, body_text: str,
                 accent_color: Optional[RGBColor] = None,
                 tone: str = "dark",
                 accent_fill: Optional[RGBColor] = None):
    """Card with decimal '01' style numbering at top, separator line under it,
    then title + body. Used in the 4-card value-prop strip (slide 4).

    Raises CardOverflowError if `h` is too small for the rendered body text
    plus healthy bottom padding. Pre-size a grid with
    `min_h_for_decimal_card(body_text, w=card_w)`.

    accent_fill (optional): paint the card's full background with a custom
        RGBColor (e.g. to highlight the Copilot card in a 4-card row). When
        absent the card is transparent on its parent surface. When set, the
        fill is laid down first so the numbering / title / body render on top.
    """
    # Cramming guard — fires before any shape is drawn.
    required_h = min_h_for_decimal_card(body_text, w=w)
    _assert_card_room(h, required_h, kind="decimal", n=n)

    accent_color = accent_color or COLORS["mint"]
    title_color = COLORS["white"] if tone == "dark" else COLORS["green"]
    body_color = COLORS["white"] if tone == "dark" else COLORS["ink_2"]

    # Optional accent fill — paint the card surface before any other elements
    if accent_fill is not None:
        _rect(slide, x, y, w, h, fill=accent_fill)
    # Top mint horizontal accent
    _hbar(slide, x, y, w - 0.12, color=accent_color, height_in=0.030)
    # Decimal number
    _text(slide, f"{n:02d}", x=x, y=y + 0.08, w=1.0, h=0.45,
          font=FONTS["sans_sb"], size=18, color=accent_color,
          bold=True, line_spacing=1.0)
    # Title
    _text(slide, title, x=x, y=y + 0.52, w=w - 0.10, h=0.60,
          font=FONTS["sans_sb"], size=13, color=title_color,
          bold=True, line_spacing=1.18)
    # Body
    body(slide, body_text, x=x, y=y + 1.18, w=w - 0.10, h=h - 1.28,
         size=9, color=body_color, tone=tone, line_spacing=1.32)


# ---------------------------------------------------------------------------
# KPI / stat blocks
# ---------------------------------------------------------------------------

def kpi_block(slide, x: float, y: float, w: float, h: float, *,
              value: str, label: str, tone: str = "dark",
              value_color: Optional[RGBColor] = None,
              size: float = 36):
    """Single KPI: big serif number on top, label underneath.
    Defaults: gold/cream on dark, deep green on light (matches reference)."""
    vc = value_color or (COLORS["gold"] if tone == "dark" else COLORS["green"])
    lc = COLORS["white"] if tone == "dark" else COLORS["ink_2"]
    _text(slide, value, x=x, y=y, w=w, h=h * 0.55,
          font=FONTS["serif_l"], size=size, color=vc,
          line_spacing=1.0, anchor="top")
    _text(slide, label, x=x, y=y + h * 0.55, w=w, h=h * 0.45,
          font=FONTS["sans"], size=9, color=lc,
          line_spacing=1.30, anchor="top")


def kpi_strip(slide, items: Sequence[Tuple[str, str]], *,
              x: float, y: float, w: float, h: float = 1.20,
              tone: str = "dark", title: Optional[str] = None,
              value_size: float = 32):
    """Horizontal strip of 2-5 KPIs separated by thin mint dividers.

    items: list of (value, label) tuples.
    Optional title appears above the strip."""
    if title:
        subhead(slide, title, x=x, y=y - 0.40, w=w, h=0.30,
                size=10, tone=tone)

    n = len(items)
    col_w = w / n
    rule_color = COLORS["mint_dark"] if tone == "dark" else COLORS["rule_light"]

    for i, (val, lbl) in enumerate(items):
        cx = x + i * col_w
        if i > 0:
            _line(slide, cx, y + 0.05, cx, y + h - 0.05,
                  color=rule_color, width_pt=0.75)
        kpi_block(slide, cx + 0.18, y, col_w - 0.22, h,
                  value=val, label=lbl, tone=tone, size=value_size)


# ---------------------------------------------------------------------------
# Tables
# ---------------------------------------------------------------------------

def table(slide, headers: Sequence[str], rows: Sequence[Sequence[str]], *,
          x: float, y: float, w: float, h: float, tone: str = "dark",
          header_fill: Optional[RGBColor] = None,
          header_text: Optional[RGBColor] = None,
          first_col_emphasis: bool = True,
          footer_row: Optional[Sequence[str]] = None,
          footer_fill: Optional[RGBColor] = None,
          col_aligns: Optional[Sequence[str]] = None,
          font_size: float = 9.5):
    """Mint-headed comparison table. Reference: slide 2, slide 3, slide 6.

    - First row = mint header bar with white-ish text.
    - Subsequent rows = striped via thin rule lines.
    - Optional footer_row = colored summary band (e.g. 'Margin Capture' row).
    - first_col_emphasis = bold the first column (row labels).
    """
    n_cols = len(headers)
    n_rows = len(rows) + (1 if footer_row else 0)
    header_h = 0.36
    footer_h = 0.55 if footer_row else 0
    body_h = h - header_h - footer_h
    row_h = body_h / max(1, len(rows))

    col_aligns = col_aligns or (["left"] + ["left"] * (n_cols - 1))
    col_w = w / n_cols
    header_fill = header_fill or COLORS["mint"]
    header_text = header_text or COLORS["green"]
    body_text = COLORS["white"] if tone == "dark" else COLORS["ink_2"]
    rule_color = COLORS["rule_dark"] if tone == "dark" else COLORS["rule_light"]

    # Header row
    _rect(slide, x, y, w, header_h, fill=header_fill)
    for i, hdr in enumerate(headers):
        _text(slide, hdr, x=x + i * col_w + 0.12, y=y,
              w=col_w - 0.24, h=header_h,
              font=FONTS["sans_sb"], size=font_size, color=header_text,
              bold=True, uppercase=True, letter_spacing=0.8,
              align=col_aligns[i] if i > 0 else "left",
              anchor="middle")

    # Body rows
    for r, row in enumerate(rows):
        ry = y + header_h + r * row_h
        # Subtle row rule
        _line(slide, x, ry, x + w, ry,
              color=rule_color, width_pt=0.5)
        for c, cell in enumerate(row):
            is_first = (c == 0)
            font = FONTS["sans_sb"] if (is_first and first_col_emphasis) else FONTS["sans_l"]
            color = (COLORS["mint"] if (is_first and tone == "dark"
                     and first_col_emphasis) else body_text)
            _text(slide, cell, x=x + c * col_w + 0.12, y=ry,
                  w=col_w - 0.24, h=row_h,
                  font=font, size=font_size, color=color,
                  bold=is_first and first_col_emphasis,
                  align=col_aligns[c] if c > 0 else "left",
                  anchor="middle", line_spacing=1.20)

    # Footer row (colored summary band)
    if footer_row:
        fy = y + header_h + body_h
        footer_fill = footer_fill or COLORS["green"]
        _rect(slide, x, fy, w, footer_h, fill=footer_fill)
        for c, cell in enumerate(footer_row):
            _text(slide, cell, x=x + c * col_w + 0.12, y=fy,
                  w=col_w - 0.24, h=footer_h,
                  font=FONTS["sans_sb"], size=font_size + 1.0,
                  color=COLORS["white"], bold=True,
                  align="left" if c == 0 else "center",
                  anchor="middle")

    # Bottom rule
    _line(slide, x, y + header_h + body_h, x + w, y + header_h + body_h,
          color=rule_color, width_pt=0.5)


# ---------------------------------------------------------------------------
# Callout boxes (e.g. "LESSON FOR NIO", "CBRE VIEW")
# ---------------------------------------------------------------------------

# Internal layout constants for callout — kept in sync between
# `predict_callout_h` (caller-side sizing) and `callout` (drawing).
#
# title_region: y-offset to body's top (title sits at +0.20, title h
#   ~0.32, spacer ~0.10 = total 0.62 to where body starts).
# bottom_pad:  inches between body's bottom edge and the bg rect's
#   bottom edge after resolve trims. 0.20" gives comfortable visual
#   breathing room (skill rule §1: ≥0.18 preferred).
# upfront_margin: predict_callout_h returns
#   `title_region + body × upfront_margin + bottom_pad` so the box
#   has 10% headroom over the raw predictor — absorbs predictor drift
#   so callers rarely need to widen it manually.
# probe_margin: resolve probe extends `body_h_pred × probe_margin`
#   past the declared body region. Larger than upfront_margin because
#   the probe should catch BIG under-predictions (wrap mismatch), not
#   just 10% drift. Clamped to safe-bottom anyway, so the multiplier
#   only matters for small callouts on tall slides.
_CALLOUT_TITLE_REGION_H = 0.62
_CALLOUT_BOTTOM_PAD = 0.20
_CALLOUT_UPFRONT_MARGIN = 1.10
_CALLOUT_PROBE_MARGIN = 1.50


def predict_callout_h(body_text: str, *, w: float,
                      title_region_h: float = _CALLOUT_TITLE_REGION_H,
                      bottom_pad: float = _CALLOUT_BOTTOM_PAD,
                      safety_margin: float = _CALLOUT_UPFRONT_MARGIN) -> float:
    """Predict the minimum `h` that a `callout(...)` needs to contain its
    body text without overflowing the bg rect.

    Use this BEFORE drawing the callout so the caller can coordinate the
    callout's height with the rest of the slide layout (e.g. shrink the
    table above to make room):

        cal_h = max(1.05, build.predict_callout_h(body_text, w=PAGE_W))
        cal_y = SAFE_BOT - cal_h
        content_bot = cal_y - 0.20       # table now sizes against this
        ...
        build.callout(s, body_text=body_text, x=PAGE_X, y=cal_y,
                      w=PAGE_W, h=cal_h, ...)

    Why not auto-grow inside `callout()`? Because growing the callout
    upward changes its top y, which collides with whatever was drawn
    above (typically a table sized against `cal_y - margin`). Growing
    downward crosses the slide's footer / wordmark band. Only the
    caller knows which trade-off is acceptable; the predict helper
    surfaces the constraint so the caller can decide.

    Parameters
    ----------
    body_text : str
        The callout's body copy.
    w : float
        Total callout width (same as you'll pass to `callout(w=...)`).
    title_region_h : float
        Vertical room reserved for the title row (default 0.62" matches
        `callout()`'s internal offset).
    bottom_pad : float
        Padding below the body text (default 0.20", from
        `_CALLOUT_BOTTOM_PAD`).
    safety_margin : float
        Multiplier on the predicted body height to absorb predictor
        drift (default 1.10 = 10% headroom). The resolve pass will trim
        back any unused space and the linked bg rect follows.

    Returns
    -------
    float — minimum callout height in inches.
    """
    body_w = w - 0.55       # matches callout's internal body width
    pred_body_h = measure_text(body_text, size=10, w=body_w,
                               font="sans_l", line_spacing=1.38)
    h = title_region_h + pred_body_h * safety_margin + bottom_pad
    # Hard floor: callout()'s inner body box is h - title_region_h - bottom_pad.
    # Keep it >= 0.18" (one line at 10pt) so a very short / empty body can never
    # back-solve a non-positive body box — which PowerPoint rejects as a corrupt
    # file (see spacing-and-rules.md section 1). This makes the "always returns a
    # safe minimum" contract true regardless of the body text.
    return max(h, title_region_h + bottom_pad + 0.18)


def callout(slide, *, title: str, body_text: str,
            x: float, y: float, w: float, h: float, tone: str = "dark",
            tag: Optional[str] = None,
            accent: Optional[RGBColor] = None):
    """The CBRE 'expert note' callout: lifted card with a left vertical mint
    bar, gold/cream uppercase title, white body, optional ribbon tag on the
    right (e.g. 'LESSON FOR NIO | BG / RO / PL / HU SHORTLIST').
    """
    title_accent = accent or COLORS["gold"]
    bar_color = COLORS["mint"] if tone == "dark" else COLORS["mint_dark"]
    fill = COLORS["green_4"] if tone == "dark" else COLORS["off_white"]
    body_color = COLORS["white"] if tone == "dark" else COLORS["ink_2"]

    # If the caller declared an `h` that the predictor says cannot hold
    # the body, print a one-line warning. Callers should use the
    # `predict_callout_h(body_text, w=...)` helper to compute h up front
    # and coordinate it with the rest of the slide's layout (e.g. the
    # table above). We DO NOT silently grow `h` here because that would
    # shift the callout into the table area above.
    body_w_for_predict = w - 0.55
    pred_body_h = measure_text(body_text, size=10, w=body_w_for_predict,
                               font="sans_l", line_spacing=1.38)
    min_h_for_body = _CALLOUT_TITLE_REGION_H + pred_body_h + _CALLOUT_BOTTOM_PAD
    if min_h_for_body > h + 0.02:
        import sys
        print(
            f"[callout] WARNING: body text needs ~{min_h_for_body:.2f}\" "
            f"of height but caller declared h={h:.2f}\". The body will "
            f"overflow the bg rect when opened in PowerPoint. Either "
            f"size the callout with `build.predict_callout_h(body_text, "
            f"w={w:.2f})` before drawing, trim the body text, or widen "
            f"the callout. (Resolve will resize the body shape and grow "
            f"the bg rect via linked_height_shapes, but the result may "
            f"cross the wordmark band at the slide bottom.)",
            file=sys.stderr,
        )

    # Card body
    bg_rect = _rect(slide, x, y, w, h, fill=fill)
    # Left vertical accent bar
    bar_rect = _rect(slide, x, y, 0.06, h, fill=bar_color)
    # Reserve right space for the tag if present
    tag_w = 0.0
    if tag:
        tag_w = max(2.0, min(3.6, len(tag) * 0.085 + 0.55))
    title_w = w - 0.50 - (tag_w + 0.20 if tag else 0)
    # Gold/cream uppercase title
    _text(slide, title, x=x + 0.30, y=y + 0.20, w=title_w, h=0.32,
          font=FONTS["sans_sb"], size=11,
          color=title_accent if tone == "dark" else COLORS["ink"],
          bold=True, uppercase=True, letter_spacing=2.0, anchor="top")
    # Body — register with the slide's resolve plan so resolve resizes
    # the shape (and the linked bg rect + bar) to the actual rendered
    # text height.
    body_x = x + 0.30
    body_y = y + _CALLOUT_TITLE_REGION_H
    body_w = body_w_for_predict
    body_h_pred = h - _CALLOUT_TITLE_REGION_H - _CALLOUT_BOTTOM_PAD
    body_shape = body(slide, body_text, x=body_x, y=body_y, w=body_w,
                      h=body_h_pred, size=10, color=body_color,
                      tone=tone, line_spacing=1.38)
    if body_shape is not None:
        # bg_hint: callout fill is green_4 (dark) or off_white (light).
        # Help render_measure pick the right local-bg colour.
        bg_rgb = (fill[0], fill[1], fill[2])
        # max_probe_h: with up-front content sizing, the body should
        # fit comfortably. But probe a bit past the declared box so
        # residual under-prediction is still caught — capped at the
        # safe-bottom to avoid the footer / wordmark band.
        max_probe = min(body_h_pred * _CALLOUT_PROBE_MARGIN,
                        SLIDE_H - 0.65 - body_y - 0.02)
        # Link the bg rect and accent bar to the body's height so when
        # resolve resizes the body, the parent container resizes too.
        _register_resolve_element(
            slide, body_shape, x=body_x, y=body_y, w=body_w,
            h_predicted=body_h_pred,
            name=f"callout.body[{body_text[:30].replace(chr(10), ' ')}]",
            bg_hint=bg_rgb,
            max_probe_h=max_probe,
            linked_height_shapes=[bg_rect, bar_rect],
        )

    if tag:
        # Tag pill on the right — title-line aligned
        tag_x = x + w - tag_w - 0.25
        # Outline tag (no fill) to match reference deck style
        _rect(slide, tag_x, y + 0.20, tag_w, 0.30,
              fill=fill, line=title_accent if tone == "dark" else COLORS["ink"],
              line_w=0.75)
        _text(slide, tag, x=tag_x + 0.08, y=y + 0.20, w=tag_w - 0.16, h=0.30,
              font=FONTS["sans_sb"], size=9,
              color=title_accent if tone == "dark" else COLORS["ink"],
              bold=True, uppercase=True, letter_spacing=1.4, align="center",
              anchor="middle")


# ===========================================================================
# HIGH-LEVEL SLIDE PATTERNS
# ===========================================================================

# ---------- 1. Cover ------------------------------------------------------

def cover(deck: Presentation, *, title: str, subtitle: Optional[str] = None,
          presenter: Optional[str] = None, org: Optional[str] = None,
          date: Optional[str] = None, tone: str = "dark",
          eyebrow_text: Optional[str] = None,
          themes: Optional[Sequence[str]] = None):
    """Cover slide. Dark by default with a giant serif title.

    `themes`, `subtitle`, `presenter`, `org`, `date`, `eyebrow_text` are all
    optional. A bare cover with a single giant headline is a finished slide.
    Add the themes band only when previewing the deck genuinely helps the
    reader; whitespace on a cover is intentional, not a void to be filled.
    """
    s = _blank_slide(deck)
    _paint_bg_native(s, tone)

    # Top row — eyebrow left, date right
    if eyebrow_text:
        eyebrow(s, eyebrow_text, tone=tone, x=0.55, y=0.55, accent="gold")
    if date:
        _text(s, date, x=SLIDE_W - 3.20, y=0.55, w=2.65, h=0.30,
              font=FONTS["mono"], size=10,
              color=COLORS["gold"] if tone == "dark" else COLORS["charcoal"],
              uppercase=True, letter_spacing=2.0, align="right", anchor="middle")

    # Title — auto-size: short titles get bigger
    title_size = 110 if len(title) <= 22 else 88 if len(title) <= 36 else 70 if len(title) <= 54 else 56
    serif_title(s, title, x=0.55, y=1.55, w=SLIDE_W - 1.10, h=2.10,
                size=title_size, tone=tone, line_spacing=1.02)

    # Subtitle sits directly under the title (no big gap)
    if subtitle:
        body(s, subtitle, x=0.55, y=3.95, w=SLIDE_W - 4.0, h=0.80,
             size=18,
             color=COLORS["mint"] if tone == "dark" else COLORS["mint_dark"],
             tone=tone, line_spacing=1.30)

    # Themes preview band — fills the middle-bottom zone
    if themes:
        band_y = 5.00
        # Cream/gold rule + label (matches the reference deck eyebrow style)
        _rect(s, 0.55, band_y, 1.40, 0.018,
              fill=COLORS["gold"] if tone == "dark" else COLORS["ink"])
        _text(s, "WHAT'S INSIDE", x=0.55, y=band_y + 0.10, w=3.0, h=0.30,
              font=FONTS["sans_sb"], size=10,
              color=COLORS["gold"] if tone == "dark" else COLORS["ink"],
              uppercase=True, letter_spacing=2.4, anchor="top")
        n = min(len(themes), 4)
        col_w = (SLIDE_W - 1.10) / n
        for i, theme in enumerate(themes[:n]):
            cx = 0.55 + i * col_w
            _text(s, f"{i+1:02d}", x=cx, y=band_y + 0.50, w=0.60, h=0.32,
                  font=FONTS["serif"], size=20,
                  color=COLORS["mint"] if tone == "dark" else COLORS["mint_dark"],
                  anchor="top")
            _text(s, theme, x=cx + 0.55, y=band_y + 0.55, w=col_w - 0.65, h=0.80,
                  font=FONTS["sans_l"], size=14,
                  color=COLORS["white"] if tone == "dark" else COLORS["ink"],
                  anchor="top", line_spacing=1.20)

    # Presenter / org as a single bottom-left band
    if presenter or org:
        lines = [v for v in (presenter, org) if v]
        # Cream divider above the presenter block
        _rect(s, 0.55, 6.55, 2.50, 0.018,
              fill=COLORS["gold"] if tone == "dark" else COLORS["ink"])
        body(s, lines, x=0.55, y=6.68, w=6.0, h=0.55,
             size=12,
             color=COLORS["white"] if tone == "dark" else COLORS["ink_2"],
             tone=tone, line_spacing=1.30)

    _paint_footer(s, tone)
    return s


# ---------- 2. Contents (Table of Contents) -------------------------------

def contents(deck: Presentation, *, items: Sequence[str],
             tone: str = "dark", eyebrow_text: Optional[str] = None,
             title: str = "Contents"):
    """Numbered table of contents, modelled after the reference deck slide 1.
    Giant serif word ('Contents') on the left, bright green vertical accent
    line, numbered list on the right."""
    s = _blank_slide(deck)
    _paint_bg_native(s, tone)

    if eyebrow_text:
        eyebrow(s, eyebrow_text, tone=tone)

    # Giant serif word on the left — wide enough to NOT wrap "Contents"
    # at 96pt. For longer words (e.g., "Table of Contents") we pick a
    # smaller declared size below — NOT TEXT_TO_FIT_SHAPE. The box still
    # uses SHAPE_TO_FIT_TEXT (autofit), so the picked size renders as-is.
    word_size = 96 if len(title) <= 8 else 76 if len(title) <= 14 else 60
    serif_title(s, title, x=0.55, y=1.55, w=5.20, h=5.0,
                size=word_size, tone=tone, line_spacing=1.0)

    # Cream/gold vertical accent line
    _vbar(s, x=5.85, y=1.35, h=5.30,
          color=COLORS["gold"] if tone == "dark" else COLORS["mint_dark"],
          width_in=0.022)

    # Numbered list
    n = len(items)
    list_top = 1.60
    list_h = 4.80
    row_h = list_h / max(1, n)
    text_color = COLORS["white"] if tone == "dark" else COLORS["ink_2"]
    num_color = COLORS["gold"] if tone == "dark" else COLORS["mint_dark"]
    size = 20 if n <= 6 else 17 if n <= 8 else 14

    for i, item in enumerate(items):
        ry = list_top + i * row_h
        _text(s, f"{i+1:>2}.", x=6.15, y=ry, w=0.70, h=row_h,
              font=FONTS["serif"], size=size + 2, color=num_color,
              anchor="middle")
        _text(s, item, x=6.95, y=ry, w=SLIDE_W - 7.50, h=row_h,
              font=FONTS["sans_l"], size=size, color=text_color,
              anchor="middle", line_spacing=1.15)

    _paint_footer(s, tone)
    return s


# ---------- 3. Section divider --------------------------------------------

def section_divider(deck: Presentation, *, number: Union[int, str],
                    title: str, eyebrow_text: Optional[str] = None,
                    tone: str = "dark",
                    lead: Optional[str] = None,
                    items: Optional[Sequence[str]] = None):
    """Section opener: giant outline number, serif title.

    `lead` and `items` are optional. The giant number + headline is enough
    on its own when the section is named clearly; pad with `lead` or `items`
    only when the divider doubles as an overview. Whitespace beside the
    numeral should be intentional, not residual.
    """
    s = _blank_slide(deck)
    _paint_bg_native(s, tone)
    if eyebrow_text:
        eyebrow(s, eyebrow_text, tone=tone)

    num_str = f"{int(number):02d}" if isinstance(number, int) or str(number).isdigit() else str(number)
    # Giant serif number — mint on dark, deep green on light
    _text(s, num_str, x=0.55, y=1.20, w=5.5, h=5.0,
          font=FONTS["serif_l"], size=300,
          color=COLORS["mint"] if tone == "dark" else COLORS["green"],
          line_spacing=1.0, anchor="top")

    # Cream/gold vertical accent
    _vbar(s, x=6.30, y=1.40, h=5.20,
          color=COLORS["gold"] if tone == "dark" else COLORS["mint_dark"],
          width_in=0.022)

    # Title
    serif_title(s, title, x=6.70, y=1.30, w=6.20, h=1.80,
                size=46, tone=tone, line_spacing=1.08)

    # Lead description
    if lead:
        body(s, lead, x=6.70, y=3.25, w=6.20, h=1.30,
             size=14, tone=tone, line_spacing=1.45)

    # Section contents list
    if items:
        list_top = 4.60 if lead else 3.40
        rule_color = COLORS["gold"] if tone == "dark" else COLORS["ink"]
        _rect(s, 6.70, list_top, 1.20, 0.018, fill=rule_color)
        _text(s, "IN THIS SECTION", x=6.70, y=list_top + 0.10, w=4.0, h=0.30,
              font=FONTS["sans_sb"], size=10, color=rule_color,
              uppercase=True, letter_spacing=2.2, anchor="top")
        text_color = COLORS["white"] if tone == "dark" else COLORS["ink"]
        num_color = COLORS["mint"] if tone == "dark" else COLORS["mint_dark"]
        # Two-column grid for items
        n = len(items)
        rows = (n + 1) // 2
        row_h = min(0.42, (6.40 - list_top - 0.55) / max(1, rows))
        for i, item in enumerate(items):
            col = i % 2
            row = i // 2
            ix = 6.70 + col * 3.20
            iy = list_top + 0.55 + row * row_h
            _text(s, f"{i+1:02d}", x=ix, y=iy, w=0.50, h=row_h,
                  font=FONTS["mono"], size=10, color=num_color, anchor="middle")
            _text(s, item, x=ix + 0.50, y=iy, w=2.65, h=row_h,
                  font=FONTS["sans_l"], size=12, color=text_color,
                  anchor="middle", line_spacing=1.18)

    _paint_footer(s, tone)
    return s


# ---------- 4. Thank you / closing ----------------------------------------

def thank_you(deck: Presentation, *, title: str = "Thank you.",
              subtitle: Optional[str] = None,
              contacts: Optional[Sequence[dict]] = None,
              tone: str = "dark"):
    """Closing slide. Optional contact cards along the bottom.
    contacts: list of dicts with keys name, title, email, phone."""
    s = _blank_slide(deck)
    _paint_bg_native(s, tone)

    serif_title(s, title, x=0.55, y=1.60, w=SLIDE_W - 1.10, h=2.40,
                size=110, tone=tone, line_spacing=1.0)
    if subtitle:
        body(s, subtitle, x=0.55, y=3.90, w=10.0, h=0.80,
             size=16, tone=tone)

    if contacts:
        n = len(contacts)
        avail_w = SLIDE_W - 1.10
        col_w = (avail_w - 0.30 * (n - 1)) / n
        cy = 5.20
        for i, ct in enumerate(contacts):
            cx = 0.55 + i * (col_w + 0.30)
            # Top cream/gold divider
            _rect(s, cx, cy, col_w - 0.30, 0.020,
                  fill=COLORS["gold"] if tone == "dark" else COLORS["ink"])
            lines = [
                ct.get("name", ""),
                ct.get("title", ""),
                ct.get("email", ""),
                ct.get("phone", ""),
            ]
            lines = [ln for ln in lines if ln]
            # Name bold, rest light
            _text(s, lines[0], x=cx, y=cy + 0.18, w=col_w - 0.20, h=0.40,
                  font=FONTS["sans_sb"], size=13,
                  color=COLORS["white"] if tone == "dark" else COLORS["green"],
                  bold=True)
            body(s, lines[1:], x=cx, y=cy + 0.65, w=col_w - 0.20, h=1.40,
                 size=10,
                 color=COLORS["white"] if tone == "dark" else COLORS["ink_2"],
                 tone=tone, line_spacing=1.35)

    _paint_footer(s, tone)
    return s


# ---------- 5. Case study (the canonical reference slide #2) --------------

def case_study(deck: Presentation, *, eyebrow_text: str, title: str,
               intro: str, framework_title: str,
               framework: Sequence[Tuple[str, str, str]],
               table_headers: Sequence[str],
               table_rows: Sequence[Sequence[str]],
               stat_strip_title: str,
               stats: Sequence[Tuple[str, str]],
               callout_title: Optional[str] = None,
               callout_body: Optional[str] = None,
               callout_tag: Optional[str] = None, tone: str = "dark"):
    """Dense editorial slide: intro + numbered framework on the left, a fact
    table + stat strip on the right. Modelled on the CATL reference slide.

    framework: list of (number_str, topic, description) for each row.

    callout_title / callout_body (both optional): when supplied, a "CBRE VIEW"
        / "LESSON FOR X" callout strap is added to the bottom-right under the
        stat strip. When omitted, no callout is drawn and the stat strip's
        labels are allowed to breathe into the freed space.

        Closing straps are opt-in by design: when every slide in a deck has a
        bottom-strap callout, the deck reads as templated. Only attach the
        callout when the slide genuinely has a single takeaway to underline.
    """
    s = _blank_slide(deck)
    _paint_bg_native(s, tone)
    eyebrow(s, eyebrow_text, tone=tone, accent="gold")

    # === Left column =====================================================
    serif_title(s, title, x=0.55, y=1.05, w=6.4, h=1.85,
                size=30, tone=tone, line_spacing=1.05)
    # Mint vertical bar next to intro paragraph (matches reference deck)
    intro_y = 3.00
    _vbar(s, x=0.55, y=intro_y, h=1.55,
          color=COLORS["mint"] if tone == "dark" else COLORS["mint_dark"],
          width_in=0.022)
    body(s, intro, x=0.75, y=intro_y - 0.02, w=6.15, h=1.60,
         size=10, tone=tone, line_spacing=1.42)

    # Framework section label — gold/cream, matches reference
    label_color = COLORS["gold"] if tone == "dark" else COLORS["ink"]
    _text(s, framework_title, x=0.55, y=4.85, w=6.4, h=0.32,
          font=FONTS["sans_sb"], size=12, color=label_color,
          bold=True, anchor="top")
    _rect(s, 0.55, 5.22, 1.20, 0.014, fill=label_color)

    # Framework rows (number | topic | description)
    fw_top = 5.45
    fw_bottom = 7.00
    fw_h = fw_bottom - fw_top
    row_h = fw_h / max(1, len(framework))
    for i, (num, topic, desc) in enumerate(framework):
        ry = fw_top + i * row_h
        # Mint serif number
        _text(s, num, x=0.55, y=ry, w=0.65, h=row_h,
              font=FONTS["serif"], size=16,
              color=COLORS["mint"] if tone == "dark" else COLORS["mint_dark"],
              anchor="top")
        # Topic — white bold
        _text(s, topic, x=1.20, y=ry + 0.02, w=1.85, h=row_h,
              font=FONTS["sans_sb"], size=10,
              color=COLORS["white"] if tone == "dark" else COLORS["green"],
              bold=True, anchor="top", line_spacing=1.20)
        # Description
        body(s, desc, x=3.10, y=ry + 0.02, w=3.85, h=row_h,
             size=9, tone=tone, line_spacing=1.32)

    # === Right column =====================================================
    right_x = 7.40
    right_w = SLIDE_W - right_x - 0.55

    # Table
    table(s, list(table_headers), list(table_rows),
          x=right_x, y=1.05, w=right_w, h=3.10,
          tone=tone, font_size=9.5,
          col_aligns=["left"] + ["left"] * (len(table_headers) - 1))

    # Stat strip title — gold/cream with underline rule
    _text(s, stat_strip_title, x=right_x, y=4.45, w=right_w, h=0.32,
          font=FONTS["sans_sb"], size=12, color=label_color,
          bold=True, anchor="top")
    _rect(s, right_x, 4.82, 1.20, 0.014, fill=label_color)

    # Stat strip — gold serif values matching the reference deck
    n_stats = len(stats)
    sw = right_w / max(1, n_stats)
    for i, (v, lbl) in enumerate(stats):
        sx = right_x + i * sw
        # Pick a smaller declared size for wide values (e.g. "EUR 7.3B")
        # so it fits on one line. Design-time sizing, not TEXT_TO_FIT_SHAPE.
        value_size = 32 if len(v) > 5 else 36
        _text(s, v, x=sx, y=5.05, w=sw - 0.05, h=0.70,
              font=FONTS["serif_l"], size=value_size,
              color=COLORS["gold"] if tone == "dark" else COLORS["green"],
              line_spacing=1.0, anchor="top")
        _text(s, lbl, x=sx, y=5.78, w=sw - 0.05, h=0.85,
              font=FONTS["sans_l"], size=9.5,
              color=COLORS["white"] if tone == "dark" else COLORS["ink_2"],
              anchor="top", line_spacing=1.30)

    # Callout — opt-in: only render when both title and body are supplied.
    # When omitted, the stat strip is allowed to breathe into the freed space
    # rather than forcing a bottom-strap on every case-study slide.
    # Use `predict_callout_h` so the callout grows downward (away from the
    # framework cards above) when the body is long. The bottom-anchor is at
    # SLIDE_H - 0.65 (safe-bottom); `cal_h` is clamped to keep the top above
    # the stat-strip rendered around y=5.50.
    if callout_title and callout_body:
        STAT_STRIP_BOT = 5.50
        SAFE_BOT_LOCAL = SLIDE_H - 0.65
        max_cal_h = SAFE_BOT_LOCAL - STAT_STRIP_BOT - 0.15
        cal_h = min(max_cal_h,
                    max(0.95, predict_callout_h(callout_body, w=right_w)))
        cal_y = SAFE_BOT_LOCAL - cal_h
        callout(s, title=callout_title, body_text=callout_body,
                x=right_x, y=cal_y, w=right_w, h=cal_h, tone=tone,
                tag=callout_tag)

    _paint_footer(s, tone)
    return s


# ---------- 6. Worksheet table (reference slide #3) -----------------------

def worksheet_table(deck: Presentation, *, eyebrow_text: str, title: str,
                    intro: str, assumptions: Sequence[Tuple[str, str]],
                    table_headers: Sequence[str],
                    table_rows: Sequence[Sequence[str]],
                    kpi_strip_items: Sequence[Tuple[str, str]],
                    tone: str = "light",
                    footnote: Optional[str] = None):
    """Big right-hand table + left-hand intro + assumptions + bottom KPI strip.

    assumptions: list of (label, value) tuples; rendered as a small bullet
    list on the left.
    """
    s = _blank_slide(deck)
    _paint_bg_native(s, tone)
    eyebrow(s, eyebrow_text, tone=tone)

    # Left column
    serif_title(s, title, x=0.55, y=0.95, w=3.85, h=2.30,
                size=28, tone=tone, line_spacing=1.05)
    body(s, intro, x=0.55, y=3.20, w=3.85, h=1.50,
         size=9.5, tone=tone, line_spacing=1.35)

    subhead(s, "Assumptions", x=0.55, y=4.80, w=3.85, h=0.30,
            size=11, tone=tone)
    a_lines = [f"-  {label}: {val}" for label, val in assumptions]
    body(s, a_lines, x=0.55, y=5.15, w=3.85, h=1.80,
         size=9, tone=tone, line_spacing=1.40)

    # Right column — table
    rx = 4.70
    rw = SLIDE_W - rx - 0.55
    n_rows = len(table_rows)
    table_h = min(5.20, 0.40 + n_rows * 0.30)
    table(s, list(table_headers), list(table_rows),
          x=rx, y=0.95, w=rw, h=table_h,
          tone=tone, font_size=9.5,
          col_aligns=["left"] + ["right"] * (len(table_headers) - 1))

    # KPI strip at bottom of right column
    kpi_y = 0.95 + table_h + 0.25
    kpi_strip(s, list(kpi_strip_items), x=rx, y=kpi_y, w=rw,
              h=min(0.95, SLIDE_H - kpi_y - 0.60),
              tone=tone, value_size=22)

    if footnote:
        _text(s, footnote, x=0.55, y=SLIDE_H - 0.55, w=SLIDE_W - 1.10, h=0.20,
              font=FONTS["sans_l"], size=9,
              color=COLORS["charcoal"], italic=True)

    _paint_footer(s, tone)
    return s


# ---------- 7. Value-prop intro (reference slide #4) ----------------------

def value_prop_intro(deck: Presentation, *, eyebrow_text: str, title: str,
                     subtitle: Optional[str] = None,
                     stats: Sequence[Tuple[str, str]] = (),
                     cards: Sequence[Tuple[str, str, str]] = (),
                     tone: str = "dark"):
    """Eyebrow + serif title + subtitle on the left; vertical stat list on
    the right; bottom row of 3-4 decimal-numbered cards (01, 02, 03, 04).

    cards: list of (number_str, card_title, card_body).
    stats: list of (value, label).
    """
    s = _blank_slide(deck)
    # Split-tone canvas: dark top (eyebrow + title + stats), light bottom (cards).
    # Mirrors the reference deck "Manufacturing Site Selection for NIO" feel.
    split_y = 3.50
    _paint_bg_native(s, "dark" if tone == "dark" else "light", log=False)
    if tone == "dark":
        _rect(s, 0, split_y, SLIDE_W, SLIDE_H - split_y,
              fill=COLORS["page_light"])
    # Log as "split" when tone='dark' (the canonical split-tone case),
    # otherwise log as plain "light".
    _log_tone(deck, "split" if tone == "dark" else "light")

    # --- Top (dark) zone ---
    eyebrow(s, eyebrow_text, tone="dark", accent="gold")

    serif_title(s, title, x=0.55, y=1.05, w=8.20, h=1.80,
                size=32, tone="dark", line_spacing=1.06)
    if subtitle:
        body(s, subtitle, x=0.55, y=2.80, w=8.20, h=0.65,
             size=12,
             color=COLORS["mint"],
             tone="dark", line_spacing=1.32)

    # Right side stat column — gold serif numerals matching the reference
    if stats:
        sx = 9.40
        sw = SLIDE_W - sx - 0.55
        n = len(stats)
        sy0 = 1.00
        sh = (split_y - 0.30 - sy0) / max(1, n)
        # Reserve a fixed band for the value so the label never collides with it
        value_band = 0.78
        for i, (v, lbl) in enumerate(stats):
            sy = sy0 + i * sh
            _text(s, v, x=sx, y=sy, w=sw, h=value_band,
                  font=FONTS["serif_l"], size=32,
                  color=COLORS["gold"],
                  line_spacing=1.0, anchor="top")
            # Label sits just below the value band, no overlap
            label_y = sy + value_band - 0.10
            label_h = max(0.20, sh - value_band)
            _text(s, lbl, x=sx, y=label_y, w=sw, h=label_h,
                  font=FONTS["sans_l"], size=10,
                  color=COLORS["white"],
                  anchor="top", line_spacing=1.20)
            # Thin mint rule between stats (not after the last one)
            if i < n - 1:
                _rect(s, sx, sy + sh - 0.06, sw * 0.55, 0.012,
                      fill=COLORS["mint"])

    # --- Bottom (light) zone — WHY CBRE + 4 cards ---
    eyebrow(s, "WHY CBRE", tone="light", accent="mint",
            x=0.55, y=split_y + 0.22, underline_w=0.85)

    if cards:
        n = len(cards)
        avail_w = SLIDE_W - 1.10
        gap = 0.20
        col_w = (avail_w - gap * (n - 1)) / n
        cy = split_y + 0.95
        ch = 6.85 - cy  # ~2.40" tall now
        for i, (num, ct, cb) in enumerate(cards):
            cx = 0.55 + i * (col_w + gap)
            # Mint top stripe
            _rect(s, cx, cy, col_w, 0.045, fill=COLORS["mint"])
            # Card body
            _rect(s, cx, cy + 0.045, col_w, ch - 0.045,
                  fill=COLORS["off_white"])
            # Decimal serif numeral
            _text(s, f"{int(num):02d}", x=cx + 0.25, y=cy + 0.12,
                  w=col_w - 0.35, h=0.55,
                  font=FONTS["serif_l"], size=30,
                  color=COLORS["mint_dark"],
                  line_spacing=1.0, anchor="top")
            # Card title
            _text(s, ct, x=cx + 0.25, y=cy + 0.65,
                  w=col_w - 0.35, h=0.60,
                  font=FONTS["sans_sb"], size=11.5,
                  color=COLORS["green"],
                  bold=True, anchor="top", line_spacing=1.18)
            # Card body — bottom portion of the card
            body_y = cy + 1.25
            body_h = max(0.40, cy + ch - body_y - 0.18)
            _text(s, cb, x=cx + 0.25, y=body_y,
                  w=col_w - 0.35, h=body_h,
                  font=FONTS["sans_l"], size=9,
                  color=COLORS["ink_2"],
                  anchor="top", line_spacing=1.30)

    _paint_footer(s, "light")
    return s


# ---------- 8. Framework with Roman numerals (slide #5 and #9) ------------

def framework_roman(deck: Presentation, *, eyebrow_text: str, title: str,
                    items: Sequence[Tuple[str, Sequence[str]]],
                    intro: Optional[str] = None,
                    side_callout: Optional[str] = None,
                    accent: str = "mint",
                    columns: int = 3, tone: str = "dark"):
    """3xN grid of Roman-numeraled cards (I, II, III, ...). The reference
    deck uses this for "5 process steps", "6 lifecycle stages" etc.

    items: list of (card_title, [bullet_1, bullet_2, ...]).
           Optionally each item can be (card_title, subtitle, [bullets]).
    intro: optional one-line intro under the title.
    side_callout: optional short paragraph in a callout on the right of the
                  title (e.g., 'CBRE leverages deep lender relationships ...').
                  Opt-in. The grid stands on its own when omitted.
    """
    s = _blank_slide(deck)
    _paint_bg_native(s, tone)
    eyebrow(s, eyebrow_text, tone=tone, accent="gold")

    # Title — left column ~8.4" wide if there's a side callout
    title_w = 8.4 if side_callout else SLIDE_W - 1.10
    serif_title(s, title, x=0.55, y=1.05, w=title_w, h=1.45,
                size=34, tone=tone, line_spacing=1.06)
    if intro:
        body(s, intro, x=0.55, y=2.55, w=title_w, h=0.65,
             size=12,
             color=COLORS["mint"] if tone == "dark" else COLORS["mint_dark"],
             tone=tone, line_spacing=1.30)

    # Optional side callout on the right — mint vertical bar + white body
    if side_callout:
        _vbar(s, x=9.20, y=1.05, h=2.05,
              color=COLORS["mint"], width_in=0.022)
        body(s, side_callout, x=9.45, y=1.00, w=SLIDE_W - 9.45 - 0.55,
             h=2.20, size=10.5,
             color=COLORS["white"] if tone == "dark" else COLORS["ink_2"],
             tone=tone, line_spacing=1.40)

    # Card grid
    grid_top = 3.30 if intro else 2.75
    grid_bottom = 6.85
    grid_h = grid_bottom - grid_top
    n_items = len(items)
    rows = (n_items + columns - 1) // columns
    gap = 0.18
    avail_w = SLIDE_W - 1.10
    col_w = (avail_w - gap * (columns - 1)) / columns
    row_h = (grid_h - gap * (rows - 1)) / rows

    # Accent handling:
    #   "cycle"  → alternate gold/mint/blue per column (debt-finance reference)
    #   "gold"/"mint"/"blue" → uniform colour across all cards
    cycle = [COLORS["gold"], COLORS["mint"], COLORS["blue"]]
    is_cycle = (accent == "cycle")

    for idx, entry in enumerate(items):
        r = idx // columns
        c = idx % columns
        cx = 0.55 + c * (col_w + gap)
        cy = grid_top + r * (row_h + gap)
        # Unpack: (title, [bullets]) or (title, subtitle, [bullets])
        if len(entry) == 2:
            ctitle, bullets = entry
            csub = None
        else:
            ctitle, csub, bullets = entry
        card_accent = cycle[c % 3] if is_cycle else accent
        roman_card(s, idx + 1, cx, cy, col_w, row_h,
                   title=ctitle, body_lines=list(bullets),
                   accent=card_accent, tone=tone, subtitle=csub)

    _paint_footer(s, tone)
    return s


# ---------- 9. Comparison table (reference slide #6) ----------------------

def comparison_table(deck: Presentation, *, eyebrow_text: str, title: str,
                     subtitle: Optional[str] = None,
                     columns: Sequence[str],
                     sections: Sequence[Tuple[str, Sequence[Sequence[str]]]],
                     footer_label: Optional[str] = None,
                     footer_values: Optional[Sequence[str]] = None,
                     tone: str = "light"):
    """N-column comparison table with mint header row, section labels on the
    left (e.g., 'Description', 'Advantages', 'Disadvantages'), and a colored
    summary row at the bottom (e.g., 'Margin Capture' percentages).

    columns: list of column header strings (first column is unused/empty
             label by convention).
    sections: list of (section_label, [[col_text, col_text, ...], ...])
              where each inner list is the column values for that row.
              For multi-line cells in a row, pass each line as its own row
              within the section.

    Simpler approach: each section has one row only; multi-bullet content is
    expressed by including \\n within a cell.
    """
    s = _blank_slide(deck)
    _paint_bg_native(s, tone)
    eyebrow(s, eyebrow_text, tone=tone)
    serif_title(s, title, x=0.55, y=0.95, w=SLIDE_W - 1.10, h=1.15,
                size=30, tone=tone, line_spacing=1.05)
    if subtitle:
        body(s, subtitle, x=0.55, y=2.10, w=SLIDE_W - 1.10, h=0.45,
             size=10, tone=tone)

    n_cols = len(columns)
    x0 = 0.55
    w = SLIDE_W - 1.10
    label_w = 1.50
    col_w = (w - label_w) / max(1, n_cols - 1)  # first column is label
    table_top = 2.70 if subtitle else 2.35
    table_bottom = SLIDE_H - 0.55 - (0.55 if footer_values else 0)
    table_h = table_bottom - table_top

    # Header row (skip first column — it's the section-label column)
    header_h = 0.40
    for c in range(1, n_cols):
        cx = x0 + label_w + (c - 1) * col_w
        _rect(s, cx + 0.04, table_top, col_w - 0.08, header_h,
              fill=COLORS["mint"])
        _text(s, columns[c], x=cx + 0.16, y=table_top, w=col_w - 0.32,
              h=header_h, font=FONTS["sans_sb"], size=10,
              color=COLORS["green"], bold=True, uppercase=False,
              align="left", anchor="middle", line_spacing=1.10)

    # Section rows
    body_top = table_top + header_h + 0.10
    body_h = table_h - header_h - 0.10
    section_h = body_h / max(1, len(sections))
    rule_color = COLORS["rule_light"] if tone == "light" else COLORS["rule_dark"]
    body_color = COLORS["ink_2"] if tone == "light" else COLORS["white"]

    for r, (label, rows) in enumerate(sections):
        ry = body_top + r * section_h
        # Section label on the left
        _text(s, label, x=x0, y=ry + 0.05, w=label_w - 0.12, h=0.40,
              font=FONTS["sans_sb"], size=10,
              color=COLORS["mint_dark"] if tone == "light" else COLORS["mint"],
              bold=True, anchor="top", line_spacing=1.15)
        # Section divider on top
        _line(s, x0, ry, x0 + w, ry, color=rule_color, width_pt=0.5)
        # Cells
        # rows might be either:
        #   - list of strings (one bullet line per column)
        #   - list of list (multiple bullets per column)
        # We assume the latter.
        for c in range(1, n_cols):
            cx = x0 + label_w + (c - 1) * col_w
            cells_text = []
            for row in rows:
                if c - 1 < len(row):
                    cells_text.append(f"- {row[c-1]}")
            body(s, cells_text, x=cx + 0.16, y=ry + 0.05,
                 w=col_w - 0.32, h=section_h - 0.10,
                 size=9, color=body_color, tone=tone, line_spacing=1.30)

    # Footer summary row
    if footer_values:
        fy = body_top + body_h
        _rect(s, x0, fy, label_w, 0.55, fill=COLORS["green"])
        _text(s, footer_label or "", x=x0 + 0.16, y=fy, w=label_w - 0.32,
              h=0.55, font=FONTS["sans_sb"], size=10,
              color=COLORS["white"], bold=True, anchor="middle")
        for c in range(1, n_cols):
            cx = x0 + label_w + (c - 1) * col_w
            _rect(s, cx + 0.04, fy, col_w - 0.08, 0.55,
                  fill=COLORS["green"])
            val = footer_values[c-1] if c-1 < len(footer_values) else ""
            _text(s, val, x=cx + 0.16, y=fy, w=col_w - 0.32, h=0.55,
                  font=FONTS["sans_sb"], size=14,
                  color=COLORS["mint"], bold=True, align="center",
                  anchor="middle")

    _paint_footer(s, tone)
    return s


# ---------- 10. Decision matrix (reference slide #7) ----------------------

def decision_matrix(deck: Presentation, *, eyebrow_text: str, title: str,
                    left_label: str, gate_label: str, right_label: str,
                    rows: Sequence[Tuple[str, str, str]],
                    callout_title: Optional[str] = None,
                    callout_body: Optional[str] = None,
                    tone: str = "dark"):
    """Three-column decision matrix: 'Favours A | Decision Gate | Favours B'.

    rows: list of (left_text, gate_text, right_text). left_text and
          right_text typically begin with a short topic bold then body
          (e.g., 'Yield > WACC | Ownership creates a positive spread...').
    """
    s = _blank_slide(deck)
    _paint_bg_native(s, tone)
    # Dynamic top-of-slide stack — eyebrow + title, no hardcoded title h.
    # Table flows from the actual end of the title rather than y=2.50.
    f = Flow(s, x=0.55, y=0.55, w=SLIDE_W - 1.10, tone=tone)
    f.eyebrow(eyebrow_text, accent="gold")
    f.gap(0.20)
    f.title(title, size=32, line_spacing=1.05)
    f.gap(0.30)

    table_top = f.y
    table_w = SLIDE_W - 1.10
    gate_w = 2.40
    side_w = (table_w - gate_w) / 2
    left_x = 0.55
    gate_x = left_x + side_w
    right_x = gate_x + gate_w
    header_h = 0.46
    has_callout = bool(callout_title or callout_body)
    body_h_bottom = 6.85 - (1.20 if has_callout else 0)
    body_top = table_top + header_h
    body_h = body_h_bottom - body_top

    body_color = COLORS["white"] if tone == "dark" else COLORS["ink_2"]
    mint = COLORS["mint"]
    gold = COLORS["gold"]
    dark = COLORS["green"]

    # === Column headers (filled bands, matches reference) ===
    # Left header — mint band, dark text
    _rect(s, left_x, table_top, side_w - 0.06, header_h, fill=mint)
    _text(s, left_label, x=left_x, y=table_top, w=side_w - 0.06, h=header_h,
          font=FONTS["sans_sb"], size=11, color=dark, bold=True,
          uppercase=True, letter_spacing=1.5, align="center", anchor="middle")
    # Gate header — darker mint band
    _rect(s, gate_x, table_top, gate_w, header_h, fill=COLORS["mint_dark"])
    _text(s, gate_label, x=gate_x, y=table_top, w=gate_w, h=header_h,
          font=FONTS["sans_sb"], size=10, color=COLORS["white"], bold=True,
          uppercase=True, letter_spacing=1.4, align="center", anchor="middle")
    # Right header — gold/cream band
    _rect(s, right_x + 0.06, table_top, side_w - 0.06, header_h, fill=gold)
    _text(s, right_label, x=right_x + 0.06, y=table_top, w=side_w - 0.06,
          h=header_h, font=FONTS["sans_sb"], size=11, color=dark, bold=True,
          uppercase=True, letter_spacing=1.5, align="center", anchor="middle")

    # === Body rows ===
    n = len(rows)
    rh = body_h / max(1, n)

    # Mint gate column body
    _rect(s, gate_x, body_top, gate_w, body_h, fill=COLORS["green_2"])

    for i, (lt, gt, rt) in enumerate(rows):
        ry = body_top + i * rh
        # Row separator rules — mint on dark, mint_dark on light
        if i > 0:
            _line(s, left_x, ry, left_x + side_w - 0.10, ry,
                  color=mint, width_pt=0.5)
            _line(s, right_x + 0.10, ry, right_x + side_w, ry,
                  color=mint, width_pt=0.5)
        # Left cell — key term in MINT (matches the mint column header,
        # signals "ownership/win" side). Body in white/ink.
        _render_row_text(s, lt, left_x + 0.10, ry + 0.05,
                         side_w - 0.32, rh - 0.10,
                         tone=tone, key_color=mint, body_color=body_color)
        # Gate cell — white/cream condition label on the dark gate band
        _text(s, gt, x=gate_x + 0.18, y=ry, w=gate_w - 0.36, h=rh,
              font=FONTS["sans_sb"], size=11, color=COLORS["white"],
              bold=True, align="center", anchor="middle", line_spacing=1.18)
        # Right cell — key term in GOLD (matches the gold column header,
        # signals "lease/consideration" side). Body in white/ink.
        _render_row_text(s, rt, right_x + 0.20, ry + 0.05,
                         side_w - 0.32, rh - 0.10,
                         tone=tone, key_color=gold, body_color=body_color)

    # Optional CBRE View callout at the bottom. Use predict_callout_h to
    # size for the actual body so the bg rect contains all rendered text.
    if has_callout:
        cb = callout_body or ""
        cal_h = max(1.00, predict_callout_h(cb, w=SLIDE_W - 1.10))
        callout(s, title=callout_title or "CBRE VIEW",
                body_text=cb,
                x=0.55, y=body_h_bottom + 0.18, w=SLIDE_W - 1.10,
                h=cal_h, tone=tone)

    _paint_footer(s, tone)
    return s


def _render_row_text(slide, text: str, x: float, y: float, w: float, h: float,
                     *, tone: str, key_color: RGBColor, body_color: RGBColor):
    """Render a decision-matrix row cell. If the text contains '|', split into
    a gold/cream serif lead and a body paragraph. Otherwise just body.
    Heights are clamped to a minimum positive value so very-small rows don't
    produce a malformed XML box.
    """
    lead_h = min(0.36, max(0.20, h * 0.42))
    if "|" in text:
        lead, rest = text.split("|", 1)
        _text(slide, lead.strip(), x=x, y=y, w=w, h=lead_h,
              font=FONTS["serif"], size=15,
              color=key_color, anchor="top", line_spacing=1.10)
        body_y = y + lead_h + 0.04
        body_h_used = max(0.18, h - lead_h - 0.06)
        body(slide, rest.strip(), x=x, y=body_y, w=w, h=body_h_used,
             size=9.5, color=body_color, tone=tone, line_spacing=1.32)
    else:
        body(slide, text, x=x, y=y, w=w, h=max(0.20, h),
             size=10, color=body_color, tone=tone, line_spacing=1.32)


# ---------- 11a. Why columns (N-column reasons slide) --------------------

def why_columns(deck: Presentation, *, eyebrow_text: str, title: str,
                columns: Sequence[Union[Tuple[str, str], Tuple[str, str, str]]],
                intro: Optional[str] = None,
                takeaway: Optional[str] = None,
                takeaway_label: Optional[str] = None,
                numbered: bool = False,
                tone: str = "light",
                accent: str = "bright_green"):
    """N-column 'reasons / pillars' slide. 3-5 columns work cleanly.

    Each column is `(header, body)` or `(header, body, tagline)`. The optional
    tagline renders in gold/cream italic at the column's bottom so all five
    end at the same y — gives the slide a baseline grid instead of a ragged
    bottom edge.

    Optional density features:
      - `intro`: 1-2 sentence anchor below the title (mint vertical bar to left)
      - `takeaway` (+ optional `takeaway_label`): bottom CBRE-VIEW callout band
        that runs flush edge-to-edge
      - `numbered=True`: small serif numeral tag (12pt) above each header

    accent: "bright_green" (matches the reference's spring-green underlines),
            "gold", "mint", or "blue".
    """
    s = _blank_slide(deck)
    _paint_bg_native(s, tone)

    accent_color = {
        "bright_green": COLORS["bright_green"],
        "gold":         COLORS["gold"] if tone == "dark" else COLORS["mint_dark"],
        "mint":         COLORS["mint"],
        "blue":         COLORS["blue"],
    }.get(accent, COLORS["bright_green"])
    # Column underline uses mint_dark on light bg (calmer than the eyebrow's
    # bright accent) so five rules don't read as a chart axis.
    column_rule_color = (
        COLORS["mint_dark"] if tone == "light" else COLORS["mint"]
    )
    # Eyebrow accent — calmer than column rules on a light bg so the eyebrow's
    # accent identity reads distinctly.
    eyebrow_rule_color = (
        COLORS["mint_dark"] if tone == "light" else COLORS["gold"]
    )

    # Eyebrow — accent rule + label
    _rect(s, 0.55, 0.55 + 0.36, 1.65, 0.018, fill=eyebrow_rule_color)
    _text(s, eyebrow_text, x=0.55, y=0.55, w=8.0, h=0.28,
          font=FONTS["sans_sb"], size=10,
          color=COLORS["ink"] if tone == "light" else COLORS["white"],
          bold=True, uppercase=True, letter_spacing=1.5, anchor="top")

    has_intro = bool(intro)
    has_takeaway = bool(takeaway)

    # ── Vertical rhythm (tightened to close dead bands) ───────────────────
    # Title sits just under the eyebrow; intro pulls up against the title;
    # columns start immediately after; takeaway band sits flush at the bottom.
    if has_intro and has_takeaway:
        title_size, title_y, title_h = 30, 1.10, 1.20
        intro_y, intro_h = 2.35, 0.85
        col_top, col_bottom = 3.40, 6.05
        cb_y, cb_h = 6.30, 0.80
    elif has_intro:
        title_size, title_y, title_h = 34, 1.10, 1.45
        intro_y, intro_h = 2.65, 0.95
        col_top, col_bottom = 3.75, 7.10
        cb_y = cb_h = 0
    elif has_takeaway:
        title_size, title_y, title_h = 38, 1.10, 1.70
        col_top, col_bottom = 2.95, 6.05
        cb_y, cb_h = 6.30, 0.80
    else:
        title_size, title_y, title_h = 40, 1.15, 1.95
        col_top, col_bottom = 3.30, 7.10
        cb_y = cb_h = 0
    if len(title) > 60:
        title_size = max(28, title_size - 6)

    serif_title(s, title, x=0.55, y=title_y, w=11.5, h=title_h,
                size=title_size, tone=tone, line_spacing=1.06)

    if has_intro:
        _vbar(s, x=0.55, y=intro_y, h=intro_h,
              color=COLORS["mint"] if tone == "dark" else COLORS["mint_dark"],
              width_in=0.022)
        body(s, intro, x=0.75, y=intro_y - 0.02, w=SLIDE_W - 1.30,
             h=intro_h + 0.05,
             size=11.5, tone=tone, line_spacing=1.40)

    # ── Column grid ───────────────────────────────────────────────────────
    n = len(columns)
    avail_w = SLIDE_W - 1.10
    gap = 0.30
    col_w = (avail_w - gap * (n - 1)) / n
    header_color = COLORS["white"] if tone == "dark" else COLORS["ink"]
    body_color   = COLORS["white"] if tone == "dark" else COLORS["ink_2"]
    num_color    = (COLORS["gold"] if tone == "dark"
                    else COLORS["mint_dark"])
    tagline_color = (COLORS["gold"] if tone == "dark"
                     else COLORS["mint_dark"])
    col_h = col_bottom - col_top

    # Pre-decide whether any column has a tagline — drives baseline grid
    any_tagline = any(len(c) >= 3 and c[2] for c in columns)
    tagline_band_h = 0.45 if any_tagline else 0.0

    for i, entry in enumerate(columns):
        header, body_text, tagline = entry if len(entry) == 3 else (entry[0], entry[1], None)
        cx = 0.55 + i * (col_w + gap)

        if numbered:
            # Demoted serif numeral — 12pt mono-feel tag above the header
            _text(s, f"{i+1:02d}", x=cx, y=col_top, w=col_w, h=0.24,
                  font=FONTS["serif"], size=12, color=num_color,
                  anchor="top", line_spacing=1.0)
            header_y = col_top + 0.30
        else:
            header_y = col_top

        # Bold sans header
        _text(s, header, x=cx, y=header_y, w=col_w, h=0.42,
              font=FONTS["sans_sb"], size=13.5, color=header_color,
              bold=True, anchor="top", line_spacing=1.16)
        # Hairline accent underline — narrower (0.85") and thinner (0.5pt)
        # so five rules don't read like a chart axis.
        rule_w = 0.85
        rule_h_in = 0.008  # ≈ 0.5pt
        _rect(s, cx, header_y + 0.48, rule_w, rule_h_in,
              fill=column_rule_color)

        # Body paragraph — font size adapts to available column height
        body_y = header_y + 0.65
        body_h_avail = col_top + col_h - body_y - tagline_band_h - 0.08
        if body_h_avail >= 1.85:
            body_size, body_ls = 10.5, 1.45
        elif body_h_avail >= 1.40:
            body_size, body_ls = 9.5, 1.36
        else:
            body_size, body_ls = 9.0, 1.30
        body(s, body_text, x=cx, y=body_y, w=col_w,
             h=max(0.30, body_h_avail),
             size=body_size, color=body_color, tone=tone,
             line_spacing=body_ls)

        # Optional tagline — sits at a fixed bottom anchor so all columns share
        # a baseline grid even when body paragraphs ragged at different depths.
        if tagline:
            tag_y = col_top + col_h - tagline_band_h + 0.05
            _text(s, tagline, x=cx, y=tag_y, w=col_w, h=tagline_band_h - 0.05,
                  font=FONTS["serif"], size=10.5, color=tagline_color,
                  italic=True, anchor="top", line_spacing=1.25)

    # ── Bottom takeaway band — flush edge-to-edge ─────────────────────────
    if takeaway:
        # Full-bleed surface fill so the band reads as a true page band, not
        # a floating card.
        fill = (COLORS["green_3"] if tone == "dark"
                else COLORS["off_white"])
        _rect(s, 0, cb_y, SLIDE_W, cb_h, fill=fill)
        # Left bar at the column-1 x-origin so the band aligns with the grid
        bar_x = 0.55
        _rect(s, bar_x, cb_y + 0.08, 0.06, cb_h - 0.16,
              fill=COLORS["mint"] if tone == "dark" else COLORS["mint_dark"])
        # Gold label — semibold uppercase
        label_color = COLORS["gold"] if tone == "dark" else COLORS["mint_dark"]
        _text(s, (takeaway_label or "CBRE VIEW"),
              x=0.75, y=cb_y + 0.14, w=2.10, h=0.30,
              font=FONTS["sans_sb"], size=11, color=label_color,
              bold=True, uppercase=True, letter_spacing=1.8, anchor="top")
        # Body — ink on light, white on dark
        body_x = 2.95
        _text(s, takeaway, x=body_x, y=cb_y + 0.14,
              w=SLIDE_W - body_x - 0.55, h=cb_h - 0.20,
              font=FONTS["sans_l"], size=10.5,
              color=COLORS["ink"] if tone == "light" else COLORS["white"],
              anchor="top", line_spacing=1.38)

    _paint_footer(s, tone)
    return s


# ---------- 11. Why two-column (reference slide #8) -----------------------

def why_two_col(deck: Presentation, *, eyebrow_text: str, title: str,
                intro: str,
                drivers_label: str,
                drivers: Sequence[Tuple[str, str, str]],
                right_eyebrow: str,
                cards: Sequence[Tuple[str, str, str]],
                tone: str = "dark"):
    """Left column: title, intro, numbered 'drivers' list (01-06).
    Right column: eyebrow + grid of 4 mint cards.

    drivers: list of (number, topic, description).
    cards: list of (card_title, when_label, when_body). Each card has a small
           heading and a definition-style body.
    """
    s = _blank_slide(deck)
    _paint_bg_native(s, tone)
    # Dual eyebrows — gold left, mint right (matches reference deck)
    eyebrow(s, eyebrow_text, tone=tone, accent="gold")

    # Left column
    lx = 0.55
    lw = 6.0
    serif_title(s, title, x=lx, y=1.05, w=lw, h=1.70,
                size=32, tone=tone, line_spacing=1.06)
    # Mint vertical bar next to intro
    _vbar(s, x=lx, y=2.90, h=1.30,
          color=COLORS["mint"] if tone == "dark" else COLORS["mint_dark"],
          width_in=0.022)
    body(s, intro, x=lx + 0.20, y=2.85, w=lw - 0.20, h=1.40,
         size=10, tone=tone, line_spacing=1.42)

    # Drivers label — gold/cream with underline
    label_color = COLORS["gold"] if tone == "dark" else COLORS["ink"]
    _text(s, drivers_label, x=lx, y=4.35, w=lw, h=0.32,
          font=FONTS["sans_sb"], size=12, color=label_color,
          bold=True, anchor="top")
    _rect(s, lx, 4.72, 1.40, 0.014, fill=label_color)
    drv_top = 4.95
    drv_h = 6.85 - drv_top
    rh = drv_h / max(1, len(drivers))
    for i, (n, topic, desc) in enumerate(drivers):
        ry = drv_top + i * rh
        _text(s, n, x=lx, y=ry + 0.02, w=0.55, h=rh,
              font=FONTS["serif"], size=14,
              color=COLORS["mint"] if tone == "dark" else COLORS["mint_dark"],
              anchor="top")
        _text(s, topic, x=lx + 0.50, y=ry + 0.02, w=1.60, h=rh,
              font=FONTS["sans_sb"], size=10,
              color=COLORS["white"] if tone == "dark" else COLORS["green"],
              bold=True, anchor="top", line_spacing=1.18)
        body(s, desc, x=lx + 2.15, y=ry + 0.02, w=lw - 2.20, h=rh,
             size=9, tone=tone, line_spacing=1.32)

    # Right column — mint eyebrow
    rx = 7.00
    rw = SLIDE_W - rx - 0.55
    eyebrow(s, right_eyebrow, tone=tone, accent="mint", x=rx, y=0.55,
            underline_w=2.20)

    n_cards = len(cards)
    cols = 2
    rows = (n_cards + cols - 1) // cols
    gap = 0.18
    card_w = (rw - gap * (cols - 1)) / cols
    grid_top = 1.30
    grid_bottom = 6.85
    card_h = (grid_bottom - grid_top - gap * (rows - 1)) / rows

    # Cards: lifted dark-green surface with mint top stripe + interior layout
    # mirrors the reference deck "MOST COMMON FUNDING STRATEGIES" board.
    for i, (ct, when_text, body_text) in enumerate(cards):
        r = i // cols
        c = i % cols
        cx = rx + c * (card_w + gap)
        cy = grid_top + r * (card_h + gap)
        # Card body
        _rect(s, cx, cy, card_w, card_h,
              fill=COLORS["green_2"] if tone == "dark" else COLORS["off_white"])
        # Mint top stripe
        _rect(s, cx, cy, card_w, 0.045, fill=COLORS["mint"])

        # Title — white bold
        _text(s, ct, x=cx + 0.22, y=cy + 0.18, w=card_w - 0.44, h=0.36,
              font=FONTS["sans_sb"], size=12,
              color=COLORS["white"] if tone == "dark" else COLORS["green"],
              bold=True, line_spacing=1.12)
        # Mint underline rule under title
        _rect(s, cx + 0.22, cy + 0.58, card_w - 0.44, 0.012,
              fill=COLORS["mint"])

        # WHEN: row — gold/cream "When:" label
        wt = when_text.lstrip()
        if wt.lower().startswith("when:"):
            wt = wt[5:].lstrip()
        _text(s, "When:", x=cx + 0.22, y=cy + 0.72, w=0.60, h=0.26,
              font=FONTS["sans_sb"], size=9.5,
              color=COLORS["gold"] if tone == "dark" else COLORS["ink"],
              bold=True, anchor="top")
        _text(s, wt, x=cx + 0.85, y=cy + 0.72, w=card_w - 1.05, h=0.55,
              font=FONTS["sans_sb"], size=9.5,
              color=COLORS["white"] if tone == "dark" else COLORS["ink_2"],
              bold=True, line_spacing=1.22, anchor="top")

        # Body
        body(s, body_text, x=cx + 0.22, y=cy + 1.35, w=card_w - 0.44,
             h=card_h - 1.50, size=9,
             color=COLORS["white"] if tone == "dark" else COLORS["ink_2"],
             tone=tone, line_spacing=1.32)

    _paint_footer(s, tone)
    return s


# ---------- 12. Single hero stat ------------------------------------------

def stat_hero(deck: Presentation, *, eyebrow_text: str, title: str,
              stat: str, label: str, footnote: Optional[str] = None,
              tone: str = "dark"):
    """Single dramatic stat on a slide: big title at top, giant number lower."""
    s = _blank_slide(deck)
    _paint_bg_native(s, tone)
    eyebrow(s, eyebrow_text, tone=tone, accent="gold")

    serif_title(s, title, x=0.55, y=1.10, w=SLIDE_W - 1.10, h=2.20,
                size=44, tone=tone, line_spacing=1.06)

    # Giant stat — gold/cream on dark, deep green on light
    _text(s, stat, x=0.55, y=3.40, w=SLIDE_W - 1.10, h=2.80,
          font=FONTS["serif_l"], size=200,
          color=COLORS["gold"] if tone == "dark" else COLORS["green"],
          line_spacing=1.0, anchor="top")
    # Label
    body(s, label, x=0.55, y=6.30, w=SLIDE_W - 1.10, h=0.50,
         size=15, tone=tone)

    if footnote:
        _text(s, footnote, x=0.55, y=SLIDE_H - 0.55, w=SLIDE_W - 1.10,
              h=0.20, font=FONTS["sans_l"], size=9,
              color=COLORS["charcoal"], italic=True)

    _paint_footer(s, tone)
    return s


# ---------- 13. Stat strip slide ------------------------------------------

def stat_strip(deck: Presentation, *, eyebrow_text: str, title: str,
               subtitle: Optional[str] = None,
               stats: Sequence[Tuple[str, str]],
               body_text: Optional[str] = None,
               coverage: Optional[Sequence[Tuple[str, str]]] = None,
               tone: str = "dark"):
    """Three-to-five hero stats side-by-side under a serif title.

    `coverage` and `body_text` are optional. A clean stat row with deliberate
    empty space below it is a legitimate output. Add a coverage band only
    when there is a real second beat to introduce — never to "fill" the
    bottom.
    """
    s = _blank_slide(deck)
    _paint_bg_native(s, tone)
    eyebrow(s, eyebrow_text, tone=tone, accent="gold")
    serif_title(s, title, x=0.55, y=1.05, w=SLIDE_W - 1.10, h=1.50,
                size=36, tone=tone, line_spacing=1.06)
    if subtitle:
        body(s, subtitle, x=0.55, y=2.55, w=SLIDE_W - 1.10, h=0.65,
             size=13, tone=tone, line_spacing=1.30,
             color=COLORS["mint"] if tone == "dark" else COLORS["mint_dark"])
    strip_y = 3.40
    kpi_strip(s, list(stats), x=0.55, y=strip_y, w=SLIDE_W - 1.10, h=2.10,
              tone=tone, value_size=72)

    # Coverage band — fills the bottom third
    if coverage:
        band_top = 5.85
        _rect(s, 0.55, band_top, SLIDE_W - 1.10, 0.014,
              fill=COLORS["gold"] if tone == "dark" else COLORS["ink"])
        n = min(len(coverage), 6)
        col_w = (SLIDE_W - 1.10) / n
        accent = COLORS["gold"] if tone == "dark" else COLORS["ink"]
        text_color = COLORS["white"] if tone == "dark" else COLORS["ink"]
        for i, (label, detail) in enumerate(coverage[:n]):
            cx = 0.55 + i * col_w
            _text(s, label, x=cx, y=band_top + 0.18, w=col_w - 0.25, h=0.32,
                  font=FONTS["sans_sb"], size=10, color=accent,
                  uppercase=True, letter_spacing=2.0, anchor="top")
            _text(s, detail, x=cx, y=band_top + 0.55, w=col_w - 0.25, h=0.85,
                  font=FONTS["sans_l"], size=11, color=text_color,
                  anchor="top", line_spacing=1.30)
    elif body_text:
        body(s, body_text, x=0.55, y=6.05, w=SLIDE_W - 1.10, h=1.05,
             size=11, tone=tone, line_spacing=1.35)
    _paint_footer(s, tone)
    return s


# ---------- 14. Statement / pull quote ------------------------------------

def statement(deck: Presentation, *, text: str,
              attribution: Optional[str] = None,
              eyebrow_text: Optional[str] = None,
              support: Optional[str] = None,
              support_label: Optional[str] = None,
              pillars: Optional[Sequence[Tuple[str, str]]] = None,
              tone: str = "dark"):
    """A dramatic pull quote / declarative statement slide.

    All support args are optional. A bare statement — eyebrow + headline (+
    quiet attribution) — is a legitimate finished slide, not an incomplete
    one. Reach for `support` / `pillars` only when the headline genuinely
    needs to lean on a second beat; deliberate emptiness beneath a strong
    statement reads as confidence, not as missing content.

    support_label  -- optional eyebrow for the support paragraph
    support        -- optional 1-2 sentence elaboration under the quote
    pillars        -- optional 3-4 (label, detail) pairs along the bottom
    """
    s = _blank_slide(deck)
    _paint_bg_native(s, tone)
    if eyebrow_text:
        eyebrow(s, eyebrow_text, tone=tone, accent="gold")

    # Quote size scales with content
    quote_size = 54 if len(text) <= 90 else 44 if len(text) <= 150 else 36
    serif_title(s, text, x=0.55, y=1.55, w=SLIDE_W - 1.10, h=2.80,
                size=quote_size, tone=tone, line_spacing=1.10)

    # Support paragraph — picks up the empty middle
    if support:
        accent = COLORS["gold"] if tone == "dark" else COLORS["ink"]
        if support_label:
            _text(s, support_label, x=0.55, y=4.55, w=4.0, h=0.30,
                  font=FONTS["sans_sb"], size=10, color=accent,
                  uppercase=True, letter_spacing=2.2, anchor="top")
            _rect(s, 0.55, 4.92, 1.20, 0.014, fill=accent)
            sup_y = 5.10
        else:
            _rect(s, 0.55, 4.60, 0.60, 0.018, fill=accent)
            sup_y = 4.82
        body(s, support, x=0.55, y=sup_y, w=SLIDE_W - 1.10, h=1.20,
             size=14, tone=tone, line_spacing=1.42)

    # Three or four supporting pillars — fills the bottom band
    if pillars:
        band_top = 6.10 if support else 5.10
        n = min(len(pillars), 4)
        col_w = (SLIDE_W - 1.10) / n
        accent = COLORS["gold"] if tone == "dark" else COLORS["ink"]
        text_color = COLORS["white"] if tone == "dark" else COLORS["ink"]
        for i, (label, detail) in enumerate(pillars[:n]):
            cx = 0.55 + i * col_w
            _rect(s, cx, band_top, col_w - 0.30, 0.014, fill=accent)
            _text(s, label, x=cx, y=band_top + 0.10, w=col_w - 0.30, h=0.28,
                  font=FONTS["sans_sb"], size=10, color=accent,
                  uppercase=True, letter_spacing=2.0, anchor="top")
            _text(s, detail, x=cx, y=band_top + 0.42, w=col_w - 0.30, h=0.62,
                  font=FONTS["sans_l"], size=11, color=text_color,
                  anchor="top", line_spacing=1.28)

    if attribution:
        _text(s, attribution, x=SLIDE_W - 5.55, y=7.10, w=4.40, h=0.25,
              font=FONTS["sans_sb"], size=9,
              color=COLORS["gold"] if tone == "dark" else COLORS["ink"],
              uppercase=True, letter_spacing=2.2, anchor="middle",
              align="right")

    _paint_footer(s, tone)
    return s


# ===========================================================================
# Editorial-bold composition helpers
# ---------------------------------------------------------------------------
# Tested building blocks for the "editorial-bold" house style: every content
# slide gets a composition that matches its rhetorical job (a statement, a
# phase timeline, a from->to transition, a ladder, asymmetric panels) instead
# of the uniform "header + card-row + callout" template. See
# references/editorial-archetypes.md for the catalogue and worked examples.
#
# Conventions:
#   - Atoms: editorial_header, num_row, chip, arrow.
#   - Compositions: phase_timeline, from_to, tier_ladder, directional_ladder,
#     intensity_bars.
#   - All take the slide first, then keyword args; honour tone="dark"|"light".
#   - Geometry is back-solved against ED_* constants below; callers can pass
#     explicit x/y/w to override.
# ===========================================================================

ED_X = SAFE_L                    # 0.55" left margin
ED_W = SLIDE_W - 2 * SAFE_L      # full content width
ED_EYEBROW_Y = 0.50
ED_SAFE_BOT = 6.85               # content bottom edge (clears wordmark band)


def editorial_header(slide, *, eyebrow_text, title, tone="dark", intro=None,
                     title_size=30, accent="gold", x=None, w=None):
    """Eyebrow -> serif title -> optional intro, drawn at absolute coords.

    Returns the y where slide content can begin (title bottom + intro + a
    0.30" gap). Use this at the top of every editorial slide, then place the
    archetype below the returned y. Keep the intro to one or two lines — drop
    it entirely when the slide body is tall (timelines, ladders, panels), so
    the content has room above ED_SAFE_BOT.
    """
    x = ED_X if x is None else x
    w = ED_W if w is None else w
    ink = COLORS["white"] if tone == "dark" else COLORS["green"]
    eyebrow(slide, eyebrow_text, tone=tone, x=x, y=ED_EYEBROW_Y, accent=accent)
    ty = 1.02
    th = measure_text(title, size=title_size, w=w, font="serif",
                      line_spacing=1.05)
    serif_title(slide, title, x=x, y=ty, w=w, h=th, size=title_size,
                tone=tone, color=ink)
    y = ty + th
    if intro:
        y += 0.14
        ih = measure_text(intro, size=12.5, w=w, font="sans_l",
                          line_spacing=1.30)
        body(slide, intro, x=x, y=y, w=w, h=ih, size=12.5, tone=tone,
             color=COLORS["mint_pale"] if tone == "dark" else COLORS["ink_2"])
        y += ih
    return y + 0.30


def num_row(slide, n, title, detail=None, *, x, y, w, tone="dark",
            accent=None):
    """One editorial numbered row: big serif numeral + bold title + detail.

    The building block for the asymmetric "left statement / right list" and
    "required items" archetypes. Lay several down the page with thin rules
    between them (draw the rules yourself with `_line`). Width must be
    positive; the numeral occupies the left 1.05"."""
    accent = accent or (COLORS["gold"] if tone == "dark"
                        else COLORS["mint_dark"])
    title_c = COLORS["white"] if tone == "dark" else COLORS["green"]
    detail_c = COLORS["mint_pale"] if tone == "dark" else COLORS["ink_2"]
    _text(slide, f"{int(n):02d}", x=x, y=y, w=0.90, h=0.55,
          font=FONTS["serif_l"], size=26, color=accent, anchor="top")
    _text(slide, title, x=x + 1.05, y=y + 0.02, w=w - 1.05, h=0.30,
          font=FONTS["sans_sb"], size=13, color=title_c, bold=True,
          anchor="top")
    if detail:
        _text(slide, detail, x=x + 1.05, y=y + 0.34, w=w - 1.05, h=0.40,
              font=FONTS["sans_l"], size=11, color=detail_c, anchor="top",
              line_spacing=1.25)


def chip(slide, text, *, x, y, w, h, fill=None, text_color=None,
         line=None, tone="dark", size=10.5):
    """Rounded-rectangle pill with centred label. Use for country/tag chips,
    ladder items, status flags. On light backgrounds pass `line` for a subtle
    border so the chip reads against the panel."""
    fill = fill or (COLORS["green_4"] if tone == "dark"
                    else COLORS["off_white"])
    text_color = text_color or (COLORS["white"] if tone == "dark"
                                else COLORS["green"])
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    _text(slide, text, x=x, y=y, w=w, h=h, font=FONTS["sans_sb"], size=size,
          color=text_color, bold=True, align="center", anchor="middle")
    return shp


_ARROW_SHAPES = {
    "right": MSO_SHAPE.RIGHT_ARROW, "left": MSO_SHAPE.LEFT_ARROW,
    "up": MSO_SHAPE.UP_ARROW, "down": MSO_SHAPE.DOWN_ARROW,
}


def arrow(slide, direction, *, x, y, w, h, color):
    """Solid directional arrow (right/left/up/down). For from->to transitions
    and directional ladders. Semantics by colour: mint=positive/forward,
    gold=lateral/refocus, blue=down/deprioritise."""
    _assert_pos_dims(w, h, "arrow")
    shp = slide.shapes.add_shape(_ARROW_SHAPES[direction],
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def phase_timeline(slide, phases, *, y=4.00, tone="dark", x=None, w=None,
                   here_label="WE ARE HERE"):
    """Horizontal phase / stage track. `phases` = list of
    (number, label, desc, done) tuples — done=True renders a mint node, the
    current phase (done=False) renders a larger gold node with a `here_label`
    tag above it. Use for "where we are in the journey" status slides.
    `y` is the baseline of the connecting line."""
    x = ED_X if x is None else x
    w = ED_W if w is None else w
    n = len(phases)
    seg = w / n
    line_y = y
    cx0 = x + seg * 0.5
    cxN = x + seg * (n - 0.5)
    _line(slide, cx0, line_y + 0.11, cxN, line_y + 0.11,
          color=COLORS["mint_dark"] if tone == "dark" else COLORS["rule_light"],
          width_pt=1.5)
    for i, (num, label, desc, done) in enumerate(phases):
        cx = x + seg * (i + 0.5)
        col_l = x + seg * i + 0.25
        col_w = seg - 0.50
        accent = COLORS["mint"] if done else COLORS["gold"]
        if not done and here_label:
            _text(slide, here_label, x=col_l, y=line_y - 1.22, w=col_w, h=0.24,
                  font=FONTS["sans_sb"], size=10, color=COLORS["gold"],
                  bold=True, uppercase=True, letter_spacing=2.2,
                  align="center", anchor="middle")
        _text(slide, str(num), x=col_l, y=line_y - 0.92, w=col_w, h=0.52,
              font=FONTS["serif_l"], size=30, color=accent,
              align="center", anchor="middle")
        _text(slide, label, x=col_l, y=line_y - 0.40, w=col_w, h=0.30,
              font=FONTS["sans_sb"], size=13,
              color=COLORS["white"] if tone == "dark" else COLORS["green"],
              bold=True, uppercase=True, letter_spacing=2.5,
              align="center", anchor="middle")
        m = 0.30 if not done else 0.22
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(cx - m / 2), Inches(line_y + 0.11 - m / 2),
            Inches(m), Inches(m))
        node.fill.solid()
        node.fill.fore_color.rgb = accent
        node.line.fill.background()
        node.shadow.inherit = False
        _text(slide, desc, x=col_l, y=line_y + 0.42, w=col_w, h=0.80,
              font=FONTS["sans_l"], size=11,
              color=COLORS["mint_pale"] if tone == "dark" else COLORS["ink_2"],
              align="center", anchor="top", line_spacing=1.30)


def from_to(slide, *, from_word, to_word, from_sub=None, to_sub=None,
            y=2.70, h=1.55, tone="dark", x=None, w=None):
    """A FROM -> TO transition: a muted source panel, a mint arrow, and an
    emphasised destination panel. Use for "shift from X to Y" slides. The TO
    panel is the filled/emphasised one (the message), FROM is outlined/muted."""
    x = ED_X if x is None else x
    w = ED_W if w is None else w
    arrow_w = 1.10
    gap = 0.15
    body_w = w - arrow_w - 2 * gap
    from_w = body_w * 0.38
    to_w = body_w * 0.62
    # FROM — muted / outlined
    _rect(slide, x, y, from_w, h,
          fill=COLORS["green_2"] if tone == "dark" else COLORS["page_light"],
          line=COLORS["rule_dark"] if tone == "dark" else COLORS["rule_light"],
          line_w=1.0)
    _text(slide, "FROM", x=x + 0.30, y=y + 0.22, w=from_w - 0.60, h=0.24,
          font=FONTS["sans_sb"], size=10, color=COLORS["charcoal"], bold=True,
          uppercase=True, letter_spacing=2.4, anchor="top")
    serif_title(slide, from_word, x=x + 0.30, y=y + 0.52, w=from_w - 0.60,
                h=0.60, size=30, tone=tone,
                color=COLORS["mint_pale"] if tone == "dark" else COLORS["ink_2"])
    if from_sub:
        _text(slide, from_sub, x=x + 0.30, y=y + 1.12, w=from_w - 0.60, h=0.30,
              font=FONTS["sans_l"], size=11,
              color=COLORS["mint_pale"] if tone == "dark" else COLORS["ink_2"],
              anchor="top")
    # Arrow
    ax = x + from_w + gap
    arrow(slide, "right", x=ax, y=y + h / 2 - 0.28, w=arrow_w, h=0.56,
          color=COLORS["mint"])
    # TO — filled / emphasised
    tx = ax + arrow_w + gap
    _rect(slide, tx, y, to_w, h,
          fill=COLORS["green"] if tone == "light" else COLORS["green_3"])
    _hbar(slide, tx, y, to_w, color=COLORS["mint"], height_in=0.05)
    _text(slide, "TO", x=tx + 0.35, y=y + 0.24, w=to_w - 0.70, h=0.24,
          font=FONTS["sans_sb"], size=10, color=COLORS["gold"], bold=True,
          uppercase=True, letter_spacing=2.4, anchor="top")
    serif_title(slide, to_word, x=tx + 0.35, y=y + 0.54, w=to_w - 0.70,
                h=0.60, size=30, tone="dark", color=COLORS["white"])
    if to_sub:
        _text(slide, to_sub, x=tx + 0.35, y=y + 1.14, w=to_w - 0.70, h=0.30,
              font=FONTS["sans_l"], size=11, color=COLORS["mint_pale"],
              anchor="top")
    return y + h


def tier_ladder(slide, tiers, *, y, gap=0.30, tone="light", x=None, w=None):
    """Stacked priority tiers, top emphasised. `tiers` = list of dicts:
        {"label": "01 · PRIMARY FOCUS", "title": "Commercial execution",
         "note": "...", "items": [...], "emphasis": True, "height": 1.55}
    Use for "what's primary vs secondary" slides. Emphasised tiers get a
    filled surface + mint bar; non-emphasised get an outline + muted bar."""
    x = ED_X if x is None else x
    w = ED_W if w is None else w
    ty = y
    for t in tiers:
        h = t.get("height", 1.55)
        emph = t.get("emphasis", True)
        if emph:
            _rect(slide, x, ty, w, h,
                  fill=COLORS["off_white"] if tone == "light"
                  else COLORS["green_2"])
            barc = COLORS["mint"]
            labelc = COLORS["mint_dark"] if tone == "light" else COLORS["mint"]
            titlec = COLORS["green"] if tone == "light" else COLORS["white"]
        else:
            _rect(slide, x, ty, w, h,
                  fill=COLORS["page_light"] if tone == "light"
                  else COLORS["green"],
                  line=COLORS["rule_light"] if tone == "light"
                  else COLORS["rule_dark"], line_w=1.0)
            barc = COLORS["rule_light"] if tone == "light" else COLORS["rule_dark"]
            labelc = COLORS["charcoal"]
            titlec = COLORS["ink_2"] if tone == "light" else COLORS["mint_pale"]
        _vbar(slide, x, ty, h, color=barc, width_in=0.07)
        _text(slide, t["label"], x=x + 0.35, y=ty + 0.22, w=6.0, h=0.26,
              font=FONTS["sans_sb"], size=11, color=labelc, bold=True,
              uppercase=True, letter_spacing=2.0, anchor="top")
        serif_title(slide, t["title"], x=x + 0.35, y=ty + 0.52, w=w * 0.58,
                    h=0.70, size=26 if emph else 24, tone=tone, color=titlec)
        if t.get("items"):
            ix = x + 0.35
            for it in t["items"]:
                _text(slide, "–  " + it, x=ix, y=ty + 1.20, w=4.0, h=0.30,
                      font=FONTS["sans"], size=12,
                      color=COLORS["ink"] if tone == "light"
                      else COLORS["white"], anchor="top")
                ix += 4.0
        if t.get("note"):
            _text(slide, t["note"], x=x + w * 0.62, y=ty, w=w * 0.38 - 0.30,
                  h=h, font=FONTS["sans_l"], size=12,
                  color=COLORS["ink_2"] if tone == "light"
                  else COLORS["mint_pale"], anchor="middle", line_spacing=1.30)
        ty += h + gap
    return ty


def directional_ladder(slide, rows, *, y, tone="dark", x=None, w=None,
                       label_col=4.15, gap=0.26):
    """Up / sideways / down rows for "strengthened / refocused / deprioritised"
    style slides. `rows` = list of (direction, label, accent, items, subtag):
        ("up",    "Strengthened",   COLORS["mint"], [...],        None)
        ("right", "Refocused",      COLORS["gold"], [...],        None)
        ("down",  "Deprioritised",  COLORS["blue"], [...], "TEMPORARY")
    Each row is a band with a coloured bar, a directional arrow, a serif
    label, and the items rendered as chips. Rows fill the area from `y` down
    to ED_SAFE_BOT."""
    x = ED_X if x is None else x
    w = ED_W if w is None else w
    n = len(rows)
    rh = (ED_SAFE_BOT - y - gap * (n - 1)) / n
    for i, (direction, label, accent, items, subtag) in enumerate(rows):
        ry = y + i * (rh + gap)
        _rect(slide, x, ry, w, rh,
              fill=COLORS["green_2"] if tone == "dark" else COLORS["off_white"])
        _vbar(slide, x, ry, rh, color=accent, width_in=0.07)
        if direction == "right":
            arrow(slide, "right", x=x + 0.40, y=ry + rh / 2 - 0.18,
                  w=0.55, h=0.36, color=accent)
        elif direction in _ARROW_SHAPES:
            arrow(slide, direction, x=x + 0.48, y=ry + rh / 2 - 0.28,
                  w=0.40, h=0.56, color=accent)
        _text(slide, label, x=x + 1.15, y=ry, w=label_col - 1.30, h=rh,
              font=FONTS["serif"], size=24,
              color=COLORS["white"] if tone == "dark" else COLORS["green"],
              anchor="middle")
        if subtag:
            _text(slide, subtag, x=x + 1.18, y=ry + rh - 0.40,
                  w=label_col - 1.30, h=0.26, font=FONTS["sans_sb"], size=9,
                  color=accent, bold=True, uppercase=True, letter_spacing=2.2,
                  anchor="top")
        ix = x + label_col
        iw_total = w - label_col - 0.30
        m = len(items)
        cgap = 0.20
        cw = (iw_total - cgap * (m - 1)) / m
        chip_h = 0.50
        for j, it in enumerate(items):
            cxp = ix + j * (cw + cgap)
            chip(slide, it, x=cxp, y=ry + rh / 2 - chip_h / 2, w=cw, h=chip_h,
                 fill=COLORS["green_4"] if tone == "dark" else COLORS["page_light"],
                 text_color=COLORS["white"] if tone == "dark" else COLORS["ink"],
                 line=None if tone == "dark" else COLORS["rule_light"], size=11)


def intensity_bars(slide, tiers, *, x, y, w, tone="light"):
    """Decreasing-width bars that visualise tiering / prioritisation inside a
    panel (e.g. management-intensity categories). `tiers` = list of
    (label, sub, fill, frac) where frac (0-1) sets the bar width. Labels sit
    in a fixed right-hand column so they never collapse to negative width.
    Returns the y below the last bar."""
    text_col_w = 1.55
    bar_max = w - text_col_w - 0.20
    bar_h = 0.58
    bar_gap = 0.22
    ty = y
    for label, sub, fill, frac in tiers:
        bw = bar_max * frac
        _rect(slide, x, ty, bw, bar_h, fill=fill)
        tx = x + bar_max + 0.20
        _text(slide, label, x=tx, y=ty + 0.06, w=text_col_w, h=0.26,
              font=FONTS["sans_sb"], size=10.5,
              color=COLORS["green"] if tone == "light" else COLORS["white"],
              bold=True, uppercase=True, letter_spacing=1.0, anchor="top")
        if sub:
            _text(slide, sub, x=tx, y=ty + 0.33, w=text_col_w, h=0.26,
                  font=FONTS["sans_l"], size=9.5,
                  color=COLORS["ink_2"] if tone == "light"
                  else COLORS["mint_pale"], anchor="top")
        ty += bar_h + bar_gap
    return ty
