"""Story-led SCENE COMPOSER for CBRE decks - the default way to build.

THE STORY MAKES THE LAYOUT. You do not pick a recipe and pour content into it.
You declare what each slide must SAY (a deck of slides; each slide a `scene` of
ordered rows; each row split into cells; each cell one styled primitive) and this
composer lays the scene on the safe CBRE grid and sizes every cell's text UP to
fill its space. Slide count and shape flex with the story; two decks should look
different; two slides should rarely look the same.

It is the same engine the cbre-il-account-briefing skill uses, generalised for any
deck. It sits on top of `build.py` (the CBRE visual system): it draws the CBRE
chrome (eyebrow, fit-to-text serif headline, lead line, footer, wordmark) and then
composes the scene from the cells, reusing the `build` primitives and helpers as
the palette each cell draws from.

Slide kinds (the `kind` field):
  cover    -> the polished build.cover (giant serif title + optional themes strip)
  section  -> build.section_divider (giant numeral + title + optional lead/items)
  scene    -> a freely composed slide: chrome + an ordered list of rows/cells
  closing  -> build.thank_you (optional contact cards)

Cell kinds (inside a scene row's `cells`):
  prose  - a full paragraph (text; optional bold lead `label`). The explainer.
  stat   - a hero value + caption label. Several stat cells in a row = a KPI strip.
  list   - items [{title,text}], numbered or bulleted.
  table  - headers + rows; the last column is the analytical read.
  panel  - a dark accent side box; title + items [{label,value}] (or text).
  quote  - a pull quote + attribution.
  heading- a small section label. rule - a thin divider line.
  callout- the CBRE expert-note box (title + body, optional tag).
  chips  - a row of rounded pills (tags/countries/status).
  card   - one roman/decimal numbered card (a row of card cells = a card grid).
  image  - a picture from a path, fit within the cell.

THE RULES (same spirit as the il-account-briefing skill):
  1. Story-led, not recipe-led. Compose each slide from what it must say.
  2. Density from substance, never from tricks. Fill a slide with MORE real,
     relevant content if it makes sense, never with spacing or ballooned fonts.
     Leftover space is a signal to write more, not to stretch what is there.
  3. Explain, do not tabulate. Prose carries the narrative; tables are evidence.
  4. No lazy repetition. A scene layout should appear at most twice, ideally once.
  5. Readable, not a billboard. Text sizes UP to a readable cap, never tiny.
  6. No em/en dashes (swept before save). Box grows to fit text; font never shrinks.

Usage (Python):
  import compose
  compose.render(plan, "Deck.pptx")        # plan = a dict (below) or a .json path

Usage (CLI):
  python compose.py plan.json Deck.pptx [--no-resolve] [--no-label-bake]

Plan shape:
  {
    "deck_meta": {"eyebrow": "CBRE | ADVISORY"},   # optional defaults
    "slides": [
      {"kind": "cover", "title": "...", "subtitle": "...", "eyebrow": "...",
       "date": "...", "themes": ["...", "..."]},
      {"kind": "scene", "tone": "dark", "eyebrow": "01 | CONTEXT",
       "headline": "...", "lead": "...", "footer": "...",
       "scene": [
         {"weight": 1.3, "cells": [{"kind": "prose", "label": "THE SHIFT", "text": "..."}]},
         {"weight": 0.8, "cells": [{"kind": "stat", "value": "46%", "label": "..."},
                                   {"kind": "stat", "value": "17", "label": "..."}]}
       ]},
      {"kind": "section", "number": 1, "title": "...", "lead": "...", "items": ["..."]},
      {"kind": "closing", "title": "Thank you.", "contacts": [{"name": "...", "email": "..."}]}
    ]
  }
"""
from __future__ import annotations

import argparse
import json
import os
import shutil
import sys
import tempfile
from pathlib import Path

# Locate the cbre-corporate-pptx build library (this file lives next to it).
_HERE = Path(__file__).resolve().parent
for _p in (_HERE, Path.home() / ".claude/skills/cbre-corporate-pptx/scripts"):
    if (_p / "build.py").exists():
        sys.path.insert(0, str(_p.resolve()))
        break
import build  # noqa: E402
from build import COLORS as C, FONTS as Fz, ED_X, ED_W, ED_SAFE_BOT, SLIDE_W  # noqa: E402
from pptx.util import Inches  # noqa: E402

# Cell kinds this composer can draw (kept in sync with the CELL dict below).
CELL_KINDS = {"prose", "stat", "list", "table", "panel", "quote", "heading",
              "rule", "callout", "chips", "card", "image",
              # nesting
              "split",
              # editorial devices (build.py helpers, bounded to a cell rect)
              "from_to", "timeline", "tiers", "directions", "bars", "sightline"}

# Smallest rect each cell kind can be drawn in without cramming. Nesting makes
# it possible to ask for a cell too small for its content; we raise instead of
# squashing, because a squashed cell is the silent bug this composer exists to
# prevent.
MIN_CELL_W = 0.90
MIN_CELL_H = {
    "rule": 0.06, "sightline": 0.10, "heading": 0.26, "image": 0.40,
    "chips": 0.42, "prose": 0.50, "stat": 0.55, "quote": 0.60, "split": 0.60,
    "list": 0.70, "bars": 0.80, "card": 0.80, "table": 0.90, "panel": 0.90,
    "directions": 1.00, "callout": 1.00, "from_to": 1.50, "tiers": 1.50,
    "timeline": 2.40,
}
_MIN_CELL_H_DEFAULT = 0.50


class SceneCellTooSmall(ValueError):
    """A scene cell's rect is too small to draw its content honestly."""


def _assert_cell_room(kind, x, y, w, h, path=""):
    need_h = MIN_CELL_H.get(kind, _MIN_CELL_H_DEFAULT)
    where = f"cell '{kind}'" + (f" at {path}" if path else "")
    if w < MIN_CELL_W or h < need_h:
        raise SceneCellTooSmall(
            f"{where} got {w:.2f}x{h:.2f}in but needs at least "
            f"{MIN_CELL_W:.2f}x{need_h:.2f}in.\n"
            f"Four legitimate fixes: (1) give the row more weight, (2) drop a "
            f"cell from the row, (3) nest less deeply, (4) split the slide in "
            f"two. Do not shrink the minimum - that is the cramming bug wearing "
            f"a number."
        )

# ---------------------------------------------------------------------------
# Text hygiene + size estimators
# ---------------------------------------------------------------------------

def _clean(s):
    if not isinstance(s, str):
        return s
    return s.replace("—", "-").replace("–", "-")

def _ink(tone):
    return C["white"] if tone == "dark" else C["ink"]

def _muted(tone):
    return C["mint"] if tone == "dark" else C["ink_2"]

def _accent(tone):
    """The primary accent, which is tone-conditional.

    Confirmed against the corporate template: wheat/gold carries the accent on
    dark grounds, bright Accent Green carries it on white. Each is near
    invisible on the other ground, so this is a rule, not a preference. Use
    _muted() for the secondary accent on the same slide."""
    return C["gold"] if tone == "dark" else C["bright_green"]

def _para_h(text, w, size, ls=1.34):
    """Conservative (slightly tall) wrapped-paragraph height in inches."""
    cpl = max(16, int(w / (0.0080 * size)))
    n = len(str(text))
    return max(1, (n + cpl - 1) // cpl) * (size / 72.0) * ls

def _fit_fill(text, w, avail_h, lo=12.0, hi=20.0, ls=1.34):
    """Largest body size in [lo, hi] whose wrapped text fits avail_h. Picking the
    largest size that fits is what fills the cell and enlarges the text (density)."""
    size = hi
    while size > lo and _para_h(text, w, size, ls) > avail_h:
        size -= 0.5
    return max(size, lo)

# ---------------------------------------------------------------------------
# Slide chrome (the consistent CBRE frame; the diversity lives in the scene)
# ---------------------------------------------------------------------------

def _chrome(deck, plan, slide):
    """Background, footer/wordmark, and the eyebrow + serif headline + lead, drawn by
    build.editorial_header so the TITLE FITS ITS TEXT and returns the exact y where
    the body begins (no reserved-whitespace gap under the title).
    Return (slide, body_top, body_bot, tone)."""
    tone = slide.get("tone", "light")
    s = build.blank(deck, tone=tone)
    default_eyebrow = (plan.get("deck_meta", {}) or {}).get("eyebrow") or "CBRE"
    title = _clean(slide.get("headline") or "")
    lead = _clean(slide.get("lead")) if slide.get("lead") else None
    y = build.editorial_header(s, eyebrow_text=_clean(slide.get("eyebrow") or default_eyebrow),
                               title=title, tone=tone, intro=lead,
                               title_size=slide.get("headline_size", 30))
    foot = slide.get("footer")
    if foot:
        build._text(s, _clean(foot), x=ED_X, y=ED_SAFE_BOT + 0.02, w=SLIDE_W - ED_X - 0.55, h=0.26,
                    font=Fz["sans_l"], size=8.5,
                    color=C["ink_2"] if tone == "light" else C["mint"], anchor="top")
    body_top = y + 0.18
    body_bot = ED_SAFE_BOT - (0.34 if foot else 0.0)
    return s, body_top, body_bot, tone

# ---------------------------------------------------------------------------
# Cell primitives (each draws one styled element to fill its rect)
# ---------------------------------------------------------------------------

def c_prose(s, cell, x, y, w, h, tone):
    yy = y
    label = _clean(cell.get("label") or "")
    if label:
        build._text(s, label, x=x, y=yy, w=w, h=0.28, font=Fz["sans_sb"], size=12,
                    color=_accent(tone), bold=True, uppercase=True, letter_spacing=1.5, anchor="top")
        yy += 0.40
    text = _clean(cell.get("text") or "")
    avail = (y + h) - yy
    size = _fit_fill(text, w, avail, lo=12.5, hi=cell.get("max_size", 16.0))
    build.body(s, [text], x=x, y=yy, w=w, h=avail, size=size, color=_ink(tone),
               tone=tone, line_spacing=1.36)

def c_stat(s, cell, x, y, w, h, tone):
    value = _clean(cell.get("value") or "")
    label = _clean(cell.get("label") or "")
    # Scale the hero number to the cell, but cap by WIDTH so a long value never
    # wraps to a stray second line. scale="hero" is the one-blockbuster-number
    # slide: it lifts the cap so a lone stat can genuinely dominate.
    hero = cell.get("scale") == "hero"
    wcap = w * 130.0 / max(4, len(value))
    vsize = max(24.0, min(150.0 if hero else 60.0, h * 30.0, wcap))
    build.serif_title(s, value, x=x, y=y, w=w, h=vsize / 72 * 1.25, size=vsize,
                      tone=tone, line_spacing=1.0)
    ly = y + vsize / 72.0 * 1.18
    if label:
        build._text(s, label, x=x, y=ly, w=w, h=max(0.3, (y + h) - ly),
                    font=Fz["sans_l"], size=12, color=_muted(tone), anchor="top", line_spacing=1.22)

def c_list(s, cell, x, y, w, h, tone):
    items = [it for it in cell.get("items", []) if it]
    if not items:
        return
    numbered = cell.get("numbered", True)
    gap = 0.16
    ih = (h - gap * (len(items) - 1)) / len(items)
    for i, it in enumerate(items):
        iy = y + i * (ih + gap)
        title = _clean(it.get("title") or "")
        text = _clean(it.get("text") or "")
        if numbered:
            build._text(s, f"{i + 1:02d}", x=x, y=iy, w=0.85, h=0.5, font=Fz["serif"],
                        size=23, color=_accent(tone), anchor="top")
            tx, tw = x + 1.0, w - 1.0
        else:
            build._rect(s, x, iy + 0.06, 0.16, 0.16, fill=_accent(tone))
            tx, tw = x + 0.34, w - 0.34
        ty = iy
        if title:
            build._text(s, title, x=tx, y=ty, w=tw, h=0.30, font=Fz["sans_sb"], size=13.5,
                        color=_ink(tone), bold=True, anchor="top")
            ty += 0.36
        if text:
            avail = (iy + ih) - ty
            tsize = _fit_fill(text, tw, avail, lo=11.0, hi=14.5)
            build.body(s, [text], x=tx, y=ty, w=tw, h=avail, size=tsize,
                       color=_ink(tone), tone=tone, line_spacing=1.30)

def c_table(s, cell, x, y, w, h, tone):
    headers = [_clean(z) for z in cell.get("headers", [])]
    rows = [[_clean(z) for z in r] for r in cell.get("rows", [])]
    if not headers and not rows:
        return
    build.table(s, headers, rows, x=x, y=y, w=w, h=max(1.0, h), tone=tone,
                font_size=max(10.5, cell.get("font_size", 12)), col_aligns=cell.get("aligns"))

def c_panel(s, cell, x, y, w, h, tone):
    """A filled accent panel (dark side box). Items stack by MEASURED height and the
    value font shrinks to fit, so a wrapped value never collides with the next label."""
    build._rect(s, x, y, w, h, fill=C["green_3"])
    build._rect(s, x, y, w, 0.055, fill=C["gold"])
    pad = 0.28
    iw = w - 2 * pad
    yy = y + pad + 0.05
    title = _clean(cell.get("title") or "")
    if title:
        build._text(s, title, x=x + pad, y=yy, w=iw, h=0.34, font=Fz["sans_sb"], size=13,
                    color=C["gold"], bold=True, uppercase=True, letter_spacing=1.6, anchor="top")
        yy += 0.40 + (0.20 if len(title) > 28 else 0.0)
    inner_bot = y + h - pad
    items = [it for it in cell.get("items", []) if it]
    if items:
        gap, lab_h = 0.20, 0.30
        avail = inner_bot - yy
        def content_h(vs):
            return sum(lab_h + _para_h(_clean(it.get("value") or ""), iw, vs, 1.22) + gap for it in items)
        vsize = 17.5
        while vsize > 10.5 and content_h(vsize) > avail:
            vsize -= 0.5
        for it in items:
            build._text(s, _clean(it.get("label") or ""), x=x + pad, y=yy, w=iw, h=0.26,
                        font=Fz["sans_sb"], size=11, color=C["mint"], bold=True,
                        uppercase=True, letter_spacing=1.0, anchor="top")
            vy = yy + lab_h
            vh = _para_h(_clean(it.get("value") or ""), iw, vsize, 1.22)
            build._text(s, _clean(it.get("value") or ""), x=x + pad, y=vy, w=iw, h=vh + 0.05,
                        font=Fz["sans_l"], size=vsize, color=C["off_white"], anchor="top", line_spacing=1.25)
            yy += lab_h + vh + gap
    elif cell.get("text"):
        txt = _clean(cell["text"])
        avail = inner_bot - yy
        size = _fit_fill(txt, iw, avail, lo=11.5, hi=15.0)
        build.body(s, [txt], x=x + pad, y=yy, w=iw, h=avail, size=size,
                   color=C["off_white"], tone="dark", line_spacing=1.34)

def c_quote(s, cell, x, y, w, h, tone):
    text = _clean(cell.get("text") or "")
    attrib = _clean(cell.get("attrib") or "")
    qh = h - (0.34 if attrib else 0.0)
    size = _fit_fill(text, w, qh, lo=15.0, hi=30.0, ls=1.14)
    build.serif_title(s, "“" + text + "”", x=x, y=y, w=w, h=qh, size=size,
                      tone=tone, line_spacing=1.12)
    if attrib:
        build._text(s, attrib, x=x, y=y + qh, w=w, h=0.3, font=Fz["sans_sb"], size=11,
                    color=_accent(tone), uppercase=True, letter_spacing=1.5, anchor="top")

def c_heading(s, cell, x, y, w, h, tone):
    build._text(s, _clean(cell.get("text") or ""), x=x, y=y, w=w, h=h, font=Fz["sans_sb"],
                size=cell.get("size", 13), color=_accent(tone), bold=True,
                uppercase=cell.get("uppercase", True), letter_spacing=1.6, anchor="middle")

def c_rule(s, cell, x, y, w, h, tone):
    build._rect(s, x, y + h / 2, w, 0.014, fill=_accent(tone))

def c_callout(s, cell, x, y, w, h, tone):
    """The CBRE expert-note box. Sized to the cell; the body autofits and the
    callout primitive keeps its bg rect in sync via the resolve pass."""
    title = _clean(cell.get("title") or "CBRE VIEW")
    text = _clean(cell.get("text") or "")
    tag = _clean(cell.get("tag")) if cell.get("tag") else None
    build.callout(s, title=title, body_text=text, x=x, y=y, w=w, h=max(1.0, h),
                  tone=tone, tag=tag)

def c_chips(s, cell, x, y, w, h, tone):
    """A row of rounded pills that wrap within the cell width."""
    items = [it for it in cell.get("items", []) if it]
    if not items:
        return
    chip_h, gap = 0.40, 0.16
    cx, cy = x, y
    line = (C["rule_light"] if tone == "light" else None)
    for it in items:
        txt = _clean(it if isinstance(it, str) else (it.get("text") or ""))
        cw = min(w, 0.34 + len(txt) * 0.095)
        if cx + cw > x + w + 0.01 and cx > x:
            cx = x
            cy += chip_h + gap
        if cy + chip_h > y + h + 0.05:
            break
        build.chip(s, txt, x=cx, y=cy, w=cw, h=chip_h, tone=tone, line=line)
        cx += cw + gap

def c_card(s, cell, x, y, w, h, tone):
    """One numbered card filling the cell. style='roman' -> roman_card (bullet body);
    style='decimal' -> decimal_card (paragraph body). A row of card cells = a grid."""
    style = cell.get("style", "decimal")
    n = cell.get("n", 1)
    title = _clean(cell.get("title") or "")
    accent = cell.get("accent", "mint")
    if style == "roman":
        items = [_clean(it if isinstance(it, str) else (it.get("text") or ""))
                 for it in (cell.get("items") or []) if it]
        if not items and cell.get("text"):
            items = [_clean(cell["text"])]
        build.roman_card(s, int(n) if str(n).isdigit() else 1, x, y, w, max(0.8, h),
                         title=title, body_lines=items or [""], tone=tone, accent=accent,
                         subtitle=_clean(cell.get("subtitle")) if cell.get("subtitle") else None)
    else:
        build.decimal_card(s, n, x, y, w, max(0.8, h), title=title,
                           body_text=_clean(cell.get("text") or ""), tone=tone)

def c_image(s, cell, x, y, w, h, tone):
    """A picture fit within the cell rect, preserving aspect ratio, centred."""
    path = cell.get("path")
    if not (path and Path(path).exists()):
        build._rect(s, x, y, w, h, fill=C["green_2"] if tone == "dark" else C["off_white"])
        build._text(s, _clean(cell.get("alt") or "image"), x=x, y=y + h / 2 - 0.2, w=w, h=0.4,
                    font=Fz["sans_l"], size=11, color=_muted(tone), align="center", anchor="middle")
        return
    pic = s.shapes.add_picture(str(path), Inches(x), Inches(y))
    cell_w, cell_h = Inches(w), Inches(h)
    scale = min(cell_w / pic.width, cell_h / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int(Inches(x) + (cell_w - pic.width) / 2)
    pic.top = int(Inches(y) + (cell_h - pic.height) / 2)

# ---------------------------------------------------------------------------
# Editorial device cells - the build.py helpers, bounded to a cell rect
#
# These are where the visual distinctiveness lives. They used to be reachable
# only from a hand-built slide, which meant the default build path could not
# express most of the documented archetypes. Each wrapper's whole job is to
# turn a cell rect into the device's own geometry contract.
# ---------------------------------------------------------------------------

def c_from_to(s, cell, x, y, w, h, tone):
    """A shift from X to Y. Takes an explicit height, so it maps 1:1 to a cell.

    Accepts `from`/`to` or the longer `from_word`/`to_word`. Both words are
    required: a from_to with an empty side used to render as two blank panels
    and an arrow, which looks like a layout bug and is impossible to spot in a
    plan. It raises instead.
    """
    frm = _clean(cell.get("from") or cell.get("from_word") or "")
    to = _clean(cell.get("to") or cell.get("to_word") or "")
    if not frm or not to:
        raise ValueError(
            "from_to cell needs both words: set `from` and `to` (or "
            "`from_word`/`to_word`). Got from=%r to=%r. The optional "
            "sub-captions are `from_sub` and `to_sub`." % (frm, to))
    build.from_to(s, from_word=frm,
                  to_word=to,
                  from_sub=_clean(cell.get("from_sub")) if cell.get("from_sub") else None,
                  to_sub=_clean(cell.get("to_sub")) if cell.get("to_sub") else None,
                  x=x, y=y, w=w, h=min(h, 2.10), tone=tone)


def c_timeline(s, cell, x, y, w, h, tone):
    """Where we are in a sequence. The device draws ~1.22in above and ~1.22in
    below its baseline, so centre the baseline in the cell."""
    phases = []
    for i, p in enumerate(cell.get("phases", [])):
        if isinstance(p, dict):
            phases.append((p.get("n") or f"{i + 1:02d}", _clean(p.get("label") or ""),
                           _clean(p.get("text") or ""), bool(p.get("done"))))
        else:
            phases.append((f"{i + 1:02d}", _clean(str(p)), "", False))
    if not phases:
        return
    has_here = any(not p[3] for p in phases)
    # Baseline: centre the drawn extent, biased down when no "here" tag is shown.
    above = 1.22 if has_here else 0.92
    line_y = y + (h - (above + 1.22)) / 2.0 + above
    build.phase_timeline(s, phases, x=x, y=line_y, w=w, tone=tone,
                         here_label=_clean(cell.get("here_label", "WE ARE HERE")))


def c_tiers(s, cell, x, y, w, h, tone):
    """Primary vs secondary. Tier heights are distributed across the cell so the
    ladder always ends inside its rect."""
    tiers = [dict(t) for t in cell.get("tiers", []) if t]
    if not tiers:
        return
    gap = 0.24
    each = (h - gap * (len(tiers) - 1)) / len(tiers)
    if each < 1.20:
        raise SceneCellTooSmall(
            f"cell 'tiers' has {len(tiers)} tiers in {h:.2f}in, giving "
            f"{each:.2f}in each (need 1.20in). Drop a tier or give the row "
            f"more weight."
        )
    for t in tiers:
        t["height"] = each
        t.setdefault("emphasis", False)
        t["label"] = _clean(t.get("label") or "")
        t["title"] = _clean(t.get("title") or "")
        if t.get("note"):
            t["note"] = _clean(t["note"])
        if t.get("items"):
            t["items"] = [_clean(i) for i in t["items"]]
    if tiers and not any(t.get("emphasis") for t in tiers):
        tiers[0]["emphasis"] = True
    build.tier_ladder(s, tiers, x=x, y=y, w=w, gap=gap, tone=tone)


_DIR_ACCENT = {"up": "mint", "right": "gold", "down": "blue"}


def c_directions(s, cell, x, y, w, h, tone):
    """Strengthened / refocused / deprioritised. Bounded by h (build.py's
    directional_ladder now accepts a height instead of filling to the slide)."""
    rows = []
    for r in cell.get("rows", []):
        direction = (r.get("direction") or "right").lower()
        accent = r.get("accent") or _DIR_ACCENT.get(direction, "gold")
        colour = C[accent] if isinstance(accent, str) else accent
        rows.append((direction, _clean(r.get("label") or ""), colour,
                     [_clean(i) for i in (r.get("items") or [])] or [""],
                     _clean(r.get("subtag")) if r.get("subtag") else None))
    if not rows:
        return
    build.directional_ladder(s, rows, x=x, y=y, w=w, h=h, tone=tone)


def c_bars(s, cell, x, y, w, h, tone):
    """Categorisation by weight. Bar pitch is fixed at 0.80in, so cap the count
    to what the cell can hold rather than overrunning it."""
    raw = [t for t in cell.get("tiers", []) if t]
    if not raw:
        return
    max_n = max(1, int(h // 0.80))
    if len(raw) > max_n:
        raise SceneCellTooSmall(
            f"cell 'bars' has {len(raw)} tiers but the cell only fits "
            f"{max_n} (0.80in each in {h:.2f}in). Drop a tier or give the row "
            f"more weight."
        )
    fallback = ["mint", "mint_dark", "rule_light", "gold"]
    tiers = []
    for i, t in enumerate(raw):
        fill = t.get("fill") or fallback[min(i, len(fallback) - 1)]
        tiers.append((_clean(t.get("label") or ""), _clean(t.get("sub") or ""),
                      C[fill] if isinstance(fill, str) else fill,
                      float(t.get("frac", 1.0 - i * 0.26))))
    build.intensity_bars(s, tiers, x=x, y=y, w=w, tone=tone)


def c_sightline(s, cell, x, y, w, h, tone):
    """The signature CBRE rule device. Horizontal = breadth, vertical = depth."""
    orientation = cell.get("orientation", "horizontal")
    length = float(cell.get("length") or (w * 0.42 if orientation == "horizontal"
                                          else min(h, 1.60)))
    build.line_of_sight(s, orientation=orientation, x=x,
                        y=y + (h / 2 if orientation == "horizontal" else 0.0),
                        length=length, tone=tone)


CELL = {
    "prose": c_prose, "stat": c_stat, "list": c_list, "table": c_table,
    "panel": c_panel, "quote": c_quote, "heading": c_heading, "rule": c_rule,
    "callout": c_callout, "chips": c_chips, "card": c_card, "image": c_image,
    "from_to": c_from_to, "timeline": c_timeline, "tiers": c_tiers,
    "directions": c_directions, "bars": c_bars, "sightline": c_sightline,
}

# ---------------------------------------------------------------------------
# Scene layout (rows fill the body by weight; cells split the row by span)
# ---------------------------------------------------------------------------

def _render_scene(s, scene, x, y, w, h, tone, _path="scene", _depth=0):
    """Lay rows down the rect by weight, split each row across its cells by span.

    A cell may itself be a nested scene (kind='split'), in which case this
    recurses into the cell's rect. That matters: a partition of a partition is
    still a partition, so nesting buys asymmetric composition (left rails,
    L-shapes, unequal quadrants) without ever making overlap possible.
    """
    if _depth > 4:
        raise SceneCellTooSmall(
            f"{_path}: nesting deeper than 4 levels. Compose the slide from "
            f"fewer, larger regions, or split it in two."
        )
    rows = [r for r in scene if r.get("cells")]
    if not rows:
        return
    gap_v = 0.22 if _depth == 0 else 0.16
    total_wt = sum(float(r.get("weight", 1.0)) for r in rows) or 1.0
    avail_h = h - gap_v * (len(rows) - 1)
    cy = y
    for ri, r in enumerate(rows):
        rh = avail_h * (float(r.get("weight", 1.0)) / total_wt)
        cells = r["cells"]
        gap_h = 0.40 if _depth == 0 else 0.28
        total_span = sum(float(c.get("span", 1.0)) for c in cells) or 1.0
        avail_w = w - gap_h * (len(cells) - 1)
        cx = x
        for ci, c in enumerate(cells):
            cw = avail_w * (float(c.get("span", 1.0)) / total_span)
            kind = c.get("kind", "prose")
            path = f"{_path}.row{ri + 1}.cell{ci + 1}"
            _assert_cell_room(kind, cx, cy, cw, rh, path)
            if kind == "split":
                _render_scene(s, c.get("scene", []), cx, cy, cw, rh, tone,
                              _path=path, _depth=_depth + 1)
            else:
                CELL.get(kind, c_prose)(s, c, cx, cy, cw, rh, tone)
            cx += cw + gap_h
        cy += rh + gap_v


# ---------------------------------------------------------------------------
# Named skeletons - geometry generators, not content templates
#
# A skeleton takes a FLAT list of cells and expands it into the rows/cells
# structure above. It decides how the slide is *carved up*; it never decides
# what the slide says. Because a skeleton expands into an ordinary scene, every
# guard on the rendering path applies to it unchanged.
# ---------------------------------------------------------------------------

def _sk_bands(cells):
    """Full-width bands, one per cell. The plain horizontal stack."""
    return [{"weight": c.get("weight", 1.0), "cells": [c]} for c in cells]


def _sk_rail(cells):
    """A full-height left rail beside a stack of everything else. The classic
    asymmetric editorial split, and the shape the old model could not express."""
    if len(cells) < 2:
        return _sk_bands(cells)
    head, rest = cells[0], cells[1:]
    head = dict(head, span=head.get("span", 0.62))
    return [{"weight": 1.0, "cells": [
        head,
        {"kind": "split", "span": 1.0,
         "scene": [{"weight": c.get("weight", 1.0), "cells": [c]} for c in rest]},
    ]}]


def _sk_hero(cells):
    """One dominant cell over a strip of supporting cells."""
    if len(cells) < 2:
        return _sk_bands(cells)
    return [
        {"weight": 1.75, "cells": [cells[0]]},
        {"weight": 1.0, "cells": list(cells[1:])},
    ]


def _sk_mosaic(cells):
    """Two per row - an even grid for parallel evidence."""
    rows = []
    for i in range(0, len(cells), 2):
        rows.append({"weight": 1.0, "cells": cells[i:i + 2]})
    return rows


def _sk_ledger(cells):
    """A narrow read on the left, the evidence wide on the right. Extra cells
    stack beneath the evidence, not beneath the read."""
    if len(cells) < 2:
        return _sk_bands(cells)
    left, right = cells[0], cells[1:]
    left = dict(left, span=0.52)
    if len(right) == 1:
        return [{"weight": 1.0, "cells": [left, dict(right[0], span=1.0)]}]
    return [{"weight": 1.0, "cells": [
        left,
        {"kind": "split", "span": 1.0,
         "scene": [{"weight": c.get("weight", 1.0), "cells": [c]} for c in right]},
    ]}]


def _sk_poster(cells):
    """One cell taking most of the slide, the rest as a quiet footer strip."""
    if len(cells) < 2:
        return _sk_bands(cells)
    return [
        {"weight": 2.6, "cells": [cells[0]]},
        {"weight": 1.0, "cells": list(cells[1:])},
    ]


SKELETONS = {
    "bands": _sk_bands, "rail": _sk_rail, "hero": _sk_hero,
    "mosaic": _sk_mosaic, "ledger": _sk_ledger, "poster": _sk_poster,
}


def _expand_shape(slide):
    """Return the scene for a slide, expanding `shape` + flat `cells` if used.

    A slide may declare either an explicit `scene` (rows/cells) or a `shape`
    plus a flat `cells` list. The two are interchangeable; `shape` is shorthand.
    """
    if slide.get("scene"):
        return slide["scene"]
    shape = slide.get("shape")
    cells = slide.get("cells")
    if shape and cells:
        fn = SKELETONS.get(shape)
        if fn is None:
            raise ValueError(
                f"unknown scene shape {shape!r}. Available: "
                f"{', '.join(sorted(SKELETONS))}."
            )
        return fn([c for c in cells if c])
    return []

def r_scene(deck, plan, slide):
    scene = _expand_shape(slide)
    if not any(r.get("cells") for r in scene):
        # Empty scene degrades to a clean placeholder callout, never a bare header.
        s = build.blank(deck, tone=slide.get("tone", "light"))
        if slide.get("eyebrow"):
            build.eyebrow(s, _clean(slide["eyebrow"]), tone=slide.get("tone", "light"))
        note = _clean(slide.get("placeholder") or "Content to be added.")
        build.callout(s, title="PLACEHOLDER", body_text=note, x=ED_X, y=1.2, w=ED_W,
                      h=max(1.05, build.predict_callout_h(note, w=ED_W)),
                      tone=slide.get("tone", "light"))
        return s
    s, top, bot, tone = _chrome(deck, plan, slide)
    _render_scene(s, scene, ED_X, top, ED_W, max(1.0, bot - top), tone)
    return s

# ---------------------------------------------------------------------------
# Dedicated chrome slides (delegate to the polished build recipes)
# ---------------------------------------------------------------------------

def r_cover(deck, plan, slide):
    build.cover(deck, title=_clean(slide.get("title", "")),
                subtitle=_clean(slide.get("subtitle")) if slide.get("subtitle") else None,
                presenter=slide.get("presenter"), org=slide.get("org"), date=slide.get("date"),
                tone=slide.get("tone", "dark"),
                eyebrow_text=_clean(slide.get("eyebrow")) if slide.get("eyebrow") else None,
                themes=[_clean(t) for t in slide.get("themes", [])] or None)

def r_section(deck, plan, slide):
    build.section_divider(deck, number=slide.get("number", 1), title=_clean(slide.get("title", "")),
                          eyebrow_text=_clean(slide.get("eyebrow")) if slide.get("eyebrow") else None,
                          tone=slide.get("tone", "dark"),
                          lead=_clean(slide.get("lead")) if slide.get("lead") else None,
                          items=[_clean(i) for i in slide.get("items", [])] or None)

def r_closing(deck, plan, slide):
    build.thank_you(deck, title=_clean(slide.get("title", "Thank you.")),
                    subtitle=_clean(slide.get("subtitle")) if slide.get("subtitle") else None,
                    contacts=slide.get("contacts"), tone=slide.get("tone", "dark"))

KIND = {"cover": r_cover, "section": r_section, "divider": r_section,
        "closing": r_closing, "thank_you": r_closing, "scene": r_scene}

# ---------------------------------------------------------------------------
# Dash sweep + render
# ---------------------------------------------------------------------------

def _sweep_dashes(deck):
    n = 0
    def walk(shapes):
        nonlocal n
        for sh in shapes:
            if sh.shape_type == 6:
                walk(sh.shapes)
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if "—" in r.text or "–" in r.text:
                            r.text = _clean(r.text); n += 1
            if sh.has_table:
                for row in sh.table.rows:
                    for cl in row.cells:
                        for p in cl.text_frame.paragraphs:
                            for r in p.runs:
                                if "—" in r.text or "–" in r.text:
                                    r.text = _clean(r.text); n += 1
    for slide in deck.slides:
        walk(slide.shapes)
    return n

def _scene_signature(slide):
    """A slide's skeleton, as comparable data: its shape name plus the cell
    kinds in each row. Two slides with the same signature look the same,
    whatever words are in them."""
    rows = _expand_shape(slide)
    if not rows:
        return None
    parts = []
    for r in rows:
        kinds = []
        for c in r.get("cells", []):
            k = c.get("kind", "prose")
            if k == "split":
                inner = [ic.get("kind", "prose")
                         for ir in c.get("scene", [])
                         for ic in ir.get("cells", [])]
                k = "split(" + "/".join(inner) + ")"
            kinds.append(k)
        parts.append(",".join(kinds))
    return (slide.get("shape") or "custom") + " | " + " / ".join(parts)


def _uses_skeleton(sl):
    """True when a scene took a named shape off the shelf rather than
    composing its own rows and cells."""
    return bool(sl.get("shape")) and not sl.get("scene")


def audit_scene_shapes(plan, *, verbose=True, strict=False,
                       discipline=True):
    """Fail the deck when slides repeat a skeleton, or lean on the shelf.

    Variety used to be enforced only by prose in the documentation while tone
    mix was enforced by code, and tone mix is the one that held. This closes
    that gap, and goes a step further.

    The named skeletons in SKELETONS exist so a genuinely conventional slide
    does not have to be rebuilt from scratch. They are NOT a menu to choose
    from. The default is to compose the scene the argument wants, out of rows
    and cells; reaching for a skeleton is the exception and has to earn it.
    Five checks:

      1. No two consecutive scenes share a skeleton.
      2. No skeleton is used more than twice.
      3. Named skeletons stay a minority of the deck (about one in three).
      4. Every named skeleton carries a `shape_why` saying why it beats a
         bespoke scene for that specific point, and no two slides may give
         the same reason, because a reason that fits two slides fits neither.
      5. Enough distinct skeletons overall (>= 70% of scenes).

    `discipline=False` keeps the two structural checks (1 and 2) and drops the
    three editorial ones (3, 4, 5). It exists for the coverage harness, which
    must exercise every skeleton in one file and so cannot obey rationing. It
    is deliberately a function argument and not a plan field, so a deck being
    authored cannot switch it off: render() always audits with discipline on.

    Deliberate parallelism is not repetition: two slides walking two comparable
    routes should look alike so the reader can compare them. Declare it with
    `parallel_to: <earlier slide number>` and the shared skeleton is reported
    as intentional rather than warned. Pairs only, so it cannot blanket-exempt
    a deck.
    """
    slides = plan.get("slides", [])
    scenes = []
    for i, sl in enumerate(slides, start=1):
        if sl.get("kind", "scene") != "scene":
            continue
        sig = _scene_signature(sl)
        if sig:
            scenes.append((i, sig, sl))

    warnings, notes = [], []
    n = len(scenes)
    by_no = {i: (sig, sl) for i, sig, sl in scenes}

    # --- declared parallelism ---------------------------------------------
    parallel, targeted = {}, {}
    for i, sig, sl in scenes:
        tgt = sl.get("parallel_to")
        if tgt is None:
            continue
        if tgt not in by_no:
            warnings.append(
                f"slide {i} declares parallel_to={tgt!r}, which is not a "
                f"scene in this deck.")
        elif tgt == i:
            warnings.append(f"slide {i} declares parallel_to itself.")
        elif by_no[tgt][0] != sig:
            warnings.append(
                f"slide {i} declares parallel_to={tgt} but the two do not "
                f"share a skeleton, so the parallel is invisible to the "
                f"reader. Either compose them alike or drop the claim.")
        elif tgt in targeted:
            warnings.append(
                f"slides {targeted[tgt]} and {i} both declare "
                f"parallel_to={tgt}. Parallelism is for pairs; a chain of "
                f"look-alike slides is the repetition this guards against.")
        else:
            targeted[tgt] = i
            parallel[i] = tgt
            notes.append(f"slides {tgt} and {i} are deliberately parallel.")

    def _paired(a, b):
        return parallel.get(a) == b or parallel.get(b) == a

    # --- 1. consecutive repeats -------------------------------------------
    for (ia, sa, _), (ib, sb, _) in zip(scenes, scenes[1:]):
        if sa == sb and not _paired(ia, ib):
            warnings.append(
                f"slides {ia} and {ib} are consecutive and share a skeleton "
                f"({sa}). Recompose one from its own point, or if they really "
                f"are two halves of one comparison declare parallel_to={ia} "
                f"on slide {ib}.")

    # --- 2. overuse ---------------------------------------------------------
    counts = {}
    for i, sig, _ in scenes:
        counts.setdefault(sig, []).append(i)
    for sig, idxs in counts.items():
        eff = [i for i in idxs
               if not any(_paired(i, j) for j in idxs if j != i)]
        if len(eff) > 2:
            warnings.append(
                f"skeleton used {len(eff)}x on slides "
                f"{', '.join(str(i) for i in eff)} ({sig}). Compose the third "
                f"from its own argument rather than reusing a shape that "
                f"already carried two other points.")

    # --- 3. skeletons stay a minority ---------------------------------------
    shelf = [i for i, _, sl in scenes if _uses_skeleton(sl)]
    if discipline and n >= 3:
        cap = max(1, n // 3)
        if len(shelf) > cap:
            warnings.append(
                f"{len(shelf)} of {n} scenes took a named shape off the shelf "
                f"(slides {', '.join(str(i) for i in shelf)}); at most {cap} "
                f"should. Named skeletons are for genuinely conventional "
                f"slides, not the default path. Compose the rest from rows and "
                f"cells that fit what each slide actually argues.")

    # --- 4. a named shape justifies itself ----------------------------------
    reasons = {}
    for i, _, sl in scenes:
        if not discipline or not _uses_skeleton(sl):
            continue
        why = (sl.get("shape_why") or "").strip()
        if not why:
            warnings.append(
                f"slide {i} uses shape={sl['shape']!r} with no `shape_why`. "
                f"Say why that shape beats a scene composed for this point; "
                f"if there is no reason, compose the scene.")
        elif len(why) < 25:
            warnings.append(
                f"slide {i} `shape_why` is too thin to be a real reason: "
                f"{why!r}.")
        else:
            key = " ".join(why.lower().split())
            if key in reasons:
                warnings.append(
                    f"slides {reasons[key]} and {i} give the same reason for "
                    f"their shape. A reason that fits two slides justified "
                    f"neither - compose one of them.")
            else:
                reasons[key] = i

    # --- 5. novelty floor ----------------------------------------------------
    distinct = len(counts)
    if discipline and n >= 4:
        floor = -(-7 * n // 10)
        allowed = floor - len(parallel)
        if distinct < allowed:
            warnings.append(
                f"only {distinct} distinct skeletons across {n} scenes "
                f"(expected at least {max(1, allowed)}). The deck is settling "
                f"into a house layout. Rebuild the thinnest slides from their "
                f"own point.")

    result = {"scenes": n, "distinct": distinct,
              "signatures": [(i, sig) for i, sig, _ in scenes],
              "skeletons": shelf, "parallel": parallel,
              "notes": notes, "warnings": warnings, "ok": not warnings}

    if verbose:
        print(f"[shape audit] {n} scenes - {distinct} distinct skeleton(s), "
              f"{len(shelf)} off the shelf"
              + ("" if discipline else "  [coverage mode: editorial checks off]"))
        for note in notes:
            print(f"  [note] {note}")
        for w in warnings:
            print(f"  [warn] {w}")
        if not warnings and n:
            print("  [ok] Every scene earns its own shape.")

    if strict and warnings:
        raise AssertionError(
            "audit_scene_shapes found repeated or unjustified layouts:\n  "
            + "\n  ".join(warnings))
    return result


def render(plan, out_path, *, resolve=None, label_and_bake=True, audit=True,
           shapes_strict=False, geometry_strict=False):
    """Compose and render a story-led scene deck. `plan` is a dict (see module
    docstring) or a path to a .json file. Saves via build.save, so on Windows it
    runs the resolve pass, inherits the org sensitivity label and bakes
    fit-to-text (disable with label_and_bake=False)."""
    if isinstance(plan, (str, Path)):
        plan = json.loads(Path(plan).read_text(encoding="utf-8"))
    if audit:
        audit_scene_shapes(plan, verbose=True, strict=shapes_strict)
    deck = build.new_deck()
    for slide in plan.get("slides", []):
        KIND.get(slide.get("kind", "scene"), r_scene)(deck, plan, slide)
    _sweep_dashes(deck)
    # Stage the build in a private temp dir. Not `tempdir/<name>`: when the
    # caller's own out_path already lives in temp, that collides with itself
    # and the copy fails with a sharing violation.
    _stage = Path(tempfile.mkdtemp(prefix="compose_"))
    tmp = _stage / Path(out_path).name
    # build.save() inherits the org sensitivity label and bakes fit-to-text by
    # default (when the resolve pass runs); label_and_bake=False disables both.
    if label_and_bake:
        build.save(deck, str(tmp), resolve=resolve, audit=audit,
                   geometry_strict=geometry_strict)
    else:
        build.save(deck, str(tmp), resolve=resolve, label_from=False, bake=False,
                   audit=audit, geometry_strict=geometry_strict)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(tmp, out_path)
    shutil.rmtree(_stage, ignore_errors=True)
    return out_path

def main():
    p = argparse.ArgumentParser(description="Render a story-led scene plan into a CBRE deck.")
    p.add_argument("plan", help="path to the plan .json")
    p.add_argument("out", help="output .pptx path")
    p.add_argument("--no-resolve", action="store_true", help="skip the render-and-measure pass")
    p.add_argument("--no-label-bake", action="store_true",
                   help="skip the sensitivity-label inherit + autofit bake (no PowerPoint)")
    a = p.parse_args()
    out = render(a.plan, a.out, resolve=(False if a.no_resolve else None),
                 label_and_bake=not a.no_label_bake)
    print(f"deck -> {out}")

if __name__ == "__main__":
    main()
