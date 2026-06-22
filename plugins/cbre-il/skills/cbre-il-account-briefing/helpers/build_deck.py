"""Stage 5 deck builder: render a FROZEN content plan into a narrative-led CBRE
I&L *account brief* where THE STORY MAKES THE LAYOUT.

Scene composer (v8). There are no archetypes and no fixed templates. Each slide
is a `scene`: an ordered list of rows; each row splits into cells; each cell
holds one styled primitive (prose, stat, list, table, panel, quote, heading,
rule). The content plan, authored from the company's story, decides per slide
how many rows, how they split, what each cell holds and what gets emphasis. The
renderer is faithful and thin: it draws the CBRE chrome (eyebrow, serif
headline, lead, footer, wordmark) and then lays the scene out on a safe grid,
sizing every cell's text UP to fill its space. That auto-fill is what makes the
deck dense and the text large; the grid only prevents overlaps and off-canvas.
Three slide kinds keep a dedicated renderer for their fixed chrome: `cover`,
`divider`, `references`. Everything else is a scene.

THE HARD RULES:
  1. Render ONLY strings that exist in content_plan.json. Never compose a new
     claim here. Stage 6 reconciliation fails any new claim.
  2. Obey cbre-corporate-pptx: no font-shrink below the floor, mix dark/light,
     build to %TEMP% then copy to deliverables/ last.
  3. DENSITY: every scene fills its body region; cell text is sized UP to fill.
     No slide left >~40% empty. An empty scene degrades to a clean intel-gap
     callout, never a bare header.
  4. No em/en dashes: the builder sweeps every run before saving.

Every successful build writes two side files next to the .pptx: `build_report.json`
(per-slide shape + reference map) and `deck_text.txt` (the rendered text, one line
per paragraph/table-cell) which the post-build text gates (G5/G6a/G6b) and the
Final Gate consume.

Run `--dry-run` to report each slide's scene shape without rendering.
Run `--dump-text <pptx> <out.txt>` to (re)extract deck text from an existing deck.
"""
from __future__ import annotations

import argparse
import json
import os
import shutil
import sys
import tempfile
from pathlib import Path

_EMU_PER_IN = 914400.0

# Locate the cbre-corporate-pptx build library without a hardcoded user path.
# Works as a plugin (sibling skill under the same plugin root), as a standalone
# ~/.claude/skills install, or via the CLAUDE_PLUGIN_ROOT / CBRE_PPTX_SCRIPTS env vars.
_SIBLING = Path(__file__).resolve().parents[2] / "cbre-corporate-pptx" / "scripts"
_CANDIDATES = [Path("scripts"), _SIBLING, Path.home() / ".claude/skills/cbre-corporate-pptx/scripts"]
if os.environ.get("CLAUDE_PLUGIN_ROOT"):
    _CANDIDATES.insert(0, Path(os.environ["CLAUDE_PLUGIN_ROOT"]) / "skills" / "cbre-corporate-pptx" / "scripts")
if os.environ.get("CBRE_PPTX_SCRIPTS"):
    _CANDIDATES.insert(0, Path(os.environ["CBRE_PPTX_SCRIPTS"]))
for _p in _CANDIDATES:
    if _p.exists():
        sys.path.insert(0, str(_p.resolve()))
        break
else:
    raise SystemExit("cbre-corporate-pptx scripts not found; set CBRE_PPTX_SCRIPTS "
                     "or install the sibling cbre-corporate-pptx skill.")
import build  # noqa: E402

# Fail loudly if the sibling library drifted (a renamed/removed symbol we rely on).
_REQUIRED = ["new_deck", "save", "blank", "eyebrow", "serif_title", "editorial_header",
             "body", "table", "callout", "section_divider", "predict_callout_h", "_text",
             "_rect", "COLORS", "FONTS", "ED_X", "ED_W", "ED_SAFE_BOT", "SLIDE_W",
             "measure_text", "_resolve_available"]
_missing = [s for s in _REQUIRED if not hasattr(build, s)]
if _missing:
    raise SystemExit(f"cbre-corporate-pptx build library is missing required symbols {_missing}; "
                     "version drift? Update build_deck.py or the sibling skill.")
from build import COLORS as C, FONTS as Fz, ED_X, ED_W, ED_SAFE_BOT, SLIDE_W  # noqa: E402

# Cell kinds the renderer can draw (used by gate_runner self-check against the schema).
CELL_KINDS = {"prose", "stat", "list", "table", "panel", "quote", "heading", "rule"}

# ---------------------------------------------------------------------------
# Text hygiene + size estimators
# ---------------------------------------------------------------------------

def _clean(s):
    if not isinstance(s, str):
        return s
    return s.replace("—", "-").replace("–", "-")

def _ink(tone):
    return C["white"] if tone == "dark" else C["ink"]

def _muted(tone):
    return C["mint"] if tone == "dark" else C["ink_2"]

def _accent(tone):
    return C["mint"] if tone == "dark" else C["mint_dark"]

def _meas(text, w, size, font_key="sans_l", ls=1.34, uppercase=False, letter_spacing=None):
    """Calibrated rendered-height estimate (inches) via the corporate per-font width
    model (build.measure_text), which is empirically tuned against PowerPoint with the
    CBRE fonts. This replaces the old single-constant character heuristic so the no-COM
    build path (Linux / Cowork, where the two-pass bake cannot run) wraps text the way
    PowerPoint actually does. Pass the cell's REAL FONTS key (e.g. 'sans_sb'); the model
    is per-font, so a wrong font under- or over-predicts the wrap."""
    if not text:
        return 0.0
    return build.measure_text(str(text), size=size, w=w, font=font_key, line_spacing=ls,
                              uppercase=uppercase, letter_spacing=letter_spacing)

def _para_h(text, w, size, ls=1.34):
    """Conservative height in inches for a wrapped body paragraph. Calibrated via
    build.measure_text (body font); callers that render a different font call _meas
    with that font directly."""
    return _meas(text, w, size, "sans_l", ls)

def _fit_fill(text, w, avail_h, lo=12.0, hi=20.0, ls=1.34):
    """Largest body size in [lo, hi] whose wrapped text fits avail_h. Picking the
    largest size that fits is what fills the cell and enlarges the text (density)."""
    size = hi
    while size > lo and _para_h(text, w, size, ls) > avail_h:
        size -= 0.5
    return max(size, lo)

def _fit_table_font(headers, rows, w, lo=10.5, hi=12.0):
    """Pick a table body size scaled to the column width and column count, so a
    dense multi-column table reads at a proportionate size instead of a fixed 12pt
    that crowds. Width/density principle, not a per-slide recipe: the largest size
    in [lo, hi] whose longest cell still wraps to <=2 lines, then a small step down
    for 3+/4+ column tables (which read better a touch smaller). The plan can still
    override via cell.font_size."""
    ncols = max(1, len(headers) or (len(rows[0]) if rows else 1))
    colw = max(0.5, w / ncols - 0.24)
    longest = 0
    for r in rows:
        for c in r:
            longest = max(longest, len(str(c)))
    size = hi
    while size > lo:
        cpl = max(8, int(colw / (0.0080 * size)))
        lines = (longest + cpl - 1) // cpl if cpl else 99
        if lines <= 2:
            break
        size -= 0.5
    if ncols >= 4:
        size = min(size, 10.5)
    elif ncols >= 3:
        size = min(size, 11.0)
    return max(lo, size)

# ---------------------------------------------------------------------------
# Slide chrome (the consistent CBRE frame; the diversity lives in the scene)
# ---------------------------------------------------------------------------

CHROME_TITLE_Y = 1.02

# True only when the COM two-pass bake actually corrects layout (Windows + PowerPoint).
# On the no-COM path the lead sits on the *predicted* title bottom, so we add a little
# more safety pad below the headline to avoid a headline/lead collision when a title
# lands at the wrap boundary. build_deck() sets this from build._resolve_available().
_BAKED = True

def _chrome(deck, slide, lay=None):
    """Eyebrow + serif headline + lead + footer, drawn here (not via editorial_header)
    so the LEAD sits at the title's TRUE rendered bottom. editorial_header places the
    lead at the *predicted* title height, which over-reserves a second line when a
    title lands right at the wrap boundary, opening a dead gap under the title.
    Returns (slide, region_top, footer_top, tone); region_top = the lead's bottom."""
    lay = lay if lay is not None else {}
    tone = slide.get("tone", "light")
    s = build.blank(deck, tone=tone)
    title = _clean(slide.get("headline") or "")
    lead = _clean(slide.get("lead")) if slide.get("lead") else None
    tsize = slide.get("headline_size", 30)
    ink = C["white"] if tone == "dark" else C["green"]
    build.eyebrow(s, _clean(slide.get("eyebrow") or "CBRE I&L ACCOUNT BRIEF"),
                  tone=tone, x=ED_X, y=getattr(build, "ED_EYEBROW_Y", 0.55), accent="gold")
    th_pred = build.measure_text(title, size=tsize, w=ED_W, font="serif", line_spacing=1.05)
    build.serif_title(s, title, x=ED_X, y=CHROME_TITLE_Y, w=ED_W, h=th_pred,
                      size=tsize, tone=tone, color=ink)
    title_bot = lay.get("true_title_bot", CHROME_TITLE_Y + th_pred)
    region_top = title_bot
    if lead:
        # On the baked path the lead sits on the title's TRUE bottom (0.14 is enough).
        # On the no-COM path title_bot is the PREDICTED bottom, so add a little more
        # clearance to absorb a title that renders a touch taller than predicted.
        ly = title_bot + (0.14 if _BAKED else 0.28)
        ih = build.measure_text(lead, size=12.5, w=ED_W, font="sans_l", line_spacing=1.30)
        build.body(s, [lead], x=ED_X, y=ly, w=ED_W, h=ih, size=12.5, tone=tone,
                   color=C["mint_pale"] if tone == "dark" else C["ink_2"], line_spacing=1.30)
        region_top = ly + ih
    foot = slide.get("footer")
    src_y = ED_SAFE_BOT + 0.02
    if foot:
        build._text(s, _clean(foot), x=ED_X, y=src_y, w=SLIDE_W - ED_X - 0.55, h=0.26,
                    font=Fz["sans_l"], size=8.5,
                    color=C["ink_2"] if tone == "light" else C["mint"], anchor="top")
    footer_top = src_y if foot else ED_SAFE_BOT
    return s, region_top, footer_top, tone

# ---------------------------------------------------------------------------
# Cell primitives (each draws one styled element to fill its rect)
# ---------------------------------------------------------------------------

# Body prose renders at a FIXED size so the rendered height is deterministic (the
# layout engine measures it from the bake; a fill-to-cap size would make pass 1 and
# pass 2 disagree). A section label is glued tight above its paragraph.
PROSE_SIZE = 13.0
PROSE_LS = 1.32
LABEL_GAP = 0.34

def c_prose(s, cell, x, y, w, h, tone):
    yy = y
    label = _clean(cell.get("label") or "")
    if label:
        build._text(s, label, x=x, y=yy, w=w, h=0.24, font=Fz["sans_sb"], size=11.5,
                    color=_accent(tone), bold=True, uppercase=True, letter_spacing=1.6, anchor="top")
        yy += LABEL_GAP
    text = _clean(cell.get("text") or "")
    avail = max(0.2, (y + h) - yy)
    # FIXED, deterministic size (set by the engine; default PROSE_SIZE). The font does
    # NOT react to the box height - that circular dependency is what broke the
    # measure-then-place two-pass and opened uneven gaps. Shrinking for an overstuffed
    # slide is decided once by the engine and passed in via _psize.
    size = cell.get("_psize", cell.get("prose_size", PROSE_SIZE))
    build.body(s, [text], x=x, y=yy, w=w, h=avail, size=size,
               color=_ink(tone), tone=tone, line_spacing=PROSE_LS)

def c_stat(s, cell, x, y, w, h, tone):
    value = _clean(cell.get("value") or "")
    label = _clean(cell.get("label") or "")
    # Size the hero number by WIDTH only - NOT by the box height. Sizing by height was
    # the same circular dependency that broke prose: the rendered height then depended
    # on the allotted height, so the two-pass never converged. Width-only is
    # deterministic, so pass 1 and pass 2 render identically.
    wcap = w * 130.0 / max(4, len(value))
    vsize = max(24.0, min(54.0, wcap))
    build.serif_title(s, value, x=x, y=y, w=w, h=vsize / 72 * 1.25, size=vsize,
                      tone=tone, line_spacing=1.0)
    ly = y + vsize / 72.0 * 1.18
    if label:
        build._text(s, label, x=x, y=ly, w=w, h=max(0.3, (y + h) - ly),
                    font=Fz["sans_l"], size=12, color=_muted(tone), anchor="top", line_spacing=1.22)

LIST_SIZE = 11.0
LIST_LS = 1.26

LIST_TITLE_SIZE = 13.5
NUM_W = 1.05

def _list_geom(items, w, cols):
    """Per-item geometry on the cols-wide grid. The detail sits below the title by the
    title's MEASURED height (titles wrap to 2 lines in a narrow column, so a fixed
    offset would overprint). Columns auto-reduce if each would be too narrow."""
    gutter = 0.6
    while cols > 1 and (w - gutter * (cols - 1)) / cols < 2.6:
        cols -= 1
    cellw = (w - gutter * (cols - 1)) / cols
    tw = cellw - NUM_W
    cells = []
    for it in items:
        title = _clean(it.get("title") or "")
        detail = _clean(it.get("text") or "")
        # The title renders in bold sans_sb, so measure it with that font: a generic
        # body-font estimate is narrower and would under-reserve the title height,
        # dropping the detail on top of a wrapped title.
        th = _meas(title, tw, LIST_TITLE_SIZE, "sans_sb", 1.12) if title else 0.0
        toff = (th + 0.10) if title else 0.0
        dh = _meas(detail, tw, LIST_SIZE, "sans_l", LIST_LS)
        cells.append({"title": title, "detail": detail, "toff": toff,
                      "h": max(0.52, toff + dh)})
    grid = [cells[i:i + cols] for i in range(0, len(cells), cols)]
    row_h = [max(c["h"] for c in gr) for gr in grid]
    return cellw, tw, gutter, cols, grid, row_h, 0.26

def c_list(s, cell, x, y, w, h, tone):
    items = [it for it in cell.get("items", []) if it]
    if not items:
        return
    cols = cell.get("cols") or (2 if len(items) >= 3 else 1)
    cellw, tw, gutter, cols, grid, row_h, rule_gap = _list_geom(items, w, cols)
    cy = y
    idx = 0
    for ri, gr in enumerate(grid):
        if ri > 0:
            build._line(s, x, cy - rule_gap * 0.5, x + w, cy - rule_gap * 0.5,
                        color=C["rule_light"] if tone == "light" else C["rule_dark"], width_pt=0.75)
        for ci, c in enumerate(gr):
            cx = x + ci * (cellw + gutter)
            build._text(s, f"{idx + 1:02d}", x=cx, y=cy - 0.04, w=0.90, h=0.55,
                        font=Fz["serif_l"], size=26, color=_accent(tone), anchor="top")
            if c["title"]:
                build._text(s, c["title"], x=cx + NUM_W, y=cy, w=tw, h=max(0.24, c["toff"] - 0.08),
                            font=Fz["sans_sb"], size=LIST_TITLE_SIZE, color=_ink(tone),
                            bold=True, anchor="top", line_spacing=1.12)
            build._text(s, c["detail"], x=cx + NUM_W, y=cy + c["toff"], w=tw,
                        h=max(0.16, row_h[ri] - c["toff"]), font=Fz["sans_l"], size=LIST_SIZE,
                        color=_muted(tone), anchor="top", line_spacing=LIST_LS)
            idx += 1
        cy += row_h[ri] + (rule_gap if ri < len(grid) - 1 else 0)

TABLE_HEADER_H = 0.42
TABLE_PAD_X = 0.14
TABLE_ROW_PAD = 0.20

def _table_col_fracs(headers, rows, w, size, override=None):
    """Balance column widths to content, but FIRST guarantee each column enough width
    to fit its header on one line (so 'FY22/23' never wraps), then share the leftover
    width proportionally to each column's longest cell."""
    n = len(headers) or (len(rows[0]) if rows else 1)
    if override and len(override) == n:
        t = sum(override)
        return [f / t for f in override]
    colmax = [len(headers[i]) if i < len(headers) else 1 for i in range(n)]
    for r in rows:
        for i in range(min(n, len(r))):
            colmax[i] = max(colmax[i], len(str(r[i])))
    hdr_char = size * 0.46 * 1.05 / 72.0          # sans_sb, uppercase
    min_w = [len(headers[i]) * hdr_char + 2 * TABLE_PAD_X + 0.06 if i < len(headers)
             else 0.6 for i in range(n)]
    leftover = w - sum(min_w)
    if leftover <= 0:
        widths = min_w
    else:
        tot = sum(colmax) or 1
        widths = [min_w[i] + leftover * (colmax[i] / tot) for i in range(n)]
    t = sum(widths)
    return [x / t for x in widths]

def _table_geom(headers, rows, w, size):
    fr = _table_col_fracs(headers, rows, w, size)
    colw = [w * f for f in fr]
    row_h = []
    for r in rows:
        hmax = 0.0
        for ci, cell in enumerate(r):
            hmax = max(hmax, _para_h(str(cell), colw[ci] - 2 * TABLE_PAD_X, size, 1.20))
        row_h.append(max(0.48, hmax + TABLE_ROW_PAD))
    return colw, row_h

def c_table(s, cell, x, y, w, h, tone):
    headers = [_clean(z) for z in cell.get("headers", [])]
    rows = [[_clean(z) for z in r] for r in cell.get("rows", [])]
    if not headers and not rows:
        return
    size = cell.get("font_size") or _fit_table_font(headers, rows, w)
    size = max(10.5, size)
    colw, row_h = _table_geom(headers, rows, w, size)
    body_text = _ink(tone)
    zebra = C["off_white"] if tone == "light" else None   # dark tables read off the rules
    rule = C["rule_light"] if tone == "light" else C["rule_dark"]
    # header band
    build._rect(s, x, y, w, TABLE_HEADER_H, fill=C["mint"])
    cx = x
    for ci, htext in enumerate(headers):
        build._text(s, htext, x=cx + TABLE_PAD_X, y=y, w=colw[ci] - 2 * TABLE_PAD_X, h=TABLE_HEADER_H,
                    font=Fz["sans_sb"], size=size, color=C["green"], bold=True,
                    uppercase=True, letter_spacing=0.8, anchor="middle")
        cx += colw[ci]
    ry = y + TABLE_HEADER_H
    for ri, r in enumerate(rows):
        rh = row_h[ri]
        if ri % 2 == 1 and zebra is not None:
            build._rect(s, x, ry, w, rh, fill=zebra)
        cx = x
        for ci, cval in enumerate(r):
            first = (ci == 0)
            build._text(s, cval, x=cx + TABLE_PAD_X, y=ry, w=colw[ci] - 2 * TABLE_PAD_X, h=rh,
                        font=Fz["sans_sb"] if first else Fz["sans_l"], size=size,
                        color=_accent(tone) if first else body_text, bold=first,
                        anchor="middle", line_spacing=1.20)
            cx += colw[ci]
        build._line(s, x, ry, x + w, ry, color=rule, width_pt=0.75)
        ry += rh
    build._line(s, x, ry, x + w, ry, color=rule, width_pt=0.75)

PANEL_PAD = 0.28
PANEL_VALUE_SIZE = 12.5
PANEL_LABEL_GAP = 0.27
PANEL_ITEM_GAP = 0.16

def _panel_title_h(title, iw):
    """Vertical advance after a panel title (its measured height + a gap). Measured
    with the title's real font (bold uppercase sans_sb) and FLOORED at the historical
    0.40 so a one-line title advances exactly as before (no layout churn), while a long
    title that wraps gets the room it needs instead of overprinting the first item.
    Shared by _panel_inner_natural (estimate) and c_panel (draw) so the two stay in
    lock-step and the Windows two-pass converges."""
    if not title:
        return 0.0
    th = _meas(title, iw, 13, "sans_sb", 1.12, uppercase=True, letter_spacing=1.6)
    return max(0.40, th + 0.16)

def _panel_inner_natural(cell, w):
    """Height the panel CONTENT (title + items / text) occupies, packed from the top
    at fixed sizes. Deterministic, so it matches what c_panel draws in both passes."""
    pad = PANEL_PAD
    iw = w - 2 * pad
    title = _clean(cell.get("title") or "")
    tot = pad + 0.05 + _panel_title_h(title, iw)
    items = [it for it in cell.get("items", []) if it]
    if items:
        for it in items:
            if _clean(it.get("label") or ""):
                tot += PANEL_LABEL_GAP
            tot += _para_h(_clean(it.get("value") or ""), iw, PANEL_VALUE_SIZE, 1.18) + PANEL_ITEM_GAP
        tot -= PANEL_ITEM_GAP
    elif cell.get("text"):
        tot += _para_h(_clean(cell["text"]), iw, 13.5, 1.34)
    return tot + pad

def c_panel(s, cell, x, y, w, h, tone):
    """A filled accent panel. Items pack from the TOP at a FIXED value size (no
    height-driven shrink), so the panel renders identically in pass 1 and pass 2 and
    the two-pass spacing converges."""
    build._rect(s, x, y, w, h, fill=C["green_3"])
    build._rect(s, x, y, w, 0.055, fill=C["gold"])
    pad = PANEL_PAD
    iw = w - 2 * pad
    yy = y + pad + 0.05
    title = _clean(cell.get("title") or "")
    if title:
        build._text(s, title, x=x + pad, y=yy, w=iw, h=0.34, font=Fz["sans_sb"], size=13,
                    color=C["gold"], bold=True, uppercase=True, letter_spacing=1.6, anchor="top")
        yy += _panel_title_h(title, iw)
    items = [it for it in cell.get("items", []) if it]
    if items:
        for it in items:
            label = _clean(it.get("label") or "")
            value = _clean(it.get("value") or "")
            if label:
                build._text(s, label, x=x + pad, y=yy, w=iw, h=0.24, font=Fz["sans_sb"], size=11,
                            color=C["mint"], bold=True, uppercase=True, letter_spacing=1.0, anchor="top")
                yy += PANEL_LABEL_GAP
            vh = _para_h(value, iw, PANEL_VALUE_SIZE, 1.18)
            build._text(s, value, x=x + pad, y=yy, w=iw, h=vh, font=Fz["sans_l"],
                        size=PANEL_VALUE_SIZE, color=C["off_white"], anchor="top", line_spacing=1.18)
            yy += vh + PANEL_ITEM_GAP
    elif cell.get("text"):
        build.body(s, [_clean(cell["text"])], x=x + pad, y=yy, w=iw, h=(y + h - pad) - yy,
                   size=13.5, color=C["off_white"], tone="dark", line_spacing=1.34)

def c_quote(s, cell, x, y, w, h, tone):
    text = _clean(cell.get("text") or "")
    attrib = _clean(cell.get("attrib") or "")
    qh = h - (0.34 if attrib else 0.0)
    size = _fit_fill(text, w, qh, lo=15.0, hi=30.0, ls=1.14)
    build.serif_title(s, "“" + text + "”", x=x, y=y, w=w, h=qh, size=size,
                      tone=tone, line_spacing=1.12)
    if attrib:
        build._text(s, attrib, x=x, y=y + qh, w=w, h=0.3, font=Fz["sans_sb"], size=11,
                    color=_accent(tone), uppercase=True, letter_spacing=1.5, anchor="top")

def c_heading(s, cell, x, y, w, h, tone):
    build._text(s, _clean(cell.get("text") or ""), x=x, y=y, w=w, h=h, font=Fz["sans_sb"],
                size=cell.get("size", 13), color=_accent(tone), bold=True,
                uppercase=cell.get("uppercase", True), letter_spacing=1.6, anchor="middle")

def c_rule(s, cell, x, y, w, h, tone):
    build._rect(s, x, y + h / 2, w, 0.014, fill=_accent(tone))

CELL = {
    "prose": c_prose, "stat": c_stat, "list": c_list, "table": c_table,
    "panel": c_panel, "quote": c_quote, "heading": c_heading, "rule": c_rule,
}

# ---------------------------------------------------------------------------
# Scene layout: MEASURE each row, then JUSTIFY the rows down the body region with
# EQUAL spacing (top gap = inter-row gaps = gap to the footer). Heights come from
# the bake (true rendered geometry) when available, never a stretch-by-weight.
# A row is a rigid unit; cells split it horizontally by span.
# ---------------------------------------------------------------------------

GAP_H = 0.40        # horizontal gutter between side-by-side cells
MAX_EVEN = 0.85     # cap on the even gap; above this a slide is genuinely underfilled

def _cell_natural_h(cell, w, tone):
    k = cell.get("kind", "prose")
    if k == "prose":
        lab = LABEL_GAP if _clean(cell.get("label") or "") else 0.0
        psize = cell.get("_psize", cell.get("prose_size", PROSE_SIZE))
        return lab + _para_h(_clean(cell.get("text") or ""), w, psize, PROSE_LS) + 0.04
    if k == "list":
        items = [it for it in cell.get("items", []) if it]
        if not items:
            return 0.3
        cols = cell.get("cols") or (2 if len(items) >= 3 else 1)
        _cw, _tw, _g, _cols, grid, row_h, rule_gap = _list_geom(items, w, cols)
        return sum(row_h) + rule_gap * (len(grid) - 1)
    if k == "table":
        headers = [_clean(z) for z in cell.get("headers", [])]
        rows = [[_clean(z) for z in r] for r in cell.get("rows", [])]
        if not headers and not rows:
            return 0.3
        size = max(10.5, cell.get("font_size") or _fit_table_font(headers, rows, w))
        _colw, row_h = _table_geom(headers, rows, w, size)
        return TABLE_HEADER_H + sum(row_h)
    if k == "stat":
        value = _clean(cell.get("value") or "")
        wcap = w * 130.0 / max(4, len(value))
        vsize = max(24.0, min(54.0, wcap))
        lh = _para_h(_clean(cell.get("label") or ""), w, 12, 1.22) if cell.get("label") else 0.0
        return vsize / 72.0 * 1.25 + 0.06 + lh
    if k == "panel":
        return _panel_inner_natural(cell, w)
    if k == "quote":
        return _para_h(_clean(cell.get("text") or ""), w, 22, 1.14) + (0.34 if cell.get("attrib") else 0.0)
    if k == "heading":
        return cell.get("size", 13) / 72.0 * 1.3 + 0.12
    if k == "rule":
        return 0.06
    return 0.6

def _row_natural_h(row, w, tone):
    cells = row.get("cells", [])
    if not cells:
        return 0.0
    total_span = sum(float(c.get("span", 1.0)) for c in cells) or 1.0
    avail_w = w - GAP_H * (len(cells) - 1)
    return max(_cell_natural_h(c, avail_w * (float(c.get("span", 1.0)) / total_span), tone)
               for c in cells)

_SHRINK_KINDS = {"prose", "quote", "heading"}
MIN_GAP = 0.14

def _shrinkable(row):
    """A row whose cells can give up height by shrinking their font (text-only)."""
    return all(c.get("kind", "prose") in _SHRINK_KINDS for c in row.get("cells", []))

def _render_scene(s, scene, x, region_top, w, footer_top, tone, lay=None):
    rows = [r for r in scene if r.get("cells")]
    if not rows:
        return
    lay = lay if lay is not None else {}
    rt = region_top
    n = len(rows)
    region = footer_top - rt

    # STEP 1 - decide ONE deterministic prose size for this slide. Shrink it (in 0.5pt
    # steps, floored at 10.5) only until the natural content fits the region. This runs
    # identically in pass 1 and pass 2 (it depends only on the plan, never on the box),
    # so the paragraph renders at the SAME size both passes -> the bake's true height is
    # the real height -> spacing converges. _psize is what c_prose renders at.
    prose_cells = [c for r in rows for c in r["cells"] if c.get("kind", "prose") == "prose"]
    psize = PROSE_SIZE
    while True:
        for c in prose_cells:
            c["_psize"] = psize
        naturals = [max(0.2, _row_natural_h(r, w, tone)) for r in rows]
        if not prose_cells or psize <= 10.5 or sum(naturals) + MIN_GAP * (n + 1) <= region:
            break
        psize -= 0.5
    lay["prose_size"] = psize

    # STEP 2 - heights: the bake's TRUE heights when we have them (pass 2), else the
    # estimate. Because the prose size is now fixed, true height is a real property.
    true_h = lay.get("true_h")
    use_h = true_h if (true_h and len(true_h) == n) else naturals

    # STEP 3 - justify the rows with EQUAL spacing. If even the floored prose still
    # overflows, pack tight (MIN_GAP) and flag overstuffed; otherwise space evenly.
    slack = region - sum(use_h)
    if slack >= 0:
        even = min(slack / (n + 1), MAX_EVEN)
        leftover = slack - even * (n + 1)
        lay["overstuffed"] = False
        lay["underfilled"] = even >= MAX_EVEN - 1e-6
    else:
        even = MIN_GAP
        leftover = 0.0
        lay["overstuffed"] = True
        lay["underfilled"] = False

    if os.environ.get("BUILD_DEBUG"):
        print(f"[scene] rt={rt:.2f} region={region:.2f} psize={psize} "
              f"naturals={[round(x,2) for x in naturals]} use_h={[round(x,2) for x in use_h]} "
              f"even={even:.3f} overstuffed={lay.get('overstuffed')}", file=sys.stderr)
    row_starts = []
    y = rt + even + leftover / 2.0
    for i, r in enumerate(rows):
        row_starts.append(y)
        cells = r["cells"]
        total_span = sum(float(c.get("span", 1.0)) for c in cells) or 1.0
        avail_w = w - GAP_H * (len(cells) - 1)
        cx = x
        for c in cells:
            cw = avail_w * (float(c.get("span", 1.0)) / total_span)
            CELL.get(c.get("kind", "prose"), c_prose)(s, c, cx, y, cw, use_h[i], tone)
            cx += cw + GAP_H
        y = y + use_h[i] + even
    lay["row_starts"] = row_starts
    lay["footer_top"] = footer_top

def r_scene(deck, plan, slide, ledger):
    scene = slide.get("scene", [])
    if not any(r.get("cells") for r in scene):
        s = build.blank(deck, tone=slide.get("tone", "light"))
        y = 1.2
        if slide.get("eyebrow"):
            build.eyebrow(s, _clean(slide["eyebrow"]), tone=slide.get("tone", "light"))
        note = "Not yet captured; a priority intel gap for the first meeting."
        build.callout(s, title="INTEL GAP", body_text=note, x=ED_X, y=y, w=ED_W,
                      h=max(1.05, build.predict_callout_h(note, w=ED_W)),
                      tone=slide.get("tone", "light"), tag="FIRST MEETING")
        return s
    lay = slide.setdefault("_lay", {})
    s, region_top, footer_top, tone = _chrome(deck, slide, lay)
    _render_scene(s, scene, ED_X, region_top, ED_W, footer_top, tone, lay)
    return s

def _read_true_layout(plan, pptx_path):
    """Read the baked deck: per scene slide, the chrome's true bottom (region top)
    and each row's true rendered height, keyed back onto slide['_lay'] for pass 2."""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    for i, slide in enumerate(plan.get("slides", [])):
        lay = slide.get("_lay")
        if slide.get("kind", "scene") != "scene" or not lay or not lay.get("row_starts"):
            continue
        if i >= len(prs.slides):
            continue
        starts = lay["row_starts"]
        ftop = lay["footer_top"]
        tops = []
        for sh in prs.slides[i].shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                tops.append((sh.top / _EMU_PER_IN, (sh.top + sh.height) / _EMU_PER_IN))
        # The headline's true bottom (the shape drawn at CHROME_TITLE_Y): pass 2 sits
        # the lead right under it, so a 1-line title that the estimator thought was
        # 2-line no longer leaves a dead gap.
        titles = [b for (t, b) in tops if abs(t - CHROME_TITLE_Y) < 0.12]
        if titles:
            lay["true_title_bot"] = max(titles)
        bounds = starts + [ftop]
        true_h = []
        for j in range(len(starts)):
            lo, hi = bounds[j] - 0.08, min(bounds[j + 1] - 0.05, ftop - 0.05)
            # Measure the row's REAL bottom (uncapped) so an overflowing paragraph is
            # detectable as overstuffed in pass 2; footer chrome is excluded by top.
            bots = [b for (t, b) in tops if lo <= t < hi and t < ftop - 0.05]
            true_h.append((max(bots) - starts[j]) if bots else 0.0)
        lay["true_h"] = true_h
        if os.environ.get("BUILD_DEBUG"):
            print(f"[read] slide#{i} starts={[round(x,2) for x in starts]} "
                  f"true_h={[round(x,2) for x in true_h]} title_bot={lay.get('true_title_bot')}",
                  file=sys.stderr)

# ---------------------------------------------------------------------------
# Dedicated chrome slides (cover / divider / references)
# ---------------------------------------------------------------------------

def r_cover(deck, plan, slide, ledger):
    m = plan.get("deck_meta", {})
    tone = "dark"
    s = build.blank(deck, tone=tone)
    build.eyebrow(s, _clean(slide.get("eyebrow") or "CBRE | INDUSTRIAL & LOGISTICS"),
                  tone=tone, accent="gold")
    if m.get("date"):
        build._text(s, _clean(m["date"]), x=SLIDE_W - 3.20, y=0.55, w=2.65, h=0.30,
                    font=Fz["mono"], size=10, color=C["gold"], uppercase=True,
                    letter_spacing=2.0, align="right", anchor="middle")
    title = _clean(m.get("target", "Company"))
    tl = len(title)
    tsize = 110 if tl <= 22 else 88 if tl <= 36 else 70 if tl <= 54 else 56
    band_w, band_h = SLIDE_W - 1.10, 2.6
    # Guard: the char-count ladder above can still wrap a long title past the title
    # band and into the subtitle. Step the size down until the measured title fits the
    # band, so the cover never overprints. Short titles already fit, so this is a no-op
    # for them (unchanged on both the COM and no-COM paths).
    while tsize > 40 and build.measure_text(title, size=tsize, w=band_w, font="serif",
                                            line_spacing=1.02) > band_h:
        tsize -= 6
    build.serif_title(s, title, x=0.55, y=1.5, w=band_w, h=band_h, size=tsize,
                      tone=tone, line_spacing=1.02)
    subtitle = _clean(slide.get("subtitle") or m.get("subtitle") or "Industrial & Logistics Account Brief")
    build.body(s, [subtitle], x=0.55, y=4.55, w=SLIDE_W - 3.0, h=0.7, size=18,
               color=C["mint"], tone=tone, line_spacing=1.3)
    contents = slide.get("contents") or []
    if contents:
        build._rect(s, 0.55, 5.55, 1.65, 0.02, fill=C["gold"])
        build._text(s, "WHAT'S INSIDE", x=0.55, y=5.66, w=4.0, h=0.28, font=Fz["sans_sb"],
                    size=10, color=C["gold"], bold=True, uppercase=True, letter_spacing=2.2, anchor="top")
        n = min(len(contents), 4)
        cw = (SLIDE_W - 1.10) / n
        for i, item in enumerate(contents[:n]):
            cx = 0.55 + i * cw
            build._text(s, f"{i + 1:02d}", x=cx, y=6.02, w=0.6, h=0.4, font=Fz["serif"],
                        size=20, color=C["mint"], anchor="top")
            build._text(s, _clean(item), x=cx + 0.55, y=6.06, w=cw - 0.7, h=0.9, font=Fz["sans_l"],
                        size=12.5, color=C["white"], anchor="top", line_spacing=1.2)
    else:
        build._rect(s, 0.55, 6.45, 2.50, 0.018, fill=C["gold"])
        build.body(s, [_clean(m.get("prepared_by", "CBRE Industrial & Logistics")), "CBRE"],
                   x=0.55, y=6.58, w=6.0, h=0.55, size=12, color=C["white"], tone=tone, line_spacing=1.3)
    return s

def r_divider(deck, plan, slide, ledger):
    build.section_divider(deck, number=slide.get("number", 1),
                          title=_clean(slide.get("title", "")),
                          eyebrow_text=_clean(slide.get("eyebrow") or "CBRE I&L ACCOUNT BRIEF"),
                          tone=slide.get("tone", "dark"),
                          lead=_clean(slide.get("lead")) if slide.get("lead") else None,
                          items=[_clean(i) for i in slide.get("items", [])] or None)

def r_references(deck, plan, slide, ledger):
    tone = "light"
    s = build.blank(deck, tone=tone)
    build.eyebrow(s, "REFERENCES", accent="mint")
    build.serif_title(s, "Sources", x=ED_X, y=0.95, w=8.0, h=1.0, size=34, tone=tone)
    from ledger import reference_rows  # the canonical, de-duplicated, numbered list
    lines = [f"[{rr['n']}] {rr['host']}  (tier {rr['tier']}, {rr['pub_date']})"
             for rr in reference_rows(ledger)]
    if not lines:
        lines = ["Reference list generated from the Source Ledger at delivery."]
    n = len(lines)
    ncol = 4 if n > 60 else (3 if n > 36 else (2 if n > 16 else 1))
    per = (n + ncol - 1) // ncol
    step = ED_W / ncol
    rsize = 8 if ncol >= 4 else (8.5 if ncol == 3 else 9)
    for ci in range(ncol):
        chunk = lines[ci * per:(ci + 1) * per]
        if chunk:
            build.body(s, chunk, x=ED_X + ci * step, y=1.95, w=step - 0.25, h=4.55, size=rsize, tone=tone)
    return s

KIND = {"cover": r_cover, "divider": r_divider, "references": r_references, "scene": r_scene}

# ---------------------------------------------------------------------------
# Dash sweep + per-slide report
# ---------------------------------------------------------------------------

def _sweep_dashes(deck):
    n = 0
    def walk(shapes):
        nonlocal n
        for sh in shapes:
            if sh.shape_type == 6:
                walk(sh.shapes)
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if "—" in r.text or "–" in r.text:
                            r.text = _clean(r.text); n += 1
            if sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            for r in p.runs:
                                if "—" in r.text or "–" in r.text:
                                    r.text = _clean(r.text); n += 1
    for slide in deck.slides:
        walk(slide.shapes)
    return n

def _deck_text_lines(slides):
    """Every rendered paragraph and table cell of text, one per line, in slide order
    (slides separated by a blank line). This is the `deck_text.txt` that the
    post-build TEXT gates consume: reconciliation (G5), the dash sweep (G6a / qa4),
    the editorial reviewer (G6b) and the Final Gate. Without it those gates have no
    input; emitting it on every build is what makes Stage 6 runnable. Walks shapes
    exactly like _sweep_dashes (groups, text frames, tables) so it sees the same
    text the deck renders. No decorative slide-number markers are injected (they
    would read as new numbers to reconciliation); slides are split by a blank line."""
    lines = []

    def emit_tf(tf):
        for p in tf.paragraphs:
            t = "".join(r.text for r in p.runs).strip()
            if t:
                lines.append(t)

    def walk(shapes):
        for sh in shapes:
            if sh.shape_type == 6:  # group
                walk(sh.shapes)
            if sh.has_text_frame:
                emit_tf(sh.text_frame)
            if sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        emit_tf(cell.text_frame)

    for slide in slides:
        walk(slide.shapes)
        lines.append("")  # blank separator between slides (harmless to the gates)
    return lines


def _write_deck_text(slides, out_txt):
    Path(out_txt).write_text("\n".join(_deck_text_lines(slides)) + "\n", encoding="utf-8")
    return out_txt


def _scene_cells(slide):
    return [c for r in slide.get("scene", []) for c in r.get("cells", [])]

def _slide_report(slide):
    kind = slide.get("kind", "scene")
    if kind == "scene":
        rows = slide.get("scene", [])
        cells = _scene_cells(slide)
        kinds = sorted({c.get("kind", "prose") for c in cells})
        bad = sorted({k for k in kinds if k not in CELL_KINDS})
        return {"slide_no": slide.get("slide_no"), "kind": kind, "known": not bad,
                "rows": len(rows), "cells": len(cells), "cell_kinds": kinds,
                "unknown_cell_kinds": bad, "sparse": len(cells) < 1}
    return {"slide_no": slide.get("slide_no"), "kind": kind, "known": kind in KIND,
            "rows": 0, "cells": 0, "cell_kinds": [], "unknown_cell_kinds": [], "sparse": False}

def _reference_map(ledger_path):
    # Same canonical numbering as the rendered references slide (fix: report and
    # slide must agree on [N], so a citation cross-check is never off by the dedupe).
    from ledger import reference_rows
    return [{"n": rr["n"], "claim_id": rr["claim_id"], "domain": rr["host"]}
            for rr in reference_rows(ledger_path)]

def _g7_index(plan):
    """The authoritative slide-identity table for the G7 visual reviewer: one row per
    PHYSICAL slide, in deck order. `deck_position` (1-based) is the slide's identity and
    equals the integer in its PNG filename (`NN.png` from to_png.ps1). The reviewer must
    cite findings by `deck_position`, never infer identity from the filename digits or
    from the on-slide eyebrow (a human chapter counter that can differ from both
    `deck_position` and `slide_no`, and may be non-numeric). `thesis` is the slide's lead,
    else its headline/title, so the reviewer knows what each slide should say."""
    out = []
    for i, s in enumerate(plan.get("slides", [])):
        thesis = s.get("lead") or s.get("headline") or s.get("title") or s.get("subtitle") or ""
        out.append({
            "png": f"{i + 1:02d}.png",
            "deck_position": i + 1,
            "slide_no": s.get("slide_no"),
            "kind": s.get("kind", "scene"),
            "eyebrow": _clean(str(s.get("eyebrow") or "")),
            "thesis": _clean(str(thesis))[:160],
        })
    return out

# ---------------------------------------------------------------------------

def build_deck(plan_path, out_path, ledger_path, resolve, dry_run=False, label_and_bake=True):
    plan = json.loads(Path(plan_path).read_text(encoding="utf-8"))
    report = [_slide_report(s) for s in plan.get("slides", [])]
    if dry_run:
        unknown = [r["slide_no"] for r in report if not r["known"]]
        sparse = [r["slide_no"] for r in report if r["sparse"]]
        print(json.dumps({"slides": report, "unknown": unknown, "sparse_slides": sparse}, indent=2))
        if unknown:
            print(f"DRY-RUN FAIL: unknown kind/cell on slides {unknown}", file=sys.stderr)
            return None
        print("DRY-RUN OK" + (f" (warning: sparse slides {sparse})" if sparse else ""))
        return None
    # Capability tier for text-fit: the COM two-pass bake only corrects layout on
    # Windows + PowerPoint; everywhere else the calibrated estimate is the final word
    # (so we add the no-bake safety pads). Surface which tier ran for the build log/G7.
    global _BAKED
    com_ok = bool(build._resolve_available())
    _BAKED = bool(label_and_bake and com_ok)
    print(f"[fit] tier={'com-twopass' if _BAKED else 'calibrated-estimate'} "
          f"(label_and_bake={label_and_bake}, com_available={com_ok})", file=sys.stderr)
    def _build_once():
        deck = build.new_deck()
        for slide in plan.get("slides", []):
            KIND.get(slide.get("kind", "scene"), r_scene)(deck, plan, slide, ledger_path)
        return deck

    tmp = Path(tempfile.gettempdir()) / Path(out_path).name
    if label_and_bake:
        # PASS 1: lay out with natural-height estimates, bake (PowerPoint computes the
        # real rendered heights), then read those true heights back per scene slide.
        deck1 = _build_once()
        tmp1 = Path(tempfile.gettempdir()) / ("_pass1_" + Path(out_path).name)
        build.save(deck1, str(tmp1), resolve=resolve)
        try:
            _read_true_layout(plan, str(tmp1))
        finally:
            try:
                os.remove(tmp1)
            except OSError:
                pass
        # PASS 2: re-lay out, spacing the rows EXACTLY evenly off the true heights.
        deck = _build_once()
        _sweep_dashes(deck)
        deck_text_lines = _deck_text_lines(deck.slides)
        build.save(deck, str(tmp), resolve=resolve)
    else:
        # No COM (smoke / CI): single pass on the natural-height estimates.
        deck = _build_once()
        _sweep_dashes(deck)
        deck_text_lines = _deck_text_lines(deck.slides)
        build.save(deck, str(tmp), resolve=resolve, label_from=False, bake=False)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(tmp, out_path)
    try:
        os.remove(tmp)
    except OSError:
        pass
    Path(out_path).with_name("build_report.json").write_text(
        json.dumps({"slides": report, "references": _reference_map(ledger_path),
                    "g7_index": _g7_index(plan)}, indent=2),
        encoding="utf-8")
    # Emit the rendered deck text beside the report so Stage 6 reconciliation, the
    # dash sweep, the editorial reviewer and the Final Gate always have an input.
    Path(out_path).with_name("deck_text.txt").write_text(
        "\n".join(deck_text_lines) + "\n", encoding="utf-8")
    return out_path

def main():
    p = argparse.ArgumentParser(description="Render the frozen content plan into the narrative-led CBRE I&L brief.")
    p.add_argument("plan", help="content_plan.json (or, with --dump-text, the .pptx to read)")
    p.add_argument("out", help="output .pptx (or, with --dump-text, the output .txt)")
    p.add_argument("--ledger", default=None)
    p.add_argument("--no-resolve", action="store_true")
    p.add_argument("--dry-run", action="store_true")
    p.add_argument("--dump-text", action="store_true",
                   help="re-extract deck text from an existing .pptx: plan=<pptx>, out=<txt>")
    p.add_argument("--no-label-bake", action="store_true",
                   help="skip the sensitivity-label inherit + autofit bake (no PowerPoint)")
    a = p.parse_args()
    if a.dump_text:
        from pptx import Presentation  # local import: only this path needs python-pptx directly
        _write_deck_text(Presentation(a.plan).slides, a.out)
        print(f"deck text -> {a.out}")
        return
    out = build_deck(a.plan, a.out, a.ledger, resolve=not a.no_resolve, dry_run=a.dry_run,
                     label_and_bake=not a.no_label_bake)
    if out:
        print(f"deck -> {out}")

if __name__ == "__main__":
    main()
