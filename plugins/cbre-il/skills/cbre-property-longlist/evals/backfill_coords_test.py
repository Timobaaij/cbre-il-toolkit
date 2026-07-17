#!/usr/bin/env python3
"""backfill_coords_test.py - the generalised, input-agnostic map-link resolver
(extract_pdf.backfill_link_coords). Fast + focused (not the 5-min extract_test). Exercises:
  Pass A  stashed __meta.map_candidates (xlsx/email)  -> coords + prov, key deleted
  precedence  numeric lat + mapLink present           -> untouched
  Pass B  PDF re-open, coords in visible PAGE TEXT    -> coords (the NEW text scan)
  Pass C  PPTX re-open, coords in slide text          -> coords
fitz/pptx sub-cases self-skip if the engine is absent (mirrors the resolver's own degradation)."""
import sys
import pathlib
import tempfile

sys.path.insert(0, str(pathlib.Path(__file__).resolve().parent.parent / "helpers"))
import extract_pdf as P  # noqa: E402


def check(name, cond):
    if not cond:
        raise AssertionError(name)


# --- Pass A: stashed candidate (xlsx/email path) ---
recs = [{"park": "X", "__meta": {"source_type": "xlsx", "source_file": "t.xlsx",
        "locator_base": "Sheet1",
        "map_candidates": ["https://maps.google.com/?q=40.4168,-3.7038"], "prov": {}}}]
P.backfill_link_coords(recs, ".")
check("A-lat", recs[0].get("lat") == 40.4168 and recs[0].get("lng") == -3.7038)
check("A-link", str(recs[0].get("mapLink", "")).startswith("https://maps.google.com"))
check("A-prov", "map link" in recs[0]["__meta"]["prov"].get("lat", ""))
check("A-consumed", "map_candidates" not in recs[0]["__meta"])  # deleted after resolve

# a CUE-LABELLED plain pair stashed (email prose) resolves; prov uses locator_base + a 'coordinates' tag
recs_e = [{"__meta": {"source_type": "email", "source_file": "Offer", "locator_base": "email 2026-07-01",
          "map_candidates": ["coordinates: 45.4642, 9.1900 (per the landlord)"], "prov": {}}}]
P.backfill_link_coords(recs_e, ".")
check("A-plain", recs_e[0].get("lat") == 45.4642 and recs_e[0].get("lng") == 9.19)
check("A-plain-prov", "email 2026-07-01" in recs_e[0]["__meta"]["prov"].get("lat", "")
      and "coordinates" in recs_e[0]["__meta"]["prov"].get("lat", ""))

# HONESTY GUARD: a period-thousands SIZE list stashed from a cell (no coord cue) must NOT become a pin
recs_area = [{"__meta": {"source_type": "xlsx", "source_file": "t.xlsx", "locator_base": "Sheet1",
             "map_candidates": ["Superficie 12.500, 18.500 m2"], "prov": {}}}]
P.backfill_link_coords(recs_area, ".")
check("area-noguess", "lat" not in recs_area[0])          # 12,500 / 18,500 sqm not misread as lat/lng
check("area-consumed", "map_candidates" not in recs_area[0]["__meta"])

# --- precedence: numeric lat + mapLink already present -> untouched ---
recs2 = [{"lat": 1.0, "lng": 2.0, "mapLink": "x", "__meta": {"source_type": "xlsx",
         "map_candidates": ["https://maps.google.com/?q=9.999,9.999"], "prov": {}}}]
P.backfill_link_coords(recs2, ".")
check("precedence", recs2[0]["lat"] == 1.0 and recs2[0]["lng"] == 2.0)
check("precedence-consumed", "map_candidates" not in recs2[0]["__meta"])

# --- Pass B: PDF re-open, coords in visible page TEXT (not an annotation) ---
try:
    import fitz  # noqa: F401
    with tempfile.TemporaryDirectory() as td:
        pdfp = pathlib.Path(td) / "deck.pdf"
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Unit A - see https://maps.google.com/?q=51.5074,-0.1278 for location")
        doc.save(str(pdfp))
        doc.close()
        rp = [{"__meta": {"source_type": "pdf", "source_file": "deck.pdf", "page_no": 0, "prov": {}}}]
        P.backfill_link_coords(rp, td)
        check("B-pdftext", rp[0].get("lat") == 51.5074 and rp[0].get("lng") == -0.1278)
        check("B-prov", "page 1" in rp[0]["__meta"]["prov"].get("lat", ""))
except ImportError:
    print("(skip Pass B: fitz absent)")

# --- Pass C: PPTX re-open, coords in slide text ---
try:
    from pptx import Presentation
    from pptx.util import Inches
    with tempfile.TemporaryDirectory() as td:
        pptp = pathlib.Path(td) / "deck.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        tb.text_frame.text = "Location https://maps.google.com/?q=41.3851,2.1734"
        prs.save(str(pptp))
        rc = [{"__meta": {"source_type": "pptx", "source_file": "deck.pptx", "page_no": 0, "prov": {}}}]
        P.backfill_link_coords(rc, td)
        check("C-pptxtext", rc[0].get("lat") == 41.3851 and rc[0].get("lng") == 2.1734)
        check("C-prov", "slide 1" in rc[0]["__meta"]["prov"].get("lat", ""))
except ImportError:
    print("(skip Pass C: python-pptx absent)")

print("BACKFILL COORDS TEST: PASS")
