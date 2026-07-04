#!/usr/bin/env python3
"""extract_test.py - extraction-hardening regression suite (improvements 3-5).

Each case here was a VERIFIED real misparse from the 2026-06-10 audit (IDs from
the findings file). The suite executes the actual extractors against synthetic
inputs - offline, no network, no real brochures.

Run: python evals/extract_test.py    (exit 0 on success, 1 on any failure)
"""
from __future__ import annotations

import json
import sys
import tempfile
from pathlib import Path

HELPERS = Path(__file__).resolve().parent.parent / "helpers"
sys.path.insert(0, str(HELPERS))
import extract_email  # noqa: E402
import extract_pdf as P  # noqa: E402
import extract_xlsx as X  # noqa: E402
import images as IMG  # noqa: E402
import match  # noqa: E402
import merge  # noqa: E402

FAILS: list[str] = []


def check(ok: bool, label: str) -> None:
    print(f"  [{'PASS' if ok else 'FAIL'}] {label}")
    if not ok:
        FAILS.append(label)


def call(module, *cmd) -> int:
    """Run a helper main() in-process; return its exit code (stdout swallowed)."""
    import io
    from contextlib import redirect_stdout
    saved = sys.argv
    sys.argv = [getattr(module, "__name__", "helper"), *[str(c) for c in cmd]]
    try:
        with redirect_stdout(io.StringIO()):
            module.main()
        rc = 0
    except SystemExit as e:
        rc = e.code if isinstance(e.code, int) else (0 if e.code is None else 1)
    except Exception:
        rc = 1
    finally:
        sys.argv = saved
    return rc


def pdf_cases() -> None:
    print("extract_pdf:")
    # E1: a size under the short ambiguous "Warehouse" label is NOT a rent
    r = P.parse_property_page("City\nPilsen\nDeveloper\nCTP\nWarehouse\n39 471 sq m\n",
                              None, "R", "CZ", "f.pdf", 1)
    check("warehouseRent" not in r and r.get("warehouseRentVal") is None,
          "E1: 'Warehouse 39 471 sq m' does not become a €39,471 rent")

    # E2: spaced monthly forms annualise x12 (these shipped 12x too low before)
    for raw in ("€4.20 / sq m / month", "€4,20 / m2 / mes", "EUR 4.20 / m2 / Monat"):
        d, n, note, _u = P._parse_rent(raw)
        check(n == 50.4 and "x12" in note, f"E2: {raw!r} -> 50.4 annual")
    d, n, _note, _u = P._parse_rent("€60 / sq m / year")
    check(n == 60.0, "E2: annual rent stays annual")

    # E1b: implausible figure with currency context ships as TEXT, never a numeric
    d, n, _note, _u = P._parse_rent("€39471 / sq m")
    check(d is not None and n is None, "E1b: implausible €/m² figure -> text only")

    # E5: ranges keep their honest text; first-plausible number, not the maximum
    r, pr = {}, {}
    P._apply_num(r, pr, "warehouseArea", "10,000 - 35,000 sq m", "page 2")
    check("warehouseArea" not in r and r.get("divisibleFrom") == "10,000 - 35,000 sq m",
          "E5: warehouseArea range -> divisibleFrom text, no invented numeric")
    r, pr = {}, {}
    P._apply_num(r, pr, "warehouseArea", "39 471 sq m (expandable to 80 000)", "page 2")
    check(r.get("warehouseArea") == 39471.0, "E5: first plausible number wins (not 80000)")
    r, pr = {}, {}
    P._apply_num(r, pr, "plotArea", "5 - 12 ha", "page 2")
    check(r.get("plotArea") == "5 - 12 ha", "E5: plotArea range keeps text (schema-legal)")

    # E6: EU dot-thousand size lists are not coordinates; real coords still found
    check(P._find_latlng("Naves de 12.500, 8.750 y 6.200 m2 disponibles") is None,
          "E6: Spanish size list does not become a lat/lng")
    check(P._find_latlng("GPS\n49.7384, 13.3736\n") == (49.7384, 13.3736),
          "E6: genuine coordinates still parse")
    check(P._find_latlng("123.4567, 200.1234") is None, "E6: out-of-range pair rejected")

    # E7: bare-space inline prose is not a labelled value
    r = P.parse_property_page("Ciudad: Madrid\nClear height 10.50 m\n"
                              "City centre is 5 km away from the site\n",
                              None, "R", "ES", "f.pdf", 1)
    check(r.get("city") == "Madrid" and r.get("clearHeight") == "10.50 m",
          "E7: 'City centre is 5 km away' prose rejected; separator/short values kept")

    # E4: TOC names map to PROPERTY pages positionally (a text-less divider page
    # used to shift every later park name by one)
    try:
        import fitz
        doc = fitz.open()
        def page(lines):
            pg = doc.new_page()
            y = 60
            for ln in lines:
                pg.insert_text((40, y), ln, fontsize=12); y += 20
        page(["1.", "Park One", "2.", "Park Two"])                       # TOC
        page(["City", "Pilsen", "Developer", "CTP",
              "Warehouse Area", "40 000 sq m"])                          # property 1
        doc.new_page()                                                    # text-less divider
        page(["City", "Brno", "Developer", "VGP",
              "Warehouse Area", "25 000 sq m"])                          # property 2
        with tempfile.TemporaryDirectory() as td:
            f = Path(td) / "toc.pdf"
            doc.save(f); doc.close()
            recs = P.extract(f, "R", "CZ")
        check(len(recs) == 2 and recs[0].get("park") == "Park One"
              and recs[1].get("park") == "Park Two",
              "E4: divider page does not shift TOC park names")
    except Exception as e:
        check(False, f"E4: TOC test errored ({e})")


def _run_spine(folder: Path, work: Path):
    """Run run.py's main() in-process on a folder; return its exit code (output
    swallowed). Mirrors the orchestrator's single spine invocation."""
    import io as _io
    from contextlib import redirect_stdout, redirect_stderr
    import run
    saved = sys.argv
    sys.argv = ["run.py", "--folder", str(folder), "--work", str(work),
                "--client", "TEDi", "--no-resume", "--quiet"]
    rc = 0
    try:
        with redirect_stdout(_io.StringIO()), redirect_stderr(_io.StringIO()):
            run.main()
    except SystemExit as e:
        rc = e.code if isinstance(e.code, int) else (0 if e.code is None else 1)
    except Exception as e:  # a crash is a test failure, surfaced by the caller
        rc = f"CRASH: {type(e).__name__}: {e}"
    finally:
        sys.argv = saved
    return rc


def interpret_cases() -> None:
    """LLM-first brochure INTERPRETATION: text-vs-raster routing, the generalised
    manifest shape, and the producer<->gate contract for a hand-written interpreted
    <region>_vision.json (resume -> merge -> trace-coverage + validate-data clean).
    Exercises run.py's NEW brochure routing END-TO-END, not just interpret_prep in
    isolation (false-green guard)."""
    print("interpret_prep + run.py interpretation routing:")
    import interpret_prep as IP
    try:
        import fitz
    except Exception as e:
        check(False, f"interpret: fitz unavailable ({e})"); return

    def _text_pdf(td: Path, name: str = "Options - Valencia.pdf") -> Path:
        """A clean born-digital flyer: every page carries selectable spec text whose
        LABELS are deliberately NOT in extract_pdf's dictionary (the TEDi case), so
        the only honest path is interpretation, not the parser."""
        doc = fitz.open()
        for opt, city, area, rent in (("Option 1", "Valencia", "12,500 m2", "4.20 / sqm / month"),
                                      ("Option 2", "Sagunto", "8,750 m2", "3.90 / sqm / month")):
            pg = doc.new_page()
            y = 60
            for ln in (f"VALENCIA REGION - {opt}", f"City {city}", "Owner/developer Goodman",
                       f"Total existing space {area}", f"Warehouse - Asking rent {rent}",
                       "This prime logistics warehouse is strategically located near the A-7 "
                       "motorway with excellent connectivity to the Port of Valencia."):
                pg.insert_text((40, y), ln, fontsize=11); y += 22
        f = td / name
        doc.save(str(f)); doc.close()
        return f

    def _image_only_pdf(td: Path, name: str = "Scan - Cataluna.pdf") -> Path:
        """An image/vector-only deck (no text layer) -> raster mode."""
        doc = fitz.open()
        for _ in range(2):
            pg = doc.new_page()
            pg.draw_rect(fitz.Rect(50, 50, 520, 420), color=(0, 0, 1), fill=(0.2, 0.4, 0.8))
        f = td / name
        doc.save(str(f)); doc.close()
        return f

    # --- 1. interpret_prep routing (unit) ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        ent_t = IP.prepare(_text_pdf(td), "Valencia", "ES", td / "vis")
        check(ent_t.get("mode") == "text" and ent_t.get("source_type") == "pdf",
              "INTERP-1: a clean-text deck routes to mode 'text'")
        check(len(ent_t["pages"]) >= 2
              and all(isinstance(p.get("page_no"), int) for p in ent_t["pages"])
              and ent_t["pages"][0]["page_no"] == 0
              and ent_t["pages"][0]["locator"] == "page 1"
              and "VALENCIA REGION" in ent_t["pages"][0]["text"],
              "INTERP-1: text deck carries per-page {page_no (0-based), locator, text}")
        # NEW: every text-deck page also carries a `render` thumbnail (a downscaled WHOLE-PAGE
        # render, so a VECTOR site plan is visible) for picking __meta.plan_page. Native
        # PyMuPDF renders here; the path must exist and be <= the render thumb edge.
        from PIL import Image as _Img1
        check(all("render" in p for p in ent_t["pages"])
              and all(p.get("render") and Path(p["render"]).exists() for p in ent_t["pages"])
              and all(max(_Img1.open(p["render"]).size) <= IP.PAGE_RENDER_THUMB_EDGE + 1
                      for p in ent_t["pages"]),
              "INTERP-1: each text-deck page carries a per-page render thumbnail (<= the edge)")
        ent_r = IP.prepare(_image_only_pdf(td), "Cataluna", "ES", td / "vis")
        check(ent_r.get("mode") == "raster" and ent_r["pages"]
              and ent_r["pages"][0].get("image"),
              "INTERP-2: an image-only deck routes to mode 'raster' (page PNGs)")

    # --- 2. run.py routes EVERY brochure through interpretation -> exit 3 + manifest ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        folder = td / "inputs"; folder.mkdir()
        work = td / "work"; work.mkdir()
        _text_pdf(folder)
        rc = _run_spine(folder, work)
        check(rc == 3, f"INTERP-3: run.py exits 3 (interpretation needed), got {rc!r}")
        mf = work / "vision" / "manifest.json"
        ok = mf.exists()
        man = json.loads(mf.read_text(encoding="utf-8")) if ok else {}
        deck = (man.get("decks") or [{}])[0]
        check(ok and deck.get("mode") == "text" and deck.get("pages")
              and all("text" in p and isinstance(p.get("page_no"), int)
                      for p in deck["pages"]),
              "INTERP-3: the generalised manifest carries per-deck 'mode' + the text payload")
        check("interpretation.md" in man.get("instructions", "")
              and "text interpretation" in man.get("instructions", "")
              and "COPIED VERBATIM" in man.get("instructions", ""),
              "INTERP-3: manifest instructions point to the interpretation contract (both modes, verbatim page_no)")

    # --- 3. producer<->gate contract: a hand-written interpreted <region>_vision.json
    #        flows resume -> merge -> gate_runner trace-coverage + validate-data clean ---
    import gate_runner
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        folder = td / "inputs"; folder.mkdir()
        work = td / "work"; work.mkdir()
        src = _text_pdf(folder)  # the real source the records cite
        # first pass: routes to interpretation, writes the manifest, exits 3
        check(_run_spine(folder, work) == 3, "INTERP-4: first pass exits 3 (manifest written)")
        # the interpretation sub-agent's output: schema records WITH prov, page_no
        # copied verbatim from the manifest (0 and 1 - both pages carry text)
        def _rec(park, city, page_no, rent_val):
            fields = {"park": park, "developer": "Goodman", "city": city, "country": "ES",
                      "region": "Valencia", "status": "Existing", "warehouseArea": 12500,
                      "warehouseRent": f"€{rent_val:g} / sq m / year",
                      "warehouseRentVal": rent_val}
            prov = {k: f"page {page_no + 1} (text interpretation)" for k in fields}
            return {**fields, "__meta": {"source_file": src.name, "source_type": "pdf",
                                         "locator_base": f"page {page_no + 1}",
                                         "page_no": page_no, "prov": prov}}
        (work / "extract" / "Valencia_vision.json").write_text(
            json.dumps([_rec("Goodman Valencia", "Valencia", 0, 50.4),
                        _rec("Goodman Sagunto", "Sagunto", 1, 46.8)]),
            encoding="utf-8")
        # second pass: resume folds the interpreted records, merges, runs the gates.
        # exit 0 (delivered) OR an exit beyond extraction means it got PAST extraction
        # cleanly (3 = still needs interpretation = a contract failure; 2 = no records).
        rc2 = _run_spine(folder, work)
        check(rc2 not in (2, 3) and not isinstance(rc2, str),
              f"INTERP-4: resume folds the interpreted records past extraction (exit {rc2!r}, not 2/3/crash)")
        canonical = work / "canonical.json"
        ledger = work / "source_ledger.csv"
        check(canonical.exists() and ledger.exists(),
              "INTERP-4: merge produced canonical.json + source_ledger.csv from interpreted records")
        if canonical.exists():
            cj = json.loads(canonical.read_text(encoding="utf-8"))
            check(len(cj.get("properties", [])) == 2,
                  "INTERP-4: both interpreted options become properties (no collapse)")
        check(call(gate_runner, "validate-data", canonical) == 0,
              "INTERP-4: validate-data PASSES on the interpreted dataset")
        check(call(gate_runner, "trace-coverage", canonical, "--ledger", ledger) == 0,
              "INTERP-4: trace-coverage PASSES on the interpreted dataset (no untraced field)")
        # the PASS above is satisfied by merge's REGENERATED ledger; assert the interpreted
        # records' OWN provenance actually reaches the ledger + the LLM-read fields are
        # Medium (not deterministic High) - so a record dropping its prov / mis-stamped
        # confidence is caught, not silently passed.
        import csv as _csv
        led_rows = list(_csv.DictReader((ledger.read_text(encoding="utf-8") if ledger.exists()
                                         else "").splitlines()))
        interp_rows = [r for r in led_rows if "text interpretation" in (r.get("source_locator") or "")]
        check(bool(interp_rows),
              "INTERP-4: interpreted-field provenance reaches the ledger ('(text interpretation)' locator)")
        check(bool(interp_rows) and all(r.get("confidence") == "Medium" for r in interp_rows),
              "INTERP-4: interpreted (LLM-read) fields are Medium confidence, not deterministic High")

    # --- 5. needs_raster ESCALATION: a text deck the sub-agent flags as garbled must
    #        escalate to raster on re-run, not wedge the run at exit 3 forever ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        folder = td / "inputs"; folder.mkdir()
        work = td / "work"; work.mkdir()
        src = _text_pdf(folder)
        check(_run_spine(folder, work) == 3, "INTERP-5: first pass exits 3 (text manifest)")
        mf = work / "vision" / "manifest.json"
        deck = (json.loads(mf.read_text(encoding="utf-8")).get("decks") or [{}])[0]
        region, sf = deck.get("region", ""), deck.get("source_file", src.name)
        check(deck.get("mode") == "text", "INTERP-5: the deck starts in mode 'text'")
        # the sub-agent finds the text unusable and writes the DOCUMENTED escalation stub
        (work / "extract" / f"{region}_vision.json").write_text(
            json.dumps([{"__meta": {"source_file": sf, "needs_raster": True}}]), encoding="utf-8")
        rc = _run_spine(folder, work)
        check(rc == 3, f"INTERP-5: re-run still needs interpretation (exit 3), not a wedge / exit 2 (got {rc!r})")
        deck2 = next((d for d in json.loads(mf.read_text(encoding="utf-8")).get("decks", [])
                      if d.get("source_file") == sf), {})
        check(deck2.get("mode") == "raster" and bool(deck2.get("pages")),
              "INTERP-5: the garbled deck ESCALATED to mode 'raster' (needs_raster consumed, not stuck)")
        check(not (work / "extract" / f"{region}_vision.json").exists(),
              "INTERP-5: the needs_raster stub is consumed (stripped), never re-fed as a record")

    # --- 6. Fix-B carry-forward: a prep-failure gap recorded on a PRIOR pass must SURVIVE
    #        the always-write of unreadable.json across re-runs (a mixed-run exit-0 re-run
    #        skips an un-preppable deck via has_vision, so it is not re-derived). The
    #        no-rasteriser failed-prep path itself only fires in a sandbox without a
    #        renderer; here we lock the carry-forward MECHANISM (the early write no longer
    #        permanently clobbers a prep gap) which is what made it a silent drop. ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        folder = td / "inputs"; folder.mkdir()
        work = td / "work"; work.mkdir()
        src = _text_pdf(folder)
        check(_run_spine(folder, work) == 3, "INTERP-6: first pass exits 3 (creates the work dir)")
        # a PRIOR pass recorded src as an un-rasterisable prep failure
        (work / "unreadable.json").write_text(json.dumps([{"file": src.name,
            "reason": "opened but could not be read as text or rasterised (needs LibreOffice / python-pptx, or the deck is damaged)"}]),
            encoding="utf-8")
        _run_spine(folder, work)  # re-run: the always-write clobbers, the fold must carry it forward
        ur = json.loads((work / "unreadable.json").read_text(encoding="utf-8"))
        check(any(isinstance(e, dict) and e.get("file") == src.name for e in ur),
              "INTERP-6: a prior prep-failure gap is carried forward (not silently clobbered) on re-run")


def xlsx_cases() -> None:
    print("extract_xlsx:")
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.append(["Property", "Location notes", "Rent free (months)",
               "Current rent (EUR/sqm/month)", "GLA (sqm)", "Status"])
    ws.append(["Beta Park", "city centre nearby", 6, 4.2, 25000, "Built"])
    ws.append(["Total", "", "", "", 25000, ""])
    with tempfile.TemporaryDirectory() as td:
        f = Path(td) / "tracker.xlsx"
        wb.save(f)
        res = X.detect_and_extract(f)
    recs = res["records"]
    check(len(recs) == 1, "F18: 'Total' footer row is not a phantom property")
    r = recs[0] if recs else {}
    check(r.get("warehouseRentVal") == 50.4 and r.get("warehouseRent") == "€50.4 / sq m / year",
          "F16: monthly unit in the HEADER -> rent annualised x12")
    check("city" not in r, "F15: 'Location notes' does not map to city")
    prov = r.get("__meta", {}).get("prov", {})
    check(all("!r" in v for v in prov.values()) and prov,
          "F17: provenance locators carry real row numbers")
    check(X._header_field("Rent free (months)") is None,
          "F15: 'Rent free (months)' never maps to the rent column")
    check(X._header_field("Headline rent") == "warehouseRentVal",
          "F15: genuine rent headers still map")

    # F25: a UNIT-SILENT rent (bare 'Rent' header + bare numeric cell) must NOT ship a
    # silent €/sq m/yr default - it ships the number but marks rentUnit ASSUMED, with a
    # provenance note AND a header_report flag so the yield/Gaps pipeline surfaces it
    # (a bare UK £8.5/sq ft figure was passing as €8.5/sq m/yr - wrong currency AND unit)
    def _run_rent(header, cell):
        wb2 = Workbook(); ws2 = wb2.active
        ws2.append(["City", "Developer", "Warehouse Area", header])
        ws2.append(["Corby", "EVO", 170000, cell])
        with tempfile.TemporaryDirectory() as td2:
            f2 = Path(td2) / "r.xlsx"
            wb2.save(f2)
            return X.detect_and_extract(f2)
    res_silent = _run_rent("Rent", 8.5)
    rs = res_silent["records"][0]
    check(rs.get("warehouseRentVal") == 8.5 and rs.get("rentUnitAssumed") is True
          and "ASSUMED" in rs.get("__meta", {}).get("prov", {}).get("rentUnit", ""),
          "F25: a unit-silent rent ships the number but flags rentUnit ASSUMED with provenance")
    check(any(hr.get("rent_unit_assumed") for hr in res_silent.get("header_report", [])),
          "F25: the unit-silent rent raises rent_unit_assumed on the header_report (Gaps surfaces it)")
    # REGRESSION GUARD: a header (or cell) that DOES state the unit is byte-identical to
    # today - the assumed path fires ONLY when both header and cell are unit-silent
    res_hdr = _run_rent("Rent (£ per sq ft)", 8.5)
    rh = res_hdr["records"][0]
    check(rh.get("rentUnit") == "£/sq ft/yr" and rh.get("warehouseRent") == "£8.5 / sq ft / year"
          and "rentUnitAssumed" not in rh
          and not any(hr.get("rent_unit_assumed") for hr in res_hdr.get("header_report", [])),
          "F25 GUARD: a header stating £/sq ft is unchanged (no ASSUMED flag, byte-identical)")
    res_cell = _run_rent("Rent", "£8.50 psf")
    rc = res_cell["records"][0]
    check(rc.get("rentUnit") == "£/sq ft/yr" and "rentUnitAssumed" not in rc,
          "F25 GUARD: a cell stating psf is unchanged (no ASSUMED flag)")

    # LANDLORD vs DEVELOPER are DISTINCT fields (the conflation fix). The dictionary
    # COLUMN_MAP no longer lists landlord/owner/asset-manager as developer synonyms and
    # there is no landlord->developer backfill, so each role maps to its OWN field.
    def _ld(header, cell):
        wbl = Workbook(); wsl = wbl.active
        wsl.append(["Property", "City", "Warehouse Area", header])
        wsl.append(["LD Park", "Corby", 50000, cell])
        with tempfile.TemporaryDirectory() as tdl:
            fl = Path(tdl) / "ld.xlsx"
            wbl.save(fl)
            return X.detect_and_extract(fl)["records"][0]
    # (1) a Landlord column maps to landlord, NOT developer
    r_land = _ld("Landlord", "NFU Mutual")
    check(r_land.get("landlord") == "NFU Mutual" and "developer" not in r_land,
          "LD-1: a Landlord column maps to landlord, NOT developer")
    # an Owner / Asset manager / Freeholder column also maps to landlord (one field)
    check(_ld("Owner", "NFU Mutual").get("landlord") == "NFU Mutual"
          and _ld("Asset Manager", "Schroders").get("landlord") == "Schroders"
          and _ld("Freeholder", "The Crown").get("landlord") == "The Crown",
          "LD-1: Owner / Asset manager / Freeholder all map to the landlord field")
    # a Developer column still maps to developer (untouched)
    r_dev = _ld("Developer", "Canmoor")
    check(r_dev.get("developer") == "Canmoor" and "landlord" not in r_dev,
          "LD-1: a Developer column still maps to developer")
    # a 'Landlord Verified' column NEVER binds to landlord (NEGATIVE-table veto)
    check(X._header_field("Landlord Verified") != "landlord",
          "LD-1: a 'Landlord Verified' column never binds to landlord (NEGATIVE veto)")
    # (2) BOTH a Landlord AND a Developer column present -> BOTH populated, no backfill,
    #     no cross-contamination
    wbb = Workbook(); wsb = wbb.active
    wsb.append(["Property", "City", "Warehouse Area", "Developer", "Landlord"])
    wsb.append(["Raven Park", "Corby", 50000, "Canmoor", "NFU Mutual"])
    with tempfile.TemporaryDirectory() as tdb:
        fb = Path(tdb) / "both.xlsx"
        wbb.save(fb)
        rb = X.detect_and_extract(fb)["records"][0]
    check(rb.get("developer") == "Canmoor" and rb.get("landlord") == "NFU Mutual",
          "LD-2: a Developer AND a Landlord column BOTH populate their own fields (no conflation)")
    # WITHIN-field fallback still works: a 'Promoter' column fills an EMPTY Developer cell
    wbp = Workbook(); wsp = wbp.active
    wsp.append(["Property", "City", "Warehouse Area", "Developer", "Promoter"])
    wsp.append(["Promo Park", "Corby", 50000, "", "Canmoor"])
    with tempfile.TemporaryDirectory() as tdp:
        fp = Path(tdp) / "promo.xlsx"
        wbp.save(fp)
        rp = X.detect_and_extract(fp)["records"][0]
    check(rp.get("developer") == "Canmoor",
          "LD-2: a 'Promoter' column still fills an empty Developer cell (within-field fallback kept)")

    # LD-3 (CONTINENTAL): the ledger landlord/developer split holds across languages
    # too (mirrors LD-1/LD-2). An owner-language header (DE Eigentümer, ES Propietario,
    # FR Propriétaire, NL Eigenaar, IT/PT Proprietario, PL Właściciel) maps to landlord,
    # NOT developer; a developer-language header (DE Entwickler, ES Promotor, FR
    # Promoteur, NL Ontwikkelaar) maps to developer. There is no landlord->developer
    # backfill in any language - an owner is never a developer.
    for owner_hdr in ("Eigentümer", "Propietario", "Propriétaire", "Eigenaar",
                      "Proprietario", "Proprietário", "Właściciel"):
        ro = _ld(owner_hdr, "Garbe")
        check(ro.get("landlord") == "Garbe" and "developer" not in ro,
              f"LD-3: a '{owner_hdr}' (owner) column maps to landlord, NOT developer")
    for dev_hdr in ("Entwickler", "Promotor", "Promoteur", "Ontwikkelaar", "Sviluppatore"):
        rd = _ld(dev_hdr, "Panattoni")
        check(rd.get("developer") == "Panattoni" and "landlord" not in rd,
              f"LD-3: a '{dev_hdr}' (developer) column maps to developer, NOT landlord")
    # both present in one DE tracker -> both populate their own field, no conflation
    wbc = Workbook(); wsc = wbc.active
    wsc.append(["Objekt", "Stadt", "Hallenfläche", "Entwickler", "Eigentümer"])
    wsc.append(["DE Park", "Hamburg", 50000, "Panattoni", "Garbe"])
    with tempfile.TemporaryDirectory() as tdc:
        fc = Path(tdc) / "de.xlsx"
        wbc.save(fc)
        rc = X.detect_and_extract(fc, "", "DE")["records"][0]
    check(rc.get("developer") == "Panattoni" and rc.get("landlord") == "Garbe",
          "LD-3: a DE tracker with Entwickler AND Eigentümer populates BOTH (developer + landlord, no conflation)")


def tracker_map_cases() -> None:
    """LLM TRACKER MAPPING: the column->field decision moves to an isolated sub-agent
    (mirrors the INTERP cases that hand-write vision.json) while the dictionary stays
    the offline fallback (column_map=None byte-identical), a NEGATIVE-table hard veto,
    a backfill of unbound columns, and a logged cross-check. Python still PARSES every
    number with the same arithmetic, so every numeric guarantee is byte-preserved."""
    print("extract_xlsx LLM tracker mapping:")
    from openpyxl import Workbook
    import run as RUN
    import hashlib

    def _corby_sheet():
        wb = Workbook(); ws = wb.active
        ws.append(["Building ID", "Marketing Name", "Town", "Latitude, Longitude",
                   "Construction status", "Status", "Landlord", "Developer",
                   "Size (sq ft)", "Size Unit", "Current quoting rent (£ per sq ft)",
                   "Historical quoting rent (£ per sq ft)", "Eaves (m)",
                   "Eaves 10m or above?", "Office content (sq ft)",
                   "Floor loading (kN/sq m)", "No. of dock level doors",
                   "Site area (acres)", "Power (KVA)"])
        ws.append([10179, "EVO 169", "Corby", "52.50304981, -0.650581854",
                   "Built", "Available", "EVO", "",
                   172867, "GIA", 8.5, 0, 15, "Yes", 13576, 50, 15, 5.2, 800])
        td = Path(tempfile.mkdtemp()); f = td / "b.xlsx"; wb.save(f); return f

    # --- (a) a hand-map mirroring the dictionary's decisions is BYTE-IDENTICAL to the
    #         dictionary path (the LLM path is correct + the arithmetic is preserved) ---
    f = _corby_sheet()
    dict_recs = X.detect_and_extract(f, "Corby", "GB")["records"]
    hand = {"columns": [
        {"index": 1, "field": "park"}, {"index": 2, "field": "city"},
        {"index": 3, "field": "latlng"}, {"index": 5, "field": "status"},
        {"index": 7, "field": "developer"},
        {"index": 8, "field": "warehouseArea", "areaUnit": "sq ft", "basis": "GIA"},
        {"index": 9, "field": None, "role": "size_basis"},
        {"index": 10, "field": "warehouseRentVal", "currency": "GBP",
         "perArea": "sq ft", "period": "annual"},
        {"index": 12, "field": "clearHeight"},
        {"index": 14, "field": "officeArea", "areaUnit": "sq ft"},
        {"index": 15, "field": "floorLoad"}, {"index": 16, "field": "loadingDocks"},
        {"index": 17, "field": "plotArea", "areaUnit": "acres"},
        {"index": 18, "field": "electricity"}]}
    llm_recs = X.detect_and_extract(f, "Corby", "GB", column_map=hand)["records"]
    eq = json.dumps(dict_recs, sort_keys=True, ensure_ascii=False) == \
        json.dumps(llm_recs, sort_keys=True, ensure_ascii=False)
    check(eq, "TRK-a: a hand-map mirroring the dictionary is byte-identical (GIA-office 159,291, £/sq ft, acres)")
    # the corby sheet has Landlord="EVO" + an EMPTY Developer cell. Landlord and
    # developer are now DISTINCT fields: the GIA-office subtraction survives, EVO lands
    # in `landlord` (NOT `developer`), and developer is left absent (no backfill - the
    # extractor never fabricates a developer from a landlord). The split holds on BOTH
    # the LLM path and the dictionary path (the byte-identity check above proves it).
    check(bool(llm_recs) and llm_recs[0].get("warehouseArea") == 159291
          and llm_recs[0].get("landlord") == "EVO"
          and "developer" not in llm_recs[0],
          "TRK-a: GIA-office subtraction survives + the Landlord maps to landlord, NOT developer (no backfill)")

    # --- (e) column_map=None == the dictionary path verbatim (the offline-eval invariant) ---
    none_recs = X.detect_and_extract(f, "Corby", "GB", column_map=None)["records"]
    check(json.dumps(none_recs, sort_keys=True) == json.dumps(dict_recs, sort_keys=True),
          "TRK-e: column_map=None is the dictionary path verbatim (byte-identical)")

    # --- (b) a continental sheet the dictionary maps THIN maps RICHLY with a hand-map
    #         AND still annualises a EUR/m2/Monat column x12 ---
    wb = Workbook(); ws = wb.active
    ws.append(["Objekt", "Stadt", "Hallenfläche (m²)", "Bürofläche (m²)",
               "Miete (€/m²/Monat)", "Status"])
    ws.append(["DC Berlin", "Berlin", "45000", "5000", "4,50", "Verfügbar"])
    td = Path(tempfile.mkdtemp()); fk = td / "kont.xlsx"; wb.save(fk)
    dthin = X.detect_and_extract(fk)["records"][0]
    cont = {"columns": [
        {"index": 0, "field": "park"}, {"index": 1, "field": "city"},
        {"index": 2, "field": "warehouseArea", "areaUnit": "sq m"},
        {"index": 3, "field": "officeArea", "areaUnit": "sq m"},
        {"index": 4, "field": "warehouseRentVal", "currency": "EUR",
         "perArea": "sq m", "period": "monthly"},
        {"index": 5, "field": "status"}]}
    rich = X.detect_and_extract(fk, column_map=cont)["records"][0]
    check("park" not in dthin and rich.get("park") == "DC Berlin",
          "TRK-b: the LLM map binds the 'Objekt' park column the dictionary missed")
    check(rich.get("warehouseRentVal") == 54.0 and rich.get("rentUnit") == "€/sq m/yr",
          "TRK-b: a EUR/m²/Monat column is still annualised x12 (4.50 -> 54.0) by Python")

    # --- (c) a hand-map binding a NEGATIVE-listed column is VETOED + falls back ---
    wb = Workbook(); ws = wb.active
    ws.append(["Property", "City", "Warehouse Area", "Rent free (months)",
               "Headline rent (EUR/sqm/yr)"])
    ws.append(["Beta", "Lyon", 25000, 6, 5.5])
    td = Path(tempfile.mkdtemp()); fv = td / "veto.xlsx"; wb.save(fv)
    bad = {"columns": [
        {"index": 0, "field": "park"}, {"index": 1, "field": "city"},
        {"index": 2, "field": "warehouseArea", "areaUnit": "sq m"},
        {"index": 3, "field": "warehouseRentVal", "currency": "EUR",
         "perArea": "sq m", "period": "annual"},   # NEGATIVE-listed: 'Rent free (months)'
        {"index": 4, "field": "warehouseRentVal", "currency": "EUR",
         "perArea": "sq m", "period": "annual"}]}
    rv = X.detect_and_extract(fv, column_map=bad)
    rvr = rv["records"][0]
    check(rvr.get("warehouseRentVal") == 5.5,
          "TRK-c: a NEGATIVE-vetoed 'Rent free (months)'->rent never ships 6 as a rent (falls back)")
    check(any("NEGATIVE-vetoed" in v for v in rv["header_report"][0].get("llm_vetoes", [])),
          "TRK-c: the veto is logged on the header_report for the reviewer")

    # cross-check (NOT a veto): an LLM binding the dictionary confidently disagrees with
    # (no NEGATIVE veto applies) is LOGGED for the reviewer. 'Status' -> the dictionary
    # field is 'status'; an LLM that calls it 'city' is recorded as a disagreement.
    wb = Workbook(); ws = wb.active
    ws.append(["Park", "Status", "Warehouse Area"])
    ws.append(["Gamma", "Lyon", 25000])
    td = Path(tempfile.mkdtemp()); fx = td / "xc.xlsx"; wb.save(fx)
    xc = {"columns": [{"index": 0, "field": "park"},
                      {"index": 1, "field": "city"},   # dictionary -> status
                      {"index": 2, "field": "warehouseArea"}]}
    xcr = X.detect_and_extract(fx, column_map=xc)
    check(any("dictionary->status" in d and "LLM->city" in d
              for d in xcr["header_report"][0].get("llm_dictionary_disagreements", [])),
          "TRK-c: an LLM-vs-dictionary disagreement is logged (LLM->city; dictionary->status)")

    # --- (d) the map cache is input-hash keyed: tracker_structure is deterministic and
    #         a changed structure recomputes a different hash ---
    def _struct_hash(structs):
        payload = json.dumps([{"region": "", "country": ""}, structs],
                             ensure_ascii=False, sort_keys=True)
        return hashlib.sha1(payload.encode("utf-8")).hexdigest()[:8]
    s1 = X.tracker_structure(f)
    s1b = X.tracker_structure(f)
    check(_struct_hash(s1) == _struct_hash(s1b),
          "TRK-d: tracker_structure is deterministic (same workbook -> same hash)")
    s_cont = X.tracker_structure(fk)
    check(_struct_hash(s1) != _struct_hash(s_cont),
          "TRK-d: a different tracker structure recomputes a different hash")
    check(bool(s1) and s1[0]["headers"][1] == "Marketing Name"
          and "Building ID" in s1[0]["unmapped_headers"]
          and len(s1[0]["sample_rows"]) == 1,
          "TRK-d: tracker_structure carries raw headers, sample rows + the dictionary miss list")

    # --- (f) a FULL-SPINE run on a raw tracker emits a tracker manifest entry + exit 3;
    #         a hand-map (then a .SKIP) lets the re-run reach past extraction (mirror INTERP-4) ---
    with tempfile.TemporaryDirectory() as tds:
        tds = Path(tds)
        folder = tds / "inputs"; folder.mkdir()
        work = tds / "work"; work.mkdir()
        wb = Workbook(); ws = wb.active
        ws.append(["Objekt", "Stadt", "Hallenfläche (m²)", "Bürofläche (m²)",
                   "Miete (€/m²/Monat)", "Status"])
        ws.append(["DC Berlin", "Berlin", "45000", "5000", "4,50", "Verfügbar"])
        ws.append(["DC Koln", "Koln", "30000", "3000", "4,00", "Verfügbar"])
        wb.save(folder / "kont.xlsx")
        rc1 = _run_spine(folder, work)
        check(rc1 == 3, f"TRK-f: a raw tracker exits 3 (mapping offered), got {rc1!r}")
        mf = work / "vision" / "manifest.json"
        man = json.loads(mf.read_text(encoding="utf-8")) if mf.exists() else {}
        job = (man.get("jobs") or [{}])[0]
        check(job.get("kind") == "tracker" and job.get("source_file") == "kont.xlsx"
              and job.get("input_hash") and job.get("output") and job.get("sheets"),
              "TRK-f: the manifest carries a 'tracker' job (kind/input_hash/output/sheets)")
        check(man.get("decks") == [] and "tracker_instructions" in man,
              "TRK-f: the brochure `decks` array is untouched (empty) + tracker_instructions present")
        # hand-map at the job output, input_hash copied verbatim
        out_path = work / "extract" / Path(job["output"]).name
        out_path.parent.mkdir(parents=True, exist_ok=True)
        out_path.write_text(json.dumps({"input_hash": job["input_hash"], "schema_version": 1,
                                        "map": cont}), encoding="utf-8")
        rc2 = _run_spine(folder, work)
        check(rc2 not in (2, 3) and not isinstance(rc2, str),
              f"TRK-f: the hand-map lets the re-run pass extraction (exit {rc2!r}, not 2/3/crash)")
        can = work / "canonical.json"
        if can.exists():
            cj = json.loads(can.read_text(encoding="utf-8"))
            props = cj.get("properties", [])
            check(len(props) == 2 and any(p.get("park") == "DC Berlin" for p in props),
                  "TRK-f: the LLM-mapped tracker yields 2 properties incl. the 'Objekt' park")
        # the .SKIP escape hatch on a SEPARATE tracker (dictionary is fine)
        with tempfile.TemporaryDirectory() as tds2:
            tds2 = Path(tds2)
            folder2 = tds2 / "inputs"; folder2.mkdir()
            work2 = tds2 / "work"; work2.mkdir()
            wb2 = Workbook(); ws2 = wb2.active
            ws2.append(["Marketing Name", "Town", "Status", "Developer",
                        "Size (sq ft)", "Current quoting rent (£ per sq ft)"])
            ws2.append(["Alpha", "Corby", "Available", "Dev", 100000, 7])
            wb2.save(folder2 / "uk.xlsx")
            check(_run_spine(folder2, work2) == 3, "TRK-f: the dictionary-mappable tracker is still OFFERED (exit 3)")
            j2 = json.loads((work2 / "vision" / "manifest.json").read_text(encoding="utf-8"))["jobs"][0]
            (work2 / "extract" / (Path(j2["output"]).stem + ".SKIP")).write_text("", encoding="utf-8")
            rc3 = _run_spine(folder2, work2)
            check(rc3 not in (2, 3) and not isinstance(rc3, str),
                  f"TRK-f: a .SKIP sentinel lets the re-run proceed on the dictionary (exit {rc3!r})")

    # --- (g) REAL-DECK best-effort: the real Corby tracker emits a tracker manifest entry ---
    real = Path(r"C:\Claude Projects\Corby Test Run\Building_Data_12_06_2026.xlsx")
    if real.exists():
        with tempfile.TemporaryDirectory() as tdr:
            tdr = Path(tdr)
            folder = tdr / "inputs"; folder.mkdir()
            work = tdr / "work"; work.mkdir()
            import shutil
            shutil.copy(real, folder / real.name)
            rcg = _run_spine(folder, work)
            man = (json.loads((work / "vision" / "manifest.json").read_text(encoding="utf-8"))
                   if (work / "vision" / "manifest.json").exists() else {})
            jobs = man.get("jobs") or []
            check(rcg == 3 and any(j.get("kind") == "tracker"
                                   and j.get("source_file") == real.name for j in jobs),
                  "TRK-g: the REAL Corby tracker emits a tracker manifest entry (full-spine exit 3)")
    else:
        check(True, "TRK-g: real Corby tracker absent on host - skipped (best-effort)")


def area_band_cases() -> None:
    """P1-A: a CONSERVATIVE area plausibility band + a sq-ft-vs-sq-m magnitude
    cross-check (the area twin of the rent band/unit smell). An out-of-band or
    magnitude-suspect AREA is KEPT (never dropped/coerced to tbd, never auto-
    converted) and SURFACED via header_report -> yield_report -> Gaps. Backward
    compat: Corby 159,291 sq ft passes SILENTLY (no flag, byte-identical)."""
    print("area band + sq-ft/sq-m magnitude cross-check (P1-A):")
    from openpyxl import Workbook
    import normalize as N

    # --- (e) unit tests: the EXACT bands + magnitude thresholds shipped ---
    # acres/ha are CONVERTED to sq ft / sq m at parse BEFORE the band runs, so the
    # band only ever sees the resolved 'sq ft' / 'sq m' (or None -> the sq m band).
    check(N.area_band_for("sq m") == (300, 600000)
          and N.area_band_for(None) == (300, 600000),
          "AREA-e: area_band_for sq m / unknown -> (300, 600000)")
    check(N.area_band_for("sq ft") == (3000, 6500000),
          "AREA-e: area_band_for sq ft -> (3000, 6500000) (= the sq m band x ~10.764)")
    check(N.area_magnitude_mismatch(172867, "sq m") and "NOT auto-converted"
          in N.area_magnitude_mismatch(172867, "sq m"),
          "AREA-e: 172,867 'sq m' is magnitude-suspect (in the sq-ft range)")
    check(N.area_magnitude_mismatch(3500, "sq ft") and "NOT auto-converted"
          in N.area_magnitude_mismatch(3500, "sq ft"),
          "AREA-e: 3,500 'sq ft' is magnitude-suspect (in the sq-m range)")
    check(N.area_magnitude_mismatch(60000, "sq m") is None
          and N.area_magnitude_mismatch(4000, "sq ft") is None,
          "AREA-e: the thresholds are exclusive (60,000 sq m / 4,000 sq ft do NOT trip)")
    check(N.area_magnitude_mismatch(159291, "sq ft") is None
          and N.area_magnitude_mismatch(40000, "sq m") is None,
          "AREA-e: a normal sq ft value and a normal sq m value never trip the magnitude check")

    # --- (d) the conservatism guard: a 300,000 sq m mega-shed AND a 350 sq m unit
    #         both pass the BAND silently (no out-of-band flag) ---
    _lo_m, _hi_m = N.area_band_for("sq m")
    check(_lo_m <= 300000 <= _hi_m and _lo_m <= 350 <= _hi_m,
          "AREA-d: a 300,000 sq m mega-shed AND a 350 sq m unit both pass the band (conservatism)")
    _lo_f, _hi_f = N.area_band_for("sq ft")
    check(_lo_f <= 3200000 <= _hi_f,
          "AREA-d: a 3.2M sq ft mega-campus passes the sq ft band (conservatism)")
    check(N.area_magnitude_mismatch(350, "sq m") is None,
          "AREA-d: the 350 sq m last-mile unit trips NEITHER the band nor the magnitude check")

    # --- (c) BACKWARD-COMPAT: the Corby 159,291 sq ft warehouseArea passes SILENTLY
    #         (no area_out_of_band / area_unit_suspect, value byte-identical) ---
    wb = Workbook(); ws = wb.active
    ws.append(["Marketing Name", "Town", "Status", "Developer",
               "Size (sq ft)", "Size Unit", "Office content (sq ft)",
               "Current quoting rent (£ per sq ft)"])
    ws.append(["EVO 169", "Corby", "Available", "EVO", 172867, "GIA", 13576, 8.5])
    td = Path(tempfile.mkdtemp()); fc = td / "corby.xlsx"; wb.save(fc)
    out = X.detect_and_extract(fc, "Corby", "GB")
    rec = out["records"][0]
    hr = out["header_report"][0]
    check(rec.get("warehouseArea") == 159291 and rec.get("areaUnit") == "sq ft",
          "AREA-c: Corby warehouseArea = GIA - office = 159,291 sq ft (unchanged)")
    check("area_out_of_band" not in hr and "area_unit_suspect" not in hr,
          "AREA-c: Corby passes SILENTLY - NO area_out_of_band / area_unit_suspect flag")
    check("AREA OUT OF BAND" not in rec["__meta"]["prov"].get("warehouseArea", "")
          and "NOT auto-converted" not in rec["__meta"]["prov"].get("warehouseArea", ""),
          "AREA-c: the Corby warehouseArea prov note carries NO band/magnitude suffix")

    # --- (b) MAGNITUDE: a 'Size (sq m)' column holding 172,867 -> KEPT (NOT converted
    #         to ~16,064), areaUnit 'sq m', header_report area_unit_suspect set ---
    wb = Workbook(); ws = wb.active
    ws.append(["Marketing Name", "Town", "Status", "Size (sq m)"])
    ws.append(["Mislabel DC", "Lyon", "Available", 172867])
    td = Path(tempfile.mkdtemp()); fm = td / "mislabel.xlsx"; wb.save(fm)
    hand = {"columns": [
        {"index": 0, "field": "park"}, {"index": 1, "field": "city"},
        {"index": 2, "field": "status"},
        {"index": 3, "field": "warehouseArea", "areaUnit": "sq m"}]}
    outm = X.detect_and_extract(fm, "Lyon", "FR", column_map=hand)
    recm = outm["records"][0]
    hrm = outm["header_report"][0]
    check(recm.get("warehouseArea") == 172867 and recm.get("areaUnit") == "sq m",
          "AREA-b: a 'sq m' value of 172,867 is KEPT (NOT auto-converted to ~16,064)")
    check(isinstance(hrm.get("area_unit_suspect"), list)
          and hrm["area_unit_suspect"][0]["value"] == 172867
          and hrm["area_unit_suspect"][0]["unit"] == "sq m",
          "AREA-b: header_report.area_unit_suspect set (park+value+unit listed)")
    check("NOT auto-converted" in recm["__meta"]["prov"].get("warehouseArea", ""),
          "AREA-b: the warehouseArea prov carries the magnitude-suspect note")

    # the yield note is raised end-to-end (header_report flag -> yield_report.md)
    with tempfile.TemporaryDirectory() as tdy:
        tdy = Path(tdy)
        folder = tdy / "inputs"; folder.mkdir()
        work = tdy / "work"; work.mkdir()
        wb = Workbook(); ws = wb.active
        ws.append(["Marketing Name", "Town", "Status", "Size (sq m)"])
        ws.append(["Mislabel DC", "Lyon", "Available", 172867])
        wb.save(folder / "mislabel.xlsx")
        rc1 = _run_spine(folder, work)
        check(rc1 == 3, f"AREA-b: a raw tracker exits 3 (mapping offered), got {rc1!r}")
        job = json.loads((work / "vision" / "manifest.json").read_text(encoding="utf-8"))["jobs"][0]
        op = work / "extract" / Path(job["output"]).name
        op.parent.mkdir(parents=True, exist_ok=True)
        op.write_text(json.dumps({"input_hash": job["input_hash"], "schema_version": 1,
                                  "map": hand}), encoding="utf-8")
        _run_spine(folder, work)
        yr = (work / "yield_report.md")
        body = yr.read_text(encoding="utf-8") if yr.exists() else ""
        check("area unit looks wrong" in body and "NOT auto-converted" in body,
              "AREA-b: a magnitude yield_report line is raised (surfaced to the Gaps Report)")

    # --- (a) OUT OF BAND: a 9-digit parse-garble -> record PRESENT, value KEPT (not
    #         dropped/tbd), header_report.area_out_of_band set, yield line raised ---
    wb = Workbook(); ws = wb.active
    ws.append(["Marketing Name", "Town", "Status", "Size (sq ft)"])
    ws.append(["Garble DC", "Leeds", "Available", 172867000])  # run-together 9-digit cell
    td = Path(tempfile.mkdtemp()); fo = td / "garble.xlsx"; wb.save(fo)
    outo = X.detect_and_extract(fo, "Leeds", "GB")
    check(bool(outo["records"]) and outo["records"][0].get("warehouseArea") == 172867000,
          "AREA-a: an out-of-band 9-digit area is PRESENT and KEPT (not dropped/coerced to tbd)")
    hro = outo["header_report"][0]
    check(isinstance(hro.get("area_out_of_band"), list)
          and hro["area_out_of_band"][0]["value"] == 172867000,
          "AREA-a: header_report.area_out_of_band set (park+value+unit listed)")
    check("AREA OUT OF BAND" in outo["records"][0]["__meta"]["prov"].get("warehouseArea", ""),
          "AREA-a: the warehouseArea prov carries the out-of-band note (low-confidence, kept)")
    with tempfile.TemporaryDirectory() as tdy:
        tdy = Path(tdy)
        folder = tdy / "inputs"; folder.mkdir()
        work = tdy / "work"; work.mkdir()
        wb = Workbook(); ws = wb.active
        ws.append(["Marketing Name", "Town", "Status", "Size (sq ft)"])
        ws.append(["Garble DC", "Leeds", "Available", 172867000])
        wb.save(folder / "garble.xlsx")
        rc1 = _run_spine(folder, work)
        job = json.loads((work / "vision" / "manifest.json").read_text(encoding="utf-8"))["jobs"][0]
        op = work / "extract" / Path(job["output"]).name
        op.parent.mkdir(parents=True, exist_ok=True)
        hand_a = {"columns": [
            {"index": 0, "field": "park"}, {"index": 1, "field": "city"},
            {"index": 2, "field": "status"},
            {"index": 3, "field": "warehouseArea", "areaUnit": "sq ft"}]}
        op.write_text(json.dumps({"input_hash": job["input_hash"], "schema_version": 1,
                                  "map": hand_a}), encoding="utf-8")
        _run_spine(folder, work)
        yr = (work / "yield_report.md")
        body = yr.read_text(encoding="utf-8") if yr.exists() else ""
        check("outside the plausibility band" in body,
              "AREA-a: an out-of-band yield_report line is raised (kept for broker review)")

    # --- merge override gate honours the band (block-the-override-only) ---
    check(merge._pick_passes_gate("warehouseArea", 159291, None, "sq ft") is True,
          "AREA-merge: an in-band sq ft override (159,291) PASSES the area gate")
    check(merge._pick_passes_gate("warehouseArea", 172867000, None, "sq ft") is False,
          "AREA-merge: an out-of-band override (9-digit) FAILS the gate (precedence stands)")
    check(merge._pick_passes_gate("warehouseArea", 350, None, "sq m") is True
          and merge._pick_passes_gate("warehouseArea", 200, None, "sq m") is False,
          "AREA-merge: the sq m band floor (300) gates a sub-floor override")


def email_cases() -> None:
    print("extract_email:")
    eml = (b"From: agent@example.com\r\nTo: broker@cbre.com\r\n"
           b"Subject: GLP Sziget II offer\r\nDate: Mon, 12 May 2025 10:11:00 +0200\r\n"
           b"Content-Type: text/plain\r\n\r\nRent EUR 62 per sqm per year.\r\n")
    with tempfile.TemporaryDirectory() as td:
        (Path(td) / "offer.eml").write_bytes(eml)
        (Path(td) / "broken.eml").write_bytes(b"\x00\xff not an email")
        out = extract_email.extract(Path(td))
    good = [d for d in out if not d.get("unreadable")]
    check(len(good) == 1 and good[0]["__meta"].get("date") == "2025-05-12",
          "H5: RFC-2822 date lands in __meta.date as ISO (newest-email-wins works)")
    check(bool(good) and good[0]["__meta"]["locator_base"] == "email 2025-05-12",
          "F9: locator matches the documented 'email <yyyy-mm-dd>' form")
    check(len(out) == 2 and any(d.get("unreadable") for d in out),
          "F23: a junk/corrupt file degrades to an explicit stub, not an empty record")


def match_cases() -> None:
    print("match:")
    def rec(park, area, src):
        return {"city": "Prague", "developer": "CTP", "park": park,
                "warehouseArea": area, "__meta": {"source_file": src}}
    check(not match.same_property(rec("", 10000, "a.pdf"), rec("CTPark Prague East", 50000, "b.pdf")),
          "M12: empty park + different area does not subset-merge")
    check(match.same_property(rec("", 50000, "a.pdf"), rec("CTPark Prague East", 50000, "b.pdf")),
          "M12: empty park + matching area still merges")
    check(match.same_property(rec("CTPark Prague East", 50000, "a.pdf"),
                              rec("CTPark Prague East", 50000, "b.pdf")),
          "cross-source identical key still merges")

    # ---- LLM match adjudication (exit 10): tiers, blocker-beats-LLM, decisions, grey, spine ----
    def r2(city, dev, park, area, src, lat=None, lng=None):
        m = {"city": city, "developer": dev, "park": park, "warehouseArea": area,
             "__meta": {"source_file": src}}
        if lat is not None:
            m["lat"], m["lng"] = lat, lng
        return m

    # (a) pair_class tiers ---------------------------------------------------------
    # GREY (was forbidden): a developer DISAGREEMENT (both known, unequal) is no longer a
    # hard block - landlord and developer are distinct fields now, so a genuine
    # developer-name difference is the LLM's to adjudicate. Same park/city + no size
    # conflict -> it clears the recall pre-filter and lands in grey.
    a_dev = r2("Corby", "Prologis", "Raven Park", 50000, "a.pdf")
    b_dev = r2("Corby", "Panattoni", "Raven Park", 50000, "b.xlsx")
    check(match.pair_class(a_dev, b_dev) == "grey",
          "MATCH-a: a developer-disagreement pair is now 'grey' (was 'forbidden'); the LLM adjudicates")
    # forbidden: a material size conflict (>15%), cross-source, stays the HARD blocker
    a_sz = r2("Corby", "Prologis", "Raven Park", 50000, "a.pdf")
    b_sz = r2("Corby", "Prologis", "Raven Park", 100000, "b.xlsx")
    check(match.pair_class(a_sz, b_sz) == "forbidden",
          "MATCH-a: a >15%-area pair STAYS 'forbidden' (the genuine hard blocker)")
    # auto: a cross-source identical key (no conflict)
    a_au = r2("Corby", "Prologis", "Raven Park", 50000, "a.pdf")
    b_au = r2("Corby", "Prologis", "Raven Park", 50000, "b.xlsx")
    check(match.pair_class(a_au, b_au) == "auto",
          "MATCH-a: a cross-source identical-key pair is 'auto'")
    # grey: same city, parks that are NOT subset/superset and NOT >= 88 fuzzy, no conflict
    a_gr = r2("Corby", "Prologis", "Apollo Court", 50000, "a.pdf")
    b_gr = r2("Corby", "Prologis", "Mercury House", 52000, "b.xlsx")
    check(match.pair_class(a_gr, b_gr) == "grey"
          and match.fuzz.token_set_ratio(match.match_key(a_gr), match.match_key(b_gr)) < 88,
          "MATCH-a: a same-city pair below token_set 88 is 'grey'")

    # (b) the blocker BEATS the LLM ------------------------------------------------
    # re-pointed at a SIZE-conflict forbidden pair (a developer disagreement is no longer
    # forbidden): a >15%-area pair can NEVER merge even on an LLM 'same' verdict.
    pid_sz = match.pair_id(a_sz, b_sz)
    check(not match.same_property(a_sz, b_sz, {pid_sz: "same"})
          and not match.same_property(a_sz, b_sz, {pid_sz: {"verdict": "same"}}),
          "MATCH-b: same_property(SIZE-conflict forbidden pair, decisions='same') is STILL False (blocker beats the LLM)")
    # and the DEMOTED developer-disagreement pair is now LLM-adjudicated: it MERGES on a
    # 'same' verdict and SPLITS on 'different' (it was an unmergeable hard block before).
    pid_dev = match.pair_id(a_dev, b_dev)
    check(match.same_property(a_dev, b_dev, {pid_dev: "same"})
          and not match.same_property(a_dev, b_dev, {pid_dev: "different"})
          and not match.same_property(a_dev, b_dev),
          "MATCH-b: a grey developer-disagreement pair MERGES on 'same', SPLITS on 'different' (offline=split)")

    # (c) dedupe consumes decisions: MERGE a grey 'same', SPLIT a grey 'different' --
    pid_gr = match.pair_id(a_gr, b_gr)
    merged = match.dedupe([a_gr, b_gr], {pid_gr: "same"})
    split = match.dedupe([a_gr, b_gr], {pid_gr: "different"})
    none_run = match.dedupe([a_gr, b_gr])  # offline fallback: a grey pair stays distinct
    check(len(merged) == 1 and len(split) == 2 and len(none_run) == 2,
          "MATCH-c: dedupe MERGES a grey pair marked 'same' and SPLITS one marked 'different' (offline=split)")

    # (d) merge.main with NO --match-decisions == today (offline fallback byte-identical) --
    #     a forbidden + an auto + a grey trio: the dedupe must equal decisions=None dedupe
    trio = [a_dev, b_dev, a_au, b_au, a_gr, b_gr]
    check(match.dedupe(trio) == match.dedupe(trio, None),
          "MATCH-d: dedupe(records) == dedupe(records, None) (default-arg offline fallback)")

    # (e) grey_pairs: deterministic, order-independent pair_id, recall-filtered -----
    far = r2("Madrid", "VGP", "Sol Park", 30000, "c.eml", lat=40.4, lng=-3.7)  # unrelated, distant
    recs = [a_gr, b_gr, far]
    gp1 = match.grey_pairs(recs)
    gp2 = match.grey_pairs(list(reversed(recs)))
    check(len(gp1) == 1 and gp1[0]["pair_id"] == pid_gr,
          "MATCH-e: grey_pairs returns exactly the one grey pair (recall-filtered)")
    check({g["pair_id"] for g in gp1} == {g["pair_id"] for g in gp2}
          and match.pair_id(a_gr, b_gr) == match.pair_id(b_gr, a_gr),
          "MATCH-e: grey_pairs is order-independent (pair_id(a,b)==pair_id(b,a))")
    check(all(match.pair_id(far, x) not in {g["pair_id"] for g in gp1} for x in (a_gr, b_gr)),
          "MATCH-e: an obviously-unrelated distant pair is NOT grey (not shown to the LLM)")

    # (g) HEADLINE REGRESSION: a developer (source A) + a landlord (source B) for the
    #     SAME building must NOT split into two cards. extract_xlsx no longer conflates
    #     the roles, so source B carries `landlord` only (NO developer column -> developer
    #     tbd/absent), and match no longer hard-blocks the pair.
    def rl(city, park, area, src, developer=None, landlord=None):
        m = {"city": city, "park": park, "warehouseArea": area,
             "__meta": {"source_file": src, "source_type": "xlsx", "prov": {}}}
        if developer is not None:
            m["developer"] = developer
        if landlord is not None:
            m["landlord"] = landlord
        return m
    # source A: developer=Canmoor (no landlord). source B: landlord=NFU Mutual, NO
    # developer column -> developer absent (the post-fix shape of a landlord-only tracker)
    a_ll = rl("Corby", "Raven Park", 50000, "tracker.xlsx", developer="Canmoor")
    b_ll = rl("Corby", "Raven Park", 50000, "brochure.pdf", landlord="NFU Mutual")
    # one side's developer is absent (not a disagreement) -> NOT forbidden; same park/city
    # -> the pair is auto (an identical-key cross-source pair: dev absent on B, no conflict)
    check(match.pair_class(a_ll, b_ll) != "forbidden",
          "MATCH-g: developer (A) + landlord (B) for one building is NOT forbidden")
    clusters = match.dedupe([a_ll, b_ll], {match.pair_id(a_ll, b_ll): "same"})
    check(len(clusters) == 1,
          "MATCH-g: with a 'same' decision the developer/landlord pair dedupes to ONE cluster")
    if clusters and len(clusters) == 1:
        mrec, _mprov, _mc = merge.merge_cluster(clusters[0])
        check(mrec.get("developer") == "Canmoor" and mrec.get("landlord") == "NFU Mutual",
              "MATCH-g: the merged record carries BOTH developer=Canmoor AND landlord=NFU Mutual")
    # CLEAN OFFLINE case: both sources name developer=Canmoor (agreement), one ALSO carries
    # landlord=NFU -> a deterministic AUTO merge (no decision needed) -> one card, both fields
    a_off = rl("Corby", "Raven Park", 50000, "tracker.xlsx", developer="Canmoor", landlord="NFU Mutual")
    b_off = rl("Corby", "Raven Park", 50000, "brochure.pdf", developer="Canmoor")
    check(match.pair_class(a_off, b_off) == "auto",
          "MATCH-g: matching developers (one with a landlord) AUTO-merge offline (no decision needed)")
    off_clusters = match.dedupe([a_off, b_off])  # offline: no decisions file at all
    check(len(off_clusters) == 1,
          "MATCH-g: the agreeing-developer pair dedupes to ONE cluster offline (auto)")
    if off_clusters and len(off_clusters) == 1:
        orec, _op, _oc = merge.merge_cluster(off_clusters[0])
        check(orec.get("developer") == "Canmoor" and orec.get("landlord") == "NFU Mutual",
              "MATCH-g: the offline auto-merged card carries BOTH developer=Canmoor AND landlord=NFU Mutual")

    # (f) FULL-SPINE: a genuine grey pair -> exit 10 + work/match_candidates.json,
    #     then a hand-written match_decisions.json -> the re-run reaches exit 0 with the
    #     expected cluster count (mirror the photo-match exit-9 + TRK-f exit-3 tests).
    from openpyxl import Workbook
    with tempfile.TemporaryDirectory() as tds:
        tds = Path(tds)
        folder = tds / "inputs"; folder.mkdir()
        work = tds / "work"; work.mkdir()
        # two SOURCES (a tracker + a CSV) describing two same-city schemes whose parks
        # are neither subset nor >= 88 fuzzy -> a grey cross-source pair per scheme.
        wb = Workbook(); ws = wb.active
        ws.append(["Marketing Name", "Town", "Status", "Developer",
                   "Size (sq ft)", "Current quoting rent (£ per sq ft)"])
        ws.append(["Apollo Court", "Corby", "Available", "Prologis", 50000, 7])
        ws.append(["Mercury House", "Corby", "Available", "Prologis", 52000, 7])
        # NB: a UNIQUE filename - merge.py memoises _SRC_RESOLVE by bare name across the
        # whole process, so reusing a name another test uses (e.g. "tracker.xlsx") would
        # serve a stale (deleted-temp-dir) path to the later test.
        wb.save(folder / "match10_tracker.xlsx")
        # a DISTINCT third scheme in a SECOND source: same city + developer (so the
        # recall pre-filter fires) but a disjoint park name and no >15% size conflict
        # -> a genuine cross-source GREY pair against each tracker row.
        (folder / "match10_extra.csv").write_text(
            "Marketing Name,Town,Status,Developer,Size (sq ft),Current quoting rent (£ per sq ft)\n"
            "Neptune Gate,Corby,Available,Prologis,51000,7\n", encoding="utf-8")
        # first pass: a tracker is OFFERED interpretation (exit 3) - SKIP it to reach matching
        rc0 = _run_spine(folder, work)
        if rc0 == 3:
            man0 = json.loads((work / "vision" / "manifest.json").read_text(encoding="utf-8"))
            for j in (man0.get("jobs") or []):
                (work / "extract").mkdir(parents=True, exist_ok=True)
                (work / "extract" / (Path(j["output"]).stem + ".SKIP")).write_text("", encoding="utf-8")
            rc0 = _run_spine(folder, work)
        mc = work / "match_candidates.json"
        check(rc0 == 10 and mc.exists(),
              f"MATCH-f: a genuine cross-source grey pair exits 10 with match_candidates.json (got {rc0!r})")
        if mc.exists():
            cand = json.loads(mc.read_text(encoding="utf-8"))
            pairs = cand.get("pairs", [])
            check(bool(pairs) and all(p.get("pair_id") and p.get("a") and p.get("b") for p in pairs)
                  and cand.get("output") == "work/match_decisions.json" and "instructions" in cand,
                  "MATCH-f: match_candidates.json carries pair_id + both full records + output + instructions")
            # hand-write a decisions file: every pair 'different' (the honest default)
            decisions = {p["pair_id"]: {"verdict": "different", "reason": "test: distinct schemes"}
                         for p in pairs}
            (work / "match_decisions.json").write_text(json.dumps(decisions), encoding="utf-8")
            rc1 = _run_spine(folder, work)
            check(rc1 not in (2, 3, 10) and not isinstance(rc1, str),
                  f"MATCH-f: with match_decisions.json the re-run passes matching (exit {rc1!r}, not 2/3/10/crash)")
            can = work / "canonical.json"
            if can.exists():
                props = json.loads(can.read_text(encoding="utf-8")).get("properties", [])
                # 'different' for every grey pair -> all 3 distinct schemes ship as separate cards
                names = {str(p.get("park", "")).lower() for p in props}
                check(len(props) == 3 and any("apollo" in n for n in names)
                      and any("mercury" in n for n in names) and any("neptune" in n for n in names),
                      f"MATCH-f: 'different' verdicts keep the 3 distinct schemes as separate cards (got {len(props)})")
            # flip ONE grey pair to 'same' (Neptune == Apollo): the merge folds them into
            # one card, proving the verdict drives clustering both ways. The blocker still
            # prevents the same-source Apollo/Mercury phases from collapsing.
            ap_id = next((p["pair_id"] for p in pairs
                          if "apollo" in str(p["a"]).lower() and "neptune" in str(p["b"]).lower()
                          or "apollo" in str(p["b"]).lower() and "neptune" in str(p["a"]).lower()), None)
            if ap_id:
                decisions[ap_id] = {"verdict": "same", "reason": "test: one scheme, two sources"}
                (work / "match_decisions.json").write_text(json.dumps(decisions), encoding="utf-8")
                rc2 = _run_spine(folder, work)
                can2 = work / "canonical.json"
                if rc2 == 0 and can2.exists():
                    props2 = json.loads(can2.read_text(encoding="utf-8")).get("properties", [])
                    check(len(props2) == 2,
                          f"MATCH-f: a 'same' verdict MERGES the grey pair (3 schemes -> 2 cards, got {len(props2)})")


def verifier_cases() -> None:
    """P1 C+B INDEPENDENT SEMANTIC VERIFIER: the two highest-risk LLM judgements (the
    tracker column->field MAPPING/BASIS and the grey-zone MATCH verdict) are checked by a
    SECOND, blind re-derivation + a deterministic Python DIFF, ADVISORY to the Gaps Report
    and NEVER auto-rejecting the author's decision. These cases hand-write the verifier
    artefacts (*_mapcheck.json / match_verify.json) as the sub-agent stand-in, mirroring
    the TRK-* / MATCH-f harness, and prove: a disagreement surfaces but does not flip the
    decision (MAP-DIFF / MATCH-DIFF), agreement is silent (MAP-AGREE / MATCH-AGREE), the
    full cycle is byte-deterministic with the artefacts cached (DETERMINISM), and the
    artefacts ABSENT is byte-identical to today (BACKWARD-COMPAT)."""
    print("independent semantic verifier (C+B):")
    from openpyxl import Workbook
    import hashlib
    from build_dashboard import render as _render_raw
    import deliver as DLV

    def _html(data):  # render() returns (html, info) - we compare the html bytes only
        return _render_raw(data)[0]

    def _corby_sheet():
        wb = Workbook(); ws = wb.active
        ws.append(["Building ID", "Marketing Name", "Town", "Latitude, Longitude",
                   "Construction status", "Status", "Landlord", "Developer",
                   "Size (sq ft)", "Size Unit", "Current quoting rent (£ per sq ft)",
                   "Historical quoting rent (£ per sq ft)", "Eaves (m)",
                   "Eaves 10m or above?", "Office content (sq ft)",
                   "Floor loading (kN/sq m)", "No. of dock level doors",
                   "Site area (acres)", "Power (KVA)"])
        ws.append([10179, "EVO 169", "Corby", "52.50304981, -0.650581854",
                   "Built", "Available", "EVO", "",
                   172867, "GIA", 8.5, 0, 15, "Yes", 13576, 50, 15, 5.2, 800])
        td = Path(tempfile.mkdtemp()); f = td / "b.xlsx"; wb.save(f); return f

    # the author (primary) map - the correct one used by the parse
    primary = {"columns": [
        {"index": 1, "field": "park"}, {"index": 2, "field": "city"},
        {"index": 3, "field": "latlng"}, {"index": 5, "field": "status"},
        {"index": 7, "field": "developer"},
        {"index": 8, "field": "warehouseArea", "areaUnit": "sq ft", "basis": "GIA"},
        {"index": 9, "field": None, "role": "size_basis"},
        {"index": 10, "field": "warehouseRentVal", "currency": "GBP",
         "perArea": "sq ft", "period": "annual"},
        {"index": 12, "field": "clearHeight"},
        {"index": 14, "field": "officeArea", "areaUnit": "sq ft"},
        {"index": 15, "field": "floorLoad"}, {"index": 16, "field": "loadingDocks"},
        {"index": 17, "field": "plotArea", "areaUnit": "acres"},
        {"index": 18, "field": "electricity"}]}

    # --- MAP-AGREE: an identical second map produces NO disagreement (no false positives) ---
    f = _corby_sheet()
    res_agree = X.detect_and_extract(f, "Corby", "GB",
                                     column_map=primary, column_map_verify=primary)
    hr_agree = res_agree["header_report"][0]
    check("semantic_disagreements" not in hr_agree,
          "MAP-AGREE: an identical second map produces NO semantic_disagreement (no false positive)")
    # the records are byte-identical to the no-verify parse (the verifier is diff-only)
    res_noverify = X.detect_and_extract(f, "Corby", "GB", column_map=primary)
    check(json.dumps(res_agree["records"], sort_keys=True)
          == json.dumps(res_noverify["records"], sort_keys=True),
          "MAP-AGREE: the verify map never touches the parse (records byte-identical with/without it)")

    # --- MAP-DIFF: a deliberately-wrong second map (basis GIA->warehouse, plus a field
    #     swap) surfaces a semantic_disagreement; the PRIMARY map still drives the parse
    #     (warehouseArea stays the GIA-office 159,291, NOT rejected) ---
    wrong = json.loads(json.dumps(primary))  # deep copy
    for col in wrong["columns"]:
        if col.get("field") == "warehouseArea":
            col["basis"] = "warehouse"          # basis disagreement: GIA vs warehouse
        if col.get("field") == "officeArea":
            col["field"] = "plotArea"           # field disagreement on col 14
    res_diff = X.detect_and_extract(f, "Corby", "GB",
                                    column_map=primary, column_map_verify=wrong)
    hr_diff = res_diff["header_report"][0]
    dis = hr_diff.get("semantic_disagreements") or []
    keys = {(d["index"], d["key"]) for d in dis}
    check(bool(dis) and (8, "basis") in keys and (14, "field") in keys,
          "MAP-DIFF: a wrong second map surfaces the basis + field disagreements (sorted)")
    check(res_diff["records"][0].get("warehouseArea") == 159291,
          "MAP-DIFF: the PRIMARY map STILL drives the parse (warehouseArea 159,291, not rejected)")
    # the diff is sorted by (index, key) -> byte-stable
    check(dis == sorted(dis, key=lambda d: (d["index"], d["key"])),
          "MAP-DIFF: the disagreement list is sorted by (index, key) (byte-stable)")

    # --- MAP-DIFF full-spine: the manifest carries the tracker_verify job; a hand-written
    #     mapcheck.json disagreeing on basis surfaces a yield_report line + the Gaps line,
    #     and the primary map still parses (159,291) ---
    with tempfile.TemporaryDirectory() as tds:
        tds = Path(tds)
        folder = tds / "inputs"; folder.mkdir()
        work = tds / "work"; work.mkdir()
        import shutil
        shutil.copy(_corby_sheet(), folder / "corby.xlsx")
        rc0 = _run_spine(folder, work)
        check(rc0 == 3, f"MAP-DIFF: a raw tracker exits 3 (mapping + verify offered), got {rc0!r}")
        man = json.loads((work / "vision" / "manifest.json").read_text(encoding="utf-8"))
        jobs = man.get("jobs") or []
        tjob = next((j for j in jobs if j.get("kind") == "tracker"), None)
        vjob = next((j for j in jobs if j.get("kind") == "tracker_verify"), None)
        check(tjob and vjob and tjob.get("input_hash") == vjob.get("input_hash")
              and vjob.get("output", "").endswith("_mapcheck.json")
              and "tracker_verify_instructions" in man,
              "MAP-DIFF: the manifest carries a tracker_verify job (same input_hash, mapcheck output) + instructions")
        if tjob and vjob:
            ex = work / "extract"; ex.mkdir(parents=True, exist_ok=True)
            # author map (correct)
            (ex / Path(tjob["output"]).name).write_text(json.dumps(
                {"input_hash": tjob["input_hash"], "schema_version": 1, "map": primary}),
                encoding="utf-8")
            # blind verify map (disagrees on the warehouseArea basis)
            (ex / Path(vjob["output"]).name).write_text(json.dumps(
                {"input_hash": vjob["input_hash"], "schema_version": 1, "map": wrong}),
                encoding="utf-8")
            rc1 = _run_spine(folder, work)
            check(rc1 not in (2, 3) and not isinstance(rc1, str),
                  f"MAP-DIFF: with both maps the re-run passes extraction (exit {rc1!r}, not 2/3/crash)")
            yr = work / "yield_report.md"
            ytext = yr.read_text(encoding="utf-8") if yr.exists() else ""
            check("independent column-mapping passes DISAGREE" in ytext and "basis" in ytext,
                  "MAP-DIFF: the disagreement is surfaced in yield_report.md (-> Gaps Report)")
            can = work / "canonical.json"
            if can.exists():
                props = json.loads(can.read_text(encoding="utf-8")).get("properties", [])
                check(props and props[0].get("warehouseArea") == 159291,
                      "MAP-DIFF: full-spine the PRIMARY map still parses 159,291 (verifier did NOT reject it)")
            # the disagreement also reaches the delivered Gaps Report
            try:
                data_md = json.loads(can.read_text(encoding="utf-8"))
                gaps = DLV.gaps_report(data_md, "corby", work_dir=work)
                check("independent column-mapping passes DISAGREE" in gaps,
                      "MAP-DIFF: the map disagreement reaches the delivered Gaps Report")
            except Exception as e:
                check(False, f"MAP-DIFF: gaps_report raised ({e})")

    # ---- the grey-MATCH verifier ----
    def r2(city, dev, park, area, src):
        return {"city": city, "developer": dev, "park": park, "warehouseArea": area,
                "__meta": {"source_file": src}}

    def _spine_to_match(folder, work):
        """Run the spine; SKIP any tracker mapping so it reaches the exit-10 match step."""
        rc = _run_spine(folder, work)
        if rc == 3:
            man = json.loads((work / "vision" / "manifest.json").read_text(encoding="utf-8"))
            (work / "extract").mkdir(parents=True, exist_ok=True)
            for j in (man.get("jobs") or []):
                (work / "extract" / (Path(j["output"]).stem + ".SKIP")).write_text("", encoding="utf-8")
            rc = _run_spine(folder, work)
        return rc

    def _match_fixture():
        td = Path(tempfile.mkdtemp())
        folder = td / "inputs"; folder.mkdir()
        work = td / "work"; work.mkdir()
        wb = Workbook(); ws = wb.active
        ws.append(["Marketing Name", "Town", "Status", "Developer",
                   "Size (sq ft)", "Current quoting rent (£ per sq ft)"])
        ws.append(["Apollo Court", "Corby", "Available", "Prologis", 50000, 7])
        ws.append(["Mercury House", "Corby", "Available", "Prologis", 52000, 7])
        wb.save(folder / "vfy_tracker.xlsx")
        (folder / "vfy_extra.csv").write_text(
            "Marketing Name,Town,Status,Developer,Size (sq ft),Current quoting rent (£ per sq ft)\n"
            "Neptune Gate,Corby,Available,Prologis,51000,7\n", encoding="utf-8")
        return folder, work

    # --- MATCH-DIFF + MATCH-AGREE: emit candidates, hand-write author 'different' verdicts
    #     AND a verifier file that DISAGREES on one pair; the disagreement lands in
    #     meta.conflicts -> Gaps; clustering is STILL driven by the author ('different' ->
    #     3 cards). Then an AGREEING verifier produces no conflict line. ---
    folder, work = _match_fixture()
    rc0 = _spine_to_match(folder, work)
    mc = work / "match_candidates.json"
    check(rc0 == 10 and mc.exists(),
          f"MATCH-DIFF: a genuine grey pair exits 10 with match_candidates.json (got {rc0!r})")
    if mc.exists():
        cand = json.loads(mc.read_text(encoding="utf-8"))
        vps = cand.get("verify_pairs") or []
        pairs = cand.get("pairs") or []
        check(bool(vps) and {p["pair_id"] for p in vps} == {p["pair_id"] for p in pairs}
              and cand.get("verify_output") == "work/match_verify.json"
              and "verify_instructions" in cand,
              "MATCH-DIFF: match_candidates.json carries verify_pairs (same ids) + verify_output + instructions")
        # author: every pair 'different' (the honest default)
        author = {p["pair_id"]: {"verdict": "different", "reason": "author: distinct"}
                  for p in pairs}
        (work / "match_decisions.json").write_text(json.dumps(author), encoding="utf-8")
        # verifier: DISAGREE on exactly one pair (call it 'same')
        flip_id = pairs[0]["pair_id"]
        verifier = dict(author)
        verifier[flip_id] = {"verdict": "same", "reason": "verifier: looks like one property"}
        (work / "match_verify.json").write_text(json.dumps(verifier), encoding="utf-8")
        rc1 = _run_spine(folder, work)
        check(rc1 not in (2, 3, 10) and not isinstance(rc1, str),
              f"MATCH-DIFF: the re-run passes matching with both files (exit {rc1!r})")
        can = work / "canonical.json"
        if can.exists():
            data = json.loads(can.read_text(encoding="utf-8"))
            conflicts = data.get("meta", {}).get("conflicts", [])
            check(any("match disagreement" in c for c in conflicts),
                  "MATCH-DIFF: the verifier disagreement lands in meta.conflicts (advisory)")
            props = data.get("properties", [])
            check(len(props) == 3,
                  f"MATCH-DIFF: clustering is STILL driven by the AUTHOR ('different' -> 3 cards, got {len(props)})")
            # and it reaches the delivered Gaps 'Source conflicts' section
            gaps = DLV.gaps_report(data, "vfy", work_dir=work)
            check("match disagreement" in gaps,
                  "MATCH-DIFF: the match disagreement reaches the Gaps 'Source conflicts' section")

        # --- MATCH-AGREE: a verifier that AGREES with the author produces NO conflict line ---
        (work / "match_verify.json").write_text(json.dumps(author), encoding="utf-8")
        rc2 = _run_spine(folder, work)
        can2 = work / "canonical.json"
        if rc2 == 0 and can2.exists():
            conflicts2 = json.loads(can2.read_text(encoding="utf-8")).get("meta", {}).get("conflicts", [])
            check(not any("match disagreement" in c for c in conflicts2),
                  "MATCH-AGREE: an agreeing verifier produces NO match-disagreement line")

    # --- DETERMINISM full-cycle: run the spine TWICE with mapcheck + match_verify present;
    #     the built.html bytes are identical across runs (no live re-dispatch on run 2) ---
    with tempfile.TemporaryDirectory() as tdd:
        tdd = Path(tdd)
        folder = tdd / "inputs"; folder.mkdir()
        work = tdd / "work"; work.mkdir()
        import shutil
        shutil.copy(_corby_sheet(), folder / "det.xlsx")
        rc0 = _run_spine(folder, work)
        if rc0 == 3:
            man = json.loads((work / "vision" / "manifest.json").read_text(encoding="utf-8"))
            ex = work / "extract"; ex.mkdir(parents=True, exist_ok=True)
            for j in (man.get("jobs") or []):
                mp = primary if j.get("kind") == "tracker" else wrong  # verify disagrees
                (ex / Path(j["output"]).name).write_text(json.dumps(
                    {"input_hash": j["input_hash"], "schema_version": 1, "map": mp}),
                    encoding="utf-8")
            _run_spine(folder, work)
        # build once via gate-free render to compare bytes deterministically
        can = work / "canonical.json"
        if can.exists():
            data = json.loads(can.read_text(encoding="utf-8"))
            h1 = hashlib.sha256(_html(data).encode("utf-8")).hexdigest()
            # second full cycle (resume): the diff is recomputed from cached files, no
            # live re-dispatch, and the canonical/built bytes must be identical
            _run_spine(folder, work)
            data2 = json.loads(can.read_text(encoding="utf-8"))
            h2 = hashlib.sha256(_html(data2).encode("utf-8")).hexdigest()
            check(h1 == h2,
                  "DETERMINISM: two full cycles with mapcheck + verify artefacts cached -> byte-identical built.html")
            check(json.dumps(data, sort_keys=True) == json.dumps(data2, sort_keys=True),
                  "DETERMINISM: canonical.json byte-identical across the two cached runs (no live re-dispatch)")

    # --- BACKWARD-COMPAT: with NO verifier artefacts the diff is EMPTY -> canonical +
    #     built.html byte-identical to a run that never had them (offline-fallback) ---
    def _spine_canonical(extra_verify: bool):
        td = Path(tempfile.mkdtemp())
        folder = td / "inputs"; folder.mkdir()
        work = td / "work"; work.mkdir()
        import shutil
        shutil.copy(_corby_sheet(), folder / "bc.xlsx")
        rc = _run_spine(folder, work)
        if rc == 3:
            man = json.loads((work / "vision" / "manifest.json").read_text(encoding="utf-8"))
            ex = work / "extract"; ex.mkdir(parents=True, exist_ok=True)
            for j in (man.get("jobs") or []):
                if j.get("kind") == "tracker":
                    (ex / Path(j["output"]).name).write_text(json.dumps(
                        {"input_hash": j["input_hash"], "schema_version": 1, "map": primary}),
                        encoding="utf-8")
                elif extra_verify:
                    (ex / Path(j["output"]).name).write_text(json.dumps(
                        {"input_hash": j["input_hash"], "schema_version": 1, "map": primary}),
                        encoding="utf-8")
                else:
                    # decline the verify job so NO mapcheck artefact exists (today's behaviour)
                    (ex / (Path(j["output"]).stem + ".SKIP")).write_text("", encoding="utf-8")
            _run_spine(folder, work)
        can = work / "canonical.json"
        return json.loads(can.read_text(encoding="utf-8")) if can.exists() else None

    base = _spine_canonical(extra_verify=False)   # no mapcheck artefact at all
    agree = _spine_canonical(extra_verify=True)    # an AGREEING mapcheck present
    if base is not None and agree is not None:
        check(hashlib.sha256(_html(base).encode("utf-8")).hexdigest()
              == hashlib.sha256(_html(agree).encode("utf-8")).hexdigest(),
              "BACKWARD-COMPAT: an agreeing/absent verifier yields byte-identical built.html (empty diff)")
        check(json.dumps(base, sort_keys=True) == json.dumps(agree, sort_keys=True),
              "BACKWARD-COMPAT: canonical.json byte-identical with no/agreeing verifier (advisory never enters data)")


def cache_cases() -> None:
    print("images / run helpers:")
    try:
        import fitz
        doc = fitz.open()
        pg = doc.new_page()
        pg.insert_text((40, 60), "hello", fontsize=12)
        with tempfile.TemporaryDirectory() as td:
            f = Path(td) / "one.pdf"
            doc.save(f); doc.close()
            cache = Path(td) / "cache"
            u1 = IMG.hero_for_pdf_page(f, 0, 40, cache_dir=cache)
            IMG.close_doc_cache()
            n_cached = len(list(cache.glob("*.uri")))
            u2 = IMG.hero_for_pdf_page(f, 0, 40, cache_dir=cache)
            IMG.close_doc_cache()
        check(u1 == u2 and n_cached == 2,  # one .uri per kind: hero + plan
              "W1: hero cache hit returns identical bytes")
    except Exception as e:
        check(False, f"W1: image cache test errored ({e})")


def prewarm_resume_cases() -> None:
    print("image pre-warm + per-page resume (convergence):")
    import io as _io, random as _r, tempfile as _tf
    try:
        import merge as _M, fitz as _fz
        from PIL import Image as _Img
    except Exception as e:
        check(False, f"prewarm: setup import failed ({e})"); return

    def _noise(seed, w=340, h=240):
        _r.seed(seed); im = _Img.new("RGB", (w, h))
        im.putdata([(_r.randint(0, 255), _r.randint(0, 255), _r.randint(0, 255)) for _ in range(w * h)])
        b = _io.BytesIO(); im.save(b, "JPEG", quality=80); return b.getvalue()
    _logo = _io.BytesIO(); _Img.new("RGB", (110, 110), (20, 80, 160)).save(_logo, "PNG"); _logo = _logo.getvalue()

    def _reset():
        IMG._PLACED_CACHE.clear(); IMG._CROPS_CACHE.clear(); IMG.close_doc_cache()

    with _tf.TemporaryDirectory(ignore_cleanup_errors=True) as td:
        td = Path(td)
        doc = _fz.open()
        for p in range(4):
            pg = doc.new_page(width=600, height=800)
            pg.insert_image(_fz.Rect(60, 60, 460, 360), stream=_noise(p))   # unique photo
            pg.insert_image(_fz.Rect(60, 700, 170, 790), stream=_logo)       # repeated -> boilerplate
        f = td / "deck.pdf"; doc.save(f); doc.close()

        # _deck_photo_index: resume after a mid-deck kill == full build (byte-identical)
        c1 = td / ".c1"
        idx_full = IMG._deck_photo_index(f, 40, c1); _reset()
        for j in sorted(c1.glob("*.json"))[::2]:
            j.unlink()                       # drop half the per-page photo caches
        for u in c1.glob("*.uri"):
            u.unlink()                       # drop the whole-deck index so it reassembles
        idx_res = IMG._deck_photo_index(f, 40, c1); _reset()
        check(idx_full == idx_res and len(idx_full) >= 1,
              "prewarm: _deck_photo_index resumes to a byte-identical index after a mid-deck kill")

        # _placed_layout: resume == full, and the repeated logo stays boilerplate
        c2 = td / ".c2"
        lay_full = IMG._placed_layout(f, c2); _reset()
        for pj in c2.glob("*.placed.json"):
            pj.unlink()                      # drop the whole-deck layout
        half = sorted(p for p in c2.glob("*.json") if not p.name.endswith(".placed.json"))[::2]
        for j in half:
            j.unlink()                       # drop half the per-page geometry caches
        lay_res = IMG._placed_layout(f, c2); _reset()
        check(lay_full == lay_res, "prewarm: _placed_layout resumes to identical geometry after a mid-deck kill")
        check(sum(1 for pg in lay_full["pages"] for e in pg if e.get("boiler")) >= 3,
              "prewarm: the repeated logo is still flagged boilerplate after a resume")

        # prewarm_images (serial path, workers=1) is a PURE ACCELERATOR: a cold merge
        # harvest and a prewarmed one yield byte-identical hero + gallery
        recs = [{"city": "M", "park": f"P{p}",
                 "__meta": {"source_file": "deck.pdf", "source_type": "pdf", "page_no": p}}
                for p in (0, 1, 2)]
        cold = td / ".cold"; cold.mkdir(); _reset()
        ph_c, _, _, _, _, gal_c = _M.attach_media(recs, td, IMG.DEFAULT_BUDGET_KB, image_cache=cold); _reset()
        warm = td / ".warm"; warm.mkdir()
        done, total = _M.prewarm_images(recs, td, warm, IMG.DEFAULT_BUDGET_KB, seconds=60, workers=1); _reset()
        ph_w, _, _, _, _, gal_w = _M.attach_media(recs, td, IMG.DEFAULT_BUDGET_KB, image_cache=warm); _reset()
        check(done == total and total > 0, f"prewarm: serial prewarm warms every unit ({done}/{total})")
        check(ph_c == ph_w and gal_c == gal_w,
              "prewarm: a prewarmed cache yields byte-identical hero+gallery vs a cold merge harvest")
        IMG.close_doc_cache()

    import run
    check(not run._filled("a consultar") and not run._filled("auf anfrage")
          and run._filled(40000) and run._filled("Pilsen"),
          "E8: run.py thin-parse probe uses the multilingual unknown list")


def enrich_cases() -> None:
    print("enrich (offline behaviour):")
    import enrich

    # DATASET path (the primary): genuine nearest from the bundled complete set,
    # fully offline, no Overpass at all
    enrich._DATASET = None  # force a fresh load of the real bundled dataset
    with tempfile.TemporaryDirectory() as td:
        enrich.CACHE_DIR = Path(td)
        canonical = {"properties": [{"id": 1, "city": "Azuqueca", "lat": 40.573, "lng": -3.245}],
                     "pois": [], "meta": {}}
        gaps0: list = []
        n, live = enrich.attach_pois(canonical, gaps0)
        by_type = {p["type"]: p for p in canonical["pois"]}
        check(live and {"air", "port", "rail"} <= set(by_type)
              and all("CBRE POI dataset" in by_type[t]["note"] for t in ("air", "port", "rail")),
              "dataset: nearest air/port/rail attached offline from the bundled set")
        rail_km = enrich._haversine_km(40.573, -3.245,
                                       by_type["rail"]["lat"], by_type["rail"]["lng"])
        check(rail_km < 60, f"dataset: rail terminal is LOCAL ({rail_km} km), not a far stand-in")

    # the remaining cases exercise the OSM fallback path - reached only when NONE of the three
    # complete datasets (poi_dataset air/port/rail, borders_dataset, cities_major_dataset) is
    # present, so pin all three absent (not just poi_dataset).
    enrich._DATASET = enrich._BORDERS = enrich._CITY_DATASET = False
    with tempfile.TemporaryDirectory() as td:
        enrich.CACHE_DIR = Path(td)
        seeded = [{"name": "Port of Hamburg", "type": "port", "lat": 53.54, "lng": 9.97}]
        canonical = {"properties": [{"id": 1, "lat": 49.7, "lng": 13.3}],
                     "pois": list(seeded), "meta": {}}
        original = enrich.nearest_pois_for
        enrich.nearest_pois_for = lambda lat, lng: (_ for _ in ()).throw(RuntimeError("offline"))
        try:
            gaps: list = []
            enrich.attach_pois(canonical, gaps)
        finally:
            enrich.nearest_pois_for = original
        check(canonical["pois"] == seeded,
              "M5: Overpass-unreachable keeps the merge-seeded POIs (Cowork default path)")
        check(any("kept" in g for g in gaps), "M5: the kept-seeds fallback is an honest gaps line")

    # beyond the calibrated 180km OSM scan, missing types come from the curated
    # major-facilities library, capped and LABELLED - never silently
    with tempfile.TemporaryDirectory() as td:
        enrich.CACHE_DIR = Path(td)
        ck = f"{round(40.4, 3)},{round(-3.7, 3)}"
        (Path(td) / "poi_osm_cache.json").write_text(json.dumps({
            ck: {"rail": {"name": "Madrid-Abronigal", "type": "rail",
                          "lat": 40.392, "lng": -3.674, "km": 1.0}}}), encoding="utf-8")
        canonical2 = {"properties": [{"id": 1, "lat": 40.4, "lng": -3.7}],
                      "pois": [], "meta": {}}
        gaps2: list = []
        n, live = enrich.attach_pois(canonical2, gaps2)
        notes = {p["type"]: p["note"] for p in canonical2["pois"]}
        check(live and "rail" in notes and "OSM" in notes["rail"],
              "supplement: the OSM-found type keeps its OSM note")
        check("port" in notes and "library" in notes["port"],
              "supplement: a beyond-scan port comes from the library, LABELLED as such")
    enrich._DATASET = enrich._BORDERS = enrich._CITY_DATASET = None  # restore the real datasets


def intake_cases() -> None:
    print("intake (S6 - no silent input loss):")
    import intake
    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        # DISTINCT bytes per file (real brochures differ; INTAKE-001 dedups byte-identical
        # files, so identical placeholder bytes would now collapse them all)
        (root / "Options - Pilsen.pdf").write_bytes(b"a")
        (root / "New stock - Pilsen.pdf").write_bytes(b"b")
        (root / "Normal Options - Madrid - FINAL.pdf").write_bytes(b"c")
        (root / "Options - St. Polten.pdf").write_bytes(b"d")
        (root / "CTPark Brno-South.pdf").write_bytes(b"e")
        sub = root / "extra"; sub.mkdir()
        (sub / "Opciones — Valencia.pdf").write_bytes(b"f")
        (sub / "tracker.xlsx").write_bytes(b"g")
        inv = intake.discover(root)
    cl = inv["clusters"]
    check(sorted(cl.get("Pilsen", {}).get("pdfs", [])) ==
          ["New stock - Pilsen.pdf", "Options - Pilsen.pdf"],
          "S6: two brochures for one region BOTH kept (no silent overwrite)")
    check("Madrid" in cl, "S6: '... - Madrid - FINAL.pdf' clusters as Madrid, not FINAL")
    check("St. Polten" in cl, "S6: dotted city name clusters correctly")
    check("CTPark Brno-South" in cl, "S6: unspaced hyphenated name is never split")
    check(cl.get("Valencia", {}).get("pdfs") == ["extra/Opciones — Valencia.pdf"],
          "S6: subfolder brochure found (em-dash separator, relative path)")
    check(inv["subfolders"] == ["extra"] and inv["xlsx"] == ["extra/tracker.xlsx"],
          "S6: scanned subfolders are named; xlsx found recursively")
    # INTAKE-001: a byte-identical duplicate input is extracted ONCE (first in sorted order
    # kept), the copy recorded as a skipped duplicate - never a wasted extraction or a
    # phantom second card; distinct-content files are untouched.
    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        (root / "aaa.pdf").write_bytes(b"same-bytes")
        (root / "zzz.pdf").write_bytes(b"same-bytes")     # an exact copy of aaa.pdf
        (root / "different.pdf").write_bytes(b"other")
        inv = intake.discover(root)
    dups = inv.get("skipped_duplicates", [])
    kept = sorted(r for c in inv["clusters"].values() for r in c.get("pdfs", []))
    check(len(dups) == 1 and dups[0]["file"] == "zzz.pdf" and dups[0]["duplicate_of"] == "aaa.pdf",
          "INTAKE-001: a byte-identical duplicate is skipped + recorded (first in sorted order kept)")
    check("aaa.pdf" in kept and "different.pdf" in kept and "zzz.pdf" not in kept,
          "INTAKE-001: original + the distinct file kept; the exact copy is dropped from extraction")

    # S7 (#7 filename->city LLM clustering) - the confidence signal + the input-hashed
    # cache override + the structural verifier + the offline fallback.
    # (a) the confidence signal: 'Options-Oporto' (no spaced dash, Oporto unknown) is LOW;
    #     a clean spaced-dash split is HIGH.
    _, _, conf_low = intake.infer_cluster("Options-Oporto.pdf", {})
    _, _, conf_high = intake.infer_cluster("Brochure - Madrid.pdf", {})
    check(conf_low == "low",
          "S7a: 'Options-Oporto' (no spaced dash, unknown tail) flags LOW confidence")
    check(conf_high == "high",
          "S7a: a clean spaced-dash 'Brochure - Madrid' flags HIGH confidence")
    # the low cluster is surfaced on inventory.json with its raw stems for the orchestrator
    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        (root / "Options-Oporto.pdf").write_bytes(b"oporto")
        inv = intake.discover(root)
    opc = inv["clusters"].get("Options-Oporto", {})
    check(opc.get("confidence") == "low" and opc.get("stems") == ["Options-Oporto"],
          "S7a: a low-confidence cluster carries confidence:'low' + its raw stems")

    # (b) intake.discover honours a matching work/intake_clusters.json -> the key is 'Oporto'
    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        (root / "Options-Oporto.pdf").write_bytes(b"oporto")
        ih = intake._brochure_input_hash(["Options-Oporto.pdf"])
        cache = {"input_hash": ih, "schema_version": 1,
                 "labels": [{"stem": "Options-Oporto", "region": "Oporto", "country": "PT"}]}
        inv = intake.discover(root, cluster_cache=cache)
    cl = inv["clusters"]
    check("Oporto" in cl and "Options-Oporto" not in cl,
          "S7b: a matching cache overrides the regex - the cluster key is 'Oporto'")
    check(cl.get("Oporto", {}).get("country") == "PT"
          and cl["Oporto"].get("confidence") == "high"
          and cl["Oporto"].get("pdfs") == ["Options-Oporto.pdf"],
          "S7b: the applied label carries its country, is now high-confidence, keeps the brochure")

    # (c) the structural verifier: a MISMATCHED hash, an empty cache, a stem that is not a
    #     real file, or a _NOISE region -> the cache is ignored and the regex result stands.
    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        (root / "Options-Oporto.pdf").write_bytes(b"oporto")
        good = {"stem": "Options-Oporto", "region": "Oporto", "country": "PT"}
        bad_hash = {"input_hash": "deadbeef", "schema_version": 1, "labels": [good]}
        empty = {}
        ghost_stem = {"input_hash": intake._brochure_input_hash(["Options-Oporto.pdf"]),
                      "schema_version": 1,
                      "labels": [{"stem": "Not-A-Real-File", "region": "Oporto"}]}
        noise_region = {"input_hash": intake._brochure_input_hash(["Options-Oporto.pdf"]),
                        "schema_version": 1,
                        "labels": [{"stem": "Options-Oporto", "region": "FINAL"}]}
        invs = {name: intake.discover(root, cluster_cache=c)
                for name, c in (("bad_hash", bad_hash), ("empty", empty),
                                ("ghost_stem", ghost_stem), ("noise_region", noise_region))}
    for name, iv in invs.items():
        check("Options-Oporto" in iv["clusters"] and "Oporto" not in iv["clusters"],
              f"S7c: a {name} cache is rejected by the verifier - the regex result stands")

    # (d) with NO cache the existing S6 clustering is byte-identical (offline fallback).
    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        for fn, b in (("Options - Pilsen.pdf", b"a"), ("New stock - Pilsen.pdf", b"b"),
                      ("Normal Options - Madrid - FINAL.pdf", b"c"),
                      ("Options - St. Polten.pdf", b"d"), ("CTPark Brno-South.pdf", b"e")):
            (root / fn).write_bytes(b)
        inv_none = intake.discover(root, cluster_cache=None)
        inv_default = intake.discover(root)
    check(json.dumps(inv_none, sort_keys=True) == json.dumps(inv_default, sort_keys=True),
          "S7d: discover(cache=None) == discover() (the cache arg defaults to the regex)")
    keys = set(inv_none["clusters"])
    check({"Pilsen", "Madrid", "St. Polten", "CTPark Brno-South"} <= keys
          and "FINAL" not in keys and "Options-Oporto" not in keys,
          "S7d: with no cache the regex clustering is byte-identical to the S6 result")


def _noise_photo_jpeg(w=800, h=450) -> bytes:
    """A synthetic 'photo': colourful noise scores high on photographic_score."""
    import random
    from PIL import Image
    rnd = random.Random(7)
    img = Image.new("RGB", (64, 36))
    img.putdata([(rnd.randrange(256), rnd.randrange(256), rnd.randrange(256))
                 for _ in range(64 * 36)])
    img = img.resize((w, h))
    buf = __import__("io").BytesIO()
    img.save(buf, format="JPEG", quality=85)
    return buf.getvalue()


def _photo_page_pdf(td: Path) -> Path:
    """One-page PDF: text + an embedded JPEG photo placed at ~quarter page."""
    import fitz
    doc = fitz.open()
    pg = doc.new_page()  # 595x842pt
    y = 60
    for ln in ("City", "Pilsen", "Developer", "CTP", "Warehouse Area", "40 000 sq m"):
        pg.insert_text((40, y), ln, fontsize=12); y += 20
    pg.insert_image(fitz.Rect(60, 320, 360, 489), stream=_noise_photo_jpeg())
    f = td / "photo_page.pdf"
    doc.save(f); doc.close()
    return f


def tedi_cases() -> None:
    print("TEDi run fixes (hero tiers / shim tiers / merge normalisation / vision resume):")
    import normalize as N
    # Fix 4a: country names normalise to ISO at merge, never fail the schema gate
    check(N.country_iso("Spain") == "ES" and N.country_iso("España") == "ES"
          and N.country_iso("cz") == "CZ" and N.country_iso("Czech Republic") == "CZ"
          and N.country_iso("??") == "??",
          "Fix4: country names -> ISO ('Spain'/'España' -> ES; '??' untouched)")
    # Fix 4b: merge owns the rent pair - agent display strings normalise, never block
    r = merge.canonicalize({"warehouseRent": "3,75 €/m²/mes"})
    check(r.get("warehouseRentVal") == 45.0 and r.get("warehouseRent") == "€45 / sq m / year",
          "Fix4: '3,75 €/m²/mes' -> val 45 annual + regenerated display")
    r = merge.canonicalize({"warehouseRentVal": 60.0, "warehouseRent": "60 EUR/m2/año"})
    check(r.get("warehouseRent") == "€60 / sq m / year",
          "Fix4: display always regenerated from the numeric (pair gate cannot block)")
    r = merge.canonicalize({"country": "Spain"})
    check(r.get("country") == "ES", "Fix4: canonicalize maps the country")

    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        pdf = _photo_page_pdf(td)
        # Fix 1 (tier A): the embedded photo beats the whole-page raster outright
        uri = IMG.hero_for_pdf_page(pdf, 0, 60)
        import base64, io as _io
        from PIL import Image
        hero = Image.open(_io.BytesIO(base64.b64decode(uri.split(",", 1)[1])))
        ar = hero.width / hero.height
        check(1.5 < ar < 2.1, f"Fix1: hero is the embedded photo (ar={ar:.2f}), not the page render")
        # Fix 1 (tier B): bbox-crop finds the photo region from geometry alone
        crop = IMG.bbox_crop_hero(pdf, 0)
        check(crop is not None and 1.4 < crop.width / crop.height < 2.2
              and IMG.photographic_score(crop) >= IMG.MODEST_PHOTO,
              "Fix1: bbox-crop tier recovers the photo region without decoding streams")
        IMG.close_doc_cache()

        # Fix 1 (shim tiers): pdfplumber-ONLY backend still extracts the photo;
        # rendering raises a CLASSIFIED error instead of a silent wrong hero
        import importlib
        saved = sys.modules.pop("pypdfium2", None)
        sys.modules["pypdfium2"] = None  # force ImportError inside the shim
        try:
            import fitz_shim
            importlib.reload(fitz_shim)
            check(fitz_shim.ENGINE == "pdfplumber-only", "Fix1: shim detects the pdfplumber-only tier")
            doc = fitz_shim.open(pdf)
            check("Warehouse Area" in doc[0].get_text(), "Fix1: shim text works without a renderer")
            infos = doc[0].get_images(full=True)
            ok_img = False
            for info in infos:
                raw = doc.extract_image(info[0])["image"]
                try:
                    im = Image.open(_io.BytesIO(raw)); im.load()
                    ok_img = ok_img or (im.width >= 320 and IMG.photographic_score(im) >= IMG.MODEST_PHOTO)
                except Exception:
                    pass
            check(ok_img, "Fix1: shim decodes the embedded JPEG straight from the PDF stream")
            try:
                doc[0].get_pixmap(dpi=72)
                check(False, "Fix1: renderer-less get_pixmap raises a clear error")
            except RuntimeError:
                check(True, "Fix1: renderer-less get_pixmap raises a clear error")
            doc.close()
        finally:
            del sys.modules["pypdfium2"]
            if saved is not None:
                sys.modules["pypdfium2"] = saved
            import fitz_shim
            importlib.reload(fitz_shim)

        # Fix 2: vision_prep is per-page resumable (and defaults to 120 dpi)
        import vision_prep
        vis = td / "vision"
        ent1 = vision_prep.prepare(pdf, "R", "CZ", vis, force=True)
        png = Path(ent1["pages"][0]["image"])
        m1 = png.stat().st_mtime_ns
        ent2 = vision_prep.prepare(pdf, "R", "CZ", vis, force=True)
        check(png.stat().st_mtime_ns == m1 and len(ent2["pages"]) == len(ent1["pages"]),
              "Fix2: second prepare() reuses the rendered PNGs (resume, no re-render)")
        check(ent1["pages"][0]["page_no"] == 0, "Fix3: manifest page_no is the canonical 0-based index")

    # Site plans + map links (second field-test round)
    import fitz
    from PIL import ImageDraw

    def _plan_png() -> bytes:  # white paper + flat colour fills + line work = plan signature
        img = Image.new("RGB", (800, 500), (255, 255, 255))
        d = ImageDraw.Draw(img)
        d.rectangle([80, 60, 720, 440], outline=(40, 40, 40), width=4)
        d.rectangle([120, 100, 500, 400], fill=(120, 180, 230))
        d.rectangle([520, 100, 690, 250], fill=(180, 220, 160))
        for x in range(140, 480, 40):
            d.line([(x, 110), (x, 390)], fill=(60, 60, 60), width=2)
        buf = _io.BytesIO(); img.save(buf, format="PNG"); return buf.getvalue()

    def _map_png() -> bytes:  # pale beige tiles, NO true white = map signature
        img = Image.new("RGB", (800, 500), (242, 238, 230))
        d = ImageDraw.Draw(img)
        for y in range(40, 480, 60):
            d.line([(0, y), (800, y + 120)], fill=(255, 230, 150), width=8)
        buf = _io.BytesIO(); img.save(buf, format="PNG"); return buf.getvalue()

    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        doc = fitz.open()
        pg = doc.new_page()  # page 1: photo + plan + linked map
        pg.insert_image(fitz.Rect(60, 80, 360, 250), stream=_noise_photo_jpeg())
        pg.insert_image(fitz.Rect(60, 300, 360, 490), stream=_plan_png())
        pg.insert_image(fitz.Rect(60, 540, 360, 690), stream=_map_png())
        pg.insert_link({"kind": fitz.LINK_URI, "from": fitz.Rect(60, 695, 360, 715),
                        "uri": "https://maps.google.com/?q=40.5729170000,-3.2452220000"})
        pg.insert_text((400, 100), "City", fontsize=12)
        pg.insert_text((400, 120), "Pilsen", fontsize=12)
        pg.insert_text((400, 140), "Developer", fontsize=12)
        pg.insert_text((400, 160), "CTP", fontsize=12)
        pg2 = doc.new_page()  # page 2: PLAN-ONLY
        pg2.insert_image(fitz.Rect(60, 80, 460, 330), stream=_plan_png())
        pg3 = doc.new_page()  # page 3: MAP-ONLY (linked)
        pg3.insert_image(fitz.Rect(60, 80, 460, 330), stream=_map_png())
        pg3.insert_link({"kind": fitz.LINK_URI, "from": fitz.Rect(60, 335, 460, 355),
                         "uri": "https://maps.google.com/?q=41.0,-3.0"})
        f = td / "plans.pdf"
        doc.save(f); doc.close()

        hero, plan = IMG.page_hero_and_plan(f, 0, 60)
        from PIL import Image as _I
        h_img = _I.open(_io.BytesIO(base64.b64decode(hero.split(",", 1)[1])))
        check(1.5 < h_img.width / h_img.height < 2.1 and plan is not None and plan != hero,
              "Plans: photo is the hero, the site plan fills the plan slot")
        hero2, plan2 = IMG.page_hero_and_plan(f, 1, 60)
        check(hero2 is not None and hero2 == plan2,
              "Plans: a plan-only page promotes the plan to hero AND keeps the plan slot")
        hero3, plan3 = IMG.page_hero_and_plan(f, 2, 60)
        check(plan3 is None, "Plans: a hyperlinked location map never fills the plan slot")
        IMG.close_doc_cache()

        recs = P.extract(f, "R", "CZ")
        r0 = recs[0] if recs else {}
        check(r0.get("lat") == 40.572917 and r0.get("lng") == -3.245222
              and "map link" in r0.get("__meta", {}).get("prov", {}).get("lat", ""),
              "Coords: exact lat/lng mined from the page's maps hyperlink (no geocoder needed)")
        check(r0.get("mapLink") == "https://maps.google.com/?q=40.5729170000,-3.2452220000"
              and "map link" in r0.get("__meta", {}).get("prov", {}).get("mapLink", ""),
              "mapLink: the brochure's own maps hyperlink ships as the property's mapLink with prov")

    # Fix 5: coverage treats geocode-ran-but-network-dead as environmental
    import json as _json
    import gate_runner
    canon = {"meta": {"enrichment": {"geocode": True, "degraded": True}},
             "properties": [{"id": 1, "country": "ES", "park": "Alpha", "developer": "CTP",
                             "city": "Madrid", "status": "Existing", "photo": "data:image/png;base64,x",
                             "warehouseArea": 40000, "warehouseRent": "€60 / sq m / year",
                             "warehouseRentVal": 60.0}],
             "pois": [], "regions": {}}
    with tempfile.TemporaryDirectory() as td:
        cp = Path(td) / "c.json"
        cp.write_text(_json.dumps(canon), encoding="utf-8")
        check(call(gate_runner, "coverage", cp) == 0,
              "Fix5: coverage passes when geocode ran but the sandbox network was dead")


def vision_validate_cases() -> None:
    print("vision validators (the whole extraction for scanned decks):")
    import vision_validate as VV
    with tempfile.TemporaryDirectory() as td:
        work = Path(td)
        (work / "vision").mkdir()
        (work / "extract").mkdir()
        (work / "vision" / "manifest.json").write_text(json.dumps({"decks": [
            {"region": "Cataluña", "source_file": "Scan.pdf",
             "pages": [{"page_no": 1}, {"page_no": 2}, {"page_no": 3}]}]}), encoding="utf-8")

        def rec(**kw):
            base = {"park": "P", "city": "C", "warehouseRentVal": 55.0,
                    "__meta": {"source_file": "Scan.pdf", "source_type": "pdf",
                               "page_no": 1, "prov": {}}}
            meta_over = kw.pop("__meta", {})
            base.update(kw)
            base["__meta"] = {**base["__meta"], **meta_over}
            return base

        def write(records):  # diacritic-STRIPPED filename must still match the deck
            (work / "extract" / "Cataluna_vision.json").write_text(
                json.dumps(records), encoding="utf-8")

        write([rec(), rec(__meta={"page_no": 2}), rec(__meta={"page_no": 3})])
        errors, warnings = VV.validate(work)
        check(not errors and not warnings, "clean transcription validates clean")

        write([rec(__meta={"page_no": 7})])  # out-of-range binding
        errors, _ = VV.validate(work)
        check(any("page_no 7" in e and "NEIGHBOUR" in e for e in errors),
              "out-of-range page_no is an ERROR (the hero mis-binding class)")

        # __meta.plan_page: a valid in-range int is clean; null is allowed; a bool / negative /
        # off-range value is an ERROR (it would render a neighbour's page as the site plan)
        write([rec(__meta={"page_no": 1, "plan_page": 2}), rec(__meta={"page_no": 2}),
               rec(__meta={"page_no": 3, "plan_page": None})])
        errors, warnings = VV.validate(work)
        check(not any("plan_page" in e for e in errors),
              "plan_page: an in-range int and a null both validate clean")
        write([rec(__meta={"page_no": 1, "plan_page": 7})])  # off-range
        errors, _ = VV.validate(work)
        check(any("plan_page 7" in e and "NEIGHBOUR" in e for e in errors),
              "plan_page: an off-range page is an ERROR (renders a neighbour's page as the plan)")
        write([rec(__meta={"page_no": 1, "plan_page": True})])  # bool rejected
        errors, _ = VV.validate(work)
        check(any("plan_page" in e and "non-negative integer" in e for e in errors),
              "plan_page: a boolean is rejected (not a page index)")

        write([rec(warehouseRentVal=4.2), rec(__meta={"page_no": 2}), rec(__meta={"page_no": 3})])
        errors, warnings = VV.validate(work)
        check(not errors and any("UN-ANNUALISED" in w for w in warnings),
              "a plausible-but-low rent warns as a likely un-annualised monthly quote")

        write([rec(warehouseRentVal=39471.0)])
        errors, _ = VV.validate(work)
        check(any("plausibility band" in e for e in errors),
              "an implausible rent is an ERROR")

        write([rec()])  # 1 record for 3 rasterised pages
        _, warnings = VV.validate(work)
        check(any("NO record" in w for w in warnings)
              and any("collapsed" in w for w in warnings),
              "uncovered pages warn about collapsed multi-property decks")


def region_code_cases() -> None:
    print("regionCode auto-derivation + loud empty-regions failure:")
    import io as _io
    from contextlib import redirect_stdout
    import gate_runner
    # requested regions + ZERO attached must HARD-block even when degraded
    with tempfile.TemporaryDirectory() as td:
        f = Path(td) / "c.json"
        f.write_text(json.dumps({"meta": {"enrichment": {"regions": True, "degraded": True}},
                                 "properties": [], "pois": [], "regions": {}}), encoding="utf-8")
        sv = sys.argv
        sys.argv = ["gate_runner", "enrichment", str(f)]
        buf = _io.StringIO()
        try:
            with redirect_stdout(buf):
                gate_runner.main()
            rc = 0
        except SystemExit as e:
            rc = e.code if isinstance(e.code, int) else 0
        finally:
            sys.argv = sv
        check(rc != 0 and "ZERO profiles" in buf.getvalue(),
              "empty workforce block HARD-blocks (never excused as DEGRADED)")


def placeholder_audit_cases() -> None:
    print("placeholder audit (a placeholder is never a silent default):")
    import base64 as _b64
    import io as _io
    import fitz
    from PIL import Image, ImageDraw
    import gate_runner
    import merge

    def _small_plan_png() -> bytes:  # line-art site plan, small
        img = Image.new("RGB", (300, 200), (255, 255, 255))
        d = ImageDraw.Draw(img)
        d.rectangle([20, 15, 280, 185], outline=(40, 40, 40), width=3)
        d.rectangle([40, 35, 200, 160], fill=(120, 180, 230))
        for x in range(50, 190, 20):
            d.line([(x, 40), (x, 155)], fill=(60, 60, 60), width=1)
        buf = _io.BytesIO(); img.save(buf, format="PNG"); return buf.getvalue()

    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        # 1. a SMALL plan (under the old photo floor) on an otherwise empty page
        #    must become the hero - it used to be discarded twice (size floor,
        #    then photo scorer) and ship a placeholder/page-render
        doc = fitz.open()
        pg = doc.new_page()
        pg.insert_image(fitz.Rect(60, 80, 180, 160), stream=_small_plan_png())
        f = td / "smallplan.pdf"
        doc.save(f); doc.close()
        hero, plan = IMG.page_hero_and_plan(f, 0, 60)
        IMG.close_doc_cache()
        ok = hero is not None and hero == plan
        if ok:
            him = Image.open(_io.BytesIO(_b64.b64decode(hero.split(",", 1)[1])))
            ok = 1.2 < him.width / him.height < 1.8  # the plan's aspect, not the A4 page's
        check(ok, "small site plan becomes the hero (plan floor, not photo floor)")

        # 2. when the pickers find nothing despite candidates existing (the
        #    degraded-engine reality), merge dumps the discard pile + the gate
        #    BLOCKS until the reviewer's sign-off is recorded
        rec = [{"park": "Audit Park", "developer": "Dev", "city": "Town", "country": "CZ",
                "status": "Existing", "warehouseArea": 20000,
                "__meta": {"source_file": "smallplan.pdf", "source_type": "pdf",
                           "locator_base": "page 1", "page_no": 0,
                           "prov": {k: "page 1" for k in ("park", "developer", "city",
                                                          "country", "status", "warehouseArea")}}}]
        (td / "rec.json").write_text(json.dumps(rec), encoding="utf-8")
        saved = IMG.page_hero_and_plan
        IMG.page_hero_and_plan = lambda *a, **k: (None, None)  # pickers blind
        try:
            sv = sys.argv
            sys.argv = ["merge", "--records", str(td / "rec.json"), "--source-dir", str(td),
                        "--out", str(td / "c.json"), "--ledger", str(td / "l.csv")]
            import io as _io2
            from contextlib import redirect_stdout
            with redirect_stdout(_io2.StringIO()):
                try:
                    merge.main()
                except SystemExit:
                    pass
            sys.argv = sv
        finally:
            IMG.page_hero_and_plan = saved
        canon = json.loads((td / "c.json").read_text(encoding="utf-8"))
        audit = canon["meta"].get("placeholderAudit", {})
        check(audit.get("1", {}).get("candidates", 0) > 0
              and (td / "render" / "placeholder_audit").exists(),
              "merge dumps the discard pile for a placeholder hero")

        def run_images_gate():
            sv = sys.argv
            sys.argv = ["gate_runner", "images", str(td / "c.json")]
            from contextlib import redirect_stdout
            buf = _io.StringIO()
            try:
                with redirect_stdout(buf):
                    gate_runner.main()
                return 0
            except SystemExit as e:
                return e.code if isinstance(e.code, int) else 0
            finally:
                sys.argv = sv
        check(run_images_gate() != 0, "images gate BLOCKS an unreviewed placeholder")
        (td / "placeholder_audit_ack.json").write_text(json.dumps(
            {"confirmed": ["1"], "by": "G-images reviewer",
             "note": "candidates inspected - none usable"}), encoding="utf-8")
        check(run_images_gate() == 0, "the recorded sign-off unblocks the gate")

        import contact_sheet
        sheet = contact_sheet.build_placeholder_audit(canon, td)
        check(sheet is not None and sheet.exists(),
              "contact sheet renders the discard-pile montage for the reviewer")


def regions_dataset_cases() -> None:
    print("regions dataset (Oxford Economics pre-fill):")
    import enrich
    enrich._REGIONS_DS = None  # fresh load of the real bundled dataset
    with tempfile.TemporaryDirectory() as td:
        enrich.CACHE_DIR = Path(td)
        # by province NAME and by NUTS code; researcher cache wins field-by-field for a
        # field the DATASET LACKS (wages removed at v11; employmentRate is dataset-absent)
        (Path(td) / "regions_cache.json").write_text(json.dumps({
            "Guadalajara": {"employmentRate": 72.0,
                            "sources": "Eurostat LFS 2026"}}), encoding="utf-8")
        canonical = {"properties": [{"id": 1, "regionCode": "Guadalajara"},
                                    {"id": 2, "regionCode": "ES424"},
                                    {"id": 3, "regionCode": "Nowhereland"}],
                     "pois": [], "regions": {}, "meta": {}}
        gaps: list = []
        updates: list = []
        n = enrich.merge_regions(canonical, gaps, updates)
        regs = canonical["regions"]
        g = regs.get("Guadalajara", {})
        check(n == 2 and g.get("unemployment") == 9.43 and g.get("labourForce") == 158300.0,
              "dataset: province-name regionCode pre-filled (unemployment, labour force)")
        check(g.get("emplTransportStorage") == 10990.0 and g.get("labourForce") == 158300.0,
              "dataset: logistics-share inputs present (transport&storage + labour force) - tile derives in-template (v11)")
        check(g.get("employmentRate") == 72.0 and "Eurostat" in g.get("sources", ""),
              "dataset: researcher values WIN field-by-field for a field the dataset lacks")
        check(regs.get("ES424", {}).get("nuts") == "ES424",
              "dataset: NUTS-code regionCode resolves directly")
        check("Nowhereland" not in regs,
              "dataset: unmatched codes stay empty for the researcher")
        check(not any("WAGE" in g_.upper() for g_ in gaps),
              "dataset: NO wage-research gap - the snapshot is fully dataset-sourced (v11)")
        check(any(u["field"] == "labourForce" for u in updates)
              and not any(u["field"] in ("avgWageGross", "minWageGross") for u in updates),
              "dataset: pre-filled figures get ledger trace rows; wages produce none (v11)")
        # the pre-filled profile passes the enrichment gate (current-year as-of)
        import gate_runner
        canonical["meta"]["enrichment"] = {"regions": True}
        f = Path(td) / "c.json"
        f.write_text(json.dumps(canonical), encoding="utf-8")
        sv = sys.argv
        sys.argv = ["gate_runner", "enrichment", str(f)]
        import io as _io
        from contextlib import redirect_stdout
        buf = _io.StringIO()
        try:
            with redirect_stdout(buf):
                gate_runner.main()
            rc = 0
        except SystemExit as e:
            rc = e.code if isinstance(e.code, int) else 0
        finally:
            sys.argv = sv
        check(rc == 0, "dataset: pre-filled profiles clear the enrichment gate first try")


def region_binding_cases() -> None:
    print("region binding by LOCATION (exact point-in-polygon; city/label fallback):")
    import enrich as _E
    _E._REGIONS_DS = None   # fresh load of the real bundled stats dataset
    _E._REGIONS_GEO = None  # fresh load of the real bundled NUTS-3 boundary asset
    ds = _E._regions_dataset()
    geo = _E._regions_geo()
    check(bool(geo and geo.get("regions")),
          "bundled NUTS-3 boundary asset (regions_geo.json.gz) loads (build_region_geo ran)")
    # point-in-polygon: the edge-of-province town nearest-centroid got WRONG is now correct
    check(_E._region_for_point(40.5658, -3.2637, geo) == "ES424",
          "PIP: Azuqueca de Henares (Guadalajara, on the Madrid border) -> ES424, NOT Madrid (the centroid bug)")
    check(_E._region_for_point(40.4168, -3.7038, geo) == "ES300", "PIP: Madrid -> ES300")
    check(_E._region_for_point(53.7997, -1.5492, geo) == "UKE42", "PIP: Leeds -> UKE42")
    canon = {"properties": [
        # coords in Azuqueca + a RESOLVING-but-wrong label 'Madrid' -> PIP overrides to ES424
        {"id": 1, "city": "Azuqueca", "region": "Madrid", "regionCode": "Madrid",
         "lat": 40.5658, "lng": -3.2637},
        # broken broad label, no coords -> CITY fallback ('Leeds' is a NUTS-3 name)
        {"id": 2, "city": "Leeds", "region": "Yorkshire And North East",
         "regionCode": "Yorkshire And North East"},
        # broken broad label, unmapped city, but coords -> PIP RESCUE
        {"id": 3, "city": "tbd", "region": "Yorkshire And North East",
         "regionCode": "Yorkshire And North East", "lat": 53.7997, "lng": -1.5492},
        # nothing usable -> left for merge_regions to gap
        {"id": 4, "city": "tbd", "region": "Nowhereland", "regionCode": "Nowhereland"}],
        "regions": {}, "pois": [], "meta": {}}
    _E.bind_region_codes(canon, ds)
    by = {p["id"]: p for p in canon["properties"]}
    check(by[1]["regionCode"] == "ES424",
          "coords AUTHORITATIVE: a wrong-but-resolving label 'Madrid' is corrected to ES424 by point-in-polygon")
    check(by[2]["regionCode"] == "UKE42",
          "city fallback (no coords): broken label + city 'Leeds' -> UKE42")
    check(by[3]["regionCode"] == "UKE42",
          "PIP rescue: broken label + unmapped city + Leeds coords -> UKE42 (the gap-report fix)")
    check(by[4]["regionCode"] == "Nowhereland",
          "nothing binds (no coords, unknown city + label) -> left for merge_regions to gap")
    gaps: list = []
    n = _E.merge_regions(canon, gaps)
    check(n == 2 and "ES424" in canon["regions"] and "UKE42" in canon["regions"],
          "merge_regions attaches the location-bound profiles (ES424 + UKE42)")
    check(any("Nowhereland" in g for g in gaps),
          "merge_regions still surfaces a genuinely unbindable regionCode as a gap")
    check("lat" not in canon["regions"]["ES424"] and "lng" not in canon["regions"]["ES424"],
          "the workforce profile carries stats only (no geometry leaked into canonical)")


def region_label_llm_cases() -> None:
    print("region-label LLM resolution (closed-set NUTS bind; verifier + PIP-primary):")
    import enrich as _E
    _E._REGIONS_DS = None   # fresh load of the real bundled stats dataset
    _E._REGIONS_GEO = None  # fresh load of the real bundled NUTS-3 boundary asset
    ds = _E._regions_dataset()

    def _seed(td, resolutions):
        ex = Path(td) / "extract"
        ex.mkdir(parents=True, exist_ok=True)
        (ex / "region_labels.json").write_text(
            json.dumps({"resolutions": resolutions}), encoding="utf-8")

    # (a) a cached mapping (a fuzzy label -> a real UKE NUTS-3 code) binds, RE-VERIFIED
    with tempfile.TemporaryDirectory() as td:
        _E.CACHE_DIR = Path(td)
        _seed(td, [{"raw_label": "Yorks & NE", "city": "Leeds", "country_cc": "GB",
                    "code": "UKE42"}])
        canon = {"properties": [{"id": 1, "city": "Leeds", "country": "GB",
                                 "region": "Yorks & NE", "regionCode": "Yorks & NE"}],
                 "regions": {}, "pois": [], "meta": {}}
        _E.bind_region_codes(canon, ds)
        check(canon["properties"][0]["regionCode"] == "UKE42",
              "a: a cached fuzzy label 'Yorks & NE' binds to the verified UKE42 NUTS-3 code")

    # (b) a cached code ABSENT from the dataset is REJECTED by the verifier -> falls to the gap
    with tempfile.TemporaryDirectory() as td:
        _E.CACHE_DIR = Path(td)
        _seed(td, [{"raw_label": "Yorks & NE", "city": "Leeds", "country_cc": "GB",
                    "code": "ZZ999"}])  # not a real dataset code
        canon = {"properties": [{"id": 1, "city": "tbd", "country": "GB",
                                 "region": "Yorks & NE", "regionCode": "Yorks & NE"}],
                 "regions": {}, "pois": [], "meta": {}}
        _E.bind_region_codes(canon, ds)
        check(canon["properties"][0]["regionCode"] == "Yorks & NE",
              "b: a cached code the dataset does not recognise is rejected (verifier holds)")
        gaps: list = []
        _E.merge_regions(canon, gaps)
        check(any("Yorks & NE" in g for g in gaps),
              "b: the rejected label still surfaces as the self-documenting difflib gap")

    # (c) a property WITH coords IGNORES the cache and the point-in-polygon bind WINS
    with tempfile.TemporaryDirectory() as td:
        _E.CACHE_DIR = Path(td)
        # cache would (wrongly) send this Leeds-coord property to a Spanish code; PIP must win
        _seed(td, [{"raw_label": "Yorks & NE", "city": "tbd", "country_cc": "GB",
                    "code": "ES424"}])
        canon = {"properties": [{"id": 1, "city": "tbd", "country": "GB",
                                 "region": "Yorks & NE", "regionCode": "Yorks & NE",
                                 "lat": 53.7997, "lng": -1.5492}],  # Leeds
                 "regions": {}, "pois": [], "meta": {}}
        _E.bind_region_codes(canon, ds)
        check(canon["properties"][0]["regionCode"] == "UKE42",
              "c: a property WITH coords ignores the cache - point-in-polygon (UKE42) wins")

    # (d) with NO region_labels.json the bind is BYTE-IDENTICAL to the deterministic path
    with tempfile.TemporaryDirectory() as td:
        _E.CACHE_DIR = Path(td)  # no extract/region_labels.json written
        canon = {"properties": [
            {"id": 1, "city": "Azuqueca", "region": "Madrid", "regionCode": "Madrid",
             "lat": 40.5658, "lng": -3.2637},
            {"id": 2, "city": "Leeds", "region": "Yorkshire And North East",
             "regionCode": "Yorkshire And North East"},
            {"id": 4, "city": "tbd", "region": "Nowhereland", "regionCode": "Nowhereland"}],
            "regions": {}, "pois": [], "meta": {}}
        _E.bind_region_codes(canon, ds)
        by = {p["id"]: p for p in canon["properties"]}
        check(by[1]["regionCode"] == "ES424" and by[2]["regionCode"] == "UKE42"
              and by[4]["regionCode"] == "Nowhereland",
              "d: with no cache the bind is byte-identical to region_binding_cases (fallback holds)")
        check(_E._region_labels_cache() == {},
              "d: an absent region_labels.json yields an empty cache (offline inert)")

    # (e) unresolved_region_labels is the SAME set merge_regions gaps on (single source of truth)
    with tempfile.TemporaryDirectory() as td:
        _E.CACHE_DIR = Path(td)  # no cache -> the deterministic path
        canon = {"properties": [
            {"id": 1, "city": "Leeds", "regionCode": "UKE42"},               # resolves (dataset)
            {"id": 2, "city": "tbd", "country": "GB",
             "regionCode": "Yorkshire And North East"},                      # coordless miss
            {"id": 3, "city": "tbd", "country": "ES",
             "regionCode": "Region Sur"},                                    # coordless miss
            {"id": 4, "city": "Leeds", "country": "GB",
             "regionCode": "Some Fuzzy Label", "lat": 53.7997, "lng": -1.5492}],  # PIP would bind
            "regions": {}, "pois": [], "meta": {}}
        # bind, then see what merge_regions reports as unresolved gaps
        unresolved = _E.unresolved_region_labels(canon, ds)
        raw_labels = sorted(rl for rl, _c, _cc in unresolved)
        _E.bind_region_codes(canon, ds)
        gaps2: list = []
        _E.merge_regions(canon, gaps2)
        # the coord-less misses are flagged; the dataset-resolving and the PIP-bound ones are not
        check(raw_labels == ["Region Sur", "Yorkshire And North East"],
              "e: unresolved_region_labels returns exactly the coord-less unresolvable labels")
        gap_text = " ".join(gaps2)
        check(all(rl in gap_text for rl in raw_labels)
              and "UKE42" not in [rl for rl, _c, _cc in unresolved],
              "e: the helper's labels are the SAME labels merge_regions gaps on (no drift)")
        # the PIP-bound property (id 4) is NOT in the unresolved set (coords win)
        check("Some Fuzzy Label" not in raw_labels,
              "e: a property with coords is excluded (point-in-polygon would override)")

    # candidate scoping: a GB-scoped list carries UK codes and excludes ES codes
    cands = _E.region_label_candidates(ds, ["GB"])
    codes = {c["code"] for c in cands}
    check(any(c.startswith("UK") for c in codes) and not any(c.startswith("ES") for c in codes),
          "candidates: a GB-scoped list carries UK codes and excludes ES codes (no cross-country bind)")
    check(any(c["code"] == "UKE" and c["name"] for c in cands),
          "candidates: the curated broad-region aliases (UKE) ride the candidate list too")


def recency_cases() -> None:
    print("labour-data recency (gate floor = current year - 1):")
    import datetime as _dt
    import io as _io
    from contextlib import redirect_stdout
    import gate_runner

    def run_gate(regions, month):
        saved = gate_runner._today
        gate_runner._today = lambda: _dt.date(2026, month, 15)
        canon = {"meta": {"enrichment": {"regions": True}},
                 "properties": [], "pois": [], "regions": regions}
        with tempfile.TemporaryDirectory() as td:
            f = Path(td) / "c.json"
            f.write_text(json.dumps(canon), encoding="utf-8")
            sv = sys.argv
            sys.argv = ["gate_runner", "enrichment", str(f)]
            buf = _io.StringIO()
            try:
                with redirect_stdout(buf):
                    gate_runner.main()
                rc = 0
            except SystemExit as e:
                rc = e.code if isinstance(e.code, int) else 0
            finally:
                sys.argv = sv
                gate_runner._today = saved
        return rc, buf.getvalue()

    def region(asof, **extra):
        return {"R1": {"name": "Region", "country": "ES", "unemployment": 8.5,
                       "unemploymentAsOf": asof, "sources": "INE 2026", **extra}}

    rc, _ = run_gate(region("2026"), 7)
    check(rc == 0, "current-year labour figure passes")
    rc, _ = run_gate(region("2025"), 7)
    check(rc == 0, "year-1 meets the floor")
    rc, _ = run_gate(region("2024"), 7)
    check(rc != 0, "year-2 BLOCKS from June onward")
    rc, _ = run_gate(region("2024"), 3)
    check(rc != 0, "year-2 in Jan-May still BLOCKS without a recencyNote")
    rc, out = run_gate(region("2024", recencyNote="No 2025 provincial unemployment release "
                                                  "published yet (INE checked 2026-03)"), 3)
    check(rc == 0 and "accepted Jan-May" in out,
          "year-2 in Jan-May passes ONLY with the failed-search recencyNote")
    rc, _ = run_gate(region("2023", recencyNote="x"), 3)
    check(rc != 0, "year-3 blocks regardless of month or note")
    rc, _ = run_gate({"R1": {"name": "R", "country": "ES", "gdpPpsEu": 88.0,
                             "gdpPpsAsOf": "2023", "sources": "Eurostat"}}, 7)
    check(rc == 0, "GDP keeps the softer advisory rule (publication lag is real)")
    # v11: wages were removed from the workforce snapshot, so a stale avgWageGross is
    # NO LONGER a gated labour figure - it must not block (proves wage left the floor)
    rc, _ = run_gate({"R1": {"name": "R", "country": "ES", "avgWageGross": 28000.0,
                             "avgWageAsOf": "2019", "sources": "INE"}}, 7)
    check(rc == 0, "v11: a stale avgWageGross no longer blocks (wage removed from the labour floor)")


def web_enrich_cases() -> None:
    print("web_enrich (orchestrator-web handoff for genuine POIs/drive-times):")
    import io as _io
    import os
    from contextlib import redirect_stdout
    import enrich
    import web_enrich
    # tests control the ORS key explicitly: pop any ambient ORS_API_KEY (this dev
    # box has one set) so the keyless/car sub-tests are deterministic; restored below
    _ORS_SAVED = os.environ.pop("ORS_API_KEY", None)

    def call(mod, *argv):
        saved = sys.argv
        sys.argv = [mod.__name__, *[str(a) for a in argv]]
        buf = _io.StringIO()
        try:
            with redirect_stdout(buf):
                mod.main()
            return 0
        except SystemExit as e:
            return e.code if isinstance(e.code, int) else 0
        finally:
            sys.argv = saved

    enrich._DATASET = False  # exercise the OSM web-handoff fallback path
    with tempfile.TemporaryDirectory() as td:
        work = Path(td)
        canonical = work / "canonical.json"
        seeded = [{"name": "Library Port", "type": "port", "lat": 50.0, "lng": 10.0}]
        canonical.write_text(json.dumps({
            "meta": {"enrichment": {}},
            "properties": [{"id": 1, "lat": 40.573, "lng": -3.245}],
            "pois": list(seeded), "regions": {},
        }), encoding="utf-8")

        # 0. with the bundled dataset present, plan emits NO Overpass requests at all
        enrich._DATASET = None
        rc0 = call(web_enrich, "plan", canonical, "--work", work, "--pois")
        plan0 = json.loads((work / "web_requests.json").read_text(encoding="utf-8"))
        check(not [r for r in plan0["requests"] if r["kind"] == "overpass"],
              "dataset: the web plan needs no Overpass round (POIs are local)")
        enrich._DATASET = False

        # 1. plan emits the exact Overpass request + the operator fetcher page
        rc = call(web_enrich, "plan", canonical, "--work", work, "--pois")
        plan = json.loads((work / "web_requests.json").read_text(encoding="utf-8"))
        check(rc == 0 and len(plan["requests"]) == 1
              and "overpass-api.de" in plan["requests"][0]["url"]
              and "aerodrome" in plan["requests"][0]["url"],
              "plan: one capped Overpass request per property location")
        page = (work / "web_enrich.html")
        check(page.exists() and "Fetch all" in page.read_text(encoding="utf-8")
              and plan["requests"][0]["save_as"] in page.read_text(encoding="utf-8"),
              "plan: self-contained browser fetcher page embeds the requests")

        # 2. the operator 'fetches' via the page - ONE seeds bundle with a port
        #    NEARBY and one far away: the GENUINE nearest must win
        fetched = work / "web_fetched"
        resp = {"elements": [
            {"type": "node", "tags": {"industrial": "port", "name": "Far Port"},
             "lat": 43.35, "lon": -3.04},
            {"type": "node", "tags": {"industrial": "port", "name": "Near Dry Port"},
             "lat": 40.60, "lon": -3.20},
            {"type": "way", "tags": {"aeroway": "aerodrome", "iata": "MAD",
                                     "name": "Madrid-Barajas"},
             "center": {"lat": 40.47, "lon": -3.56}},
        ]}
        (work / "web_seeds.json").write_text(
            json.dumps({plan["requests"][0]["save_as"]: json.dumps(resp)}), encoding="utf-8")
        # an OOM/remark-errored response must be SKIPPED, never cached as empty
        (work / "web_seeds.json").write_text(json.dumps({
            plan["requests"][0]["save_as"]: json.dumps(
                {"elements": [], "remark": "runtime error: Query run out of memory"})}),
            encoding="utf-8")
        call(web_enrich, "ingest", "--work", work)
        poisoned = json.loads((work / "poi_osm_cache.json").read_text(encoding="utf-8")) \
            if (work / "poi_osm_cache.json").exists() else {}
        check(plan["requests"][0]["key"] not in poisoned,
              "ingest: an OOM-remark response never poisons the cache as 'no POIs'")
        (work / "web_seeds.json").write_text(
            json.dumps({plan["requests"][0]["save_as"]: json.dumps(resp)}), encoding="utf-8")
        check(call(web_enrich, "ingest", "--work", work) == 0,
              "ingest: the browser bundle (web_seeds.json) is consumed")
        (work / "web_seeds.json").unlink()

        # 3. offline attach_pois consumes the warm cache: genuine nearest, no library.
        # The seeded-cache OSM fallback is reached only with NO complete dataset present, so
        # pin the border/city datasets absent too (poi_dataset is already _DATASET=False here).
        enrich.CACHE_DIR = work
        _sB, _sC = enrich._BORDERS, enrich._CITY_DATASET
        enrich._BORDERS = enrich._CITY_DATASET = False
        original = enrich.nearest_pois_for
        enrich.nearest_pois_for = lambda lat, lng: (_ for _ in ()).throw(RuntimeError("offline"))
        try:
            data = json.loads(canonical.read_text(encoding="utf-8"))
            gaps: list = []
            n, live = enrich.attach_pois(data, gaps)
        finally:
            enrich.nearest_pois_for = original
            enrich._BORDERS, enrich._CITY_DATASET = _sB, _sC
        names = sorted(p["name"] for p in data["pois"])
        check(live and "Near Dry Port" in names and "Library Port" not in names,
              "offline --pois attaches the GENUINE nearest from the seeded cache")

        # 4. OSRM round: plan emits one table call; a fake response bakes real minutes
        canonical.write_text(json.dumps(data), encoding="utf-8")
        rc = call(web_enrich, "plan", canonical, "--work", work, "--osrm")
        plan2 = json.loads((work / "web_requests.json").read_text(encoding="utf-8"))
        osrm_reqs = [r for r in plan2["requests"] if r["kind"] == "osrm-table"]
        check(rc == 0 and len(osrm_reqs) == 1 and "/table/v1/driving/" in osrm_reqs[0]["url"],
              "plan: one OSRM table request per property")
        ndest = len(osrm_reqs[0]["dests"])
        fake = {"durations": [[0] + [1800 + 60 * j for j in range(ndest)]],
                "distances": [[0] + [25000 + 1000 * j for j in range(ndest)]]}
        (fetched / osrm_reqs[0]["save_as"]).write_text(json.dumps(fake), encoding="utf-8")
        check(call(web_enrich, "ingest", "--work", work) == 0, "ingest: drive times cached")
        gaps2: list = []
        done = enrich.osrm_prebake(data, gaps2, "https://router.project-osrm.org")
        dists = (data["properties"][0].get("preBaked") or {}).get("distances", {})
        check(done == 1 and dists and all("min" in v and "km" in v for v in dists.values()),
              "offline --osrm bakes REAL routed minutes from the seeded cache")

        # 4b. TRUCKING round (openrouteservice key set): the plan emits an ORS
        #     driving-hgv MATRIX request (POST, Authorization header) instead of
        #     car OSRM; ingest tags the cache entries hgv| so car values can
        #     never masquerade as trucking
        data["properties"][0].pop("preBaked", None)
        canonical.write_text(json.dumps(data), encoding="utf-8")
        rc = call(web_enrich, "plan", canonical, "--work", work, "--osrm", "--ors-key", "TESTKEY")
        plan_h = json.loads((work / "web_requests.json").read_text(encoding="utf-8"))
        ors = [r for r in plan_h["requests"] if r["kind"] == "ors-matrix"]
        check(rc == 0 and len(ors) == 1 and "/v2/matrix/driving-hgv" in ors[0]["url"]
              and ors[0]["method"] == "POST" and ors[0]["headers"]["Authorization"] == "TESTKEY"
              and json.loads(ors[0]["body"])["metrics"] == ["distance", "duration"],
              "ORS: plan emits one trucking matrix POST per property (key in header)")
        nd = len(ors[0]["dests"])
        fake_h = {"durations": [[2400 + 60 * j for j in range(nd)]],
                  "distances": [[31000 + 1000 * j for j in range(nd)]]}
        (work / "web_seeds.json").write_text(
            json.dumps({ors[0]["save_as"]: json.dumps(fake_h)}), encoding="utf-8")
        check(call(web_enrich, "ingest", "--work", work) == 0, "ORS: ingest caches hgv pairs")
        (work / "web_seeds.json").unlink()
        rc_cache = json.loads((work / "osrm_cache.json").read_text(encoding="utf-8"))
        check(any(k.startswith("hgv|") for k in rc_cache),
              "ORS: trucking entries are profile-tagged in the cache")
        gaps3: list = []
        done = enrich.osrm_prebake(data, gaps3, "https://router.project-osrm.org",
                                   None, ors_key="TESTKEY")
        dists_h = (data["properties"][0].get("preBaked") or {}).get("distances", {})
        check(done == 1 and dists_h and any(v.get("min") == 40 for v in dists_h.values()),
              "ORS: offline prebake bakes TRUCKING minutes from the hgv cache")

        # 5. geocode round: an unlocated city joins the same artifact; the
        #    Nominatim response seeds the geocode cache
        data["properties"].append({"id": 2, "city": "Alovera", "country": "ES"})
        canonical.write_text(json.dumps(data), encoding="utf-8")
        rc = call(web_enrich, "plan", canonical, "--work", work, "--geocode")
        plan3 = json.loads((work / "web_requests.json").read_text(encoding="utf-8"))
        geo = [r for r in plan3["requests"] if r["kind"] == "nominatim"]
        check(rc == 0 and len(geo) == 1 and "countrycodes=es" in geo[0]["url"],
              "plan: unlocated city -> one country-scoped Nominatim request")
        (work / "web_seeds.json").write_text(json.dumps({
            geo[0]["save_as"]: json.dumps([{"lat": "40.5869", "lon": "-3.2328",
                                            "address": {"country_code": "es"}}])}),
            encoding="utf-8")
        check(call(web_enrich, "ingest", "--work", work) == 0, "ingest: geocode cached")
        gc = json.loads((work / "geocode_cache.json").read_text(encoding="utf-8"))
        check(gc.get("alovera|es", {}).get("latlng") == [40.5869, -3.2328],
              "geocode cache carries the browser-fetched coordinates")

        # 6. SELF-CHAINING single round: the page embeds the chain spec (geocode
        #    urls + POI set + caps + routing config) and ingest accepts the
        #    self-describing v2 bundle - geocode + routes in ONE click
        enrich._DATASET = None  # the chain embeds the real bundled POI dataset
        data["properties"].append({"id": 3, "city": "Toledo", "country": "ES"})
        canonical.write_text(json.dumps(data), encoding="utf-8")
        rc = call(web_enrich, "plan", canonical, "--work", work,
                  "--geocode", "--osrm", "--ors-key", "TESTKEY")
        page = (work / "web_enrich.html").read_text(encoding="utf-8")
        check(rc == 0 and '"do_routes": true' in page.replace('"do_routes":true', '"do_routes": true')
              and "geocode_url" in page and '"caps"' in page and "runChain" in page,
              "chain: the page embeds geocode urls, the POI set, caps and the chain driver")
        seeds_v2 = {"v": 2,
                    "geocode": [{"key": "toledo|es",
                                 "body": json.dumps([{"lat": "39.8628", "lon": "-4.0273",
                                                      "address": {"country_code": "es"}}])}],
                    "routes": [{"kind": "ors-matrix", "src": [39.8628, -4.0273],
                                "dests": [[40.4, -3.6]],
                                "body": json.dumps({"durations": [[5400]],
                                                    "distances": [[95000]]})}],
                    "static": {}}
        (work / "web_seeds.json").write_text(json.dumps(seeds_v2), encoding="utf-8")
        check(call(web_enrich, "ingest", "--work", work) == 0, "chain: v2 bundle ingests")
        gc2 = json.loads((work / "geocode_cache.json").read_text(encoding="utf-8"))
        rc2 = json.loads((work / "osrm_cache.json").read_text(encoding="utf-8"))
        check(gc2.get("toledo|es", {}).get("latlng") == [39.8628, -4.0273]
              and rc2.get("hgv|39.8628,-4.0273;40.4,-3.6", {}).get("min") == 90,
              "chain: one bundle seeds BOTH the geocode and the trucking route caches")

        # ================= Batch E: web-enrich correctness =================
        import os as _os

        def call_cap(mod, *argv):
            """Run a helper main(), returning (rc, captured_stdout)."""
            saved = sys.argv
            sys.argv = [mod.__name__, *[str(a) for a in argv]]
            buf = _io.StringIO()
            try:
                with redirect_stdout(buf):
                    mod.main()
                rc = 0
            except SystemExit as e:
                rc = e.code if isinstance(e.code, int) else 0
            finally:
                sys.argv = saved
            return rc, buf.getvalue()

        # --- P3-9: _ingest_route returns (n, err); an ORS/OSRM error envelope is
        #     flagged and NEVER cached, while empty-but-valid stays distinct ---
        check(web_enrich._ingest_route("ors-matrix", [50.0, 5.0], [[51.0, 6.0]],
              {"durations": [[1200]], "distances": [[40000]]}, {}) == (1, None),
              "P3-9: a valid ORS body returns (1, None) - happy path unchanged")
        _ce = {}
        _re = web_enrich._ingest_route("ors-matrix", [50.0, 5.0], [[51.0, 6.0]],
                                       {"error": {"code": 2004, "message": "rate limit"}}, _ce)
        check(_re[0] == 0 and _re[1] and not _ce,
              "P3-9: an ORS error envelope returns (0, err) and never poisons the cache")
        _co = {}
        _ro = web_enrich._ingest_route("osrm-table", [50.0, 5.0], [[51.0, 6.0]],
                                       {"code": "InvalidQuery", "message": "bad"}, _co)
        check(_ro[0] == 0 and str(_ro[1]).startswith("OSRM InvalidQuery") and not _co,
              "P3-9: an OSRM InvalidQuery is flagged, not cached")
        _cv = {}
        check(web_enrich._ingest_route("osrm-table", [50.0, 5.0], [[51.0, 6.0]],
              {"code": "Ok", "durations": [[None, None]]}, _cv) == (0, None) and not _cv,
              "P3-9: an empty-but-valid response stays distinct (0, None), nothing cached")
        with tempfile.TemporaryDirectory() as _td9:
            _w9 = Path(_td9)
            (_w9 / "web_requests.json").write_text(json.dumps({"requests": []}), encoding="utf-8")
            (_w9 / "web_seeds.json").write_text(json.dumps({"v": 2, "geocode": [], "routes": [
                {"kind": "ors-matrix", "src": [40.0, -3.0], "dests": [[41.0, -3.5]],
                 "body": json.dumps({"error": {"code": 2004, "message": "quota"}})},
                {"kind": "ors-matrix", "src": [40.0, -3.0], "dests": [[40.5, -3.2]],
                 "body": json.dumps({"durations": [[1500]], "distances": [[30000]]})}],
                "static": {}}), encoding="utf-8")
            _rc9, _out9 = call_cap(web_enrich, "ingest", "--work", _w9)
            _rcache = json.loads((_w9 / "osrm_cache.json").read_text(encoding="utf-8")) \
                if (_w9 / "osrm_cache.json").exists() else {}
            check("[skip]" in _out9 and "1 unreadable" in _out9
                  and len(_rcache) == 1 and any(v.get("min") == 25 for v in _rcache.values()),
                  "P3-9: a v2 errored route is [skip]-logged + counted bad; only the good pair caches")
        # _ingest_route must be RAISE-SAFE (the static caller has no try/except): a
        # malformed-but-valid-JSON body (null/non-numeric durations, bad dests/src) and
        # an error envelope carrying a null durations row must yield (0, ...), never raise
        for _bad in ({"error": {"code": 2004}, "durations": [None]},
                     {"durations": [["NA"]]}, {"code": "Ok", "durations": "junk"}):
            _cb = {}
            _nb, _eb = web_enrich._ingest_route("osrm-table", [40.0, -3.0], [[40.5, -3.2]], _bad, _cb)
            check(_nb == 0 and not _cb, f"P3-9: malformed body {list(_bad)[:1]} is raise-safe, nothing cached")
        check(web_enrich._ingest_route("ors-matrix", None, [[40.5, -3.2]], {"x": 1}, {})[0] == 0,
              "P3-9: a bad src is raise-safe (returns 0, never unpacks-crashes)")
        # the STATIC WebFetch-fallback branch (no try/except) must [skip]-log an error
        # body by save_as and survive a malformed one, never crashing the whole ingest
        with tempfile.TemporaryDirectory() as _td9s:
            _w9s = Path(_td9s)
            _fdir = _w9s / "web_fetched"
            _fdir.mkdir()
            (_w9s / "web_requests.json").write_text(json.dumps({"requests": [
                {"id": "ors_p1", "kind": "ors-matrix", "src": [40.0, -3.0],
                 "dests": [[40.5, -3.2]], "save_as": "ors_p1.json", "url": "x"},
                {"id": "ors_p2", "kind": "ors-matrix", "src": [40.0, -3.0],
                 "dests": [[40.6, -3.3]], "save_as": "ors_p2.json", "url": "y"}]}),
                encoding="utf-8")
            (_fdir / "ors_p1.json").write_text(json.dumps(
                {"error": {"code": 2004, "message": "quota"}, "durations": [None]}), encoding="utf-8")
            (_fdir / "ors_p2.json").write_text(json.dumps(
                {"durations": [[1500]], "distances": [[30000]]}), encoding="utf-8")
            _rc9s, _out9s = call_cap(web_enrich, "ingest", "--work", _w9s)
            _rc9scache = json.loads((_w9s / "osrm_cache.json").read_text(encoding="utf-8")) \
                if (_w9s / "osrm_cache.json").exists() else {}
            check("[skip] ors_p1.json" in _out9s and len(_rc9scache) == 1
                  and any(v.get("min") == 25 for v in _rc9scache.values()),
                  "P3-9: the static caller [skip]-logs an error body by save_as and survives it")

        # --- P3-8: a STRING \"2\" version still folds (no silent static degrade);
        #     a versionless seed-shaped bundle WARNS rather than dropping silently ---
        with tempfile.TemporaryDirectory() as _td8:
            _w8 = Path(_td8)
            (_w8 / "web_requests.json").write_text(json.dumps({"requests": []}), encoding="utf-8")
            (_w8 / "web_seeds.json").write_text(json.dumps({"v": "2", "geocode": [
                {"key": "leon|es", "body": json.dumps(
                    [{"lat": "42.6", "lon": "-5.57", "address": {"country_code": "es"}}])}],
                "routes": [], "static": {}}), encoding="utf-8")
            _rc8, _out8 = call_cap(web_enrich, "ingest", "--work", _w8)
            _gc8 = json.loads((_w8 / "geocode_cache.json").read_text(encoding="utf-8")) \
                if (_w8 / "geocode_cache.json").exists() else {}
            check(_gc8.get("leon|es", {}).get("latlng") == [42.6, -5.57],
                  "P3-8: a string \"2\" version still folds the bundle (no silent static degrade)")
            (_w8 / "web_seeds.json").write_text(json.dumps({"geocode": [
                {"key": "vigo|es", "body": json.dumps(
                    [{"lat": "42.24", "lon": "-8.72", "address": {"country_code": "es"}}])}],
                "routes": []}), encoding="utf-8")
            _rcw, _outw = call_cap(web_enrich, "ingest", "--work", _w8)
            _gcw = json.loads((_w8 / "geocode_cache.json").read_text(encoding="utf-8")) \
                if (_w8 / "geocode_cache.json").exists() else {}
            check("[warn]" in _outw and "vigo|es" not in _gcw,
                  "P3-8: a versionless seed-shaped bundle warns and is NOT folded silently")
            # bool subclasses int: v:true must be treated as unrecognised (warn, not fold),
            # never coerced to 1; an old v:1 seed-shaped bundle must also warn, not drop silently
            for _badv in (True, 1):
                (_w8 / "web_seeds.json").write_text(json.dumps({"v": _badv, "geocode": [
                    {"key": "jaen|es", "body": json.dumps(
                        [{"lat": "37.77", "lon": "-3.79", "address": {"country_code": "es"}}])}],
                    "routes": []}), encoding="utf-8")
                _rcb, _outb = call_cap(web_enrich, "ingest", "--work", _w8)
                _gcb = json.loads((_w8 / "geocode_cache.json").read_text(encoding="utf-8")) \
                    if (_w8 / "geocode_cache.json").exists() else {}
                check("[warn]" in _outb and "jaen|es" not in _gcb,
                      f"P3-8: a seed-shaped bundle with v={_badv!r} warns and is NOT folded")

        # --- P2-7: a complete v2 bundle never says 'fetch again' for chain-covered
        #     kinds; a leftover Overpass request gets the POI-specific nudge; the
        #     no-bundle fallback still reports every missing request verbatim ---
        with tempfile.TemporaryDirectory() as _td7:
            _w7 = Path(_td7)
            (_w7 / "web_requests.json").write_text(json.dumps({"requests": [
                {"id": "geo_x", "kind": "nominatim", "key": "x|es", "save_as": "geo_x.json",
                 "url": "https://nominatim.openstreetmap.org/search?q=x"},
                {"id": "ors_p1", "kind": "ors-matrix", "src": [40.0, -3.0],
                 "dests": [[40.5, -3.2]], "save_as": "ors_p1.json", "url": "https://api"}]}),
                encoding="utf-8")
            (_w7 / "web_seeds.json").write_text(json.dumps({"v": 2, "geocode": [
                {"key": "x|es", "body": json.dumps(
                    [{"lat": "40.0", "lon": "-3.0", "address": {"country_code": "es"}}])}],
                "routes": [{"kind": "ors-matrix", "src": [40.0, -3.0], "dests": [[40.5, -3.2]],
                            "body": json.dumps({"durations": [[1500]], "distances": [[30000]]})}],
                "static": {}}), encoding="utf-8")
            _rc7, _out7 = call_cap(web_enrich, "ingest", "--work", _w7)
            check("0 response(s) not fetched yet" in _out7
                  and "Fetch the remaining requests" not in _out7
                  and "outstanding" not in _out7,
                  "P2-7: a complete v2 bundle reports 0 missing and never says 'fetch again'")
            (_w7 / "web_requests.json").write_text(json.dumps({"requests": [
                {"id": "poi_x", "kind": "overpass", "key": "40,-3", "src": [40.0, -3.0],
                 "save_as": "poi_x.json", "url": "https://overpass"}]}), encoding="utf-8")
            _rc7b, _out7b = call_cap(web_enrich, "ingest", "--work", _w7)
            check("1 response(s) not fetched yet" in _out7b and "Overpass" in _out7b
                  and "Fetch the remaining requests" not in _out7b,
                  "P2-7: a leftover Overpass request after a v2 bundle gets the POI nudge")
        with tempfile.TemporaryDirectory() as _td7b:
            _w7c = Path(_td7b)
            (_w7c / "web_requests.json").write_text(json.dumps({"requests": [
                {"id": "geo_y", "kind": "nominatim", "key": "y|es", "save_as": "geo_y.json",
                 "url": "https://nominatim"}]}), encoding="utf-8")
            _rc7c, _out7c = call_cap(web_enrich, "ingest", "--work", _w7c)
            check("1 response(s) not fetched yet" in _out7c
                  and "Fetch the remaining requests" in _out7c,
                  "P2-7: the no-bundle WebFetch fallback still reports every missing request")

        # --- P1-8: standalone `plan --osrm` honours ORS_API_KEY from the env
        #     (trucking matrix), an explicit --ors-key still wins, and no key
        #     anywhere keeps the keyless car path ---
        with tempfile.TemporaryDirectory() as _td6:
            _w6 = Path(_td6)
            _c6 = _w6 / "canonical.json"
            _c6.write_text(json.dumps({"meta": {"enrichment": {}}, "properties": [
                {"id": 1, "lat": 40.5, "lng": -3.2}],
                "pois": [{"name": "P", "type": "port", "lat": 40.6, "lng": -3.1}],
                "regions": {}}), encoding="utf-8")
            _saved_env = _os.environ.get("ORS_API_KEY")
            try:
                _os.environ["ORS_API_KEY"] = "ENVKEY"
                call(web_enrich, "plan", _c6, "--work", _w6, "--osrm")  # NO --ors-key
                _p6 = json.loads((_w6 / "web_requests.json").read_text(encoding="utf-8"))
                _ors6 = [r for r in _p6["requests"] if r["kind"] == "ors-matrix"]
                check(len(_ors6) == 1 and _ors6[0]["headers"]["Authorization"] == "ENVKEY",
                      "P1-8: standalone plan --osrm reads ORS_API_KEY from the env (trucking matrix)")
                call(web_enrich, "plan", _c6, "--work", _w6, "--osrm", "--ors-key", "EXPLICIT")
                _p6b = json.loads((_w6 / "web_requests.json").read_text(encoding="utf-8"))
                _ors6b = [r for r in _p6b["requests"] if r["kind"] == "ors-matrix"]
                check(bool(_ors6b) and _ors6b[0]["headers"]["Authorization"] == "EXPLICIT",
                      "P1-8: an explicit --ors-key still wins over the env default")
                _os.environ.pop("ORS_API_KEY", None)  # no key anywhere -> keyless car
                call(web_enrich, "plan", _c6, "--work", _w6, "--osrm")
                _p6c = json.loads((_w6 / "web_requests.json").read_text(encoding="utf-8"))
                check([r for r in _p6c["requests"] if r["kind"] == "osrm-table"]
                      and not [r for r in _p6c["requests"] if r["kind"] == "ors-matrix"],
                      "P1-8: no key in env or args keeps the keyless car (OSRM) path")
            finally:
                if _saved_env is None:
                    _os.environ.pop("ORS_API_KEY", None)
                else:
                    _os.environ["ORS_API_KEY"] = _saved_env

        # --- P2-8: the chain page carries known + the PASS B re-query; offline,
        #     enrich flags an ambiguous far-outlier as an honest gap (cannot re-query) ---
        enrich._DATASET = None  # embed the real bundled POI dataset in the chain
        with tempfile.TemporaryDirectory() as _td28:
            _w28 = Path(_td28)
            _c28 = _w28 / "canonical.json"
            _c28.write_text(json.dumps({"meta": {"enrichment": {}}, "properties": [
                {"id": 1, "lat": 40.4, "lng": -3.7, "country": "ES"},
                {"id": 2, "city": "Boston", "country": "??"}], "pois": [], "regions": {}}),
                encoding="utf-8")
            call(web_enrich, "plan", _c28, "--work", _w28, "--geocode", "--osrm", "--ors-key", "K")
            _page28 = (_w28 / "web_enrich.html").read_text(encoding="utf-8")
            check('"known"' in _page28 and "geocode(reB)" in _page28 and "dominant" in _page28,
                  "P2-8: the chain page carries known + the PASS B dominant-country re-query")
            # offline honesty: a wrong-continent same-name pin must be flagged on the
            # CANONICAL warm-cache path - ALL cities are cache hits (so PASS A makes no
            # network call and `dead` stays False), yet the far-outlier re-query fails
            # offline; the gap must still fire (it must NOT depend on a cache miss)
            enrich.CACHE_DIR = _w28
            _seed = {
                "madrid|es": {"latlng": [40.42, -3.70], "cc": "ES"},
                "toledo|es": {"latlng": [39.86, -4.03], "cc": "ES"},
                "getafe|es": {"latlng": [40.31, -3.73], "cc": "ES"},
                "boston|??": {"latlng": [42.36, -71.06], "cc": "US"}}
            enrich._save_cache(enrich.GEOCODE_CACHE, _seed)
            _props28 = [{"id": 1, "city": "Madrid", "country": "ES"},
                        {"id": 2, "city": "Toledo", "country": "ES"},
                        {"id": 3, "city": "Getafe", "country": "ES"},
                        {"id": 4, "city": "Boston", "country": "??"}]
            _data28 = {"meta": {"enrichment": {}}, "properties": [dict(p) for p in _props28],
                       "pois": [], "regions": {}}
            _orig_geo1 = enrich._geocode_one
            enrich._geocode_one = lambda *a, **k: (_ for _ in ()).throw(RuntimeError("offline"))
            _gaps28: list = []
            try:
                enrich.geocode(_data28, _gaps28)
            finally:
                enrich._geocode_one = _orig_geo1
            _boston = next(p for p in _data28["properties"] if p["id"] == 4)
            check(any("could not be re-checked" in g and "Boston" in g for g in _gaps28)
                  and abs(_boston.get("lat", 0) - 42.36) < 0.01,
                  "P2-8: warm-cache offline far-outlier is flagged (pin kept, verify note)")
            # live PASS B still CORRECTS an ambiguous outlier and emits the constrained gap
            enrich._save_cache(enrich.GEOCODE_CACHE, _seed)
            _data28b = {"meta": {"enrichment": {}}, "properties": [dict(p) for p in _props28],
                        "pois": [], "regions": {}}
            enrich._geocode_one = lambda req, city, cc: ([39.96, -4.83], "ES")  # in-country hit
            _gaps28b: list = []
            try:
                enrich.geocode(_data28b, _gaps28b)
            finally:
                enrich._geocode_one = _orig_geo1
            _boston_b = next(p for p in _data28b["properties"] if p["id"] == 4)
            check(any("constrained to ES" in g and "Boston" in g for g in _gaps28b)
                  and abs(_boston_b.get("lat", 0) - 39.96) < 0.01,
                  "P2-8: live PASS B re-queries the outlier into the dominant country")
    enrich._DATASET = None  # restore the real dataset
    if _ORS_SAVED is not None:  # restore the ambient key the suite popped
        os.environ["ORS_API_KEY"] = _ORS_SAVED


def pptx_slide_cases() -> None:
    """Slide-sourced records get the same hero/plan/audit treatment as PDF pages
    (a vision transcription of a PPTX deck used to ship 100% placeholders), and
    vision_prep's LibreOffice tier rasterises slides python-pptx cannot."""
    print("PPTX slide harvest (heroes + vision rasters for slide decks):")
    import io as _io
    import os as _os
    import tempfile
    from PIL import Image
    import images as IMG
    import merge as M
    import vision_prep as VP
    from pptx import Presentation
    from pptx.util import Inches

    def _noise_jpeg(w=640, h=400):  # pure pixel noise = unambiguous photo signature
        img = Image.frombytes("RGB", (w, h), _os.urandom(w * h * 3))
        buf = _io.BytesIO()
        img.save(buf, format="JPEG", quality=90)
        return buf.getvalue()

    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        prs = Presentation()
        blank = prs.slide_layouts[6]
        s1 = prs.slides.add_slide(blank)  # slide 1: a real photo
        s1.shapes.add_picture(_io.BytesIO(_noise_jpeg()),
                              Inches(0.5), Inches(0.5), Inches(6), Inches(3.75))
        s2 = prs.slides.add_slide(blank)  # slide 2: text only (python-pptx cannot render)
        tb = s2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
        tb.text_frame.text = "Bratislava logistics park - vector slide"
        deck = td / "Deck Bratislava.pptx"
        prs.save(str(deck))

        saved = IMG._SOFFICE
        IMG._SOFFICE = None  # force the no-LibreOffice fallback tier (deterministic)
        try:
            hero, _plan = IMG.slide_hero_and_plan(deck, 0, 60)
            check(isinstance(hero, str) and hero.startswith("data:image/"),
                  "PPTX: a photo slide yields a real hero (embedded-picture tier)")
            rec = {"park": "BA Park", "city": "Bratislava",
                   "__meta": {"source_file": deck.name, "source_type": "pptx", "page_no": 0}}
            photo, _p, prec, _r, tried, gal = M.attach_media([rec], td, 60)
            check(photo != IMG.placeholder() and prec is rec
                  and bool(tried) and tried[0][2] == "pptx",
                  "PPTX: attach_media harvests a slide hero for a vision-style pptx record")
            check(isinstance(gal, list) and gal and gal[0] == photo and len(gal) <= IMG.GALLERY_MAX,
                  "GALLERY: attach_media returns a gallery, hero first, capped at GALLERY_MAX")
            rec2 = {"park": "BA Park 2", "city": "Bratislava",
                    "__meta": {"source_file": deck.name, "source_type": "pptx", "page_no": 1}}
            photo2, _p2, prec2, _r2, tried2, gal2 = M.attach_media([rec2], td, 60)
            check(photo2 == IMG.placeholder() and prec2 is None
                  and bool(tried2) and tried2[0][2] == "pptx",
                  "PPTX: a pictureless slide degrades to the placeholder WITH an audit-trail entry")
            ent = VP.prepare(deck, "R", "SK", td / "vis", force=True)
            with_img = [p for p in ent["pages"] if p.get("image")]
            without = [p for p in ent["pages"] if not p.get("image")]
            check(len(ent["pages"]) == 2 and len(with_img) == 1 and len(without) == 1
                  and "LibreOffice" in without[0]["reason"],
                  "PPTX: vision_prep fallback rasterises picture slides and NAMES the LibreOffice gap")
        finally:
            IMG._SOFFICE = saved

        # live LibreOffice tier - exercised only on machines that have soffice
        IMG._SOFFICE = False  # re-probe
        if IMG._find_soffice():
            ent = VP.prepare(deck, "R", "SK", td / "vis2", force=True)
            check(len(ent["pages"]) == 2 and all(p.get("image") for p in ent["pages"]),
                  "PPTX: LibreOffice tier rasterises EVERY slide (text/vector included)")
            hero2, _ = IMG.slide_hero_and_plan(deck, 1, 60)
            check(hero2 is not None,
                  "PPTX: LibreOffice tier gives even a text slide a hero via the page ladder")
        else:
            print("  [note] LibreOffice absent here - the soffice tier runs live on machines that have it")
        IMG.close_doc_cache()


def vision_dedup_cases() -> None:
    """Batch 2: manifest-driven per-file vision supersede, the coordinate-proximity
    dedupe net, and numeric reconciliation against the twin text layer."""
    print("vision supersede + coordinate net + numeric reconciliation:")
    import tempfile
    import run as RUN
    import match as MA
    import vision_validate as VV

    # 1. manifest-driven supersede: the rasterised file is superseded OUTRIGHT,
    # its never-rasterised clean twin is not (mixed regions stay safe)
    with tempfile.TemporaryDirectory() as td:
        work = Path(td)
        (work / "extract").mkdir()
        (work / "vision").mkdir()
        (work / "extract" / "Cataluna_vision.json").write_text("[]", encoding="utf-8")
        (work / "vision" / "manifest.json").write_text(json.dumps({"decks": [
            {"region": "Cataluña", "source_file": "Garbled Deck.pdf", "pages": []}]}),
            encoding="utf-8")
        check(RUN._vision_supersedes(work, "Cataluña", "Garbled Deck.pdf"),
              "supersede: a manifest-rasterised file is superseded OUTRIGHT (diacritic-tolerant)")
        check(not RUN._vision_supersedes(work, "Cataluña", "Clean Twin.pdf"),
              "supersede: a never-rasterised clean twin keeps its records")
        check(not RUN._vision_supersedes(work, "Madrid", "Garbled Deck.pdf"),
              "supersede: a region without a transcription is never superseded")

    # 2. coordinate net: a tbd city defeats the text key; first-party pins decide
    A = {"city": "tbd", "developer": "CTP", "park": "tbd", "lat": 40.5000, "lng": -3.2000,
         "warehouseArea": 40000, "__meta": {"source_file": "a.pdf"}}
    B = {"city": "Azuqueca", "developer": "CTP", "park": "CTPark Azuqueca",
         "lat": 40.5008, "lng": -3.2005, "warehouseArea": 40500,
         "__meta": {"source_file": "b_vision.json"}}
    check(MA.same_property(A, B),
          "coord net: ~100 m apart, same dev, same size -> ONE property despite tbd city/park")
    far = dict(A, lat=40.52)  # ~2.2 km away
    check(not MA.same_property(far, B), "coord net: distant pins stay distinct")
    other_dev = dict(A, developer="Panattoni")
    check(not MA.same_property(other_dev, B), "coord net: developer DISAGREEMENT stays distinct")
    phase2 = dict(A, warehouseArea=80000)
    check(not MA.same_property(phase2, B),
          "coord net: same pin but materially different size = a distinct phase")
    unknown_dev = dict(A, developer="tbd")
    check(MA.same_property(unknown_dev, B),
          "coord net: an unknown developer never counts as a disagreement")

    # 3. numeric reconciliation vs the twin text layer
    import fitz
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        inputs = td / "inputs"; inputs.mkdir()
        work = td / "work"
        (work / "extract").mkdir(parents=True)
        (work / "vision").mkdir()
        doc = fitz.open()
        pg = doc.new_page()
        pg.insert_text((72, 100), "Warehouse 39 471 sq m", fontsize=12)
        pg.insert_text((72, 130), "Rent: 4,50 EUR / m2 / mes", fontsize=12)
        src = inputs / "Twin.pdf"
        doc.save(src); doc.close()
        (work / "vision" / "manifest.json").write_text(json.dumps({"decks": [
            {"region": "R", "source_file": "Twin.pdf",
             "pages": [{"page_no": 0, "locator": "page 1"}]}]}), encoding="utf-8")

        def _vrec(rent, area):
            return [{"park": "P", "city": "X", "developer": "D", "status": "Existing",
                     "warehouseRentVal": rent, "warehouseArea": area,
                     "__meta": {"source_file": "Twin.pdf", "page_no": 0,
                                "locator_base": "page 1",
                                "prov": {"warehouseRentVal": "page 1 (x12 annualised)"}}}]

        (work / "extract" / "R_vision.json").write_text(
            json.dumps(_vrec(54.0, 39471)), encoding="utf-8")
        errs, warns = VV.validate(work, source_dir=inputs)
        check(not errs and not any("misread" in w for w in warns),
              "reconcile: 54 EUR/yr (= page's 4,50/mes x12) + the page's area pass clean")
        (work / "extract" / "R_vision.json").write_text(
            json.dumps(_vrec(63.0, 39471)), encoding="utf-8")
        _e, warns = VV.validate(work, source_dir=inputs)
        check(any("warehouseRentVal 63" in w and "misread" in w for w in warns),
              "reconcile: a rent the page never shows is flagged as a suspected misread")
        (work / "extract" / "R_vision.json").write_text(
            json.dumps(_vrec(54.0, 38471)), encoding="utf-8")
        _e, warns = VV.validate(work, source_dir=inputs)
        check(any("warehouseArea 38471" in w for w in warns),
              "reconcile: an area the page never shows is flagged")
        _e, warns = VV.validate(work, source_dir=None)
        check(not any("misread" in w for w in warns),
              "reconcile: disengages silently without the inputs folder")


def precedence_cases() -> None:
    """Quality-aware brochure precedence: 'PDF preferred for fields' holds only
    while the PDF parse is reliable - a mostly-poor file loses to a clean twin."""
    print("source precedence (quality-aware brochure demotion):")
    import merge as MG

    def _rec(src, st, **fields):
        r = {"__meta": {"source_file": src, "source_type": st, "locator_base": "x"}}
        r.update(fields)
        return r

    # the garbled PDF's records are mostly poor (no core fields); the PPTX twin is clean
    garbled = _rec("garbled.pdf", "pdf", clearHeight="10 m", park="Alpha")
    clean = _rec("clean.pptx", "pptx", clearHeight="12 m", park="Alpha", city="Pilsen",
                 developer="CTP", status="Existing", warehouseArea=40000,
                 warehouseRentVal=55.0)
    MG.compute_file_quality([garbled, clean])
    merged, _prov, _conf = MG.merge_cluster([garbled, clean])
    check(merged.get("clearHeight") == "12 m",
          "an unreliable PDF loses the spec contest to its clean PPTX twin")

    # control: when BOTH parses are healthy, the static rank holds (PDF wins specs)
    good_pdf = _rec("good.pdf", "pdf", clearHeight="10 m", park="Alpha", city="Pilsen",
                    developer="CTP", status="Existing", warehouseArea=40000,
                    warehouseRentVal=60.0)
    MG.compute_file_quality([good_pdf, clean])
    merged, _prov, _conf = MG.merge_cluster([good_pdf, clean])
    check(merged.get("clearHeight") == "10 m",
          "a healthy PDF keeps its field authority over the PPTX twin (static rank)")
    MG._FILE_UNRELIABLE.clear()  # never leak state into other cases


def units_cases() -> None:
    """Corby round: wide tracker map, header units (source convention KEPT),
    combined lat/lng, collision priorities, rich-tracker authority, yield report,
    UK rent parsing, area-unit tagging, merge unit dominance, cache default."""
    print("units + wide tracker (Corby round):")
    import tempfile
    from openpyxl import Workbook
    import extract_xlsx as X
    import merge as MG
    import extract_pdf as P
    import normalize as NN

    wb = Workbook()
    ws = wb.active
    ws.append(["Building ID", "Marketing Name", "Town", "Latitude, Longitude",
               "Construction status", "Status", "Landlord", "Developer",
               "Size (sq ft)", "Size Unit", "Current quoting rent (£ per sq ft)",
               "Historical quoting rent (£ per sq ft)", "Eaves (m)",
               "Eaves 10m or above?", "Office content (sq ft)",
               "Floor loading (kN/sq m)", "No. of dock level doors",
               "Site area (acres)", "Power (KVA)"])
    ws.append([10179, "EVO 169", "Corby", "52.50304981, -0.650581854",
               "Built", "Available", "EVO", "",
               172867, "GIA", 8.5, 0, 15, "Yes", 13576, 50, 15, 5.2, 800])
    with tempfile.TemporaryDirectory() as td:
        f = Path(td) / "building_data.xlsx"
        wb.save(f)
        res = X.detect_and_extract(f, "Corby", "GB")
    r = res["records"][0] if res["records"] else {}
    check(r.get("lat") == 52.50304981 and r.get("lng") == -0.650581854,
          "xlsx: combined 'Latitude, Longitude' splits into exact first-party coords")
    check(r.get("warehouseArea") == 159291 and r.get("areaUnit") == "sq ft",
          "xlsx: Size Unit=GIA -> WAREHOUSE = GIA - office (172,867 - 13,576 = 159,291), kept in sq ft")
    check(r.get("warehouseRentVal") == 8.5 and r.get("rentUnit") == "£/sq ft/yr"
          and r.get("warehouseRent") == "£8.5 / sq ft / year",
          "xlsx: a £/sq ft quoting rent keeps its OWN convention (never €8.5/m²)")
    # SPLIT (the headline correctness fix): Landlord and Developer are DISTINCT fields.
    # This sheet has Landlord="EVO" + an EMPTY Developer cell -> EVO lands in `landlord`,
    # and developer is left ABSENT (never backfilled from the landlord - the old
    # conflation shipped a landlord as the developer, then split one property into two).
    check(r.get("landlord") == "EVO" and "developer" not in r,
          "xlsx: a Landlord column maps to landlord, NOT developer; an empty Developer is left absent (no backfill)")
    check(r.get("status") == "Available",
          "xlsx: availability Status beats Construction status (alias priority)")
    check(r.get("clearHeight") == "15 m",
          "xlsx: 'Eaves (m)' wins clear height and the header unit suffixes the value")
    check(r.get("officeArea") == "13576 sq ft" and r.get("floorLoad") == "50 kN/sq m",
          "xlsx: office content / floor loading map with their header units")
    check(r.get("plotArea") == round(5.2 * 43560),
          "xlsx: site area in acres converts to sq ft (UK pairing), prov-noted")
    check(r.get("__meta", {}).get("tracker_rich") is True,
          "xlsx: a >=8-field sheet is flagged as a RICH tracker")
    hr = (res.get("header_report") or [{}])[0]
    check(hr.get("mapped_columns", 0) < hr.get("populated_columns", 0)
          and "Size Unit" in hr.get("unmapped_headers", []),
          "xlsx: the yield report names every populated-but-unmapped column")
    # GUARD: a plain 'Size' column with NO GIA/gross marker stays warehouse area - we
    # subtract office ONLY when the size is explicitly a gross total (the user's rule:
    # check the size unit BEFORE assuming warehouse; default to warehouse only after)
    wb3 = Workbook(); ws3 = wb3.active
    ws3.append(["Marketing Name", "Town", "Status", "Developer", "Size (sq ft)",
                "Office content (sq ft)", "Current quoting rent (£ per sq ft)", "Eaves (m)"])
    ws3.append(["Plain Park", "Corby", "Available", "Dev", 100000, 5000, 7, 12])
    with tempfile.TemporaryDirectory() as td:
        f3 = Path(td) / "plain.xlsx"
        wb3.save(f3)
        rp = (X.detect_and_extract(f3, "Corby", "GB").get("records") or [{}])[0]
    check(rp.get("warehouseArea") == 100000 and rp.get("officeArea") == "5000 sq ft",
          "xlsx GUARD: a size with NO GIA marker stays warehouse area (no office subtraction)")

    # merge: dataset unit dominance + display ownership + tracker authority
    check(MG.dominant_units([{"areaUnit": "sq ft", "rentUnit": "£/sq ft/yr"},
                             {"areaUnit": "sq ft"}, {}]) == ("sq ft", "£/sq ft/yr"),
          "merge: dataset units = the dominant source convention")
    c = MG.canonicalize({"warehouseRentVal": 8.5, "rentUnit": "£/sq ft/yr"})
    check(c.get("warehouseRent") == "£8.5 / sq ft / year",
          "merge: display regenerated in the record's OWN rent convention")
    c = MG.canonicalize({"warehouseRent": "£8.50 psf"})
    check(c.get("warehouseRentVal") == 8.5 and c.get("rentUnit") == "£/sq ft/yr",
          "merge: a display-only '£8.50 psf' derives numeric + unit (no € invention)")
    # v14: office rent + area NUMERICS for the total-rent split (derived + traced)
    oc = MG.canonicalize({"officeRent": "€ 9 / sq m / year", "officeArea": "5000 sq m",
                          "warehouseRentVal": 60.0, "warehouseArea": 40000, "rentUnit": "€/sq m/yr"})
    check(oc.get("officeRentVal") == 9.0 and oc.get("officeAreaVal") == 5000,
          "v14: officeRentVal + officeAreaVal derived from the office rent/area strings")
    oc2 = MG.canonicalize({"officeRent": "€ 0.75 / sq m / month",
                           "warehouseRentVal": 60.0, "warehouseArea": 40000})
    check(oc2.get("officeRentVal") == 9.0,
          "v14: a monthly office rent annualises x12 (0.75 -> 9.0)")
    oc3 = MG.canonicalize({"officeArea": "10% of GLA",
                           "warehouseRentVal": 60.0, "warehouseArea": 40000})
    check("officeAreaVal" not in oc3 or oc3.get("officeAreaVal") in (None,),
          "v14: a '% of GLA' office area is NOT mis-parsed as an absolute area")
    # v16: brochure DESCRIPTION harvest - font-size grouped, identity-anchored, boilerplate-rejected
    import fitz as _fitz
    with tempfile.TemporaryDirectory() as td:
        f = Path(td) / "brochure.pdf"
        doc = _fitz.open()
        pg = doc.new_page()
        pg.insert_text((60, 80), "BIG PARK LOGISTICS", fontsize=26)            # ALL-CAPS callout
        pg.insert_text((60, 150), "Alpha Park is strategically located in Corby, within an "
                                  "established industrial estate.", fontsize=11)
        pg.insert_text((60, 172), "The scheme offers excellent connectivity to the A14 and M1, "
                                  "providing efficient distribution across the Midlands.", fontsize=11)
        pg2 = doc.new_page()
        pg2.insert_text((60, 80), "Misrepresentation Act 1967. Every care has been taken; all "
                                  "measurements are approximate.", fontsize=6)  # legal footer
        doc.save(str(f))
        doc.close()
        ddesc, _dp = P.best_description_in_deck(f)
    check(bool(ddesc) and ddesc.startswith("Alpha Park is strategically located")
          and "Misrepresentation" not in ddesc and "BIG PARK LOGISTICS" not in ddesc,
          "v16: best_description_in_deck picks the body description (not the caps callout or legal footer)")
    # bundled sandbox-accelerator wheel (PyMuPDF only): the sandbox lacks PyMuPDF so the
    # native wheel is bundled; Pillow IS available in the sandbox, so NO Pillow wheel ships
    # (keeps the skill under the org upload-size cap). ensure() must NO-OP when the package
    # is already importable (so the host / shim tests are never disturbed).
    import _vendor_wheels as VW
    check(VW.ensure("json") == "system",
          "vendor: ensure() no-ops ('system') when the package is already importable")
    _vd = Path(VW.__file__).resolve().parent.parent / "vendor"
    _whls = sorted(_vd.glob("*.whl")) if _vd.is_dir() else []
    if _whls:
        import zipfile as _zf
        _tops = [{n.split("/")[0] for n in _zf.ZipFile(w).namelist()} for w in _whls]
        check(all(t & {"fitz", "pymupdf"} for t in _tops),
              "vendor: every bundled wheel is a PyMuPDF zip - Pillow ships from the sandbox, "
              "no Pillow wheel (%d shipped)" % len(_whls))
        check(any("fitz" in t for t in _tops),
              "vendor: the native PyMuPDF wheel is present (sandbox lacks PyMuPDF)")
        # _compatible must SELECT the bundled wheel on the cp310/linux/x86_64 sandbox it
        # targets, and must TOLERATE a Path arg: a Path has no .endswith, so passing one
        # would raise -> get swallowed by ensure() -> reject every wheel -> the wheel never
        # loads and decks needlessly route to vision (the real Cowork bug this guards).
        import sysconfig as _sc
        _wn = next(w.name for w in _whls if "pymupdf" in w.name.lower())
        _orig = _sc.get_platform
        try:
            _sc.get_platform = lambda: "linux-x86_64"   # simulate the Cowork sandbox arch
            check(VW._compatible(_wn) is True,
                  "vendor: the bundled cp310-abi3 manylinux wheel is compatible on linux-x86_64 (abi3 forward-compat)")
            check(VW._compatible(_vd / _wn) is True,
                  "vendor: _compatible tolerates a pathlib.Path arg (str-coerced, no .endswith crash)")
            check(VW._compatible("pymupdf-1.27-cp310-abi3-win_amd64.whl") is False,
                  "vendor: a non-linux-arch wheel is rejected")
        finally:
            _sc.get_platform = _orig
    rich = {"clearHeight": "15 m", "park": "Alpha", "city": "Corby", "developer": "EVO",
            "status": "Available", "warehouseArea": 172867,
            "__meta": {"source_file": "tracker.xlsx", "source_type": "xlsx",
                       "tracker_rich": True}}
    pdf_rec = {"clearHeight": "12 m", "park": "Alpha", "city": "Corby", "developer": "EVO",
               "status": "Available", "warehouseArea": 172867, "warehouseRentVal": 55.0,
               "__meta": {"source_file": "brochure.pdf", "source_type": "pdf"}}
    MG._FILE_UNRELIABLE.clear()
    merged, _p, _c = MG.merge_cluster([pdf_rec, rich])
    check(merged.get("clearHeight") == "15 m",
          "merge: a RICH tracker outranks the brochure on structured specs")
    plain = dict(rich, __meta={"source_file": "tracker.xlsx", "source_type": "xlsx"})
    merged, _p, _c = MG.merge_cluster([pdf_rec, plain])
    check(merged.get("clearHeight") == "12 m",
          "merge: a non-rich tracker keeps the classic brochure-first spec order")

    # extract_pdf: UK rent quoting + area-unit tagging
    d, n, _note, u = P._parse_rent("£8.50 per sq ft")
    check(n == 8.5 and u == "£/sq ft/yr" and d == "£8.5 / sq ft / year",
          "pdf: '£8.50 per sq ft' parses in its own convention")
    d, n, _note, u = P._parse_rent("£75 psf")
    check(n is None and d == "£75 psf",
          "pdf: an implausible psf figure ships as text, never a numeric")
    d, n, _note, u = P._parse_rent("€4,20 / m2 / mes")
    check(n == 50.4 and (u or "€/sq m/yr") == "€/sq m/yr",
          "pdf: continental monthly quotes still annualise to €/m²/yr")
    rec, prov = {}, {}
    P._apply_num(rec, prov, "warehouseArea", "131,536 sq ft (12,220 sq m)", "page 1")
    check(rec.get("warehouseArea") == 131536 and rec.get("areaUnit") == "sq ft",
          "pdf: the unit NEAREST the matched number tags the area (sq ft first)")
    rec, prov = {}, {}
    P._apply_num(rec, prov, "warehouseArea", "12 220 m² (131,536 sq ft)", "page 1")
    check(rec.get("warehouseArea") == 12220 and rec.get("areaUnit") == "sq m",
          "pdf: ...and the reverse order tags sq m")

    # band sanity: a £8.5 psf rent is healthy, not 'poor'
    import _common as CC
    check(not CC.record_is_poor({"city": "Corby", "developer": "EVO", "status": "Available",
                                 "warehouseArea": 170000, "warehouseRentVal": 8.5,
                                 "rentUnit": "£/sq ft/yr"}),
          "probe: a normal £/sq ft rent never marks a record poor")

    # merge image cache defaults next to the canonical (manual helper calls
    # survive a capped shell exactly like run.py calls)
    import sys as _sys
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        recs = [{"park": "Alpha", "city": "Corby", "developer": "EVO",
                 "status": "Available", "warehouseArea": 170000,
                 "__meta": {"source_file": "x.pdf", "source_type": "pdf",
                            "locator_base": "page 1", "prov": {}}}]
        rf = td / "recs.json"
        rf.write_text(json.dumps(recs), encoding="utf-8")
        (td / "src").mkdir()
        argv = _sys.argv
        _sys.argv = ["merge.py", "--records", str(rf), "--source-dir", str(td / "src"),
                     "--out", str(td / "canonical.json")]
        try:
            MG.main()
        finally:
            _sys.argv = argv
        check((td / ".image_cache").exists(),
              "merge: the hero-image cache defaults next to the canonical (no flag needed)")
        meta_units = json.loads((td / "canonical.json").read_text(encoding="utf-8"))["meta"]["units"]
        check(meta_units == {"area": "sq m", "rent": "€/sq m/yr"},
              "merge: meta.units defaults metric when no source states a unit")


def description_pick_cases() -> None:
    """#5: brochure DESCRIPTION moves from the EN-keyword/font-size heuristic to an
    LLM pick on the photo-match path, with the heuristic kept as (a) the offline
    fallback + (b) the font_grouped_blocks hint and a DETERMINISTIC verbatim quote
    gate so a hallucinated description can never reach canonical.json.
      (a) font_grouped_blocks yields the groups the refactored heuristic scores AND
          best_description_in_deck's result is byte-identical to before (refactor-safety);
      (b) a --photo-descriptions whose quote IS in the deck -> ACCEPTED, prov '(text interpretation)';
      (c) a --photo-descriptions whose quote is ABSENT (hallucination) -> REJECTED -> heuristic;
      (d) a STALE text_hash -> rejected -> heuristic;
      (e) NO --photo-descriptions -> the description path is byte-identical to today."""
    print("brochure description pick (#5: LLM pick + quote-verify gate + offline fallback):")
    import tempfile as _tf, sys as _sys, json as _json, io as _io
    import fitz as _fitz
    import merge as MG
    import extract_pdf as P
    import images as IMG
    from PIL import Image as _PI

    # two short sentences (each fits one PDF line - PyMuPDF insert_text does not wrap,
    # so an over-long line is clipped at the page edge and loses its full stop)
    S1 = "Northgate 220 is strategically located in Doncaster, within a logistics hub."
    S2 = "The scheme offers excellent connectivity to the A1 and M18 motorway network."

    def _photo_bytes(w=900, h=520):
        im = _PI.new("RGB", (w, h))
        im.putdata([(((x * 255) // w), ((y * 255) // h), (((x + y) * 127) // (w + h)))
                    for y in range(h) for x in range(w)])
        b = _io.BytesIO(); im.save(b, "JPEG", quality=85); return b.getvalue()

    def _build_deck(path):
        doc = _fitz.open()
        pg = doc.new_page()
        pg.insert_text((60, 70), "NORTHGATE LOGISTICS PARK", fontsize=26)   # ALL-CAPS callout
        pg.insert_text((60, 150), S1, fontsize=11)
        pg.insert_text((60, 175), S2, fontsize=11)
        pg.insert_image(_fitz.Rect(60, 230, 540, 510), stream=_photo_bytes())  # a real hero
        pg2 = doc.new_page()
        pg2.insert_text((60, 80), "Misrepresentation Act 1967. Every care has been taken; all "
                        "measurements are approximate. Subject to contract.", fontsize=6)
        doc.save(str(path)); doc.close()

    # (a) refactor-safety: font_grouped_blocks yields the heuristic's groups, and
    # best_description_in_deck is byte-identical to scoring the legacy in-line loop.
    with _tf.TemporaryDirectory() as td:
        f = Path(td) / "deck.pdf"; _build_deck(f)
        blocks = P.font_grouped_blocks(f)
        check(bool(blocks) and all({"page", "size", "text"} <= set(b) for b in blocks)
              and all(isinstance(b["page"], int) and b["page"] >= 1 for b in blocks),
              "a: font_grouped_blocks returns 1-based {page,size,text} groups")

        def _legacy_pick(path):  # the loop best_description_in_deck used BEFORE the refactor
            best = None
            for grp in P.font_grouped_blocks(path):
                pno = grp["page"] - 1
                txt = P._desc_anchor(P._desc_clean(grp["text"]))
                if len(txt) < 120 or txt.count(". ") < 1 or P._DESC_BAD.search(txt):
                    continue
                words = txt.split()
                if sum(1 for w in words if w[:1].islower()) / max(len(words), 1) < 0.5:
                    continue
                score = len(txt) + (300 if P._DESC_GOOD.search(txt) else 0) + max(0, 10 - pno) * 25
                if best is None or score > best[0]:
                    best = (score, pno, P._desc_cap(txt))
            import normalize as _N
            return (_N.clean_value(best[2]), best[1]) if best else (None, None)

        heur = P.best_description_in_deck(f)
        check(heur == _legacy_pick(f) and heur[0] and heur[0].startswith("Northgate 220 is strategically"),
              "a: best_description_in_deck is byte-identical to the legacy font-size pick (refactor-safe)")

    def _merge_with(photo_desc_obj, td):
        """Run merge with a 1-property tracker + a photo-matched brochure; return the
        property's (description, prov['description'].locator). Each call clears the
        module source-resolve memo + the doc cache so a reused brochure basename
        across temp dirs resolves freshly (a fresh process does this implicitly)."""
        MG._SRC_RESOLVE.clear()
        IMG.close_doc_cache()
        f = Path(td) / "Northgate brochure.pdf"; _build_deck(f)
        rec_tr = {"park": "Northgate 220", "city": "Doncaster", "developer": "Acme",
                  "status": "Built", "warehouseArea": 50000,
                  "__meta": {"source_file": "tracker.xlsx", "source_type": "xlsx",
                             "locator_base": "Sheet1", "prov": {"park": "Sheet1!r2"},
                             "tracker_rich": True}}
        rf = Path(td) / "recs.json"; rf.write_text(_json.dumps([rec_tr]), encoding="utf-8")
        pmap = Path(td) / "pm.json"
        pmap.write_text(_json.dumps({MG.match.match_key(rec_tr): "Northgate brochure.pdf"}),
                        encoding="utf-8")
        out = Path(td) / "canonical.json"
        ledger = Path(td) / "source_ledger.csv"
        argv = _sys.argv
        cli = ["merge.py", "--records", str(rf), "--source-dir", str(td),
               "--out", str(out), "--ledger", str(ledger), "--photo-map", str(pmap)]
        if photo_desc_obj is not None:
            pdf = Path(td) / "photo_descriptions.json"
            pdf.write_text(_json.dumps(photo_desc_obj), encoding="utf-8")
            cli += ["--photo-descriptions", str(pdf)]
        _sys.argv = cli
        try:
            MG.main()
        finally:
            _sys.argv = argv
        cj = _json.loads(out.read_text(encoding="utf-8"))
        p0 = cj["properties"][0]
        import csv as _csv
        loc = ""
        with open(ledger, newline="", encoding="utf-8") as fh:
            for row in _csv.DictReader(fh):
                if row.get("field") == "description":
                    loc = row.get("source_locator", "")
        return p0.get("description"), loc

    # the deck text-hash the manifest would stamp (so a valid cache entry matches)
    with _tf.TemporaryDirectory() as td:
        f = Path(td) / "Northgate brochure.pdf"; _build_deck(f)
        good_hash = MG._deck_text_hash(P.font_grouped_blocks(f))

    # (b) ACCEPTED: a slightly reworded but quote-grounded LLM description, quote present
    with _tf.TemporaryDirectory() as td:
        llm_desc = ("Northgate 220 is a prime logistics scheme in Doncaster with direct "
                    "access to the A1(M) and M18 motorway network.")
        obj = {"Northgate brochure.pdf": {"description": llm_desc, "page": 1,
               "quote": "Northgate 220 is strategically located in Doncaster",
               "text_hash": good_hash}}
        desc, loc = _merge_with(obj, td)
        check(desc == llm_desc, "b: a quote-grounded LLM description is ACCEPTED over the heuristic")
        check("text interpretation" in loc and loc.startswith("page 1"),
              "b: the accepted pick's prov locator carries 'page 1 (brochure description, text interpretation)'")

    # (b-locator) verify the prov locator carries the '(text interpretation)' tag - read it
    # straight off _verified_photo_description so the test does not depend on csv parsing
    with _tf.TemporaryDirectory() as td:
        f = Path(td) / "deck.pdf"; _build_deck(f)
        v = MG._verified_photo_description(f, {
            "description": "Reworded prose grounded in the deck.", "page": 1,
            "quote": "Northgate 220 is strategically located", "text_hash": good_hash})
        check(v is not None and v[0] == "Reworded prose grounded in the deck." and v[1] == 0,
              "b: _verified_photo_description accepts a present-quote pick (page 1 -> 0-based)")

    # (c) REJECTED: a hallucinated description whose quote is NOT in the deck -> heuristic
    with _tf.TemporaryDirectory() as td:
        f = Path(td) / "deck.pdf"; _build_deck(f)
        v = MG._verified_photo_description(f, {
            "description": "Fabricated: a 500,000 sq ft mega-shed in Marseille near the port.",
            "page": 1, "quote": "a 500,000 sq ft mega-shed in Marseille", "text_hash": good_hash})
        check(v is None, "c: a hallucinated description (quote absent from the deck) is REJECTED")
    with _tf.TemporaryDirectory() as td:
        obj = {"Northgate brochure.pdf": {
            "description": "Fabricated: a mega-shed in Marseille near the port.",
            "page": 1, "quote": "a mega-shed in Marseille", "text_hash": good_hash}}
        desc, loc = _merge_with(obj, td)
        check(desc and desc.startswith("Northgate 220 is strategically located"),
              "c: the rejected hallucination FALLS BACK to best_description_in_deck (gate holds)")
        check("text interpretation" not in loc,
              "c: the fallback prov locator is the plain heuristic tag, not '(text interpretation)'")

    # (d) STALE text_hash (the deck changed since the pick) -> rejected -> heuristic
    with _tf.TemporaryDirectory() as td:
        f = Path(td) / "deck.pdf"; _build_deck(f)
        v = MG._verified_photo_description(f, {
            "description": "A valid-quote pick but stamped with a stale hash.", "page": 1,
            "quote": "Northgate 220 is strategically located", "text_hash": "deadbeefdeadbeef"})
        check(v is None, "d: a stale text_hash is REJECTED even when the quote is present")
    with _tf.TemporaryDirectory() as td:
        obj = {"Northgate brochure.pdf": {
            "description": "A valid-quote pick but stamped with a stale hash.", "page": 1,
            "quote": "Northgate 220 is strategically located", "text_hash": "deadbeefdeadbeef"}}
        desc, loc = _merge_with(obj, td)
        check(desc and desc.startswith("Northgate 220 is strategically located"),
              "d: a stale-hash pick falls back to the heuristic (re-pick discipline)")

    # (e) NO --photo-descriptions -> byte-identical to today's heuristic path
    with _tf.TemporaryDirectory() as td:
        desc_off, _loc_off = _merge_with(None, td)
    with _tf.TemporaryDirectory() as td:
        f = Path(td) / "Northgate brochure.pdf"; _build_deck(f)
        heur_direct, _hp = P.best_description_in_deck(f)
    check(desc_off and desc_off == heur_direct,
          "e: with NO --photo-descriptions the description == best_description_in_deck (offline byte-identical)")


def batch_a_cases() -> None:
    """Batch A (ship-readiness): Europe-wide xlsx genericity + unreadable-input honesty.
    P0-6 multilingual headers, P1-2 wide-unknown guard, P1-3/P1-4 cell unit + monthly,
    P2-11 land price, P2-1 prior-ledger refusal, P1-1 unreadable + CSV, P2-3 linked URLs."""
    print("batch A (Europe-wide xlsx + unreadable-input honesty):")
    import run as RUN
    from openpyxl import Workbook

    def build(headers, rows, ext=".xlsx"):
        wb = Workbook(); ws = wb.active
        ws.append(headers); [ws.append(r) for r in rows]
        td = Path(tempfile.mkdtemp()); f = td / ("t" + ext); wb.save(f)
        return f, X.detect_and_extract(f, "R", "")

    # P0-6: a continental DE tracker MAPS (was 0 records / questionnaire -> no dashboard)
    _, de = build(["Stadt", "Entwickler", "Hallenfläche (m²)", "Miete (€/m²/Monat)", "Status"],
                  [["Berlin", "GLP", "45000", "4,50", "Verfügbar"]])
    d0 = de["records"][0] if de["records"] else {}
    check(d0.get("city") == "Berlin" and d0.get("warehouseArea") == 45000
          and d0.get("warehouseRentVal") == 54.0 and d0.get("rentUnit") == "€/sq m/yr",
          "P0-6: DE tracker maps (Stadt/Hallenfläche/Miete €/m²/Monat x12) - English-first preserved")

    # P0-6 byte-stability: the real UK Corby tracker is unchanged by the ledger merge
    cor = X.detect_and_extract(Path(r"C:\Claude Projects\Corby Test Run\Building_Data_12_06_2026.xlsx"), "Corby", "GB")["records"]
    check(len(cor) == 4 and cor[0]["warehouseArea"] == 159291 and cor[0]["rentUnit"] == "£/sq ft/yr",
          "P0-6: UK Corby tracker stable (4 records; warehouse = GIA - office = 172,867 - 13,576 = 159,291; £/sq ft)")

    # P1-3 cell-unit fallback / P1-4 cell-monthly x12
    _, p13 = build(["Property", "City", "Quoting rent"], [["Alpha", "Leeds", "£8.50 / sq ft / yr"]])
    check(p13["records"][0].get("rentUnit") == "£/sq ft/yr" and p13["records"][0].get("warehouseRentVal") == 8.5,
          "P1-3: unit-silent header + £/sq ft CELL -> £/sq ft/yr (not the €/m² default)")
    _, p14 = build(["Property", "City", "Rent"], [["Beta", "Lyon", "EUR 4.50 / sq m / month"]])
    check(p14["records"][0].get("warehouseRentVal") == 54.0,
          "P1-4: monthly-in-CELL annualised x12 -> 54.0 / year")

    # P2-11 land price + the collision guards (Warehouse, Land, £/sq ft Price)
    _, p2 = build(["Scheme", "Land price", "Site area"], [["Plot A", "€1,200,000", "5 acres"]])
    check(bool(p2["records"][0].get("landPrice")) and "1,200,000" in p2["records"][0]["landPrice"],
          "P2-11: Land price column -> landPrice")
    check(X._header_field("Warehouse") != "warehouseRentVal" and X._header_field("Land") != "country"
          and X._header_field("Price (£/sq ft)") != "landPrice",
          "P0-6 collisions: bare Warehouse !-> rent, Land !-> country, £/sq ft Price !-> landPrice")

    # P1-2: a wide sheet with unrecognised headers is NOT a questionnaire (no phantom reqs)
    _, wide = build(["Spalte A", "Spalte B", "Spalte C", "Spalte D", "Spalte E"],
                    [["1", "2", "3", "4", "5"], ["6", "7", "8", "9", "10"]])
    check(not wide["records"] and not wide["requirements"]
          and any(h.get("suspected_tracker") for h in wide["header_report"]),
          "P1-2: wide unknown-headers -> suspected_tracker, no phantom requirements")

    # P2-1: a prior-run Source Ledger is refused by its column signature
    _, led = build(["property_id", "record_type", "field", "value", "source_file", "source_type"],
                   [["1", "property", "city", "Corby", "x.pdf", "pdf"]])
    check(not led["records"] and not led["requirements"],
          "P2-1: a Source Ledger sheet is refused (no phantom requirements)")

    # P1-1: corrupt/encrypted/empty inputs are TYPED, never a crash; CSV tracker reads
    import fitz
    td = Path(tempfile.mkdtemp())
    (td / "corrupt.pdf").write_bytes(b"%PDF-1.4 nope \x00")
    (td / "empty.pdf").write_bytes(b"")
    d = fitz.open(); d.new_page(); d.save(td / "enc.pdf", encryption=fitz.PDF_ENCRYPT_AES_256, owner_pw="o", user_pw="u"); d.close()
    check(P.extract(td / "corrupt.pdf", "R", "GB") == [] and P.extract(td / "empty.pdf", "R", "GB") == [],
          "P1-1: extract_pdf returns [] (no traceback) on corrupt / 0-byte")
    check(RUN._classify_unreadable(td / "corrupt.pdf") == "corrupt / unreadable"
          and RUN._classify_unreadable(td / "empty.pdf") == "empty file (0 bytes)"
          and RUN._classify_unreadable(td / "enc.pdf") == "encrypted / password-protected",
          "P1-1: _classify_unreadable types corrupt / empty / encrypted")
    (td / "t.csv").write_text("Property,City,Warehouse area (sq m),Quoting rent\nA,Berlin,40000,5.5\nB,Lyon,25000,6\nC,Madrid,30000,5\n", encoding="utf-8")
    csv_rec = X.detect_and_extract(td / "t.csv")["records"]
    check(len(csv_rec) == 3 and csv_rec[0]["warehouseArea"] == 40000,
          "P1-1: a .csv tracker is read (stdlib) -> 3 records")

    # P2-3: a brochure URL + a hyperlink target are harvested, never silently lost
    wb = Workbook(); ws = wb.active
    ws.append(["Property", "City", "Warehouse area", "Brochure URL"])
    ws.append(["Alpha", "Berlin", 40000, "https://landlord.com/alpha.pdf"])
    ws["E2"].hyperlink = "https://agent.com/alpha-deck.pptx"
    f = Path(tempfile.mkdtemp()) / "links.xlsx"; wb.save(f)
    tg = {x["target"] for x in X.detect_and_extract(f)["linked_sources"]}
    check("https://landlord.com/alpha.pdf" in tg and "https://agent.com/alpha-deck.pptx" in tg,
          "P2-3: brochure URL (string) + hyperlink (object) both harvested into linked_sources")


def batch_b_cases() -> None:
    """Batch B: P0-2 cross-source CONTAINMENT dedup (the tracker+brochure double-count).
    The mandatory regression guards: a postal-address tracker park merges with the
    brochure scheme name; genuinely distinct 'Alpha Park'/'Beta Park' of one developer
    stay separate; phases of materially different size stay separate."""
    print("batch B (cross-source containment dedup):")

    def rec(city, dev, park, area, src):
        return {"city": city, "developer": dev, "park": park, "warehouseArea": area,
                "__meta": {"source_file": src}}

    tr = rec("Corby", "Canmoor", "Unit 1, Raven Park, Earlstree Industrial Estate, Corby, NN17 4XD", 177750, "data.xlsx")
    vis = rec("Corby", "Canmoor", "Raven Park", 177750, "raven.pdf")
    check(match.same_property(tr, vis),
          "P0-2: tracker postal-address park merges with the brochure scheme name (no double-count)")
    check(len(match.dedupe([tr, vis])) == 1,
          "P0-2: 1 tracker + 1 brochure record for one property -> exactly 1 cluster")
    check(not match.same_property(rec("Pilsen", "CTP", "Alpha Park", 40000, "a.pdf"),
                                  rec("Pilsen", "CTP", "Beta Park", 40000, "b.pdf")),
          "P0-2 NEGATIVE: distinct Alpha Park / Beta Park (same dev, city, area) stay separate")
    check(not match.same_property(rec("Corby", "Canmoor", "Raven Park", 177750, "x.pdf"),
                                  rec("Corby", "Canmoor", "Raven Park North", 90000, "y.pdf")),
          "P0-2: same-name phases of materially different size stay separate")
    check(not match.same_property(rec("Corby", "Canmoor", "Raven Park", 177750, "x.pdf"),
                                  rec("Leeds", "GLP", "Raven Park", 177750, "y.pdf")),
          "P0-2: same park name in a DIFFERENT city + developer does NOT merge")

    # P0-1: deck-level hero harvest (for a 0-record brochure matched to a property) +
    # the merge photo-override. The brochure<->property MATCHING is agentic (no rules);
    # these test the deterministic CONSUME mechanics.
    import tempfile as _tf
    import images as IMG
    import fitz
    from PIL import Image as _PILImg
    import io as _io2

    def _noise_jpeg(w=900, h=560):
        # a smooth colour gradient scores as photographic (colourful, non-white,
        # spread across shades) - a self-contained stand-in for a real hero photo
        im = _PILImg.new("RGB", (w, h))
        im.putdata([(((x * 255) // w), ((y * 255) // h), (((x + y) * 127) // (w + h)))
                    for y in range(h) for x in range(w)])
        b = _io2.BytesIO(); im.save(b, format="JPEG", quality=85); return b.getvalue()

    with _tf.TemporaryDirectory() as td:
        td = Path(td)
        doc = fitz.open(); pg = doc.new_page()
        pg.insert_image(fitz.Rect(40, 40, 340, 240), stream=_noise_jpeg())
        photo_pdf = td / "Acme Park brochure.pdf"; doc.save(photo_pdf); doc.close()
        doc2 = fitz.open(); doc2.new_page().insert_text((72, 72), "no photos here, only text")
        text_pdf = td / "textonly.pdf"; doc2.save(text_pdf); doc2.close()
        cache = td / "c"
        hero = IMG.best_hero_in_deck(photo_pdf, 110, cache)
        check(isinstance(hero, str) and hero.startswith("data:image/") and len(hero) > 3000,
              "P0-1: best_hero_in_deck finds the embedded photo across a deck")
        check(IMG.best_hero_in_deck(text_pdf, 110, cache) is None,
              "P0-1: best_hero_in_deck returns None for a photo-less deck (honest placeholder upstream)")
        # GALLERY: the photo-match hero MUST equal gallery_for_deck()[0] - merge sets
        # photo=best_hero_in_deck AND gallery=gallery_for_deck, and the images gate blocks
        # on gallery[0]!=photo. They must be ONE ranking (best_hero delegates to the index).
        gdeck, _gt = IMG.gallery_for_deck(photo_pdf, 110, cache)
        check(gdeck and gdeck[0] == hero,
              "GALLERY: best_hero_in_deck == gallery_for_deck()[0] (photo-match hero/gallery[0] consistent)")
        check(IMG.gallery_for_deck(text_pdf, 110, cache)[0] == [],
              "GALLERY: a photo-less deck yields an empty gallery (placeholder upstream)")
        IMG.close_doc_cache()

        # merge photo-override: a tracker property with no image source gets the matched
        # brochure's deck hero (and its provenance), via --photo-map
        import merge as MG, sys as _sys, json as _json
        rec_tr = {"park": "Acme Park", "city": "Leeds", "developer": "Acme", "status": "Built",
                  "warehouseArea": 50000, "__meta": {"source_file": "tracker.xlsx",
                  "source_type": "xlsx", "locator_base": "Sheet1", "prov": {"park": "Sheet1!r2"},
                  "tracker_rich": True}}
        rf = td / "recs.json"; rf.write_text(_json.dumps([rec_tr]), encoding="utf-8")
        pmap = td / "pm.json"
        pmap.write_text(_json.dumps({MG.match.match_key(rec_tr): "Acme Park brochure.pdf"}), encoding="utf-8")
        out = td / "canonical.json"
        argv = _sys.argv
        _sys.argv = ["merge.py", "--records", str(rf), "--source-dir", str(td),
                     "--out", str(out), "--photo-map", str(pmap)]
        try:
            MG.main()
        finally:
            _sys.argv = argv
        cj = _json.loads(out.read_text(encoding="utf-8"))
        p0 = cj["properties"][0]
        check(p0["photo"].startswith("data:image/") and p0["photo"] != IMG.placeholder(),
              "P0-1: merge attaches the matched brochure's deck hero to the tracker property")
        check(p0["photo"] == hero,
              "P0-1: the attached hero is exactly best_hero_in_deck's result (deterministic)")


def batch_c_cases() -> None:
    """Batch C: the 45s speed cliff. P0-4 skips a pathological vector-art page's
    geometry (the 20-35s pdfplumber page.images access); P0-5 checkpoints the geometry
    to disk; P2-2 drops a coreless own-line record (column-garble) to vision."""
    print("batch C (speed cliff + own-line guard):")
    import tempfile as _tf, glob as _glob, io as _io3
    import images as IMG, extract_pdf as P, fitz
    from PIL import Image as _PI

    def _grad(w=900, h=560):
        im = _PI.new("RGB", (w, h))
        im.putdata([(((x * 255) // w), ((y * 255) // h), (((x + y) * 127) // (w + h)))
                    for y in range(h) for x in range(w)])
        b = _io3.BytesIO(); im.save(b, "JPEG", quality=85); return b.getvalue()

    with _tf.TemporaryDirectory() as td:
        td = Path(td)
        doc = fitz.open()
        pg0 = doc.new_page()
        for i in range(405):  # >_PATHOLOGICAL_IMAGES distinct tiny images = vector-art page
            sm = _PI.new("RGB", (3, 3), (i % 256, (i * 7) % 256, (i * 13) % 256))
            sb = _io3.BytesIO(); sm.save(sb, "PNG")
            x = (i % 20) * 5 + 2; y = (i // 20) * 5 + 2
            pg0.insert_image(fitz.Rect(x, y, x + 4, y + 4), stream=sb.getvalue())
        pg1 = doc.new_page()
        pg1.insert_image(fitz.Rect(40, 40, 340, 240), stream=_grad())
        f = td / "patho.pdf"; doc.save(f); doc.close()
        cache = td / "c"
        layout = IMG._placed_layout(f, cache)
        check(layout["pages"][0] == [],
              "P0-4: a >400-image vector-art page has its geometry SKIPPED (no page.images cliff)")
        check(len(_glob.glob(str(cache / "*.placed.json"))) == 1,
              "P0-5: per-document geometry checkpointed to disk")
        IMG._PLACED_CACHE.clear()
        layout2 = IMG._placed_layout(f, cache)
        check(layout2.get("pages") and layout2["pages"][0] == [],
              "P0-5: geometry reloaded from disk after an in-memory cache clear")
        h1, _pl = IMG.page_hero_and_plan(f, 1, 110, cache_dir=cache)
        check(isinstance(h1, str) and h1.startswith("data:image/"),
              "P0-4: the photo page still yields a hero despite the pathological sibling")
        IMG.close_doc_cache()

    with _tf.TemporaryDirectory() as td:
        td = Path(td)
        doc = fitz.open(); pg = doc.new_page()
        for yy, s in ((100, "Clear height"), (120, "10 m"), (140, "Floor load"), (160, "50 kN")):
            pg.insert_text((72, yy), s)
        f = td / "coreless.pdf"; doc.save(f); doc.close()
        check(P.extract(f, "R", "CZ") == [],
              "P2-2: a coreless own-line page falls to vision (dropped), never a junk card")
        doc = fitz.open(); pg = doc.new_page()
        for yy, s in ((100, "City"), (120, "Pilsen"), (140, "Clear height"), (160, "10 m")):
            pg.insert_text((72, yy), s)
        f2 = td / "cored.pdf"; doc.save(f2); doc.close()
        recs = P.extract(f2, "R", "CZ")
        check(bool(recs) and recs[0].get("city") == "Pilsen",
              "P2-2: an own-line page WITH a core field is still kept")
        IMG.close_doc_cache()


def batch_d_cases() -> None:
    """Batch D: P2-9 enrichment-gate awareness (a requested-but-crashed layer blocks),
    P2-6 curated NUTS-1/2 region aliases -> honest aggregate (never one province
    mislabelled), P2-5 the '??' sentinel filtered from the hero country KPI."""
    print("batch D (gate awareness + NUTS aliases + KPI country):")
    import tempfile as _tf, io as _io4, argparse as _ap
    from contextlib import redirect_stdout as _rso
    import enrich as _E, gate_runner as _G, build_dashboard as _BD

    # P2-6: broad region labels resolve to an HONEST aggregate of their NUTS-3 provinces
    ds = _E._regions_dataset()
    if ds:
        em = _E._dataset_region(ds, "East Midlands")
        check(em and em["nuts"] == "UKF" and em["population"] > 4_000_000
              and isinstance(em.get("unemployment"), (int, float))
              and "aggregated across" in em.get("sources", ""),
              "P2-6: 'East Midlands' -> honest NUTS-1 aggregate (sum pop + weighted unemployment)")
        check(all(_E._dataset_region(ds, n) and _E._dataset_region(ds, n)["nuts"] == "ES51"
                  for n in ("Cataluna", "Catalonia", "Catalunya")),
              "P2-6: native + EN synonyms (Cataluna/Catalonia/Catalunya) all resolve, exact-match")
        check(_E._dataset_region(ds, "Guadalajara") and _E._dataset_region(ds, "Guadalajara")["nuts"] == "ES424",
              "P2-6: a PROVINCE label still resolves to its NUTS-3 profile (unchanged)")
        check(_E._dataset_region(ds, "Nowhereland") is None,
              "P2-6: an unknown broad label stays None (the validate-data block remains the backstop)")

    def _enr(meta, requested):
        p = Path(_tf.mktemp(suffix=".json"))
        p.write_text(json.dumps({"meta": {"client": "X", "hero": {}, "enrichment": meta},
                                 "properties": [], "pois": [], "regions": {}}), encoding="utf-8")
        with _rso(_io4.StringIO()):
            return _G.cmd_enrichment(_ap.Namespace(canonical=str(p), requested=requested))

    check(_enr({}, "geocode,pois") == 1,
          "P2-9: enrichment REQUESTED but the stage left no record (crash) -> BLOCKS")
    check(_enr({}, "") == 0,
          "P2-9: a genuine no-enrichment run still passes (requested set empty)")
    check(_enr({"geocode": True}, "geocode") == 0,
          "P2-9: requested AND stamped -> proceeds (no false block)")

    k = _BD.compute_kpis([{"country": "GB"}, {"country": "??"}, {"country": "GB"}, {"country": "IE"}], {}, None)
    check(k["kpi_countries"] == "2" and "??" not in k["kpi_countries_sub"],
          "P2-5: the '??' sentinel is filtered from the hero country KPI (count + list)")


def batch_f_cases() -> None:
    """Batch F: the quiet 'Done' Gaps note fires for EVERY substantive section
    deliver.gaps_report emits (run._gaps_to_chase), so the broker is never steered away
    from a report with real content - incl. a recorded SOURCE CONFLICT (P3-10)."""
    print("batch F (P3-10 done-note gaps detection):")
    import tempfile as _tf6
    import run as _RUN6
    _full = {"warehouseArea": 1, "warehouseRent": "x", "status": "x", "city": "x",
             "developer": "x", "lat": 1.0, "lng": 1.0, "clearHeight": "x",
             "earlyAccess": "x", "motorway": "x"}  # every CORE field non-tbd

    def _gc(obj, fp=None, pd=None, ui=None, yn=None):
        p = Path(_tf6.mktemp(suffix=".json"))
        p.write_text(json.dumps(obj), encoding="utf-8")
        return _RUN6._gaps_to_chase(p, fp or [], pd or [], ui or [], yn or [])

    clean = {"meta": {}, "properties": [dict(_full)], "pois": [], "regions": {}}
    check(_gc(clean) is False,
          "P3-10: a fully-sourced run with no gaps -> no Gaps note")
    check(_gc({**clean, "meta": {"conflicts": ["id 0 warehouseRent: discarded 9.0"]}}) is True,
          "P3-10: a recorded SOURCE CONFLICT flags the Gaps note (the adversarial-caught miss)")
    check(_gc({**clean, "meta": {"enrichmentGaps": ["x"]}}) is True,
          "P3-10: an enrichment gap flags the note")
    check(_gc({"meta": {}, "properties": [{"city": "x"}], "pois": [], "regions": {}}) is True,
          "P3-10: a tbd CORE field flags the note")
    check(_gc(clean, yn=["unmapped column: Foo"]) is True,
          "P3-10: unmapped tracker columns (yield_notes) flag the note")
    check(_gc(clean, ui=[("x.pdf", "encrypted")]) is True,
          "P3-10: an unreadable input flags the note")


def batch_g_cases() -> None:
    """Batch G (live-run feedback): enrichmentGaps reflect the FINAL state, not the union of
    attempts (#2); a keyless drive-times request surfaces the CAR-vs-truck downgrade as a gap
    (#3a); vision numeric reconciliation is WHOLE-DECK not hero-page (#4, the false-positive
    fix - still catches a genuine misread); a photo-stripped canonical_review.json is emitted
    for the data reviewers (#5a)."""
    print("batch G (live-run feedback fixes):")
    import io as _io7, os as _os7, tempfile as _tf7
    from contextlib import redirect_stdout as _rso7
    import enrich as _E7, run as _RUN7, vision_validate as _VV7

    def _enrich(canon_path, *flags):
        saved = sys.argv
        sys.argv = ["enrich", str(canon_path), *flags, "--cache-dir", str(canon_path.parent)]
        try:
            with _rso7(_io7.StringIO()):
                _E7.main()
        except SystemExit:
            pass
        finally:
            sys.argv = saved

    # #2: a stale gap from pass 1 must NOT linger after pass 2 resolves it
    with _tf7.TemporaryDirectory() as td:
        work = Path(td); canon = work / "canonical.json"
        canon.write_text(json.dumps({"meta": {}, "properties": [
            {"id": 1, "city": "Nowhere City ZZZ", "country": "??"}], "pois": [], "regions": {}}),
            encoding="utf-8")
        _enrich(canon, "--geocode")                       # pass 1: geocode fails (no net/cache)
        g1 = json.loads(canon.read_text(encoding="utf-8"))["meta"]["enrichmentGaps"]
        had_geo_gap = any("geocode" in x.lower() for x in g1)
        d = json.loads(canon.read_text(encoding="utf-8"))
        d["properties"][0]["lat"], d["properties"][0]["lng"] = 52.0, 1.0   # now located
        canon.write_text(json.dumps(d), encoding="utf-8")
        _enrich(canon, "--geocode")                       # pass 2: nothing to geocode (resolved)
        meta2 = json.loads(canon.read_text(encoding="utf-8"))["meta"]
        check(had_geo_gap and not any("could not geocode" in x for x in meta2["enrichmentGaps"]),
              "G#2: a geocode gap a later pass resolved no longer lingers (final state, not union)")
        check("enrichmentGapsByLayer" in meta2,
              "G#2: per-layer gap buckets drive the flat enrichmentGaps")

    # #3a: a keyless --osrm surfaces the CAR-routing downgrade; a keyed one does not
    with _tf7.TemporaryDirectory() as td:
        work = Path(td); canon = work / "canonical.json"
        def _seed():
            canon.write_text(json.dumps({"meta": {}, "properties": [
                {"id": 1, "lat": 52.0, "lng": 1.0, "city": "X", "country": "GB"}],
                "pois": [{"name": "P", "type": "port", "lat": 52.1, "lng": 1.1}], "regions": {}}),
                encoding="utf-8")
        saved_key = _os7.environ.pop("ORS_API_KEY", None)
        try:
            _seed(); _enrich(canon, "--osrm")             # no key in env or args -> car
            gk = json.loads(canon.read_text(encoding="utf-8"))["meta"]["enrichmentGaps"]
            check(any("CAR routing" in x for x in gk),
                  "G#3a: a keyless drive-times request surfaces the car-vs-truck downgrade as a gap")
            _seed(); _enrich(canon, "--osrm", "--ors-key", "TESTKEY")   # HGV -> no car gap
            gh = json.loads(canon.read_text(encoding="utf-8"))["meta"]["enrichmentGaps"]
            check(not any("CAR routing" in x for x in gh),
                  "G#3a GUARD: a keyed (HGV) drive-times request emits NO car-downgrade gap")
        finally:
            if saved_key is not None:
                _os7.environ["ORS_API_KEY"] = saved_key

    # #4: numeric reconciliation is WHOLE-DECK - a value on a different page than the hero
    # does not false-flag; a value absent from the whole deck still flags
    with _tf7.TemporaryDirectory() as td:
        work = Path(td)
        (work / "vision").mkdir(); (work / "extract").mkdir()
        (work / "vision" / "manifest.json").write_text(json.dumps({"decks": [
            {"region": "R", "source_file": "Scan.pdf",
             "pages": [{"page_no": 1}, {"page_no": 2}]}]}), encoding="utf-8")
        _orig_rs, _orig_lt = _VV7._resolve_source, _VV7._load_page_texts
        _VV7._resolve_source = lambda sd, name: Path("dummy.pdf")
        # page 0 = SPEC page (rent 45, area 12000); page 1 = the PHOTO page the hero binds to
        _VV7._load_page_texts = lambda src: ["Quoting rent 45 EUR per sq m; warehouse 12000 sq m",
                                             "Site photograph - building B"]
        try:
            def _wr(records):
                (work / "extract" / "R_vision.json").write_text(json.dumps(records), encoding="utf-8")
            base = {"park": "P", "city": "C", "__meta": {"source_file": "Scan.pdf",
                    "source_type": "pdf", "page_no": 1, "prov": {}}}  # hero on the PHOTO page
            _wr([{**base, "warehouseRentVal": 45.0, "warehouseArea": 12000}])
            _e, _w = _VV7.validate(work, source_dir=work)
            check(not any("digit misread" in x for x in _w),
                  "G#4: a spec value on a different page than the hero does NOT false-flag (whole-deck)")
            _wr([{**base, "warehouseRentVal": 33.0}])      # 33 appears nowhere in the deck
            _e, _w = _VV7.validate(work, source_dir=work)
            check(any("digit misread" in x and "deck" in x for x in _w),
                  "G#4 GUARD: a value absent from the WHOLE deck still flags a suspected misread")
        finally:
            _VV7._resolve_source, _VV7._load_page_texts = _orig_rs, _orig_lt

    # #5a/v11: the FREEZE regenerates a photo-stripped review twin (strips data URIs,
    # keeps data fields, tiny) - exercised through the REAL `gate_runner freeze` path so
    # a manual re-freeze after an out-of-band data fix can never leave the DATA reviewers
    # reading a stale twin (the wasted duplicate-review-round the freeze emission closes).
    with _tf7.TemporaryDirectory() as td:
        work = Path(td); canon = work / "canonical.json"
        big = "data:image/jpeg;base64," + "A" * 6000
        canon.write_text(json.dumps({"meta": {}, "properties": [
            {"id": 1, "photo": big, "plan": big, "warehouseRent": "ok"}], "pois": [], "regions": {}}),
            encoding="utf-8")
        import gate_runner as _GR7
        import io as _io7
        from contextlib import redirect_stdout as _rso7
        _sv7 = sys.argv
        sys.argv = ["gate_runner", "freeze", str(canon)]
        try:
            with _rso7(_io7.StringIO()):
                _GR7.main()
        except SystemExit:
            pass
        finally:
            sys.argv = _sv7
        p = json.loads((work / "canonical_review.json").read_text(encoding="utf-8"))["properties"][0]
        check(p["photo"].startswith("<image stripped") and p["plan"].startswith("<image stripped")
              and p["warehouseRent"] == "ok"
              and (work / "canonical_review.json").stat().st_size < 2000,
              "G#5a/v11: gate_runner freeze regenerates a stripped, tiny canonical_review.json twin")


def batch_i_cases() -> None:
    """Batch I (ORS truck-key UX): osrm_prebake keyless PREFERS a browser-supplied hgv|
    (truck) cache entry over re-routing as car - so a key pasted into the fetcher page
    yields truck times on a keyless re-run, the key never leaving the browser - and the
    web_enrich fetcher page carries the optional truck-key field + a browser-only override."""
    print("batch I (ORS truck-key UX):")
    import io as _io9, tempfile as _tf9
    from contextlib import redirect_stdout as _rso9
    import enrich as _E9, web_enrich as _W9  # noqa: F401 (web_enrich imported for plan)

    with _tf9.TemporaryDirectory() as td:
        work = Path(td); _E9.CACHE_DIR = work
        p = {"id": 1, "lat": 52.0, "lng": 1.0, "city": "X", "country": "GB"}
        poi = {"name": "Port A", "type": "port", "lat": 52.1, "lng": 1.1}
        canon = {"meta": {}, "properties": [p], "pois": [poi], "regions": {}}
        hk = _E9._pair_key(True, 52.0, 1.0, 52.1, 1.1)        # a browser-supplied truck entry
        _E9._save_cache(_E9.OSRM_CACHE, {hk: {"min": 40, "km": 31}})
        gaps, updates = [], []
        with _rso9(_io9.StringIO()):
            done = _E9.osrm_prebake(canon, gaps, "https://router.project-osrm.org",
                                    updates, ors_key="")  # KEYLESS re-run
        dist = (p.get("preBaked") or {}).get("distances", {})
        check(done == 1 and dist.get("Port A", {}).get("min") == 40,
              "I: keyless osrm_prebake reuses a browser-supplied hgv| (truck) cache entry")
        check(any("driving-hgv" in str(u.get("source_locator", "")) for u in updates),
              "I: the ledger labels a browser-supplied truck time as driving-hgv, not car")

    with _tf9.TemporaryDirectory() as td:
        work = Path(td); canon = work / "canonical.json"
        canon.write_text(json.dumps({"meta": {}, "properties": [{"id": 1, "lat": 52.0, "lng": 1.0}],
                         "pois": [{"name": "P", "type": "port", "lat": 52.1, "lng": 1.1}],
                         "regions": {}}), encoding="utf-8")
        saved = sys.argv
        sys.argv = ["web_enrich", "plan", str(canon), "--work", str(work), "--osrm"]
        try:
            with _rso9(_io9.StringIO()):
                _W9.main()
        except SystemExit:
            pass
        finally:
            sys.argv = saved
        page = (work / "web_enrich.html").read_text(encoding="utf-8")
        check('id="orskey"' in page and "CHAIN.ors_key = _k.trim()" in page and "keyrow" in page,
              "I: web_enrich.html carries the optional truck-key field + a browser-only key override")


def batch_j_cases() -> None:
    """Batch J: render_qa's browser-free structural FLOOR for G-visual - passes a complete,
    token-clean dashboard but BLOCKS a structurally broken one (leaked {{token}} / empty
    PROPS / a non-embedded photo), so a no-renderer Cowork run has a real mechanical floor
    (then DEGRADED on the visual dimension) instead of nothing."""
    print("batch J (G-visual static structural floor):")
    import render_qa as _RQ
    good = ("<html><body>CBRE shortlist; L.map(\n"
            'const PROPS = [{"id": 1, "photo": "data:image/jpeg;base64,AAAA"}];\n'
            "const POIS = [];\n"
            "const REGIONS = {};\n"
            "</body></html>\n")
    check(all(ok for ok, _ in _RQ.static_dom_floor(good)),
          "J: a complete, token-clean dashboard passes the static floor")
    check(not all(ok for ok, _ in _RQ.static_dom_floor(good.replace("CBRE", "CBRE {{eyebrow}}"))),
          "J: an unreplaced {{config}} token BLOCKS the floor")
    check(not all(ok for ok, _ in _RQ.static_dom_floor(
            good.replace('[{"id": 1, "photo": "data:image/jpeg;base64,AAAA"}]', "[]"))),
          "J: an empty PROPS array BLOCKS the floor")
    check(not all(ok for ok, _ in _RQ.static_dom_floor(
            good.replace('"data:image/jpeg;base64,AAAA"', '"http://example.com/x.jpg"'))),
          "J: a non-embedded (non data:) property photo BLOCKS the floor")
    # a MALFORMED PROPS array (non-object entries) must FAIL the floor, never CRASH it -
    # the floor classifies a broken file from ANY source (adversarial-caught)
    for malformed in ('[null, {"id": 1, "photo": "data:image/jpeg;base64,AAAA"}]',
                      "[1, 2, 3]", '["a", "b"]', "[true]"):
        try:
            res = _RQ.static_dom_floor(good.replace(
                '[{"id": 1, "photo": "data:image/jpeg;base64,AAAA"}]', malformed))
            ok = not all(o for o, _ in res)   # must report a failure, not raise
        except Exception:
            ok = False                        # a crash is the defect we are guarding against
        check(ok, f"J: a malformed PROPS array {malformed[:14]}... FAILS the floor without crashing")


def batch_h_cases() -> None:
    """Batch H (offline city gazetteer): geocode resolves real European city coordinates +
    country with ZERO network from the bundled assets/cities_dataset.json - the map works in
    Cowork without the exit-8 round-trip, diacritics fold, and an unknown city never gets a
    false pin (honesty preserved)."""
    print("batch H (offline city gazetteer):")
    import io as _io8, tempfile as _tf8
    from contextlib import redirect_stdout as _rso8
    import enrich as _E8

    ll, cc = _E8._gazetteer_lookup("Corby", "GB")
    check(ll is not None and cc == "GB" and 52.0 < ll[0] < 53.0,
          "H: a UK town resolves to real coordinates offline from the gazetteer")
    ll2, cc2 = _E8._gazetteer_lookup("Łódź", "PL")  # 'Łódź' - the Ł/ź fold must match
    check(ll2 is not None and cc2 == "PL",
          "H: a diacritic city name (Lodz, written with the Polish letters) folds and resolves")
    ll3, cc3 = _E8._gazetteer_lookup("Stockholm", "")              # country unknown -> pop-max
    check(ll3 is not None and cc3 == "SE",
          "H: a country-unknown city resolves AND supplies its country (pop-max disambiguation)")
    check(_E8._gazetteer_lookup("Nowhere City ZZZ", "GB")[0] is None,
          "H: an unknown city returns None - never a false pin (honesty)")

    with _tf8.TemporaryDirectory() as td:
        work = Path(td); _E8.CACHE_DIR = work
        _g, _r = _E8._geocode_one, _E8._reverse_cc           # force the network DEAD
        _E8._geocode_one = lambda *a, **k: (_ for _ in ()).throw(RuntimeError("offline"))
        _E8._reverse_cc = lambda *a, **k: (_ for _ in ()).throw(RuntimeError("offline"))
        try:
            canon = {"meta": {}, "properties": [{"id": 1, "city": "Corby", "country": "GB"},
                     {"id": 2, "city": "Stockholm", "country": "??"}], "pois": [], "regions": {}}
            gaps, updates = [], []
            with _rso8(_io8.StringIO()):
                n = _E8.geocode(canon, gaps, updates)
        finally:
            _E8._geocode_one, _E8._reverse_cc = _g, _r
        p1, p2 = canon["properties"]
        check(n == 2 and isinstance(p1.get("lat"), float) and isinstance(p2.get("lat"), float),
              "H: offline geocode fills BOTH cities' coordinates from the gazetteer (no network)")
        check(p2.get("country") == "SE",
              "H: offline geocode fills an unknown country (??->SE) from the gazetteer")
        check(any(u.get("source_type") == "dataset" for u in updates),
              "H: gazetteer coords are cited to the dataset, not Nominatim (G-trace honesty)")


def feedback3_vision_dpi_cases() -> None:
    """Feedback #3: the vision render DPI is raised so small-font SPEC TABLES + KPI badges are
    legible on the FIRST pass (120 caused empty reads + clear-height/area misreads); the stage
    is per-page resumable so the extra cost is bounded by the shell cap, not lost."""
    print("feedback #3 (vision render DPI for dense tables):")
    import inspect
    import vision_prep as _VP
    dpi_default = inspect.signature(_VP.prepare).parameters["dpi"].default
    check(dpi_default >= 180,
          f"F3: vision_prep default dpi raised for legible small-font tables (now {dpi_default}, was 120)")


def feedback1_hero_cases() -> None:
    """Feedback #1: photographic_score demotes non-property imagery via a single-dominant-
    colour penalty - a logo / solid 'photo pending' holding card / road-map screenshot loses
    to a real photo, while a continuous-tone photo (no dominant colour) is untouched. Pure
    pixel statistics, so language- and client-agnostic."""
    print("feedback #1 (hero scoring - demote logos/holding-cards/road-maps):")
    import os as _os
    import images as _IMG
    from PIL import Image as _Img, ImageDraw as _Dw

    solid = _Img.new("RGB", (640, 400), (0, 140, 60))                      # logo / solid card
    photo = _Img.frombytes("RGB", (640, 400), _os.urandom(640 * 400 * 3))  # real photo (no dominant colour)
    s_solid, s_photo = _IMG.photographic_score(solid), _IMG.photographic_score(photo)
    check(s_solid < _IMG.MODEST_PHOTO,
          "F1: a solid logo / 'photo pending' holding card scores below MODEST_PHOTO (demoted)")
    check(s_photo >= _IMG.MODEST_PHOTO and s_photo > s_solid,
          "F1: a real photo clears MODEST_PHOTO and outranks the holding card (unaffected by the penalty)")
    road = _Img.new("RGB", (640, 400), (242, 239, 233))                    # map-paper background
    dd = _Dw.Draw(road)
    for x in range(0, 640, 40):
        dd.line([(x, 0), (x, 400)], fill=(255, 214, 90), width=3)          # 'roads'
    check(_IMG.photographic_score(road) < s_photo,
          "F1: a road-map screenshot (one dominant map-paper colour) loses to a real photo")


def feedback2_region_cases() -> None:
    """Feedback #2: bilingual / dual-name provinces resolve. name_index indexes every variant
    (split on / , ( )) so a property carrying the plain, the local-language OR the joined form
    binds; an ambiguous fragment is dropped (never a false bind); an unresolved code is
    self-documenting (closest known names)."""
    print("feedback #2 (bilingual / dual-name region resolution):")
    import json as _json, tempfile as _tf
    from pathlib import Path as _P
    import build_regions_dataset as _BR
    import enrich as _E

    v = _BR._name_variants("Valencia / València")
    check("valencia" in v and "valencia / valencia" in v,
          "F2: a '/'-joined name yields both the plain and the joined variant")
    check({"bolzano", "bozen"} <= _BR._name_variants("Bolzano (Bozen)"),
          "F2: a parenthetical local name is indexed (Bolzano + Bozen)")
    # uniqueness: a fragment shared by two provinces is DROPPED (no invented bind)
    idx = _BR._build_name_index({"A1": {"name": "North / Nord"}, "A2": {"name": "North / Sud"}})
    check("north" not in idx and idx.get("nord") == ["A1"] and idx.get("sud") == ["A2"],
          "F2: an ambiguous shared fragment is dropped; the unambiguous variants still bind")

    _E._REGIONS_DS = None  # force a fresh load via the gz-aware loader (dataset ships .json.gz)
    ds = _E._regions_dataset()
    for q, code in [("Valencia", "ES523"), ("Alicante", "ES521"), ("Alacant", "ES521"),
                    ("Valencia / València", "ES523")]:
        r = _E._dataset_region(ds, q)
        check(r is not None and r.get("nuts") == code, f"F2: '{q}' resolves to {code}")
    check(_E._dataset_region(ds, "Definitely Not A Province") is None,
          "F2: a genuine non-province still returns None (honest miss, no false bind)")

    with _tf.TemporaryDirectory() as td:           # mis-spelled code -> self-documenting gap
        _E.CACHE_DIR = _P(td)
        canon = {"meta": {}, "properties": [{"id": 1, "regionCode": "Valenca"}], "regions": {}}
        gaps: list = []
        _E.merge_regions(canon, gaps)
        check(any("closest known" in g for g in gaps),
              "F2: an unresolved/mis-spelled regionCode prints the closest known names")


def llm_hero_cases() -> None:
    """LLM-PICKS-THE-HERO scaffolding (the deterministic classifier + G-images gate stay the
    VERIFIER). Covers: interpret_prep emits per-page `candidates` (stable 0-based index +
    written thumbnail); merge binds __meta.heroRef to candidates[index]; a null/absent
    heroRef falls back to the deterministic ladder; the G-images gate still BLOCKS a bound
    non-photo hero (and nonphoto_hero_ok clears it); candidates_for_page and
    embedded_by_index agree on the SAME stable order. Plus a best-effort REAL-DECK check.
    The live multimodal pick runs in Cowork (like text-interpretation quality) and cannot be
    exercised here - this proves the scaffolding it rides on."""
    print("LLM-picks-the-hero (candidate exposure + ref binding; classifier/gate verify):")
    import interpret_prep as IP
    import gate_runner
    try:
        import fitz
    except Exception as e:
        check(False, f"llm-hero: fitz unavailable ({e})"); return
    from PIL import Image, ImageDraw

    def _plan_png() -> bytes:  # flat white paper + colour fills + line work = plan, not photo
        img = Image.new("RGB", (820, 560), (255, 255, 255))
        d = ImageDraw.Draw(img)
        d.rectangle([80, 60, 740, 500], outline=(40, 40, 40), width=4)
        d.rectangle([120, 100, 520, 460], fill=(120, 180, 230))
        d.rectangle([540, 100, 710, 280], fill=(180, 220, 160))
        for x in range(140, 500, 40):
            d.line([(x, 110), (x, 450)], fill=(60, 60, 60), width=2)
        b = __import__("io").BytesIO(); img.save(b, format="PNG"); return b.getvalue()

    def _hero_pdf(td: Path, name: str = "Options - Valencia.pdf") -> tuple:
        """A born-digital text deck whose page 0 carries TWO hero-size candidates: a flat
        PLAN (index 0, larger area -> first in the largest-first order) and a real PHOTO
        (index 1). Returns (path, photo_index) where photo_index is the candidate the LLM
        SHOULD pick. The labels are deliberately outside extract_pdf's dictionary so the
        deck routes to interpretation."""
        doc = fitz.open()
        pg = doc.new_page()
        y = 60
        for ln in ("VALENCIA REGION - Option 1", "City Valencia", "Owner/developer Goodman",
                   "Total existing space 12,500 m2", "Warehouse - Asking rent 4.20 / sqm / month",
                   "This prime logistics warehouse is strategically located near the A-7 "
                   "motorway with excellent connectivity to the Port of Valencia."):
            pg.insert_text((40, y), ln, fontsize=11); y += 20
        pg.insert_image(fitz.Rect(60, 300, 470, 760), stream=_plan_png())        # bigger area
        pg.insert_image(fitz.Rect(60, 80, 360, 250), stream=_noise_photo_jpeg())  # the photo
        f = td / name
        doc.save(str(f)); doc.close()
        # merge._resolve_source memoises by BARE FILENAME across the module's lifetime; each
        # block here reuses one temp dir + filename then deletes it, so clear the memo so the
        # freshly-created file resolves (a stale path to a deleted temp file = a placeholder).
        merge._SRC_RESOLVE.clear()
        # the photo's candidate index = its position in the hero-size filtered, largest-first list
        cfp = IMG.candidates_for_page(f, 0)
        pidx = next((c["index"] for c in cfp if IMG.classify_image(c["img"]) == "photo"), None)
        nidx = next((c["index"] for c in cfp if IMG.classify_image(c["img"]) != "photo"), None)
        IMG.close_doc_cache()
        return f, pidx, nidx

    # --- 1. interpret_prep emits `candidates` (stable index + written thumbnail) ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        f, pidx, nidx = _hero_pdf(td)
        ent = IP.prepare(f, "Valencia", "ES", td / "vis")
        pg0 = ent["pages"][0]
        cands = pg0.get("candidates")
        check(isinstance(cands, list) and len(cands) == 2
              and [c["index"] for c in cands] == [0, 1]
              and all(isinstance(c.get("w"), int) and isinstance(c.get("h"), int) for c in cands),
              "LLM-1: interpret_prep emits per-page candidates with a stable 0-based index")
        check(all(c.get("image") and Path(c["image"]).exists() for c in cands)
              and all(Image.open(c["image"]).size[0] <= IP.CANDIDATE_THUMB_EDGE + 1
                      and Image.open(c["image"]).size[1] <= IP.CANDIDATE_THUMB_EDGE + 1
                      for c in cands),
              "LLM-1: each candidate has a written thumbnail (<= the thumb edge)")

        # --- 5. candidates_for_page and embedded_by_index agree on the SAME order ---
        cfp = IMG.candidates_for_page(f, 0)
        check([c["index"] for c in cfp] == [c["index"] for c in cands]
              and len(cfp) == len(cands),
              "LLM-5: interpret_prep indices == candidates_for_page indices (stable invariant)")
        IMG.close_doc_cache()

    # --- 2. merge binds __meta.heroRef to candidates[index] (bytes == the compressed image) ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        f, pidx, nidx = _hero_pdf(td)
        check(isinstance(pidx, int) and isinstance(nidx, int) and pidx != nidx,
              "LLM-2: fixture has a distinct photo candidate and a non-photo candidate")
        rec = {"park": "Goodman Valencia", "city": "Valencia", "country": "ES",
               "__meta": {"source_file": f.name, "source_type": "pdf", "page_no": 0,
                          "heroRef": pidx, "prov": {}}}
        cache = td / "imgc"
        photo, plan_uri, prec, plr, tried, gal = merge.attach_media(
            [rec], f.parent, IMG.DEFAULT_BUDGET_KB, image_cache=cache)
        expected = IMG.embedded_by_index(f, 0, pidx, IMG.DEFAULT_BUDGET_KB)
        check(isinstance(photo, str) and photo == expected,
              "LLM-2: heroRef binds EXACTLY candidates[index] (embedded_by_index bytes)")
        check(IMG.classify_data_uri(photo) == "photo",
              "LLM-2: the heroRef-bound hero classifies as a real photo")
        check(gal and gal[0] == photo,
              "LLM-2: gallery[0] == the heroRef-bound hero (carousel invariant kept)")
        check("hero chosen by interpretation" in rec["__meta"]["prov"].get("photo", ""),
              "LLM-2: the prov locator records the interpretation pick")
        # planRef binds the plan slot the same way (the LLM picks the site-plan index)
        rec2 = {"park": "P", "city": "Valencia", "country": "ES",
                "__meta": {"source_file": f.name, "source_type": "pdf", "page_no": 0,
                           "heroRef": pidx, "planRef": nidx, "prov": {}}}
        _p2, plan2, *_rest2 = merge.attach_media(
            [rec2], f.parent, IMG.DEFAULT_BUDGET_KB, image_cache=td / "imgc2")
        check(isinstance(plan2, str)
              and plan2 == IMG.embedded_by_index(f, 0, nidx, IMG.DEFAULT_BUDGET_KB),
              "LLM-2: planRef binds EXACTLY candidates[index] into the plan slot")
        check("site plan chosen by interpretation" in rec2["__meta"]["prov"].get("plan", ""),
              "LLM-2: the prov locator records the plan-interpretation pick")
        IMG.close_doc_cache()

    # --- 3. heroRef null/absent -> deterministic classifier pick (existing behaviour) ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        f, pidx, nidx = _hero_pdf(td)
        rec_null = {"park": "X", "__meta": {"source_file": f.name, "source_type": "pdf",
                                            "page_no": 0, "heroRef": None, "prov": {}}}
        photo_n, *_ = merge.attach_media([rec_null], f.parent, IMG.DEFAULT_BUDGET_KB,
                                         image_cache=td / "c1")
        rec_absent = {"park": "X", "__meta": {"source_file": f.name, "source_type": "pdf",
                                              "page_no": 0, "prov": {}}}
        photo_a, *_ = merge.attach_media([rec_absent], f.parent, IMG.DEFAULT_BUDGET_KB,
                                         image_cache=td / "c2")
        # the deterministic ladder picks the page's real photo (Tier A), so a null/absent ref
        # both classify as photo - and the absent-ref hero equals the page ladder's own pick
        det_hero, _det_plan = IMG.page_hero_and_plan(f, 0, IMG.DEFAULT_BUDGET_KB, cache_dir=td / "c2")
        check(IMG.classify_data_uri(photo_n) == "photo"
              and IMG.classify_data_uri(photo_a) == "photo"
              and photo_a == det_hero,
              "LLM-3: heroRef null/absent falls back to the deterministic classifier pick")
        IMG.close_doc_cache()

    # --- 4. the G-images gate still BLOCKS a bound NON-PHOTO hero; nonphoto_hero_ok clears it ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        f, pidx, nidx = _hero_pdf(td)
        rec_bad = {"park": "Bad pick", "city": "Valencia", "country": "ES",
                   "__meta": {"source_file": f.name, "source_type": "pdf", "page_no": 0,
                              "heroRef": nidx, "prov": {}}}
        bad_hero, *_ = merge.attach_media([rec_bad], f.parent, IMG.DEFAULT_BUDGET_KB,
                                          image_cache=td / "cbad")
        check(IMG.classify_data_uri(bad_hero) != "photo",
              "LLM-4: a heroRef pointing at the non-photo candidate yields a non-photo hero")
        # run the real images gate on a canonical carrying that hero
        canon = {"meta": {"client": "X", "hero": {}}, "regions": {}, "pois": [],
                 "properties": [{"id": 1, "photo": bad_hero, "gallery": [bad_hero]}]}
        cpath = td / "canonical.json"
        cpath.write_text(json.dumps(canon), encoding="utf-8")
        check(call(gate_runner, "images", cpath) == 1,
              "LLM-4: G-images BLOCKS the LLM's bad (non-photo) heroRef pick (verifier bites)")
        (td / "placeholder_audit_ack.json").write_text(
            json.dumps({"nonphoto_hero_ok": ["1"]}), encoding="utf-8")
        check(call(gate_runner, "images", cpath) == 0,
              "LLM-4: a recorded nonphoto_hero_ok sign-off clears the bound non-photo hero")
        IMG.close_doc_cache()

    # --- 6. REAL-DECK best-effort: candidates per page + a simulated heroRef binds ---
    real_dir = Path(r"C:\Users\TBaaij\CBRE, Inc\European I&L Occupier Team - "
                    r"IL Occupier Data\06 GtM Projects\TEDi\TEDi Spain\Brochures")
    real_decks = [real_dir / "Sale opportunity_El Morell_CBRE.pdf",
                  real_dir / "CBRE_Valencia_Options_TEDI.pdf"]
    present = [d for d in real_decks if d.exists()]
    if present:
        with tempfile.TemporaryDirectory() as td:
            td = Path(td)
            any_cands = False
            any_bind = False
            for deck in present:
                ent = IP.prepare(deck, "TEDi", "ES", td / "vis")
                if ent.get("mode") != "text":
                    continue  # a raster deck has no candidates payload (expected); skip
                for pg in ent.get("pages", []):
                    cands = pg.get("candidates", [])
                    if not cands:
                        continue
                    any_cands = True
                    # indices are stable + agree with the binder; pick the first candidate
                    cfp = IMG.candidates_for_page(deck, pg["page_no"])
                    if [c["index"] for c in cfp] != [c["index"] for c in cands]:
                        continue
                    uri = IMG.embedded_by_index(deck, pg["page_no"], cands[0]["index"],
                                                IMG.DEFAULT_BUDGET_KB)
                    if isinstance(uri, str) and uri.startswith("data:image/"):
                        any_bind = True
                        break
                if any_bind:
                    break
            check(any_cands,
                  f"LLM-6 (real deck): candidates extracted from {len(present)} TEDi deck(s)")
            check(any_bind,
                  "LLM-6 (real deck): a simulated heroRef binds a real candidate image")
            IMG.close_doc_cache()
    else:
        check(True, "LLM-6 (real deck): skipped - TEDi decks not on this host")


def plan_page_cases() -> None:
    """SITE-PLAN-PAGE binding: a site plan that is VECTOR line-art RENDERED into a whole page
    (not a placed embedded raster - pulled as an image it goes solid black, so planRef /
    page_plan find nothing). The whole-page render is bound to the PLAN SLOT, never the hero.
    Covers: __meta.plan_page render+verify+bind; the deterministic fallback finds a synthetic
    LINE-ART plan page with NO LLM hint; per-property scoping (a neighbour's plan page never
    binds on a multi-property deck); raster-mode plan_page; BACKWARD-COMPAT (no plan_page + no
    plan-like page -> plan stays None); and that the detector is CONSERVATIVE (a photo/blank/
    photo-bearing page never binds the plan slot)."""
    print("site-plan-page render+bind (vector plans a placed-image crop cannot reach):")
    import merge as M
    import images as IMG
    try:
        import fitz
    except Exception as e:
        check(False, f"plan_page: fitz unavailable ({e})"); return
    import io as _io, random, base64
    from PIL import Image

    def _vector_plan_page(pg) -> None:
        """Draw a SITE PLAN as VECTOR content straight into a fitz page (no embedded raster):
        an outer outline, saturated flat unit fills + internal grid line work = the plan
        signature. page_raster renders it; page_embedded_images / candidates_for_page find
        NOTHING (the plan is vector, not a placed image), which is exactly why planRef misses
        it and the whole-page render is needed."""
        pg.draw_rect(fitz.Rect(60, 50, 540, 380), color=(0.1, 0.1, 0.1), width=3)
        pg.draw_rect(fitz.Rect(90, 80, 300, 350), color=(0.1, 0.1, 0.1), fill=(0.45, 0.65, 0.85), width=2)
        pg.draw_rect(fitz.Rect(320, 80, 510, 220), color=(0.1, 0.1, 0.1), fill=(0.6, 0.8, 0.55), width=2)
        pg.draw_rect(fitz.Rect(320, 240, 510, 350), color=(0.1, 0.1, 0.1), fill=(0.85, 0.8, 0.55), width=2)
        for x in range(110, 290, 30):
            pg.draw_line(fitz.Point(x, 90), fitz.Point(x, 340), color=(0.2, 0.2, 0.2), width=1)
        for y in range(100, 340, 30):
            pg.draw_line(fitz.Point(95, y), fitz.Point(295, y), color=(0.2, 0.2, 0.2), width=1)

    def _photo_jpeg(seed=7, w=800, h=450) -> bytes:
        rnd = random.Random(seed)
        img = Image.new("RGB", (64, 36))
        img.putdata([(rnd.randrange(256), rnd.randrange(256), rnd.randrange(256)) for _ in range(64 * 36)])
        img = img.resize((w, h))
        b = _io.BytesIO(); img.save(b, format="JPEG", quality=85); return b.getvalue()

    def _is_plan_uri(uri) -> bool:
        return isinstance(uri, str) and uri.startswith("data:image/")

    BUD = IMG.DEFAULT_BUDGET_KB

    # --- 1. __meta.plan_page render+verify+bind: a property page (photo) + a VECTOR plan page;
    #        the LLM names plan_page=1; merge renders + binds it to the plan slot (NOT the hero) ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        doc = fitz.open()
        pg0 = doc.new_page(width=600, height=800)  # property page: text + a real photo
        pg0.insert_text((40, 50), "City Pilsen Developer CTP Warehouse Area 40000 sq m", fontsize=11)
        pg0.insert_image(fitz.Rect(60, 90, 460, 360), stream=_photo_jpeg(1))
        pg1 = doc.new_page(width=600, height=420)   # the VECTOR site-plan page
        _vector_plan_page(pg1)
        f = td / "PlanDeck.pdf"
        doc.save(str(f)); doc.close()
        M._SRC_RESOLVE.clear(); IMG.close_doc_cache()
        rec = {"park": "P", "city": "Pilsen", "country": "CZ",
               "__meta": {"source_file": f.name, "source_type": "pdf", "page_no": 0,
                          "plan_page": 1, "prov": {}}}
        cl = [rec]
        foreign = M.build_foreign_pages([cl], td)
        photo, plan_uri, prec, plr, tried, gal = M.attach_media(
            cl, td, BUD, image_cache=td / "c1", foreign_pages=foreign[0])
        check(_is_plan_uri(plan_uri) and IMG.classify_data_uri(plan_uri) != "photo",
              "PLAN-1: __meta.plan_page renders the vector plan page and binds it to the plan slot")
        check(plan_uri == IMG.page_render_plan(f, 1, BUD, cache_dir=td / "c1"),
              "PLAN-1: the bound plan == page_render_plan(plan_page) bytes (cached, deterministic)")
        check(_is_plan_uri(photo) and photo != plan_uri,
              "PLAN-1: the page-0 photo is the hero; the vector plan never becomes the hero")
        check("site plan page render chosen by interpretation"
              in rec["__meta"]["prov"].get("plan", ""),
              "PLAN-1: the prov locator records the plan_page interpretation pick")
        IMG.close_doc_cache()

    # --- 2. DETERMINISTIC FALLBACK: NO plan_page hint, but the property owns a vector plan
    #        page -> best_plan_page_render finds it and binds the plan slot ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        doc = fitz.open()
        pg0 = doc.new_page(width=600, height=800)
        pg0.insert_text((40, 50), "City Brno Developer VGP Warehouse Area 25000 sq m", fontsize=11)
        pg0.insert_image(fitz.Rect(60, 90, 460, 360), stream=_photo_jpeg(2))
        pg1 = doc.new_page(width=600, height=420)
        _vector_plan_page(pg1)
        f = td / "FallbackDeck.pdf"
        doc.save(str(f)); doc.close()
        M._SRC_RESOLVE.clear(); IMG.close_doc_cache()
        # NO plan_page; the property claims both its pages via image_pages (its own scope)
        rec = {"park": "Q", "city": "Brno", "country": "CZ",
               "__meta": {"source_file": f.name, "source_type": "pdf", "page_no": 0,
                          "image_pages": [0, 1], "prov": {}}}
        cl = [rec]
        foreign = M.build_foreign_pages([cl], td)
        photo, plan_uri, prec, plr, tried, gal = M.attach_media(
            cl, td, BUD, image_cache=td / "c2", foreign_pages=foreign[0])
        check(_is_plan_uri(plan_uri) and IMG.classify_data_uri(plan_uri) != "photo",
              "PLAN-2: the deterministic fallback finds the vector LINE-ART plan page (no LLM hint)")
        uri_det, pno_det = IMG.best_plan_page_render(f, [0, 1], BUD, td / "c2")
        check(pno_det == 1 and plan_uri == uri_det,
              "PLAN-2: the fallback picks page 1 (the plan page), not page 0 (the photo page)")
        check("detected" in prec.get("__meta", {}).get("prov", {}).get("plan", "")
              or "detected" in (plr or {}).get("__meta", {}).get("prov", {}).get("plan", ""),
              "PLAN-2: the prov locator records the deterministic detection")
        IMG.close_doc_cache()

    # --- 3. PER-PROPERTY SCOPING: a multi-property deck where the plan page is B's, never A's.
    #        Property A (page 0) must NOT bind B's plan page (page 2) - the foreign-page guard
    #        scopes the fallback to A's own pages only. ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        doc = fitz.open()
        pgA = doc.new_page(width=600, height=800)   # A's property page (photo, page 0)
        pgA.insert_text((40, 50), "Property A City Pilsen", fontsize=11)
        pgA.insert_image(fitz.Rect(60, 90, 460, 360), stream=_photo_jpeg(3))
        pgB = doc.new_page(width=600, height=800)   # B's property page (photo, page 1)
        pgB.insert_text((40, 50), "Property B City Brno", fontsize=11)
        pgB.insert_image(fitz.Rect(60, 90, 460, 360), stream=_photo_jpeg(4))
        pgBplan = doc.new_page(width=600, height=420)  # B's VECTOR plan page (page 2)
        _vector_plan_page(pgBplan)
        f = td / "MultiPropPlan.pdf"
        doc.save(str(f)); doc.close()
        M._SRC_RESOLVE.clear(); IMG.close_doc_cache()
        # A over-claims page 2 (B's plan) in image_pages; B anchors page 1 and claims [1, 2].
        clA = [{"park": "A", "__meta": {"source_file": f.name, "source_type": "pdf",
                                        "page_no": 0, "image_pages": [0, 2], "prov": {}}}]
        clB = [{"park": "B", "__meta": {"source_file": f.name, "source_type": "pdf",
                                        "page_no": 1, "image_pages": [1, 2], "prov": {}}}]
        clusters = [clA, clB]
        foreign = M.build_foreign_pages(clusters, td)
        # page 2 is claimed by BOTH but anchored by NEITHER (no record has page_no == 2);
        # it is contested -> foreign to both, so NEITHER property binds it via the fallback.
        _phA, planA, *_a = M.attach_media(clA, td, BUD, image_cache=td / "c3", foreign_pages=foreign[0])
        _phB, planB, *_b = M.attach_media(clB, td, BUD, image_cache=td / "c3", foreign_pages=foreign[1])
        check(planA is None,
              "PLAN-3: property A never binds the CONTESTED plan page (per-property scoping holds)")
        check(planB is None,
              "PLAN-3: property B never binds the contested plan page either (no cross-leak)")
        # but B alone, uniquely claiming its own plan page, DOES bind it (the scope allows own pages)
        clB2 = [{"park": "B", "__meta": {"source_file": f.name, "source_type": "pdf",
                                         "page_no": 1, "image_pages": [1, 2], "prov": {}}}]
        f2only = M.build_foreign_pages([clB2], td)
        _ph, planB2, *_ = M.attach_media(clB2, td, BUD, image_cache=td / "c3b", foreign_pages=f2only[0])
        check(_is_plan_uri(planB2),
              "PLAN-3: B as the SOLE claimant of its own plan page DOES bind it (own-page scope)")
        IMG.close_doc_cache()

    # --- 3b. The LLM plan_page HINT is per-property scoped too: a plan_page pointing at a page
    #         uniquely owned by ANOTHER property of the same deck is FOREIGN and must NOT bind
    #         (the hint path uses the same foreign guard as the deterministic fallback). ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        doc = fitz.open()
        pa = doc.new_page(width=600, height=800)    # A anchor (photo, page 0)
        pa.insert_text((40, 50), "Property A City Pilsen", fontsize=11)
        pa.insert_image(fitz.Rect(60, 90, 460, 360), stream=_photo_jpeg(3))
        pb = doc.new_page(width=600, height=800)    # B anchor (photo, page 1)
        pb.insert_text((40, 50), "Property B City Brno", fontsize=11)
        pb.insert_image(fitz.Rect(60, 90, 460, 360), stream=_photo_jpeg(4))
        pbplan = doc.new_page(width=600, height=420)   # B's VECTOR plan page (page 2)
        _vector_plan_page(pbplan)
        f = td / "MultiPropPlanHint.pdf"
        doc.save(str(f)); doc.close()
        M._SRC_RESOLVE.clear(); IMG.close_doc_cache()
        # B uniquely claims page 2 (its plan) via image_pages; A does NOT claim it but
        # erroneously sets plan_page = 2 (a neighbour's page). The guard must reject A's hint.
        clA = [{"park": "A", "__meta": {"source_file": f.name, "source_type": "pdf",
                                        "page_no": 0, "image_pages": [0], "plan_page": 2, "prov": {}}}]
        clB = [{"park": "B", "__meta": {"source_file": f.name, "source_type": "pdf",
                                        "page_no": 1, "image_pages": [1, 2], "plan_page": 2, "prov": {}}}]
        foreign = M.build_foreign_pages([clA, clB], td)
        offl = M.plan_offlimits_pages([clA, clB], td)
        _pa2, planA, *_a = M.attach_media(clA, td, BUD, image_cache=td / "c3h",
                                          foreign_pages=foreign[0], plan_offlimits=offl[0])
        _pb2, planB, *_b = M.attach_media(clB, td, BUD, image_cache=td / "c3h",
                                          foreign_pages=foreign[1], plan_offlimits=offl[1])
        check(planA is None,
              "PLAN-3b: an LLM plan_page at a NEIGHBOUR's page is rejected (hint path is scoped, no leak)")
        check(_is_plan_uri(planB),
              "PLAN-3b: the owning property B binds its own plan_page via the hint (own-page scope)")
        IMG.close_doc_cache()

    # --- 4. RASTER-mode plan_page: a record from a raster deck names plan_page; merge renders
    #        + binds it identically (the binding is engine-agnostic, driven by the integer) ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        doc = fitz.open()
        pg0 = doc.new_page(width=600, height=420)   # the vector plan page IS page 0 here
        _vector_plan_page(pg0)
        f = td / "RasterPlan.pdf"
        doc.save(str(f)); doc.close()
        M._SRC_RESOLVE.clear(); IMG.close_doc_cache()
        rec = {"park": "R", "__meta": {"source_file": f.name, "source_type": "pdf",
                                       "page_no": 0, "plan_page": 0, "prov": {}}}
        cl = [rec]
        foreign = M.build_foreign_pages([cl], td)
        photo, plan_uri, *_ = M.attach_media(cl, td, BUD, image_cache=td / "c4", foreign_pages=foreign[0])
        check(_is_plan_uri(plan_uri),
              "PLAN-4: a raster-mode record's plan_page renders + binds the plan slot")
        IMG.close_doc_cache()

    # --- 5. BACKWARD-COMPAT: NO plan_page + NO plan-like page -> plan stays None (== today).
    #        A property whose only page is a photo page (no vector plan, an embedded photo) must
    #        keep plan None: the deterministic fallback is CONSERVATIVE and never fabricates. ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        doc = fitz.open()
        pg0 = doc.new_page(width=600, height=800)
        pg0.insert_text((40, 50), "City Lyon Developer Dev Warehouse Area 30000 sq m", fontsize=11)
        pg0.insert_image(fitz.Rect(60, 90, 460, 360), stream=_photo_jpeg(5))
        f = td / "PhotoOnly.pdf"
        doc.save(str(f)); doc.close()
        M._SRC_RESOLVE.clear(); IMG.close_doc_cache()
        rec = {"park": "S", "city": "Lyon", "__meta": {"source_file": f.name, "source_type": "pdf",
                                                       "page_no": 0, "image_pages": [0], "prov": {}}}
        cl = [rec]
        foreign = M.build_foreign_pages([cl], td)
        photo, plan_uri, *_ = M.attach_media(cl, td, BUD, image_cache=td / "c5", foreign_pages=foreign[0])
        check(plan_uri is None,
              "PLAN-5: NO plan_page + a photo-only page -> the plan slot stays None (backward-compat)")
        u, p = IMG.best_plan_page_render(f, [0], BUD, td / "c5")
        check(u is None and p is None,
              "PLAN-5: best_plan_page_render returns (None, None) on a photo page (never fabricates)")
        IMG.close_doc_cache()

    # --- 6. the conservative detector primitives directly (photo / blank / photo-bearing pages
    #        never bind; the vector plan page does) ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        # a vector plan page (0 embedded images) binds; a blank page does not
        doc = fitz.open()
        pgp = doc.new_page(width=600, height=420); _vector_plan_page(pgp)
        doc.new_page(width=600, height=420)  # page 1: BLANK
        f = td / "PlanAndBlank.pdf"
        doc.save(str(f)); doc.close()
        M._SRC_RESOLVE.clear(); IMG.close_doc_cache()
        crop, sig, kind = IMG._rendered_plan_crop(f, 0)
        check(kind == "plan" and 0.15 <= sig.get("white", 0) <= 0.90,
              "PLAN-6: the vector plan page renders with the plan signature (kind 'plan', balanced white)")
        u0, _ = IMG.best_plan_page_render(f, [0], BUD, td / "c6")
        check(_is_plan_uri(u0), "PLAN-6: the vector plan page (no embedded image) binds the fallback")
        cropb, sigb, kindb = IMG._rendered_plan_crop(f, 1)
        check(cropb is None,
              "PLAN-6: a BLANK page renders no ink bbox -> (None) (never a plan)")
        u1, _ = IMG.best_plan_page_render(f, [1], BUD, td / "c6")
        check(u1 is None, "PLAN-6: a blank page never binds the plan slot")
        IMG.close_doc_cache()

    # --- 7. determinism: the rendered-plan URI is cached per (source, page, budget) and a
    #        second render is byte-identical (a pure function of the inputs) ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        doc = fitz.open(); pg0 = doc.new_page(width=600, height=420); _vector_plan_page(pg0)
        f = td / "DetPlan.pdf"
        doc.save(str(f)); doc.close()
        M._SRC_RESOLVE.clear(); IMG.close_doc_cache()
        a = IMG.page_render_plan(f, 0, BUD, cache_dir=td / "cd")
        b = IMG.page_render_plan(f, 0, BUD, cache_dir=td / "cd")   # cache hit
        c = IMG.page_render_plan(f, 0, BUD, cache_dir=td / "cd2")  # fresh cache, recompute
        check(_is_plan_uri(a) and a == b == c,
              "PLAN-7: the rendered-plan URI is deterministic + cache-stable across runs")
        IMG.close_doc_cache()


def image_pages_carousel_cases() -> None:
    """IMAGE-PAGES CAROUSEL: __meta.image_pages widens the per-property carousel to ALL
    of a property's pages, with the deterministic unique-claimant guard (build_foreign_pages)
    enforcing that every deck page feeds AT MOST ONE property. Mirrors the INTERP-4
    hand-written-records pattern: hand-build records with __meta, run merge.attach_media
    (fed the precomputed foreign_pages), assert the gallery. Five topologies:
      (a) multi-page SINGLE-property -> gallery len > 1, hero at [0], <= GALLERY_MAX, only this deck;
      (b) multi-property ONE page each -> each gallery holds only its own page (no neighbour);
      (c) multi-property MULTI page each -> each gallery holds only its uniquely-claimed pages,
          AND a page in A's image_pages that is B's anchor (or contested) is EXCLUDED from A;
      (d) single-page single-property unchanged;
      (e) NO image_pages anywhere -> gallery byte-identical to the page_no-only behaviour."""
    print("image_pages carousel (widened scope + unique-claimant anti-leak guard):")
    import tempfile
    import merge as M
    import images as IMG
    try:
        import fitz
    except Exception as e:
        check(False, f"image_pages: fitz unavailable ({e})"); return
    from PIL import Image

    def _seeded_photo(seed: int, w: int = 360, h: int = 240) -> bytes:
        """A distinct synthetic 'photo' per seed (so the gallery dedup-by-URI keeps each
        page's photo as its own entry). Colourful noise scores high as a photo."""
        import random
        rnd = random.Random(seed)
        img = Image.new("RGB", (w, h))
        img.putdata([(rnd.randrange(256), rnd.randrange(256), rnd.randrange(256))
                     for _ in range(w * h)])
        buf = __import__("io").BytesIO()
        img.save(buf, format="JPEG", quality=85)
        return buf.getvalue()

    def _deck(td: Path, name: str, n_pages: int) -> Path:
        """An n-page PDF, ONE distinct photo per page (seeded by page index)."""
        doc = fitz.open()
        for p in range(n_pages):
            pg = doc.new_page(width=600, height=800)
            pg.insert_text((40, 50), f"Page {p} option text", fontsize=11)
            pg.insert_image(fitz.Rect(60, 80, 460, 360), stream=_seeded_photo(p + 1))
        f = td / name
        doc.save(str(f)); doc.close()
        M._SRC_RESOLVE.clear()  # memoises by bare filename; clear so the fresh temp file resolves
        IMG.close_doc_cache()
        return f

    def _rec(f: Path, page_no: int, image_pages=None, park: str = "P") -> dict:
        m = {"source_file": f.name, "source_type": "pdf", "page_no": page_no, "prov": {}}
        if image_pages is not None:
            m["image_pages"] = image_pages
        return {"park": park, "__meta": m}

    BUD = IMG.DEFAULT_BUDGET_KB

    # --- (a) multi-page SINGLE-property: image_pages = pages 0..3, appendix page 4 OMITTED ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        f = _deck(td, "SingleMulti.pdf", 5)  # pages 0-4; 4 is the appendix (not listed)
        cl = [_rec(f, 0, image_pages=[0, 1, 2, 3], park="Solo")]
        clusters = [cl]
        foreign = M.build_foreign_pages(clusters, td)
        photo, _pl, _pr, _plr, _tr, gal = M.attach_media(
            cl, td, BUD, image_cache=td / "ca", foreign_pages=foreign[0])
        # the appendix page-4 photo, never listed, never a candidate
        app_uri = IMG.embedded_by_index(f, 4, 0, BUD)
        check(len(gal) > 1 and gal[0] == photo and len(gal) <= IMG.GALLERY_MAX,
              "IMGP-a: multi-page single-property gallery has >1 image, hero at [0], <= GALLERY_MAX")
        check(app_uri not in gal,
              "IMGP-a: the OMITTED appendix page's photo is never in the carousel")
        check(not foreign[0],
              "IMGP-a: single-property deck (count==1) has NO foreign pages")
        IMG.close_doc_cache()

    # --- (b) multi-property ONE page each: A=page0, B=page1; no neighbour leak ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        f = _deck(td, "MultiOne.pdf", 2)
        clA = [_rec(f, 0, image_pages=[0], park="A")]
        clB = [_rec(f, 1, image_pages=[1], park="B")]
        clusters = [clA, clB]
        foreign = M.build_foreign_pages(clusters, td)
        _phA, *_a, galA = M.attach_media(clA, td, BUD, image_cache=td / "cb", foreign_pages=foreign[0])
        _phB, *_b, galB = M.attach_media(clB, td, BUD, image_cache=td / "cb", foreign_pages=foreign[1])
        a0, b1 = IMG.embedded_by_index(f, 0, 0, BUD), IMG.embedded_by_index(f, 1, 0, BUD)
        check(a0 in galA and b1 not in galA,
              "IMGP-b: property A's gallery holds ONLY its own page (no neighbour photo)")
        check(b1 in galB and a0 not in galB,
              "IMGP-b: property B's gallery holds ONLY its own page (no neighbour photo)")
        IMG.close_doc_cache()

    # --- (c) multi-property MULTI page each + a contested page B-anchors: excluded from A ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        f = _deck(td, "MultiMulti.pdf", 5)
        # A anchors page 0, claims [0,1,2] AND over-claims page 3 (which is B's ANCHOR).
        # B anchors page 3, claims [3,4]. Page 3 must go to its anchor owner B, never A.
        clA = [_rec(f, 0, image_pages=[0, 1, 2, 3], park="A")]
        clB = [_rec(f, 3, image_pages=[3, 4], park="B")]
        clusters = [clA, clB]
        foreign = M.build_foreign_pages(clusters, td)
        _phA, *_a, galA = M.attach_media(clA, td, BUD, image_cache=td / "cc", foreign_pages=foreign[0])
        _phB, *_b, galB = M.attach_media(clB, td, BUD, image_cache=td / "cc", foreign_pages=foreign[1])
        p1, p2, p3, p4 = (IMG.embedded_by_index(f, n, 0, BUD) for n in (1, 2, 3, 4))
        check(p1 in galA and p2 in galA,
              "IMGP-c: property A's gallery holds its uniquely-claimed pages 1 and 2")
        check(p3 not in galA,
              "IMGP-c: page 3 (B's ANCHOR, over-claimed by A) is EXCLUDED from A (no leak)")
        check(p3 in galB and p4 in galB and p1 not in galB and p2 not in galB,
              "IMGP-c: property B's gallery holds only its own pages 3 and 4")
        check(bool(foreign[0]) and set().union(*foreign[0].values()) == {3},
              "IMGP-c: the guard marks page 3 (and only page 3) FOREIGN to A")
        IMG.close_doc_cache()

    # --- (d) single-page single-property unchanged (image_pages omitted == today) ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        f = _deck(td, "Single.pdf", 1)
        cl = [_rec(f, 0, park="Solo")]  # no image_pages
        foreign = M.build_foreign_pages([cl], td)
        photo, *_d, gal = M.attach_media(cl, td, BUD, image_cache=td / "cd", foreign_pages=foreign[0])
        check(gal and gal[0] == photo and len(gal) >= 1 and len(gal) <= IMG.GALLERY_MAX,
              "IMGP-d: single-page single-property gallery unchanged (hero first)")
        IMG.close_doc_cache()

    # --- (e) NO image_pages anywhere -> byte-identical to the page_no-only behaviour ---
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        f = _deck(td, "NoIP.pdf", 3)
        clA = [_rec(f, 0, park="A")]
        clB = [_rec(f, 1, park="B")]
        clusters = [clA, clB]
        foreign = M.build_foreign_pages(clusters, td)
        check(all(not fm for fm in foreign),
              "IMGP-e: with NO image_pages, foreign_pages is empty for every cluster")
        # the page_no-only baseline: attach_media with foreign_pages omitted (today's default)
        _pa1, *_e1, gal_with = M.attach_media(clA, td, BUD, image_cache=td / "ce", foreign_pages=foreign[0])
        _pa2, *_e2, gal_base = M.attach_media(clA, td, BUD, image_cache=td / "ce")
        check(gal_with == gal_base,
              "IMGP-e: gallery is byte-identical with/without the guard when image_pages is absent")
        IMG.close_doc_cache()


def geofix_cases() -> None:
    """Regression guards for the geographic-enrichment fixes (committee bugs #1-#5 + #D and the
    adversarial-review follow-ups R1-C1 / R4-I1). With the REAL bundled datasets loaded: nearest
    border + city are nearest-of-complete-set with population shown; the geocode/cache/map-link
    fixes hold. Without these, a refactor could silently re-open the confirmed bugs."""
    print("geo fixes (nearest border/city + population, geocode #3/#4/#5, map-link #D):")
    import enrich
    import tempfile
    from pathlib import Path

    class _FakePage:  # minimal page exposing get_links() for _page_link_coords URL parsing
        def __init__(self, uri): self._u = uri
        def get_links(self): return [{"uri": self._u}]

    # nearest BORDER + CITY from the complete datasets (bugs #1/#2), population in the city note
    enrich._DATASET = enrich._BORDERS = enrich._CITY_DATASET = None  # real bundled datasets
    with tempfile.TemporaryDirectory() as td:
        enrich.CACHE_DIR = Path(td)
        canonical = {"properties": [{"id": 1, "lat": 51.10, "lng": 17.03}],  # Wrocław, PL
                     "pois": [], "meta": {}}
        enrich.attach_pois(canonical, [])
        by_type = {p["type"]: p for p in canonical["pois"]}
        b, c = by_type.get("border"), by_type.get("city")
        check(bool(b) and "CBRE border dataset" in b["note"] and b.get("country"),
              "geofix: nearest BORDER from the complete border dataset (not the curated library)")
        check(bool(c) and "CBRE cities dataset" in c["note"] and isinstance(c.get("population"), int)
              and c["population"] >= 100000 and "pop " in c["note"],
              "geofix: nearest CITY is a >=100k city from the complete dataset, population in the note")

    # bug #5: _cache_lookup refuses a wrong-country cache entry under a KNOWN country
    _cache = {"toledo|us": {"latlng": [41.6528, -83.5379], "cc": "US"}}
    check(enrich._cache_lookup(_cache, "Toledo", "ES") == (None, ""),
          "geofix #5: a KNOWN country + exact miss does NOT adopt a different-country cache entry")
    check(enrich._cache_lookup(_cache, "Toledo", "??")[0] == [41.6528, -83.5379],
          "geofix #5: an UNKNOWN country still uses the cross-country prefix cache (seed pattern)")

    # bug #3: an ambiguous bare name misses without a dominant, resolves to the right country with one
    enrich._GAZETTEER = None; enrich._GAZETTEER_MULTI = None
    check(enrich._gazetteer_lookup("Halle", "??") == (None, ""),
          "geofix #3: an ambiguous bare name (Halle DE/BE) w/o a dominant country is an honest MISS")
    _ll, _cc = enrich._gazetteer_lookup("Halle", "??", dominant="DE")
    check(_cc == "DE" and _ll and abs(_ll[0] - 51.5) < 1.0,
          "geofix #3: with a dominant country the ambiguous name resolves to the RIGHT country")

    # review R1-C1: a lone ambiguous unknown-country property is left tbd+gap, NEVER a silent
    # global wrong-country pin (even if the network would return one)
    enrich._GAZETTEER = None; enrich._GAZETTEER_MULTI = None
    with tempfile.TemporaryDirectory() as td:
        enrich.CACHE_DIR = Path(td)
        _sg = enrich._geocode_one
        enrich._geocode_one = lambda *a, **k: ([50.7333, 4.2333], "BE")  # must NOT be used
        try:
            data = {"properties": [{"id": 1, "city": "Halle", "country": "??"}], "pois": [], "meta": {}}
            gaps: list = []
            enrich.geocode(data, gaps)
        finally:
            enrich._geocode_one = _sg
        check(not isinstance(data["properties"][0].get("lat"), (int, float))
              and any("ambiguous" in g.lower() for g in gaps),
              "geofix R1-C1: a lone ambiguous unknown-country city stays tbd+gap, not a global wrong pin")

    # review R4-I1 / bug #4: an already-located property's unknown country is filled from a seeded
    # CACHE even when the city is absent from the gazetteer (the documented offline seed workflow)
    with tempfile.TemporaryDirectory() as td:
        enrich.CACHE_DIR = Path(td)
        enrich._save_cache(enrich.GEOCODE_CACHE, {"zzznotacity|": {"latlng": [52.49, -0.70], "cc": "GB"}})
        _sr = enrich._reverse_cc
        enrich._reverse_cc = lambda *a, **k: (_ for _ in ()).throw(RuntimeError("offline"))
        try:
            data = {"properties": [{"id": 1, "city": "Zzznotacity", "country": "??",
                                    "lat": 52.49, "lng": -0.70}], "pois": [], "meta": {}}
            enrich.geocode(data, [])
        finally:
            enrich._reverse_cc = _sr
        check(data["properties"][0].get("country") == "GB",
              "geofix #4: a located property's country is filled from the seeded cache (gazetteer-absent city)")

    # map-link (#D): a q/geo/place link yields (coords, mapLink); a directions link resolves to
    # the DESTINATION (daddr), NEVER the origin (saddr) - review R1-I2/R2-I1
    for uri, want in [("https://maps.google.com/?q=52.2298,21.0118", 52.2298),
                      ("geo:48.8566,2.3522", 48.8566),
                      ("https://www.google.com/maps/place/40.4168,-3.7038", 40.4168)]:
        ll, mu = P._page_link_coords(_FakePage(uri))
        check(ll is not None and abs(ll[0] - want) < 1e-3 and mu == uri,
              f"geofix #D: a maps link with coords yields (coords, mapLink): {uri[:38]}")
    ll2, _ = P._page_link_coords(_FakePage(
        "https://maps.google.com/maps?saddr=48.1351,11.5820&daddr=52.5200,13.4050"))
    check(ll2 is not None and abs(ll2[0] - 52.52) < 1e-3,
          "geofix #D: a directions link resolves to the DESTINATION (daddr), never the origin (saddr)")

    enrich._DATASET = enrich._BORDERS = enrich._CITY_DATASET = None  # leave the real datasets active


def audit2_extract_cases() -> None:
    """2026-07 audit, Batch A - extraction data-honesty fixes (S1-1/9/10/11/43)."""
    print("audit2 extract (Batch A):")
    import normalize as _N2

    # S1-1: a list separator (comma-space) must not glue two number groups into one
    # fabricated area ('Unit 5, 25 000' had become 525000). EU/US grouping still reads.
    for raw, want, lbl in (
        ("Unit 5, 25 000 sq m", 25000.0, "'Unit 5, 25 000' -> 25000, not a glued 525000"),
        ("25 000 sq m", 25000.0, "EU space grouping '25 000' still reads 25000"),
        ("131,536 sq ft", 131536.0, "US comma thousands '131,536' still reads 131536"),
        ("12.500 m2", 12500.0, "EU dot thousands '12.500' still reads 12500"),
        ("39 471 sq m (expandable to 80 000)", 39471.0, "first plausible number still wins"),
    ):
        r, pr = {}, {}
        P._apply_num(r, pr, "warehouseArea", raw, "page 1")
        check(r.get("warehouseArea") == want, f"S1-1: {lbl} (got {r.get('warehouseArea')})")

    # S1-9: a rent RANGE carrying a currency symbol must not ship one end as a number
    d, n, _note, _u = P._parse_rent("€55 - 60 / sq m / year")
    check(d is not None and n is None, f"S1-9: rent range '€55 - 60' ships as text, no number (got {n})")
    d, n, _note, _u = P._parse_rent("€60 / sq m / year")
    check(n == 60.0, "S1-9: a non-range currency rent still parses (regression guard)")
    d, n, _note, _u = P._parse_rent("€4.50 / sq m (2024-2025 lease)")
    check(n == 4.5, f"S1-9: a DATE span in parens does not suppress a real single rent (got {n})")

    # S1-10: a unit-silent rent (no currency, no per-area) is flagged rentUnitAssumed
    d, n, _note, u = P._parse_rent("55")
    check(n == 55.0 and P._rent_unit_assumed("55", n) is True,
          "S1-10: a bare '55' rent ships the number AND is flagged unit-assumed")
    check(P._rent_unit_assumed("€55 / sq m", 55.0) is False,
          "S1-10: a currency+unit rent is NOT flagged assumed")

    # S1-11: a dot-thousands pair with no coordinate cue and no unit is NOT coordinates
    check(P._find_latlng("Superficies disponibles: 12.500, 8.750") is None,
          "S1-11: a cue-less dot-thousands pair is not accepted as lat/lng")
    check(P._find_latlng("GPS\n49.7384, 13.3736\n") == (49.7384, 13.3736),
          "S1-11: a genuine 4-decimal coordinate pair still parses")
    check(P._find_latlng("50.075, 14.437") == (50.075, 14.437),
          "S1-11: a bare 3-decimal coordinate pair (no size word) IS accepted (no false reject)")

    # S1-43: the spelled-out word 'euro'/'euros' is recognised as EUR
    check(_N2.currency_of("45 euros / m2") == "€", "S1-43: 'euros' recognised as EUR currency")
    check(bool(P.RENT_CONTEXT.search("45 euros")), "S1-43: 'euros' is a rent context")


def audit2b_matcher_cases() -> None:
    """2026-07 audit, Batch B - matcher (S2-8 one-park-missing over-merge on shared unknown dev)."""
    print("audit2 matcher (Batch B):")
    import match as M
    base = {"city": "Brno", "warehouseArea": 30000.0, "country": "CZ"}
    # S2-8: exactly one park missing + a SHARED UNKNOWN ('tbd') developer must NOT auto-merge
    a = {**base, "park": "Some Park", "developer": "tbd"}
    b = {**base, "developer": "tbd"}
    check(M._cross_source_auto(a, b) is False,
          "S2-8: one-park-missing + shared 'tbd' developer does NOT auto-merge (over-merge guard)")
    # but a shared KNOWN developer (same city, area within 5%) still auto-merges
    a2 = {**base, "park": "Some Park", "developer": "CTP"}
    b2 = {**base, "developer": "CTP"}
    check(M._cross_source_auto(a2, b2) is True,
          "S2-8: one-park-missing + shared KNOWN developer still auto-merges")


def audit2c_gate_cases() -> None:
    """2026-07 audit, Batch C - gate floors (S4-14 identity trace-coverage, S6-16 env console)."""
    print("audit2 gates (Batch C):")
    import csv as _csv
    import gate_runner as G
    import render_qa as RQ
    # S6-16: an offline map-tile / OSRM console error is [ENV], not a code defect
    check(RQ._is_env_error("Failed to load resource: tile.openstreetmap.org/9/1/2.png") is True,
          "S6-16: a map-tile fetch error is classified [ENV]")
    check(RQ._is_env_error("Uncaught TypeError: openModal is not a function") is False,
          "S6-16: a real JS error is NOT [ENV]")
    # S4-14: a fabricated identity field with no ledger row must BLOCK trace-coverage
    d = Path(tempfile.mkdtemp(prefix="cbre_gate_"))
    canon = d / "canonical.json"; led = d / "ledger.csv"
    canon.write_text(json.dumps({"meta": {}, "properties": [
        {"id": 1, "developer": "Fabricated Ltd", "city": "Nowhere", "park": "Ghost Park",
         "country": "??", "status": "Existing"}]}), encoding="utf-8")
    cols = ["property_id", "record_type", "field", "value", "source_file", "source_locator",
            "source_type", "extractor", "confidence", "conflict_note", "verified"]

    def _row(f, v, st="pdf"):
        return {"property_id": 1, "record_type": "property", "field": f, "value": v,
                "source_file": "b.pdf", "source_locator": "page 1", "source_type": st,
                "extractor": "pdf", "confidence": "Medium", "conflict_note": "", "verified": ""}

    def _write(rows):
        with open(led, "w", newline="", encoding="utf-8") as fh:
            w = _csv.DictWriter(fh, fieldnames=cols); w.writeheader()
            for r in rows:
                w.writerow(r)
    # city/park/status traced, country '??' gap-documented, but developer UNTRACED
    _write([_row("city", "Nowhere"), _row("park", "Ghost Park"), _row("status", "Existing"),
            {**_row("country", "??"), "source_file": "", "source_locator": "", "source_type": "gap"}])
    check(call(G, "trace-coverage", canon, "--ledger", led) != 0,
          "S4-14: a fabricated developer (no ledger row) BLOCKS trace-coverage")
    _write([_row("city", "Nowhere"), _row("park", "Ghost Park"), _row("status", "Existing"),
            _row("developer", "Fabricated Ltd"),
            {**_row("country", "??"), "source_file": "", "source_locator": "", "source_type": "gap"}])
    check(call(G, "trace-coverage", canon, "--ledger", led) == 0,
          "S4-14: once developer traces (and '??' country is a gap sentinel), trace-coverage PASSES")
    # S4-52: an empty / header-only ledger must BLOCK, not report ALL-PASS
    import ledger as _L
    empty_led = d / "empty.csv"
    empty_led.write_text(",".join(cols) + "\n", encoding="utf-8")
    check(call(_L, "validate", empty_led) != 0, "S4-52: an empty (header-only) ledger BLOCKS")


def audit2d_cases() -> None:
    """2026-07 audit, Batch D - build KPI/diagnostic (S5-15/S5-50) + email HTML body (S1-12)."""
    print("audit2 build/email (Batch D):")
    import build_dashboard as B
    import extract_email as E
    # S5-15: the 'tbd' unknown-region sentinel must not inflate kpi_regions
    k = B.compute_kpis([
        {"region": "Pilsen", "country": "CZ", "developer": "A"},
        {"region": "Brno", "country": "CZ", "developer": "B"},
        {"region": "tbd", "country": "CZ", "developer": "C"}],
        {}, {"area": "sq m", "rent": "€/sq m/yr"})
    check(k["kpi_regions"] == "2", f"S5-15: 'tbd' region excluded from kpi_regions (got {k['kpi_regions']})")
    # S5-50: a canonical with no 'properties' key raises a CLEAN error, not an opaque KeyError
    try:
        B.render({"meta": {}})
        check(False, "S5-50: render without 'properties' should raise")
    except ValueError:
        check(True, "S5-50: render without 'properties' raises a clean ValueError (not KeyError)")
    except KeyError:
        check(False, "S5-50: render still raises an opaque KeyError")
    # S1-12: an HTML-only email body is recovered (tags stripped, entities unescaped), not lost
    st = E._strip_html("<p>Rent is <b>&euro;55</b> per sq m</p>")
    check("Rent is" in st and "€55" in st and "<" not in st,
          "S1-12: _strip_html recovers HTML body prose (tags stripped, entities unescaped)")
    import intake as I
    # S0-41: a single real segment before noise ('City - FINAL') must not leak the noise
    reg, _c, _conf = I.infer_cluster("Pilsen - FINAL.pdf", {})
    check(reg == "Pilsen", f"S0-41: 'Pilsen - FINAL' -> region 'Pilsen' (noise stripped, got {reg!r})")
    # S0-13: the own-output ledger pattern is ANCHORED - a client file merely CONTAINING the
    # substring is not dropped, while our real deliverable names still are
    check(I._is_own_output("Client_Source_Ledger_notes.pdf") is False,
          "S0-13: a client file containing '_Source_Ledger' substring is NOT dropped")
    check(I._is_own_output("Normal_Source_Ledger.xlsx") is True,
          "S0-13: our real Source Ledger deliverable IS skipped")


def audit3_dupid_cases() -> None:
    print("audit3 gate (#32 dup-id single-pass Counter):")
    import io as _io
    from contextlib import redirect_stdout
    import gate_runner as G
    prop = lambda i: {"id": i, "country": "ES", "park": f"Park{i}", "developer": "CTP",
                      "city": "Madrid", "status": "Existing",
                      "photo": "data:image/png;base64,x", "warehouseArea": 40000,
                      "warehouseRent": "€60 / sq m / year", "warehouseRentVal": 60.0}
    # ids [1, 2, 1] - duplicate at positions 0 and 2; the report must preserve ORDER
    # and MULTIPLICITY (one entry per duplicated occurrence) -> '[1, 1]', not '[1]'
    data = {"meta": {}, "properties": [prop(1), prop(2), prop(1)], "pois": [], "regions": {}}
    with tempfile.TemporaryDirectory() as td:
        cp = Path(td) / "c.json"
        cp.write_text(json.dumps(data), encoding="utf-8")
        sv = sys.argv
        sys.argv = ["gate_runner", "validate-data", str(cp)]
        buf = _io.StringIO()
        try:
            with redirect_stdout(buf):
                try:
                    G.main()
                    rc = 0
                except SystemExit as e:
                    rc = e.code if isinstance(e.code, int) else (0 if e.code is None else 1)
        finally:
            sys.argv = sv
        out = buf.getvalue()
        check(rc != 0 and "duplicate property ids: [1, 1]" in out,
              "#32: validate-data reports duplicate ids as [1, 1] (order + multiplicity preserved by Counter)")


def reconcile_kpi_cases() -> None:
    print("audit3 gate (#24/#34 reconcile reuses compute_kpis, no 2nd render):")
    import build_dashboard as B
    import gate_runner as G
    import _common as C
    prop = lambda i, dev, cc: {"id": i, "country": cc, "park": f"Park{i}", "developer": dev,
                               "city": "Madrid", "status": "Existing",
                               "photo": "data:image/png;base64,x", "warehouseArea": 40000,
                               "warehouseRent": "€60 / sq m / year", "warehouseRentVal": 60.0}
    # 5 properties (KPI properties=5, unique among the strip: countries=1, regions=0, developers=2)
    data = {"meta": {"units": {"area": "sq m", "rent": "€/sq m/yr"}},
            "properties": [prop(1, "A", "ES"), prop(2, "A", "??"), prop(3, "B", "ES"),
                           prop(4, "B", "ES"), prop(5, "A", "ES")],
            "pois": [], "regions": {}}
    out_html, tokens = B.render(data)
    kpi_expected = tokens["kpi_properties"]
    props = [C.fill_render_sentinels(dict(p)) for p in data["properties"]]
    kpi_direct = B.compute_kpis(props, data.get("regions", {}),
                                (data.get("meta") or {}).get("units"))["kpi_properties"]
    check(kpi_direct == kpi_expected and kpi_direct == str(len(data["properties"])),
          "#24: compute_kpis kpi_properties == render() token == property count (the value reconcile now checks)")
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        cp = td / "c.json"
        cp.write_text(json.dumps(data), encoding="utf-8")
        built = td / "built.html"
        built.write_text(out_html, encoding="utf-8")
        check(call(G, "reconcile", built, "--canonical", cp) == 0,
              "#24: reconcile PASSES the correctly-rendered built file")
        target = f'<div class="kpi-value">{kpi_expected}</div>'
        tampered = td / "tampered.html"
        tampered.write_text(out_html.replace(target, f'<div class="kpi-value">{int(kpi_expected) + 1}</div>'),
                            encoding="utf-8")
        check(call(G, "reconcile", tampered, "--canonical", cp) != 0,
              "#24: reconcile still BLOCKS when the hero KPI properties count is wrong")


def selfcheck_i18n_placeholder_cases() -> None:
    print("audit3 gate (#33 self-check guards the EN KPI-sub placeholders):")
    import io as _io
    import types as _types
    from contextlib import redirect_stdout
    import gate_runner as G
    import i18n as I18N

    def _run():
        buf = _io.StringIO()
        with redirect_stdout(buf):
            rc = G.cmd_self_check(_types.SimpleNamespace())
        return rc, buf.getvalue()

    saved_a = I18N.EN.get("kpi_wh_area_sub_fmt")
    saved_u = I18N.EN.get("kpi_rent_sub_fmt")
    try:
        rc_ok, _ = _run()
        check(rc_ok == 0, "#33: self-check PASSES a healthy EN baseline")
        I18N.EN["kpi_wh_area_sub_fmt"] = "per building"  # {area} placeholder dropped
        rc_area, out_area = _run()
        check(rc_area != 0 and "kpi_wh_area_sub_fmt" in out_area,
              "#33: self-check BLOCKS when kpi_wh_area_sub_fmt loses its {area} placeholder")
        I18N.EN["kpi_wh_area_sub_fmt"] = saved_a
        I18N.EN["kpi_rent_sub_fmt"] = "per year"  # {unit} placeholder dropped
        rc_unit, out_unit = _run()
        check(rc_unit != 0 and "kpi_rent_sub_fmt" in out_unit,
              "#33: self-check BLOCKS when kpi_rent_sub_fmt loses its {unit} placeholder")
    finally:
        if saved_a is not None:
            I18N.EN["kpi_wh_area_sub_fmt"] = saved_a
        if saved_u is not None:
            I18N.EN["kpi_rent_sub_fmt"] = saved_u


def conflict_singleton_cases() -> None:
    print("audit3 merge (#44 conflict_candidates singleton short-circuit):")
    recA = {"city": "Corby", "developer": "Prologis", "park": "Apollo", "warehouseRentVal": 55.0,
            "__meta": {"source_type": "xlsx", "date": "2026-01-01", "prov": {}, "source_file": "a.xlsx"}}
    recB = {"city": "Corby", "developer": "Prologis", "park": "Apollo", "warehouseRentVal": 70.0,
            "__meta": {"source_type": "pdf", "date": "", "prov": {}, "source_file": "b.pdf"}}
    real = merge.conflict_candidates([[recA, recB]])
    check(len(real) == 1 and real[0]["field"] == "warehouseRentVal",
          "#44: a 2-record cluster with a genuine rent disagreement yields exactly 1 conflict")
    check(merge.conflict_candidates([[recA]]) == [],
          "#44: a singleton cluster yields NO conflicts (the short-circuit is byte-identical to the len<2 skip)")
    check(merge.conflict_candidates([]) == [],
          "#44: an empty cluster list yields NO conflicts")
    mixed = merge.conflict_candidates([[recA, recB], [recA]])
    check([c["conflict_id"] for c in mixed] == [c["conflict_id"] for c in real],
          "#44: interleaving a singleton does not change the real cluster's conflicts (same ids, same order)")


def geometry_prebatch_cases() -> None:
    print("audit3 image (#20 serial prewarm: one deck-wide geometry open, not per-page):")
    import io as _io
    import random as _r
    try:
        import merge as _M
        import fitz as _fz
        import pdfplumber as _pp
        from PIL import Image as _Img
    except Exception as e:
        check(False, f"#20: setup import failed ({e})")
        return

    def _noise(seed, w=340, h=240):
        _r.seed(seed)
        im = _Img.new("RGB", (w, h))
        im.putdata([(_r.randint(0, 255), _r.randint(0, 255), _r.randint(0, 255)) for _ in range(w * h)])
        b = _io.BytesIO()
        im.save(b, "JPEG", quality=80)
        return b.getvalue()

    def _clear():
        IMG._PLACED_CACHE.clear()
        IMG.close_doc_cache()

    with tempfile.TemporaryDirectory(ignore_cleanup_errors=True) as td:
        td = Path(td)
        doc = _fz.open()
        for p in range(4):
            pg = doc.new_page(width=600, height=800)
            pg.insert_image(_fz.Rect(60, 60, 460, 360), stream=_noise(p))
        f = td / "deck.pdf"
        doc.save(f)
        doc.close()
        recs = [{"city": "M", "park": f"P{p}",
                 "__meta": {"source_file": "deck.pdf", "source_type": "pdf", "page_no": p}}
                for p in range(4)]
        # reference geometry from a pristine per-page build (BEFORE installing the counter)
        ref = td / ".ref"
        ref.mkdir()
        _clear()
        lay_ref = IMG._placed_layout(f, ref)
        _clear()
        # count pdfplumber.open() during a SERIAL prewarm on a fresh geometry cache
        warm = td / ".warm"
        warm.mkdir()
        orig = _pp.open
        cnt = {"n": 0}

        def _counting(*a, **k):
            cnt["n"] += 1
            return orig(*a, **k)

        _pp.open = _counting
        # isolate GEOMETRY opens: the hero path's crop tier (_page_crops) keeps its OWN
        # per-page pdfplumber open for the page width, which #20's minimal fix deliberately
        # does NOT touch (it reaches into crop-scale math, a byte-identity risk). Neutralise
        # it here so the count reflects only the placedpage geometry the fix targets.
        _crops_orig = IMG._page_crops
        IMG._page_crops = lambda *a, **k: []
        try:
            _M.prewarm_images(recs, td, warm, IMG.DEFAULT_BUDGET_KB, seconds=60, workers=1)
        finally:
            _pp.open = orig
            IMG._page_crops = _crops_orig
        _clear()
        lay_warm = IMG._placed_layout(f, warm)
        _clear()
        check(lay_warm == lay_ref,
              "#20: serial-prewarmed deck geometry is byte-identical to the per-page reference")
        check(cnt["n"] <= 2,
              f"#20: serial prewarm opens the deck ONCE for geometry, not once per page (geometry opens={cnt['n']}, was 4)")


def match_memo_cases() -> None:
    print("audit3 match (#29+#30 memoise norm + token_set_ratio, byte-identical):")
    a = {"city": "Corby", "developer": "Prologis", "park": "Apollo Court",
         "warehouseArea": 50000, "__meta": {"source_file": "a.pdf", "source_type": "pdf"}}
    b = {"city": "Corby", "developer": "Prologis", "park": "Mercury House",
         "warehouseArea": 52000, "__meta": {"source_file": "b.xlsx", "source_type": "xlsx"}}
    key_a, key_b = match.match_key(a), match.match_key(b)
    check(key_a == "corby|prologis|apollo court" and key_b == "corby|prologis|mercury house",
          "#29: match_key is byte-identical after memoising norm")
    check(match.norm("Praha 5, Česko") == "praha 5 cesko" and match.norm(None) == ""
          and match.norm(123) == "123",
          "#29: cached norm returns identical strings for scalar inputs incl None/int")
    raw = match.fuzz.token_set_ratio(key_a, key_b)  # the un-cached scorer, still exposed
    check(match._tsr(key_a, key_b) == raw,
          "#30: cached _tsr equals the raw scorer for the same key pair (backend-agnostic)")
    check(raw < match.MATCH_THRESHOLD,
          "#30: the pinned pair scores below the match threshold (still 'grey')")
    check(match.pair_class(a, b) == "grey" and match.dedupe([a, b]) == match.dedupe([a, b], None),
          "#29+#30: pair_class/dedupe verdicts unchanged (grey; offline split == default-arg fallback)")


def image_reuse_memo_cases() -> None:
    print("audit3 image (#21+#38+#39 in-process decode / slide-list reuse):")
    import io as _io
    try:
        import fitz as _fz
        from PIL import Image as _Img
    except Exception as e:
        check(False, f"#21/#38/#39: setup import failed ({e})")
        return
    png = _io.BytesIO()
    _Img.new("RGB", (400, 300), (30, 120, 200)).save(png, "PNG")
    png = png.getvalue()
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        doc0 = _fz.open()
        pg = doc0.new_page(width=600, height=800)
        pg.insert_image(_fz.Rect(60, 60, 460, 360), stream=png)
        pdf = td / "deck.pdf"
        doc0.save(pdf)
        doc0.close()

        IMG.close_doc_cache()
        doc = IMG._get_doc(pdf)
        a = IMG.page_embedded_images(doc, 0)
        b = IMG.page_embedded_images(doc, 0)
        check(len(a) >= 1 and [(d["w"], d["h"]) for d in a] == [(d["w"], d["h"]) for d in b],
              "#39: page_embedded_images returns an equivalent list on repeat")
        check(a is b, "#39: page_embedded_images REUSES the decoded list (same object; decoded once per run)")
        c1 = IMG.candidates_for_page(pdf, 0)
        c2 = IMG.candidates_for_page(pdf, 0)
        check(len(c1) >= 1 and [(d["w"], d["h"]) for d in c1] == [(d["w"], d["h"]) for d in c2],
              "#21: candidates_for_page equivalent on repeat")
        check(c1[0]["img"] is c2[0]["img"],
              "#21: candidates_for_page reuses the decoded PIL object (no re-decode across calls)")
        # close_doc_cache drops the memo WITHOUT corrupting results (id(doc) recycle guard)
        IMG.close_doc_cache()
        doc2 = IMG._get_doc(pdf)
        a2 = IMG.page_embedded_images(doc2, 0)
        check([(d["w"], d["h"]) for d in a2] == [(d["w"], d["h"]) for d in a],
              "#39: after close_doc_cache the memo is dropped but results stay equivalent")
        IMG.close_doc_cache()

        # #38 slide-list + slide_pictures reuse (guarded: needs python-pptx)
        try:
            from pptx import Presentation as _Prs
            from pptx.util import Inches as _In
            prs = _Prs()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(_io.BytesIO(png), _In(1), _In(1), _In(3), _In(2))
            pptx = td / "deck.pptx"
            prs.save(str(pptx))
        except Exception:
            pptx = None
        if pptx is not None:
            IMG.close_doc_cache()
            p1 = IMG.slide_pictures(pptx, 0)
            p2 = IMG.slide_pictures(pptx, 0)
            check(len(p1) >= 1 and [(d["w"], d["h"]) for d in p1] == [(d["w"], d["h"]) for d in p2],
                  "#39: slide_pictures equivalent on repeat")
            check(p1 is p2, "#39: slide_pictures REUSES the decoded pictures list (decoded once per run)")
            _ps = getattr(IMG, "_pptx_slides", None)
            if _ps is not None:
                check(_ps(pptx) is _ps(pptx),
                      "#38: _pptx_slides reuses the enumerated slide list (no re-enumeration per call)")
            IMG.close_doc_cache()


def counts_once_cases() -> None:
    print("audit3 build (#55 KPI counts derived once, byte-identical):")
    import build_dashboard as B
    props = [
        {"country": "GB", "developer": "CTP", "region": "Midlands", "regionCode": "MID"},
        {"country": "GB", "developer": "CTP", "region": "Midlands", "regionCode": "MID"},
        {"country": "IE", "developer": "Panattoni", "region": "Dublin", "regionCode": "DUB"},
        {"country": "IE", "developer": "CTP", "region": "Dublin", "regionCode": "DUB"},
    ]
    k = B.compute_kpis(props, {}, {"area": "sq m", "rent": "€/sq m/yr"}, {})
    check(k["kpi_countries"] == "2", "#55: kpi_countries de-dups GB/GB,IE/IE -> 2")
    check(k["kpi_developers"] == "2", "#55: kpi_developers de-dups CTP x3 + Panattoni -> 2 (no redundant re-dedup)")
    check(k["kpi_countries_sub"] == "GB · IE",
          "#55: kpi_countries_sub is the sorted distinct enumeration (byte-exact, single country_set)")


def intake_memory_cases() -> None:
    print("audit3 intake (#36 dedup-hash size-prefilter + MemoryError guard):")
    import intake
    import unittest.mock as _mock
    # (a) size-prefilter: shrink the cap so a normal small file trips the oversize branch
    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        (root / "small.pdf").write_bytes(b"a-real-small-brochure")   # 21 bytes > 8
        (root / "tiny.pdf").write_bytes(b"x")                         # 1 byte <= 8
        orig_cap = getattr(intake, "_DEDUP_MAX_BYTES", None)
        intake._DEDUP_MAX_BYTES = 8
        try:
            inv = intake.discover(root)
        finally:
            if orig_cap is not None:
                intake._DEDUP_MAX_BYTES = orig_cap
        oversize = inv.get("skipped_hash_oversize", [])
        disc = any("small.pdf" in (c.get("pdfs") or []) for c in inv.get("clusters", {}).values())
        check("skipped_hash_oversize" in inv and "small.pdf" in oversize and disc,
              "S36a: an oversize file skips the dedup hash but is STILL discovered as a brochure")
        check("tiny.pdf" not in oversize,
              "S36a: a file under the cap is hashed normally (not flagged oversize)")
    # (b) MemoryError backstop: force the read/hash to raise MemoryError, assert no crash
    with tempfile.TemporaryDirectory() as td2:
        root2 = Path(td2)
        (root2 / "boom.pdf").write_bytes(b"anything")
        with _mock.patch.object(intake.hashlib, "sha256", side_effect=MemoryError):
            try:
                inv2 = intake.discover(root2)
                ok_b = isinstance(inv2, dict) and any(
                    "boom.pdf" in (c.get("pdfs") or []) for c in inv2.get("clusters", {}).values())
            except MemoryError:
                ok_b = False
        check(ok_b, "S36b: a MemoryError during the dedup hash is CAUGHT - file kept, run not crashed")


def pptx_slide_fail_cases() -> None:
    print("audit3 pptx (#19 one bad slide must not mis-route the whole deck to raster):")
    try:
        from pptx import Presentation as _Prs
        from pptx.util import Inches as _In
        import interpret_prep as IP
        import extract_pptx as PPTX
    except Exception as e:
        check(False, f"#19: setup import failed ({e})")
        return
    filler = ("City Pilsen Developer CTP Warehouse Area 40000 sq m clear height 12 m "
              "slide {n} plus filler text to comfortably exceed the eighty character threshold")
    with tempfile.TemporaryDirectory() as td:
        prs = _Prs()
        for n in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            tb = slide.shapes.add_textbox(_In(0.5), _In(0.5), _In(9), _In(2))
            tb.text_frame.text = filler.format(n=n + 1)
        deck = Path(td) / "deck.pptx"
        prs.save(str(deck))

        orig = PPTX.slide_text
        calls = {"n": 0}

        def _flaky(slide):
            calls["n"] += 1
            if calls["n"] == 2:
                raise RuntimeError("bad slide XML")
            return orig(slide)

        PPTX.slide_text = _flaky
        try:
            texts = IP._pptx_slide_texts(deck)
            mode = IP._decide_mode(texts)
        finally:
            PPTX.slide_text = orig
        check(len(texts) == 3 and mode == "text",
              "#19: one raising slide yields '' for that slide only; the deck still routes text (2 clean slides survive)")
        check(len(texts) == 3 and texts[1] == "" and texts[0] != "" and texts[2] != "",
              "#19: only the raising slide is blanked; its neighbours keep their text")
        st_fn = getattr(PPTX, "slide_texts", None)
        if st_fn is not None:
            calls["n"] = 0
            PPTX.slide_text = _flaky
            try:
                direct = st_fn(_Prs(str(deck)))
            finally:
                PPTX.slide_text = orig
            check(len(direct) == 3 and direct[1] == "" and direct[0] and direct[2],
                  "#19: extract_pptx.slide_texts blanks ONLY the bad slide (['t','','t'])")


def intake_resume_cases() -> None:
    print("audit3 resume (#27 intake stamp is recursive over nested/in-place input edits):")
    import run as RUN
    import os as _os
    saved_resume = RUN.RESUME
    RUN.RESUME = True
    try:
        with tempfile.TemporaryDirectory() as td:
            root = Path(td)
            inputs = root / "inputs"
            (inputs / "sub").mkdir(parents=True)
            top = inputs / "top.pdf"
            top.write_bytes(b"top-file")
            nested = inputs / "sub" / "nested.pdf"
            nested.write_bytes(b"nested-file")
            out = root / "inventory.json"
            out.write_text("{}", encoding="utf-8")
            t0 = 1_000_000
            for p in (top, nested, inputs / "sub", inputs):
                _os.utime(p, (t0, t0))
            _os.utime(out, (t0 + 10, t0 + 10))   # out newer than every input -> current
            check(RUN._is_current(out, [inputs]) is True,
                  "S#27a: no input change -> intake stamp reports current (stage skipped/resumed)")
            # (b) an in-place edit of a NESTED input (dir-node mtimes do NOT bubble up on NTFS)
            _os.utime(nested, (t0 + 20, t0 + 20))
            check(RUN._is_current(out, [inputs]) is False,
                  "S#27b: an in-place edit of a NESTED input invalidates the intake stamp (re-runs)")
            # (c) reset the nested edit, then an in-place edit of a TOP-LEVEL input
            _os.utime(nested, (t0, t0))
            _os.utime(top, (t0 + 20, t0 + 20))
            check(RUN._is_current(out, [inputs]) is False,
                  "S#27c: an in-place edit of a TOP-LEVEL input invalidates the intake stamp (re-runs)")
    finally:
        RUN.RESUME = saved_resume


def deliver_resume_cases() -> None:
    print("audit3 resume (#25 Stage-7 deliver has a currency guard: skip current, redo on change):")
    try:
        import fitz as _fz
        import run as RUN
    except Exception as e:
        check(False, f"#25: setup import failed ({e})")
        return
    import io as _io
    import os as _os
    from contextlib import redirect_stdout, redirect_stderr

    def _mk_text_pdf(folder: Path) -> Path:
        # a clean born-digital flyer whose LABELS are NOT in extract_pdf's dictionary
        # (the TEDi case) -> the only honest path is interpretation (routes to exit 3)
        doc = _fz.open()
        for opt, city, area, rent in (("Option 1", "Valencia", "12,500 m2", "4.20 / sqm / month"),
                                      ("Option 2", "Sagunto", "8,750 m2", "3.90 / sqm / month")):
            pg = doc.new_page()
            y = 60
            for ln in (f"VALENCIA REGION - {opt}", f"City {city}", "Owner/developer Goodman",
                       f"Total existing space {area}", f"Warehouse - Asking rent {rent}",
                       "This prime logistics warehouse is strategically located near the A-7 "
                       "motorway with excellent connectivity to the Port of Valencia."):
                pg.insert_text((40, y), ln, fontsize=11)
                y += 22
        f = folder / "Options - Valencia.pdf"
        doc.save(f)
        doc.close()
        return f

    def _spine_resume(folder, work):
        saved = sys.argv
        sys.argv = ["run.py", "--folder", str(folder), "--work", str(work),
                    "--client", "TEDi", "--quiet"]   # resume ON (NO --no-resume)
        rc = 0
        try:
            with redirect_stdout(_io.StringIO()), redirect_stderr(_io.StringIO()):
                RUN.main()
        except SystemExit as e:
            rc = e.code if isinstance(e.code, int) else (0 if e.code is None else 1)
        except Exception as e:
            rc = f"CRASH: {type(e).__name__}: {e}"
        finally:
            sys.argv = saved
        return rc

    def _rec(park, city, page_no, rent_val, src_name):
        fields = {"park": park, "developer": "Goodman", "city": city, "country": "ES",
                  "region": "Valencia", "status": "Existing", "warehouseArea": 12500,
                  "warehouseRent": f"€{rent_val:g} / sq m / year", "warehouseRentVal": rent_val}
        prov = {k: f"page {page_no + 1} (text interpretation)" for k in fields}
        return {**fields, "__meta": {"source_file": src_name, "source_type": "pdf",
                                     "locator_base": f"page {page_no + 1}",
                                     "page_no": page_no, "prov": prov}}

    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        folder = td / "inputs"
        folder.mkdir()
        work = td / "work"
        work.mkdir()
        src = _mk_text_pdf(folder)
        rc1 = _spine_resume(folder, work)
        if rc1 != 3:
            check(False, f"#25: first pass should route to interpretation (exit 3, got {rc1!r})")
            return
        (work / "extract" / "Valencia_vision.json").write_text(
            json.dumps([_rec("Goodman Valencia", "Valencia", 0, 50.4, src.name),
                        _rec("Goodman Sagunto", "Sagunto", 1, 46.8, src.name)]), encoding="utf-8")
        rc2 = _spine_resume(folder, work)
        deliverables = work / "deliverables"
        dash = sorted(deliverables.glob("*.html")) if deliverables.exists() else []
        if rc2 != 0 or not dash:
            check(False, f"#25: second pass should deliver a dashboard (exit {rc2!r}, dash={bool(dash)})")
            return
        dash = dash[0]
        m_after = dash.stat().st_mtime_ns
        # (a) a no-change resume SKIPS deliver -> dashboard mtime unchanged
        rc3 = _spine_resume(folder, work)
        check(rc3 == 0 and dash.stat().st_mtime_ns == m_after,
              "#25a: a no-change resume SKIPS deliver (dashboard mtime unchanged)")
        # (b) canonical changes -> build + deliver RE-FIRE (dashboard re-delivered, mtime advances)
        canonical = work / "canonical.json"
        fut = canonical.stat().st_mtime + 100
        _os.utime(canonical, (fut, fut))
        rc4 = _spine_resume(folder, work)
        check(rc4 == 0 and dash.stat().st_mtime_ns > m_after,
              "#25b: after canonical changes, deliver RE-FIRES (never a stale skip; mtime advances)")


def interpret_resume_cases() -> None:
    print("audit3 resume (#28+#37 interpret_prep text-mode refreshes region/country on a re-cluster):")
    try:
        import fitz as _fz
        import interpret_prep as IP
    except Exception as e:
        check(False, f"#28+#37: setup import failed ({e})")
        return
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        doc = _fz.open()
        pg = doc.new_page()
        pg.insert_text((40, 60), "City Prague. Warehouse Area 10000 sq ft. Clear Height 12m. "
                                 "Rent 5 EUR per sq m per year. Developer Acme. Park Logistics One.",
                       fontsize=11)
        pdf = td / "deck.pdf"
        doc.save(pdf)
        doc.close()
        out = td / "vision"
        e1 = IP.prepare(pdf, "RegionA", "CZ", out, resume=True)
        check(e1.get("mode") == "text", "#28+#37: a clean text deck routes to mode 'text' (baseline)")
        rp = e1["pages"][0].get("render") if e1.get("pages") else None
        m0 = Path(rp).stat().st_mtime_ns if rp and Path(rp).exists() else None
        # no-change resume: reuse the cached entry + thumbnails
        e2 = IP.prepare(pdf, "RegionA", "CZ", out, resume=True)
        reuse_ok = (m0 is None) or (rp and Path(rp).exists() and Path(rp).stat().st_mtime_ns == m0)
        check(e2.get("region") == "RegionA" and reuse_ok,
              "#28+#37: a no-change resume reuses the cached text entry + thumbnails (skipped/resumed)")
        # re-cluster: SAME bytes, corrected region/country -> the entry must re-reflect them
        e3 = IP.prepare(pdf, "RegionB", "PL", out, resume=True)
        check(e3.get("region") == "RegionB" and e3.get("country") == "PL",
              "#28+#37: a re-clustered deck (same bytes, corrected region) re-reflects the NEW "
              "region/country, not the stale cached one")
        text_reused = (e3.get("pages") and e1.get("pages")
                       and e3["pages"][0].get("text") == e1["pages"][0].get("text"))
        render_reused = (m0 is None) or (rp and Path(rp).exists() and Path(rp).stat().st_mtime_ns == m0)
        check(bool(text_reused) and render_reused,
              "#28+#37: the re-cluster refreshes ONLY the manifest metadata; the byte-derived "
              "text + thumbnails stay cached")


def load_canonical_cache_cases() -> None:
    print("audit3 perf (#22+#35 load_canonical mtime-keyed parse-cache, deepcopy-isolated):")
    import _common as C
    with tempfile.TemporaryDirectory() as td:
        cp = Path(td) / "canonical.json"
        A = {"meta": {}, "properties": [{"id": 1, "city": "Prague"}], "pois": [], "regions": {}}
        C.atomic_write_text(cp, json.dumps(A))
        if hasattr(C, "_CANON_CACHE"):
            C._CANON_CACHE.clear()
        # (1) CACHE HIT (RED->GREEN): an unchanged canonical is parsed ONCE across two loads
        orig = C.load_json
        cnt = {"n": 0}

        def _counting(p):
            cnt["n"] += 1
            return orig(p)

        C.load_json = _counting
        try:
            d1 = C.load_canonical(cp)
            d2 = C.load_canonical(cp)
        finally:
            C.load_json = orig
        check(d1 == A and d2 == A, "#22: load_canonical returns the correct content")
        check(cnt["n"] == 1,
              f"#22: an unchanged canonical is parsed ONCE across two loads (in-process cache hit; parses={cnt['n']})")
        # (2) MUTATION ISOLATION (#35): mutating a returned dict must not poison the cache/another load
        d1["properties"][0]["city"] = "MUTATED"
        d3 = C.load_canonical(cp)
        check(d3["properties"][0]["city"] == "Prague",
              "#35: a caller mutating the returned dict does NOT affect a later load (deepcopy isolation)")
        # (3) STALE-READ / INVALIDATION (#35): an atomic rewrite (new mtime/size) must be seen
        B = {"meta": {}, "properties": [{"id": 1, "city": "Brno"}, {"id": 2, "city": "Ostrava"}],
             "pois": [], "regions": {}}
        C.atomic_write_text(cp, json.dumps(B))
        d4 = C.load_canonical(cp)
        check(d4 == B and len(d4["properties"]) == 2,
              "#35: a rewritten canonical (new mtime/size) invalidates the cache (no stale read)")
        # (4) SAME-size rewrite forced to the SAME mtime_ns: st_ino still invalidates (closes the
        #     coarse-Windows-mtime collision hole; the skill only ever writes via atomic os.replace)
        import os as _os
        C._CANON_CACHE.clear()
        S1 = json.dumps({"meta": {}, "properties": [{"id": 1, "city": "AAAA"}], "pois": [], "regions": {}})
        C.atomic_write_text(cp, S1)
        m = cp.stat().st_mtime_ns
        _ = C.load_canonical(cp)
        S2 = json.dumps({"meta": {}, "properties": [{"id": 1, "city": "BBBB"}], "pois": [], "regions": {}})
        check(len(S2) == len(S1), "#35: (precondition) the two canonicals are the same byte size")
        C.atomic_write_text(cp, S2)
        _os.utime(cp, ns=(m, m))  # force the SAME mtime_ns as the cached entry
        e2 = C.load_canonical(cp)
        check(e2["properties"][0]["city"] == "BBBB",
              "#35: a SAME-size rewrite forced to the SAME mtime_ns still invalidates (st_ino closes the coarse-mtime hole)")


def main() -> int:
    pdf_cases()
    audit2_extract_cases()
    audit2b_matcher_cases()
    audit2c_gate_cases()
    audit2d_cases()
    interpret_cases()
    xlsx_cases()
    tracker_map_cases()
    area_band_cases()
    email_cases()
    match_cases()
    verifier_cases()
    cache_cases()
    prewarm_resume_cases()
    enrich_cases()
    geofix_cases()
    intake_cases()
    web_enrich_cases()
    vision_validate_cases()
    region_code_cases()
    placeholder_audit_cases()
    regions_dataset_cases()
    region_binding_cases()
    region_label_llm_cases()
    recency_cases()
    tedi_cases()
    pptx_slide_cases()
    vision_dedup_cases()
    precedence_cases()
    units_cases()
    description_pick_cases()
    batch_a_cases()
    batch_b_cases()
    batch_c_cases()
    batch_d_cases()
    batch_f_cases()
    batch_g_cases()
    batch_h_cases()
    batch_i_cases()
    batch_j_cases()
    feedback1_hero_cases()
    feedback2_region_cases()
    feedback3_vision_dpi_cases()
    llm_hero_cases()
    image_pages_carousel_cases()
    plan_page_cases()
    audit3_dupid_cases()
    reconcile_kpi_cases()
    selfcheck_i18n_placeholder_cases()
    conflict_singleton_cases()
    geometry_prebatch_cases()
    match_memo_cases()
    image_reuse_memo_cases()
    counts_once_cases()
    intake_memory_cases()
    pptx_slide_fail_cases()
    intake_resume_cases()
    deliver_resume_cases()
    interpret_resume_cases()
    load_canonical_cache_cases()
    if FAILS:
        print(f"\nEXTRACT TEST: FAIL ({len(FAILS)})")
        for f in FAILS:
            print(f"  - {f}")
        return 1
    print("\nEXTRACT TEST: PASS (all audit misparses fixed; guards verified)")
    return 0


if __name__ == "__main__":
    # the suite prints non-ASCII test labels (e.g. the multilingual landlord/owner
    # headers: Eigentumer, PropriEtaire, Wlasciciel with the Polish l-stroke). A
    # Windows host console defaults to cp1252, which cannot encode those - a bare
    # print() would raise UnicodeEncodeError and crash the suite (a DISPLAY bug, not a
    # test failure). Force UTF-8 (errors=replace as a belt-and-braces) so the battery
    # runs identically under cp1252 (mcp__shell) and a UTF-8 shell.
    for _s in (sys.stdout, sys.stderr):
        try:
            _s.reconfigure(encoding="utf-8", errors="replace")
        except Exception:
            pass
    sys.exit(main())
