#!/usr/bin/env python3
"""cowork_sim.py - a COWORK-FIDELITY simulator for the run.py state machine.

WHY THIS EXISTS. The rest of the eval battery runs in ONE process, to completion, with
warm in-memory state and cooperative inputs. Cowork does the opposite, and every bug we
have shipped to a live run lived in that gap:

  * FRESH PROCESS per invocation  -> anything not on disk is lost
  * a HARD ~40-45s shell cap that KILLS run.py mid-work -> it must checkpoint + resume,
    and each resumed run must be STRICTLY CHEAPER (the 2026-07-30 placeholder-audit bug:
    an uncached whole-deck re-parse every re-run -> the merge never finished)
  * exit-code ROUND-TRIPS answered by an LLM sub-agent that may DECLINE, answer PARTIALLY,
    echo keys differently, or write malformed JSON (the 2026-07-30 region-label bug: a
    declined label was re-asked forever because the "already asked" guard dropped declines)

So this simulator drives the REAL run.py as a SUBPROCESS, under a real timeout, and plays
the orchestrator - answering each exit code the way Claude would, optionally adversarially.
It then asserts the two properties the battery cannot see:

  CONVERGENCE  - the loop reaches exit 0 in a bounded number of rounds, and never repeats
                 a (exit_code, work-dir state) fingerprint (a repeat IS a livelock).
  PROGRESS     - a resumed round is strictly cheaper / adds state. A round that costs the
                 same and adds nothing is work being redone (the bug-1 signature).

It is a DIAGNOSTIC, not a pass/fail gate for the battery: it needs several minutes and
native PyMuPDF. Run it deliberately, before shipping a bundle to Cowork.

Usage:
  python evals/cowork_sim.py                          # base scenario, 40s cap
  python evals/cowork_sim.py --cap 12                 # brutal cap: forces many resumes
  python evals/cowork_sim.py --answers decline        # sub-agent declines everything legal
  python evals/cowork_sim.py --answers partial        # answers only the first job each round
  python evals/cowork_sim.py --answers malformed      # writes truncated JSON once, then clean
  python evals/cowork_sim.py --decks 6 --props 4      # bigger corpus
  python evals/cowork_sim.py --keep                   # keep the temp corpus for inspection

Exit 0 = converged with no defect flagged. Exit 1 = a defect was flagged (read the verdict).
"""
from __future__ import annotations

import argparse
import hashlib
import json
import shutil
import subprocess
import sys
import tempfile
import time
from pathlib import Path

SKILL = Path(__file__).resolve().parent.parent
RUN_PY = SKILL / "helpers" / "run.py"

# ---------------------------------------------------------------- corpus

CITIES = [
    ("Corby", "Northamptonshire", "GB", "Prologis"),
    ("Venlo", "Limburg", "NL", "Montea"),
    ("Zaragoza", "Aragon", "ES", "Merlin Properties"),
    ("Pilsen", "Plzensky kraj", "CZ", "CTP"),
    ("Duisburg", "Nordrhein-Westfalen", "DE", "Garbe"),
    ("Lodz", "Lodzkie", "PL", "Panattoni"),
]


def _photo_bytes(seed: int, w: int = 900, h: int = 600) -> bytes:
    """A noisy colour gradient - reads as a photograph to the classifier (a flat fill
    would be classified 'plan'/'text' and never bind as a hero)."""
    import io
    import random

    from PIL import Image

    rnd = random.Random(seed)
    img = Image.new("RGB", (w, h))
    px = img.load()
    for y in range(h):
        base = (y * 255) // max(1, h)
        for x in range(w):
            px[x, y] = (min(255, base // 2 + rnd.randrange(90)),
                        min(255, (x * 200) // w + rnd.randrange(70)),
                        min(255, 60 + rnd.randrange(150)))
    b = io.BytesIO()
    img.save(b, format="JPEG", quality=82)
    return b.getvalue()


def build_images(folder: Path, n: int) -> list[str]:
    """Standalone photo files - what a broker drops in beside (or INSTEAD of) a deck."""
    folder.mkdir(parents=True, exist_ok=True)
    names = []
    for i in range(n):
        city = CITIES[i % len(CITIES)][0]
        name = f"{city}-warehouse-{i + 1}.jpg"
        (folder / name).write_bytes(_photo_bytes(500 + i))
        names.append(name)
    return names


def build_xlsx(folder: Path, n_rows: int) -> list[str]:
    """An availability tracker - the classic no-brochure input."""
    from openpyxl import Workbook

    folder.mkdir(parents=True, exist_ok=True)
    wb = Workbook()
    ws = wb.active
    ws.title = "Availability"
    ws.append(["Park", "City", "Country", "Developer", "Warehouse area (sq m)",
               "Rent (EUR/sq m/year)", "Clear height (m)", "Availability"])
    for i in range(n_rows):
        city, region, cc, dev = CITIES[i % len(CITIES)]
        ws.append([f"{dev} Park {city}", city, cc, dev, 15000 + i * 4200,
                   round(52 + i * 3.5, 2), round(10.5 + (i % 4) * 0.5, 1), f"Q{1 + i % 4} 2027"])
    name = "availability-tracker.xlsx"
    wb.save(str(folder / name))
    return [name]


def build_json(folder: Path, n_rows: int) -> list[str]:
    """A JSON property dataset. THE 2026-07-30 LIVE BREAK: a broker supplied a JSON plus
    photos and the run failed because the photos were not inside a brochure. intake.py
    classifies only .pdf/.pptx/.xlsx/.xlsm/.csv/images/.msg/.eml - a .json input matches
    NO branch, so it is dropped without even being recorded as unreadable."""
    folder.mkdir(parents=True, exist_ok=True)
    rows = []
    for i in range(n_rows):
        city, region, cc, dev = CITIES[i % len(CITIES)]
        rows.append({"park": f"{dev} Park {city}", "city": city, "country": cc,
                     "region": region, "developer": dev,
                     "warehouseArea": 15000 + i * 4200,
                     "warehouseRent": f"€{52 + i * 3.5:.2f} / sq m / year",
                     "clearHeight": f"{10.5 + (i % 4) * 0.5} m",
                     "photo": f"{city}-warehouse-{i + 1}.jpg"})
    name = "properties.json"
    (folder / name).write_text(json.dumps({"properties": rows}, indent=2), encoding="utf-8")
    return [name]


def build_txt(folder: Path, n_rows: int) -> list[str]:
    """A pasted broker email saved as .txt - also unclassified by intake."""
    folder.mkdir(parents=True, exist_ok=True)
    out = ["Hi - here are the options we discussed:", ""]
    for i in range(n_rows):
        city, region, cc, dev = CITIES[i % len(CITIES)]
        out += [f"{dev} Park {city} ({city}, {region})",
                f"  Warehouse area: {15000 + i * 4200:,} sq m",
                f"  Rent: EUR {52 + i * 3.5:.2f} per sq m per year",
                f"  Clear height: {10.5 + (i % 4) * 0.5} m",
                f"  Developer: {dev}", ""]
    name = "broker-email.txt"
    (folder / name).write_text("\n".join(out), encoding="utf-8")
    return [name]


def build_pptx(folder: Path, n_decks: int, n_props: int) -> list[str]:
    """A PPTX brochure - the 'preferred IMAGE source' path."""
    import io

    from pptx import Presentation
    from pptx.util import Inches, Pt

    folder.mkdir(parents=True, exist_ok=True)
    names = []
    for d in range(n_decks):
        city, region, cc, dev = CITIES[d % len(CITIES)]
        prs = Presentation()
        for p in range(n_props):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"{dev} Park {city} - Unit {chr(65 + p)}{d + 1}"
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.5), Inches(3.5))
            tf = tb.text_frame
            for ln in [f"{city}, {region}",
                       f"Warehouse area: {14000 + p * 5200:,} sq m",
                       f"Rent: EUR {4.80 + 0.3 * p:.2f} per sq m per month",
                       f"Clear height: {10.5 + 0.5 * p:.1f} m",
                       f"Developer: {dev}"]:
                para = tf.add_paragraph()
                para.text = ln
                para.font.size = Pt(14)
            slide.shapes.add_picture(io.BytesIO(_photo_bytes(d * 40 + p + 3)),
                                     Inches(5.2), Inches(1.6), width=Inches(4.2))
        name = f"{dev.split()[0]}-{city}-deck.pptx".replace(" ", "-")
        prs.save(str(folder / name))
        names.append(name)
    return names


def build_shape(folder: Path, shape: str, n_decks: int, n_props: int) -> list[str]:
    """Assemble an INPUT SHAPE. Real brokers hand over all of these; the pipeline is
    written around the brochure case, so every non-brochure shape is a risk surface."""
    n_rows = max(2, n_decks * n_props)
    if shape == "pdf":
        return build_corpus(folder, n_decks, n_props)
    if shape == "pptx":
        return build_pptx(folder, n_decks, n_props)
    if shape == "xlsx":
        return build_xlsx(folder, n_rows)
    if shape == "xlsx+images":
        return build_xlsx(folder, n_rows) + build_images(folder, n_rows)
    if shape == "images":
        return build_images(folder, n_rows)
    if shape == "json+images":
        return build_json(folder, n_rows) + build_images(folder, n_rows)
    if shape == "json":
        return build_json(folder, n_rows)
    if shape == "txt+images":
        return build_txt(folder, n_rows) + build_images(folder, n_rows)
    if shape == "mixed":
        return (build_corpus(folder, max(1, n_decks - 1), n_props)
                + build_xlsx(folder, n_rows) + build_images(folder, 2))
    if shape == "empty":
        folder.mkdir(parents=True, exist_ok=True)
        return []
    raise SystemExit(f"unknown --inputs shape: {shape}")


def build_corpus(folder: Path, n_decks: int, n_props: int) -> list[str]:
    """Synthesize born-digital brochure PDFs: per property a spec page carrying CORE
    fields (city / developer / area / rent - see extract_pdf._has_core) plus an embedded
    photo, and one vector site-plan page per deck. Returns the file names."""
    import fitz

    folder.mkdir(parents=True, exist_ok=True)
    names = []
    for d in range(n_decks):
        doc = fitz.open()
        city, region, cc, dev = CITIES[d % len(CITIES)]
        # cover
        cover = doc.new_page(width=595, height=842)
        cover.insert_text((60, 90), f"{dev} Park {city}", fontsize=26)
        cover.insert_text((60, 130), f"{city}, {region}", fontsize=14)
        cover.insert_image(fitz.Rect(60, 170, 535, 480), stream=_photo_bytes(d * 17 + 1))
        for p in range(n_props):
            pg = doc.new_page(width=595, height=842)
            unit = f"Unit {chr(65 + p)}{d + 1}"
            area = 12000 + p * 6500 + d * 1500
            rent = 4.75 + 0.35 * p + 0.2 * d
            lines = [
                f"{dev} Park {city} - {unit}",
                f"{city}, {region}",
                "United Kingdom" if cc == "GB" else region,
                "",
                f"Warehouse area: {area:,} sq m",
                f"Office area: {900 + p * 120:,} sq m",
                f"Rent: EUR {rent:.2f} per sq m per month",
                f"Clear height: {10.5 + 0.5 * p:.1f} m",
                f"Dock doors: {12 + 2 * p}",
                f"Car parking: {80 + 10 * p} spaces",
                f"Developer: {dev}",
                f"Available: Q{1 + (p % 4)} 202{6 + (p % 3)}",
                f"Motorway: A{14 + d} ({2 + p} km)",
                "BREEAM: Very Good",
            ]
            y = 70
            for ln in lines:
                pg.insert_text((60, y), ln, fontsize=11)
                y += 22
            pg.insert_image(fitz.Rect(300, 430, 545, 620), stream=_photo_bytes(d * 100 + p + 7))
        # a vector SITE PLAN page (line art - the render tier, never an embedded raster)
        plan = doc.new_page(width=842, height=595)
        plan.insert_text((60, 40), "SITE PLAN    Scale 1:1250", fontsize=13)
        plan.draw_rect(fitz.Rect(70, 70, 780, 520), color=(0.1, 0.1, 0.1), width=2)
        for i in range(n_props):
            x0 = 90 + i * (660 // max(1, n_props))
            plan.draw_rect(fitz.Rect(x0, 120, x0 + (600 // max(1, n_props)), 400),
                           color=(0.1, 0.1, 0.1), fill=(0.72, 0.78, 0.85), width=2)
            plan.insert_text((x0 + 10, 110), f"Unit {chr(65 + i)}{d + 1}", fontsize=10)
        for gx in range(90, 780, 40):
            plan.draw_line(fitz.Point(gx, 430), fitz.Point(gx, 500), color=(0.3, 0.3, 0.3), width=1)
        name = f"{dev.split()[0]}-{city}-brochure.pdf".replace(" ", "-")
        doc.save(str(folder / name))
        doc.close()
        names.append(name)
    return names


# ---------------------------------------------------------------- state fingerprint

def fingerprint(work: Path, out_dir: Path) -> tuple[str, int, int]:
    """(sha1 of the on-disk state, file count, total bytes). The fingerprint is the
    ONLY reliable cross-process progress signal: a repeat means the round changed
    nothing, i.e. the loop is not converging."""
    entries = []
    total = 0
    for root in (work, out_dir):
        if not root.exists():
            continue
        for f in sorted(root.rglob("*")):
            if f.is_file():
                try:
                    sz = f.stat().st_size
                except OSError:
                    continue
                entries.append(f"{f.relative_to(root).as_posix()}:{sz}")
                total += sz
    h = hashlib.sha1("|".join(entries).encode()).hexdigest()
    return h, len(entries), total


def accounting_check(inputs: list[str], work: Path, out_dir: Path) -> list[str]:
    """THE SILENT-DROP ORACLE. The skill's core promise is that nothing is invented AND
    nothing is silently ignored: every input file must end up either (a) used - it appears
    in the inventory / a record / the ledger - or (b) SURFACED as a gap (unreadable.json,
    the Gaps Report, an inventory 'skipped'/'unclassified' entry). A file that appears
    NOWHERE is a silent drop: the broker believes their data was considered when it never
    was. That is worse than a crash, and it is exactly what a JSON-plus-photos handover
    hits, because intake.py classifies only .pdf/.pptx/.xlsx/.xlsm/.csv/images/.msg/.eml.

    Returns a list of verdict strings ([] == every input accounted for)."""
    blob_parts: list[str] = []
    for root in (work, out_dir):
        if not root.exists():
            continue
        for f in root.rglob("*"):
            if not f.is_file():
                continue
            if f.suffix.lower() in (".json", ".md", ".csv", ".txt", ".html", ".yaml", ".yml"):
                try:
                    blob_parts.append(f.read_text(encoding="utf-8", errors="replace"))
                except Exception:
                    pass
    blob = "\n".join(blob_parts)
    missing = [n for n in inputs if n not in blob]
    out: list[str] = []
    if missing:
        out.append(
            "SILENT DROP: input file(s) appear NOWHERE in the work dir, the ledger or the "
            f"Gaps Report - neither used nor surfaced as a gap: {', '.join(sorted(missing))}. "
            "The broker cannot tell their data was ignored. An unsupported input must still "
            "be recorded as an honest unreadable/unclassified gap.")
    return out


# ---------------------------------------------------------------- the orchestrator's answers

class Responder:
    """Plays the LLM orchestrator: satisfies each exit-code round-trip. `mode` controls
    how ADVERSARIAL the sub-agent's answers are - the point of the whole exercise, since
    a guard that only tolerates a perfect answer is a livelock waiting to happen."""

    def __init__(self, work: Path, mode: str, log):
        self.work = work
        self.mode = mode
        self.log = log
        self.round_seen: dict[str, int] = {}
        self.malformed_done = False

    # -- helpers -------------------------------------------------------
    def _manifests(self) -> list[Path]:
        vis = self.work / "vision"
        return sorted(vis.glob("*.json")) if vis.exists() else []

    def _write(self, rel: str, obj) -> Path:
        f = self.work / rel
        f.parent.mkdir(parents=True, exist_ok=True)
        if self.mode == "malformed" and not self.malformed_done:
            self.malformed_done = True
            f.write_text(json.dumps(obj)[: max(8, len(json.dumps(obj)) // 3)], encoding="utf-8")
            self.log(f"      (wrote TRUNCATED {rel} - simulating a killed sub-agent write)")
            return f
        f.write_text(json.dumps(obj, ensure_ascii=False, indent=2), encoding="utf-8")
        return f

    def _touch(self, rel: str) -> None:
        f = self.work / rel
        f.parent.mkdir(parents=True, exist_ok=True)
        f.write_text("", encoding="utf-8")

    # -- per exit code -------------------------------------------------
    def handle(self, code: int, stderr: str) -> bool:
        """Answer the request. True = answered (re-run), False = cannot answer (stop)."""
        fn = getattr(self, f"exit_{code}", None)
        if fn is None:
            return False
        return fn(stderr)

    def exit_3(self, stderr: str) -> bool:
        """Interpretation: decks / tracker `jobs` / `region_labels`."""
        answered = False
        for man in self._manifests():
            try:
                payload = json.loads(man.read_text(encoding="utf-8"))
            except Exception:
                continue
            if not isinstance(payload, dict):
                continue
            # --- region labels (the 2026-07-30 loop bug lives exactly here) ---
            jobs = payload.get("region_labels") or []
            if jobs:
                res = []
                for i, j in enumerate(jobs):
                    cands = j.get("candidates") or []
                    decline = (self.mode == "decline"
                               or (self.mode == "partial" and i > 0)
                               or not cands)
                    res.append({
                        "raw_label": j.get("raw_label"), "city": j.get("city"),
                        "country_cc": j.get("country_cc"),
                        "code": None if decline else cands[0].get("code"),
                        "matched_name": None if decline else cands[0].get("name"),
                        "confidence": "low" if decline else "high",
                        "reason": "no confident candidate" if decline else "closest name match",
                    })
                out = payload.get("output") or "work/extract/region_labels.json"
                self._write(out.replace("work/", "", 1), {"resolutions": res})
                n_dec = sum(1 for r in res if r["code"] is None)
                self.log(f"      region_labels: answered {len(res)} ({n_dec} DECLINED)")
                answered = True
            # --- tracker mapping jobs: decline via the .SKIP sentinel ---
            # CONTRACT MISMATCH (found 2026-07-30): the manifest instruction (run.py:1368)
            # says "an empty file at the output path with a .SKIP suffix" -> literally
            # <output>.json.SKIP, but run.py:893/902 only ever looks for <stem>.SKIP
            # (.json REPLACED). An orchestrator following the instruction loops forever.
            # We write BOTH spellings: `skip_literal` reproduces the bug, `skip_stem`
            # is what the code actually reads.
            for j in (payload.get("jobs") or []):
                out = (j.get("output") or "").replace("work/", "", 1)
                if not out:
                    continue
                if self.mode != "literal_skip":
                    self._touch(str(Path(out).with_suffix(".SKIP")))  # what the CODE reads
                self._touch(out + ".SKIP")                            # what the DOC says
                self.log(f"      tracker job: wrote .SKIP sentinel(s) for {out}")
                answered = True
            # --- brochure decks needing text/vision interpretation ---
            decks = payload.get("decks") or []
            for di, deck in enumerate(decks):
                if self.mode == "partial" and di > 0:
                    continue  # answer only the FIRST deck: a re-run must re-ask ONLY the rest
                # B51: the manifest hands the routing label as `cluster_label` (it is FILENAME-
                # derived and must never be copied into a record's `region`). Resolve it the way
                # a real orchestrator does - new key first, legacy `region` as the fallback for a
                # warm work dir whose manifest predates the rename.
                region = deck.get("cluster_label") or deck.get("region") or f"deck{di}"
                recs = self._records_from_deck(deck)
                if not recs:
                    continue
                self._write(f"extract/{region}_vision.json", recs)
                self.log(f"      deck '{region}': wrote {len(recs)} record(s)")
                answered = True
        return answered

    def _records_from_deck(self, deck: dict) -> list[dict]:
        """Do what the interpretation sub-agent does: structure each page's TEXT into a
        record per templates/record_schema.json. Deterministic because the simulator
        authored the corpus - the point is to exercise run.py's RESUME/CONVERGENCE
        plumbing, not to test an LLM."""
        import re
        src = deck.get("source_file", "")
        stype = deck.get("source_type", "pdf")
        cc_by_city = {c[0]: c[2] for c in CITIES}
        region_by_city = {c[0]: c[1] for c in CITIES}
        recs: list[dict] = []
        plan_page = None
        for pg in deck.get("pages") or []:
            text = pg.get("text") or ""
            pno = pg.get("page_no")
            if "SITE PLAN" in text.upper():
                plan_page = pno
                continue
            m_area = re.search(r"Warehouse area:\s*([\d,\.]+)\s*sq m", text)
            if not m_area:
                continue  # cover / non-property page
            area = int(m_area.group(1).replace(",", "").split(".")[0])
            city = next((c for c in cc_by_city if c in text), "")
            m_rent = re.search(r"Rent:\s*EUR\s*([\d.]+)\s*per sq m per month", text)
            m_ch = re.search(r"Clear height:\s*([\d.]+)\s*m", text)
            m_dev = re.search(r"Developer:\s*(.+)", text)
            m_park = (text.strip().splitlines() or [""])[0].strip()
            rec: dict = {
                "park": m_park,
                "city": city or "tbd",
                "country": cc_by_city.get(city, "tbd"),
                "region": region_by_city.get(city, "tbd"),
                "warehouseArea": area,
                "areaUnit": "sq m",
                "__meta": {
                    "source_file": src,
                    "source_type": stype,
                    "page_no": pno,
                    "locator_base": pg.get("locator") or f"page {(pno or 0) + 1}",
                },
            }
            if m_dev:
                rec["developer"] = m_dev.group(1).strip()
            if m_ch:
                rec["clearHeight"] = f"{m_ch.group(1)} m"
            if m_rent:  # schema: warehouseRentVal is ANNUAL - annualise the monthly quote
                annual = round(float(m_rent.group(1)) * 12, 2)
                rec["warehouseRent"] = f"€{annual:.2f} / sq m / year"
                rec["warehouseRentVal"] = annual
                rec["rentUnit"] = "€/sq m/yr"
            # the LLM's image picks: heroRef = a real candidate index on THIS page
            cands = pg.get("candidates") or []
            if cands and self.mode != "decline":
                rec["__meta"]["heroRef"] = cands[0].get("index", 0)
            elif self.mode == "decline":
                rec["__meta"]["heroRef"] = None  # "no real photo here" - a LEGITIMATE answer
            recs.append(rec)
        if plan_page is not None and recs and self.mode != "decline":
            recs[0]["__meta"]["plan_page"] = plan_page
        return recs

    def exit_9(self, stderr: str) -> bool:
        man = self.work / "photo_match_manifest.json"
        try:
            payload = json.loads(man.read_text(encoding="utf-8"))
        except Exception:
            return False
        items = payload.get("brochures") or payload.get("items") or []
        props = payload.get("properties") or []
        out = {"confident": [], "uncertain": [], "unrelated": []}
        for i, b in enumerate(items):
            key = b.get("file") or b.get("source_file") or str(b)
            if self.mode == "decline" or not props:
                out["uncertain"].append({"file": key, "reason": "no confident match"})
            else:
                pid = props[i % len(props)]
                out["confident"].append({"file": key,
                                         "property_id": pid.get("id", i + 1)
                                         if isinstance(pid, dict) else i + 1})
        self._write("photo_map.json", out)
        self.log(f"      photo_map: {len(out['confident'])} confident, "
                 f"{len(out['uncertain'])} uncertain")
        return True

    def exit_10(self, stderr: str) -> bool:
        try:
            cand = json.loads((self.work / "match_candidates.json").read_text(encoding="utf-8"))
        except Exception:
            return False
        pairs = cand.get("pairs") or []
        confs = cand.get("field_conflicts") or []
        # ECHO THE ID VERBATIM. Both guards look the answer up BY ID (`md.get(pair_id)` /
        # `fd.get(conflict_id)`), so an answer that omits the id can never be matched and the
        # round-trip loops forever. Getting this wrong in an early version of this responder
        # produced a 25-round livelock that looked like a skill bug and was not - the ids are
        # right there in match_candidates.json. A real sub-agent is told the same thing.
        # 'decline' = 'different' / keep the default: both are LEGITIMATE answers that bind
        # nothing, which is the shape that broke the region-label guard.
        self._write("match_decisions.json", {"decisions": [
            {"pair_id": p.get("pair_id", i),
             "a": p.get("a"), "b": p.get("b"),
             "verdict": "different" if self.mode == "decline" else "same",
             "reason": "simulated"} for i, p in enumerate(pairs)]})
        self._write("field_decisions.json", {"decisions": [
            {"conflict_id": c.get("conflict_id"), "field": c.get("field"),
             "pick": "default", "reason": "simulated"} for c in confs]})
        self.log(f"      adjudicated {len(pairs)} pair(s), {len(confs)} conflict(s)")
        return True

    def exit_11(self, stderr: str) -> bool:
        d = self.work / "i18n"
        if not d.exists():
            return False
        for req in d.glob("*_request.json"):
            code = req.name.split("_request")[0]
            self._touch(f"i18n/{code}.SKIP")
            self.log(f"      i18n: wrote {code}.SKIP (fall back to English)")
        return True

    def exit_12(self, stderr: str) -> bool:
        self._touch("i18n/data_translate.SKIP")
        self.log("      data translation: wrote data_translate.SKIP")
        return True


# ---------------------------------------------------------------- the driver

def run_round(folder: Path, work: Path, client: str, cap: float, flags: list[str]):
    cmd = [sys.executable, str(RUN_PY), "--folder", str(folder), "--work", str(work),
           "--client", client] + flags
    t0 = time.time()
    killed = False
    try:
        pr = subprocess.run(cmd, capture_output=True, text=True, timeout=cap,
                            encoding="utf-8", errors="replace")
        rc, out, err = pr.returncode, pr.stdout or "", pr.stderr or ""
    except subprocess.TimeoutExpired as e:
        killed = True
        rc = -9
        out = (e.stdout or "") if isinstance(e.stdout, str) else ""
        err = (e.stderr or "") if isinstance(e.stderr, str) else ""
    return rc, out, err, time.time() - t0, killed


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--cap", type=float, default=40.0,
                    help="wall-clock kill per invocation, seconds (Cowork is ~40-45)")
    ap.add_argument("--rounds", type=int, default=25, help="max rounds before declaring a loop")
    ap.add_argument("--decks", type=int, default=4)
    ap.add_argument("--props", type=int, default=3)
    ap.add_argument("--answers", default="clean",
                    choices=["clean", "decline", "partial", "malformed", "literal_skip"],
                    help="literal_skip = obey run.py:1368's wording exactly (<output>.json.SKIP) "
                         "and nothing else - reproduces the tracker decline livelock")
    ap.add_argument("--inputs", default="pdf",
                    choices=["pdf", "pptx", "xlsx", "xlsx+images", "images", "json",
                             "json+images", "txt+images", "mixed", "empty"],
                    help="the INPUT SHAPE a broker hands over (brokers send all of these; "
                         "the pipeline is written around the brochure case)")
    ap.add_argument("--flags", default="", help="extra run.py flags, comma-separated "
                                               "(e.g. regions,geocode)")
    ap.add_argument("--keep", action="store_true", help="keep the temp corpus")
    args = ap.parse_args()

    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

    flags = [f"--{f.strip()}" for f in args.flags.split(",") if f.strip()]
    tmp = Path(tempfile.mkdtemp(prefix="cowork_sim_"))
    folder, work, out_dir = tmp / "inputs", tmp / "work", tmp / "out"
    lines: list[str] = []

    def log(msg=""):
        print(msg, flush=True)
        lines.append(msg)

    log(f"COWORK SIM  inputs={args.inputs}  cap={args.cap}s  answers={args.answers}  "
        f"corpus={args.decks}x{args.props}  flags={flags or 'none'}")
    log(f"  temp: {tmp}")
    try:
        names = build_shape(folder, args.inputs, args.decks, args.props)
        folder.mkdir(parents=True, exist_ok=True)
        log(f"  corpus: {len(names)} file(s) [{', '.join(names[:6])}"
            f"{' ...' if len(names) > 6 else ''}], "
            f"{sum(f.stat().st_size for f in folder.iterdir()) // 1024} KB")
    except ImportError as e:
        log(f"  SKIPPED: shape '{args.inputs}' needs a library that is absent ({e})")
        return 0
    except Exception as e:
        log(f"  FATAL: corpus build failed ({type(e).__name__}: {e})")
        return 1

    resp = Responder(work, args.answers, log)
    seen: dict[tuple, int] = {}
    hist: list[dict] = []
    verdict: list[str] = []
    converged = False

    log("\n  rnd  exit  secs  killed  files  bytes   note")
    for rnd in range(1, args.rounds + 1):
        rc, out, err, secs, killed = run_round(folder, work, "SimCo", args.cap, flags)
        fp, nf, nb = fingerprint(work, out_dir)
        note = ""
        hist.append({"rnd": rnd, "rc": rc, "secs": secs, "killed": killed,
                     "files": nf, "bytes": nb, "fp": fp})
        key = (rc, fp)
        if key in seen and not killed:
            note = f"LOOP (identical state as round {seen[key]})"
            verdict.append(
                f"NON-CONVERGENCE: round {rnd} exited {rc} with a work-dir state identical "
                f"to round {seen[key]}. The round changed NOTHING yet the same request was "
                f"re-emitted - this is the 2026-07-30 region-label bug's shape (an "
                f"'already answered' guard narrower than the set of legitimate answers).")
        seen.setdefault(key, rnd)

        # PROGRESS: a killed round must still leave more state behind than the last one
        if killed and len(hist) >= 2:
            prev = hist[-2]
            if nf <= prev["files"] and nb <= prev["bytes"]:
                note = (note + " " if note else "") + "STALL (killed, no new state)"
                verdict.append(
                    f"NO FORWARD PROGRESS: round {rnd} was killed at the {args.cap}s cap and "
                    f"left no more state on disk than round {prev['rnd']} "
                    f"({nf} files/{nb} bytes vs {prev['files']}/{prev['bytes']}). Under a real "
                    f"cap this never finishes - the 2026-07-30 placeholder-audit bug's shape "
                    f"(expensive work redone uncached every run).")
        # REGRESSION: only when BOTH the file count and the byte total fall. Bytes alone
        # legitimately shrink as the request manifest empties out (each answered deck is
        # removed from it) - flagging that was a simulator false positive, not a skill bug.
        if len(hist) >= 2 and nb < hist[-2]["bytes"] and nf < hist[-2]["files"]:
            note = (note + " " if note else "") + "REGRESSION (state shrank)"
            verdict.append(f"NON-MONOTONIC: round {rnd} REDUCED on-disk state "
                           f"({nf} files/{nb} bytes vs {hist[-2]['files']}/{hist[-2]['bytes']}) "
                           f"- a re-run destroyed work a previous run completed.")

        log(f"  {rnd:>3}  {rc:>4}  {secs:>4.1f}  {str(killed):>6}  {nf:>5}  {nb:>6}   {note}")

        if rc == 0:
            converged = True
            log(f"\n  CONVERGED at round {rnd}.")
            break
        if killed:
            continue  # the cap fired: resume is the whole point, just re-run
        if rc in (2, 4, 5, 6, 7):
            log(f"      run.py exited {rc} (a hard failure/gate block, not a round-trip).")
            log("      stderr tail: " + " | ".join(
                [l for l in err.strip().splitlines()[-4:] if l.strip()]))
            # exit 2 on a corpus that genuinely contains NO source the pipeline can read is the
            # CORRECT outcome, not a defect - the test that matters there is whether every input
            # was ACCOUNTED FOR (checked below), i.e. named rather than silently dropped. Only
            # flag exit 2 when a readable source WAS present and still produced nothing.
            _readable = {".pdf", ".pptx", ".xlsx", ".xlsm", ".csv", ".msg", ".eml"}
            _have = {Path(n).suffix.lower() for n in names}
            if rc == 2 and not (_have & _readable):
                log("      -> CORRECT: this corpus holds no source type the pipeline reads, so "
                    "an honest refusal is the right answer. Checking it was SURFACED, not dropped.")
            else:
                verdict.append(f"BLOCKED: exit {rc} with a readable source present "
                               f"({sorted(_have & _readable) or 'none'}). Either the corpus is "
                               f"unrealistic or a gate/validator is wrong.")
            break
        if not resp.handle(rc, err):
            log(f"      UNHANDLED exit {rc} - the simulator cannot answer this round-trip.")
            log("      stdout tail: " + " | ".join(
                [l for l in out.strip().splitlines()[-3:] if l.strip()]))
            verdict.append(f"UNHANDLED exit {rc}: extend the Responder to cover it "
                           f"(this is a simulator gap, not necessarily a skill bug).")
            break
    else:
        verdict.append(f"NON-CONVERGENCE: hit the {args.rounds}-round bound without exit 0.")

    # determinism: one more clean run must be a no-op and must not change the output
    if converged:
        html = sorted(out_dir.glob("*.html")) + sorted(work.glob("**/*.html"))
        sha_before = (hashlib.sha256(html[0].read_bytes()).hexdigest() if html else None)
        rc2, _o2, _e2, secs2, killed2 = run_round(folder, work, "SimCo", args.cap, flags)
        sha_after = (hashlib.sha256(html[0].read_bytes()).hexdigest() if html else None)
        log(f"  re-run after convergence: exit {rc2} in {secs2:.1f}s "
            f"(killed={killed2}); output sha {'UNCHANGED' if sha_before == sha_after else 'CHANGED'}")
        if rc2 != 0:
            verdict.append(f"NOT IDEMPOTENT: a re-run of a CONVERGED work dir exited {rc2}, "
                           f"not 0. The orchestrator's final verify would re-open the pipeline.")
        if sha_before and sha_before != sha_after:
            verdict.append("NOT DETERMINISTIC: re-running a converged work dir changed the "
                           "output HTML bytes.")
        # only meaningful when the cold run was big enough to HAVE savings: on a 2s corpus
        # (one tracker row, no images) fixed start-up dominates and the ratio is noise, which
        # produced a false "no resume saving" flag. 4s is comfortably above start-up cost.
        _cold_max = max((h["secs"] for h in hist), default=0.0)
        if not killed2 and _cold_max >= 4.0 and secs2 > _cold_max * 0.9:
            verdict.append(f"NO RESUME SAVING: the re-run of a fully-cached work dir took "
                           f"{secs2:.1f}s vs a {_cold_max:.1f}s cold round - work is being "
                           f"redone rather than read from cache.")

    # every input file must be USED or SURFACED - never silently ignored
    acct = accounting_check(names, work, out_dir)
    for a in acct:
        verdict.append(a)
    if names and not acct:
        log(f"  accounting: all {len(names)} input file(s) used or surfaced.")

    log("\n" + "=" * 78)
    if verdict:
        log(f"DEFECTS FLAGGED: {len(verdict)}")
        for i, v in enumerate(verdict, 1):
            log(f"  {i}. {v}")
    else:
        log("NO DEFECTS FLAGGED: converged, monotonic, idempotent, deterministic.")
    log("=" * 78)

    report = tmp / "sim_report.txt"
    try:
        report.write_text("\n".join(lines), encoding="utf-8")
        log(f"report: {report}")
    except Exception:
        pass
    if not args.keep:
        try:
            shutil.rmtree(tmp, ignore_errors=True)
        except Exception:
            pass
    else:
        log(f"kept: {tmp}")
    return 1 if verdict else 0


if __name__ == "__main__":
    sys.exit(main())
