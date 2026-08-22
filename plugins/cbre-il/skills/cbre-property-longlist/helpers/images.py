#!/usr/bin/env python3
"""images.py - harvest, compress and base64-embed property imagery.

Hero images come from the brochure pages (embedded rasters). Each image is
resized and JPEG-compressed to a per-image byte budget so the assembled
single-file dashboard stays a sane size (the reference was 11 MB; we target
~80-120 KB/image). Returns data: URIs ready for the canonical 'photo'/'plan'
fields. EMF/WMF vectors that Pillow cannot read are skipped (caller falls back
to a page raster or the placeholder).
"""
from __future__ import annotations

import base64
import hashlib
import io
import json
import math
import os
import re
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
try:
    import fitz  # PyMuPDF
except Exception:  # sandbox without PyMuPDF: pypdfium2/pdfplumber shim (same 9-call surface)
    import fitz_shim as fitz
try:
    fitz.TOOLS.mupdf_display_errors(False)  # silence MuPDF C-level warnings (broker quiet mode)
except Exception:
    pass
try:
    from PIL import Image
    _HAS_PIL = True
except ImportError:  # Cowork sandbox without Pillow and no pip: degrade, never crash at import.
    # Pillow has no pure-python shim - but the skill MAY ship a bundled Pillow wheel; try
    # it (a NO-OP unless it matches this interpreter). If still absent, the whole image
    # ladder no-ops to the pre-baked placeholder asset. `from __future__ import annotations`
    # (top of file) keeps every `Image.Image` type hint a lazy string, so rebinding is safe.
    try:
        import _vendor_wheels as _vw
        _vw.ensure("PIL", "pillow")
        from PIL import Image
        _HAS_PIL = True
    except Exception:
        Image = None  # type: ignore[assignment]
        _HAS_PIL = False

if _HAS_PIL:
    try:
        import pillow_heif  # noqa: registers HEIC opener
        pillow_heif.register_heif_opener()
    except Exception:
        pass

HERO_MAX_EDGE = 1280
PLAN_MAX_EDGE = 1100
DEFAULT_BUDGET_KB = 110
GALLERY_MAX = 6  # max photos attached per property for the carousel (hero + up to 5 more),
#                  filled BEST-FIRST by photographic_score; the rest are noted in the Gaps Report.
#                  The single self-contained HTML embeds every image as base64, so this cap keeps
#                  the file portable (shareable/emailable) even on an image-heavy deck.
#                  DELIBERATELY 6, re-confirmed when the carousel's page scope was widened: the
#                  cap was never what starved a card (no property was short of SLOTS, they were
#                  short of CANDIDATES), 6 is the ergonomic size for a MANUAL prev/next carousel,
#                  a brochure holds 2-10 distinct photographs (median 5) so the 7th is a
#                  near-duplicate elevation, and 6-per-property lands the self-contained HTML at
#                  ~19 MB - already the edge of emailable, where 8 would push ~24 MB.
MIN_HERO_W, MIN_HERO_H = 320, 200  # reject logos/icons (PHOTO hero floor)
# --- CAROUSEL (gallery) floors - STRICTER than the hero's, and for different reasons ------- #
# The hero floor exists to reject logos/icons; anything above it can lead a card because the
# alternative is the placeholder. A CAROUSEL entry has a real alternative - not being there -
# so it is held to what the card actually RENDERS. The card thumb is
# `repeat(auto-fill,minmax(340px,1fr))` in a 1400px container (~432 CSS px wide) at
# `aspect-ratio:16/10` => ~432x270 CSS px. 640x400 is that box at 1.5x device-pixel-ratio: never
# upscaled on a standard display, still sharp on a HiDPI laptop, 2x the linear size of the
# 323x215 thumbnails this floor exists to stop, and it rejects banner strips (3242x250) on the
# height term without a separate rule.
MIN_GALLERY_W, MIN_GALLERY_H = 640, 400
# DETAIL floor (see detail_score): the discriminator between a PHOTOGRAPH and flat vector /
# gradient DECORATION. Neither photographic_score nor classify_image separates them - a flat
# yellow geometric background scored 7.8 (over MODEST_PHOTO) and smooth gradient art classifies
# as 'photo' (high luminance entropy, no flat 8x8 blocks). Measured over 159 candidates in 15
# real decks: decorative art of EVERY kind sits at 0.03-1.84, real photographs at 3.33-20+, with
# an empty band between. 2.5 sits inside that band with ~1.4x margin on both sides.
MIN_GALLERY_DETAIL = 2.5
# plans get their OWN, lower floor: a usable site plan is often small, and the
# photo floor silently discarded it before any scoring happened (a real
# placeholder-shipped-despite-usable-plan failure). Below even this = icon.
MIN_PLAN_W, MIN_PLAN_H = 220, 160
# score at/above which an embedded image (or a cropped image REGION) is a real,
# if unspectacular, photo and ALWAYS beats the whole-page raster. With the
# flatness multiplier (see photographic_score), real photos measure 9-27 and
# plans/maps/logos 0-13 on a real Spanish deck (TEDi calibration); the multiplier
# (not this floor) is what reranks photo-over-plan when both are present.
# Without this floor, a busy page could out-score its own photo and ship a
# cluttered full-page tile as the hero.
MODEST_PHOTO = 6.0


def _atomic_save_png(im, path):
    """tmp + os.replace for the audit/candidate dumps. These are the evidence the G-images
    reviewer signs off, and a kill mid-encode otherwise freezes a PARTIAL set at the final
    names - which every resume guard then accepts. _common is imported lazily so images.py
    keeps importing on a host missing an optional dep. No direct-save fallback: that would
    reinstate the non-atomic write this exists to remove, and both call sites already skip
    the file on an exception. (B16)"""
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    import _common as _C
    return _C.atomic_save_image(im, path, "PNG")


def to_data_uri(jpeg_bytes: bytes, mime: str = "image/jpeg") -> str:
    return f"data:{mime};base64," + base64.b64encode(jpeg_bytes).decode("ascii")


def compress(img: Image.Image, max_edge: int = HERO_MAX_EDGE,
             budget_kb: int = DEFAULT_BUDGET_KB) -> bytes:
    """Resize to max_edge and step quality down until under budget."""
    img = img.convert("RGB")
    w, h = img.size
    scale = min(1.0, max_edge / max(w, h))
    if scale < 1.0:
        img = img.resize((max(1, int(w * scale)), max(1, int(h * scale))), Image.LANCZOS)
    for q in (78, 70, 62, 54, 46, 38):
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=q, optimize=True)
        data = buf.getvalue()
        if len(data) <= budget_kb * 1024 or q == 38:
            return data
    return data  # type: ignore[return-value]


def _open(image_bytes: bytes) -> Image.Image | None:
    if Image is None:
        return None  # no Pillow: no decoder; callers treat None as 'skip this candidate'
    try:
        return Image.open(io.BytesIO(image_bytes))
    except Exception:
        return None


def page_embedded_images(doc: "fitz.Document", page_index: int) -> list[dict]:
    """All decodable raster images on a page, largest first. Memoised per (doc, page) so a
    page's rasters are decoded ONCE per run (the interpret thumbnail pass and merge's hero
    binding both ask for them). The memo is cleared with the doc cache. (#39/#21)"""
    key = (id(doc), page_index)
    hit = _EMBED_CACHE.get(key)
    if hit is not None:
        return hit
    out = []
    page = doc[page_index]
    for info in page.get_images(full=True):
        xref = info[0]
        try:
            ext = doc.extract_image(xref)
            img = _open(ext["image"])
            if img is None:
                continue
            w, h = img.size
            out.append({"img": img, "w": w, "h": h, "area": w * h, "xref": xref})
        except Exception:
            continue
    out.sort(key=lambda d: -d["area"])
    _EMBED_CACHE[key] = out
    return out


def candidates_for_page(path: Path, page_index: int) -> list[dict]:
    """The page's HERO-SIZE embedded image candidates, in a STABLE 0-based order -
    the ONE list both sides of the LLM-hero contract share. interpret_prep writes a
    thumbnail per entry (the sub-agent LOOKS at them and references the `index`); merge
    re-derives the SAME list (embedded_by_index) to bind the chosen `index`. Both call
    THIS function on the same filtered order, so the index can never disagree.

    The order is page_embedded_images() (largest-area first) FILTERED to the hero size
    floor (>= MIN_HERO_W x MIN_HERO_H, which screens out logos/icons), and the 0-based
    POSITION in that filtered list IS the candidate's `index`. A PDF page reads its
    embedded rasters; a .pptx slide reads its slide pictures; any other kind (or an open
    failure) yields [] so the caller falls back to the deterministic ladder gracefully.

    Returns [{index, img, w, h}] - `img` is a live PIL image (callers that only need the
    metadata, e.g. a manifest, ignore it). Never raises."""
    try:
        path = Path(path)
        if path.suffix.lower() == ".pptx":
            raw = slide_pictures(path, page_index)
        else:
            raw = page_embedded_images(_get_doc(path), page_index)
    except Exception:
        return []
    out: list[dict] = []
    for im in raw:
        if im.get("w", 0) >= MIN_HERO_W and im.get("h", 0) >= MIN_HERO_H:
            out.append({"index": len(out), "img": im["img"], "w": im["w"], "h": im["h"]})
    return out


def embedded_by_index(path: Path, page_index: int, index: int,
                      budget_kb: int = DEFAULT_BUDGET_KB,
                      cache_dir: Path | str | None = None) -> str | None:
    """The candidate at `index` of candidates_for_page(path, page_index), compressed to a
    hero data URI - the binding half of the LLM-hero contract (merge passes the sub-agent's
    chosen __meta.heroRef here). None when the index is out of range, Pillow is absent, or
    anything fails (the caller then falls back to the deterministic ladder). Never raises.

    DISK-CACHED per (deck, page, budget, index). This was the ONLY image producer with no
    cache at all, while every sibling (page_hero_and_plan, page_render_plan, _deck_page_photos,
    slide_hero_and_plan) had one - and it is the tier that WINS over all others, i.e. the
    intended happy path whenever the interpretation sub-agent supplies a heroRef. Measured at
    81% of a warm 40-property merge (2.56s of 3.16s; 64 ms/property, 6 JPEG-ladder rungs each,
    of which the LANCZOS resize is ~60 ms).
    THE INDEX IS PART OF THE KEY, deliberately: two properties anchored on the SAME page with
    DIFFERENT heroRef values must not collide, or they would silently swap heroes. planRef
    rides this same function, and gets its own key for free the same way."""
    if Image is None:
        return None  # no Pillow: no decoder; merge falls back to the deterministic ladder
    if not (isinstance(index, int) and not isinstance(index, bool) and index >= 0):
        return None
    cf = _cache_file(path, page_index, budget_kb, f"heroref{index}", cache_dir)
    cached = _cache_read(cf)
    if cached is not None:
        return cached or None
    try:
        cands = candidates_for_page(path, page_index)
        if index >= len(cands):
            return None
        uri = to_data_uri(compress(cands[index]["img"], HERO_MAX_EDGE, budget_kb))
    except Exception:
        return None
    _cache_write(cf, uri)
    return uri


_PIX_MODE = {(1, False): "L", (1, True): "LA", (3, False): "RGB",
             (3, True): "RGBA", (4, False): "CMYK"}


def page_raster(doc: "fitz.Document", page_index: int, dpi: int = 150) -> Image.Image:
    """Render a whole page to a raster (EMF/vector fallback).

    Wraps the pixmap's RAW SAMPLES directly instead of round-tripping them through a PNG
    encode+decode. This is the hottest primitive in the skill - every render tier goes
    through it (_page_crops, _rendered_plan_crop, the hero tier B/C ladder, vision/interpret
    prep, contact_sheet) - and the PNG detour cost ~6x: measured 86.8 ms -> 14.6 ms per A4
    page at 150 dpi, with `fz_write_pixmap_as_png` alone accounting for ~40% of a warm merge
    profile. At 10 decks the cold interpretation prep was 48 s, i.e. OVER the ~40-45 s shell
    cap, forcing an extra kill-and-resume round purely to pay for PNG encoding.

    BYTE-IDENTICAL, verified before switching: the samples ARE the decoded PNG's pixels, so
    pixels, the classifier verdict/photographic score AND the downstream `to_data_uri(compress(
    ...))` all match byte for byte at 90/150/180 dpi. That last one is what protects the
    skill's byte-identity contract and the cached `.uri` units. An unexpected colorspace
    (CMYK+alpha, a separation space) or a shim Pixmap without `.samples` falls back to the
    original PNG path rather than guessing a mode."""
    if Image is None:
        return None  # no Pillow: no raster; callers None-check or wrap in try/except
    pix = doc[page_index].get_pixmap(dpi=dpi)
    try:
        ncomp = pix.colorspace.n if pix.colorspace is not None else 1
        mode = _PIX_MODE.get((ncomp, bool(pix.alpha)))
        if mode is not None:
            return Image.frombytes(mode, (pix.width, pix.height), pix.samples)
    except Exception:
        pass  # fall through to the universally-correct PNG path
    return Image.open(io.BytesIO(pix.tobytes("png")))


def photographic_score(img) -> float:
    """A ~0..60 'is this a real photograph?' score from PURE PIXEL STATISTICS, used
    to pick a property hero over a floor plan, logo, branded divider or map. Nothing
    about any client, brand, colour, language or region is encoded - so it works the
    same on a Spanish, German, Polish or UK deck. Combines:
      * colourfulness (Hasler-Susstrunk) - photos are colourful; line art/logos/plans are not,
      * non-white fraction - floor plans and paper docs are mostly white,
      * an aspect-ratio sanity penalty - logos/banners are extreme, photos are not,
      * a FLATNESS multiplier - the share of pixels in the 5 dominant quantised
        colours. A COLOURFUL site plan (blue floor fill + green landscaping) beat a
        real aerial photo on colourfulness alone in a real run; but plans/maps are
        dominated by flat fills (flat5 0.55-0.95 measured) while photos spread
        across shades (0.19-0.37), so (1 - flat5) reranks photo-over-plan without
        rejecting 3D renders (~0.7) when they are the page's only marketing image,
      * a SINGLE-DOMINANT-COLOUR penalty - the share of the ONE most common colour: a
        logo, a solid 'photo pending' holding card, a road-map screenshot and a
        text-on-white page are each dominated by one colour, while a photo, a render and
        a satellite/aerial are not, so it demotes the non-property imagery the harvester
        used to mistake for a hero."""
    try:
        im = img.convert("RGB")
    except Exception:
        return 0.0
    w0, h0 = im.size
    if w0 < 2 or h0 < 2:
        return 0.0
    im.thumbnail((96, 96))  # small sample - the stats are scale-stable and this keeps it fast
    px = list(im.getdata())
    n = len(px) or 1
    s_rg = s_yb = s_rg2 = s_yb2 = 0.0
    white = 0
    from collections import Counter
    bins: Counter = Counter()
    for r, g, b in px:
        a = r - g
        c = 0.5 * (r + g) - b
        s_rg += a; s_yb += c; s_rg2 += a * a; s_yb2 += c * c
        if r >= 235 and g >= 235 and b >= 235:
            white += 1
        bins[(r >> 4, g >> 4, b >> 4)] += 1
    m_rg, m_yb = s_rg / n, s_yb / n
    std = math.sqrt(max(0.0, s_rg2 / n - m_rg * m_rg) + max(0.0, s_yb2 / n - m_yb * m_yb))
    mean = math.sqrt(m_rg * m_rg + m_yb * m_yb)
    colourfulness = std + 0.3 * mean
    white_frac = white / n
    flat5 = sum(c for _, c in bins.most_common(5)) / n
    flat1 = bins.most_common(1)[0][1] / n if bins else 0.0
    # SINGLE-DOMINANT-COLOUR penalty (added for the live-run feedback): a logo, a solid
    # 'photo pending' holding card, a ROAD-map screenshot (dominant map-paper colour) and a
    # text page on white all have ONE quantised colour over a large share of the image; a
    # real photo, a 3D render and a SATELLITE/aerial never do. So this demotes exactly the
    # non-property imagery Tier A used to pick, WITHOUT touching a genuine hero. Calibrated
    # conservatively - no penalty below 0.45 (a photo with a big sky or wall is safe), full
    # demotion by 0.80 (logos/solids) - and it is pure pixel statistics, so language/client agnostic.
    single_pen = 1.0 if flat1 <= 0.45 else max(0.0, 1.0 - (flat1 - 0.45) / 0.35)
    ratio = w0 / max(1, h0)
    aspect_pen = 1.0 if 0.45 <= ratio <= 2.4 else (0.6 if 0.3 <= ratio <= 3.5 else 0.3)
    return colourfulness * (1.0 - white_frac) * aspect_pen * (1.0 - flat5) * single_pen


# --- hero-kind classifier (the hero ladder + G-images gate) ------------------- #
# A card's hero must be the REAL photographic content on the page - a warehouse/site
# PHOTO, an AERIAL, or a 3-D RENDER - never a road MAP, a flat PLAN diagram, or a
# whole-slide screenshot/text raster. photographic_score (colourfulness) alone mis-ranks
# these on real CBRE decks: a Google-Maps screenshot scored 18.85 and a real warehouse
# photo 3.66. The discriminator is CONTINUOUS TONE, not colour: a photo/aerial/render
# spreads luminance smoothly (high luminance entropy, few flat-fill blocks); a map, plan,
# icon or text raster is piecewise-flat (low entropy and/or large uniform blocks) or a
# single dominant colour. Pure pixel statistics - client/brand/language-agnostic.
# Thresholds calibrated on the real TEDi ES decks (photos/aerials/renders lumEnt >= ~6.6,
# site plan ~6.06, road map ~4.3, icons <= ~4).
PHOTO_LUM_ENT = 6.3     # luminance entropy (bits) at/above which an image is continuous-tone
PHOTO_FLAT_MAX = 0.18   # max share of near-uniform 8x8 luminance blocks for a photo
MAP_PALE_MIN = 0.40     # share of pale low-saturation non-white pixels = road-map palette
LOGO_FLAT1_MAX = 0.55   # single dominant quantised-colour share = logo / solid / holding card
TEXT_WHITE_MIN = 0.55   # mostly-white + low tonal variety = a text/doc raster


def _hero_signals(img) -> dict:
    """Cheap continuous-tone statistics on a 96px thumbnail for classify_image."""
    im = img.convert("RGB")
    im.thumbnail((96, 96))
    w, h = im.size
    n = w * h or 1
    from collections import Counter
    bins: Counter = Counter()
    white = pale = 0
    for r, g, b in im.getdata():
        bins[(r >> 4, g >> 4, b >> 4)] += 1
        if r >= 235 and g >= 235 and b >= 235:
            white += 1
        else:
            mx, mn = max(r, g, b), min(r, g, b)
            if mx >= 150 and (mx - mn) <= 40:  # bright + desaturated = map land/road paper
                pale += 1
    L = im.convert("L")
    hist = L.histogram()
    lum_ent = -sum((c / n) * math.log2(c / n) for c in hist if c)
    # near-uniform 8x8 luminance blocks on a 48x48 grid (flat fills of plans/maps/icons)
    Ls = list(L.resize((48, 48)).getdata())
    flat = tot = 0
    for by in range(0, 48, 8):
        for bx in range(0, 48, 8):
            blk = [Ls[(by + dy) * 48 + (bx + dx)] for dy in range(8) for dx in range(8)]
            tot += 1
            if max(blk) - min(blk) <= 12:
                flat += 1
    return {"white": white / n, "pale": pale / n, "lum_ent": lum_ent,
            "flat": flat / (tot or 1),
            "flat1": (bins.most_common(1)[0][1] / n if bins else 0.0)}


def classify_image(img, sig: dict | None = None) -> str:
    """Coarse hero-relevant kind from pure pixel statistics:
      'photo' - a real photograph, aerial or 3-D render (continuous tone): the ONLY kind
                that should lead a card. A near-monochrome warehouse interior, a satellite
                aerial and a marketing render all qualify; colour is NOT required.
      'plan'  - a flat-fill site/floor diagram (line art, saturated fills).
      'map'   - a road-map screenshot (pale desaturated land/road palette).
      'text'  - a text/doc raster (mostly white, little tonal variety) - incl. a slide screenshot.
      'logo'  - a solid / holding card / single-colour mark.
    Never raises (a stats failure returns 'photo' so the harvest is never blocked)."""
    if Image is None:
        return "photo"
    try:
        s = sig or _hero_signals(img)
    except Exception:
        return "photo"
    if s["flat1"] >= LOGO_FLAT1_MAX:
        return "logo"
    if s["lum_ent"] >= PHOTO_LUM_ENT and s["flat"] <= PHOTO_FLAT_MAX:
        return "photo"
    if s["white"] >= TEXT_WHITE_MIN and s["lum_ent"] < 5.0:
        return "text"
    if s["pale"] >= MAP_PALE_MIN:
        return "map"
    return "plan"


# hero ladder: a photo/aerial/render leads; a plan beats a map beats text/logo; only drop
# a tier when nothing higher exists on the property's pages (plans + maps still go in the
# Site Plan toggle, they just stop being the first impression).
HERO_TIER = {"photo": 0, "plan": 2, "map": 3, "text": 4, "logo": 5}


def detail_score(img) -> float:
    """Mean absolute adjacent-pixel LUMINANCE STEP on a 512-long-edge greyscale - 'how much
    fine detail does this image actually carry?'. Pure pixel statistics, so it is client /
    brand / language / region agnostic like every other signal here, and it answers the ONE
    question classify_image and photographic_score cannot: PHOTOGRAPH vs DECORATION.

    A photograph carries texture everywhere - foliage, brick, tarmac, cladding seams, sensor
    noise - so neighbouring pixels differ constantly. Synthetic marketing decoration (a
    gradient mesh, a flat geometric background, a light-streak render, a brand wash) is smooth
    by construction: neighbouring pixels are nearly equal everywhere except at a handful of
    banding edges. classify_image reads a smooth gradient as continuous tone and calls it
    'photo'; photographic_score reads it as colourful and scores it 50-75. This does not.

    Measured over 159 candidate rasters in 15 real industrial decks: decoration of every
    classified kind lands 0.03-1.84, real photographs 3.33 and up, with nothing in between.
    Sampled every other row/column - the statistic is scale-stable and this keeps it cheap.
    Never raises: 0.0 on any failure, which reads as 'flat' and fails the floor (fail-closed -
    an image we cannot measure must not silently enter a client-facing carousel)."""
    if Image is None:
        return 0.0
    try:
        im = img.convert("L")
        w, h = im.size
        if w < 2 or h < 2:
            return 0.0
        s = min(1.0, 512 / max(w, h))
        if s < 1.0:
            im = im.resize((max(2, int(w * s)), max(2, int(h * s))), Image.LANCZOS)
        px = im.load()
        W, H = im.size
        total = 0
        n = 0
        for y in range(0, H, 2):
            for x in range(1, W, 2):
                total += abs(px[x, y] - px[x - 1, y])
                n += 1
        return total / n if n else 0.0
    except Exception:
        return 0.0


def gallery_admissible(entry: dict) -> bool:
    """Is one `_deck_photo_index` entry fit for the CAROUSEL a client sees?

    Three floors, each closing a hole a real run shipped through:
      * kind == 'photo'  - the carousel is the PHOTO carousel. A plan/map has its own Site
        Plan slot and toggle in the modal, so admitting it here only duplicates it - and the
        admission path that let a non-photo in (`score >= MODEST_PHOTO`) is exactly how a flat
        decorative background graphic reached a card as a 'photo'.
      * MIN_GALLERY_W/H  - what the card actually renders (four 323x215 thumbnails shipped).
      * MIN_GALLERY_DETAIL - photograph vs decoration (gradient art classifies as 'photo').

    FAIL-CLOSED on an entry written by an older cache that carries no measurements: an image
    whose admissibility cannot be established does not enter a client-facing carousel."""
    if not isinstance(entry, dict):
        return False
    if entry.get("kind") != "photo":
        return False
    w, h, d = entry.get("w"), entry.get("h"), entry.get("detail")
    if not all(isinstance(v, (int, float)) and not isinstance(v, bool) for v in (w, h, d)):
        return False
    return w >= MIN_GALLERY_W and h >= MIN_GALLERY_H and d >= MIN_GALLERY_DETAIL


def uri_gallery_admissible(uri: str) -> bool:
    """gallery_admissible for an already-compressed 'data:image/...;base64,...' URI - the form
    the bound HERO takes by the time merge can compare it against the carousel floors. Decodes
    once and re-derives the same three measurements. False on any decode/stats failure."""
    if Image is None or not (isinstance(uri, str) and "base64," in uri):
        return False
    try:
        import base64
        raw = base64.b64decode(uri.split("base64,", 1)[1])
        img = Image.open(io.BytesIO(raw))
        img.load()
    except Exception:
        return False
    try:
        w, h = img.size
        return gallery_admissible({"kind": classify_image(img), "w": w, "h": h,
                                   "detail": detail_score(img)})
    except Exception:
        return False


def is_photo_kind(img) -> bool:
    """True when an image is a real photo/aerial/render (the only valid silent hero)."""
    return classify_image(img) == "photo"


def classify_data_uri(uri: str) -> str:
    """classify_image for a 'data:image/...;base64,...' hero URI (used by the G-images
    gate to BLOCK a map/plan/screenshot hero). Returns 'photo' on any decode/stats
    failure - a gate must never crash or block on an unreadable URI."""
    if Image is None or not (isinstance(uri, str) and "base64," in uri):
        return "photo"
    try:
        import base64
        raw = base64.b64decode(uri.split("base64,", 1)[1])
        return classify_image(Image.open(io.BytesIO(raw)))
    except Exception:
        return "photo"


_DOC_CACHE: dict[str, "fitz.Document"] = {}
# in-process memo of a page's decoded embedded rasters, keyed (id(doc), page_index): the
# same page is decoded by BOTH the interpret thumbnail pass and merge's hero binding -
# decode once per run. MUST be cleared in close_doc_cache (id(doc) is recycled once the
# handle is closed/dropped), which it is. (#39/#21)
_EMBED_CACHE: dict[tuple, list] = {}


def _get_doc(pdf_path: Path) -> "fitz.Document":
    """Open a brochure PDF once and reuse the handle. Image harvesting hits the
    same PDF once per property, so caching avoids reopening/reparsing it N times
    in the merge loop. Call close_doc_cache() when a merge run is done."""
    key = str(Path(pdf_path).resolve())
    doc = _DOC_CACHE.get(key)
    if doc is None:
        doc = fitz.open(pdf_path)
        _DOC_CACHE[key] = doc
    return doc


def close_doc_cache() -> None:
    """Close and forget every cached PDF handle (call at the end of merge), and
    drop the per-document geometry/crop memos that go with them."""
    for doc in _DOC_CACHE.values():
        try:
            doc.close()
        except Exception:
            pass
    _DOC_CACHE.clear()
    try:
        _PLACED_CACHE.clear()
        _CROPS_CACHE.clear()
        _PPTX_CACHE.clear()
        _EMBED_CACHE.clear()          # (#39) id(doc)-keyed decode memo - MUST drop with the handles
        _SLIDEPIC_CACHE.clear()       # (#39) per-slide decoded pictures
        _PPTX_SLIDES_CACHE.clear()    # (#38) enumerated slide lists
    except Exception:
        pass


def _engine_tag() -> str:
    """The ACTIVE image tier, as part of every cache key. Without it the cache was
    engine-blind and a DEGRADED pass poisoned every later pass permanently:

    under the `fitz_shim` tier `get_pixmap` raises, so `page_hero_and_plan` cached the literal
    negative "NONE" for a page whose only image is Flate-encoded (a PNG-sourced photo the shim's
    DCT/JPX-only decoder cannot see). `run._is_current` keys purely on input MTIMES, so the tier
    is not a resume input - a later NATIVE-PyMuPDF run resume-skipped merge and served the
    poisoned negative. The run then printed "native PyMuPDF ... full-fidelity extraction" while
    merge wrote the ledger line "no usable photo in any source (placeholder shown)" about a
    source that demonstrably holds one, and neither placeholder gate could fire (the audit had
    recorded `candidates: 0`, and the rate check needs >=50%). `soffice_pdf` avoided this by
    never caching a negative.

    Only the entries that are actually wrong are invalidated, because a tier change genuinely
    changes what is extractable.

    Three components, and the last two were missing (B17): the module NAME alone collapsed
    fitz_shim's two backends - pdfplumber and pypdfium2 decode different things - into one
    key, and with no VERSION a PyMuPDF upgrade that changes what is extractable looked
    identical to the cache."""
    name = getattr(fitz, "__name__", "fitz")
    backend = getattr(fitz, "_BACKEND", "") or getattr(fitz, "_backend", "")
    if backend:
        name = f"{name}.{backend}"
    ver = (getattr(fitz, "__version__", "") or getattr(fitz, "VersionBind", "") or "0")
    # ...and the GEOMETRY backend, for exactly the same reason as the decode tier: which
    # engine answers "where is this image placed / how wide is this page" decides whether
    # the tier-B hero crop and the whole placed-image plan tier can see anything at all.
    # A cache filled while geometry was unavailable holds negatives that are wrong the
    # moment it becomes available, and mtimes alone would never invalidate them.
    return f"{name}|{'pil' if _HAS_PIL else 'nopil'}|{ver}|geom:{_geom_backend()}"


def _cache_file(pdf_path, page_index, budget_kb, kind, cache_dir, ext=".uri"):
    """Cache path for a (deck, page, budget, kind, ENGINE) unit. ext is `.uri` for the visual
    hero/plan/gallery data URIs and `.json` for intermediate per-page geometry/photo
    caches (so the two never collide in a `*.uri` count and stay self-describing)."""
    if not cache_dir:
        return None
    try:
        import hashlib
        st = Path(pdf_path).stat()
        # v4 (2026-08-20): `page_plan` now rejects a candidate the classifier reads as a PHOTO.
        # A cached v3 plan URI was computed WITHOUT that screen, and one of them is a bright-sky
        # aeroplane photo sitting in a property's Site Plan slot - a stale positive the engine tag
        # cannot invalidate, because the engine did not change, the RULE did. Bumping the prefix
        # is the documented human decision for exactly this (merge's resume note says so).
        key = hashlib.sha1(f"v4|{Path(pdf_path).name}|{st.st_size}|{st.st_mtime_ns}|"
                           f"{page_index}|{budget_kb}|{kind}|{_engine_tag()}".encode()).hexdigest()
        cdir = Path(cache_dir)
        cdir.mkdir(parents=True, exist_ok=True)
        return cdir / f"{key}{ext}"
    except Exception:
        return None  # cache trouble must never break the harvest


def _cache_read(cf):
    """data URI, '' for a cached negative, or None when there is no cache entry.
    A TRUNCATED entry (a kill mid-write before the atomic rename below existed) must
    not be served: base64 from compress() is always %4==0, so a payload failing that
    is incomplete -> treat as a miss and recompute."""
    if cf is None or not cf.exists():
        return None
    try:
        v = cf.read_text(encoding="ascii")
        if v == "NONE":
            return ""
        if v.startswith("data:image/") and "base64," in v:
            b64 = v.split("base64,", 1)[1]
            if len(b64) > 32 and len(b64) % 4 == 0:
                return v
        return None
    except Exception:
        return None


def _cache_write(cf, val):
    if cf is not None:
        try:  # atomic: write a temp then rename, so a kill mid-write can never leave
            # a truncated .uri that the prefix check would happily serve as a hero
            import os
            tmp = cf.with_suffix(cf.suffix + ".tmp")
            tmp.write_text(val if val else "NONE", encoding="ascii")
            os.replace(tmp, cf)
        except Exception:
            pass


def _cache_read_json(cf):
    """A cached JSON value (list/dict), or None when absent / truncated (a kill mid-write
    before the atomic rename) - a parse failure is a miss, so the unit is recomputed."""
    if cf is None or not cf.exists():
        return None
    try:
        return json.loads(cf.read_text(encoding="utf-8"))
    except Exception:
        return None


def _cache_write_json(cf, obj):
    if cf is not None:
        try:  # atomic temp+rename, exactly like _cache_write - a shell-cap kill mid-write
            import os
            tmp = cf.with_suffix(cf.suffix + ".tmp")
            tmp.write_text(json.dumps(obj, ensure_ascii=False), encoding="utf-8")
            os.replace(tmp, cf)
        except Exception:
            pass


def page_hero_and_plan(pdf_path: Path, page_index: int,
                       budget_kb: int = DEFAULT_BUDGET_KB,
                       cache_dir: Path | str | None = None) -> tuple[str | None, str | None]:
    """(hero_uri, plan_uri) for one brochure page - the broker's combination rules:
      * a photo exists           -> hero = the photo; plan slot = the site plan (or None)
      * no photo, a plan exists  -> hero = the PLAN, and the plan slot carries it too
      * neither                  -> hero = legacy page-render comparison (or None); plan None

    Hero ladder (engine-agnostic): A) best decodable EMBEDDED image - a real photo
    (>= MODEST_PHOTO with the flatness multiplier) wins outright, the page raster
    can never beat it; B) the most photographic image REGION cropped from the page
    render via pdfplumber geometry (location-map and boilerplate boxes excluded) -
    covers backends that cannot decode the image streams; C) page render vs
    sub-modest embedded; D) None (placeholder upstream). The site plan comes from
    page_plan(). Both results are memoised on disk per (source, page, budget) -
    re-rastering + the JPEG ladder were the dominant re-run wall-clock cost."""
    if Image is None:
        return (None, None)  # no Pillow: no hero/plan; merge fills the placeholder
    cf_h = _cache_file(pdf_path, page_index, budget_kb, "hero", cache_dir)
    cf_p = _cache_file(pdf_path, page_index, budget_kb, "plan", cache_dir)
    ch, cp = _cache_read(cf_h), _cache_read(cf_p)
    if ch is not None and cp is not None:
        return (ch or None), (cp or None)

    hero = plan = None
    # Tier A: embedded images ranked by KIND TIER first (a real photo/aerial/render leads;
    # a plan/map/text never beats a photo even when it scores higher on colourfulness),
    # then by photographic_score within a tier. A PHOTO-kind embedded image is the hero
    # OUTRIGHT - even a near-monochrome warehouse interior the colourfulness score
    # under-rates. A non-photo best (a slide's map/plan) is NOT made the hero here: it
    # falls through to the crop/plan/render tiers so the page's real photo (or the plan
    # slot) wins, and the G-images gate BLOCKS a non-photo hero that still survives.
    doc = _get_doc(pdf_path)
    cands = []  # (tier, -score, img)
    for im in page_embedded_images(doc, page_index):
        if im["w"] >= MIN_HERO_W and im["h"] >= MIN_HERO_H:
            cands.append((HERO_TIER.get(classify_image(im["img"]), 9),
                          -photographic_score(im["img"]), im["img"]))
    cands.sort(key=lambda c: (c[0], c[1]))
    best = cands[0][2] if cands else None
    best_is_photo = bool(cands) and cands[0][0] == HERO_TIER["photo"]
    if best is not None and best_is_photo:
        hero = to_data_uri(compress(best, HERO_MAX_EDGE, budget_kb))
    # Tier B: photographic region cropped from the page render (stream-decode-proof, and
    # it rescues a photo baked into a flattened slide alongside text/plan/map)
    if hero is None:
        crop = bbox_crop_hero(pdf_path, page_index, cache_dir=cache_dir)
        if crop is not None:
            hero = to_data_uri(compress(crop, HERO_MAX_EDGE, budget_kb))
    # the site plan, independent of the hero
    plan_img = page_plan(pdf_path, page_index, cache_dir=cache_dir)
    if plan_img is not None:
        plan = to_data_uri(compress(plan_img, PLAN_MAX_EDGE, budget_kb))
    # plan-only page: the plan IS the hero (above a map / slide screenshot / placeholder)
    # and stays in the plan slot too
    if hero is None and plan is not None:
        hero = plan
    # Tier C: a TEXT-BEARING page is a SLIDE - its whole-page render is a "screenshot of
    # the slide", NEVER a hero. Only an IMAGE-ONLY page (no text layer) may use its render,
    # and only when the render itself is a real photo (a full-bleed photo page). Otherwise
    # fall to the best embedded as a LAST RESORT (a plan/map - the gate then BLOCKS the
    # non-photo hero for sign-off), else leave hero None -> the honest placeholder.
    if hero is None:
        try:
            page_text = (doc[page_index].get_text() or "").strip()
        except Exception:
            page_text = ""
        if len(page_text) < 200:  # image-only page: its render may BE the property image
            try:
                raster = page_raster(doc, page_index)
            except Exception:  # no renderer in this sandbox tier
                raster = None
            if (raster is not None and raster.width >= MIN_HERO_W
                    and classify_image(raster) == "photo"):
                hero = to_data_uri(compress(raster, HERO_MAX_EDGE, budget_kb))
        if hero is None and best is not None:
            hero = to_data_uri(compress(best, HERO_MAX_EDGE, budget_kb))

    _cache_write(cf_h, hero)
    _cache_write(cf_p, plan)
    return hero, plan


def hero_for_pdf_page(pdf_path: Path, page_index: int,
                      budget_kb: int = DEFAULT_BUDGET_KB,
                      cache_dir: Path | str | None = None) -> str | None:
    """Back-compat wrapper: the hero half of page_hero_and_plan()."""
    return page_hero_and_plan(pdf_path, page_index, budget_kb, cache_dir)[0]


def best_hero_in_deck(path: Path, budget_kb: int = DEFAULT_BUDGET_KB,
                      cache_dir: Path | str | None = None, max_pages: int = 80) -> str | None:
    """Best photographic HERO across a WHOLE deck (PDF pages OR PPTX slides). Used when
    a 0-record brochure has been MATCHED to a known property (from a tracker, email or
    other deck) but carries no page_no, so we scan the deck for its best photo. CHEAP by
    design - the embedded-image / slide-picture tier ONLY (no pdfplumber geometry, no
    page render) - so it never hits the _placed_layout speed cliff. Returns a hero data
    URI or None (no usable photo -> the property keeps its honest placeholder).

    DELEGATES to _deck_photo_index so the result is EXACTLY gallery_for_deck()[0]. The
    photo-match path in merge sets photo=this AND gallery=gallery_for_deck(), and the
    images gate requires gallery[0]==photo: deriving the hero from a SEPARATE ranking
    (the earlier code compared UNROUNDED scores while the index rounds to 3dp + tie-breaks
    by page/sig) could diverge on a near-tie and HARD-BLOCK the gate. Shares the index's
    on-disk cache, so the photo-match hero + gallery are one computation."""
    if Image is None:
        return None  # no Pillow: no hero; the property keeps its honest placeholder
    idx = _deck_photo_index(Path(path), budget_kb, cache_dir, max_pages)
    return idx[0]["uri"] if idx else None


def _photo_sig(img) -> str:
    """Cheap content signature (8x8 greyscale) to dedup the SAME photo repeated across
    pages/slides (a reused hero, a logo) so the gallery never shows a near-duplicate.
    Deterministic; '' on failure (treated as unique)."""
    import hashlib
    try:
        return hashlib.sha1(img.convert("L").resize((8, 8)).tobytes()).hexdigest()
    except Exception:
        return ""


def _deck_photo_index(path: Path, budget_kb: int, cache_dir, max_pages: int = 80) -> list[dict]:
    """Every REAL photo in a deck (>= hero size AND >= MODEST_PHOTO score), one entry per
    distinct image (sig-deduped, highest score kept, lowest page kept), compressed to a
    data URI and tagged with its 0-BASED page/slide, ranked best-first. Cached on disk as
    JSON per source so a re-run/resume is free. The page tag is what lets a MULTI-PROPERTY
    deck contribute only a given property's photos to its gallery, never a neighbour's.
    The MODEST_PHOTO + size floors keep the index to genuine photos (no logos/icons/maps),
    so even a 400-image vector page yields a short list."""
    if Image is None:
        return []
    path = Path(path)
    cf = _cache_file(path, "deck", budget_kb, "galleryidx2", cache_dir)
    if cf is not None and cf.exists():
        try:
            return json.loads(cf.read_text(encoding="utf-8"))
        except Exception:
            pass
    # RESUMABLE: scan page-by-page through a per-page cache (each page's photos are
    # compressed + cached on their own), so a shell-cap kill mid-deck loses at most the
    # page in flight - the next run continues instead of re-scanning the whole deck (the
    # all-or-nothing whole-deck scan was an infinite re-run trap on a deck too big for one
    # ~40s window). Assembly below reproduces the old dedup/ranking exactly.
    try:
        if path.suffix.lower() == ".pptx":
            n_pages = min(len(list(_get_pptx(path).slides)), max_pages)
        else:
            n_pages = min(_get_doc(path).page_count, max_pages)
    except Exception:
        n_pages = 0
    by_sig: dict[str, dict] = {}
    for pno in range(n_pages):
        for e in _deck_page_photos(path, pno, budget_kb, cache_dir):
            prev = by_sig.get(e["sig"])
            if prev is None:
                by_sig[e["sig"]] = dict(e)
            else:  # same image across pages: EARLIEST page, BEST score (its uri)
                prev["page"] = min(prev["page"], e["page"])
                if e["score"] > prev["score"]:
                    prev["score"], prev["uri"] = e["score"], e["uri"]
                    for k in ("w", "h", "detail"):  # the measurements travel WITH the uri
                        if k in e:
                            prev[k] = e[k]
    # rank by KIND TIER first (a real photo/aerial/render leads best_hero_in_deck and the
    # gallery) then by score - so the photo-match hero is a photo, not a map/plan. (kind
    # defaults to 'photo' for an older cache written before the classifier, preserving its
    # score order.) The carousel filters this ranked list by gallery_admissible, so its
    # best-first order IS the carousel's quality order.
    ranked = sorted(by_sig.values(),
                    key=lambda d: (HERO_TIER.get(d.get("kind", "photo"), 9),
                                   -d["score"], d["page"], d["sig"]))
    index = [{"page": d["page"], "score": d["score"], "sig": d["sig"], "uri": d["uri"],
              "kind": d.get("kind", "photo"), "w": d.get("w"), "h": d.get("h"),
              "detail": d.get("detail")} for d in ranked]
    _cache_write_json(cf, index)  # whole-deck index (cheap once the per-page caches exist)
    return index


def _deck_page_photos(path: Path, page_index: int, budget_kb: int, cache_dir) -> list[dict]:
    """One page/slide's qualifying photos as [{page, score, sig, uri, kind, w, h, detail}] -
    the highest-scoring instance per distinct image (sig), >= hero size AND >= MODEST_PHOTO,
    compressed to a data URI. Cached per (deck, page, budget) so _deck_photo_index resumes
    mid-deck.

    ADMISSION IS THE HERO'S, deliberately unchanged: this index feeds best_hero_in_deck (the
    photo-match hero) as well as the carousel, so tightening it here would silently change
    which hero a card gets. The CAROUSEL's stricter floors are applied downstream, in
    gallery_for_pages / gallery_for_deck via gallery_admissible - which is why every entry
    now also carries the measurements that decision needs. `w`/`h` are the RENDERED dimensions
    (post HERO_MAX_EDGE downscale, i.e. what the URI actually contains and the card actually
    shows), not the source raster's."""
    cf = _cache_file(path, page_index, budget_kb, "gidxpage2", cache_dir, ext=".json")
    cached = _cache_read_json(cf)
    if cached is not None:
        return cached
    best: dict[str, dict] = {}  # sig -> {"_sc": float, "_img": Image}
    try:
        if Path(path).suffix.lower() == ".pptx":
            items = slide_pictures(path, page_index)
        else:
            items = page_embedded_images(_get_doc(path), page_index)
        for im in items:
            if im["w"] >= MIN_HERO_W and im["h"] >= MIN_HERO_H:
                kind = classify_image(im["img"])
                sc = photographic_score(im["img"])
                # keep a real PHOTO even when the colourfulness score under-rates it
                # (grey/industrial), and keep a plan/map that clears the score (the Site Plan
                # toggle keeps them - the tier rank just never lets them LEAD)
                if kind != "photo" and sc < MODEST_PHOTO:
                    continue
                sig = _photo_sig(im["img"])
                prev = best.get(sig)
                if prev is None or sc > prev["_sc"]:
                    best[sig] = {"_sc": sc, "_img": im["img"], "_kind": kind}
    except Exception:
        pass
    out: list[dict] = []
    for sig, d in best.items():
        try:
            uri = to_data_uri(compress(d["_img"], HERO_MAX_EDGE, budget_kb))
        except Exception:
            continue
        sw, sh = d["_img"].size
        scale = min(1.0, HERO_MAX_EDGE / max(1, max(sw, sh)))
        out.append({"page": page_index, "score": round(d["_sc"], 3), "sig": sig,
                    "uri": uri, "kind": d["_kind"],
                    "w": max(1, int(sw * scale)), "h": max(1, int(sh * scale)),
                    "detail": round(detail_score(d["_img"]), 3)})
    _cache_write_json(cf, out)
    return out


def _excluded_sigs(path: Path, exclude_by_page: dict) -> set:
    """Map each excluded candidate index (the interpreter's __meta.exclude_refs, keyed by
    0-based page) to its image SIG. The index is the SAME candidates_for_page space heroRef
    uses, and _photo_sig(candidates[idx].img) equals the sig the deck index stored for that
    image (both derive from the memoised page_embedded_images), so the drop is EXACT. Bad
    keys/indices are skipped (a validator bounds them upstream). (exclude_refs)"""
    sigs: set = set()
    for pg, refs in (exclude_by_page or {}).items():
        try:
            page = int(pg)
        except (TypeError, ValueError):
            continue
        try:
            cands = candidates_for_page(Path(path), page)
        except Exception:
            cands = []
        for r in (refs or []):
            if isinstance(r, int) and not isinstance(r, bool) and 0 <= r < len(cands):
                try:
                    sigs.add(_photo_sig(cands[r]["img"]))
                except Exception:
                    pass
    return sigs


def gallery_for_pages(path: Path, page_nos, budget_kb: int = DEFAULT_BUDGET_KB,
                      cache_dir: Path | str | None = None, max_n: int = GALLERY_MAX,
                      exclude_by_page: dict | None = None) -> list[str]:
    """Up to max_n photo data URIs from the given 0-based pages of a deck, best-first.
    PAGE-SCOPED so a multi-property deck contributes only THIS property's photos. Returns
    a 2-tuple (uris, total_available) so the caller can note in the Gaps Report when more
    photos existed than the cap allowed.

    exclude_by_page = {page: [candidate indices]} (the interpreter's __meta.exclude_refs -
    candidates it judged DECORATIVE/non-building via vision). Those candidates are dropped
    from the carousel by SIG match (no cache-key change; empty/None = byte-identical to today).

    QUALITY: filtered by gallery_admissible, so a decorative graphic, a plan/map duplicate of
    the Site Plan slot, and a thumbnail-scale raster never reach a client-facing carousel.
    `total_available` counts the ADMISSIBLE photos on those pages - i.e. what the cap actually
    withheld, not how many rasters the page happened to hold."""
    pages = set(page_nos or [])
    idx = _deck_photo_index(Path(path), budget_kb, cache_dir)
    items = [e for e in idx if e["page"] in pages] if pages else idx
    if exclude_by_page:
        excl = _excluded_sigs(Path(path), exclude_by_page)
        if excl:
            items = [e for e in items if e.get("sig") not in excl]
    items = [e for e in items if gallery_admissible(e)]
    return [e["uri"] for e in items[:max_n]], len(items)


def gallery_for_deck(path: Path, budget_kb: int = DEFAULT_BUDGET_KB,
                     cache_dir: Path | str | None = None, max_n: int = GALLERY_MAX) -> list[str]:
    """Up to max_n best photo data URIs across a WHOLE deck (the photo-match case: a single
    brochure matched to a tracker property IS that property). (uris, total_available).

    Filtered by gallery_admissible, exactly like gallery_for_pages. The caller composes the
    final carousel (the hero leads it) - see merge._compose_gallery: this returns CANDIDATES,
    not a finished gallery, so [0] is not guaranteed to be the deck's hero."""
    idx = _deck_photo_index(Path(path), budget_kb, cache_dir)
    items = [e for e in idx if gallery_admissible(e)]
    return [e["uri"] for e in items[:max_n]], len(items)


def page_image_audit(pdf_path: Path, page_index: int, out_dir: Path, tag: str,
                     cache_dir: Path | str | None = None) -> list[str]:
    """The PLACEHOLDER AUDIT: dump EVERY image candidate on the page - all embedded
    images regardless of size, plus the geometry crops - as labelled thumbnails.
    A placeholder is never a silent default: when the pickers found nothing, a
    human/reviewer must be able to SEE the discard pile and sign off that nothing
    in it was a usable photo or plan (the failure this audits was a real site plan
    filtered out twice - by the size floor, then by the photo scorer - with nobody
    ever shown what was discarded). Returns the written file paths.

    cache_dir is the pickers' persistent image cache and MUST be passed on a merge
    run: on a resumed run the pickers return from disk without touching the geometry
    layer, so the audit's _page_crops call is what pays for _placed_layout - uncached
    (cache_dir=None) that is a full pdfplumber re-parse of the deck, redone and thrown
    away EVERY re-run, which under the ~40s shell-cap kill/re-run cycle is exactly the
    infinite re-run trap the per-page geometry cache exists to prevent."""
    if Image is None:
        return []  # no Pillow: no audit montage (placeholders are surfaced honestly elsewhere)
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    # RESUME SHORT-CIRCUIT: the thumbnail names are deterministic per (tag, page), so once this
    # page's audit is on disk the whole pass - page_raster + every embedded-image decode + the
    # PNG re-writes - is pure waste. Measured 1.2-1.7s per audited page, paid on EVERY resumed
    # merge round; stacked across placeholder properties that alone can exhaust merge's ~10-15s
    # window under the ~40s shell cap. (The cache_dir fix covered only the GEOMETRY layer.)
    existing = sorted(str(p.resolve()) for p in out_dir.glob(f"{tag}_p{page_index + 1}_*.png"))
    if existing:
        return existing
    files: list[str] = []

    def _save(img, kind: str, idx: int):
        try:
            im = img.convert("RGB")
            w, h = im.size
            im.thumbnail((480, 480))
            name = f"{tag}_p{page_index + 1}_{kind}{idx}_{w}x{h}.png"
            _atomic_save_png(im, out_dir / name)
            files.append(str((out_dir / name).resolve()))
        except Exception:
            pass

    try:
        doc = _get_doc(pdf_path)
        for i, im in enumerate(page_embedded_images(doc, page_index), start=1):
            _save(im["img"], "embedded", i)
    except Exception:
        pass
    try:
        for i, c in enumerate(_page_crops(pdf_path, page_index, cache_dir=cache_dir), start=1):
            _save(c["crop"], "region", i)
    except Exception:
        pass
    return files


_MAP_URI = re.compile(r"maps\.google|google\.[a-z.]{2,8}/maps|goo\.gl/maps|openstreetmap\.org|bing\.com/maps", re.I)
_PLACED_CACHE: dict[str, dict] = {}
# a page with more placed images than this is vector-art / a tiled background, not a
# normal brochure page (a real page has a handful of photos/plans/logos). pdfplumber's
# page.images costs 20-35s to build on such a page - OVER the ~45s shell cap. We detect
# it CHEAPLY via fitz get_images (~0.01s/page) and SKIP that page's geometry: Tier A
# (embedded photo) + Tier C (page render) in page_hero_and_plan still yield a hero; only
# the site-plan slot for that one page is sacrificed. Verified on a real deck: a 2,652-
# image page (page.images = 23.5s) where exactly 1 image is >=1.2% of the page.
_PATHOLOGICAL_IMAGES = 400


def _link_near_box(hl: dict, x0, top, x1, bot) -> bool:
    """A link overlapping the box, or sitting as its caption just below/above
    (brochures put the 'click for location' link under the map image)."""
    lx0, lt = float(hl["x0"]), float(hl["top"])
    lx1, lb = float(hl["x1"]), float(hl["bottom"])
    if min(x1, lx1) - max(x0, lx0) <= 0:
        return False  # no horizontal overlap
    if min(bot, lb) - max(top, lt) > 0:
        return True  # overlaps the image itself
    return min(abs(lt - bot), abs(top - lb)) <= 60  # caption proximity


def _placed_cache_file(pdf_path, cache_dir):
    """Disk path for the per-document geometry cache, keyed on bytes+mtime+ENGINE (P0-5).
    The engine belongs in the key because the pathological-page SKIP set is derived from
    `fitz.get_images`, so a shim tier can produce a different geometry set for the same deck -
    see `_engine_tag` for the poisoning this prevents."""
    if not cache_dir:
        return None
    try:
        import hashlib
        st = Path(pdf_path).stat()
        h = hashlib.sha1(f"placed-v2|{Path(pdf_path).name}|{st.st_size}|{st.st_mtime_ns}"
                         f"|{_engine_tag()}".encode()).hexdigest()
        cdir = Path(cache_dir)
        cdir.mkdir(parents=True, exist_ok=True)
        return cdir / f"{h}.placed.json"
    except Exception:
        return None


def _extract_geom_for_page(page) -> list[dict]:
    """RAW image-geometry entries for ONE pdfplumber page: [{bbox, key, frac, aspect, map}].
    `key` is a STABLE per-image-object identity STRING (same PDF image object reused across
    pages -> same key) for cross-page boilerplate detection - never the placed box, because
    templated brochures put photo/plan/map at identical positions on every page. The boiler
    bool is resolved later (it needs the whole deck). Pure per-page, so it parallelises."""
    import hashlib
    entries: list[dict] = []
    map_links = [hl for hl in page.hyperlinks if _MAP_URI.search(str(hl.get("uri", "")))]
    pw, ph = float(page.width), float(page.height)
    for im in page.images:
        x0, top = float(im["x0"]), float(im["top"])
        x1, bot = float(im["x1"]), float(im["bottom"])
        w, h = x1 - x0, bot - top
        if w <= 4 or h <= 4:
            continue
        stream = im.get("stream")
        oid = getattr(stream, "objid", None)
        if oid is not None:
            ks = f"o:{oid}"
        else:
            try:
                raw = stream.rawdata or b""
                ks = f"h:{len(raw)}:{hashlib.sha1(raw[:64]).hexdigest()[:16]}"
            except Exception:
                ks = f"b:{round(x0)},{round(top)},{round(x1)},{round(bot)}"
        entries.append({"bbox": [x0, top, x1, bot], "key": ks,
                        "frac": (w * h) / (pw * ph), "aspect": w / h,
                        "map": any(_link_near_box(hl, x0, top, x1, bot) for hl in map_links)})
    return entries


# --- MEDIA GEOMETRY: PyMuPDF is the PRIMARY backend, pdfplumber the FALLBACK ------- #
# This layer used to be written pdfplumber-FIRST even though PyMuPDF is the skill's primary
# engine, so on a host with the BETTER engine and no fallback installed the whole layer went
# silently dead: `_placed_layout` returned 0 boxes on every page of every deck (killing the
# tier-B hero crop AND the entire placed-image site-plan tier), and `_page_crops` opened
# pdfplumber solely to read the page WIDTH IN POINTS and `return []` when it could not - one
# number, gating two media tiers. PyMuPDF answers all three questions natively (`page.rect`,
# `page.get_image_info()` / `get_image_bbox()`, `page.get_links()`), so the dependency
# direction is now the right way round. pdfplumber is still fully wired for a genuine
# no-PyMuPDF sandbox - the point is not to DEPEND on the fallback when the primary is present.
_GEOM_BACKEND: str | None = None


def _geom_backend() -> str:
    """'fitz' when the ACTIVE engine serves image geometry natively, else 'plumber'.
    Probed from the engine's own Page surface (the `fitz_shim` Page has `get_links` but
    no `get_image_info`/`rect`, so a real sandbox shim still routes to pdfplumber).
    Memoised - it cannot change within a process."""
    global _GEOM_BACKEND
    if _GEOM_BACKEND is None:
        try:
            pg = getattr(fitz, "Page", None)
            _GEOM_BACKEND = ("fitz" if (pg is not None
                                        and hasattr(pg, "get_image_info")
                                        and hasattr(pg, "get_links")
                                        and hasattr(pg, "rect")) else "plumber")
        except Exception:
            _GEOM_BACKEND = "plumber"
    return _GEOM_BACKEND


_MEDIA_CAPS: dict | None = None


def media_capabilities() -> dict:
    """What this host can ACTUALLY do to media, probed - never guessed.

    Every media tier in this module degrades to an honest `None`/`[]` when its capability is
    absent, which is correct behaviour and exactly why the absence is invisible: a run whose
    image layer is dead looks, in every artefact it writes, like a run whose sources hold no
    images. That indistinguishability is the defect this function exists to close - it is read
    by `gate_runner media-harvest` and printed by `run.py` beside the `PDF engine:` line, so a
    lost CAPABILITY is stated as a fact instead of being inferred from a thin gallery.

    Each key is a HARD probe against a synthesised one-page document, not a version sniff or an
    `hasattr` guess (the `fitz_shim` Page carries several of these names and raises when called):
      pillow    - PIL is importable, so a raster can be decoded/cropped/compressed at all
      renderer  - the engine can rasterise a page (hero tier-C, every page render, plan renders)
      geometry  - placed-image boxes + page width are answerable (hero tier-B, plan crop tier)
      drawings  - `page.get_drawings()` works, i.e. the VECTOR site-plan route can fire
      text      - a page's text layer is readable (plan titles, furniture, the spec gate)
    plus `engine` (the cache-key engine tag) and `geometry_backend` ('fitz'/'plumber').

    Memoised; never raises."""
    global _MEDIA_CAPS
    if _MEDIA_CAPS is not None:
        return dict(_MEDIA_CAPS)
    caps = {"pillow": Image is not None, "renderer": False, "geometry": False,
            "drawings": False, "text": False}
    doc = None
    try:
        doc = fitz.open()                       # a new, empty in-memory document
        page = doc.new_page()
        try:
            page.draw_line((10, 10), (100, 100))
        except Exception:
            pass
        try:
            caps["renderer"] = page.get_pixmap(dpi=18) is not None
        except Exception:
            pass
        try:
            page.get_image_info(xrefs=True)
            page.get_links()
            float(page.rect.width)
            caps["geometry"] = True
        except Exception:
            pass
        try:
            caps["drawings"] = isinstance(page.get_drawings(), (list, tuple))
        except Exception:
            pass
        try:
            page.get_text()
            caps["text"] = True
        except Exception:
            pass
    except Exception:
        pass
    finally:
        try:
            if doc is not None:
                doc.close()
        except Exception:
            pass
    if not caps["geometry"]:
        # a host whose engine cannot answer geometry natively may still answer via the
        # pdfplumber FALLBACK - the capability is present, just on the slower backend
        try:
            import pdfplumber  # noqa: F401
            caps["geometry"] = True
        except Exception:
            pass
    caps["engine"] = _engine_tag()
    caps["geometry_backend"] = _geom_backend()
    _MEDIA_CAPS = caps
    return dict(caps)


# The capabilities whose ABSENCE materially costs media (the ones worth shouting about).
MEDIA_CRITICAL_CAPS = ("pillow", "renderer", "geometry", "drawings", "text")


def deck_media_facts(path: Path) -> dict:
    """What a deck HOLDS, cheaply: `{pages, large_images}` - the page count and how many distinct
    embedded rasters clear the hero size floor (MIN_HERO_W x MIN_HERO_H). METADATA ONLY: image
    dimensions are read off the xref table, nothing is decoded, so this is a millisecond-scale
    read even on a 40-page deck - it is a SIGNAL input, not a harvest.

    Its purpose is the arithmetic no gate could previously do: compare what a source HOLDS with
    what a property TOOK from it. `{}` on any failure (an honest absence - the caller then simply
    has no signal for that deck, never a false one)."""
    try:
        path = Path(path)
        if path.suffix.lower() == ".pptx":
            slides = _get_pptx(path).slides
            n = len(list(slides))
            big = 0
            for i in range(n):
                for im in slide_pictures(path, i):
                    if im.get("w", 0) >= MIN_HERO_W and im.get("h", 0) >= MIN_HERO_H:
                        big += 1
            return {"pages": n, "large_images": big}
        doc = _get_doc(path)
        n = doc.page_count
        seen: set = set()
        for p in range(n):
            try:
                for info in (doc.get_page_images(p, full=True) or ()):
                    # (xref, smask, width, height, bpc, colorspace, ...) - PyMuPDF's own order
                    xref, w, h = info[0], info[2], info[3]
                    if int(w) >= MIN_HERO_W and int(h) >= MIN_HERO_H:
                        seen.add(int(xref))   # by xref: a boilerplate banner counts ONCE
            except Exception:
                continue
        return {"pages": n, "large_images": len(seen)}
    except Exception:
        return {}


def _fitz_map_links(page) -> list[dict]:
    """The page's maps-service hyperlinks, emitted in pdfplumber's HYPERLINK SHAPE
    ({uri, x0, top, x1, bottom}, page-relative top-left coords) so `_link_near_box` is
    reused byte-for-byte by both backends and the location-map exclusion cannot diverge."""
    out: list[dict] = []
    try:
        r = page.rect
        ox, oy = float(r.x0), float(r.y0)
        for lk in (page.get_links() or []):
            uri = str(lk.get("uri") or "")
            if not _MAP_URI.search(uri):
                continue
            fr = lk.get("from")
            if fr is None:
                continue
            out.append({"uri": uri,
                        "x0": float(fr.x0) - ox, "top": float(fr.y0) - oy,
                        "x1": float(fr.x1) - ox, "bottom": float(fr.y1) - oy})
    except Exception:
        return []
    return out


def _fitz_placed_boxes(page) -> list[tuple]:
    """[(xref, (x0, top, x1, bottom))] for the page's PLACED image boxes, in page-relative
    top-left points. `get_image_info(xrefs=True)` yields one entry per PLACEMENT (the same
    image placed twice gives two boxes - exactly pdfplumber's `page.images` semantics);
    `get_images` + `get_image_bbox` is the fallback for an engine build without it."""
    r = page.rect
    ox, oy = float(r.x0), float(r.y0)
    out: list[tuple] = []
    infos = None
    try:
        infos = page.get_image_info(xrefs=True)
    except Exception:
        infos = None
    if infos:
        for it in infos:
            bb = it.get("bbox")
            if not bb:
                continue
            try:
                out.append((int(it.get("xref") or 0),
                            (float(bb[0]) - ox, float(bb[1]) - oy,
                             float(bb[2]) - ox, float(bb[3]) - oy)))
            except (TypeError, ValueError):
                continue
        return out
    for im in (page.get_images(full=True) or []):
        try:
            bb = page.get_image_bbox(im)
            out.append((int(im[0]), (float(bb.x0) - ox, float(bb.y0) - oy,
                                     float(bb.x1) - ox, float(bb.y1) - oy)))
        except Exception:
            continue
    return out


# A page with more placed boxes than this is vector art / a tiled background rather than a
# normal brochure page. The pdfplumber backend SKIPS such a page entirely (its `page.images`
# access costs 20-35s, over the shell cap) and sacrifices that page's plan slot. The fitz
# backend has no such cost (`get_image_info` is ~ms), so it keeps the page and merely drops
# the icon-scale boxes - the same floor `_page_crops` applies anyway - which bounds the cache
# without sacrificing the site-plan slot on exactly the vector-heavy pages a masterplan lives on.
_GEOM_DENSE_FLOOR = 0.012


def _extract_geom_for_page_fitz(page) -> list[dict] | None:
    """RAW image-geometry entries for ONE PyMuPDF page, under the SAME contract as
    `_extract_geom_for_page` ([{bbox, key, frac, aspect, map}]) so every consumer is
    untouched. `key` is the PDF image OBJECT number (xref) - the same cross-page-stable
    identity pdfplumber's `stream.objid` gave, so the >=3-page boilerplate detector behaves
    identically. None when this engine cannot answer (caller falls back to pdfplumber)."""
    try:
        r = page.rect
        pw, ph = float(r.width), float(r.height)
        if pw <= 0 or ph <= 0:
            return []
        links = _fitz_map_links(page)
        boxes = _fitz_placed_boxes(page)
        dense = len(boxes) > _PATHOLOGICAL_IMAGES
        entries: list[dict] = []
        for xref, (x0, top, x1, bot) in boxes:
            # CLIP to the page. A placement rectangle is reported UNCLIPPED, so a masked or
            # bleed-off tile legitimately extends past the page (measured: boxes from -587
            # to +1938 pt on an 1190 pt page). Unclipped they overstate `frac`, and a
            # negative edge crashed the crop downstream. Clipping is also the RIGHT
            # semantic here: every consumer asks "how much of the PAGE does this cover".
            x0, x1 = max(0.0, min(x0, x1)), min(pw, max(x0, x1))
            top, bot = max(0.0, min(top, bot)), min(ph, max(top, bot))
            w, h = x1 - x0, bot - top
            if w <= 4 or h <= 4:
                continue
            frac = (w * h) / (pw * ph)
            if dense and frac < _GEOM_DENSE_FLOOR:
                continue
            ks = (f"o:{xref}" if xref
                  else f"b:{round(x0)},{round(top)},{round(x1)},{round(bot)}")
            entries.append({"bbox": [x0, top, x1, bot], "key": ks,
                            "frac": frac, "aspect": w / h,
                            "map": any(_link_near_box(hl, x0, top, x1, bot) for hl in links)})
        return entries
    except Exception:
        return None


def _page_width_pts(pdf_path: Path, page_index: int) -> float | None:
    """The page's width in POINTS - `page.rect.width` natively, pdfplumber only when the
    active engine cannot answer. This single number is the only reason `_page_crops` ever
    opened pdfplumber, and returning [] when it could not is what disabled the tier-B hero
    crop and the placed-image plan tier on every PyMuPDF-only host."""
    try:
        doc = _get_doc(pdf_path)
        if 0 <= page_index < doc.page_count:
            w = float(doc[page_index].rect.width)
            if w > 0:
                return w
    except Exception:
        pass
    try:
        import pdfplumber
        with pdfplumber.open(str(pdf_path)) as pl:
            if page_index < len(pl.pages):
                w = float(pl.pages[page_index].width)
                return w if w > 0 else None
    except Exception:
        pass
    return None


def _page_pathological(pdf_path: Path, page_index: int) -> bool:
    """True for a vector-art/tiled page (thousands of placed images) - its 20-35s
    pdfplumber page.images access is SKIPPED (Tier A/C still give a hero). O(1) fitz check.
    Only consulted on the PDFPLUMBER geometry backend: the fitz backend pays no such cost
    and keeps the page (see `_extract_geom_for_page_fitz`)."""
    try:
        doc = _get_doc(pdf_path)
        return (page_index < doc.page_count
                and len(doc[page_index].get_images(full=True)) > _PATHOLOGICAL_IMAGES)
    except Exception:
        return False


def _placed_page(pdf_path: Path, page_index: int, cache_dir: Path | str | None = None) -> list[dict]:
    """One page's RAW geometry (boiler not yet resolved), cached per (deck, page) so the
    per-document layout RESUMES mid-deck after a shell-cap kill and PARALLELISES across
    pages in the pre-warm. A pathological page caches []."""
    cf = _cache_file(pdf_path, page_index, 0, "placedpage", cache_dir, ext=".json")
    cached = _cache_read_json(cf)
    if cached is not None:
        return cached
    entries: list[dict] = []
    try:
        if _geom_backend() == "fitz":
            doc = _get_doc(pdf_path)
            if 0 <= page_index < doc.page_count:
                entries = _extract_geom_for_page_fitz(doc[page_index])
            if entries is None:                 # engine could not answer -> pdfplumber
                entries = []
                if not _page_pathological(pdf_path, page_index):
                    import pdfplumber
                    with pdfplumber.open(str(pdf_path)) as pl:
                        if page_index < len(pl.pages):
                            entries = _extract_geom_for_page(pl.pages[page_index])
        elif not _page_pathological(pdf_path, page_index):
            import pdfplumber
            with pdfplumber.open(str(pdf_path)) as pl:
                if page_index < len(pl.pages):
                    entries = _extract_geom_for_page(pl.pages[page_index])
    except Exception:
        entries = []
    _cache_write_json(cf, entries)
    return entries


def _placed_layout(pdf_path: Path, cache_dir: Path | str | None = None) -> dict:
    """Per-document image GEOMETRY via pdfplumber (engine-agnostic): for every page the
    placed image boxes, each flagged when a maps-service hyperlink sits on/under it
    (= the location map) or when the same image object repeats on >=3 pages (= boilerplate).

    RESUMABLE: each page's raw geometry is cached individually (_placed_page), then the deck
    is assembled by resolving the boiler flag across the per-page caches. The whole-deck
    layout is checkpointed ONLY when every page is present, so a shell-cap kill mid-deck
    loses at most the page in flight and the next run continues (the old all-or-nothing
    whole-deck parse was an infinite re-run trap on a deck too big for one ~40s window).
    The geometry is intermediate bbox data, so determinism/chrome are unaffected."""
    key = str(Path(pdf_path).resolve())
    if key in _PLACED_CACHE:
        return _PLACED_CACHE[key]
    cf = _placed_cache_file(pdf_path, cache_dir)
    if cf is not None and cf.exists():
        try:
            disk = json.loads(cf.read_text(encoding="utf-8"))
            if isinstance(disk, dict) and "pages" in disk:
                _PLACED_CACHE[key] = disk
                return disk
        except Exception:
            pass
    # page count (cheap, fitz) + the per-page pathological set, which only the PDFPLUMBER
    # backend needs (it is a cost dodge for `page.images`, not a correctness rule)
    n = 0
    skip_pages: set = set()
    _plumber_geom = _geom_backend() != "fitz"
    try:
        doc = _get_doc(pdf_path)
        n = doc.page_count
        if _plumber_geom:
            for p in range(n):
                if len(doc[p].get_images(full=True)) > _PATHOLOGICAL_IMAGES:
                    skip_pages.add(p)
    except Exception:
        pass
    # read whatever per-page caches already exist; compute the rest in ONE deck-wide pass
    raw: list = [_cache_read_json(_cache_file(pdf_path, p, 0, "placedpage", cache_dir, ext=".json"))
                 for p in range(n)]
    missing = [p for p in range(n) if raw[p] is None]
    if missing and _geom_backend() == "fitz":
        # NATIVE: the doc is already open in the shared handle cache, so there is no
        # deck-wide open to amortise and no pathological-page cost to dodge.
        try:
            doc = _get_doc(pdf_path)
            for p in list(missing):
                entries = (_extract_geom_for_page_fitz(doc[p]) if p < doc.page_count else [])
                if entries is None:
                    continue                    # leave it missing -> the plumber pass below
                raw[p] = entries
                _cache_write_json(_cache_file(pdf_path, p, 0, "placedpage", cache_dir,
                                              ext=".json"), entries)
            missing = [p for p in range(n) if raw[p] is None]
        except Exception:
            pass
    if missing:
        try:
            import pdfplumber
            with pdfplumber.open(str(pdf_path)) as pl:
                for p in missing:
                    entries = ([] if (p in skip_pages or p >= len(pl.pages))
                               else _extract_geom_for_page(pl.pages[p]))
                    raw[p] = entries
                    _cache_write_json(_cache_file(pdf_path, p, 0, "placedpage", cache_dir, ext=".json"), entries)
        except Exception:
            pass
    # assemble: boiler = an image object placed on >=3 pages (deck >=4 pages), from the
    # pages we have - identical to the old single-scan result once every page is present
    content_pages: dict = {}
    for p in range(n):
        for e in (raw[p] or []):
            content_pages.setdefault(e["key"], set()).add(p)
    boiler = ({k for k, ps in content_pages.items() if len(ps) >= 3} if n >= 4 else set())
    out = {"pages": []}
    for p in range(n):
        out["pages"].append([{"bbox": e["bbox"], "frac": e["frac"], "aspect": e["aspect"],
                              "map": e["map"], "boiler": e["key"] in boiler}
                             for e in (raw[p] or [])])
    _PLACED_CACHE[key] = out
    complete = n > 0 and all(raw[p] is not None for p in range(n))
    if cf is not None and complete and out["pages"]:
        _cache_write_json(cf, out)  # whole-deck checkpoint ONLY when every page is in
    return out


def _unit_cached(spec) -> bool:
    """True when a pre-warm work unit's atomic cache already exists (so it is skipped)."""
    kind, path_str, page, budget, cache_str = spec
    if kind == "placedpage":
        cf = _cache_file(path_str, page, 0, "placedpage", cache_str, ext=".json")
    elif kind == "gidxpage":
        # the WORK-UNIT name is stable; the cache TAG carries the entry-shape version, so a
        # pre-warm never reports a v1 cache as covering the v2 unit it would now compute
        cf = _cache_file(path_str, page, budget, "gidxpage2", cache_str, ext=".json")
    else:  # 'hero' (PDF) / 'slidehero' (PPTX) -> the hero .uri is the primary artefact
        cf = _cache_file(path_str, page, budget,
                         "slide_hero" if kind == "slidehero" else "hero", cache_str)
    return cf is not None and cf.exists()


def _prewarm_unit(spec):
    """ProcessPool worker: compute + CACHE one image unit; the on-disk atomic cache IS the
    result (the return is just ok/err). Top-level + stdlib-importable so it pickles to a
    child on both fork (Linux/Cowork) and spawn (Windows). Never raises."""
    try:
        kind, path_str, page, budget, cache_str = spec
        path = Path(path_str)
        if kind == "placedpage":
            _placed_page(path, page, cache_str)
        elif kind == "gidxpage":
            _deck_page_photos(path, page, budget, cache_str)
        elif kind == "slidehero":
            slide_hero_and_plan(path, page, budget, cache_dir=cache_str)
        else:  # 'hero'
            page_hero_and_plan(path, page, budget, cache_dir=cache_str)
        return True
    except Exception:
        return False


def _crop_stats(crop) -> tuple[float, float]:
    """(white_frac, balance) of a crop - the plan signature is a BALANCED mix of
    white paper and drawn ink (balance peaks at white_frac 0.5); photos and map
    tiles sit at the extremes."""
    im = crop.convert("RGB")
    im.thumbnail((96, 96))
    px = list(im.getdata())
    n = len(px) or 1
    white = sum(1 for r, g, b in px if r >= 235 and g >= 235 and b >= 235) / n
    return white, 4.0 * white * (1.0 - white)


_CROPS_CACHE: dict[tuple, list] = {}


def _page_crops(pdf_path: Path, page_index: int, dpi: int = 150,
                cache_dir: Path | str | None = None) -> list[dict]:
    """The page's content-image regions cropped from the render: each candidate
    placed box (>=4%% of the page, sane aspect, not boilerplate) with its crop,
    photographic score and plan stats, plus the map flag. [] when no renderer or
    no geometry - callers fall through their ladders. Memoised per (path, page)
    since the hero tier-B and the plan picker both consume it."""
    memo_key = (str(Path(pdf_path).resolve()), page_index, dpi)
    if memo_key in _CROPS_CACHE:
        return _CROPS_CACHE[memo_key]
    _CROPS_CACHE[memo_key] = []  # set early so a failure is not recomputed per caller
    layout = _placed_layout(pdf_path, cache_dir)
    if page_index >= len(layout["pages"]):
        return []
    # fraction floor only screens out icon-scale boxes; the ABSOLUTE pixel floors
    # (MIN_PLAN / MIN_HERO) do the real gating downstream. 0.04 was high enough
    # to discard genuinely small site plans before they were ever scored.
    boxes = [b for b in layout["pages"][page_index]
             if b["frac"] >= 0.012 and 0.45 <= b["aspect"] <= 3.0
             and not b.get("boiler")]
    if not boxes:
        return []
    try:
        raster = page_raster(_get_doc(pdf_path), page_index, dpi=dpi)
    except Exception:
        return []  # renderer-less sandbox tier
    # raster width / page width in points -> px per point. NATIVE (`page.rect.width`) with a
    # pdfplumber fallback; this one number used to be a hard pdfplumber dependency and its
    # failure silently disabled BOTH the tier-B hero crop and the placed-image plan tier.
    pw = _page_width_pts(pdf_path, page_index)
    if not pw:
        return []
    scale = raster.width / pw
    out = []
    for b in boxes:
        x0, top, x1, bot = b["bbox"]
        # CLAMP into the raster, in that order (clamp the low edge first, then hold the high
        # edge at or above it). A placement box can legitimately sit partly off-page, and the
        # old expression turned that into `lower < upper` -> a hard ValueError that killed the
        # whole page's crops. A box already inside the raster is unaffected.
        cx0 = min(max(0, int(x0 * scale)), raster.width)
        cy0 = min(max(0, int(top * scale)), raster.height)
        cx1 = max(cx0, min(raster.width, int(x1 * scale)))
        cy1 = max(cy0, min(raster.height, int(bot * scale)))
        crop = raster.crop((cx0, cy0, cx1, cy1))
        # the PLAN floor, not the photo floor: small site plans must reach the
        # scorers (photo candidacy re-applies MIN_HERO downstream)
        if crop.width < MIN_PLAN_W or crop.height < MIN_PLAN_H:
            continue
        white, balance = _crop_stats(crop)
        out.append({"crop": crop, "map": b["map"], "score": photographic_score(crop),
                    "white": white, "balance": balance,
                    "rank": balance * math.sqrt(crop.width * crop.height)})
    _CROPS_CACHE[memo_key] = out
    return out


def bbox_crop_hero(pdf_path: Path, page_index: int, dpi: int = 150,
                   cache_dir: Path | str | None = None):
    """Tier-B hero: the most photographic embedded-image REGION cropped out of
    the page raster (geometry via pdfplumber - works even when the backend cannot
    DECODE the image streams, the failure that shipped whole cluttered pages as
    heroes). Location-map and boilerplate boxes are excluded. Returns a PIL image
    scoring >= MODEST_PHOTO, or None."""
    if Image is None:
        return None  # no Pillow: no geometry crops
    cands = [c for c in _page_crops(pdf_path, page_index, dpi, cache_dir)
             if not c["map"] and c["crop"].width >= MIN_HERO_W
             and c["crop"].height >= MIN_HERO_H]  # photo candidacy keeps the photo floor
    best = max(cands, key=lambda c: c["score"], default=None)
    return best["crop"] if best and best["score"] >= MODEST_PHOTO else None


def page_plan(pdf_path: Path, page_index: int, dpi: int = 150,
              cache_dir: Path | str | None = None):
    """The page's SITE PLAN, or None. A plan is a content image region that is
    not the page's PHOTO (only the box that actually wins the hero is excluded -
    a COLOURFUL site plan scores photo-ish ~13 and must stay a plan candidate),
    not the hyperlinked location map, not boilerplate, with the plan signature: a
    balanced mix of white paper and drawn ink (white_frac 0.15-0.90; balance x
    area ranks the best). Calibrated on a real Spanish deck where the colourful
    site plan, the grey line plan, the Google location map and the photos all had
    to separate correctly."""
    if Image is None:
        return None  # no Pillow: no geometry crops
    cands = [c for c in _page_crops(pdf_path, page_index, dpi, cache_dir) if not c["map"]]
    if not cands:
        return None
    hero = max(cands, key=lambda c: c["score"])
    if hero["score"] < MODEST_PHOTO:
        hero = None  # no photo on this page - every region stays a plan candidate
    # ...and it must not be a PHOTOGRAPH. "Not the page's hero" was the only photo screen here,
    # which is no screen at all on a page carrying SEVERAL photos: the runner-up photo then sits
    # in the plan slot. It bound a bright-sky aeroplane photo (a park landmark) as a property's
    # site plan the moment the placed-image geometry backend started answering again, because a
    # sky puts the white fraction squarely inside the plan band. `classify_image` is the right
    # test and the one `_page_has_dominant_photo` already uses for the same question: it judges
    # TONE (continuous-tone vs drawn ink), so a COLOURFUL site plan - which the colourfulness
    # score over-rates at photo-ish ~13, the case this band was calibrated on - still classifies
    # 'plan' and still qualifies. Precision: a wrong image in the trace-less Site Plan slot is
    # worse than an honest gap.
    plans = [c for c in cands
             if c is not hero and 0.15 <= c["white"] <= 0.90
             and classify_image(c["crop"]) != "photo"]
    best = max(plans, key=lambda c: c["rank"], default=None)
    return best["crop"] if best else None


# --- WHOLE-PAGE RENDERED SITE PLAN (vector line-art a placed-image crop cannot reach) -- #
# A site plan is often VECTOR graphics drawn straight into the page, not a placed raster.
# Pulled as an embedded image it rasterises to solid black, so page_plan()/planRef (which
# crop placed-image boxes only) find nothing. page_raster() renders the vector content
# correctly, so we RENDER the page and crop to its ink bbox. The LLM names the plan page
# (__meta.plan_page) and a deterministic render+classify detector is the universal
# fallback + verifier. The render binds the PLAN SLOT ONLY - a vector plan is NEVER made
# the card hero (it would trip the G-images gate and change the card look; intentional).
PLAN_RENDER_DPI = 150


def _rendered_plan_crop(path: Path, page_index: int, dpi: int = PLAN_RENDER_DPI,
                        cache: Path | str | None = None) -> tuple:
    """Render the whole page and crop to its INK bounding box (Pillow getbbox on a
    thresholded non-white copy, with a small margin), then classify the crop. Returns
    (crop, signals, kind) - (None, {}, None) when Pillow / the renderer is unavailable or
    the page is effectively blank. Pure helper for page_render_plan + best_plan_page_render;
    cache is accepted for signature parity (the rendered URI is cached by the callers)."""
    if Image is None:
        return (None, {}, None)
    try:
        path = Path(path)
        if path.suffix.lower() == ".pptx":
            pdf = soffice_pdf(path, cache)
            if pdf is None:
                return (None, {}, None)
            doc = _get_doc(pdf)
        else:
            doc = _get_doc(path)
        if not (0 <= page_index < doc.page_count):
            return (None, {}, None)
        raster = page_raster(doc, page_index, dpi=dpi)
    except Exception:
        return (None, {}, None)  # renderer-less sandbox tier / open failure -> honest None
    if raster is None:
        return (None, {}, None)
    try:
        rgb = raster.convert("RGB")
        # ink mask: anything that is NOT near-white. point() on the greyscale gives a
        # crisp foreground; getbbox() returns the tight box of the drawn content so the
        # paper margins (which read as 'text'/blank) do not dominate the classification.
        grey = rgb.convert("L")
        ink = grey.point(lambda p: 0 if p >= 238 else 255)
        bbox = ink.getbbox()
        if bbox is None:
            return (None, {}, None)  # all white = a blank page, never a plan
        x0, y0, x1, y1 = bbox
        mw = max(8, int((x1 - x0) * 0.03))
        mh = max(8, int((y1 - y0) * 0.03))
        x0 = max(0, x0 - mw); y0 = max(0, y0 - mh)
        x1 = min(rgb.width, x1 + mw); y1 = min(rgb.height, y1 + mh)
        crop = rgb.crop((x0, y0, x1, y1))
        if crop.width < 2 or crop.height < 2:
            return (None, {}, None)
        sig = _hero_signals(crop)
        kind = classify_image(crop, sig)
        return (crop, sig, kind)
    except Exception:
        return (None, {}, None)


PLAN_TITLE_MIN = 1.0  # plan_signal.plan_title_score at/above which a real plan TITLE rescues a page

# --- VECTOR SITE-PLAN SIGNAL (page.get_drawings) ----------------------------------------- #
# The placed-image geometry above can only see PLACED RASTER boxes. Most masterplans are VECTOR
# line art drawn straight into the page, so they have NO placed box at all: the geometry tier
# yields nothing for them even with a working backend, and the whole-page render+classify tier
# mis-reads the designed ones (a full-bleed colour-background plan classifies 'photo'; a plan on
# grey or dark paper falls outside the white-balance band). `page.get_drawings()` exposes the
# vector art DIRECTLY, which is the only route by which such a page is discoverable at all.
#
# THE PRECISION PROBLEM, and how it is answered. Vector density alone is NOT a plan signal: a
# regional road map, a drive-time map and a locator map are all just as vector-dense (measured:
# 2,900 / 2,600 / 6,600 drawing objects on real ones). What a masterplan has and a map does not is
# LABELLED DRAWING FURNITURE - dock doors, yard depth, parking bays, a gatehouse, access gates, a
# site boundary, an accommodation schedule. So the vector route requires BOTH a full-page vector
# body AND `plan_signal.plan_furniture_score` corroboration, and it stays ranked BELOW the two
# existing routes so a real white-paper drawing still wins its own cluster.
VEC_MIN_ITEMS = 800     # primitive path items. A site layout is hundreds of segments for the
#                         building outline alone, then parking bays, dock doors and a boundary;
#                         below this a page's vector content is rules, icons and a logo.
VEC_MIN_CELLS = 0.30    # fraction of a 24x24 page grid the drawing touches - it must be the
#                         page's SUBJECT, not a chart in a corner.
VEC_MIN_SPAN = 0.35     # union-bbox area / page area - same requirement, stated as extent.
_VEC_GRID = 24


def page_vector_art(path: Path, page_index: int, cache_dir: Path | str | None = None) -> dict:
    """The page's VECTOR line-art signal: `{items, cells, span}`.
      items - count of primitive path items (lines/curves/rects) actually drawn on the page;
      cells - fraction of a 24x24 page grid touched by a drawing's bbox (is it the page's subject);
      span  - union bbox of all drawings / page area (the same question as extent).
    `{}` when the engine cannot answer - `get_drawings` is PyMuPDF-only, so on the sandbox shim
    this is empty and the vector route can never fire (today's behaviour, exactly). Cached per
    (deck, page) so a shell-cap kill resumes and a re-run is byte-deterministic. Never raises."""
    cf = _cache_file(path, page_index, 0, "vecart", cache_dir, ext=".json")
    cached = _cache_read_json(cf)
    if isinstance(cached, dict):
        return cached
    out: dict = {}
    try:
        path = Path(path)
        if path.suffix.lower() == ".pptx":
            pdf = soffice_pdf(path, cache_dir)
            if pdf is None:
                return {}
            doc = _get_doc(pdf)
        else:
            doc = _get_doc(path)
        if not (0 <= page_index < doc.page_count):
            return {}
        page = doc[page_index]
        pr = page.rect
        parea = float(pr.width) * float(pr.height)
        if parea <= 0:
            return {}
        items = 0
        cells: set = set()
        ux0 = uy0 = float("inf")
        ux1 = uy1 = float("-inf")
        for dr in (page.get_drawings() or []):
            r = dr.get("rect")
            if r is None:
                continue
            r = r & pr                      # clip to the page: bleed is not page coverage
            if r.is_empty:
                continue
            items += len(dr.get("items") or ())
            ux0 = min(ux0, r.x0); uy0 = min(uy0, r.y0)
            ux1 = max(ux1, r.x1); uy1 = max(uy1, r.y1)
            cx0 = int((r.x0 - pr.x0) / pr.width * _VEC_GRID)
            cx1 = int((r.x1 - pr.x0) / pr.width * _VEC_GRID)
            cy0 = int((r.y0 - pr.y0) / pr.height * _VEC_GRID)
            cy1 = int((r.y1 - pr.y0) / pr.height * _VEC_GRID)
            for cx in range(max(0, cx0), min(_VEC_GRID - 1, cx1) + 1):
                for cy in range(max(0, cy0), min(_VEC_GRID - 1, cy1) + 1):
                    cells.add((cx, cy))
        span = 0.0 if ux1 < ux0 else ((ux1 - ux0) * (uy1 - uy0)) / parea
        out = {"items": items, "cells": len(cells) / float(_VEC_GRID * _VEC_GRID),
               "span": round(span, 4)}
    except Exception:
        return {}                            # no get_drawings / open failure -> honest silence
    _cache_write_json(cf, out)
    return out


def _vector_body(vector) -> bool:
    """True when the page carries a full-page VECTOR DRAWING BODY (not rules, icons or a chart)."""
    if not isinstance(vector, dict) or not vector:
        return False
    return (vector.get("items", 0) >= VEC_MIN_ITEMS
            and vector.get("cells", 0.0) >= VEC_MIN_CELLS
            and vector.get("span", 0.0) >= VEC_MIN_SPAN)


PLAN_FURNITURE_MIN = 2       # distinct drawing annotations needed to call a vector body a PLAN
PLAN_FURNITURE_SPEC_MIN = 4  # ...and to out-argue the spec gate on a half-spec/half-plan spread


PLAN_FIGURE_MIN = 1000            # below this a printed number is a door count, not an area
_PLAN_FIGURE_RE = re.compile(r"\d[\d,. ]*")


def page_area_figures(path: Path, page_index: int, cache: Path | str | None = None) -> set:
    """Every integer >= PLAN_FIGURE_MIN printed on a page, thousands separators tolerated - the
    page's own SCHEDULE OF ACCOMMODATION, read as a set of figures. Empty on any failure."""
    out: set = set()
    for m in _PLAN_FIGURE_RE.finditer(_page_plaintext(path, page_index, cache) or ""):
        raw = m.group(0).replace(",", "").replace(" ", "").rstrip(".")
        if raw.isdigit():
            v = int(raw)
            if v >= PLAN_FIGURE_MIN:
                out.add(v)
    return out


def _plan_rank(kind, furnished, titled, balance, own_figures=False) -> tuple:
    """How two ELIGIBLE plan pages of ONE property are ordered. Higher wins.

    THE PROPERTY'S OWN FIGURES FIRST. A park brochure routinely carries a masterplan per UNIT,
    and the longlist usually holds only one of them - so "is this page a site plan" is not the
    whole question, "is it THIS unit's site plan" is. The page that prints the property's own
    warehouse / office / plot figures is that unit's spread; every other masterplan in the deck
    is a NEIGHBOURING UNIT's, which is a wrong bind in a trace-less slot. Measured: on a
    two-unit park deck the property is the 734,636 sq ft unit, and the OTHER unit's masterplan
    (436,000 sq ft) was winning the slot purely because it also carried a plan title. Absent /
    unknown figures make this component 0 on every page, i.e. exactly the ranking below it.

    Then FURNITURE, above the pixel classifier's own verdict. That order was measured, not
    assumed, and it is the fix for a live WRONG BIND. On a real 13-page deck the property's
    regional DRIVE-TIME MAP classified `kind='plan'` (a pale, line-heavy page) with the best
    white/ink balance on the deck (0.980), while the actual unit site plan - dock doors, 55 m
    yard depth, 595 car parking spaces, gatehouse, schedule of accommodation, all printed ON the
    drawing - classified `kind='map'`. Ranked classifier-first, the drive-time map won the
    trace-less Site Plan slot outright. `kind` is a pixel heuristic about tone and ink; the
    furniture count is direct TEXTUAL evidence that the page is a site layout, and it is the
    stronger signal of the two whenever they disagree.

    Precision is unaffected: this only orders pages that ALREADY passed `_plan_page_eligible`
    (a photo-dominated page, a spec page and an unlabelled location map are all rejected before
    they ever reach a rank), so it can never admit a page - only choose between admissible ones.
    Then the classifier verdict, then a plan TITLE, then the most balanced white/ink page."""
    return (1 if own_figures else 0, 1 if furnished else 0, 1 if kind == "plan" else 0,
            1 if titled else 0, balance)

# Bumped whenever the DETECTOR changes. It keys `best_plan_page_render`'s whole-verdict cache, so a
# cached "no plan on this deck" from an older detector is never served to a newer one.
_PLAN_DETECTOR_SIG = 5


def _plan_page_eligible(kind, sig, has_photo, title_score, is_spec, has_marker,
                        vector=None, furniture=0) -> tuple:
    """Unified plan-page acceptance, shared by page_render_plan (LLM-hint tier) AND
    best_plan_page_render (deterministic fallback) so the two never diverge. Returns (ok, titled).

    Precision gates FIRST (a WRONG image in the authoritative, trace-less Site Plan slot is worse than
    a miss):
      * a SPEC / availability page (>=2 own-line labels) is never a site plan;
      * a page a real PHOTO dominates (`has_photo`: a sizable embedded raster the CLASSIFIER reads as
        continuous-tone 'photo' - measured by tone, not colourfulness, so a low-colour warehouse photo
        is caught) is a property/photo page, never a plan - this REPLACES the old blunt 'image-light
        only' gate, so a designed plan carrying a small logo/legend/north-arrow now qualifies;
      * a continuous-tone whole-page render (classify 'photo') is a photo/aerial, never a plan.
    Then accept on EITHER:
      * VISUAL: the classifier calls it 'plan' inside the balanced-white band - a real drawing needs
        no title or marker (a vector plan often has no extractable text at all); OR
      * TITLE-RESCUE: a page the classifier mis-read as 'map' but that carries BOTH a real plan TITLE
        ('site plan'/'Lageplan'/'plan de masse'/...) AND a to-scale DRAWING MARKER ('scale 1:500'/
        'drawing no'/...). The marker gate (2nd-review fix) is what separates a genuine to-scale
        colour/aerial-overlay site plan (has one) from a photographic AERIAL or a genuine LOCATION /
        overview map (titled "Site Plan" but carrying no drawing furniture) - the latter used to
        wrong-bind. A 'text'/'logo' page (a bare 'SITE PLAN' divider/agenda that merely NAMES a plan)
        is never rescued, and only inside the SAME tight white band.
      * VECTOR (added 2026-08-20, the ONLY route to a vector masterplan): the page's BODY is a
        full-page VECTOR DRAWING (`_vector_body`: >= VEC_MIN_ITEMS primitive path items spread over
        >= VEC_MIN_CELLS of the page grid with >= VEC_MIN_SPAN extent) AND it carries >=
        PLAN_FURNITURE_MIN distinct site-plan DRAWING ANNOTATIONS (`plan_signal.
        plan_furniture_score`: dock doors, yard depth, parking bays, gatehouse, access gates, site
        boundary, accommodation schedule). Both halves are required and neither is sufficient:
        vector density alone also describes a regional road map, a drive-time map and a locator map
        (measured at 2,900 / 2,600 / 6,600 drawing objects on real ones), and the furniture count
        alone also describes a spec sheet. Together they are close to the definition of a labelled
        site layout, which is why this route may overrule the two verdicts that are DEMONSTRABLY
        wrong on a designed plan - the whole-page classifier ('photo' for a full-bleed colour plan,
        'map' for a pale one), the white-balance band (a plan on grey or dark paper) and
        `has_photo` (a decorative gradient tile that classifies 'photo' at 31% of the page) - while
        every one of those verdicts still stands on its own for a page with NO vector body.
        The SPEC gate is the one it may not simply overrule: a half-spec/half-plan SPREAD (the
        contract already tells the interpretation agent to nominate exactly this shape) needs the
        strictly higher PLAN_FURNITURE_SPEC_MIN before the drawing out-argues the spec labels.
    `titled` reports the title hit (used for ranking + near-miss)."""
    titled = (title_score or 0.0) >= PLAN_TITLE_MIN
    white = (sig or {}).get("white", 0.0)
    in_band = 0.15 <= white <= 0.90
    furn = int(furniture or 0)
    vec_body = _vector_body(vector)
    vector_ok = vec_body and furn >= PLAN_FURNITURE_MIN
    if is_spec:
        # a full-page LABELLED vector drawing on a page that also carries >=2 own-line spec labels
        # is the half-spec/half-plan spread, not a spec sheet - but it must clear the higher bar
        return ((vec_body and furn >= PLAN_FURNITURE_SPEC_MIN), titled)
    if has_photo:
        return (vector_ok, titled)
    if kind == "photo":
        return (vector_ok, titled)
    visual_ok = (kind == "plan") and in_band
    title_ok = titled and in_band and (kind == "map") and bool(has_marker)
    return ((visual_ok or title_ok or vector_ok), titled)


def _page_plaintext(path: Path, page_index: int, cache: Path | str | None = None) -> str:
    """Extracted plain text of a page for the plan title-signal + spec gate. '' on any failure or a
    page with no text layer (a vector plan whose labels are drawn strokes yields '' -> the page falls
    through to the pure-VISUAL path unchanged; text never rescues nor blocks such a page on its own).
    Mirrors _rendered_plan_crop's PDF/PPTX doc handling; docs are cached in _DOC_CACHE."""
    try:
        path = Path(path)
        if path.suffix.lower() == ".pptx":
            pdf = soffice_pdf(path, cache)
            if pdf is None:
                return ""
            doc = _get_doc(pdf)
        else:
            doc = _get_doc(path)
        if not (0 <= page_index < doc.page_count):
            return ""
        return doc[page_index].get_text() or ""
    except Exception:
        return ""


def _page_has_dominant_photo(path: Path, page_index: int) -> bool:
    """True if the page carries a sizable EMBEDDED raster the CLASSIFIER reads as continuous-tone
    'photo' - a property/photo page, never a site plan. Uses classify_image (TONE), NOT the
    colourfulness photographic_score (a low-colour warehouse photo scores ~3.7, below any useful
    colourfulness floor), and scans ALL embedded rasters >= 200x200 (not only the >=320x200 hero
    list), so a sub-hero-width portrait photo is still seen. A logo/legend/north-arrow classifies
    'logo'/'plan', not 'photo', so it does NOT disqualify (the point of relaxing the image-light gate).

    FAILS CLOSED (2nd-review fix #3): on any scan/decode error this returns True (assume a photo may
    dominate -> an honest MISS), never False. Every other failure in the plan path falls toward a
    miss; a False here would be the one precision-UNSAFE default, letting a real photo/aerial page
    with a plan caption slip into the trace-less Site Plan slot."""
    try:
        path = Path(path)
        if path.suffix.lower() == ".pptx":
            raws = slide_pictures(path, page_index)
        else:
            raws = page_embedded_images(_get_doc(path), page_index)
        for im in raws:
            if im.get("w", 0) >= 200 and im.get("h", 0) >= 200 and classify_image(im.get("img")) == "photo":
                return True
        return False
    except Exception:
        return True


PLAN_NEAR_BLANK_WHITE = 0.985  # a hinted page whiter than this is effectively empty (no real content)


def page_render_plan(path: Path, page_index: int, budget_kb: int = DEFAULT_BUDGET_KB,
                     cache_dir: Path | str | None = None) -> str | None:
    """The LLM-HINTED plan page (__meta.plan_page): the interpretation sub-agent LOOKED at this page's
    render thumbnail and judged it THIS property's site plan. TRUST that visual judgment - render the
    page, ink-crop it, and bind it. The ONLY deterministic screen here is a BLANK / near-blank page
    (no real content to show); the perceptual 'is it really a site plan' question is confirmed by an
    INDEPENDENT LLM verify (the plan_verify dispatch, consulted in merge), NEVER by a pixel classifier.

    WHY no classifier gate: a full-bleed COLOUR plan, a plan on ONE HALF of a 2-page spread, an
    aerial-overlay site plan all render fine but do NOT match a white-paper line-art signature, so a
    deterministic classify_image / white-band / marker gate mis-judged them - that gate attached 0/4
    real plans on the live Corby decks (e.g. a genuine site plan rejected at white=0.14 vs a 0.15
    cutoff). The classifier is the wrong tool for a PERCEPTION judgment; the LLM already made it.

    Returns a plan data URI or None. Cached per (source, page, budget) under kind='planpage' (a resume
    is byte-deterministic). Degrades to None without Pillow / a renderer (honest null)."""
    if Image is None:
        return None
    if not (isinstance(page_index, int) and not isinstance(page_index, bool) and page_index >= 0):
        return None
    cf = _cache_file(path, page_index, budget_kb, "planpage", cache_dir)
    cached = _cache_read(cf)
    if cached is not None:
        return cached or None
    crop, sig, _kind = _rendered_plan_crop(path, page_index, cache=cache_dir)
    uri = None
    if crop is not None and (sig or {}).get("white", 0.0) <= PLAN_NEAR_BLANK_WHITE:
        uri = to_data_uri(compress(crop, PLAN_MAX_EDGE, budget_kb))  # real content -> trust the pick
    _cache_write(cf, uri)
    return uri


def best_plan_page_render(path: Path, page_nos, budget_kb: int = DEFAULT_BUDGET_KB,
                          cache_dir: Path | str | None = None, near_miss: list | None = None,
                          own_figures: set | None = None) -> tuple:
    """DETERMINISTIC fallback (no LLM hint): over the given (per-property) pages, render+ink-crop+
    classify and pick the most plan-like page via the shared `_plan_page_eligible` predicate - so a
    designed plan carrying a small logo/legend (previously disqualified by the blunt image-light
    gate) now binds, and a full-bleed COLOUR-background plan the classifier mis-reads is rescued when
    its page carries a real plan TITLE - while a real PHOTO page, a SPEC page and an untitled
    location-map never bind (a wrong image in the Site Plan slot is worse than a miss). Prefers a
    TITLED plan, then the most balanced white/ink page. Returns (uri, page_no) or (None, None); the
    page set is SORTED so the result is a pure function of (source, pages); each bound page's URI is
    cached per (source, page, budget) under kind='planpage'. Degrades to (None, None) without Pillow.
    `near_miss` (optional list) collects pages that LOOKED plan-ish but a guard rejected (a positive
    plan signal that did not bind), so a genuinely missed plan surfaces to the Gaps Report."""
    if Image is None:
        return (None, None)
    import plan_signal as _PS
    pages = sorted({p for p in (page_nos or [])
                    if isinstance(p, int) and not isinstance(p, bool) and p >= 0})
    # WHOLE-VERDICT CACHE. The per-page 'planpage' URI cache below was consulted only AFTER the
    # expensive work, and the DECISION (which page won, or that none did) was never persisted at
    # all - so every resumed round re-rendered, re-classified and re-text-scanned every candidate
    # page and threw the answer away. Measured at 40 properties: 22.8s then 23.6s (warm SLOWER
    # than cold) for 0 plans bound, i.e. 87% of a warm merge profile spent re-deriving 'no'. The
    # verdict is a pure function of (deck bytes, sorted page set, budget) - _cache_file already
    # keys on the deck's size+mtime_ns, so a deck edit invalidates it.
    # The NEAR-MISS list is cached WITH the verdict on purpose: it feeds the Gaps Report's
    # "possible site plans not captured" lines, and a resumed run that skipped the scan would
    # otherwise drop them silently - trading a hang for a quiet loss of honesty.
    # The signature below is part of the verdict-cache key on purpose: a cached "no plan on this
    # deck" is only valid for the DETECTOR that produced it, and adding the vector route changes
    # the answer. Without it a warm work dir would keep serving the pre-vector negative for ever -
    # the same class of stale-negative poisoning `_engine_tag` exists to prevent.
    import hashlib as _hl
    # own_figures is part of the verdict-cache key: it changes WHICH eligible page wins, so a
    # verdict computed without it (or with another property's figures) must not be served.
    _ofk = ",".join(str(x) for x in sorted(own_figures or ()))
    _vk = _hl.sha1((",".join(map(str, pages)) + f"|d{_PLAN_DETECTOR_SIG}|f{_ofk}").encode()).hexdigest()[:12]
    _vf = _cache_file(path, 0, budget_kb, f"planverdict{_vk}", cache_dir, ext=".json")
    _v = _cache_read_json(_vf)
    if isinstance(_v, dict) and "page" in _v:
        if near_miss is not None:
            near_miss.extend(_v.get("near_miss") or [])
        _vp = _v.get("page")
        if _vp is None:
            return (None, None)
        _vu = _cache_read(_cache_file(path, _vp, budget_kb, "planpage", cache_dir))
        if _vu:
            return (_vu, _vp)
        # the winning page's URI cache is gone (pruned/corrupt) -> fall through and rebuild
    _nm: list = []   # collected locally so it can be cached, then handed to the caller
    best = None  # ((titled, balance), page_no, uri)
    for pno in pages:
        crop, sig, kind = _rendered_plan_crop(path, pno, cache=cache_dir)
        if crop is None:
            continue
        text = _page_plaintext(path, pno, cache=cache_dir)
        is_spec = _PS.looks_like_spec_page(text)
        title = _PS.plan_title_score(text)
        has_photo = _page_has_dominant_photo(path, pno)
        has_marker = _PS.has_drawing_marker(text)
        vector = page_vector_art(path, pno, cache_dir)
        furniture = _PS.plan_furniture_score(text)
        ok, titled = _plan_page_eligible(kind, sig, has_photo, title, is_spec, has_marker,
                                         vector=vector, furniture=furniture)
        white = sig.get("white", 0.0)
        in_band = 0.15 <= white <= 0.90
        balance = 4.0 * white * (1.0 - white)
        vec_body = _vector_body(vector)
        furnished = furniture >= PLAN_FURNITURE_MIN
        if not ok:
            # NEAR-MISS: a page carrying a positive plan signal (classify 'plan', or a plan title)
            # that a precision guard rejected -> surface it so a real missed plan is visible.
            if is_spec and titled:
                _nm.append({"page": pno,
                            "why": "classified as a spec page but carries a plan title"})
            elif is_spec and kind == "plan" and in_band:
                # a real plan DRAWING on a page that also carries >=2 own-line labels (a legend /
                # title-block) - the spec gate rejected it; surface it so it is not silently lost.
                _nm.append({"page": pno,
                            "why": "classified as a spec page but has the site-plan visual signature"})
            elif vec_body and not furnished and not has_photo:
                # VECTOR NEAR-MISS: the page's whole body IS a full-page vector drawing - the
                # shape a masterplan has - but it carries no site-plan drawing labels to tell it
                # apart from a vector LOCATION MAP / infographic, so the vector route refused it.
                # This is the honest disclosure the precision bar buys: a positive plan-shaped
                # signal that did NOT bind, surfaced rather than silently dropped.
                # It is deliberately ranked ABOVE the generic branch below: on such a page
                # "classified 'plan' outside the white-balance band" is a true but MISLEADING
                # explanation - the band is not what refused it, the missing labels are, and
                # only the second sentence tells a reviewer what to go and look at.
                _nm.append({"page": pno,
                            "why": ("full-page vector drawing (%d path items) but no site-plan "
                                    "labels (dock doors / yard / parking / boundary) to confirm "
                                    "it is a site plan rather than a map"
                                    % int((vector or {}).get("items", 0)))})
            elif not is_spec and (kind == "plan" or titled):
                why = ("a real photo dominates the page" if has_photo
                       else "classified '%s' outside the plan white-balance band" % kind if kind == "plan"
                       else "classified '%s'; a plan title is present but not visually confirmed" % kind if titled
                       else "")
                if why:
                    _nm.append({"page": pno, "why": why})
            continue
        cf = _cache_file(path, pno, budget_kb, "planpage", cache_dir)
        uri = _cache_read(cf)
        if uri is None:
            uri = to_data_uri(compress(crop, PLAN_MAX_EDGE, budget_kb))
            _cache_write(cf, uri)
        # does this page print THIS property's own schedule figures? (see _plan_rank)
        own = bool(own_figures) and bool(set(own_figures) & page_area_figures(path, pno, cache_dir))
        rank = _plan_rank(kind, furnished, titled, balance, own)  # see _plan_rank for the order
        if uri and (best is None or rank > best[0]):
            best = (rank, pno, uri)
    _cache_write_json(_vf, {"page": (best[1] if best else None), "near_miss": _nm})
    if near_miss is not None:
        near_miss.extend(_nm)
    if best is None:
        return (None, None)
    return (best[2], best[1])


_PLACEHOLDER: str | None = None


def placeholder() -> str:
    """A neutral CBRE-green 16:9 placeholder data URI (cached)."""
    global _PLACEHOLDER
    if _PLACEHOLDER is None:
        if Image is None:  # no Pillow: serve the pre-baked placeholder asset (a valid
            # data:image/jpeg URI the chrome + the images gate accept). Integrity-tracked,
            # so a truncated copy is caught by preflight, not shipped.
            _PLACEHOLDER = (Path(__file__).resolve().parent.parent
                            / "assets" / "placeholder.uri").read_text(encoding="utf-8").strip()
        else:
            img = Image.new("RGB", (1280, 720), (0, 63, 45))  # CBRE dark green
            _PLACEHOLDER = to_data_uri(compress(img, 1280, 40))
    return _PLACEHOLDER


# --------------------------------------------------------------------------- #
# PPTX slides - the PDF page ladder's twin. A vision transcription of a slide
# deck carries source_type "pptx" (per the vision contract), but the harvest
# above is PDF-only, so every such record silently degraded to the placeholder
# AND bypassed the placeholder audit (tried_pages was PDF-fed only). These give
# slides the same hero/plan/audit treatment.
# --------------------------------------------------------------------------- #

_PPTX_CACHE: dict[str, object] = {}
_PPTX_SLIDES_CACHE: dict[str, list] = {}   # (#38) enumerated slide list per deck
_SLIDEPIC_CACHE: dict[tuple, list] = {}    # (#39) decoded pictures per (deck, slide)


def _get_pptx(path: Path):
    key = str(Path(path).resolve())
    prs = _PPTX_CACHE.get(key)
    if prs is None:
        from pptx import Presentation
        prs = Presentation(str(path))
        _PPTX_CACHE[key] = prs
    return prs


def _pptx_slides(path: Path) -> list:
    """Enumerated slide list for a deck, memoised: python-pptx materialises every slide
    part on list(prs.slides), and slide_pictures is called once per slide/page across the
    gallery/candidate/audit passes. Cleared with the pptx cache in close_doc_cache. (#38)"""
    key = str(Path(path).resolve())
    sl = _PPTX_SLIDES_CACHE.get(key)
    if sl is None:
        sl = list(_get_pptx(path).slides)
        _PPTX_SLIDES_CACHE[key] = sl
    return sl


def slide_pictures(pptx_path: Path, slide_index: int) -> list[dict]:
    """All decodable raster pictures on one slide, largest first (undecodable
    WMF/EMF vectors are skipped, never abort the harvest).

    A picture is ANY shape that carries an image blob - not only shape_type PICTURE.
    Filtering on PICTURE alone loses every photo dropped into a picture PLACEHOLDER
    (shape_type 14, which is how a slide built from a corporate template holds its
    imagery) and every picture nested inside a GROUP. That is not hypothetical: a
    CBRE options deck put the aerial in a placeholder and the site plan as a loose
    picture on each slide, so the harvest returned the PLAN for every option, and
    the two slides that placed BOTH images in placeholders shipped a blank
    placeholder card. `shape.image` reads a placeholder-held picture perfectly, so
    the fix is to ask for the blob rather than to trust the shape type."""
    key = (str(Path(pptx_path).resolve()), slide_index)
    hit = _SLIDEPIC_CACHE.get(key)
    if hit is not None:
        return hit
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        slides = _pptx_slides(pptx_path)
        if not (0 <= slide_index < len(slides)):
            return []
        out = []
        seen: set = set()          # dedupe by image bytes: one picture, one candidate

        def _collect(shapes, depth=0):
            for shape in shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP and depth < 4:
                    try:
                        _collect(shape.shapes, depth + 1)
                    except Exception:
                        pass
                    continue
                try:
                    blob = shape.image.blob      # works for PICTURE *and* picture PLACEHOLDER
                except Exception:
                    continue                     # no image on this shape (table, text box, ...)
                sig = hashlib.sha1(blob).hexdigest()
                if sig in seen:
                    continue
                seen.add(sig)
                try:
                    img = _open(blob)
                except Exception:
                    img = None
                if img is not None:
                    w, h = img.size
                    out.append({"img": img, "w": w, "h": h, "area": w * h})

        _collect(slides[slide_index].shapes)
        out.sort(key=lambda d: -d["area"])
        _SLIDEPIC_CACHE[key] = out
        return out
    except Exception:
        return []


_SOFFICE: object = False  # False = not probed yet; then str | None


def _find_soffice() -> str | None:
    """Locate a headless LibreOffice (the only reliable slide renderer -
    python-pptx cannot rasterise). Memoised; None when absent."""
    global _SOFFICE
    if _SOFFICE is not False:
        return _SOFFICE  # type: ignore[return-value]
    import shutil
    cand = shutil.which("soffice") or shutil.which("libreoffice")
    if not cand and sys.platform == "win32":
        for p in (Path(r"C:\Program Files\LibreOffice\program\soffice.exe"),
                  Path(r"C:\Program Files (x86)\LibreOffice\program\soffice.exe")):
            if p.exists():
                cand = str(p)
                break
    _SOFFICE = cand
    return cand


def soffice_pdf(src: Path, cache_dir: Path | str | None) -> Path | None:
    """Convert a slide deck to PDF via headless LibreOffice (slides map 1:1 to
    pages), memoised on disk per (name, size, mtime) so one deck converts once
    per run history. None when LibreOffice is absent or conversion fails -
    callers fall back to embedded slide pictures."""
    exe = _find_soffice()
    if exe is None:
        return None
    src = Path(src)
    try:
        st = src.stat()
    except OSError:
        return None
    import hashlib
    import subprocess
    import tempfile
    cdir = Path(cache_dir) if cache_dir else Path(tempfile.gettempdir()) / "cbre_longlist_soffice"
    try:
        cdir.mkdir(parents=True, exist_ok=True)
    except Exception:
        return None
    key = hashlib.sha1(f"soffice|{src.name}|{st.st_size}|{st.st_mtime_ns}".encode()).hexdigest()[:16]
    target = cdir / f"{src.stem}.{key}.pdf"
    if target.exists() and target.stat().st_size > 0:
        return target
    try:
        with tempfile.TemporaryDirectory(prefix="soffice_") as td:
            # a private user profile so a running desktop LibreOffice (or a
            # parallel conversion) can never lock the conversion out
            profile = Path(td) / "profile"
            profile.mkdir()
            # TIMEOUT must be SHORTER than the shell window, not longer than the whole run.
            # At 180s a single conversion could outlive the ~40-45s cap, so the shell was killed
            # mid-convert, the PDF only lands on success, and the next round started from zero -
            # an unbounded, invisible loop. It is also the one prewarm unit that can hold the
            # process at interpreter exit. Override with CBRE_SOFFICE_TIMEOUT when converting a
            # genuinely huge deck outside a capped sandbox.
            try:
                _soffice_timeout = float(os.environ.get("CBRE_SOFFICE_TIMEOUT") or 0) or 25.0
            except ValueError:
                _soffice_timeout = 25.0
            subprocess.run(
                [exe, "--headless", "--norestore",
                 f"-env:UserInstallation={profile.as_uri()}",
                 "--convert-to", "pdf", "--outdir", td, str(src)],
                check=True, capture_output=True, timeout=_soffice_timeout)
            produced = Path(td) / (src.stem + ".pdf")
            if not produced.exists() or produced.stat().st_size == 0:
                return None
            import shutil
            shutil.move(str(produced), str(target))
        return target
    except Exception:
        return None


def slide_hero_and_plan(pptx_path: Path, slide_index: int,
                        budget_kb: int = DEFAULT_BUDGET_KB,
                        cache_dir: Path | str | None = None) -> tuple[str | None, str | None]:
    """(hero_uri, plan_uri) for one slide, same combination rules as the PDF
    page ladder. PREFERRED: LibreOffice renders the deck to PDF (slides map 1:1
    to pages) and the FULL page ladder runs on the converted page - embedded
    tiers, geometry crops, the plan picker, everything. Fallback (no
    LibreOffice): the slide's decodable embedded pictures - the best
    photo-scoring picture is the hero, a plan-signature picture fills the plan
    slot, and a plan-only slide promotes the plan to hero. Memoised on disk per
    (source, slide, budget) like the PDF path."""
    if Image is None:
        return (None, None)  # no Pillow: no slide hero/plan; merge fills the placeholder
    cf_h = _cache_file(pptx_path, slide_index, budget_kb, "slide_hero", cache_dir)
    cf_p = _cache_file(pptx_path, slide_index, budget_kb, "slide_plan", cache_dir)
    ch, cp = _cache_read(cf_h), _cache_read(cf_p)
    if ch is not None and cp is not None:
        return (ch or None), (cp or None)

    hero = plan = None
    pdf = soffice_pdf(pptx_path, cache_dir)
    if pdf is not None:
        try:
            hero, plan = page_hero_and_plan(pdf, slide_index, budget_kb, cache_dir=cache_dir)
        except Exception:
            hero = plan = None
    if hero is None and plan is None:
        pics = slide_pictures(pptx_path, slide_index)
        best, best_score = None, -1.0
        for im in pics:
            if im["w"] >= MIN_HERO_W and im["h"] >= MIN_HERO_H:
                sc = photographic_score(im["img"])
                if sc > best_score:
                    best, best_score = im["img"], sc
        if best is not None and best_score >= MODEST_PHOTO:
            hero = to_data_uri(compress(best, HERO_MAX_EDGE, budget_kb))
        plan_cands = []
        for im in pics:
            if im["img"] is best and hero is not None:
                continue  # the hero photo is never also the plan
            if im["w"] < MIN_PLAN_W or im["h"] < MIN_PLAN_H:
                continue
            try:
                white, balance = _crop_stats(im["img"].convert("RGB"))
            except Exception:
                continue
            if 0.15 <= white <= 0.90:  # the plan signature (page_plan's bounds)
                plan_cands.append((balance * math.sqrt(im["w"] * im["h"]), im["img"]))
        if plan_cands:
            plan_cands.sort(key=lambda t: -t[0])
            plan = to_data_uri(compress(plan_cands[0][1], PLAN_MAX_EDGE, budget_kb))
        if hero is None and plan is not None:
            hero = plan  # plan-only slide: the plan IS the hero (broker rule)
        elif hero is None and best is not None:
            # sub-modest but real picture still beats the placeholder (the PDF
            # ladder's renderer-less tier C does the same)
            hero = to_data_uri(compress(best, HERO_MAX_EDGE, budget_kb))

    _cache_write(cf_h, hero)
    _cache_write(cf_p, plan)
    return hero, plan


def slide_image_audit(pptx_path: Path, slide_index: int, out_dir: Path, tag: str,
                      cache_dir: Path | str | None = None) -> list[str]:
    """page_image_audit's PPTX twin: dump every embedded slide picture (plus the
    LibreOffice-rendered slide, when a renderer is available) as labelled
    thumbnails, so a slide-sourced placeholder is a reviewed conclusion too."""
    if Image is None:
        return []  # no Pillow: no audit montage
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    # RESUME SHORT-CIRCUIT - see page_image_audit. Worth more here: a miss re-runs soffice_pdf.
    existing = sorted(str(p.resolve()) for p in out_dir.glob(f"{tag}_s{slide_index + 1}_*.png"))
    if existing:
        return existing
    files: list[str] = []

    def _save(img, kind: str, idx: int):
        try:
            im = img.convert("RGB")
            w, h = im.size
            im.thumbnail((480, 480))
            name = f"{tag}_s{slide_index + 1}_{kind}{idx}_{w}x{h}.png"
            _atomic_save_png(im, out_dir / name)
            files.append(str((out_dir / name).resolve()))
        except Exception:
            pass

    for i, im in enumerate(slide_pictures(pptx_path, slide_index), start=1):
        _save(im["img"], "picture", i)
    try:
        pdf = soffice_pdf(pptx_path, cache_dir)
        if pdf is not None:
            _save(page_raster(_get_doc(pdf), slide_index), "sliderender", 1)
    except Exception:
        pass
    return files
