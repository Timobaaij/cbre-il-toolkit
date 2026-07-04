#!/usr/bin/env python3
"""vision_prep.py - Stage 1 vision fallback (DETERMINISTIC half).

When a brochure is image/vector-only (no extractable text, so extract_pdf /
extract_pptx return 0 records), this rasterises the unreadable pages to PNGs and
emits a manifest. It does NOT transcribe - transcription is an agentic step: the
orchestrator dispatches an isolated VISION sub-agent that reads each PNG and
writes candidate records (templates/record_schema.json) to
work/extract/<region>_vision.json, which merge.py then folds in like any other
extractor output. See reference/vision-fallback.md for the sub-agent contract.

A page is sent to vision only when text parsing cannot read it (own-line AND
inline label parsing both find <2 labels) - so a normal text deck produces an
empty manifest and nothing is rasterised. PDF pages render via images.page_raster;
PPTX falls back to the slide's largest embedded picture (a purely vector/text-less
slide cannot be rasterised by python-pptx and is reported for manual export).

CLI:
  python vision_prep.py <file.pdf|.pptx> --region R --country C --out-dir work/vision
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
try:
    import fitz  # PyMuPDF
except Exception:  # sandbox without PyMuPDF: pypdfium2/pdfplumber shim
    import fitz_shim as fitz
try:
    fitz.TOOLS.mupdf_display_errors(False)
except Exception:
    pass
import images as IMG
import extract_pdf as P


def _needs_vision(text: str) -> bool:
    """True when neither own-line nor inline text parsing can read this page
    (so it is an image/scan/vector page that only a vision pass can transcribe)."""
    return len(P._find_labels(text)) < 2 and len(P._find_inline_labels(text)) < 2


def _largest_slide_picture(slide):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    best = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img = IMG._open(shape.image.blob)
            except Exception:
                img = None
            if img is not None and (best is None or img.width * img.height > best[0]):
                best = (img.width * img.height, img)
    return best[1] if best else None


def _stamp_current(path: Path, out_dir: Path, dpi: int) -> bool:
    """True when out_dir's stamp for this source matches its current bytes + dpi,
    so existing page PNGs can be REUSED. Under a short shell cap (Cowork ~45s) a
    kill mid-rasterisation used to restart from page one on every call - the
    single biggest wall-clock sink of a real run. With the stamp, a resumed call
    renders only the pages whose PNGs are still missing."""
    import json as _json
    stamp = out_dir / f"{path.stem}.stamp.json"
    st = path.stat()
    cur = {"size": st.st_size, "mtime_ns": st.st_mtime_ns, "dpi": dpi}
    try:
        if stamp.exists() and _json.loads(stamp.read_text(encoding="utf-8")) == cur:
            return True
    except Exception:
        pass
    # source or dpi changed: stale PNGs must not be reused - drop them, restamp
    # (_p*.png = PDF pages, _s*.png = PPTX slides)
    for pat in (f"{path.stem}_p*.png", f"{path.stem}_s*.png"):
        for old in out_dir.glob(pat):
            try:
                old.unlink()
            except Exception:
                pass
    try:
        stamp.write_text(_json.dumps(cur), encoding="utf-8")
    except Exception:
        pass
    return False


def prepare(path: Path, region: str, country: str, out_dir, dpi: int = 180,
            force: bool = False) -> dict:
    """Rasterise the pages of one brochure and return a manifest entry
    {source_file, source_type, region, country, pages:[{page_no, locator, image, reason}]}.

    page_no in each entry is the CANONICAL 0-BASED page index: the transcription
    agent must copy it VERBATIM into __meta.page_no (the PNG filename suffix is
    1-based for humans - deriving page_no from it binds the hero to the
    neighbouring property). dpi defaults to 180: 120 left small-font SPEC TABLES and
    KPI badges illegible (a live run had to re-render a dense flyer by hand, and caused
    clear-height / area misreads), so the first pass now renders at a resolution those
    are readable at. The extra render cost is bounded per page and the stage is per-page
    RESUMABLE (an existing PNG with a matching source stamp is reused, never re-rendered),
    so a deck that crosses the shell cap simply continues on the next run rather than
    failing - the cap-vs-legibility trade is paid in resume cycles, not lost work.

    Default: only pages with no extractable labels (image/scan/vector). With
    force=True: rasterise EVERY page - used when the orchestrator routes a file
    to vision (it yielded no usable records, or a thin/low-quality parse), so a
    page that has text but whose record was dropped/thin is still sent to
    vision instead of being silently lost."""
    path = Path(path)
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    st = path.suffix.lower().lstrip(".")
    pages: list[dict] = []

    if st == "pdf":
        reusable = _stamp_current(path, out_dir, dpi)
        doc = fitz.open(path)
        for pno in range(doc.page_count):
            name = f"{path.stem}_p{pno + 1}.png"
            target = out_dir / name
            reused = reusable and target.exists() and target.stat().st_size > 0
            # P0-7: on a warm-resume page in a FORCED pass, SKIP the costly get_text()
            # (pdfplumber on the no-PyMuPDF shim is ~27s/deck and was re-paid on EVERY
            # resume) - it only feeds the advisory `reason`, never the merge. The
            # routing (non-force) path still reads text to decide whether to rasterise.
            if reused and force:
                pages.append({"page_no": pno, "locator": f"page {pno + 1}",
                              "image": str(target.resolve()),
                              "reason": "text present but no usable record parsed - transcribe the whole page"})
                continue
            try:
                text = doc[pno].get_text()
            except Exception:
                text = ""
            needs = _needs_vision(text)
            if not force and not needs:
                continue
            reason = ("no extractable text/labels (image/scan/vector page)" if needs
                      else "text present but no usable record parsed - transcribe the whole page")
            if reused:
                pages.append({"page_no": pno, "locator": f"page {pno + 1}",
                              "image": str(target.resolve()), "reason": reason})
                continue  # already rendered on a prior (killed) pass
            try:
                img = IMG.page_raster(doc, pno, dpi=dpi)
                img.convert("RGB").save(target, "PNG")
                pages.append({"page_no": pno, "locator": f"page {pno + 1}",
                              "image": str(target.resolve()), "reason": reason})
            except Exception as e:
                pages.append({"page_no": pno, "locator": f"page {pno + 1}", "image": None,
                              "reason": f"could not rasterise: {e}"})
        doc.close()

    elif st == "pptx":
        # slide texts drive the needs-vision test in BOTH tiers below
        slide_texts: list[str] = []
        try:
            from pptx import Presentation
            import extract_pptx as PPTX
            prs = Presentation(str(path))
            slide_texts = PPTX.slide_texts(prs)   # one bad slide -> '' (prs kept), not a dead deck (#19)
        except Exception:
            prs = None
        # PREFERRED tier: LibreOffice renders the deck to PDF (slides map 1:1 to
        # pages) and every slide rasterises like a PDF page - including the
        # vector/text-only slides python-pptx cannot render at all (a real deck
        # produced 0 PNGs and stranded the whole region without this).
        pdf = IMG.soffice_pdf(path, out_dir)
        if pdf is not None:
            reusable = _stamp_current(path, out_dir, dpi)
            doc = fitz.open(str(pdf))
            for i in range(doc.page_count):
                text = slide_texts[i] if i < len(slide_texts) else ""
                needs = _needs_vision(text)
                if not force and not needs:
                    continue
                reason = ("no extractable text/labels (image/vector slide)" if needs
                          else "text present but no usable record parsed - transcribe the whole slide")
                name = f"{path.stem}_s{i + 1}.png"
                target = out_dir / name
                if reusable and target.exists() and target.stat().st_size > 0:
                    pages.append({"page_no": i, "locator": f"slide {i + 1}",
                                  "image": str(target.resolve()), "reason": reason})
                    continue  # already rendered on a prior (killed) pass
                try:
                    img = IMG.page_raster(doc, i, dpi=dpi)
                    img.convert("RGB").save(target, "PNG")
                    pages.append({"page_no": i, "locator": f"slide {i + 1}",
                                  "image": str(target.resolve()), "reason": reason})
                except Exception as e:
                    pages.append({"page_no": i, "locator": f"slide {i + 1}", "image": None,
                                  "reason": f"could not rasterise: {e}"})
            doc.close()
            return {"source_file": path.name, "source_type": st, "region": region,
                    "country": country, "pages": pages}
        # FALLBACK tier (no LibreOffice): each slide's largest embedded picture
        if prs is None:
            return {"source_file": path.name, "source_type": st, "region": region,
                    "country": country, "pages": [],
                    "note": "python-pptx unavailable and no LibreOffice - cannot rasterise pptx"}
        for i, slide in enumerate(prs.slides):
            # reuse the already-harvested per-slide texts (guarded); a bad slide's '' reads
            # as needs-vision, so it still rasterises rather than raising out of prepare() (#19)
            if not force and not _needs_vision(slide_texts[i] if i < len(slide_texts) else ""):
                continue
            img = _largest_slide_picture(slide)
            if img is None:
                pages.append({"page_no": i, "locator": f"slide {i + 1}", "image": None,
                              "reason": "vector/text-less slide - python-pptx cannot rasterise "
                                        "(install LibreOffice for full slide rendering, or export it manually)"})
                continue
            name = f"{path.stem}_s{i + 1}.png"
            img.convert("RGB").save(out_dir / name, "PNG")
            pages.append({"page_no": i, "locator": f"slide {i + 1}",
                          "image": str((out_dir / name).resolve()),
                          "reason": "no extractable text/labels (image slide)"})

    return {"source_file": path.name, "source_type": st, "region": region,
            "country": country, "pages": pages}


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("file")
    ap.add_argument("--region", required=True)
    ap.add_argument("--country", required=True)
    ap.add_argument("--out-dir", default="vision")
    ap.add_argument("--dpi", type=int, default=180,
                    help="raster dpi (180: legible small-font spec TABLES + KPI badges on the "
                         "first pass; the stage is per-page resumable so the extra cost is bounded)")
    ap.add_argument("--force", action="store_true",
                    help="rasterise EVERY page (region routed to vision), not only label-less pages")
    args = ap.parse_args()
    sys.stdout.reconfigure(encoding="utf-8")
    ent = prepare(Path(args.file), args.region, args.country, args.out_dir, args.dpi, force=args.force)
    print(json.dumps(ent, ensure_ascii=False, indent=2))
    print(f"OK {len(ent['pages'])} page(s) need vision transcription -> {args.out_dir}")


if __name__ == "__main__":
    main()
